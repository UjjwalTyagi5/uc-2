import io
import structlog
from pptx import Presentation
from file_classifier.extractors.base import BaseExtractor, ExtractionResult

logger = structlog.get_logger()

MAX_SLIDES = 15


class PptxExtractor(BaseExtractor):
    def extract(self, file_bytes: bytes, filename: str) -> ExtractionResult:
        logger.info("Extracting PPTX content", filename=filename)
        buf = io.BytesIO(file_bytes)
        prs = Presentation(buf)

        slides = list(prs.slides)
        total_slides = len(slides)
        parts = []
        parts.append(f"## Presentation: {total_slides} slides\n")

        for i, slide in enumerate(slides[:MAX_SLIDES]):
            slide_parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)

                if shape.has_table:
                    tbl = shape.table
                    rows = []
                    for row in tbl.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        header = rows[0]
                        table_str = " | ".join(header) + "\n"
                        table_str += " | ".join("---" for _ in header) + "\n"
                        for row in rows[1:20]:
                            table_str += " | ".join(row) + "\n"
                        slide_parts.append(table_str)

            if slide_parts:
                parts.append(f"### Slide {i + 1}")
                parts.append("\n".join(slide_parts))
                parts.append("")

        text = "\n".join(parts)

        return ExtractionResult(
            text_content=text,
            metadata={"total_slides": total_slides, "slides_extracted": min(total_slides, MAX_SLIDES)},
        )

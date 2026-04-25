"""Convert quotation files into LLM-consumable content.

Strategy (hybrid):
  PDF / images  → render pages as PNG, return base-64 list
  XLSX / XLS    → parse with openpyxl / xlrd, return markdown text
  DOCX          → python-docx text extraction
  DOC           → OLE text extraction fallback, else render via fitz
  PPTX / PPT    → render slides as images via fitz
  TXT / CSV     → read as text
  MSG           → extract body text
  HTM / HTML    → strip tags, return text
"""

from __future__ import annotations

import base64
import html as html_mod
import io
import pathlib
import re
from typing import Optional

from loguru import logger

from .models import DocumentContent

_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp"}
_TEXT_EXTS = {".txt", ".csv"}
_HTML_EXTS = {".htm", ".html"}


def load_document(
    file_path: str | pathlib.Path,
    max_pages: int = 20,
    work_dir: Optional[pathlib.Path] = None,
) -> DocumentContent:
    """Load a file and return an LLM-ready :class:`DocumentContent`.

    Parameters
    ----------
    work_dir:
        Directory under which any intermediate files (e.g. rendered page
        images) should be written.  Defaults to the source file's own
        directory so that all temp output stays inside the per-RAS work
        folder and is cleaned up together with it.  Currently all rendering
        is done in-memory; this parameter exists to guarantee the right
        location is used if disk-backed rendering is ever added.
    """
    fp = pathlib.Path(file_path)
    ext = fp.suffix.lower()
    # Use the file's own directory as the default so temp files always stay
    # inside the per-RAS work folder.
    wdir = work_dir or fp.parent
    logger.debug("Loading document: {} (ext={}, work_dir={})", fp.name, ext, wdir)

    if ext == ".pdf":
        return _load_pdf(fp, max_pages, wdir)
    if ext in (".xlsx", ".xls"):
        return _load_spreadsheet(fp, wdir)
    if ext == ".docx":
        return _load_docx(fp, wdir)
    if ext == ".doc":
        return _load_doc_legacy(fp, max_pages, wdir)
    if ext in (".pptx", ".ppt"):
        return _load_presentation(fp, max_pages, wdir)
    if ext in _IMAGE_EXTS:
        return _load_image(fp)
    if ext in _TEXT_EXTS:
        return _load_text(fp)
    if ext in _HTML_EXTS:
        return _load_html(fp)
    if ext == ".msg":
        return _load_msg(fp)

    logger.warning("Unsupported extension '{}' — attempting fitz render", ext)
    return _try_fitz_render(fp, max_pages, wdir)


# ── PDF ──


def _load_pdf(fp: pathlib.Path, max_pages: int, work_dir: pathlib.Path) -> DocumentContent:
    """Load a PDF returning BOTH extracted text and page images.

    A PDF can be mixed — some pages have embedded text, others are scanned.
    Sending both gives the vision-capable LLM the full picture in one call:
    extracted text is cheaper to process while images cover scanned sections.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(str(fp))
    page_count = len(doc)
    total = min(page_count, max_pages)

    text_pages: list[str] = []
    images: list[str] = []

    for i in range(total):
        page = doc[i]
        page_text = page.get_text("text").strip()
        if page_text:
            text_pages.append(page_text)
        pix = page.get_pixmap(dpi=200)
        images.append(base64.b64encode(pix.tobytes("png")).decode())

    doc.close()

    text = "\n\n".join(text_pages) if text_pages else None
    avg_chars = sum(len(p) for p in text_pages) / max(len(text_pages), 1) if text_pages else 0
    logger.debug(
        "PDF loaded: {} page(s) — text avg {:.0f} chars/page, {} image(s) (work_dir={})",
        total, avg_chars, len(images), work_dir,
    )
    return DocumentContent(text=text, images=images, source_path=str(fp), page_count=total)


# ── Spreadsheets ──


def _load_spreadsheet(fp: pathlib.Path, work_dir: pathlib.Path) -> DocumentContent:
    ext = fp.suffix.lower()
    parts: list[str] = []

    if ext == ".xlsx":
        parts = _parse_xlsx(fp)
    elif ext == ".xls":
        parts = _parse_xls(fp)

    if not parts or all(p.strip() == "" for p in parts):
        logger.info("Spreadsheet text empty — falling back to image render")
        return _try_fitz_render(fp, max_pages=20, work_dir=work_dir)

    text = "\n\n".join(parts)
    return DocumentContent(text=text, source_path=str(fp), page_count=1)


def _parse_xlsx(fp: pathlib.Path) -> list[str]:
    from openpyxl import load_workbook

    wb = load_workbook(str(fp), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        lines: list[str] = [f"### Sheet: {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append("| " + " | ".join(cells) + " |")
        if len(lines) > 1:
            parts.append("\n".join(lines))
    wb.close()
    return parts


def _parse_xls(fp: pathlib.Path) -> list[str]:
    import xlrd

    wb = xlrd.open_workbook(str(fp))
    parts: list[str] = []
    for sheet in wb.sheets():
        lines: list[str] = [f"### Sheet: {sheet.name}"]
        for rx in range(sheet.nrows):
            cells = [
                str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)
            ]
            if any(c.strip() for c in cells):
                lines.append("| " + " | ".join(cells) + " |")
        if len(lines) > 1:
            parts.append("\n".join(lines))
    return parts


# ── DOCX ──


def _load_docx(fp: pathlib.Path, work_dir: pathlib.Path) -> DocumentContent:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(fp))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    table_parts: list[str] = []
    for tbl in doc.tables:
        rows: list[str] = []
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            table_parts.append("\n".join(rows))

    text = "\n".join(paragraphs)
    if table_parts:
        text += "\n\n" + "\n\n".join(table_parts)

    if not text.strip():
        logger.info("DOCX text empty — falling back to image render")
        return _try_fitz_render(fp, max_pages=20, work_dir=work_dir)

    return DocumentContent(text=text, source_path=str(fp), page_count=1)


# ── DOC (legacy OLE) ──


def _load_doc_legacy(
    fp: pathlib.Path,
    max_pages: int,
    work_dir: pathlib.Path,
) -> DocumentContent:
    try:
        return _try_fitz_render(fp, max_pages, work_dir)
    except Exception:
        logger.debug("fitz failed on .doc — trying OLE text extraction")

    try:
        import olefile

        ole = olefile.OleFileIO(str(fp))
        if ole.exists("WordDocument"):
            stream = ole.openstream("WordDocument").read()
            text = stream.decode("utf-8", errors="ignore")
            text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
            if text.strip():
                return DocumentContent(
                    text=text, source_path=str(fp), page_count=1
                )
        ole.close()
    except Exception:
        pass

    logger.warning("Could not extract content from .doc file: {}", fp.name)
    return DocumentContent(text="[Document could not be read]", source_path=str(fp))


# ── PPT / PPTX ──


def _load_presentation(
    fp: pathlib.Path,
    max_pages: int,
    work_dir: pathlib.Path,
) -> DocumentContent:
    ext = fp.suffix.lower()
    if ext == ".pptx":
        text = _parse_pptx_text(fp)
        if text and text.strip():
            return DocumentContent(text=text, source_path=str(fp), page_count=1)
    return _try_fitz_render(fp, max_pages, work_dir)


def _parse_pptx_text(fp: pathlib.Path) -> Optional[str]:
    try:
        from pptx import Presentation

        prs = Presentation(str(fp))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = [f"--- Slide {idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        slide_text.append("| " + " | ".join(cells) + " |")
            if len(slide_text) > 1:
                parts.append("\n".join(slide_text))
        return "\n\n".join(parts) if parts else None
    except Exception:
        return None


# ── Images (JPG, PNG, TIF, …) ──


def _load_image(fp: pathlib.Path) -> DocumentContent:
    raw = fp.read_bytes()

    ext = fp.suffix.lower()
    if ext in (".tif", ".tiff"):
        from PIL import Image

        img = Image.open(io.BytesIO(raw))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        raw = buf.getvalue()

    b64 = base64.b64encode(raw).decode()
    return DocumentContent(images=[b64], source_path=str(fp), page_count=1)


# ── Plain text ──


def _load_text(fp: pathlib.Path) -> DocumentContent:
    text = fp.read_text(encoding="utf-8", errors="replace")
    return DocumentContent(text=text, source_path=str(fp), page_count=1)


# ── HTML ──


def _load_html(fp: pathlib.Path) -> DocumentContent:
    raw = fp.read_text(encoding="utf-8", errors="replace")
    text = html_mod.unescape(re.sub(r"<[^>]+>", " ", raw))
    text = re.sub(r"\s+", " ", text).strip()
    return DocumentContent(text=text, source_path=str(fp), page_count=1)


# ── MSG (Outlook email) ──


def _load_msg(fp: pathlib.Path) -> DocumentContent:
    try:
        import extract_msg

        msg = extract_msg.Message(str(fp))
        parts = [f"Subject: {msg.subject or ''}"]
        if msg.body:
            parts.append(msg.body)
        msg.close()
        return DocumentContent(
            text="\n\n".join(parts), source_path=str(fp), page_count=1
        )
    except Exception:
        logger.warning("Failed to parse .msg file: {}", fp.name)
        return DocumentContent(
            text="[Email could not be read]", source_path=str(fp)
        )


# ── Fallback: render with fitz ──


def _try_fitz_render(
    fp: pathlib.Path,
    max_pages: int,
    work_dir: pathlib.Path,
) -> DocumentContent:
    import fitz

    doc = fitz.open(str(fp))
    page_count = len(doc)
    images: list[str] = []
    total = min(page_count, max_pages)
    for i in range(total):
        pix = doc[i].get_pixmap(dpi=200)
        # Render in-memory as base64; work_dir is available here if disk
        # caching is ever needed (files would land inside the per-RAS folder).
        images.append(base64.b64encode(pix.tobytes("png")).decode())
    doc.close()
    if not images:
        return DocumentContent(
            text="[Document could not be rendered]", source_path=str(fp)
        )
    return DocumentContent(images=images, source_path=str(fp), page_count=total)

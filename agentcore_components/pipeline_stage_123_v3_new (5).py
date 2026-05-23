"""Full Pipeline V2 (Stages 1-8) — AgentCore custom-code component.

Standalone single-file V2 of the procurement benchmarking pipeline.
Identical to V1 (pipeline_stage_123.py) for Stages 1-4 and 8; Stages 5-7
swap in LLM-driven category classification + brand-free embedding text
+ 3-stage retrieval (SQL → Pinecone within pool → LLM relevance rank).

Designed to be uploaded to MiCore as its own Custom Component — no
imports from sibling project files. The Node registers as
"Full Pipeline V2 (LLM Category + Smart Benchmark)" alongside V1.

Run sql/migrations/quotation_extracted_items_add_classification_cols.sql
once on the target DB before using this Node.

Logging
-------
This file does NOT call logger.add(...) or open a log file. loguru uses its
default stderr sink, and Node-level events are also surfaced via self._safe_log()
so they appear in the AgentCore canvas log panel in real time. No log
directory is created on disk by this component.

Verbosity:
  - logger.info  : per-stage progress, page/image counts, LLM call summary
  - logger.warning : recoverable issues (page truncation, transient retries,
                     extraction failures for one file, no quotation found)
  - logger.error : unrecoverable per-PR failures (DB write, render error)
                   captured with stack trace via logger.opt(exception=True)
  - logger.debug : low-level diagnostics (suppressed at default level)
"""
from __future__ import annotations

import json
import random
import re
import time
from dataclasses import dataclass, field
from decimal import Decimal
from difflib import SequenceMatcher
from typing import Any, Optional

from loguru import logger

# Pre-import Azure modules for blob operations
try:
    from azure.identity import DefaultAzureCredential  # noqa: F401
    from azure.storage.blob import BlobServiceClient    # noqa: F401
except ImportError:
    pass

# Module-level cached Azure credential. We keep DefaultAzureCredential so
# the same code works on dev VMs (az login → AzureCliCredential) and in
# production AgentCore hosts (system-assigned Managed Identity).
# Caching matters: without it, every blob op constructs a fresh credential
# and re-runs the whole chain, which hammers IMDS (169.254.169.254) and
# causes intermittent "Failed to invoke the Azure CLI" errors when MI
# briefly fails and the chain falls through. One cached instance reuses
# its internal token cache across all calls.
#
# Token refresh resilience: every ~50-60 min the SDK refreshes the bearer
# token. If IMDS is briefly unreachable at that exact moment, the cached
# credential starts returning auth errors on every call. _with_az_retry()
# catches those, forces a fresh credential build, and retries the op once
# — recovering transparently from IMDS hiccups instead of cascading the
# failure up to the pipeline.
import threading as _threading_for_cred
_AZ_CRED = None
_AZ_CRED_LOCK = _threading_for_cred.Lock()

def _get_az_credential(force_rebuild: bool = False):
    """Return the shared DefaultAzureCredential. Pass force_rebuild=True to
    discard the cached instance and create a fresh one — used by the retry
    wrapper after a transient IMDS / token-refresh failure."""
    global _AZ_CRED
    if force_rebuild:
        with _AZ_CRED_LOCK:
            old = _AZ_CRED
            _AZ_CRED = None
            if old is not None:
                try:
                    if hasattr(old, "close"):
                        old.close()
                except Exception:
                    pass
    if _AZ_CRED is not None:
        return _AZ_CRED
    with _AZ_CRED_LOCK:
        if _AZ_CRED is None:
            from azure.identity import DefaultAzureCredential
            _AZ_CRED = DefaultAzureCredential(
                exclude_environment_credential=True,
                exclude_interactive_browser_credential=True,
            )
    return _AZ_CRED


def _is_az_auth_error(exc: BaseException) -> bool:
    """True if exc looks like an Azure credential / token-acquisition failure
    (worth retrying with a fresh credential) vs a real blob-storage error
    (which a retry would not fix)."""
    try:
        from azure.core.exceptions import ClientAuthenticationError
        if isinstance(exc, ClientAuthenticationError):
            return True
    except Exception:
        pass
    msg = str(exc).lower()
    return any(s in msg for s in (
        "failed to invoke the azure cli",
        "managedidentitycredential",
        "defaultazurecredential failed",
        "no credential in this chain",
        "credentialunavailable",
        "azurecli",
    ))


def _with_az_retry(fn, *args, on_rebuild=None, max_attempts: int = 2, **kwargs):
    """Run fn(*args, **kwargs); on a transient Azure auth failure, force-rebuild
    the cached credential, invoke on_rebuild() (typically to invalidate any
    BlobServiceClient that was created with the now-stale credential), and
    retry once. Non-auth exceptions propagate immediately."""
    last_exc: BaseException | None = None
    for attempt in range(max_attempts):
        try:
            return fn(*args, **kwargs)
        except Exception as exc:
            if not _is_az_auth_error(exc) or attempt == max_attempts - 1:
                raise
            last_exc = exc
            try:
                logger.warning(
                    f"Azure auth transient failure (attempt {attempt + 1}/{max_attempts}): "
                    f"{type(exc).__name__}: {exc} — rebuilding credential and retrying"
                )
            except Exception:
                pass
            _get_az_credential(force_rebuild=True)
            if on_rebuild is not None:
                try:
                    on_rebuild()
                except Exception:
                    pass
            time.sleep(min(2 ** attempt, 5))
    if last_exc is not None:
        raise last_exc

from agentcore.custom import Node
from agentcore.io import BoolInput, HandleInput, IntInput, MessageTextInput, MultilineInput, Output
from agentcore.schema.data import Data
from agentcore.schema.message import Message

_MAX_RETRIES   = 3
_BASE_DELAY    = 2.0
_STAGE_INGESTION  = 1
_STAGE_EMBED_DOC  = 2
_STAGE_BLOB_UPLOAD = 3
_STAGE_EXCEPTION  = 99


# ─────────────────────────────────────────────────────────────────────────────
# Connection pool — keeps total concurrent DB connections below Azure SQL's
# per-DB cap (Standard tier defaults: ~75-100 concurrent sessions). Without
# this, parallel PR workers × inner Stage 4 classification pool ×
# per-helper _connect() calls would each open their own TCP connection and
# trip "max connections reached" on Azure side.
#
# How it works:
#   • Bounded — `max_size` caps concurrent OPEN connections globally
#     (idle + busy combined). New requests block on threading.Condition()
#     until a slot frees instead of opening yet another socket.
#   • Reuse — idle connections are checked back in on .close() and handed
#     to the next acquirer (validated with SELECT 1; stale ones are
#     discarded silently).
#   • Per-cs partitioning — different connection strings (source vs target,
#     different DBs) don't compete; idle entries are tagged with their cs
#     and only returned to acquirers asking for the same one.
#   • Transparent close() — _PooledConnection.close() returns the underlying
#     connection to the pool. Existing `conn.close()` calls in finally
#     blocks throughout the file Just Work.
#   • Failure-aware — a connection that raised an exception in cursor() /
#     commit() / rollback() is marked broken and discarded on release
#     instead of being recycled.
#
# Defaults: max_size=20, login_timeout=30. Resized at startup via the
# canvas knob `max_db_connections` (see PipelineStage123NodeV2.inputs).
# ─────────────────────────────────────────────────────────────────────────────

import queue as _queue_mod
import threading as _threading

_DEFAULT_POOL_SIZE = 20


class _PooledConnection:
    """Wraps a pyodbc.Connection so .close() returns it to the pool
    instead of tearing down the socket. Forwards every other attribute /
    method to the underlying connection, marking the wrapper broken if
    cursor/commit/rollback raise so it won't be recycled."""

    __slots__ = ("_pool", "_cs", "_raw", "_broken")

    def __init__(self, pool: "_ConnectionPool", cs: str, raw: Any) -> None:
        object.__setattr__(self, "_pool",    pool)
        object.__setattr__(self, "_cs",      cs)
        object.__setattr__(self, "_raw",     raw)
        object.__setattr__(self, "_broken",  False)

    def cursor(self, *args, **kwargs):
        try:
            return self._raw.cursor(*args, **kwargs)
        except Exception:
            object.__setattr__(self, "_broken", True)
            raise

    def commit(self):
        try:
            return self._raw.commit()
        except Exception:
            object.__setattr__(self, "_broken", True)
            raise

    def rollback(self):
        try:
            return self._raw.rollback()
        except Exception:
            object.__setattr__(self, "_broken", True)
            raise

    def close(self) -> None:
        raw = self._raw
        if raw is None:
            return
        object.__setattr__(self, "_raw", None)
        self._pool.release(raw, self._cs, broken=self._broken)

    def __getattr__(self, name):
        raw = object.__getattribute__(self, "_raw")
        if raw is None:
            raise AttributeError(f"Pooled connection is closed (attr '{name}')")
        return getattr(raw, name)

    def __setattr__(self, name, value):
        if name in self.__slots__:
            object.__setattr__(self, name, value)
        else:
            raw = object.__getattribute__(self, "_raw")
            if raw is None:
                raise AttributeError(f"Pooled connection is closed (attr '{name}')")
            setattr(raw, name, value)

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        if exc is not None:
            object.__setattr__(self, "_broken", True)
        self.close()
        return False


class _ConnectionPool:
    """Thread-safe bounded connection pool for pyodbc. See header comment
    above for full behaviour."""

    def __init__(self, max_size: int = _DEFAULT_POOL_SIZE) -> None:
        self.max_size = max(1, int(max_size))
        self._idle: _queue_mod.Queue = _queue_mod.Queue()
        self._busy = 0
        self._cv   = _threading.Condition()
        self._opened_total = 0
        self._reused_total = 0

    def resize(self, new_max: int) -> None:
        """Change the cap. Existing checkouts continue to count; new acquires
        respect the new cap immediately."""
        with self._cv:
            self.max_size = max(1, int(new_max))
            self._cv.notify_all()

    def stats(self) -> dict:
        with self._cv:
            return {
                "max_size":     self.max_size,
                "busy":         self._busy,
                "idle":         self._idle.qsize(),
                "opened_total": self._opened_total,
                "reused_total": self._reused_total,
            }

    def _acquire_slot(self) -> None:
        with self._cv:
            while self._busy >= self.max_size:
                self._cv.wait()
            self._busy += 1

    def _release_slot(self) -> None:
        with self._cv:
            self._busy -= 1
            self._cv.notify()

    def acquire(self, cs: str, *, login_timeout: int = 30) -> _PooledConnection:
        """Get a connection (from idle pool or fresh). Caller must call
        .close() on the returned wrapper — that returns it to the pool."""
        self._acquire_slot()
        # Try idle reuse first — drain stale and wrong-cs entries
        try:
            while True:
                try:
                    cached_cs, conn = self._idle.get_nowait()
                except _queue_mod.Empty:
                    break
                if cached_cs != cs:
                    # Different connection string parked here — close and try next
                    try: conn.close()
                    except Exception: pass
                    continue
                try:
                    cur = conn.cursor()
                    cur.execute("SELECT 1")
                    cur.fetchone()
                    cur.close()
                    with self._cv:
                        self._reused_total += 1
                    return _PooledConnection(self, cs, conn)
                except Exception:
                    # Stale — discard and try next idle entry
                    try: conn.close()
                    except Exception: pass
                    continue
            # No idle entry — open a fresh one
            import pyodbc
            raw = pyodbc.connect(cs, timeout=login_timeout)
            # Set a default query timeout so a slow / locked SQL Server
            # statement can't silently block a worker thread forever.
            # 300s is generous for bulk ETL inserts but still bounded.
            try:
                raw.timeout = 300
            except Exception:
                pass
            with self._cv:
                self._opened_total += 1
            return _PooledConnection(self, cs, raw)
        except Exception:
            self._release_slot()
            raise

    def release(self, conn: Any, cs: str, *, broken: bool = False) -> None:
        """Return a connection to the pool; close + free slot if broken."""
        if conn is None:
            self._release_slot()
            return
        if broken:
            try: conn.close()
            except Exception: pass
            self._release_slot()
            return
        try:
            self._idle.put_nowait((cs, conn))
        except _queue_mod.Full:
            try: conn.close()
            except Exception: pass
        finally:
            self._release_slot()


# Global pool — sized at startup from the canvas knob max_db_connections.
# Module-level so the same pool is shared across all PR workers.
_CONN_POOL = _ConnectionPool(max_size=_DEFAULT_POOL_SIZE)

# Belt-and-suspenders — also enable ODBC driver-manager pooling so that
# anything that bypasses _CONN_POOL (third-party libs that import pyodbc
# directly) still benefits from lower-level reuse.
try:
    import pyodbc as _pyodbc_for_pooling
    _pyodbc_for_pooling.pooling = True
except Exception:
    pass


def _start_stack_dump_thread(interval_s: int = 300) -> None:
    """Background thread that periodically logs the current stack of every
    live thread. Makes silent hangs diagnosable — when MiCore appears
    stuck mid-batch with no log lines, look in the logs around the next
    interval to see exactly where each worker is blocked (socket read,
    SQL execute, threading wait, etc.).

    Idempotent — guarded by a function attribute so repeated component
    rebuilds within the same process don't spawn multiple dumpers.
    """
    if getattr(_start_stack_dump_thread, "_started", False):
        return
    object.__setattr__(_start_stack_dump_thread, "_started", True)

    import sys as _sys, traceback as _tb, threading as _th, time as _time

    def _dumper() -> None:
        # Startup beacon — confirms the dumper is actually running so users
        # know "no error logs" isn't because the dumper itself failed.
        try:
            logger.info(
                f"STACK-DUMP: stack dumper started, will log every {interval_s}s — "
                "look for '=== Thread stack dump' lines when pipeline appears stuck"
            )
        except Exception:
            pass
        while True:
            try:
                _time.sleep(interval_s)
                frames = _sys._current_frames()
                names = {t.ident: t.name for t in _th.enumerate()}
                lines = [
                    f"=== Thread stack dump | {len(frames)} live thread(s) ===",
                ]
                for tid, frame in frames.items():
                    nm = names.get(tid, "?")
                    lines.append(f"--- Thread {tid} ({nm}) ---")
                    lines.append("".join(_tb.format_stack(frame)).rstrip())
                # Use a single multi-line message + force a flush via warning
                # level so MiCore's log capture doesn't drop / buffer it.
                logger.warning("\n".join(lines))
            except Exception as exc:
                # Never let the diagnostic thread itself crash the pipeline.
                try:
                    logger.warning("stack-dumper: {}", exc)
                except Exception:
                    pass

    t = _th.Thread(target=_dumper, name="stack-dumper", daemon=True)
    t.start()


# Start the diagnostic thread immediately on module load. Daemon thread so
# it won't keep the process alive on shutdown. 90s interval — short enough
# that a stuck worker shows up quickly while debugging large batches.
try:
    _start_stack_dump_thread(interval_s=90)
except Exception:
    pass


class ExtractionAbortError(RuntimeError):
    """Raised when Stage 5 (Extraction) cannot produce any items.

    The Stage 4-8 catch block treats every subclass as an expected pipeline
    outcome (not a crash): logs it as a warning without traceback, records a
    row in ras_pipeline_exceptions (stage_id=5), sets
    ras_tracker.current_stage_fk = 99 (EXCEPTION), and skips stages 6/7/8.

    Subclasses identify the specific reason so the exception_message in DB
    is self-explanatory.
    """
    _suppress_traceback = True


class NoQuotationFoundError(ExtractionAbortError):
    """No attachments classified as 'Quotation'."""


class NoLineItemsError(ExtractionAbortError):
    """RAS context has zero rows in purchase_req_detail."""


class NoRASContextError(ExtractionAbortError):
    """No purchase_req_mst row found for this PR."""


class AllExtractionsFailedError(ExtractionAbortError):
    """All quotation files failed to extract any items."""

_TEXT_EXT  = {".txt", ".csv", ".xml", ".json", ".html", ".htm", ".eml"}
_PDF_EXT   = {".pdf"}
_EXCEL_EXT = {".xlsx", ".xls", ".xlsm", ".xlsb"}
_WORD_EXT  = {".doc", ".docx"}
_PPTX_EXT  = {".ppt", ".pptx"}
_MSG_EXT   = {".msg"}

SUPPORTED_PARENTS   = (".xlsx", ".xls", ".docx", ".doc", ".pptx", ".ppt", ".pdf", ".msg")
ARCHIVE_EXTENSIONS  = (".zip", ".7z", ".rar", ".tar", ".tar.gz", ".tgz", ".tar.bz2")
EXTRACT_EXTENSIONS  = (
    ".pdf", ".xls", ".xlsx", ".doc", ".docx", ".ppt", ".pptx",
    ".jpg", ".jpeg", ".png", ".tif", ".tiff",
    ".zip", ".7z", ".rar", ".tar", ".tar.gz", ".tgz", ".tar.bz2",
    ".txt", ".msg",
)
SKIP_EXTENSIONS = (".emf", ".bin")
FILE_SIGNATURES = [
    (b"%PDF",               ".pdf"),
    (b"PK\x03\x04",         ".zip"),
    (b"7z\xBC\xAF\x27\x1C", ".7z"),
    (b"Rar!\x1A\x07",       ".rar"),
    (b"\xFF\xD8\xFF",       ".jpg"),
    (b"\x89PNG\r\n\x1A\n",  ".png"),
    (b"II*\x00",            ".tif"),  # TIFF little-endian
    (b"MM\x00*",            ".tif"),  # TIFF big-endian
]


# ── Embedded-doc extractor (inlined from embed_doc_extraction.extractor) ──

class FileExtractor:
    def __init__(self) -> None:
        self.extracted_count: int = 0
        self.parent_prefix:   str = ""

    def should_extract(self, filename: str) -> bool:
        fl = filename.lower()
        if any(fl.endswith(s) for s in SKIP_EXTENSIONS):
            return False
        return any(fl.endswith(e) for e in EXTRACT_EXTENSIONS)

    def sanitize(self, name: str) -> str:
        import os
        name = str(name).replace("\x00", "").replace("\x01", "").strip()
        name = os.path.basename(name)
        for ch in '<>:"/\\|?*':
            name = name.replace(ch, "_")
        name = "".join(ch for ch in name if ch.isprintable() and ord(ch) > 31)
        return name.strip(" .") or "file"

    def with_prefix(self, filename: str) -> str:
        filename = self.sanitize(filename)
        prefix   = self.sanitize(self.parent_prefix).strip()
        return f"{prefix}__{filename}" if prefix else filename

    def unique_path(self, directory: str, filename: str) -> str:
        import os
        fp = os.path.join(directory, filename)
        if not os.path.exists(fp):
            return fp
        base, ext = os.path.splitext(filename)
        c = 1
        while os.path.exists(os.path.join(directory, f"{base}_{c}{ext}")):
            c += 1
        return os.path.join(directory, f"{base}_{c}{ext}")

    def _find_payload(self, data: bytes):
        best = None
        for sig, ext in FILE_SIGNATURES:
            pos = data.find(sig)
            if pos != -1 and (best is None or pos < best[0]):
                best = (pos, ext)
        return best if best else (None, None)

    def _classify_zip(self, path: str) -> str:
        """If a saved .zip is actually an Office Open XML file, return the
        correct extension (.xlsx / .docx / .pptx). Mirrors doc-intel's
        embed_doc_extraction.extractor._classify_zip — without it, an XLSX
        payload extracted from an OLE container ends up on disk as `.zip`,
        which the classifier (limited to _SUPPORTED_CLASSIFY_EXTS) treats
        as 'Others' and any embedded Quotation is lost.
        """
        import zipfile
        try:
            with zipfile.ZipFile(path, "r") as z:
                names = set(z.namelist())
                if "[Content_Types].xml" in names:
                    if any(n.startswith("xl/")   for n in names): return ".xlsx"
                    if any(n.startswith("word/") for n in names): return ".docx"
                    if any(n.startswith("ppt/")  for n in names): return ".pptx"
        except Exception:
            pass
        return ".zip"

    def _maybe_fix_zip_extension(self, target: str) -> str:
        """If `target` ends in .zip but is actually OOXML, rename and return new path."""
        import os
        if not target.lower().endswith(".zip"):
            return target
        correct = self._classify_zip(target)
        if correct == ".zip":
            return target
        new_path = self.unique_path(
            os.path.dirname(target),
            os.path.splitext(os.path.basename(target))[0] + correct,
        )
        try:
            os.replace(target, new_path)
            logger.debug("Reclassified .zip → {}: {}", correct, os.path.basename(new_path))
            return new_path
        except Exception as exc:
            logger.warning("Could not rename {} → {}: {}", target, new_path, exc)
            return target

    def _save_file(self, data: bytes, filename: str, out: str) -> str:
        import os, shutil
        target = self.unique_path(out, self.with_prefix(filename))
        with open(target, "wb") as f:
            f.write(data)
        self.extracted_count += 1
        # Fix OOXML-misclassified-as-.zip BEFORE archive expansion so the
        # file isn't extracted as a plain ZIP archive when it's actually an
        # XLSX/DOCX/PPTX worth classifying directly.
        target = self._maybe_fix_zip_extension(target)
        if target.lower().endswith(ARCHIVE_EXTENSIONS):
            self.extract_archive(target, out)
        return target

    def extract_archive(self, path: str, out: str) -> None:
        import os, shutil, tarfile, tempfile, zipfile
        if not os.path.exists(path):
            return
        fl = path.lower()
        try:
            if fl.endswith(".zip"):
                with zipfile.ZipFile(path, "r") as z:
                    for m in z.infolist():
                        if m.is_dir():
                            continue
                        fn = os.path.basename(m.filename.replace("\\", "/").rstrip("/"))
                        if fn:
                            t = self.unique_path(out, self.with_prefix(fn))
                            with z.open(m) as s, open(t, "wb") as d:
                                shutil.copyfileobj(s, d)
                            self.extracted_count += 1
            elif fl.endswith(".7z"):
                import py7zr
                with tempfile.TemporaryDirectory() as tmp:
                    with py7zr.SevenZipFile(path, "r") as ar:
                        ar.extractall(path=tmp)
                    for root, _, files in os.walk(tmp):
                        for f in files:
                            t = self.unique_path(out, self.with_prefix(f))
                            shutil.copy2(os.path.join(root, f), t)
                            self.extracted_count += 1
            elif fl.endswith(".rar"):
                import rarfile
                with tempfile.TemporaryDirectory() as tmp:
                    with rarfile.RarFile(path, "r") as rar:
                        rar.extractall(tmp)
                    for root, _, files in os.walk(tmp):
                        for f in files:
                            t = self.unique_path(out, self.with_prefix(f))
                            shutil.copy2(os.path.join(root, f), t)
                            self.extracted_count += 1
            elif fl.endswith((".tar", ".tar.gz", ".tgz", ".tar.bz2")):
                with tarfile.open(path, "r:*") as tar:
                    for m in tar.getmembers():
                        if not m.isfile():
                            continue
                        fn = os.path.basename(m.name)
                        if fn:
                            t  = self.unique_path(out, self.with_prefix(fn))
                            src = tar.extractfile(m)
                            if src:
                                with src as s, open(t, "wb") as d:
                                    shutil.copyfileobj(s, d)
                                self.extracted_count += 1
        except Exception as exc:
            logger.warning(f"Archive extraction error ({os.path.basename(path)}): {exc}")
        try:
            os.remove(path)
        except Exception:
            pass

    def extract_from_ooxml(self, file_path: str, out: str) -> bool:
        import os, shutil, tempfile, zipfile
        try:
            with tempfile.TemporaryDirectory() as tmp:
                with zipfile.ZipFile(file_path, "r") as zr:
                    zr.extractall(tmp)
                for prefix in ("xl", "word", "ppt"):
                    for sub in ("media", "embeddings"):
                        folder = os.path.join(tmp, prefix, sub)
                        if not os.path.exists(folder):
                            continue
                        for fn in os.listdir(folder):
                            src = os.path.join(folder, fn)
                            if not os.path.isfile(src):
                                continue
                            if fn.lower().endswith(".bin"):
                                self._extract_from_bin(src, out)
                            elif self.should_extract(fn):
                                dst = self.unique_path(out, self.with_prefix(fn))
                                shutil.copy2(src, dst)
                                self.extracted_count += 1
                                # OOXML embeddings often come as .zip-stamped
                                # bytes that are really nested xlsx/docx/pptx.
                                dst = self._maybe_fix_zip_extension(dst)
                                if dst.lower().endswith(ARCHIVE_EXTENSIONS):
                                    self.extract_archive(dst, out)
            return True
        except Exception as exc:
            logger.error(f"OOXML error ({os.path.basename(file_path)}): {exc}")
            return False

    def _extract_from_bin(self, bin_path: str, out: str) -> None:
        try:
            import olefile
            if not olefile.isOleFile(bin_path):
                return
            ole = olefile.OleFileIO(bin_path)
            for entry in ole.listdir():
                try:
                    sname = entry[-1] if entry else ""
                    data  = ole.openstream(entry).read()
                    if not data or len(data) < 16:
                        continue
                    offset, ext = self._find_payload(data)
                    if offset is None:
                        continue
                    payload   = data[offset:]
                    suggested = None
                    if sname.lower() in ("\x01ole10native", "ole10native"):
                        try:
                            end = data.find(b"\x00", 4)
                            if end != -1 and end - 4 < 260:
                                suggested = data[4:end].decode(errors="ignore")
                        except Exception:
                            pass
                    suggested = suggested or f"embedded_file{ext}"
                    target = self.unique_path(out, self.with_prefix(self.sanitize(suggested)))
                    with open(target, "wb") as f:
                        f.write(payload)
                    self.extracted_count += 1
                    # Reclassify .zip → .xlsx/.docx/.pptx if the payload was
                    # actually an OOXML doc; then expand if it remains a zip.
                    target = self._maybe_fix_zip_extension(target)
                    if target.lower().endswith(ARCHIVE_EXTENSIONS):
                        self.extract_archive(target, out)
                except Exception:
                    pass
            ole.close()
        except Exception as exc:
            logger.warning(f"OLE bin error: {exc}")

    def extract_from_ole(self, file_path: str, out: str) -> bool:
        try:
            import olefile
            if not olefile.isOleFile(file_path):
                return False
            ole = olefile.OleFileIO(file_path)
            for entry in ole.listdir():
                try:
                    sname = entry[-1] if entry else ""
                    data  = ole.openstream(entry).read()
                    if not data or len(data) < 16:
                        continue
                    offset, ext = self._find_payload(data)
                    if offset is None:
                        continue
                    payload   = data[offset:]
                    suggested = None
                    if sname.lower() in ("\x01ole10native", "ole10native"):
                        try:
                            end = data.find(b"\x00", 4)
                            if end != -1 and end - 4 < 260:
                                suggested = data[4:end].decode(errors="ignore")
                        except Exception:
                            pass
                    suggested = suggested or f"embedded_file{ext}"
                    target = self.unique_path(out, self.with_prefix(self.sanitize(suggested)))
                    with open(target, "wb") as f:
                        f.write(payload)
                    self.extracted_count += 1
                    # Same OOXML-misclassified-as-.zip fix as in _extract_from_bin.
                    target = self._maybe_fix_zip_extension(target)
                    if target.lower().endswith(ARCHIVE_EXTENSIONS):
                        self.extract_archive(target, out)
                except Exception:
                    pass
            ole.close()
            return True
        except Exception as exc:
            logger.error(f"OLE error ({os.path.basename(file_path)}): {exc}")
            return False

    def extract_from_pdf(self, file_path: str, out: str) -> bool:
        import os  # used by the error handler below
        try:
            import fitz
            doc = fitz.open(file_path)
            try:
                if doc.embfile_count() > 0:
                    for i in range(doc.embfile_count()):
                        info = doc.embfile_info(i)
                        name = self.sanitize(info.get("filename", f"pdf_attachment_{i}"))
                        data = doc.embfile_get(i)
                        if data:
                            self._save_file(data, name, out)
                for pn, page in enumerate(doc, 1):
                    annots = page.annots()
                    if not annots:
                        continue
                    for annot in annots:
                        if annot.type[0] == 17:
                            try:
                                fd      = annot.get_file()
                                name    = self.sanitize(fd.get("filename", f"page{pn}_attachment"))
                                content = fd.get("content", b"")
                                if content:
                                    self._save_file(content, name, out)
                            except Exception:
                                pass
            finally:
                doc.close()
            return True
        except Exception as exc:
            logger.error(f"PDF error ({os.path.basename(file_path)}): {exc}")
            return False

    def extract_from_msg(self, file_path: str, out: str) -> bool:
        try:
            import tempfile, os
            import extract_msg
            msg = extract_msg.Message(file_path)
            if not msg.attachments:
                msg.close()
                return True
            for att in msg.attachments:
                name = "attachment"
                try:
                    name = (
                        getattr(att, "longFilename", None)
                        or getattr(att, "shortFilename", None)
                        or getattr(att, "filename", None)
                        or "attachment"
                    )
                    name = self.sanitize(name)
                    if getattr(att, "type", None) == "msg" and not name.lower().endswith(".msg"):
                        name += ".msg"
                    data = getattr(att, "data", None) or b""
                    if not data:
                        with tempfile.TemporaryDirectory() as tmp:
                            att.save(customPath=tmp)
                            for fn in os.listdir(tmp):
                                src = os.path.join(tmp, fn)
                                if os.path.isfile(src):
                                    with open(src, "rb") as f:
                                        data = f.read()
                                    if name == "attachment":
                                        name = fn
                                    break
                    if data:
                        self._save_file(data, name, out)
                except Exception as exc:
                    logger.warning(f"MSG attachment '{name}' failed: {exc}")
            msg.close()
            return True
        except Exception as exc:
            logger.error(f"MSG error ({os.path.basename(file_path)}): {exc}")
            return False

    def detect_type(self, path: str):
        lower = path.lower()
        if lower.endswith((".xlsx", ".docx", ".pptx")): return "ooxml"
        if lower.endswith((".xls",  ".doc",  ".ppt")):  return "ole"
        if lower.endswith(".pdf"):                       return "pdf"
        if lower.endswith(".msg"):                       return "msg"
        try:
            with open(path, "rb") as f:
                h = f.read(8)
            if h.startswith(b"PK\x03\x04"):                        return "ooxml"
            if h.startswith(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"): return "ole"
            if h.startswith(b"%PDF"):                               return "pdf"
        except Exception:
            pass
        return None

    def process_file(self, file_path: str, out: str) -> bool:
        import os
        ft = self.detect_type(file_path)
        if not ft:
            return False
        os.makedirs(out, exist_ok=True)
        if ft == "ooxml": return self.extract_from_ooxml(file_path, out)
        if ft == "ole":   return self.extract_from_ole(file_path, out)
        if ft == "pdf":   return self.extract_from_pdf(file_path, out)
        if ft == "msg":   return self.extract_from_msg(file_path, out)
        return False


# ── Stage 4-8 helpers (inlined) ─────────────────────────────────────────────

# ── Stage IDs ─────────────────────────────────────────────────────────────────
_STAGE_CLASSIFICATION  = 4
_STAGE_EXTRACTION      = 5
_STAGE_EMBEDDINGS      = 6
_STAGE_PRICE_BENCHMARK = 7
_STAGE_COMPLETE        = 8


# ── Classification constants ───────────────────────────────────────────────────
_CLASSIFICATION_MAP = {"E-Auction": "E-Auction Results", "Other": "Others"}
_SUPPORTED_CLASSIFY_EXTS = {
    ".xlsx", ".xls", ".csv", ".pdf", ".docx", ".doc",
    ".pptx", ".ppt", ".txt", ".html", ".htm",
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp",
}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"}
_MAX_CLASSIFY_CHARS = 100000  # safety cap inside each extractor only — the
                              # *active* classification truncation happens in
                              # _classify_file via the configurable cls_max_chars
                              # (default 24000 ≈ 6000 tokens, matches doc-intel
                              # truncate_to_token_limit max_content_tokens=6000).

_EXT_TO_TYPE = {
    ".xlsx": "excel",  ".xls": "excel",  ".xlsm": "excel",  ".xlsb": "excel",  ".csv": "csv",
    ".pdf":  "pdf",    ".docx": "word",  ".doc": "legacy_doc",
    ".pptx": "pptx",  ".ppt":  "legacy_doc",
    ".txt":  "text",   ".html": "html",  ".htm": "html",
    ".xml":  "text",   ".json": "text",  ".eml": "text",
    ".png":  "image",  ".jpg":  "image", ".jpeg": "image",
    ".tiff": "image",  ".tif":  "image", ".bmp":  "image",
    ".msg":  "msg",
}

VALID_CLASSIFICATIONS = {"RFQ", "Quotation", "MPBC", "BER", "E-Auction", "Other"}

# ── Classification system prompt (full original recovered from file_classifier) ──
CLASSIFICATION_SYSTEM_PROMPT = """You are an expert document classifier for Motherson's procurement team. You will be given the content of a file (extracted text, tables, sheet structure, or an image) and you must classify it into EXACTLY ONE of six categories.

============================================================
CRITICAL: MULTILINGUAL SUPPORT
============================================================
Documents may be in ANY language: English, German, Czech, Hindi, French, Spanish, Chinese, Japanese, Hungarian, and others. You MUST:
- Match fields by MEANING, not by exact English labels
- Recognize translated equivalents. Common examples:
  • German: "Preisspiegel" / "Angebotsvergleich" = Bid Comparison (MPBC), "Angebot" = Quotation/Offer, "Anfrage" = Inquiry/RFQ, "Lieferant" = Supplier, "Preis" = Price, "Lieferzeit" = Delivery Time, "Zahlungsbedingungen" = Payment Terms, "Genehmigung" = Approval, "Begründung" = Justification
  • Czech: "Nabídka" = Offer/Quotation, "Poptávka" = RFQ/Inquiry, "Dodavatel" = Supplier, "Cena" = Price, "Technický požadavek" / "Lastenheft" = Technical Specification
  • Hindi: "कोटेशन" = Quotation, "मूल्य" = Price, "आपूर्तिकर्ता" = Supplier
- Apply the same field-matching logic regardless of language
- If a document is in a language you recognize, translate field names mentally and match against the English category definitions

============================================================
CLASSIFICATION METHOD
============================================================
Classification is FIELD-BASED VALIDATION. Each category has a list of mandatory fields. You must verify the presence of those fields in the document content (allowing for naming/labeling variations, synonyms, AND translations in any language — the underlying meaning matters, not the exact label). The category whose mandatory fields are most completely satisfied wins.

============================================================
CATEGORIES & MANDATORY FIELDS
============================================================

1. **MPBC** — Motherson Purchase BID Comparison
   Official Motherson template that consolidates quotations from multiple suppliers (typically 3+) into a side-by-side comparison for procurement evaluation. Sheet usually titled "MPBC" or "Motherson Purchase BID Comparison". May be a multi-sheet workbook with sheets like "1. MPBC", "2. mandatory cells", "BER", "Supp X risk", "exchange_rates", etc.

   Mandatory fields (yellow-marked in the official MPBC template — must find most of these):
   PROJECT / HEADER:
     • Project Detail
     • RAS Number
     • Sheet No.
     • Indenter, Originator or RFQ Responsible
     • Contact Number
   PER-SUPPLIER (repeated for each of typically 3+ suppliers):
     • Name of the Supplier
     • Preference (Preferred source / Second Source / etc.)
     • Quote meet spec (Yes/No)
     • Supplier's Contact Person
     • Offer No.
     • Offer Date
     • Offer Validity
     • Supplier Country / Classification
     • Country of Origin / Classification
   PRICING:
     • Currency for Comparison
     • Total Amount in supplier currency
     • Total Amount in EUR
     • Saving %age
   SOURCING DECISION:
     • Insourcing
     • Hybrid
     • Single sourcing
     • Supplier justification (BER) signed
     • In case of Cust. funded: Revenue in EUR
     • Target against Budget
   COMMERCIAL TERMS (per supplier):
     • Incoterms
     • Packing and Forwarding
     • Freight
     • Insurance
     • Taxes
     • Customs / duties
     • Installation
     • Delivery Period (Weeks)
     • Advance payment (any payment before delivery to MOTHERSON)
     • SCF (Supply Chain Finance Program)
     • Payment Terms
   APPROVAL / FINAL:
     • GSP Purchase Saving %age
     • GSP Purchase Saving amount in EUR
     • Approved Supplier / Classification
     • Approved Cost (total amount)
     • Landed Cost (for Reference)
     • Approvals (Department / Name / Date / Approver Signature for Purchasing, Technics, Controlling, Sales, Plant Manager)

   Strongest distinguishing signals:
     • Title / sheet name contains "MPBC" or "Motherson Purchase BID Comparison"
     • Supplier 1 / Supplier 2 / Supplier 3 columns appearing side-by-side
     • Presence of RAS Number, GSP (Global Strategic Procurement) terminology
     • Multi-sheet Excel with sheets like "1. MPBC", "2. mandatory cells", "BER", "Supp X risk"
     • MULTILINGUAL: German "Preisspiegel" (price mirror) / "Angebotsvergleich" (offer comparison) / "Bid comparison (Purchased parts)" = MPBC. If 3+ suppliers are compared side-by-side in ANY language with pricing, it is MPBC even without the exact Motherson MPBC template fields like RAS Number.
   Allow naming variations / synonyms / translations; the MEANING of fields matters, not exact labels.

2. **Quotation** — A SINGLE vendor's price offer / response to an RFQ (vendor → buyer such as Motherson)
   Synonyms commonly used as the document title: "Quotation", "Quote", "Estimate", "Offer", "Proposal", "Price Bid". Treat all of these as Quotation candidates.

   Mandatory fields (naming may vary widely; identify by MEANING — Indian, German, and other vendor styles all appear):
     • Vendor Name (vendor's company name appears in letterhead at top AND in signature/footer like "For [Company Name]")
     • Vendor Address (postal address of the vendor)
     • Date of Quotation (e.g., "Date:", "Quotation Dt.", "Offer Date", "Est. Date")
     • Total Amount / Value (grand total / net amount / total in supplier currency)
     • Payment Terms (e.g., "30 days", "100% advance", "Against delivery", "100% after PO confirmation")
     • Delivery Time (e.g., "15-20 days", "within 2-3 weeks", "Delivery Period", "Lead Time", "Dispatch Time")
     • Validity (e.g., "Quotation valid until ___", "Offer Validity: 60 days", "Valid for one month")
     • Specifications (technical / product / scope specifications — may be a separate "TECH SPEC" sheet in Excel. For service quotations: scope of work, service descriptions, routes, equipment specs, or operational details)
     • Item Description (line items with descriptions — columns like "Description of Goods", "Item Description", "Sr. No. + Description + Qty + Rate + Amount", "BOQ items". For service quotations: line items may be lump-sum / flat-rate entries like "Line #1: Service description... Pauschal/Lump sum [amount]")

   Additional supporting signals frequently observed (not all needed, but strengthen confidence):
     • Vendor letterhead at top: company logo / name + address + phone + email + GST No. / GSTIN / PAN No. / VAT No.
     • A reference / quotation number (e.g., "Quotation No.: PPAQ004560", "Ref: Q/875", "Offer No.", "Est. No.: SCL/Haryana/20-21/21", "Ref: MIPL/MATE/...")
     • Addressed to: "To, M/s [Customer Name]", "Client:", "Kind Attn: Mr./Ms. ___"
     • Formal letter language: "We are pleased to offer / quote", "With reference to your enquiry", "Thanking you", "Yours faithfully"
     • Closing: "For [Vendor Company]" + "Authorised Signatory" / "Proprietor" / "Sales Manager"
     • Tax columns: HSN/SAC Code, CGST %, SGST %, IGST %, Excise Duty, GST @ 18 %
     • Bank details / E. & O. E. / "Please mention our quotation number on your purchase order"

   CRITICAL DISAMBIGUATION:
     • A Quotation is from ONE vendor. If THREE OR MORE vendor names appear as parallel columns/sections being compared → it is MPBC, not Quotation.
     • A Quotation MAY be a multi-sheet Excel (e.g., "COMMERCIAL", "TECH. SPEC.", "Summary", "NPV", "Vendors" sheets) but with only ONE vendor. Do not misclassify these as MPBC — check vendor count, not sheet count.
     • If price columns are EMPTY (template for vendor to fill) → it is RFQ, not Quotation.
     • Filename is unreliable — e.g., "(875) CK Motherson Auto Hitech.docx" was issued BY Hi-Tech to CK Motherson; the customer name in the filename does not make it MPBC or RFQ.
     • PDF extraction may produce duplicated characters from font issues (e.g., "TTOO" instead of "TO", "QQUUOOTTAATTIIOONN" instead of "QUOTATION") — interpret semantically.
     • **A Quotation can use the buyer's RFQ TEMPLATE FORMAT.** When a supplier (e.g., HAITIAN, Engel, Sumitomo, KraussMaffei, Arburg) fills in an RFQ specification template with their technical responses ("STANDARD", "OPTIONAL", "OPTIONAL - AVAILABLE", option codes) AND populates pricing fields (Basic Price, Option Price, Net Price, Total Machine Cost, Delivery Price, etc. with actual monetary values in USD/EUR/INR), the document has become a Quotation — the supplier's price offer. Key differentiator: if ACTUAL PRICES are filled in (not blank "to be specified" placeholders) and a supplier contact person / company name is prominently featured, treat it as Quotation even if the original RFQ structure ("Pls Specify", "Required") is retained.
     . There can be many vendors/supplier not just HAITIAN, Engel, Sumitomo, KraussMaffei, Arburg etc. these are just examples.
     • **A Quotation can be for ANY type of goods or services** — not just products or equipment. Transport, crane, logistics, packaging, freight forwarding, installation, civil works, consulting, IT services, facility management, and similar SERVICE quotations are still Quotation. Look for a vendor offering a price (even a lump-sum / "Pauschal" / flat rate) for defined scope of work. European-style quotations may use "Pauschal" (lump sum) line items instead of per-unit pricing — this is still a valid Quotation.
     • **Repeated page headers do NOT disqualify a Quotation.** Multi-page PDFs often repeat the vendor letterhead, bank details, and document header on every page. Focus on the actual content (pricing, scope, terms) rather than the repetition of header/footer blocks.
     • German Quotation synonyms beyond "Angebot": "Kostenvoranschlag" (cost estimate), "Preisangebot" (price offer), "Offerte" (offer). German "Angebots-Nr." = Offer/Quotation Number. "Pauschal" = lump sum / flat rate.
     • **Vendor price lists, rate cards, and Annual Rate Contract (ARC) reports are Quotation** when they contain a SINGLE vendor's agreed pricing. These documents may be system-generated (with internal fields like "Token No", "Type: ARC", "Qty in RAS") but still represent a vendor's price offer with per-item pricing (Basic Price, Discount %, Net Price), credit/payment terms, freight terms, and tax details (HSN/SAC Code, GST %). Classify as Quotation — the system-generated format does not change the commercial nature of the pricing data.

3. **RFQ** — Request For Quotation (Motherson → vendors)
   A specification / scope document issued BY Motherson (the buyer) TO vendors, asking them to submit a quotation. The document defines WHAT Motherson needs and asks the vendor to respond with specs and pricing. May be an Excel template with columns for "Supplier Spec/Confirmation" that are blank (or have been partially filled by a responding vendor).

   Mandatory fields:
     • Project Name (e.g., "INJECTION MOULDING MACHINE SPECIFICATIONS", "30KLD STP", or a specific project title)
     • RFQ Number (reference / inquiry number — e.g., "MATEB/IMG 1300/2020/SEPT/1-Rev 01", "MATE SADDLES/SP/01-24/Rev03")
     • Date (issue date)
     • Specifications — the BULK of the document; detailed technical / scope specifications structured as tables with:
       - Sl. No / Item number
       - Description / technical parameter (e.g., "Screw diameter [mm]", "Clamping force [kN]")
       - "Required" / "Not Required" / "Std" flags (Motherson's requirement)
       - "To be specified by the supplier" / "Supplier Spec/Confirmation" / "Pls specify" columns (blank or for vendor to fill)
       - Multiple specification sections: e.g., Injection unit, Clamping unit, Electrical/Hydraulic, Controller, General features, Spares, Special features
     • Commercials — a section (typically at the bottom) asking for pricing breakdown:
       - Option Price, Basic Price, Gross price, Special discount, Net price
       - Spares Package, Sea worthy packing, CIF Cost + Insurance
       - Startup cost, Delivery Price, Price in INR
       - Payment terms, Taxes, Delivery period, Warranty, Guarantee

   Strong supporting signals:
     • Sheet name contains "RFQ" (e.g., " RFQ", "RFQ - STP")
     • "MATE-B Req" or "MATE Spec" column (Motherson Automotive Technologies & Engineering — a Motherson entity)
     • "Supplier Spec/Confirmation" or "Supplier Remarks" columns (empty = template, filled = vendor response captured in same template)
     • "Techno Commercial Comparison" as a title (still an RFQ when it's the request template, even if some vendor responses have been captured)
     • Accompanying sheet " Parts & Mold data " with part specifications (part name, size, weight, wall thickness)
     • Columns with "Fill this column by '0' if STD or Enter the cost if Optional" — template instruction language
     • "Company Name:", "Contact Person:", "Telephone:", "Email:" fields (for vendor to fill)

   CRITICAL DISAMBIGUATION:
     • An RFQ is the TEMPLATE / REQUEST document. Even if one or two vendors have filled in their spec responses, the document may still be an RFQ IF:
         (a) pricing fields remain EMPTY or say "to be specified", AND
         (b) no supplier contact/company is featured prominently.
       HOWEVER, if a supplier has populated pricing (Basic Price, Net Price, Total cost, Option prices with actual USD/EUR/INR values) AND the supplier's contact details appear prominently (name, phone, email), the document has become a QUOTATION — the supplier's completed price offer using the RFQ template format. Classify as Quotation, not RFQ.
     • If the document compares 3+ vendors side-by-side in a CONSOLIDATED EVALUATION format → it is MPBC, not RFQ.
     • If the document is a standalone vendor letter with prices (no Motherson spec template structure) → it is Quotation, not RFQ.
     • RFQ documents are typically MUCH longer / more detailed than Quotations (100+ rows of specifications with Required/Not Required flags). BUT length alone does not determine the category — a long spec sheet with filled-in prices is still a Quotation.

4. **BER** — Bid Exception Report
   SPECIFICALLY the Motherson "BID EXCEPTION REPORT" template form. This is NOT a catch-all for any waiver, justification, or single-source document. The document must use the Motherson BER template structure.

   Mandatory fields (ALL of these must be present or nearly all — this is a strict template match):
     • "BID EXCEPTION REPORT" appearing explicitly as a header / title (this exact phrase or very close equivalent is REQUIRED — a generic "Waiver of Competition" or "Single Source Justification" title is NOT sufficient for BER classification)
     • Reasoning for not obtaining at least three bids/quotes (NOTE: language may differ in 20–40% of cases)
     • Order Value field (e.g., "Order Value: 671,57 €", "Budget Line Ref. + amount")
     • Description of the product or service to be ordered (a labeled section)
     • Justification for waiver of competitive bidding (a labeled section)
     • Justification of the product or service to be ordered (a separate labeled section)

   Strong supporting signals (Motherson BER template-specific — presence of 2-3 is a near-certain indicator):
     • Header "Capital Equipment & Indirect Purchasing"
     • Budget Line Ref. (e.g., "Budget Line Ref.: ID 413128")
     • Reference to "LCC Suppliers" / "Low Cost Country" / "Motherson internal company"
     • Checkbox-style options A through E for reasoning (A: less than three potential bidders; B: sole-source / proprietary item; C: national/local supply contract; D: similar item purchased in past 6 months; E: Other reasons)
     • Three approval rows: "Prepared by:" + "Purchasing approval:" + "Managing Director / COO / EVP approval:"
     • "Approver comments:" section
     • Footer with template revision history (e.g., "Compiled by Mac Cheema", "Revised by Lousie Osgood")
     • Sheet named "BER" in an Excel workbook

   CRITICAL: Do NOT classify as BER if the document is:
     • A generic "Waiver of Competition" form (different template, different structure)
     • A "Single Source Justification" from a non-Motherson template
     • A "Double Source Waiver" or supplier approval form
     • Any competitive bidding exception document that does NOT use the specific Motherson "BID EXCEPTION REPORT" template with the A-E checkbox structure
     These should be classified as "Other" instead.

5. **E-Auction** — E-Auction results / reports / trackers
   Documents generated from or summarizing an online reverse-auction event where vendors bid in real time. May be raw auction output from an e-procurement platform OR a summary/tracker/presentation consolidating auction results.

   Mandatory fields (from client spec — match by meaning, naming may vary):
     • Event ID (e.g., "Doc3071019131", "Auction ID")
     • Event Name (e.g., "Live Japanese Reverse eAuction - R_268526-2026-EPP boxes...")
     • Publish Date
     • Open Date
     • Close Date
     • BID Id (e.g., "ID3240307353", "ID3244896110")
     • BID Status (e.g., "Accepted", "Default")
     • Participant (vendor / bidder name — e.g., "MORAplast, s.r.o(Sarka Bris...)")
     • Basic Price (per unit price)
     • Extended Price (total price for volume)

   Strong supporting signals (presence of 3+ is near-certain):
     • "eAuction", "e-Auction", "Reverse Auction", "Japanese Auction", "Japanese Reverse eAuction" in title/headers
     • Sheet names like "Overview Sheet", "Full Bid Data Sheet"
     • Event Type field (e.g., "Japanese Auction")
     • Owner field with eAuction email (e.g., "eAuction.GSP@motherson.com")
     • Rank column (bidder ranking: 1, 2, 3...)
     • Savings column / Total Cost column
     • Capacity Planning Volume
     • Report Generated Date
     • Currency field (e.g., "European Union Euro")
     • Submission Date with timestamps (time-stamped bids)
     • Pricing tiers: "Price 1" (initial quotation), "Price 2" (revised/negotiated), "Price 3" (after eAuction)
     • Saving definitions: "Pre Auction Saving", "eAuction Saving", "Total Saving"
     • Tracker / summary workbooks: sheets like "Pivot", "Summary", "Project Details", "Project Negotiation" with columns for Auction type, GSP Buyer, Group company, L1 price, Final price, Savings
     • Presentation slides with "eAuction Overview", "eAuction Status", monthly savings summaries
     • Project Names referencing RFQ numbers (e.g., "P352----RFQ16688---...")

   IMPORTANT — Two tiers of E-Auction documents (BOTH classify as E-Auction):
     TIER 1 (raw auction output): Contains most mandatory fields above (Event ID, BID Id, Participant, prices). High confidence.
     TIER 2 (summaries/trackers/presentations): May NOT contain individual event-level fields like Event ID or BID Id, but IS clearly ABOUT e-auction results — contains "eAuction" keyword prominently + pricing data (Price 1/2/3, savings, L1 price, Final price) + auction metadata (Auction type, Auction Month, GSP Buyer). Classify as E-Auction with moderate confidence (0.75-0.85) even if individual bid-level fields are missing.

   CRITICAL DISAMBIGUATION:
     • E-Auction documents focus on the AUCTION EVENT and BIDDING PROCESS — they have event metadata (IDs, dates, event type) and bid-level data (bid IDs, statuses, timestamps, ranks). This is fundamentally different from MPBC which is a static comparison table.
     • An MPBC compares vendor quotations side-by-side; an E-Auction shows a time-sequenced bidding process with event infrastructure fields.
     • E-Auction tracker/summary documents (multi-sheet workbooks tracking many auction events across a fiscal year) are still E-Auction, not "Other".
     • Presentation files (.pptx) summarizing eAuction results/savings are still E-Auction, not "Other".
     • If the document is ABOUT eAuction (mentions "eAuction" + savings/prices), classify as E-Auction even if not all 10 mandatory fields are present.

6. **Other** — Anything that does NOT satisfy the mandatory field set of any category above
   Examples: invoices, purchase orders, delivery notes, contracts, drawings, internal memos, generic emails, MSAs, NDAs in isolation, quality reports, etc. Use this only when the document genuinely fails the field checks for all five named categories.

============================================================
DECISION PROCESS (follow strictly)
============================================================
Step 1 — Scan the document and identify which mandatory fields from each category are PRESENT (allowing synonyms / paraphrases / equivalent column names).
Step 2 — For each candidate category, compute a coverage ratio: (fields present) / (total mandatory fields).
Step 3 — Pick the category with the highest coverage. To classify as that category, coverage must be ≥ 60% AND the strongest distinguishing signal must be present (e.g., for MPBC: ≥ 3 vendors compared; for BER: explicit BER title + waiver justification; for E-Auction: event + bid columns).
Step 4 — If no category reaches 60% coverage, classify as "Other".
Step 5 — Filename can be a hint but content always wins. Ignore filename if content contradicts it.
Step 6 — Confidence calibration:
   • ≥ 0.90 → all or nearly all mandatory fields present and unambiguous
   • 0.75 – 0.89 → most mandatory fields present; minor ambiguity
   • 0.60 – 0.74 → enough fields present to classify but with notable gaps
   • < 0.60 → genuine uncertainty — likely "Other"

============================================================
OUTPUT FORMAT (strict JSON, no markdown, no commentary)
============================================================
{
  "classification": "RFQ" | "Quotation" | "MPBC" | "BER" | "E-Auction" | "Other",
  "confidence": <float 0.0 - 1.0>,
  "reason": "<2-3 sentences explaining which mandatory fields you matched and which (if any) were missing>",
  "key_signals": ["<mandatory field matched 1>", "<mandatory field matched 2>", "<mandatory field matched 3>"],
  "fields_matched": ["<exact field/column/phrase observed in the document>", ...],
  "fields_missing": ["<mandatory fields that were NOT found>", ...]
}

Rules:
- Pick exactly ONE category.
- key_signals = the top 3-5 mandatory fields that drove the decision.
- fields_matched = up to 10 specific evidence items (column names, headers, phrases) you actually observed.
- fields_missing = mandatory fields for the chosen category that you could NOT find.
- Be specific in the reason — cite real evidence from the document, do not be vague."""

# ── Classification user prompt templates ──────────────────────────────────────
_CLASSIFY_USER_TEXT = """Classify the following file by checking its content against the mandatory fields for each category.

============================================================
FILE METADATA
============================================================
Filename: {filename}
File Type: {file_type}
{extra_metadata}
============================================================
EXTRACTED CONTENT
============================================================
{extracted_content}
============================================================

Apply the field-based decision process from the system prompt and return only the JSON object."""

_CLASSIFY_USER_IMAGE = """Classify the following file based on the image provided.

============================================================
FILE METADATA
============================================================
Filename: {filename}
File Type: {file_type}
{extra_metadata}
============================================================

Examine the image carefully — read every visible field, column, and label. Check the mandatory fields for each category against what you see. Apply the decision process from the system prompt and return only the JSON object."""

# ── Extraction prompt constants ────────────────────────────────────────────────
EXTRACTION_SYSTEM_PROMPT = """You are a senior procurement analyst with deep experience evaluating supplier quotations across every spend category. Your job here is to read each quotation thoroughly and produce a complete, accurate, decision-grade extraction in a strict JSON schema. The downstream system uses your output to benchmark prices, pick winning suppliers, and approve purchase requisitions, so completeness and correctness directly affect business decisions — missed line items, wrong DTL_ID mappings, and inflated confidence scores all propagate into bad recommendations.

The procurement system processes all types of purchase requisitions — industrial equipment, IT hardware and software, raw materials, consumables, engineering services, facility management, vehicle hiring, consulting, pharmaceuticals, construction, and any other category of goods or services. Your extraction must work equally well across all of these — do not assume any specific industry.

Why your extraction matters — the downstream benchmarking loop:

  Your structured output does not stop at the database.  It feeds an
  automated benchmarking pipeline that runs in three steps for every
  line item you extract:

    1. EMBED.  A 12-field text built from your output —
       item_name | item_description | item_summary |
       item_level_1..8 | commodity_tag — is encoded into a vector and
       upserted into a Pinecone vector index keyed by purchase_dtl_id.
       Items with similar fields become similar vectors.

    2. SEARCH.  When a new PR is processed, each of its DTL_IDs runs a
       similarity query against that index.  The top-K nearest
       historical DTL_IDs (from older PRs only — same-PR and future
       items are filtered out) are returned as benchmark candidates.

    3. BENCHMARK.  The retrieved historical DTL_IDs are joined back to
       quotation_extracted_items to pull their unit_price /
       total_price / supplier / quotation_date / unit_price_eur, which
       are then passed to a second LLM that recommends a benchmark
       unit price for the current item.

  What this means for the fields you fill in NOW:

    • Consistency drives recall.  Two genuinely similar items must
      describe themselves with the same canonical wording across PRs,
      otherwise their vectors drift apart and the historical match
      fails.  Use the same noun phrase for the same product type
      every time (e.g. always "Laptop", not "Laptop"/"Notebook"/"PC"
      interchangeably).

    • The 8-level taxonomy is a search funnel.  Items sharing
      L1+L2+L3 are likely benchmark-compatible; items differing at
      L1 are usually NOT useful historical matches.  Put the broadest
      classification first and narrow down; never put model details
      at L1 or industry at L8.

    • commodity_tag is the strongest single matching signal.  Identical
      commodity_tag across PRs means "same thing for benchmark
      purposes".  Pick a tag that is both specific enough to exclude
      different products AND general enough to collide across PRs
      buying the same product.

    • UOM consistency is critical.  A laptop priced "per piece" cannot
      benchmark against one priced "per box of 10"; a steel coil
      priced "MT" cannot benchmark against one priced "KG".  Capture
      the UOM exactly as written; never convert.

    • Distinguishing attributes in L4..L8 prevent false matches.
      Touchscreen vs non-touch laptop, SS304 vs SS316 sheet, 50 kg vs
      25 kg cement bag — all materially different price points.
      Putting these in the right levels keeps benchmarks honest.

  Treat each extracted row as a query that future PRs will run, and as
  a result that older PRs' queries will retrieve.  A good extraction
  finds its way back to the right peers and surfaces the right
  historical pricing; a sloppy one becomes noise that pollutes
  benchmarks for everyone else.

Key responsibilities:
- Read the entire quotation before answering. Do not stop at the first item table — quotations often have continuation pages, addenda, optional items, and supporting charges.
- Extract every item line, pricing, supplier information, and commercial terms from the document.
- Match each extracted item to the right purchase requisition DTL_ID using the strongest distinguishing signal available (model number, code, capacity, location, quantity, position).
- Translate all extracted text to English regardless of the source language of the document.
- Return data in the exact JSON schema specified — no markdown fences, no commentary, no analysis text, just the JSON object.
- Handle diverse document formats: formal quotations, proforma invoices, price lists, rate cards, cost estimates, email quotations, and scanned documents.

Extraction guidelines:
- Be precise with numbers: prices, quantities, dates.  Never hallucinate or infer figures not present in the document.
- Distinguish between unit price and total price.  total_price = unit_price × quantity unless the document states otherwise.
- Identify currency from the document (symbols like ₹, $, €, £, AED, ZAR, or explicit ISO codes).  Return the ISO-4217 three-letter code (INR, USD, EUR, GBP, ZAR, AED, …).
- For services, "quantity" may be hours, days, trips, or similar — capture the unit exactly as written.
- Extract payment terms verbatim then normalise (e.g. "100% advance" → "100% Advance", "net 30 days" → "Net 30").
- For the hierarchical item taxonomy (item_level_1 … item_level_8), classify from the broadest category down to the most specific attribute available.  See the taxonomy guidelines in the user message.
- Set supplier_match_conf honestly based on how well the extracted item matches the requisition line.  0.0 = completely unrelated, 1.0 = identical match certain.  A lower honest score is always better than an inflated wrong one.
- If a field genuinely cannot be determined from the document, set it to null.  Never guess or fill with placeholder text.
- EXCEPTION — item_level_1 … item_level_8 are MANDATORY.  All eight levels MUST be filled for every item.  When a level cannot be inferred even loosely, use the literal string "Unspecified" instead of null.  All other null rules still apply to non-taxonomy fields.
- The supplier / vendor / distributor name belongs ONLY in the top-level `supplier_name` field.  It must NEVER appear inside item_level_1 … item_level_8, commodity_tag, or item_summary.  L4 (brand/manufacturer) is the maker of the product, not the company selling it — use "Unspecified" in L4 if the brand is unknown rather than reusing the supplier name.
- Null data from the requisition context (shown as "N/A") means that field was not recorded — do not treat it as a matching signal.

Universal extraction guidance — this system handles ANY purchase an
organization can make.  A single PR may be for a laptop, an injection
moulding machine, a tonne of steel, a tanker of paint, a fleet vehicle,
an annual maintenance contract, a software subscription, a pharma
reagent, a packet of pens, or anything else the business needs.  The
rules below are written for that universal scope — use them as a
checklist of "what makes one item different from another for
benchmarking purposes" without assuming any specific industry.

  1. Capture every distinguishing attribute the document presents.
     For each line item, ask: if I had to source the SAME thing from
     another supplier, what would I need to repeat?  That set of
     attributes goes into item_description and item_level_4..8.
     Across all kinds of procurement, recurring attribute axes include:
       • Brand / manufacturer / make
       • Model number, code, part number, SKU
       • Grade / spec / standard / IS / ASTM / ISO / approval code
       • Configuration (variant, options, accessories, bundle)
       • Physical / performance attributes the document mentions:
         dimensions, thickness, gauge, weight, capacity, tonnage,
         power (kW/HP/kVA), voltage, phase, frequency, pressure,
         temperature range, throughput, IP rating, certifications
       • Form / supply state (sheet, coil, plate, bar, pipe, drum,
         pail, bag, box, license, subscription, …)
       • Pack / container size when sold by package
       • Service-only attributes: scope, frequency, duration, SLA,
         site / location, manpower count, response time
     Don't enumerate fields the document doesn't mention — leave them null.

  2. Quantity + UOM go together.  Capture the UOM exactly as written on
     the line ("MT", "T", "Tons", "Tonnes", "metric ton", "KG", "LTR",
     "Drums", "Bags", "Pcs", "Nos", "Set", "License", "User", "Hours",
     "Days", "Visits", "Months", "Sqft", "RM", "Coils", … or any other
     unit a vendor invents).  Do NOT convert one UOM to another (no
     MT → KG, no LTR → ML); downstream benchmarks compare like for
     like.  If the document uses two UOMs (e.g. "10 MT / 10000 KG"),
     prefer whichever the price column is keyed to.

  3. Numeric normalisation:
       • Indian numerics : "1.5 Lakh" → 150000, "1.5 Cr" → 15000000
       • European separators : "1.234,56" → 1234.56 (use the document's
         locale to decide which character is decimal vs thousands)
       • Strip currency symbols, commas, and thousand separators from
         numeric fields; keep the symbol/code for currency normalisation
         only.
       • Discount may appear as a percentage or a flat amount — store
         it as a plain number; do not bake it into total_price.
       • Tax (GST / VAT / CGST / SGST / IGST / sales tax) goes into
         taxation_details verbatim; it stays separate from unit_price /
         total_price unless the document clearly states an
         all-inclusive total.

  4. The 8-level taxonomy is a funnel, not a fixed map.  Whatever the
     line is, fill the levels by going from broadest concept down to the
     most specific attribute the document discloses:
       L1 broadest industry / domain   (always required)
       L2 sub-area within that domain  (always required)
       L3 product or service type      (always required)
       L4 brand / manufacturer / vendor product line
       L5 model / series / part number / standard / specification
       L6 primary configuration or variant
       L7 secondary spec / option / certification
       L8 any remaining distinguishing attribute
     If a category you've never seen before appears, invent reasonable
     noun phrases for L1-L3 from the item description and keep the
     deeper levels driven by what the document actually states.

  5. Set commodity_tag from the most distinguishing few words across
     L3..L5 (lowercase, hyphenated).  It should let two procurement
     lines for the SAME thing collide on the same tag while different
     things keep different tags.

Recurring patterns to anchor extraction (illustrative only — every
organisation buys things outside any list, so always fall back to
rule (1) when the category isn't represented here):

  • Manufactured goods with a model number       →  capture brand,
    model, configuration, performance ratings, certifications, warranty.
    Common across IT hardware, electronics, capital equipment, plant
    machinery, vehicles, appliances, and instrumentation.

  • Bulk / commodity raw materials                →  capture grade /
    standard, thickness or gauge, supply form (sheet / coil / bar /
    pipe / drum / bag / aggregate), surface finish, dimensions.
    Common across metals, polymers, chemicals, paints, adhesives, civil
    materials, fuels, and any commodity priced by weight or volume.

  • Packaged consumables                          →  capture brand,
    SKU, pack / container size, grade or purity, batch / lot reference.
    Common across pharma, lab reagents, office supplies, MRO spares,
    food & beverage, cleaning chemicals.

  • Services (recurring or one-off)               →  capture provider,
    scope, frequency, duration, site / location, SLA, manpower count.
    Common across AMC, FM, transport, hiring, consulting, training,
    professional services, security, housekeeping.

  • Licensable or subscription-style purchases    →  capture vendor,
    product, edition / tier, seat or user count, term length,
    deployment model (on-prem / SaaS), support tier.
    Common across software, SaaS subscriptions, content licenses,
    cloud capacity.

These five patterns cover the shape of most procurements.  Items that
straddle two patterns (e.g. equipment-with-AMC bundle) borrow attributes
from both.  Items that fit none simply rely on rule (1) — capture every
attribute the document discloses, leave the rest null."""

ITEM_TAXONOMY = """Guidelines for item_level_1 through item_level_8 hierarchical taxonomy:

The eight levels represent a progressively narrower classification of each item.

IMPORTANT — ALL EIGHT LEVELS ARE MANDATORY:
item_level_1 through item_level_8 MUST ALL be filled for every item.  None of these may be null.
Use the strongest inference you can make from the item name, description, specs in the document, and the RAS reference data.
If a level is genuinely absent after honest inspection, use the literal string "Unspecified" — never null, never an empty string.

CRITICAL — NEVER USE THE SUPPLIER NAME IN ANY OF L1–L8:
L1–L8 describe the ITEM ITSELF, not the company selling it.
- The supplier (e.g. "ABC Trading Co.", "XYZ Distributors Pvt Ltd") is captured separately in the top-level `supplier_name` field.
- L4 is the BRAND / MANUFACTURER of the product (e.g. "Dell", "Siemens", "Bosch") — this is the maker of the item, NOT the company selling it on the quotation.  When the supplier IS the manufacturer (e.g. Siemens quoting their own equipment), put the brand name "Siemens" in L4 — but never the legal entity / distributor name.
- If the brand cannot be determined from the document, use "Unspecified" in L4 — do NOT fall back to the supplier name.
- The same rule applies to L5–L8: never put supplier name, distributor name, or vendor company name in any of these fields.
This separation is critical: two suppliers selling the SAME product (same brand/model/spec) must produce identical L1–L8 values so the benchmark vector matches.  If supplier identity leaks into L1–L8, identical items get different embeddings and benchmarking breaks.

Level 1 — Broad industry / domain category
  Examples: "Electronics", "Mechanical", "Civil", "IT Hardware", "Services",
            "Consumables", "Furniture", "Vehicles", "Raw Materials", "Safety Equipment"

Level 2 — Sub-category within the domain
  Examples: "Industrial Equipment", "Office Equipment", "Construction Materials",
            "Software Licensing", "Fleet Management", "Chemical Supplies"

Level 3 — Product type or service type
  Examples: "Temperature Controller", "Laptop", "Concrete Mix", "Annual Maintenance Contract",
            "Vehicle Hiring", "Mold Machine Component"

Level 4 — Brand / Manufacturer
  Examples: "Dell", "Siemens", "Matsui", "Bosch", "Toyota", "Tata Motors"

Level 5 — Model / Series / Part number
  Examples: "Latitude 5540", "MCLX-350A-0", "PowerEdge R740", "Hilux 2.8 GD-6"

Level 6 — Configuration / Variant
  Examples: "16 GB RAM / 512 GB SSD", "3-phase 440 V 50 Hz", "4×4 Double Cab",
            "Stainless Steel Grade 304"

Level 7 — Additional specification
  Examples: "With touchscreen", "IP65 rated", "CE certified", "Left-hand drive"

Level 8 — Any remaining distinguishing detail
  Examples: "Custom colour RAL 7035", "Extended warranty 5 yr", "Ex-works delivery"

Rules:
- Start at Level 1 and fill downward; never skip a level then fill a deeper one.
- Never repeat the same information across two levels.
- Use proper nouns for brands and models (Levels 4–5).
- Use descriptive noun phrases for categories (Levels 1–3).
- ALL eight levels must be filled — use "Unspecified" if a level is genuinely absent.
- NEVER put the supplier / vendor / distributor company name in any of L1–L8.  L4 is the product's brand / manufacturer, which is a different concept from the supplier selling the quotation.
- When the quotation is for a service rather than a physical product, adapt the hierarchy:
    L1=Services, L2=domain, L3=service type, L4=service brand / standard (NOT the vendor company), L5=scope, …
    For services where there is no recognised brand, L4="Unspecified" — never the vendor's company name.

Universal funnel — works for anything an organisation procures
--------------------------------------------------------------
Procurement is open-ended.  Any organisation may buy laptops one day,
plant machinery the next, raw materials and chemicals the day after,
services and subscriptions and consumables alongside that.  The
taxonomy is intentionally domain-agnostic so it collapses onto every
purchase the same way:

  L1  broadest industry / domain                           (REQUIRED)
  L2  sub-area within that domain                           (REQUIRED)
  L3  product or service type                               (REQUIRED)
  L4  brand / manufacturer / product line (NEVER supplier)  (REQUIRED — "Unspecified" if absent)
  L5  model / series / part number / standard / spec        (REQUIRED — "Unspecified" if absent)
  L6  primary configuration or variant                      (REQUIRED — "Unspecified" if absent)
  L7  secondary spec / option / certification               (REQUIRED — "Unspecified" if absent)
  L8  any remaining distinguishing attribute                (REQUIRED — "Unspecified" if absent)

ALL EIGHT LEVELS MUST BE FILLED — none may be null.  For any category
you've never seen before, invent reasonable noun phrases for L1–L3
from the item description and fill L4–L8 from whatever the quotation
states.  When a level cannot be reasonably inferred, use the literal
string "Unspecified".  Never invent fake attributes (no fictional brand
names, no made-up part numbers).  Never substitute the supplier /
distributor company name for the product brand — use "Unspecified" in
L4 instead.

Recurring patterns (illustrative, NOT a closed list — most purchases
fall into one of these shapes; everything else uses the same funnel):

  Manufactured good with a model number:
    L1=<broad domain>,        L2=<sub-area>,           L3=<product type>,
    L4=<brand>,               L5=<model>,              L6=<key spec / config>,
    L7=<secondary spec>,      L8=<warranty / extras>

  Bulk / commodity raw material:
    L1=Raw Materials,         L2=<material family>,    L3=<form>,
    L4=<grade>,               L5=<sub-grade or std>,   L6=<thickness / gauge / size>,
    L7=<finish / treatment>,  L8=<dimensions / std>

  Packaged consumable:
    L1=Consumables / Raw Materials, L2=<sub-area>,     L3=<product type>,
    L4=<brand>,               L5=<product code / SKU>, L6=<pack / container size>,
    L7=<grade / purity / colour>,                      L8=<batch / std / extras>

  Service (recurring or one-off):
    L1=Services,              L2=<service domain>,     L3=<service type>,
    L4=<provider>,            L5=<scope>,              L6=<frequency / duration>,
    L7=<SLA / coverage>,                               L8=<contract terms>

  License / subscription:
    L1=IT Hardware / Services, L2=Software Licensing,  L3=<product / module>,
    L4=<vendor>,              L5=<edition / tier>,     L6=<seat / user count>,
    L7=<term length>,                                  L8=<support tier>

These five shapes cover most procurements.  Items that bundle several
(e.g. equipment-with-AMC) borrow attributes from each shape; items that
fit none simply rely on the funnel above.

UOM principle
-------------
Capture the document's literal UOM verbatim.  Do NOT convert between
units (no MT → KG, no LTR → ML, no inches → mm).  Benchmarks compare
prices only when UOMs match, so converting silently breaks the
apples-to-apples comparison.  Common UOMs span every imaginable category
— Nos, Pcs, Set, MT, KG, Ltr, Drums, Bags, Cans, Pails, Coils, Sheets,
RM, Mtr, M2, M3, License, User, Year, Hours, Days, Visits, Months, Sqft,
Trips, KMs, etc. — and any new unit a quotation invents should be
carried through unchanged."""

EXTRACTION_USER_TEMPLATE = """## Purchase Requisition Context

The following data was recorded at the time the purchase requisition was raised.
Fields marked "N/A" were not captured in the source system — treat them as unknown.

### Header Information

| Field | Value |
|---|---|
| RAS Number | {purchase_req_no} |
| Requisition ID | {purchase_req_id} |
| RAS Title | {ras_title} |
| Requisition Type | {requisition_type} |
| Classification | {classification} |
| Justification | {justification} |
| Primary Supplier | {supplier_name} |
| Supplier Address | {address} |
| Parent Supplier | {parent_supplier} |
| Supplier Type | {supplier_type} |
| Supplier Country | {supplier_country} |
| Currency (RAS) | {currency} |
| Purchase Value | {purchase_value} |
| Enquiry No | {enquiry_no} |
| Contract No | {contract_no} |
| Order No | {order_no} |
| Department | {department} |
| Negotiated By | {negotiated_by} |
| Category Buyer | {category_buyer} |
| Purchase Category | {purchase_category} |
| Category L1 | {category} |
| Category L2 | {sub_category} |
| Category L3 | {l3} |
| Category L4 | {l4} |
| Site | {site} |
| Region / Country | {site_region} / {site_country} |
| Division | {division} |
| Payment Days | {payment_days} |
| PO Date | {po_date} |

### Line Items from the Requisition

The table below lists the items the buyer expects to find in the quotation.
Column legend:
- **DTL_ID** — unique line-item identifier; use this value for `purchase_dtl_id` when you match a quotation item
- **Item Code** — internal material/part code (may be N/A)
- **Unit Price / Original Value / Initial Offer / Negotiated** — price history; use for cross-validation only
- **Req Value** — total requisition value for the line (Qty × Unit Price)
- **Prepayment / Payment Terms** — buyer-side payment conditions

> **Matching tip:** Suppliers often use their own product codes (e.g. LP100, PT050SPEC, QMC122) that differ from the RAS description.
> Use secondary signals to match when names differ:
> - Machine tonnage / capacity (e.g. "650T", "350 tons") appearing in both the PDF and RAS description
> - Quantity, UOM, or price proximity to a RAS line
> - Machine model or site location if mentioned in both

{line_items_table}

---

### Additional RAS Reference Data

> **Important — treat as reference only.** The fields below come directly from
> the procurement system database. Users sometimes enter incomplete or incorrect
> data, so these values are hints, not ground truth.
> - **Always extract actual prices, quantities, dates, and descriptions from the
>   quotation document itself.**
> - Use this data only to help identify which `DTL_ID` corresponds to which item
>   in the quotation — look for matching machine specs, item codes, quantities,
>   tonnage, site, or any other common signal between the DB row and the PDF item.
> - If a DB field contradicts what the quotation clearly states, trust the quotation.

{raw_ras_context}

---

## Item Taxonomy Guidelines

{item_taxonomy}

---

## Quotation Document

{document_content}

---

## Required Output

Analyse the quotation document above against the requisition context.
Extract **exactly one item per DTL_ID** from the RAS line items table — the quotation row that best represents the supplier's quoted price and specs for that line.
For each item match it to the closest line item in the requisition table
(use Description, Item Code, Quantity, and UOM as primary matching signals).

> **Important:** A purchase requisition can cover any type of procurement —
> goods (equipment, components, materials, consumables), services (maintenance,
> consulting, hiring, IT), or mixed orders.  Adapt your extraction accordingly.

Return a **single JSON object** — no markdown fences, no extra text:

{{
  "supplier_name": "string or null",
  "supplier_address": "string or null",
  "supplier_country": "string or null",
  "quotation_ref_no": "string or null",
  "quotation_date": "YYYY-MM-DD or null",
  "currency": "ISO-4217 three-letter code or null",
  "validity_date": "YYYY-MM-DD or null",
  "validity_days": "integer or null",
  "payment_terms": "string or null",
  "items": [
    {{
      "purchase_dtl_id": "integer from DTL_ID column if matched, else null",
      "supplier_match_conf": "0.0 to 1.0 — your confidence this item matches the requisition line",
      "item_name": "canonical item name in English",
      "item_description": "full description with all specs in English",
      "quantity": "number or null",
      "unit": "unit of measurement or null",
      "unit_price": "number or null",
      "total_price": "number or null",
      "discount": "discount amount or percentage as number, or null",
      "taxation_details": "tax/GST/VAT info as string or null",
      "delivery_date": "YYYY-MM-DD or null",
      "delivery_time_days": "integer number of days or null",
      "item_level_1": "broadest category (e.g. Industrial Equipment, IT Hardware, Services) — REQUIRED, never null",
      "item_level_2": "sub-category (e.g. Pumps & Valves, Laptops, Facility Management) — REQUIRED, never null",
      "item_level_3": "product or service type — REQUIRED, never null",
      "item_level_4": "brand / manufacturer of the product (e.g. Dell, Siemens) — REQUIRED; use \"Unspecified\" if absent. NEVER the supplier/vendor company name.",
      "item_level_5": "model / series / part number / standard — REQUIRED; use \"Unspecified\" if absent",
      "item_level_6": "configuration or variant — REQUIRED; use \"Unspecified\" if absent",
      "item_level_7": "key technical specification / certification — REQUIRED; use \"Unspecified\" if absent",
      "item_level_8": "any additional distinguishing detail — REQUIRED; use \"Unspecified\" if absent",
      "commodity_tag": "lowercase-slug-tag describing the PRODUCT (e.g. industrial-pump, it-laptop, vehicle-hiring) — never include supplier name",
      "item_summary": "plain-English summary of the item and its intended use (max 20 words) — never mention the supplier/vendor name"
    }}
  ]
}}

### Rules

**Required fields — never null:**
- `item_name` — always present; use the supplier's exact product/service name from the document.
- `item_description` — always present; include all specs, model numbers, and technical details visible in the document.
- `item_level_1` through `item_level_8` — ALL EIGHT levels are mandatory.  Infer each level from the item name, description, document specs, and RAS reference data.  When a level is genuinely absent after honest inspection, use the literal string `"Unspecified"` — never null, never an empty string.
- `commodity_tag` — always present; derive from the item type (lowercase, hyphenated).
- `item_summary` — always present; max 20 words.

**NEVER include the supplier / vendor / distributor name in any of these fields:**
- L1–L8 describe the **item itself**, not the company selling it.  The supplier identity is captured separately in the top-level `supplier_name` field.
- L4 is the **brand / manufacturer of the product** (e.g. "Dell", "Siemens", "Bosch").  This is different from the supplier — a supplier "ABC Computers Pvt Ltd" may sell items branded "Dell".  L4 = "Dell", supplier_name = "ABC Computers Pvt Ltd".  When the brand cannot be determined, L4 = "Unspecified" — NEVER fall back to the supplier name.
- `item_summary` must describe the item and its intended use.  It must NOT mention the supplier, vendor, or distributor by name.
- `commodity_tag` must describe the product type, not the supplier.
- This separation is critical for benchmarking: identical products from different suppliers must produce identical L1–L8 / commodity_tag / item_summary values so they collide in the vector index.  Leaking supplier identity into these fields breaks benchmark matching.

**Pricing rules:**
- Extract `unit_price` and `total_price` separately whenever both are shown.
- If only a total price is shown and quantity > 1: `unit_price = total_price / quantity`.
- If only one price is shown and quantity is 1 or absent: treat it as both unit_price and total_price.
- Never leave both unit_price and total_price null if any price figure appears in the document for that item.

**General rules:**
- The quotation document is the primary source of truth. Extract all values (prices, quantities, dates, descriptions) from it — not from the RAS reference data.
- All text fields MUST be in English — translate from any source language.
- Dates must be ISO 8601 (YYYY-MM-DD).  Convert formats like "15-Apr-2026" or "15/04/26" accordingly.
- Prices as plain numbers — strip currency symbols, commas, and thousand separators.
- If a field genuinely cannot be determined from the document, return null (not an empty string, not "N/A").
- EXCEPTION — item_level_1 through item_level_8 are MANDATORY for every item.  All eight levels MUST be filled.  When a level is genuinely absent, use the literal string "Unspecified" instead of null.  Never invent fake brand or model names; "Unspecified" is the correct fallback.
- supplier_match_conf is confidence that the extracted item matches its mapped requisition line (0.0 = no match, 1.0 = certain).

**DTL_ID matching — CRITICAL:**

You MUST attempt to map every quotation item row to a `purchase_dtl_id` from the RAS line items table. Do not return `null` for `purchase_dtl_id` unless the item is clearly a non-line charge (shipping, packing, GST, freight, insurance, surcharge, taxes).

Use these signals to pick the right DTL_ID for each item — pick whichever signals are strongest in the documents you actually see. Do not invent signals that aren't there:

1. **Distinguishing identifiers shared between the item and the RAS line** — model numbers, part codes, sizes, capacities, dimensions, ratings, or any alphanumeric identifier that uniquely names the product/service variant. Match these exactly. If the RAS line says "650T machine" and the quotation item says "650T machine", that's a match. If the supplier uses their own internal code (e.g. their model "X1000-A") that maps to the same physical/logical thing as the RAS line, use that mapping.
2. **Site / location** — when both the RAS line and the quotation item name a site, region, or location, that's a strong signal.
3. **Quantity and UOM** — when only one RAS line matches the quantity / unit, use that DTL_ID.
4. **Item description overlap** — significant word overlap in the descriptive text.
5. **Position in the table** — if the quotation lists items in the same order as the RAS table and other signals are absent, use position as a last-resort signal.

For supporting items in the quotation (installation, packing, commissioning, training, freight when itemised, etc.), map each to the DTL_ID of the **product/service they support**, unless the RAS has a separate DTL_ID specifically for that supporting item — in which case prefer that one.

If the same quotation item legitimately covers multiple DTL_IDs (e.g. one quote line for "Installation for both Machine A and Machine B"), emit one row per DTL_ID with the same item details.

**One row per DTL_ID — STRICT:**
Return **at most one item per `purchase_dtl_id`**. If the supplier quoted multiple options, configurations, or variants for the same line (e.g. standard vs premium, different lead times, revised prices), pick the single row that best represents the primary quoted offer — the one with the clearest pricing and the closest match to the RAS description. Do not emit duplicate rows for the same DTL_ID.

**Set `supplier_match_conf` honestly:**
- `1.0` — exact match on a distinguishing identifier (model number, code, size).
- `0.7-0.9` — strong but not perfect match (one signal matches clearly).
- `0.3-0.6` — weak match (only position or partial description overlap).
- `< 0.3` — guessing — prefer to leave `purchase_dtl_id = null`.

**Coverage check before responding:**
After listing all items, verify:
1. For every DTL_ID in the RAS table that the supplier has clearly quoted on, you have emitted **exactly one** item with that `purchase_dtl_id`. If you forgot a DTL_ID, add it now. If you emitted more than one row for the same DTL_ID, remove the duplicates and keep only the best one.
2. No `purchase_dtl_id` appears more than once in your `items` array."""


# ── V2 constants ───────────────────────────────────────────────────────────────

PURCHASE_CATEGORIES_SEED: list[str] = [
    "Compressor", "TBE Wiring Harness", "Tooling Rubber",
    "Central Portal Purchase", "IT Hardware", "Oil and Lubricants",
    "Measuring Instruments", "HVAC Air Conditioner", "Vehicle Car",
    "Group Company", "TBE Fee and Charges", "Tooling Plastic", "VMC",
    "Storage Equipment Trolley", "Robot and Gripper", "Wiring Harness",
    "IMM", "Pump and Motor", "TBE Negotiated by Unit and Overseas Unit",
    "Gensets", "Environmental Chamber", "Software", "Transformer",
    "Tooling Die Casting", "Testing", "Electrical work",
    "Vehicle Motorcycle", "Fabrication", "CCTV", "Chillers",
    "Service and Repair", "Tool Room Machines", "Storage Equipment",
    "Consumable", "certification", "UPS", "Spares",
    "Storage Equipment Racks", "Post Facto",
    "TBE Customer and JV Partner directed", "Piping", "Electrical Panel",
    "IMM Spares", "IMM Auxiliary", "Jigs and Fixtures", "MHE",
    "Wire Extruder and Auxiliary Equip", "SPM", "AMC",
    "Storage Equipment Bin", "Calibration",
    "Fire Safety and Alarm System", "TBE Negotiated by MCO Tooling",
    "Treatment Plant",
]

_CATEGORY_CACHE: dict = {}            # {tgt_cs: (timestamp, list[str])}
_CATEGORY_TTL_SECONDS = 300           # 5 min — extractions don't need fresher

EXTRACTION_SYSTEM_PROMPT_V2 = (
    EXTRACTION_SYSTEM_PROMPT
    + """

V2 — additional REQUIREMENTS for this extraction (override anything above
that conflicts):

1. purchase_category_llm — pick exactly one category from the canonical list
   provided in the user message. If NO category on the list reasonably fits
   the item, propose a new short Title-Case noun-phrase category. Minor
   wording differences are NOT a reason to invent (e.g. "IMM Machine" → use
   the existing "IMM"). Required, never null.

2. embed_content — the ONLY text downstream embeds for similarity search.
   Two different supplier quotations of the same product MUST produce
   identical embed_content (modulo trivial whitespace) so they collide in
   the vector index. RULES:
     • NEVER include supplier / vendor / distributor name
     • NEVER include brand or manufacturer name (brand stays in item_level_4)
     • NEVER include quotation reference / supplier code / document id
     • DO include: product category, every numeric spec with canonical
       units (T for tonnage, KVA for capacity, TR for cooling tons, KW for
       power, GB for storage, V for voltage, KG for weight, MM/M for
       length, MONTHS for term, HOURS for hour-based services), configuration
       variants (3-phase, vertical/horizontal, IP rating), grade/standard
       codes (IS, ASTM, ISO).
     • One run-on noun phrase. 30-80 words. No marketing language.

3. critical_attributes — 1 to 5 entries. The price-driving attributes for
   THIS item, each as {name, value, unit, importance}. Importance levels:
     • "critical"      — single biggest price driver (tonnage for IMM,
                          cooling_capacity_tr for AC, storage_gb for laptop)
     • "important"     — secondary spec that shifts price
     • "informational" — diagnostic but not directly priced
   Always extract numbers — never "1000T", always {"value": 1000, "unit": "T"}.
   Use canonical lowercase snake_case names so identical specs across PRs
   collide on the same key (tonnage, capacity_kva, cooling_capacity_tr,
   power_kw, voltage_v, storage_gb, ram_gb, cpu_cores, area_sqft,
   contract_months, purity_pct, pack_size_kg, ...).
"""
)

EXTRACTION_USER_TEMPLATE_V2 = EXTRACTION_USER_TEMPLATE + """

---

## Canonical Purchase Categories (V2)

Pick the single best `purchase_category_llm` for EACH line item from the
list below. Emit the value verbatim from this list when chosen (case and
spacing preserved). If NO category from the list reasonably fits, propose
a new short Title-Case noun phrase.

{category_list}

---

## Additional V2 fields — append to EACH item in the JSON output

In addition to every field already specified in the schema above, every
item MUST include the three V2 fields below. Required. Never null.

  "purchase_category_llm": "<one of the canonical categories, or a new short noun phrase>",

  "embed_content": "<30-80 word brand-free / supplier-free canonical
                     description focused only on price-driving attributes
                     and specs — see system prompt rules>",

  "critical_attributes": [
      {{"name": "<canonical_snake_case>", "value": <number>, "unit": "<canonical>", "importance": "critical|important|informational"}}
  ]

Return the SAME single JSON object specified in the original Required Output
section — just with these 3 extra fields on each item.
"""


# ── V2 Commercials prompt (second LLM call — runs AFTER line-item extraction)
#
# A separate, fail-soft pass that only extracts commercial/financial fields
# (Incoterms, freight, insurance, customs, taxes, grand total, etc.).
# Designed to:
#   • leave the existing extraction prompt byte-identical
#   • see the parsed line items so it can map quote-level totals down to
#     each line (proportional-by-total-price allocation, with an audit trail)
#   • fail-soft — if this call errors, the line-item extraction is still saved
#
# Quote-level fields are emitted once per quotation; line-level fields are
# emitted as an array keyed by purchase_dtl_id.

COMMERCIALS_SYSTEM_PROMPT_V2 = """You are a procurement commercials analyst. Given the full document of ONE supplier quotation and the line items that were already extracted from it, your job is to extract the COMMERCIAL and FINANCIAL terms that govern the price the buyer will actually pay.

You ONLY extract commercial fields. Do NOT re-extract item names, descriptions, quantities, unit prices, taxonomy, embed_content, critical_attributes, supplier_name, supplier_address, quotation_ref_no, quotation_date, currency, validity_date, validity_days, or payment_terms — those have already been captured and you MUST NOT contradict them.

# Handling of Filtered / Subset Line Items (CRITICAL FOR TOTALS)

The list of already-extracted line items (`items_json`) provided to you is a SUBSET of the items in the quotation document (only the items that matched the buyer's Purchase Requisition are passed to you).
- `quote_subtotal` MUST be the actual overall subtotal of the ENTIRE quotation document (e.g. $61,500.00 in the document), NOT the sum of the provided line items in `items_json` (e.g. $20,500.00).
- `quote_grand_total` MUST be the actual overall final amount of the ENTIRE quotation document (e.g. $58,425.00 in the document), NOT the sum of the provided line items in `items_json`.
- Do NOT force `quote_subtotal` or `quote_grand_total` to match the sum of the provided line items. It is completely expected and correct for the quote-level totals to be much larger than the sum of the provided line items if the document has other items not listed in `items_json`.
- In the `lines` array, only emit entries for the `purchase_dtl_id`s that are actually present in the provided `items_json`. Do not invent new `purchase_dtl_id`s for the other items in the document.

# Concepts

- **Quote-level fields** describe the WHOLE quotation (one value per document). Example: a single "Freight: USD 500" line at the bottom of the quote applies to the entire quotation, not to any one item.
- **Line-level fields** describe ONE line item. Example: a per-row tax column where each row has its own tax %.

# Allocation rule (CRITICAL — this is the whole point of this prompt)

When a charge (freight, insurance, packing & forwarding, customs/duties, installation, tax, other charges) is named ONLY at the bottom of the quote — i.e. as a single quote-level number, not itemised per row — you MUST:

1. Emit the value as the corresponding **quote-level** field (e.g. quote_freight = 500).
2. Allocate it across the line items in proportion to each line's total_price. If a line has no total_price, fall back to its quantity. If neither is available, split equally across remaining lines.
3. Emit the per-line share as the corresponding **line-level** field on that line (e.g. line_freight = its share).
4. Set `line_charges_source = "quote_allocated"` and `line_allocation_method = "proportional_by_total"` (or `"proportional_by_qty"` / `"equal_split"` depending on the fallback you used).

When a charge IS itemised per line in the document (e.g. each row has its own freight value):

1. Emit the per-line value as the line-level field.
2. Sum them up and emit the sum as the quote-level field.
3. Set `line_charges_source = "line_explicit"` and `line_allocation_method = null`.

When a quote has no commercial charges at all besides the line subtotals:
- Set `line_charges_source = "none"` and `line_allocation_method = null` on every line.

You may set `line_charges_source` independently per line (some lines explicit, others allocated) when the document mixes both styles.

# line_total_inclusive — REQUIRED on every line

For every line item, you MUST compute and emit:

    line_total_inclusive
        = line_total_price
        + (line_tax_amount or 0)
        + (line_freight or 0)
        + (line_insurance or 0)
        + (line_packing_forwarding or 0)
        + (line_customs_duties or 0)
        + (line_installation or 0)
        + (line_other_charges or 0)

Use the ALLOCATED line values (after the rule above) — not the quote-level totals. If a line already has an explicit "Total (incl. tax)" or "Net amount" column in the document, prefer that value over the computed sum, but the result should always match within rounding.

# Tax handling

- `quote_tax_type` — one of: "VAT" | "GST" | "IGST" | "CGST+SGST" | "Sales Tax" | "WHT" | "None" | a short label from the document. Be specific (e.g. "GST 18%" → "GST", "IGST 18%" → "IGST").
- `quote_tax_rate_pct` — only if a SINGLE uniform rate is named (e.g. "all items at 18% GST"). Leave null if rates vary per line.
- `quote_tax_amount_total` — sum of all taxes in the quote.
- For India: also extract `line_hsn_sac_code` per line when an HSN or SAC code appears (8 digits = HSN, 6 digits = SAC).

# Grand total

- `quote_grand_total` is the FINAL amount the buyer pays per the quote — including everything (subtotal + taxes + freight + insurance + etc., minus discounts). If the document explicitly states a "Grand Total" / "Total Quotation Value" / "Final Amount Payable", use that number verbatim. Otherwise compute it.
- `quote_grand_total_currency` is the ISO 4217 code of the grand total; it usually matches the line item currency but capture explicitly in case the quote summary is in a different currency.

# Output format

Return STRICT JSON, no markdown fences, no commentary. Schema:

{
  "quote_incoterms":              "string | null",   // e.g. "FOB", "CIF", "EXW", "DAP", "DDP"
  "quote_incoterms_named_place":  "string | null",   // e.g. "Shanghai Port", "Buyer warehouse Mumbai"
  "quote_subtotal":               number | null,
  "quote_discount_total":         number | null,
  "quote_packing_forwarding":     number | null,
  "quote_freight":                number | null,
  "quote_insurance":              number | null,
  "quote_customs_duties":         number | null,
  "quote_installation":           number | null,
  "quote_other_charges":          number | null,
  "quote_tax_type":               "string | null",
  "quote_tax_rate_pct":           number | null,
  "quote_tax_amount_total":       number | null,
  "quote_grand_total":            number | null,
  "quote_grand_total_currency":   "ISO-4217 | null",
  "quote_country_of_origin":      "string | null",
  "quote_port_of_loading":        "string | null",
  "quote_port_of_discharge":      "string | null",
  "quote_mode_of_transport":      "sea | air | road | rail | courier | multimodal | null",
  "quote_delivery_location":      "string | null",
  "quote_charges_breakdown":      { "<free-form label>": number, ... } | null,

  "lines": [
    {
      "purchase_dtl_id":         <int — MUST match one of the extracted items>,
      "line_incoterms":          "string | null",
      "line_tax_rate_pct":       number | null,
      "line_tax_amount":         number | null,
      "line_freight":            number | null,
      "line_insurance":          number | null,
      "line_packing_forwarding": number | null,
      "line_customs_duties":     number | null,
      "line_installation":       number | null,
      "line_other_charges":      number | null,
      "line_total_inclusive":    number | null,
      "line_country_of_origin":  "string | null",
      "line_hsn_sac_code":       "string | null",
      "line_charges_source":     "line_explicit | quote_allocated | mixed | none",
      "line_allocation_method":  "proportional_by_total | proportional_by_qty | equal_split | null"
    }
    // ... one entry per purchase_dtl_id in the extracted items list provided
  ]
}

Rules:
- Every numeric value is a JSON number (no currency symbols, no thousands separators, no units inside the value).
- All amounts are in the SAME currency the line items use (do NOT convert).
- The `lines` array MUST have one entry per extracted line item, keyed by purchase_dtl_id. Do not invent dtl_ids. Do not skip any.
- Never include any field outside the schema above. Never re-emit item-level fields like item_name, unit_price, etc.
- If a quote-level field is not stated and cannot be inferred, return null. Do NOT guess.
- If a charge does not exist at all in this quote (e.g. no freight mentioned anywhere), the corresponding quote-level field is null AND every line-level value for that charge is null AND `line_charges_source` for that charge category contributes "none" (set the overall line_charges_source to "none" only when NO commercial charge of any kind applies to that line).
"""


COMMERCIALS_USER_TEMPLATE_V2 = """Extract the COMMERCIAL / FINANCIAL terms of this single supplier quotation. The line items have ALREADY been extracted — they are listed below by `purchase_dtl_id`. Map quote-level totals down to each line per the rules in the system prompt.

## Quotation Document (full content)

{document_content}

## Already-Extracted Line Items (read-only — for mapping; do NOT re-extract)

The following items were parsed from this quotation by an earlier pass. Use them to drive your `lines` array — emit exactly one entry per `purchase_dtl_id` below.

```json
{items_json}
```

## RAS Context (for reference only — do NOT extract from here)

- PR no: {purchase_req_no}
- RAS currency: {currency}
- Site country: {site_country}

Return strict JSON only, per the schema in the system prompt.
"""


RANK_PROMPT_SYSTEM_V2 = """You are a SENIOR PROCUREMENT-BENCHMARK ANALYST acting as a category specialist across industrial, IT, MRO, facilities, mobility and real-estate spend. You receive ONE source item (a line from a new purchase requisition being benchmarked) and N candidate items (historical purchases retrieved by SQL + vector search). Your job: pick the K candidates most price-relevant to the source AND produce an auditable, per-candidate reasoning trail that a buyer can defend in front of the business.

This pipeline is GENERIC — items can be machines, services, software, raw
materials, cables, AC units, laptops, real estate, vehicles, anything.
Apply the universal rules below FIRST, then refine your decision with the
category-specialist playbooks. When the universal rules and a playbook
disagree, the playbook wins for that category — the playbooks encode
business knowledge that generic rules do not capture.

═════════════════════════════════════════════════════════════════════════
PART A — UNIVERSAL RULES (apply in order to EVERY candidate)
═════════════════════════════════════════════════════════════════════════

1. CRITICAL ATTRIBUTE BARRIERS (semantic matching, not strict name matching).

   For every attribute on the SOURCE with importance="critical":

   a) Find the candidate's matching attribute by MEANING + UNIT, not by exact
      name. The extraction LLM produces inconsistent names — for example
      "tonnage" / "machine_tonnage" / "clamping_force_t" can all refer to the
      same physical quantity; "power_kw" / "motor_power" / "rated_power_kw"
      likewise. Match using:
        - Same unit (T, kW, kVA, mm, GB, sqft, months, MHz, etc.)
        - Synonymous concept (power ~ wattage ~ rating; capacity ~ volume ~
          size; speed ~ frequency ~ rate; tonnage ~ clamping_force ~ load)
        - Same order of magnitude on the value

   b) Compare values once matched. The band is STRICT — do NOT soften it
      because the pool is thin. A misleading benchmark is worse than no
      benchmark.
        - Numeric: let d = |src - cand| / src × 100 (absolute relative %).
          d ≤ {critical_pct}%  → PASS  (this candidate is benchmark-valid
                                        on this attribute).
          d > {critical_pct}%  → REJECT — name this attribute as the
                                 disqualifier and do not pick this candidate.
        - String / category (e.g. drive_type, fuel_type, sub_type): if the
          values differ semantically (ignore casing / whitespace / format)
          → REJECT. These are family boundaries and the playbooks in
          Part B enumerate them.

      Worked example for a SOURCE injection moulding machine at 1400T,
      critical_pct=10:
        - candidate 1350T  →   3.6% delta → PASS  on tonnage.
        - candidate 1500T  →   7.1% delta → PASS  on tonnage.
        - candidate 1550T  →  10.7% delta → REJECT (just over the band).
        - candidate 1700T  →  21.4% delta → REJECT (well outside).
        - candidate  900T  →  35.7% delta → REJECT.
      The strict band is intentional — different tonnage = different product
      line = different OEM pricing tier. A 21% delta returned as a "weak
      match" would mislead the buyer.

   c) If the candidate has NO attribute semantically matching a source
      critical attribute → REJECT. Missing critical info is a mismatch, not
      a free pass.

   d) Extra critical attributes on the candidate that the source lacks are
      a PENALTY by default, NOT an automatic reject. The extraction LLM is
      noisy — the same physical machine extracted from two different PDFs
      will produce different attribute lists. Treat extra attrs as a soft
      penalty (record in `penalties`) unless they UNAMBIGUOUSLY name a
      different product variant. Examples of attrs that DO indicate a
      different variant (still reject):
        - injection_units_count ≥ 2  → multi-component moulding machine.
        - graphics_card_vram_gb on a laptop the source didn't spec → likely
          mobile workstation rather than business laptop.
        - magnetic_platen=true on a moulding machine where source did not
          spec one → quick-mould-change retrofit, not the base press.
      When in doubt, treat as penalty and let confidence + reasoning reflect it.

   e) Absence principle: only REJECT when the candidate has 3+ "important"
      attributes that the source completely lacks AND those attrs are
      clearly domain-specific (machine_model, tube_diameter, voltage_rating,
      injection_units, port_type, etc.) — the candidate is then almost
      certainly a different product family. One or two missing important
      attrs are penalties only (Rule 4), not disqualifiers. This threshold
      is intentionally higher than the original draft so noisy extraction
      does not shrink the comparable pool to zero on its own.

2. PRODUCT TYPE CONSISTENCY (item_level hierarchy).

   The item_level_1..8 chain identifies what the product IS, not just its
   specs. Reject candidates whose item_level chain diverges from the source
   even when numeric attributes match. Examples:
     - source level_3="Injection Moulding Machine",
       candidate level_3="Injection Moulding Machine Upgrade" → REJECT.
     - source level_2="Laptop", candidate level_2="Desktop" → REJECT.
     - source level_3="Diesel Generator", candidate level_3="Generator
       Maintenance Service" → REJECT.

   Also reject when item_name / item_description clearly indicates the
   candidate is an upgrade, modification, accessory, spare part, retrofit,
   installation service, AMC, or warranty — and the source is the parent
   product itself (or vice versa).

3. SPEC-TO-PRICE RATIO SANITY.

   When the source has at least one critical numeric attribute, compute
   ratio = candidate.unit_price_eur / candidate.<that_attribute>:
     - If ratio is more than 5× off the surviving-pool median → REJECT
       (category drift, e.g. a "1800T servo upgrade" priced at 13k EUR vs
       a "1800T machine" priced at 400k EUR).
     - If within 5× but outside ±{ratio_pct}% of the median → downrank.

4. IMPORTANT ATTRIBUTES ARE PENALTIES.
   For every "important" attribute mismatch of more than {important_pct}%
   (matched semantically per Rule 1a), lower the candidate's rank and
   record it under `penalties`. An "important" attribute present on the
   candidate but entirely absent from the source also counts as one such
   penalty. Three or more important mismatches together (was two — relaxed
   alongside rule 1e for the same noisy-extraction reason) are
   disqualifying.

5. BRAND SIMILARITY IS NOT RELEVANCE. Same brand alone is not a boost. Specs first.

6. RECENCY TIE-BREAK. Among candidates with identical verdict + same
   distance on the lock-on spec, prefer the more recently created PR
   (closer to source date = lower inflation correction).

7. DETERMINISM (so the same SOURCE + CANDIDATES set produces the same
   output across runs):
     a) Process candidates in ascending `purchase_dtl_id` order regardless
        of input order.
     b) Round every `delta_pct` to one decimal place using round-half-to-even.
     c) When two candidates remain truly tied after all rules above, break
        the tie by ascending `purchase_dtl_id`.
     d) Do not introduce randomness in confidence scores — confidence is a
        function of the rule evaluation, not a free parameter.

═════════════════════════════════════════════════════════════════════════
PART B — CATEGORY SPECIALIST PLAYBOOKS
═════════════════════════════════════════════════════════════════════════
Identify the source category from `purchase_category_llm` (preferred) or the
item_level_1..3 chain. Apply the matching playbook below. If none matches,
fall back to Part A only. "LOCK-ON" = treat as critical even if extraction
marked it merely "important"; mismatch beyond the stated band → REJECT.

▸ MOULDING MACHINE  (injection moulding presses)
   LOCK-ON: tonnage / clamping_force (T) — THE single biggest price driver.
            Hard band ±10%; outside that, reject regardless of other matches.
   FAMILY:  drive type — hydraulic / servo-hydraulic / hybrid / all-electric
            are DISTINCT families; do NOT cross even at same tonnage.
            Number of injection units (1 / 2 / 3 — multi-component machines
            are a separate family).
   SECONDARY: shot_weight (g/oz), screw_diameter (mm), tie-bar distance /
              platen size (mm), L/D ratio, control system.
   PRICE INTUITION: EUR per tonne is roughly stable within a family;
   anomalies usually mean accessory / upgrade / used machine — flag.

▸ MOULDING MACHINE AUXILIARIES
   FIRST: identify the sub-type and do NOT cross — chiller, MTC (mould
          temperature controller), dryer, hopper-loader, granulator, take-out
          robot, conveyor are distinct sub-categories.
   LOCK-ON per sub-type:
     - Chiller: cooling capacity (kW or TR), refrigerant, temp range.
     - MTC: heating capacity (kW), max temperature, water/oil medium.
     - Dryer: throughput (kg/hr), max temperature, dehumidifier vs hot-air.
     - Loader / hopper: capacity (kg or L), vacuum vs gravimetric.
     - Take-out robot: payload (kg), reach/strokes (mm), axes.
     - Granulator: throat size, throughput (kg/hr), motor power (kW).

▸ MOULD CHANGE EQUIPMENT
   SUB-TYPES (do not cross): quick-mould-change clamp, magnetic platen,
   mould cart, die-lifter / roller-bar, mould storage rack, mould pre-heater.
   LOCK-ON: max mould weight (T or kg), platen-size compatibility (mm), and
            host-press clamping-force range (T) it is rated for.

▸ TOOLING CAPEX (moulds, dies, fixtures, jigs)
   FIRST: identify tool type — injection mould, compression mould, stamping
          die, progressive die, fixture/jig, gauge — do NOT cross.
   LOCK-ON for MOULDS: cavity_count, part_weight (g), shot_weight (g), mould
            steel grade (P20 / H13 / 1.2738 / S136 / NAK80 — different
            price brackets), runner system (hot vs cold), mould-base size,
            cycle-life expectancy (number of shots).
   LOCK-ON for DIES: stages (single / progressive / transfer), die size (mm),
            press tonnage compatibility (T), strip width.
   PRICE INTUITION: cavity_count and steel grade dominate mould price; never
   compare a 1-cavity prototype mould with a 16-cavity production mould.

▸ ROBOT
   FAMILY (do not cross): 6-axis articulated, SCARA, delta/parallel,
   cartesian/gantry, cobot (collaborative). They occupy different price tiers.
   LOCK-ON: payload (kg), reach (mm), repeatability (±mm). Axes count.
   SECONDARY: control system / brand-family, IP rating, mounting (floor/
              ceiling/wall), end-of-arm tooling included or not.

▸ WIRE PROCESSING EQUIPMENT
   SUB-TYPES (do not cross): cut-&-strip, crimper / applicator, twister,
   harness-board, seal-insertion, splice.
   LOCK-ON: wire cross-section range (mm² or AWG min–max), throughput
            (cycles/min or m/min), station count, terminal-applicator type.

▸ ASSEMBLY / SECONDARY EQUIPMENT
   OPERATION TYPE (do not cross): ultrasonic welding, hot-plate welding,
   spin/vibration welding, screwdriving, heat-staking, riveting, press-fit,
   leak-test, marking/laser-engraving, glue-dispensing.
   LOCK-ON per type: process force (kN/N), travel (mm), cycle time,
            automation level (manual / semi-automatic / fully-automatic).

▸ MATERIAL HANDLING EQUIPMENT (MHE)
   FAMILY (do not cross): forklift, reach-truck, pallet-jack/PPT, stacker,
   AGV/AMR, conveyor, EOT/jib crane, hoist, scissor lift, elevator/lift.
   LOCK-ON per family:
     - Forklift / reach: load capacity (T or kg), lift height (mm),
       FUEL TYPE — diesel / LPG / electric / lithium-ion are different
       price bands; do NOT cross fuel types.
     - Crane / hoist: capacity (T), span (m), lifting height (m), duty class.
     - Conveyor: type (belt / roller / chain / overhead), length (m),
       width (mm), throughput.
     - AGV / AMR: payload (kg), navigation (laser/magnetic/SLAM), battery.

▸ PAINT SHOP
   SUB-TYPES (do not cross): spray booth, curing/baking oven, pre-treatment
   line, paint-mixing room, paint robot, conveyor for paint line.
   LOCK-ON:
     - Booth: working envelope (L×W×H mm), filtration (downdraft / crossdraft
       / sidedraft), exhaust airflow.
     - Oven: max temperature, throughput, energy source (gas / electric / IR).
     - Pre-treatment line: number of stages, throughput (m²/hr).

▸ PLANT EQUIPMENT  (generic catch-all for utility machinery)
   Identify sub-family from item_level_3/4 and match STRICTLY within that
   sub-family. LOCK-ON: rated power (kW) or capacity (the dominant numeric
   spec for that sub-family). NEVER match across sub-families even if other
   numbers happen to align.

▸ IT AND COMMUNICATIONS
   DEVICE CLASS (do not cross): laptop, desktop, workstation, thin-client,
   tablet, server, network switch, router, firewall, wireless AP, UPS,
   storage array, printer, MFP, monitor, KVM.
   LOCK-ON per class:
     - Laptop / desktop / workstation: CPU family AND generation (Intel
       i5-12xxx ≠ i5-9xxx), RAM (GB), storage (GB/TB, SSD vs HDD), display
       size for laptops, dGPU presence (workstation vs business laptop).
     - Server: CPU sockets + cores + generation, RAM, raw storage capacity,
       form factor (1U/2U/4U/tower/blade), redundant PSU.
     - Network switch: port count + speed (1G / 10G / 25G / 40G / 100G),
       managed/unmanaged, PoE budget (W).
     - UPS: VA / kVA, runtime at load, topology (offline / line-interactive
       / online double-conversion), single- vs three-phase.
   PRICE INTUITION: never match a 13-inch business laptop with a 17-inch
   mobile workstation just because both say "laptop".

▸ TELE COMMUNICATION
   DEVICE TYPE (do not cross): IP-PBX, IP phone / softphone licence,
   video-conferencing endpoint, video-conferencing room kit, GSM/SIM
   connectivity, leased-line bandwidth, mobile handset.
   LOCK-ON: user count / port count / call capacity / bandwidth (Mbps),
            contract duration (months) for services.

▸ LABORATORY AND TEST EQUIPMENT
   INSTRUMENT CLASS (do not cross): CMM, tensile/UTM, hardness tester,
   profilometer, surface roughness, optical comparator, microscope,
   spectrometer, chromatograph, oscilloscope, multimeter, environmental
   chamber, salt-spray, vibration table.
   LOCK-ON: measurement range, accuracy / resolution, working volume,
            actuator type (servo / hydraulic / pneumatic) where relevant.

▸ FIRE PROTECTION
   SYSTEM TYPE (do not cross): portable extinguisher, sprinkler system,
   hydrant, FM-200 / Novec / CO2 suppression, smoke / heat detector,
   hose reel, fire-alarm panel.
   LOCK-ON: capacity (kg / L for extinguishers; coverage area m² for
            sprinkler/suppression; protected volume m³ for gas systems),
            agent type, standard (BS / NFPA / IS).

▸ OFFICE SUPPLIES EQUIPMENT
   EQUIPMENT CLASS (do not cross): printer, copier, MFP, scanner, shredder,
   projector / display board, office chair, desk, lockers, partitions.
   LOCK-ON: function set (single / MFP), PPM, colour vs mono, duplex,
            paper size (A4/A3) for printers; ergonomic class for chairs.

▸ PACKAGING
   DELIVERABLE TYPE (do not cross): corrugated box, wooden pallet/crate,
   plastic returnable crate / FLC / KLT, stretch film, strapping, printed
   label, VCI/anti-corrosion bag, foam dunnage, packaging MACHINE.
   LOCK-ON consumables: material + dimensions (mm) + GSM or thickness
                         (microns/mm) + load rating where applicable.
   LOCK-ON machines: machine type (filler / sealer / labeler / wrapper /
                     case-erector), throughput (units/min), pack format.

▸ FACILITIES MANAGEMENT  (predominantly SERVICES — not equipment)
   SERVICE SCOPE (do not cross): housekeeping, security/guarding, pest
   control, landscaping/horticulture, HVAC AMC, electrical AMC, plumbing,
   waste management, canteen.
   LOCK-ON: area serviced (sqft/sqm), manpower headcount, shifts/day,
            contract duration (months), inclusions (chemicals/consumables).
   PRICE INTUITION: always normalise to per-sqft-per-month or per-headcount
   before comparing absolutes; raw totals are misleading.

▸ FLEET MANAGEMENT
   FIRST: distinguish PURCHASE vs LEASE vs MANAGED-FLEET SERVICE — these are
          three different deliverables, never cross.
   VEHICLE CLASS (do not cross): car (hatch/sedan/SUV), van/MPV, light truck,
   heavy truck, bus, two-wheeler.
   LOCK-ON: vehicle class + variant, fuel (petrol/diesel/CNG/EV), seating
            or payload capacity, transmission, contract duration if lease.

▸ MACHINE LEASE
   COMPOSITE: this is "machine + financing". Match BOTH legs.
     1) The underlying machine — apply that category's playbook (often
        Moulding Machine, MHE, Robot, etc.).
     2) Lease attributes — lease type (operating vs finance), duration
        (months), monthly vs total quoted basis, residual value, included
        services (AMC bundled or not). Compare on monthly per-unit basis.

▸ PROPERTY / LAND LEASE
   DELIVERABLE TYPE (do not cross): warehouse, factory shed, office,
   industrial land (raw), commercial, residential.
   LOCK-ON: area (sqft or sqm), location/city (and tier — comparing tier-1
            with tier-3 is invalid), lease duration (months/years), fitted /
            bare-shell, CAM charges in or out.
   PRICE INTUITION: always normalise to per-sqft-per-month rent; comparing
   absolute monthly rents across different areas is meaningless.

▸ MAINTENANCE REPAIR AND OPERATION (MRO)
   FIRST: distinguish SPARE PART, CONSUMABLE, and SERVICE — three different
          deliverables, never cross.
   For SPARES / CONSUMABLES: parent machine model AND OEM part number (or
            functional equivalent) compatibility is paramount, then material
            and grade.
   For SERVICES: scope of work (preventive vs breakdown vs overhaul),
            duration / hours, response SLA, AMC vs CMC, parts inclusive.

═════════════════════════════════════════════════════════════════════════
PART C — OUTPUT
═════════════════════════════════════════════════════════════════════════

The output uses a TWO-ARRAY schema designed to make completeness mandatory
and to avoid the lazy-skip failure mode that occurs when the model is asked
to write a long paragraph for every candidate. Read the rules below before
the schema.

Hard rules for the output:

  1. SELECT THE TOP 20–25 MOST SIMILAR CANDIDATES — FOLLOW THIS EXACT
     PROCEDURE. Do not let intuition override the steps. The goal is to
     NEVER miss a numerically-closest candidate, even if its variant or
     family differs from the source. A missed close-match is the single
     worst failure mode of this step.

     ─── STEP 1A: Identify the source's PRIMARY LOCK-ON attribute ───

     Use this category → lock-on map. These are the price-dominant
     attributes that procurement specialists prioritize. Find the entry
     whose key matches source.purchase_category_llm and use its
     attribute list as the lock-on candidates:

       Moulding Machine                → tonnage / clamping_force_t /
                                          clamping_force_tonnage /
                                          clamping_force_kn
                                          (units T or kN; primary spec)
       Moulding Machine Auxiliaries    → capacity_kw or throughput_kg_hr
                                          (depending on sub-type — see
                                          Part B playbook)
       Mould Change Equipment          → max_mould_weight (T or kg)
       Tooling CapEx                   → cavity_count (PRIMARY) +
                                          mould_steel_grade (FAMILY)
       Robot                           → payload (kg) AND reach (mm)
                                          BOTH equal weight — combine
                                          both into similarity score
       Wire Processing Equipment       → wire_cross_section_max
                                          (mm² or AWG)
       Assembly / Secondary Equipment  → process_force_kn or
                                          cycle_time_s
       Material Handling Equipment     → load_capacity_t AND
                                          lift_height_mm BOTH equal weight
       Paint Shop                      → working_envelope_mm OR
                                          throughput_units_per_hr
       Plant Equipment                 → rated_power_kw OR capacity
                                          (look at the dominant numeric
                                          attr in source.critical_attributes)
       IT and Communications           → device-class-dependent:
                                          laptop/desktop → ram_gb +
                                            storage_gb + cpu_generation
                                          server → cpu_cores + ram_gb
                                          network switch → port_count
                                          UPS → kva_rating
       Tele Communication              → user_count OR port_count
       Laboratory and Test Equipment   → measurement_range (instrument-
                                          class dependent)
       Fire Protection                 → capacity_kg OR coverage_m2
                                          (system-type dependent)
       Office Supplies Equipment       → equipment-class-dependent:
                                          printer → ppm_pages_per_min
                                          chair → ergonomic_class
       Packaging                       → for consumables: material +
                                                          dimensions_mm
                                          for machines: throughput_units_per_min
       Facilities Management           → area_sqft (PRIMARY) +
                                          headcount + duration_months
       Fleet Management                → vehicle_class + duration_months
       Machine Lease                   → underlying machine's lock-on
                                          (apply that machine's playbook) +
                                          duration_months
       Property / Land Lease           → area_sqft (PRIMARY) +
                                          location_tier + duration_months
       Maintenance Repair and Operation→ scope (services) OR parent_machine
                                          (spares); part_number where present

     Look in source.critical_attributes for an attribute whose name
     semantically matches one of the lock-on names above. That attribute's
     value is your PRIMARY_LOCK_ON_VALUE. If multiple lock-on attributes
     are listed (e.g. Robot uses payload AND reach), combine them with
     equal weight into the similarity score (see STEP 1B).

     ─── STEP 1B: Score EVERY candidate (do NOT skip any yet) ───

     For each of the N candidates provided, compute a similarity rank
     using this priority order. Higher priority always wins over lower
     priority — do NOT collapse them into one weighted sum.

       PRIORITY 1 — item_level_3 match (HARD filter for ranking).
         Candidates whose item_level_3 equals source.item_level_3 rank
         ABOVE candidates whose item_level_3 differs. A 1500T Injection
         Moulding Machine ranks above a 1400T Rubber Injection Moulding
         Machine for an Injection Moulding source, even though the
         rubber one has zero tonnage delta.

       PRIORITY 2 — Numeric proximity on PRIMARY_LOCK_ON_VALUE.
         Compute delta_pct = |src − cand| / src × 100. Smaller delta_pct
         ranks higher. This is the DOMINANT signal within the same
         item_level_3 bucket.
         For multi-attr lock-ons (Robot, MHE, FM, etc.), use the AVERAGE
         of the delta_pct values across the lock-on attrs.

       PRIORITY 3 — Same Part B playbook variant.
         Within identical numeric proximity, prefer candidates that share
         the source's variant attrs (drive_type, injection_units,
         fuel_type, sub_type, device_class).
         CRITICAL — a variant difference does NOT exclude a candidate
         from top 25. A 1400T 1K machine MUST appear in the top 25 of
         a 1400T 2K source's pool because tonnage is identical. The
         variant difference is handled at the VERDICT step (later, where
         it becomes a FAIL with rule=VARIANT_MISMATCH).

       PRIORITY 4 — Recency (more recent first).
         Tie-breaker only when priorities 1-3 are identical.

     ─── STEP 1C: Take the top 25 by the priority-ordered rank ───

     The output of STEP 1B is an ordered list of all N candidates. Take
     the first 25 (or all N if N < 20).

     MANDATORY GUARDS — verify before proceeding:

       (G1) The candidate with the SMALLEST delta_pct on the primary
            lock-on (the absolute-closest candidate) MUST be in your
            top 25, no matter what its variant, family, or product
            type. Missing the absolute closest is the worst failure.
       (G2) The 2nd, 3rd, 4th, 5th smallest delta_pct candidates MUST
            also be in your top 25 (5 closest by tonnage = mandatory).
       (G3) ALL candidates with delta_pct ≤ {critical_pct}% on the
            primary lock-on AND same item_level_3 MUST be in your top
            25 — they are PASS-eligible and a missed PASS is a critical
            error.
       (G4) For ties on primary lock-on (e.g. two candidates at 1300T
            for a 1400T source), include BOTH — never arbitrarily
            drop one of a tied pair.

     If any of G1–G4 is violated, REDO the selection.

     ─── STEP 1D: Worked example (Moulding Machine, source 1400T 2K) ───

     Candidate pool snippet:
       A: 1350T 2K hydraulic, level_3="Injection Moulding Machine"
            → delta 3.6%, same variant, same level_3 → rank ≈ 1
       B: 1500T 2K servo,     level_3="Injection Moulding Machine"
            → delta 7.1%, same variant, same level_3 → rank ≈ 2
       C: 1300T 1K hydraulic, level_3="Injection Moulding Machine"
            → delta 7.1%, diff variant, same level_3 → rank ≈ 3
            (B and C tie on delta; B's same-variant wins Priority 3,
             but BOTH belong in top 25)
       D: 1700T 2K servo,     level_3="Injection Moulding Machine"
            → delta 21.4%, same variant, same level_3 → rank ≈ 6-10
       E: 100T all-electric,  level_3="Injection Moulding Machine"
            → delta 92.9%, far family → near bottom, NOT in top 25
       F: 1400T Rubber Inj.,  level_3="Rubber Injection Moulding Machine"
            → delta 0% but Priority 1 demotes it → MAY appear in top 25
            if no closer same-level_3 candidate fills the slot, else
            excluded

     Candidate C is the critical case: a 1K machine with same tonnage as
     a 2K source MUST appear in top 25. The LLM later issues a FAIL
     verdict with rule=VARIANT_MISMATCH and a note explaining the 1K vs
     2K price difference. EXCLUDING C from the top 25 would hide a
     candidate that is only off by ONE variant attribute — that is a
     serious failure of this step.

     ─── STEP 1E: Sanity check (run before emitting verdicts) ───

       - Did the absolute-closest-tonnage candidate make my top 25? Yes / No
       - Did all candidates with delta_pct ≤ {critical_pct}% AND same
         level_3 make my top 25? Yes / No
       - Are there exactly 20–25 candidates in my verdicts list? Yes / No

     If any answer is "No", REDO. The top 25 selection is the foundation
     of the rest of the output; do not move on with a flawed selection.

     Hard cap: `len(verdicts) MUST be between 20 and 25` (or exactly
     N if N < 20, where N is the total candidates received).

  2. WITHIN the top 20–25, decide each candidate's verdict:

       verdict = "PASS"  → the candidate is genuinely benchmarkable (same
                           family, critical numeric within ±{critical_pct}%,
                           no killer family/variant disqualifier).
       verdict = "FAIL"  → close enough to consider, BUT disqualified by
                           tonnage band, product type, variant family,
                           ratio, or absence rule.

     Both PASS and FAIL share the compact-row shape (one row in `verdicts`
     each); detailed reasoning differs (see rules 3 and 4).

  3. Each `verdicts` row is COMPACT, with fields:
       purchase_dtl_id, verdict (PASS|FAIL), rule, attr, src, cand, unit,
       delta_pct, note.

     For verdict=FAIL rows, the `note` field is MANDATORY and must be
     15–35 words explaining the BUSINESS CONSEQUENCE of including this
     candidate as a benchmark. Mechanical reasoning like "tonnage X vs Y,
     delta exceeds band" is INSUFFICIENT — that just restates the
     structured fields. The note must explain WHY a buyer should not use
     this row: the price-tier difference, the application difference, or
     the OEM positioning difference. Examples of acceptable FAIL notes:
       - "100T precision machines serve small parts (~30–50k EUR class);
          using as anchor for a 1400T quote (~350–450k EUR class) would
          understate benchmark by an order of magnitude."
       - "4K combi presses carry ~50–70% price premium over single-shot
          machines at same tonnage due to four injection aggregates plus
          rotary platen — would inflate benchmark substantially."
       - "Rubber injection presses use vertical clamping and elastomer-
          specific systems; their EUR/tonne pricing is 30–40% lower than
          thermoplastic horizontals, so cross-class anchoring distorts
          benchmark."
       - "Retrofit/upgrade kits are services priced at a fraction of base
          machine cost; anchoring a new-press quote on a retrofit would
          understate benchmark by 10–20×."

     For verdict=PASS rows, `note` may be empty (the rich reasoning lives
     in `selected_detail`).

  4. The `selected_detail` array contains the FULL reasoning ONLY for
     candidates whose `verdict == "PASS"` in `verdicts`. One detail row
     per PASS verdict. If no candidate passes, `selected_detail` is `[]`.

  5. Rejected items live ONLY in `verdicts` with verdict=FAIL. Do NOT
     create a separate "rejected_detail" array — the FAIL note (rule 3
     above) is sufficient prose for the buyer.

  6. Order: sort `verdicts` by ascending purchase_dtl_id. Sort
     `selected_detail` by rank (1 = best).

LABELING RULE — CRITICAL — read carefully:

  The `rule` and `attr` fields MUST name the ACTUAL disqualifier, not just
  any attribute you happened to look at. Do NOT name `tonnage` (or any other
  numeric lock-on) as the disqualifier_attr when that numeric value is
  INSIDE the strict band — that is mathematically inconsistent and would
  read as a contradiction to a buyer.

  Decision tree for `rule` selection:

    1. Is the numeric lock-on spec OUTSIDE the ±{critical_pct}% band?
         YES → rule = "CRITICAL_BARRIER", attr = that numeric spec
                       (e.g. "tonnage"), src/cand/unit/delta_pct populated.
                       Note may be empty.
         NO  → continue to step 2.

    2. Is there a Part B playbook family boundary breached (different
       drive_type / injection_units / fuel_type / sub_type / device_class /
       etc.) even though the numeric lock-on is within band?
         YES → rule = "VARIANT_MISMATCH", attr = the family attribute
                       (e.g. "injection_units", "drive_type",
                       "fuel_type"), src/cand = the family values
                       (e.g. src=2, cand=1 for injection_units; or
                       src="hydraulic", cand="all-electric"). unit may be
                       "count" or null; delta_pct may be a number for
                       countable family attrs (1 → 3 is +200%) or null for
                       categorical (hydraulic vs all-electric).
                       note may amplify but is not required.
         NO  → continue to step 3.

    3. Is item_level_3/4 chain divergent (e.g. "Injection Moulding
       Machine" vs "Rubber Injection Moulding Machine" or vs
       "Retrofit/Service/Options")?
         YES → rule = "PRODUCT_TYPE", attr = "item_level_3" or the
                       diverging level, src/cand = the level values.

    4. Is the spec-to-price ratio more than 5× off the surviving-pool median?
         YES → rule = "RATIO", attr = the numeric spec used for the ratio.

    5. Does the candidate carry 3+ domain-specific important attrs the
       source completely lacks?
         YES → rule = "ABSENCE", attr = a representative missing attr name.

    6. None of the above but you still reject for some other reason?
         → rule = "OTHER", attr = whatever is most informative, note
                   REQUIRED.

  Concrete worked examples (apply this exact shape — don't deviate):

    • Candidate 1300T vs source 1400T (7.1% delta — INSIDE ±10% band) BUT
      candidate has 1 injection unit vs source's 2 → rule=VARIANT_MISMATCH,
      attr="injection_units", src=2, cand=1, unit="count", delta_pct=null,
      note="single-injection-unit vs 2K source — different machine variant
      per Moulding Machine playbook".

    • Candidate 100T vs source 1400T (92.9% delta — OUTSIDE band) →
      rule=CRITICAL_BARRIER, attr="tonnage", src=1400, cand=100, unit="T",
      delta_pct=92.9, note="".

    • Candidate 1700T servo-hydraulic vs source 1700T all-electric →
      rule=VARIANT_MISMATCH, attr="drive_type", src="all-electric",
      cand="servo-hydraulic", unit=null, delta_pct=null, note="".

    • Candidate is a "Rubber Injection Moulding Machine" line item →
      rule=PRODUCT_TYPE, attr="item_level_3", src="Injection Moulding
      Machine", cand="Rubber Injection Moulding Machine", note="".

For each `selected_detail` row, produce a rich `agent_reasoning` payload
so a buyer can audit the match. Because every `verdict=PASS` row has
ALREADY passed the strict ±critical_pct band, the STRONG / MODERATE / WEAK
verdict reflects the WHOLE picture — important attr penalties, spec-to-
price-ratio band, candidate-has-extra-attrs penalties, age, brand:

  STRONG_MATCH    — All source critical attrs within band, zero important
                    attribute penalties, spec-to-price ratio inside
                    ±{ratio_pct}% of pool median. confidence ≥ 0.80.
  MODERATE_MATCH  — All critical attrs within band, but 1–2 soft penalties.
                    confidence 0.55–0.79.
  WEAK_MATCH      — All critical attrs within band, but multiple soft
                    penalties AND/OR spec-to-price in DOWNRANK zone AND/OR
                    age >24 months. confidence 0.30–0.54.

A candidate OUTSIDE the ±critical_pct band is NEVER a WEAK_MATCH. It is a
FAIL in `verdicts`. WEAK is reserved for candidates that pass the band but
carry multiple soft penalties.

Detail row fields (only for verdict=PASS rows):
  - matched_attrs: every critical AND important attribute you compared,
    with src/cand/unit/delta_pct (1 decimal, round-half-to-even),
    importance, and per-attribute PASS / PENALTY verdict.
  - spec_to_price_ratio: when meaningful (critical numeric attr +
    unit_price_eur both exist). Include ratio, pool_median, and
    PASS/DOWNRANK/REJECT verdict. Else verdict="N/A", value/median=null.
  - penalties: short bullet strings with numbers + units. Be specific.
  - confidence: 0.0–1.0. Map from verdict per the bands above.
  - verdict: STRONG_MATCH | MODERATE_MATCH | WEAK_MATCH.
  - reasoning: 5–8 sentence buyer-facing justification structured as four
    named beats with the labels visible in the prose:
      PRODUCT FIT: same product family (cite drive_type / sub_type /
        fuel_type / device_class) AND same item_level chain. Quote source
        and candidate level_3.
      SPEC ALIGNMENT: every critical and important attribute by name with
        src, cand, unit, delta %. Numbers, not adjectives.
      PRICE COMPARABILITY: EUR-per-spec ratio vs pool median; inside
        ±{ratio_pct}% sanity band? If not applicable, state why.
      RECOMMENDED USE: how the buyer should use this row (primary anchor,
        secondary corroboration, last-resort sanity check) and the
        quantified confidence. End with a declarative sentence.

Output strict JSON, no fences, no prose, no trailing commas:

{{
  "category_playbook": "<the Part B playbook you applied, or 'generic'>",
  "thresholds_applied": {{
    "critical_pct":  {critical_pct},
    "important_pct": {important_pct},
    "ratio_pct":     {ratio_pct}
  }},

  "verdicts": [
    {{
      "purchase_dtl_id": <int>,
      "verdict":         "PASS|FAIL",
      "rule":            "CRITICAL_BARRIER|PRODUCT_TYPE|RATIO|ABSENCE|VARIANT_MISMATCH|OTHER|null",
      "attr":            "<attr / level / ratio name or null>",
      "src":             <value or string or null>,
      "cand":            <value or string or null>,
      "unit":            "<unit or null>",
      "delta_pct":       <number to 1 decimal or null>,
      "note":            "<10-15 word note when judgement applies; otherwise empty string>"
    }}
    // ... ONE row per candidate, sorted by ascending purchase_dtl_id ...
  ],

  "selected_detail": [
    {{
      "purchase_dtl_id": <PASS id from verdicts>,
      "rank":            1,
      "category_playbook": "<playbook applied to this candidate>",
      "matched_attrs": [
        {{
          "name":       "<canonical attr name>",
          "src":        <value or string>,
          "cand":       <value or string>,
          "unit":       "<unit or null>",
          "delta_pct":  <number to 1 decimal or null>,
          "importance": "critical|important|nice_to_have",
          "verdict":    "PASS|PENALTY"
        }}
      ],
      "spec_to_price_ratio": {{
        "attribute":   "<attr name or null>",
        "value":       <number or null>,
        "pool_median": <number or null>,
        "verdict":     "PASS|DOWNRANK|REJECT|N/A"
      }},
      "penalties":  ["<short bullet with numbers + units>", "..."],
      "confidence": <0.0-1.0>,
      "verdict":    "STRONG_MATCH|MODERATE_MATCH|WEAK_MATCH",
      "reasoning":  "PRODUCT FIT: ... SPEC ALIGNMENT: ... PRICE COMPARABILITY: ... RECOMMENDED USE: ..."
    }}
    // ... ONE row per PASS verdict, sorted by rank ascending ...
  ]
}}
"""


RANK_PROMPT_USER_V2 = """Pick the 20 to 25 candidates MOST similar to the SOURCE item below. Emit EXACTLY 20–25 verdict rows. If the input pool has fewer than 20 candidates, emit a verdict for every candidate. Do NOT emit fewer than 20 when N >= 20.

## SOURCE ITEM

```json
{source_json}
```

## CANDIDATES ({n_candidates} retrieved by SQL+Pinecone)

```json
{candidates_json}
```

Thresholds: critical_pct={critical_pct}, important_pct={important_pct}, ratio_pct={ratio_pct}.

Process — follow EVERY step:

  STEP 0 — COMPLETENESS SCAN (run this FIRST, before any rejection):
    Mentally enumerate EVERY candidate. For each one, identify its value
    for the source's primary lock-on spec (e.g. tonnage for a moulding
    machine, payload+reach for a robot, capacity+lift_height for a forklift,
    CPU+RAM for a laptop). DO NOT proceed to rejection until you have
    located the lock-on spec on every candidate. If a candidate's lock-on
    spec is genuinely absent, that is itself a finding (Rule 1c reject),
    but you must NOTICE it rather than skip the candidate.

    Concretely for a 1800T source: of the N candidates given, identify the
    1750T / 1800T / 1850T / 1700T / similar-tonnage entries explicitly
    BEFORE deciding any rejections. Missing a within-band candidate is
    the single worst failure mode of this prompt — guard against it here.

  STEP 1 — PLAYBOOK SELECT.
    Read SOURCE.purchase_category_llm (and item_level chain) and pick the
    matching Part B playbook. State it in the top-level `category_playbook`.

  STEP 2 — RULE WALK.
    Process candidates in ASCENDING purchase_dtl_id order. For each, walk
    Part A rules 1 → 7, then the playbook lock-ons. Apply rule 1b STRICTLY
    — beyond ±critical_pct on any source critical attribute is a REJECT
    with no exceptions.

  STEP 3 — PICK TOP 20–25 MOST SIMILAR CANDIDATES.
    Out of the {n_candidates} candidates provided, select the 20 to 25
    that are MOST SIMILAR to the source by:
      (a) closest on the source's primary lock-on numeric spec (the spec
          the Part B playbook calls out for this category, e.g. tonnage
          for moulding machines),
      (b) same product family / matching item_level_3 chain,
      (c) same Part B playbook variant when one applies (drive type,
          injection-units count, fuel type, device class, etc.).

    Emit `verdicts[]` rows ONLY for these 20–25 candidates. Candidates
    outside the top 20–25 are too distant to be benchmark-relevant — they
    do NOT get a verdict, and they do NOT appear anywhere in the output.

    If {n_candidates} < 20, emit a verdict for every candidate
    (`len(verdicts) == {n_candidates}`).

    Hard cap: `len(verdicts) MUST be between 20 and 25` (or exactly
    {n_candidates} if smaller than 20). Going over 25 risks output
    truncation; under 20 means you skipped close matches.

  STEP 4 — VERDICTS (compact rows for each of the top 20–25).
    For each selected candidate in ascending purchase_dtl_id order, emit
    ONE compact `verdicts` row with the fields:
      purchase_dtl_id, verdict (PASS or FAIL), rule, attr, src, cand,
      unit, delta_pct (1 decimal, round-half-to-even), note.

    For verdict=FAIL rows, `note` is MANDATORY (15–35 words) and must
    explain the BUSINESS CONSEQUENCE of using this candidate as a
    benchmark — price-tier difference, application difference, or OEM
    positioning difference. Do NOT just restate the structured fields.
    See the FAIL note examples in the system prompt.

    For verdict=PASS rows, `note` may be empty (rich reasoning goes in
    `selected_detail`).

    DO NOT write paragraphs in this step. The whole row should fit in
    ~30–60 tokens.

  STEP 5 — DETAIL (only for verdict=PASS rows).
    For each row in `verdicts` whose `verdict == "PASS"`, emit ONE row in
    `selected_detail` with the full 5–8 sentence `reasoning` using the
    named PRODUCT FIT / SPEC ALIGNMENT / PRICE COMPARABILITY / RECOMMENDED
    USE structure, plus matched_attrs, spec_to_price_ratio, penalties,
    confidence, and verdict (STRONG/MODERATE/WEAK_MATCH).

    If NO candidate passes (none meet the strict ±{critical_pct}% band on
    the lock-on spec), `selected_detail` is `[]`. Do NOT invent a
    comparable just to fill the array. Returning an empty `selected_detail`
    is the correct answer when no candidate in the top 20–25 truly passes.

  STEP 6 — ORDERING.
    Sort `verdicts` by ascending purchase_dtl_id. Sort `selected_detail`
    by rank ascending (1 = best); among ties, by ascending purchase_dtl_id.
    This guarantees deterministic output across runs.

Return strict JSON exactly as specified in the system prompt — no fences,
no prose, no trailing commas.
"""


# ── Dataclasses ────────────────────────────────────────────────────────────────

@dataclass
class LineItemContext:
    purchase_dtl_id:  int
    purchase_req_id:  int
    item_no:          int
    quantity:         Optional[Decimal]
    item_type:        Optional[str]
    item_description: Optional[str]
    unit_price:       Optional[Decimal] = None
    uom:              Optional[str]     = None
    supplier_name:    Optional[str]     = None
    discount:         Optional[Decimal] = None
    req_value:        Optional[Decimal] = None
    currency:         Optional[str]     = None
    delivery_date:    Optional[str]     = None
    payment_details:  Optional[str]     = None
    original_value:   Optional[Decimal] = None
    initial_offer:    Optional[Decimal] = None
    negotiation:      Optional[Decimal] = None
    comments:         Optional[str]     = None
    prepayment:       Optional[str]     = None
    item_code:        Optional[str]     = None


@dataclass
class RASContext:
    purchase_req_no:   str
    purchase_req_id:   int
    supplier_name:     Optional[str]
    justification:     Optional[str]
    currency:          Optional[str]
    enquiry_no:        Optional[str]     = None
    classification:    Optional[str]     = None
    department:        Optional[str]     = None
    negotiated_by:     Optional[str]     = None
    address:           Optional[str]     = None
    contract_no:       Optional[str]     = None
    order_no:          Optional[str]     = None
    purchase_value:    Optional[Decimal] = None
    category:          Optional[str]     = None
    sub_category:      Optional[str]     = None
    site_country:      Optional[str]     = None
    site_region:       Optional[str]     = None
    site:              Optional[str]     = None
    division:          Optional[str]     = None
    requisition_type:  Optional[str]     = None
    parent_supplier:   Optional[str]     = None
    supplier_type:     Optional[str]     = None
    supplier_country:  Optional[str]     = None
    payment_days:      Optional[str]     = None
    po_date:           Optional[str]     = None
    category_buyer:    Optional[str]     = None
    l3:                Optional[str]     = None
    l4:                Optional[str]     = None
    purchase_category: Optional[str]     = None
    ras_title:         Optional[str]     = None
    line_items:        list              = field(default_factory=list)
    req_start_date:    Optional[object]  = None
    c_datetime:        Optional[object]  = None
    raw_mst:           dict             = field(default_factory=dict)
    raw_dtl_rows:      list             = field(default_factory=list)
    raw_vw_rows:       list             = field(default_factory=list)


@dataclass
class DocumentContent:
    text:        Optional[str]      = None
    images:      Optional[list]     = None
    source_path: str                = ""
    page_count:  int                = 0
    ocr_source:  bool               = False

    @property
    def is_image_based(self) -> bool:
        return bool(self.images)


# ── Connection helpers ─────────────────────────────────────────────────────────

def _conn_str(conn_data: Data) -> str:
    d      = conn_data.data or {}
    driver = d.get("driver", "ODBC Driver 18 for SQL Server")
    server = d.get("host", d.get("server", ""))
    port   = d.get("port", 1433)
    db     = d.get("database_name", d.get("database", ""))
    user   = d.get("username", "")
    pwd    = d.get("password", "")
    return (
        f"DRIVER={{{driver}}};SERVER={server},{port};"
        f"DATABASE={db};UID={user};PWD={pwd};TrustServerCertificate=yes;"
    )


def _is_transient(exc: Exception) -> bool:
    kw = ["connection reset", "timeout", "throttl", "resource limit",
          "broken pipe", "transport-level", "login failed"]
    return any(k in str(exc).lower() for k in kw)


def _is_deadlock_error(exc: Exception) -> bool:
    """True for SQL Server deadlock-victim errors.

    Triggered when concurrent workers run _cleanup_for_pr against
    different PRs but the DELETEs join through shared rows in
    ras_tracker / attachment_classification and SQL Server picks the
    transaction as a deadlock victim.

    Error 1205 / SQLSTATE 40001 — the SQL Server message asks us
    to rerun the transaction, which is exactly what the retry wrapper
    in _process_pr does on the cleanup path.
    """
    msg = str(exc)
    if "40001" in msg or "1205" in msg:
        return True
    return "deadlock" in msg.lower()


def _connect(cs: str):
    """Get a pooled DB connection. The returned object behaves like a normal
    pyodbc.Connection — .cursor() / .commit() / .rollback() / .close() all
    work — but .close() returns it to the shared _CONN_POOL instead of
    tearing down the socket. Retries on transient errors before giving up."""
    for attempt in range(_MAX_RETRIES + 1):
        try:
            return _CONN_POOL.acquire(cs, login_timeout=30)
        except Exception as exc:
            if not _is_transient(exc) or attempt == _MAX_RETRIES:
                raise
            time.sleep(_BASE_DELAY * (2 ** attempt) * (1 + random.random() * 0.2))
    raise RuntimeError("unreachable")


# ── Blob helpers ───────────────────────────────────────────────────────────────

def _get_blob_config_by_name(connector_name: str) -> dict:
    name = (connector_name or "").strip()
    if not name:
        raise ValueError("blob_connector_name is empty.")
    try:
        import asyncio
        import concurrent.futures as _cf
        from sqlalchemy import select

        async def _fetch():
            from agentcore.services.deps import get_db_service
            from agentcore.services.database.models.connector_catalogue.model import ConnectorCatalogue
            db_service = get_db_service()
            async with db_service.with_session() as session:
                stmt = (
                    select(ConnectorCatalogue)
                    .where(ConnectorCatalogue.name == name)
                    .where(ConnectorCatalogue.provider == "azure_blob")
                )
                result = await session.execute(stmt)
                row = result.scalars().first()
                if row is None:
                    raise ValueError(f"No azure_blob connector named {name!r}.")
                cfg         = row.provider_config or {}
                account_url = (cfg.get("account_url") or row.host or "").strip()
                container   = (cfg.get("container_name") or row.database_name or "").strip()
                from urllib.parse import urlparse
                if not urlparse(account_url).netloc:
                    raise ValueError(f"Invalid Storage Account URL: {account_url!r}")
                if not container:
                    raise ValueError(f"No container_name in connector {name!r}.")
                return {"account_url": account_url, "container_name": container}

        try:
            asyncio.get_running_loop()
            with _cf.ThreadPoolExecutor() as pool:
                return pool.submit(asyncio.run, _fetch()).result(timeout=10)
        except RuntimeError:
            return asyncio.run(_fetch())
    except Exception as exc:
        raise RuntimeError(f"Blob connector lookup failed: {exc}") from exc


def _download_blob(blob_path: str, blob_cfg: dict) -> bytes:
    if "blob_client" in blob_cfg:
        # Direct client passed from re_commercials.py (pre-created with Azure CLI auth)
        client = blob_cfg["blob_client"]
        blob = client.get_blob_client(container=blob_cfg["container_name"], blob=blob_path)
        return blob.download_blob().readall()
    else:
        # Connector catalogue config (from AgentCore): use the cached
        # DefaultAzureCredential so dev (az login) and prod (Managed
        # Identity) both work without rebuilding the credential chain
        # on every blob op. Wrapped in _with_az_retry so a transient
        # token-refresh failure rebuilds the credential and retries once.
        from azure.storage.blob import BlobServiceClient

        def _do() -> bytes:
            credential = _get_az_credential()
            client = BlobServiceClient(
                account_url=blob_cfg["account_url"], credential=credential
            )
            blob = client.get_blob_client(
                container=blob_cfg["container_name"], blob=blob_path
            )
            return blob.download_blob().readall()

        return _with_az_retry(_do)


# ── File type detection ────────────────────────────────────────────────────────

def _detect_file_type(filename: str) -> str:
    import os
    ext = os.path.splitext(filename.lower())[1]
    return _EXT_TO_TYPE.get(ext, "unknown")


# ── Content extraction for classification ─────────────────────────────────────

def _extract_for_classification(file_bytes: bytes, filename: str) -> tuple:
    """Returns (text_content: str, image_b64: str|None, metadata_str: str)."""
    import os, io
    ext      = os.path.splitext(filename.lower())[1]
    ftype    = _EXT_TO_TYPE.get(ext, "unknown")
    meta_str = ""

    if ftype == "excel":
        return _extract_excel_classify(file_bytes, filename, meta_str)
    if ftype == "csv":
        return _extract_csv_classify(file_bytes, meta_str)
    if ftype == "pdf":
        return _extract_pdf_classify(file_bytes, filename, meta_str)
    if ftype == "word":
        return _extract_word_classify(file_bytes, filename, meta_str)
    if ftype == "pptx":
        return _extract_pptx_classify(file_bytes, filename, meta_str)
    if ftype == "legacy_doc":
        return _extract_fitz_classify(file_bytes, filename, meta_str)
    if ftype == "image":
        return _extract_image_classify(file_bytes, filename, meta_str)
    if ftype == "text":
        return _extract_text_classify(file_bytes, meta_str)
    if ftype == "html":
        return _extract_html_classify(file_bytes, meta_str)
    if ftype == "msg":
        return _extract_msg_classify(file_bytes, filename, meta_str)
    return "[Unsupported file type]", None, meta_str


def _extract_excel_classify(file_bytes, filename, meta_str):
    import io, pandas as pd
    _MAX_SHEETS = 8
    buf  = io.BytesIO(file_bytes)
    parts: list[str] = []
    try:
        for engine in ("openpyxl", "xlrd"):
            try:
                buf.seek(0)
                xls = pd.ExcelFile(buf, engine=engine)
                break
            except Exception:
                continue
        else:
            return "[Excel could not be opened]", None, meta_str

        sheets = xls.sheet_names
        meta_str = f"- sheet_count: {len(sheets)}\n- sheet_names: {sheets}\n- multi_sheet: {len(sheets) > 1}\n"
        parts.append(f"## Workbook Structure\nTotal Sheets: {len(sheets)}\nSheet Names: {sheets}\n")
        for idx, sheet in enumerate(sheets):
            try:
                df = pd.read_excel(xls, sheet_name=sheet, nrows=100, header=None)
                df = df.iloc[:, :30].dropna(how="all").dropna(axis=1, how="all")
                non_empty = len(df)
                if idx < _MAX_SHEETS and non_empty > 0:
                    parts.append(f"### Sheet {idx+1}: '{sheet}'")
                    parts.append(f"Non-empty rows in sample: {non_empty}")
                    try:
                        first_row = df.iloc[0]
                        if all(isinstance(x, str) for x in first_row if x != ""):
                            hdrs = [str(x).strip() for x in first_row]
                            df = df.iloc[1:]
                        else:
                            hdrs = [f"col_{c}" for c in range(df.shape[1])]
                        parts.append(df.fillna("").astype(str).to_markdown(index=False, headers=hdrs))
                    except Exception:
                        parts.append(df.fillna("").astype(str).to_string(index=False))
                    parts.append("")
                else:
                    parts.append(f"### Sheet {idx+1}: '{sheet}' (summary only — {non_empty} non-empty rows)\n")
            except Exception:
                pass
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Excel extraction error: {exc}]", None, meta_str


def _extract_csv_classify(file_bytes, meta_str):
    import io, pandas as pd
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), nrows=40, header=None).iloc[:, :30]
        try:
            text = df.fillna("").astype(str).to_markdown(index=False)
        except Exception:
            text = df.fillna("").astype(str).to_string(index=False)
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[CSV extraction error: {exc}]", None, meta_str


def _extract_pdf_classify(file_bytes, filename, meta_str):
    import io
    _MAX_PAGES = 10
    try:
        import pdfplumber
        buf  = io.BytesIO(file_bytes)
        with pdfplumber.open(buf) as pdf:
            total = len(pdf.pages)
            meta_str = f"- total_pages: {total}\n- pages_processed: {min(total, _MAX_PAGES)}\n"
            if total == 0:
                return "[Empty or corrupt PDF — no pages found]", None, meta_str

            # smart selection: first N-2 pages + last 2 pages (captures pricing summaries at end)
            if total <= _MAX_PAGES:
                page_idxs = list(range(total))
            else:
                page_idxs = sorted(set(list(range(_MAX_PAGES - 2)) + list(range(total - 2, total))))

            text_parts: list[str] = []
            table_parts: list[str] = []
            for i in page_idxs:
                page = pdf.pages[i]
                t = (page.extract_text() or "").strip()
                if t:
                    text_parts.append(f"--- Page {i+1} ---\n{t}")
                # extract tables from this page
                for t_idx, table in enumerate(page.extract_tables() or []):
                    if not table:
                        continue
                    header = table[0]
                    rows   = table[1:]
                    tstr   = " | ".join(str(c) for c in header) + "\n"
                    tstr  += " | ".join("---" for _ in header) + "\n"
                    for row in rows[:20]:
                        tstr += " | ".join(str(c) for c in row) + "\n"
                    table_parts.append(f"Table {t_idx+1} (Page {i+1}):\n{tstr}")

            combined = "\n\n".join(text_parts)
            if table_parts:
                combined += "\n\n### Extracted Tables\n" + "\n\n".join(table_parts)

            if len(combined.strip()) < 50:
                return _extract_pdf_as_image_classify(file_bytes, filename, meta_str, total)
            return combined[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception:
        return _extract_pdf_fitz_classify(file_bytes, filename, meta_str)


def _extract_pdf_fitz_classify(file_bytes, filename, meta_str):
    import io, base64
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            total = len(doc)
            meta_str = f"- total_pages: {total}\n"
            pages = list(range(min(total, 8)))
            texts = [doc[i].get_text("text").strip() for i in pages]
            combined = "\n\n".join(f"--- Page {i+1} ---\n{t}" for i, t in zip(pages, texts) if t)
            if len(combined.strip()) < 50:
                pix = doc[0].get_pixmap(dpi=150)
                b64 = base64.b64encode(pix.tobytes("png")).decode()
                pix = None
                return "[Scanned PDF - content sent as image]", b64, meta_str
            return combined[:_MAX_CLASSIFY_CHARS], None, meta_str
        finally:
            doc.close()
    except Exception as exc:
        return f"[PDF extraction error: {exc}]", None, meta_str


def _extract_pdf_as_image_classify(file_bytes, filename, meta_str, total):
    import io, base64
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            img = pdf.pages[0].to_image(resolution=200)
            buf = io.BytesIO()
            try:
                img.original.save(buf, format="PNG")
                b64 = base64.b64encode(buf.getvalue()).decode()
            finally:
                try: img.original.close()
                except Exception: pass
        return "[Scanned PDF - content sent as image]", b64, meta_str
    except Exception:
        return _extract_pdf_fitz_classify(file_bytes, filename, meta_str)


def _extract_word_classify(file_bytes, filename, meta_str):
    import io, zipfile, base64
    try:
        from docx import Document
        from PIL import Image
        doc   = Document(io.BytesIO(file_bytes))
        parts: list[str] = []
        para_count = 0
        for para in doc.paragraphs:
            if para_count >= 200:
                break
            text = para.text.strip()
            if not text:
                continue
            if para.style and para.style.name.startswith("Heading"):
                lvl = para.style.name.replace("Heading ", "").strip()
                try:
                    lvl = int(lvl)
                except ValueError:
                    lvl = 1
                parts.append(f"{'#' * lvl} {text}")
            else:
                parts.append(text)
            para_count += 1
        for i, tbl in enumerate(doc.tables):
            rows = []
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(cells)
            if rows:
                hdr = rows[0]
                tstr  = " | ".join(hdr) + "\n"
                tstr += " | ".join("---" for _ in hdr) + "\n"
                for row in rows[1:20]:
                    tstr += " | ".join(row) + "\n"
                parts.append(f"\n### Table {i+1}\n{tstr}")
        text = "\n\n".join(parts)
        meta_str = f"- paragraphs: {len(doc.paragraphs)}\n- tables: {len(doc.tables)}\n"
        if len(text.strip()) < 50:
            # Fallback: walk all <w:t> elements in the XML tree. Catches content
            # wrapped inside SDT (Structured Document Tags), text boxes, content
            # controls, or any other non-standard container that python-docx's
            # .paragraphs / .tables iterators skip.
            try:
                _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                xml_texts: list[str] = []
                for t_el in doc.element.body.iter(f"{{{_wns}}}t"):
                    if t_el.text:
                        xml_texts.append(t_el.text)
                xml_text = " ".join(xml_texts).strip()
                if len(xml_text) >= 50:
                    text = xml_text
                    meta_str += "- sdt_fallback: true\n"
                    logger.info(
                        "Word SDT fallback recovered {} chars for {!r}",
                        len(xml_text), filename,
                    )
            except Exception:
                pass
        if len(text.strip()) < 50:
            # try to extract largest embedded image from word/media/
            try:
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                    imgs = sorted([n for n in z.namelist() if n.startswith("word/media/")])
                    if imgs:
                        largest = max(imgs, key=lambda n: z.getinfo(n).file_size)
                        img_bytes = z.read(largest)
                        img = Image.open(io.BytesIO(img_bytes))
                        try:
                            if img.mode not in ("RGB", "L"):
                                img = img.convert("RGB")
                            if max(img.size) > 2048:
                                img.thumbnail((2048, 2048), Image.LANCZOS)
                            buf2 = io.BytesIO()
                            img.save(buf2, format="PNG")
                            b64 = base64.b64encode(buf2.getvalue()).decode()
                        finally:
                            try: img.close()
                            except Exception: pass
                        return "[Word document is image-based - content sent as image for visual analysis]", b64, meta_str
            except Exception:
                pass
            return _extract_fitz_classify(file_bytes, filename, meta_str)
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Word extraction error: {exc}]", None, meta_str


def _extract_pptx_classify(file_bytes, filename, meta_str):
    import io
    _MAX_SLIDES = 15
    try:
        from pptx import Presentation
        buf = io.BytesIO(file_bytes)
        prs = Presentation(buf)
        slides = list(prs.slides)
        total  = len(slides)
        parts  = [f"## Presentation: {total} slides\n"]
        for i, slide in enumerate(slides[:_MAX_SLIDES]):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_parts.append(t)
                if shape.has_table:
                    tbl = shape.table
                    rows = [[c.text.strip() for c in row.cells] for row in tbl.rows]
                    if rows:
                        hdr   = rows[0]
                        tstr  = " | ".join(hdr) + "\n"
                        tstr += " | ".join("---" for _ in hdr) + "\n"
                        for row in rows[1:20]:
                            tstr += " | ".join(row) + "\n"
                        slide_parts.append(tstr)
            if slide_parts:
                parts.append(f"### Slide {i+1}")
                parts.append("\n".join(slide_parts))
                parts.append("")
        meta_str = f"- total_slides: {total}\n- slides_extracted: {min(total, _MAX_SLIDES)}\n"
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return _extract_fitz_classify(file_bytes, filename, meta_str)


def _extract_fitz_classify(file_bytes, filename, meta_str):
    import base64
    try:
        import fitz
        doc = fitz.open(stream=file_bytes)
        try:
            total = len(doc)
            pages = list(range(min(total, 4)))
            images: list[str] = []
            for i in pages:
                pix = doc[i].get_pixmap(dpi=150)
                images.append(base64.b64encode(pix.tobytes("png")).decode())
                pix = None
        finally:
            doc.close()
        if not images:
            return "[Document could not be rendered]", None, meta_str
        return "[Legacy document - content sent as images]", images[0], meta_str
    except Exception as exc:
        return f"[Document extraction error: {exc}]", None, meta_str


def _extract_image_classify(file_bytes, filename, meta_str):
    import io, base64
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        try:
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")
            if max(img.size) > 2048:
                img.thumbnail((2048, 2048))
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode()
            meta_str = f"- size: {img.size[0]}x{img.size[1]}\n"
        finally:
            try: img.close()
            except Exception: pass
        return "[Image file - content sent as image]", b64, meta_str
    except Exception as exc:
        return f"[Image extraction error: {exc}]", None, meta_str


def _extract_text_classify(file_bytes, meta_str):
    try:
        import chardet
        enc = (chardet.detect(file_bytes).get("encoding") or "utf-8")
        try:
            text = file_bytes.decode(enc)
        except Exception:
            text = file_bytes.decode("utf-8", errors="replace")
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Text extraction error: {exc}]", None, meta_str


def _extract_html_classify(file_bytes, meta_str):
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(file_bytes, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
    except Exception:
        raw  = file_bytes.decode("utf-8", errors="replace")
        text = re.sub(r"<[^>]+>", " ", raw)
        text = re.sub(r"\s+", " ", text).strip()
    return text[:_MAX_CLASSIFY_CHARS], None, meta_str


def _extract_msg_classify(file_bytes, filename, meta_str):
    import os, tempfile
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject: parts.append(f"Subject: {msg.subject}")
            if msg.sender:  parts.append(f"From: {msg.sender}")
            if msg.body:    parts.append(msg.body)
            msg.close()
        finally:
            os.unlink(tmp_path)
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[MSG extraction error: {exc}]", None, meta_str


# ── Trivial image heuristic ────────────────────────────────────────────────────

def _is_trivial_image(file_bytes: bytes, filename: str) -> tuple:
    import os
    if os.path.splitext(filename.lower())[1] not in _IMAGE_EXTS:
        return False, ""
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        img.load()
        w, h = img.size
        long_e  = max(w, h)
        short_e = max(min(w, h), 1)
        if long_e < 200:           return True, f"tiny {w}x{h}"
        if short_e < 150:          return True, f"short edge {w}x{h}"
        if long_e / short_e > 5:   return True, f"extreme aspect {w}x{h}"
        if long_e / short_e > 3 and short_e < 250:
            return True, f"banner {w}x{h}"
        if long_e < 400 and short_e < 200:
            return True, f"sub-thumbnail {w}x{h}"
        gray = img.convert("L")
        hist = gray.histogram()
        total = sum(hist)
        if total and sum(hist[230:]) / total > 0.97:
            return True, "mostly white"
        try:
            colors = img.convert("RGB").getcolors(maxcolors=256)
            if colors is not None and len(colors) < 16:
                return True, f"flat palette ({len(colors)} colors)"
        except Exception:
            pass
    except Exception as exc:
        return True, f"unreadable ({exc.__class__.__name__})"
    return False, ""


# ── Classification ─────────────────────────────────────────────────────────────

def _get_prompt_text(msg_input, default: str) -> str:
    """Extract text from an optional wired Message input; fall back to default."""
    if msg_input is None:
        return default
    text = getattr(msg_input, "text", None)
    if text and str(text).strip():
        return str(text).strip()
    return default


# ── LLM rate-limit retry helpers ──────────────────────────────────────────────

def _is_rate_limit_error(exc: Exception) -> bool:
    """Return True for 429 / token-quota / rate-limit errors from any LLM provider."""
    msg = str(exc).lower()
    return any(k in msg for k in (
        "429", "rate limit", "rate_limit", "ratelimit",
        "too many requests", "resource exhausted",
        "quota exceeded", "token quota", "tokens exceeded",
        "insufficient_quota", "token limit exceeded",
        "requests per minute", "tokens per minute",
    ))


def _is_timeout_error(exc: Exception) -> bool:
    """Return True if the LLM call timed out (our wrapper or the SDK itself).

    A silent stuck call is the #1 cause of MiCore appearing frozen mid-batch
    on large workloads. Treating timeouts as retryable lets the worker
    recover instead of hanging forever on a dropped TCP connection or a
    slow Azure OpenAI gateway response.
    """
    if isinstance(exc, TimeoutError):
        return True
    msg = str(exc).lower()
    return any(k in msg for k in (
        "timed out", "timeout", "read timed out", "request timed out",
        "operation timed out", "deadline exceeded",
    ))


def _call_with_timeout(fn, *args, timeout_s: int = 120, label: str = "remote-call", **kwargs):
    """Run any callable with a hard wall-clock timeout.

    The #1 cause of "silent stuck" pipelines is a network call with no
    timeout: AzureChatOpenAI, AzureOpenAIEmbeddings, and the MiCore
    Pinecone helpers (search_via_service / ingest_via_service) all have
    no default request timeout, so a hung socket read blocks the worker
    forever with no exception and no log line.

    Uses a daemon thread + queue rather than ThreadPoolExecutor.
    ThreadPoolExecutor's `with` block calls shutdown(wait=True) on exit,
    which means a stuck inner thread would deadlock the timeout itself
    — the caller would wait forever for the executor to clean up. With
    multiple workers all triggering this concurrently the entire pipeline
    freezes (worked with 1 worker, hung at 2+).

    The daemon-thread approach: launch the work on a daemon, wait on a
    queue with timeout, and ORPHAN the inner thread on timeout. The
    orphaned thread continues running (consuming whatever it was doing)
    but the caller is unblocked and the surrounding retry logic fires.
    Daemons don't keep the process alive on shutdown.
    """
    import threading as _th
    import queue as _q

    result_q: _q.Queue = _q.Queue(maxsize=1)

    def _runner() -> None:
        try:
            result_q.put(("ok", fn(*args, **kwargs)))
        except BaseException as exc:  # noqa: BLE001 — re-raised below
            result_q.put(("err", exc))

    t = _th.Thread(target=_runner, name=label, daemon=True)
    t.start()
    try:
        kind, value = result_q.get(timeout=timeout_s)
    except _q.Empty:
        raise TimeoutError(
            f"{label} exceeded {timeout_s}s wall-clock — aborted; "
            f"inner thread orphaned (will finish in background or die with process)"
        )
    if kind == "err":
        raise value
    return value


def _invoke_llm_with_timeout(llm, messages: list, timeout_s: int = 120):
    """Thin wrapper kept for backward compatibility — delegates to _call_with_timeout."""
    return _call_with_timeout(llm.invoke, messages, timeout_s=timeout_s, label="llm-call")


def _embed_with_timeout(embed_model, text: str, timeout_s: int = 60):
    """Embed text with a hard timeout. Default 60s — embedding endpoints
    are typically faster than chat completions; if we're past 60s on a
    single embed, the call is hung."""
    return _call_with_timeout(embed_model.embed_query, text, timeout_s=timeout_s, label="embed-query")


def _call_pinecone_with_retry(
    fn,
    *args,
    timeout_s: int,
    label: str,
    max_attempts: int = 3,
    **kwargs,
):
    """Run a Pinecone call with hard timeout AND retry on transient errors.

    Mirrors the agentcore pinecone-service tenacity policy: 3 attempts,
    exponential backoff (1s → 2s → 4s, capped at 10s), retries only on
    transient server/connectivity errors as classified by
    `_is_pinecone_server_error` (connection, timeout, 5xx, network, etc.).
    Non-transient errors (auth, schema, invalid args) raise immediately
    without burning retry attempts.

    The caller's existing exception handling (fall-back-to-benchmark-stub
    in Stage 7, log-and-continue in Stage 6) is preserved — retries
    happen inside this helper, and the final exception still surfaces if
    all attempts fail.
    """
    import time as _time
    last_exc: Exception | None = None
    for attempt in range(1, max_attempts + 1):
        try:
            return _call_with_timeout(
                fn, *args, timeout_s=timeout_s, label=label, **kwargs
            )
        except Exception as exc:
            last_exc = exc
            transient = _is_pinecone_server_error(exc)
            if attempt >= max_attempts or not transient:
                if transient:
                    logger.warning(
                        f"[{label}] all {max_attempts} attempt(s) failed — "
                        f"surfacing {type(exc).__name__}: {exc}"
                    )
                raise
            sleep_s = min(10, 2 ** (attempt - 1))  # 1, 2, 4 …
            logger.warning(
                f"[{label}] attempt {attempt}/{max_attempts} failed "
                f"({type(exc).__name__}: {exc}); retrying in {sleep_s}s…"
            )
            _time.sleep(sleep_s)
    # Defensive — loop above always either returns or raises.
    raise last_exc if last_exc else RuntimeError(f"[{label}] retry loop exited unexpectedly")


def _pinecone_search_with_timeout(search_fn, *args, timeout_s: int = 60, **kwargs):
    """Pinecone query with hard timeout + 3-attempt retry on transient errors."""
    return _call_pinecone_with_retry(
        search_fn, *args, timeout_s=timeout_s, label="pinecone-search", **kwargs,
    )


def _pinecone_ingest_with_timeout(ingest_fn, *args, timeout_s: int = 120, **kwargs):
    """Pinecone upsert with hard timeout + 3-attempt retry on transient errors.
    Larger timeout budget than search because upserts can include big batches
    of vectors + metadata."""
    return _call_pinecone_with_retry(
        ingest_fn, *args, timeout_s=timeout_s, label="pinecone-ingest", **kwargs,
    )


def _pinecone_ensure_with_retry(ensure_fn, *args, timeout_s: int = 60, **kwargs):
    """Pinecone ensure-index with hard timeout + 3-attempt retry on transient
    errors. Used at the start of Stage 6 so a flaky Pinecone start doesn't
    immediately fail the PR — the index check is harmless to repeat."""
    return _call_pinecone_with_retry(
        ensure_fn, *args, timeout_s=timeout_s, label="pinecone-ensure-index", **kwargs,
    )


def _is_json_format_error(exc: Exception) -> bool:
    """Detect when the LLM rejects the response_format=json_object kwarg.

    Older models / some Azure deployments don't support strict JSON mode and
    raise a 400-class BadRequestError mentioning response_format / json_object.
    """
    msg = str(exc).lower()
    return ("response_format" in msg) or ("json_object" in msg) or (
        "json mode" in msg
    )


def _call_llm_with_retry(
    llm,
    messages: list,
    prompts: dict | None = None,
    *,
    force_json: bool = False,
):
    """Invoke the LLM with automatic retry + cooldown on rate-limit / token-quota errors.

    Retry count and cooldown are read from the prompts dict so they can be set
    from the controller (llm_max_retries, llm_retry_cooldown keys).
    Cooldown is progressive: cooldown × (attempt+1) so each retry waits longer.

    When force_json=True, binds response_format={"type": "json_object"} on the
    LLM (mirrors doc-intel file_classifier/classifier/llm_client.py). If the
    model rejects the kwarg, falls back transparently to plain invoke so older
    deployments keep working.
    """
    import time
    max_retries = int((prompts or {}).get("llm_max_retries",   3))
    cooldown_s  = int((prompts or {}).get("llm_retry_cooldown", 60))
    timeout_s   = int((prompts or {}).get("llm_call_timeout_s", 120))

    # Try to bind strict JSON output. Most modern Azure / OpenAI chat models
    # support it; older ones reject the kwarg and we fall back below on first
    # invoke. .bind() itself is lazy so it won't fail here.
    active_llm = llm
    json_mode_enabled = False
    if force_json:
        try:
            active_llm = llm.bind(response_format={"type": "json_object"})
            json_mode_enabled = True
        except Exception:
            active_llm = llm  # bind unsupported → use plain LLM

    for attempt in range(max_retries + 1):
        try:
            return _invoke_llm_with_timeout(active_llm, messages, timeout_s=timeout_s)
        except Exception as exc:
            # If JSON mode tripped a model-not-supported error, drop it and
            # retry once with the unbound LLM. This is a recovery path, not a
            # rate-limit retry, so it doesn't consume a retry slot.
            if json_mode_enabled and _is_json_format_error(exc):
                logger.warning(
                    "LLM rejected response_format=json_object — "
                    "retrying without JSON mode: {}",
                    exc,
                )
                active_llm = llm
                json_mode_enabled = False
                continue
            # Treat rate-limit AND timeout as retryable — both indicate a
            # transient remote condition (quota / hung socket) rather than
            # a programming error.
            is_rl = _is_rate_limit_error(exc)
            is_to = _is_timeout_error(exc)
            if (is_rl or is_to) and attempt < max_retries:
                wait = cooldown_s * (attempt + 1)  # 60s → 120s → 180s
                kind = "rate-limit" if is_rl else "timeout"
                logger.warning(
                    "LLM {} (attempt {}/{}) — cooling down {}s before retry: {}",
                    kind, attempt + 1, max_retries, wait, exc,
                )
                time.sleep(wait)
            else:
                raise


def _render_pdf_pages_to_b64(
    file_bytes: bytes,
    filename: str,
    *,
    max_pages: int = 1,
    timeout_s: int = 30,
) -> list[str]:
    """Render the first `max_pages` pages of a PDF to base64-encoded PNGs.

    Used by vision-escalation in _classify_file. Two safety nets so a
    single malformed PDF can't silently hang a worker forever (which was
    a likely cause of 500-PR batches "getting stuck at one point only"):

      1. The entire render runs inside _call_with_timeout with a
         configurable wall-clock budget (default 30s). pdfplumber's
         pdfminer-based parser can spin in pure-Python loops on corrupt
         cross-reference tables; the timeout interrupts those.

      2. If pdfplumber fails or times out, fall back to PyMuPDF (fitz)
         which uses a completely different parser and often handles
         malformed PDFs that pdfplumber chokes on.

    Returns an empty list on any failure — caller should treat empty as
    "vision escalation not possible, keep text classification".
    """
    import io as _io, base64 as _b64

    def _render_with_pdfplumber() -> list[str]:
        import pdfplumber
        out: list[str] = []
        with pdfplumber.open(_io.BytesIO(file_bytes)) as pdf:
            n = min(len(pdf.pages), max_pages)
            for i in range(n):
                img = pdf.pages[i].to_image(resolution=150)
                buf = _io.BytesIO()
                img.original.save(buf, format="PNG")
                out.append(_b64.b64encode(buf.getvalue()).decode())
        return out

    def _render_with_fitz() -> list[str]:
        import fitz
        out: list[str] = []
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            n = min(len(doc), max_pages)
            for i in range(n):
                pix = doc[i].get_pixmap(dpi=150)
                out.append(_b64.b64encode(pix.tobytes("png")).decode())
                pix = None
        finally:
            doc.close()
        return out

    # Try pdfplumber first (matches prior behaviour); fall back to fitz.
    for label, fn in (("pdfplumber", _render_with_pdfplumber), ("fitz", _render_with_fitz)):
        try:
            return _call_with_timeout(fn, timeout_s=timeout_s, label=f"pdf-render-{label}")
        except TimeoutError as exc:
            logger.warning(
                "PDF render for vision sample timed out ({}) on {}: {} — trying next renderer",
                label, filename, exc,
            )
        except Exception as exc:
            logger.warning(
                "PDF render for vision sample failed ({}) on {}: {} — trying next renderer",
                label, filename, exc,
            )
    return []


def _classify_file(llm, file_bytes: bytes, filename: str, prompts: dict | None = None) -> tuple:
    """Returns (doc_type: str, confidence: float)."""
    import os
    from langchain_core.messages import HumanMessage, SystemMessage

    p = prompts or {}
    sys_prompt      = p.get("cls_system", CLASSIFICATION_SYSTEM_PROMPT)
    user_text_tmpl  = p.get("cls_user_text", _CLASSIFY_USER_TEXT)
    user_image_tmpl = p.get("cls_user_image", _CLASSIFY_USER_IMAGE)

    ext = os.path.splitext(filename.lower())[1]
    if ext not in _SUPPORTED_CLASSIFY_EXTS:
        return "Others", 0.0

    trivial, reason = _is_trivial_image(file_bytes, filename)
    if trivial:
        logger.info(f"Trivial image {filename!r} ({reason}) → Others")
        return "Others", 0.0

    file_type = _EXT_TO_TYPE.get(ext, "unknown")
    text_content, image_b64, meta_str = _extract_for_classification(file_bytes, filename)

    cls_max = (prompts or {}).get("cls_max_chars")
    if cls_max and text_content and len(text_content) > cls_max:
        text_content = text_content[:cls_max]

    try:
        if image_b64:
            user_prompt = user_image_tmpl.format(
                filename=filename, file_type=file_type, extra_metadata=meta_str
            )
            content = [
                {"type": "text", "text": user_prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}", "detail": "high"}},
            ]
            messages = [SystemMessage(content=sys_prompt), HumanMessage(content=content)]
        else:
            user_prompt = user_text_tmpl.format(
                filename=filename, file_type=file_type,
                extra_metadata=meta_str, extracted_content=text_content,
            )
            messages = [SystemMessage(content=sys_prompt), HumanMessage(content=user_prompt)]

        # force_json=True mirrors doc-intel's response_format={"type":"json_object"}.
        # Without this the LLM may emit prose / markdown and json.loads() below
        # fails, dropping the file to "Others" — a major source of misclassification.
        response = _call_llm_with_retry(llm, messages, prompts, force_json=True)
        raw = (getattr(response, "content", None) or str(response)).strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        result = json.loads(raw)

        raw_cls = result.get("classification", "Other")
        if raw_cls not in VALID_CLASSIFICATIONS:
            raw_cls = "Other"
        confidence = float(result.get("confidence", 0.0))
        
        # --- ESCALATION LOGIC ---
        # If classification is 'Other' or confidence is low, and it is a PDF
        # (and didn't already use image_b64 for classification), escalate to
        # the vision model using a rendered page image.
        needs_escalation = (raw_cls == "Other" or confidence < 0.75)
        if needs_escalation and not image_b64 and file_type == "pdf":
            max_pages_vision = max(1, int((prompts or {}).get("cls_max_pages_vision", 1)))
            render_timeout_s = max(5, int((prompts or {}).get("cls_render_timeout_s", 30)))
            vision_samples_b64 = _render_pdf_pages_to_b64(
                file_bytes, filename,
                max_pages=max_pages_vision,
                timeout_s=render_timeout_s,
            )

            # Only call the vision LLM when we actually have image samples.
            # The earlier indentation had this call running even when
            # rendering failed — wasted an LLM call on a redundant text-only
            # classification using the previous `messages` value.
            if vision_samples_b64:
                logger.info(
                    f"Escalating classification for {filename!r} to vision using "
                    f"{len(vision_samples_b64)} page(s) (conf={confidence:.2f}, cls={raw_cls})"
                )
                user_prompt = user_image_tmpl.format(
                    filename=filename, file_type=file_type, extra_metadata=meta_str
                )
                content = [{"type": "text", "text": user_prompt}]
                for b64 in vision_samples_b64:
                    content.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"},
                    })
                messages = [SystemMessage(content=sys_prompt), HumanMessage(content=content)]
                response = _call_llm_with_retry(llm, messages, prompts, force_json=True)
                raw = (getattr(response, "content", None) or str(response)).strip()
                raw = re.sub(r"^```(?:json)?\s*", "", raw)
                raw = re.sub(r"\s*```$", "", raw)
                result = json.loads(raw)

                raw_cls = result.get("classification", "Other")
                if raw_cls not in VALID_CLASSIFICATIONS:
                    raw_cls = "Other"
                confidence = float(result.get("confidence", 0.0))
            else:
                logger.info(
                    f"Vision escalation skipped for {filename!r} — render produced no "
                    f"samples (conf={confidence:.2f}, cls={raw_cls}); keeping text classification"
                )

        doc_type   = _CLASSIFICATION_MAP.get(raw_cls, raw_cls)
        return doc_type, confidence

    except json.JSONDecodeError as exc:
        # Surface JSON parse failures explicitly so the user can see in the
        # canvas log that the LLM output was unparseable (not just a generic
        # classification miss).
        logger.opt(exception=True).warning(
            "Classification JSON parse failed for {!r}: {} — raw response was: {!r}",
            filename, exc, (raw[:500] if 'raw' in locals() else '<no response>'),
        )
        return "Others", 0.0
    except Exception as exc:
        logger.opt(exception=True).warning(
            "Classification failed for {!r}: {}", filename, exc
        )
        return "Others", 0.0


# ── Stage 4: run classification for a PR ──────────────────────────────────────

def _run_classification(llm, tgt_cs: str, blob_cfg: dict, pr_no: str, prompts: dict | None = None) -> None:
    import os
    from concurrent.futures import ThreadPoolExecutor, as_completed

    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT ac.[attachment_classify_uuid_pk], ac.[file_path], ac.[attachment_id]
              FROM [ras_procurement].[attachment_classification] ac
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
        """, pr_no)
        parent_rows = cur.fetchall()

        cur.execute("""
            SELECT ec.[embedded_attachment_classification_id], ec.[file_path],
                   ac.[attachment_classify_uuid_pk]
              FROM [ras_procurement].[embedded_attachment_classification] ec
              JOIN [ras_procurement].[attachment_classification] ac
                ON ec.[attachment_classification_id] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
        """, pr_no)
        embedded_rows = cur.fetchall()
    finally:
        conn.close()

    # Each file: download blob + LLM call + single-row DB update — fully independent.
    # All three operations use separate connections / different rows, so concurrent
    # execution is safe. LangChain LLMs use httpx connection pooling (thread-safe).
    # ProcessPoolExecutor is NOT used because llm/embed_model are not picklable.

    def _classify_parent(row):
        att_pk, blob_path, att_id = str(row[0]), row[1], row[2]
        if not blob_path:
            return
        try:
            file_bytes     = _download_blob(blob_path, blob_cfg)
            filename       = os.path.basename(blob_path)
            doc_type, conf = _classify_file(llm, file_bytes, filename, prompts)
            _update_parent_classification(tgt_cs, att_id, doc_type, conf)
            logger.info(f"[{pr_no}] Parent {filename!r}: {doc_type} (conf={conf:.2f})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Parent classification failed ({blob_path!r}): {exc}")

    def _classify_embedded(row):
        emb_pk, blob_path, parent_pk = str(row[0]), row[1], str(row[2])
        if not blob_path:
            return
        try:
            file_bytes     = _download_blob(blob_path, blob_cfg)
            filename       = os.path.basename(blob_path)
            doc_type, conf = _classify_file(llm, file_bytes, filename, prompts)
            _update_embedded_classification(tgt_cs, parent_pk, blob_path, doc_type, conf)
            logger.info(f"[{pr_no}] Embedded {filename!r}: {doc_type} (conf={conf:.2f})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Embedded classification failed ({blob_path!r}): {exc}")

    tasks = (
        [(_classify_parent,   r) for r in parent_rows] +
        [(_classify_embedded, r) for r in embedded_rows]
    )
    if not tasks:
        return

    n_workers = max(1, int((prompts or {}).get("cls_parallel_sources", 8)))
    with ThreadPoolExecutor(max_workers=min(len(tasks), n_workers)) as inner_pool:
        futures = [inner_pool.submit(fn, row) for fn, row in tasks]
        for f in as_completed(futures):
            f.result()  # exceptions already caught inside fn; this just awaits all


def _update_parent_classification(tgt_cs: str, att_id: str, doc_type: str, conf: float) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[attachment_classification]
               SET [doc_type]            = ?,
                   [classification_conf] = ?,
                   [updated_at]          = SYSUTCDATETIME()
             WHERE [attachment_id] = ?
        """, doc_type, conf, att_id)
        conn.commit()
    finally:
        conn.close()


def _update_embedded_classification(tgt_cs: str, parent_pk: str, blob_path: str, doc_type: str, conf: float) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[embedded_attachment_classification]
               SET [doc_type]            = ?,
                   [classification_conf] = ?,
                   [updated_at]          = SYSUTCDATETIME()
             WHERE [attachment_classification_id] = ?
               AND [file_path] = ?
        """, doc_type, conf, parent_pk, blob_path)
        conn.commit()
    finally:
        conn.close()


# ── RAS context builder ────────────────────────────────────────────────────────

def _build_ras_context(tgt_cs: str, pr_no: str) -> Optional[RASContext]:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT TOP 1
                   prm.[PURCHASE_REQ_ID], prm.[SUPPLIER_NAME], prm.[JUSTIFICATION],
                   prm.[CURRENCY], prm.[ENQUIRY_NO], prm.[CLASSIFICATION],
                   prm.[Department], prm.[NEGOTIATED_BY], prm.[ADDRESS],
                   prm.[CONTRACT_NO], prm.[ORDER_NO], prm.[PURCHASE_VALUE],
                   prm.[REQ_START_DATE], prm.[C_DATETIME]
              FROM [ras_procurement].[purchase_req_mst] prm
             WHERE prm.[PURCHASE_REQ_NO] = ?
        """, pr_no)
        mst = cur.fetchone()
        if not mst:
            return None

        req_id = mst[0]
        def _d(v): return Decimal(str(v)) if v is not None else None

        cur.execute("""
            SELECT prd.[PURCHASE_DTL_ID], prd.[PURCHASE_REQ_ID], prd.[ITEM_NO],
                   prd.[QUANTITY], prd.[ITEM_TYPE], prd.[ITEMDESCRIPTION],
                   prd.[PRICE], prd.[UOM], prd.[DISCOUNT], prd.[REQ_VALUE],
                   prd.[CURRENCY], prd.[DELIVERY_DATE], prd.[SUPPLIER_NAME],
                   prd.[PAYMENT_DETAILS], prd.[ORIGINAL_VALUE], prd.[Initial_Offer],
                   prd.[Negotiation], prd.[CommentsforItem], prd.[PREPAYMENT], prd.[ITEM_CODE]
              FROM [ras_procurement].[purchase_req_detail] prd
             WHERE prd.[PURCHASE_REQ_ID] = ?
             ORDER BY prd.[ITEM_NO]
        """, req_id)
        dtl_rows = cur.fetchall()

        line_items: list[LineItemContext] = []
        for r in dtl_rows:
            raw_del = r[11]
            del_str = (
                raw_del.date().isoformat() if hasattr(raw_del, "date")
                else str(raw_del) if raw_del else None
            )
            line_items.append(LineItemContext(
                purchase_dtl_id=r[0], purchase_req_id=r[1], item_no=r[2] or 0,
                quantity=_d(r[3]), item_type=r[4], item_description=r[5],
                unit_price=_d(r[6]), uom=r[7], discount=_d(r[8]),
                req_value=_d(r[9]), currency=r[10], delivery_date=del_str,
                supplier_name=r[12] or mst[1], payment_details=r[13],
                original_value=_d(r[14]), initial_offer=_d(r[15]),
                negotiation=_d(r[16]), comments=r[17], prepayment=r[18], item_code=r[19],
            ))

        # BI dashboard enrichment (best-effort)
        cat=sub_cat=site_c=site_r=div=site=req_t=l3=l4=pur_cat=ras_t=par_s=sup_t=sup_c=pay_d=po_d=cat_b=None
        try:
            cur.execute("""
                SELECT TOP 1 vw.[L1], vw.[Sub_Category_Type], vw.[Site_Country],
                             vw.[Site_Region], vw.[Division], vw.[L3], vw.[L4],
                             vw.[Purchase_Category], vw.[Title], vw.[Site],
                             vw.[Requisition_Type], vw.[Parent_Supplier], vw.[Supplier_Type],
                             vw.[Suplier_country], vw.[Payment_Days], vw.[PO_Date],
                             vw.[Category_Buyer], vw.[L2]
                  FROM vw_get_ras_data_for_bidashboard vw
                 WHERE vw.[PURCHASE_REQ_ID] = ?
            """, req_id)
            vw = cur.fetchone()
            if vw:
                cat,sub_cat,site_c,site_r,div,l3,l4 = vw[0],vw[1]or vw[17],vw[2],vw[3],vw[4],vw[5],vw[6]
                pur_cat,ras_t,site,req_t = vw[7],vw[8],vw[9],vw[10]
                par_s,sup_t,sup_c = vw[11],vw[12],vw[13]
                pay_d,po_d,cat_b  = vw[14],vw[15],vw[16]
        except Exception:
            pass

        def _to_date(v):
            if v is None: return None
            return v.date() if hasattr(v, "date") else v

        return RASContext(
            purchase_req_no=pr_no, purchase_req_id=req_id,
            supplier_name=mst[1], justification=mst[2], currency=mst[3],
            enquiry_no=mst[4], classification=mst[5], department=mst[6],
            negotiated_by=mst[7], address=mst[8], contract_no=mst[9],
            order_no=mst[10], purchase_value=_d(mst[11]),
            req_start_date=_to_date(mst[12]), c_datetime=_to_date(mst[13]),
            category=cat, sub_category=sub_cat, site_country=site_c,
            site_region=site_r, site=site, division=div, requisition_type=req_t,
            parent_supplier=par_s, supplier_type=sup_t, supplier_country=sup_c,
            payment_days=pay_d, po_date=po_d, category_buyer=cat_b,
            l3=l3, l4=l4, purchase_category=pur_cat, ras_title=ras_t,
            line_items=line_items,
        )
    finally:
        conn.close()


# ── Document loader for extraction ─────────────────────────────────────────────

def _load_document(file_bytes: bytes, filename: str, max_pages: int = 20) -> DocumentContent:
    import os
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pdf":                                  return _load_pdf_for_extract(file_bytes, max_pages)
    if ext in (".xlsx", ".xls", ".xlsm", ".xlsb"):    return _load_spreadsheet_for_extract(file_bytes, ext)
    if ext == ".docx":                                 return _load_docx_for_extract(file_bytes)
    if ext == ".doc":                                  return _load_doc_legacy_for_extract(file_bytes, max_pages)
    if ext in (".pptx", ".ppt"):                       return _load_pptx_for_extract(file_bytes, max_pages)
    if ext in _IMAGE_EXTS:                             return _load_image_for_extract(file_bytes, ext)
    if ext in (".txt", ".csv", ".xml", ".json", ".eml"):
        return DocumentContent(text=file_bytes.decode("utf-8", errors="replace")[:50000], page_count=1)
    if ext in (".html", ".htm"):                       return _load_html_for_extract(file_bytes)
    if ext == ".msg":                                  return _load_msg_for_extract(file_bytes)
    # Fallback: fitz render (handles most office formats via PyMuPDF)
    return _fitz_render_for_extract(file_bytes, max_pages)


def _load_pdf_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    """Hybrid PDF loader (matches doc-intel branch non-OCR strategy).

    Digital page (≥50 extracted chars): captured as text + 72 DPI image for layout context.
    Scanned page (<50 chars): captured as 150 DPI image only.
    Returns both text and images so the extraction LLM has full context without Azure Doc Intel.
    """
    import base64
    _SCANNED_THRESHOLD = 50
    try:
        import fitz
        doc        = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            page_count = len(doc)
            if page_count > max_pages:
                # Match doc-intel: warn explicitly so the user can see truncation in the canvas log.
                logger.warning(
                    "PDF has {} page(s) but max_pages={} — only first {} page(s) will be processed",
                    page_count, max_pages, max_pages,
                )
            total      = min(page_count, max_pages)
            text_pages: list[str] = []
            images:     list[str] = []
            text_count = scanned_count = 0

            for i in range(total):
                page      = doc[i]
                page_text = page.get_text("text").strip()
                if len(page_text) >= _SCANNED_THRESHOLD:
                    # Digital page — extract text; low-res image preserves layout context
                    text_pages.append(f"[Page {i + 1}]\n{page_text}")
                    text_count   += 1
                    pix = page.get_pixmap(dpi=72)
                else:
                    # Scanned/image page — high-res render is the only usable content
                    scanned_count += 1
                    pix = page.get_pixmap(dpi=150)
                images.append(base64.b64encode(pix.tobytes("png")).decode())
                pix = None
        finally:
            doc.close()
        logger.info(
            "PDF loaded: {}/{} page(s) — {} digital (72 DPI text+image), {} scanned (150 DPI image)",
            total, page_count, text_count, scanned_count,
        )
        text = "\n\n".join(text_pages) if text_pages else None
        return DocumentContent(text=text, images=images, page_count=total)
    except Exception as exc:
        logger.opt(exception=True).error("PDF load failed: {}", exc)
        return DocumentContent(text=f"[PDF error: {exc}]", page_count=0)


def _load_spreadsheet_for_extract(file_bytes: bytes, ext: str) -> DocumentContent:
    """Load Excel workbook as a markdown table string.

    .xlsx / .xlsm — openpyxl (handles macro-enabled workbooks too)
    .xls           — xlrd (legacy binary format)
    .xlsb          — openpyxl attempt; falls back to fitz image render
    """
    import io
    parts: list[str] = []

    def _openpyxl_parts(raw: bytes) -> list[str]:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        result: list[str] = []
        for ws in wb.worksheets:
            lines = [f"### Sheet: {ws.title}"]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    lines.append("| " + " | ".join(cells) + " |")
            if len(lines) > 1:
                result.append("\n".join(lines))
        wb.close()
        return result

    try:
        if ext in (".xlsx", ".xlsm"):
            parts = _openpyxl_parts(file_bytes)
        elif ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(file_contents=file_bytes)
            for sheet in wb.sheets():
                lines = [f"### Sheet: {sheet.name}"]
                for rx in range(sheet.nrows):
                    cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
                    if any(c.strip() for c in cells):
                        lines.append("| " + " | ".join(cells) + " |")
                if len(lines) > 1:
                    parts.append("\n".join(lines))
        else:  # .xlsb — binary Excel; openpyxl may handle it, else render as images
            try:
                parts = _openpyxl_parts(file_bytes)
            except Exception:
                return _fitz_render_for_extract(file_bytes, 20)

        text = "\n\n".join(parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, 20)
    except Exception as exc:
        return DocumentContent(text=f"[Spreadsheet error: {exc}]", page_count=0)


def _load_docx_for_extract(file_bytes: bytes) -> DocumentContent:
    import io
    try:
        from docx import Document
        doc   = Document(io.BytesIO(file_bytes))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        table_parts: list[str] = []
        for tbl in doc.tables:
            rows = ["| " + " | ".join(c.text.strip() for c in row.cells) + " |"
                    for row in tbl.rows]
            if rows:
                table_parts.append("\n".join(rows))
        text = "\n".join(paras)
        if table_parts:
            text += "\n\n" + "\n\n".join(table_parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, 20)
    except Exception as exc:
        return DocumentContent(text=f"[DOCX error: {exc}]", page_count=0)


def _load_doc_legacy_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    try:
        return _fitz_render_for_extract(file_bytes, max_pages)
    except Exception:
        pass
    try:
        import io, olefile
        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        if ole.exists("WordDocument"):
            import re as _re
            stream = ole.openstream("WordDocument").read()
            text = stream.decode("utf-8", errors="ignore")
            text = _re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
            if text.strip():
                return DocumentContent(text=text, page_count=1)
        ole.close()
    except Exception:
        pass
    return DocumentContent(text="[Document could not be read]")


def _load_pptx_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    import io
    try:
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(file_bytes))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_t: list[str] = [f"--- Slide {idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_t.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        slide_t.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
            if len(slide_t) > 1:
                parts.append("\n".join(slide_t))
        text = "\n\n".join(parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, max_pages)
    except Exception:
        return _fitz_render_for_extract(file_bytes, max_pages)


def _load_image_for_extract(file_bytes: bytes, ext: str) -> DocumentContent:
    import io, base64
    raw = file_bytes
    if ext in (".tif", ".tiff"):
        try:
            from PIL import Image
            buf = io.BytesIO()
            with Image.open(io.BytesIO(raw)) as _img:
                _img.save(buf, format="PNG")
            raw = buf.getvalue()
        except Exception:
            pass
    return DocumentContent(images=[base64.b64encode(raw).decode()], page_count=1)


def _load_msg_for_extract(file_bytes: bytes) -> DocumentContent:
    import os, tempfile
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject: parts.append(f"Subject: {msg.subject}")
            if msg.body:    parts.append(msg.body)
            msg.close()
        finally:
            os.unlink(tmp_path)
        return DocumentContent(text="\n\n".join(parts), page_count=1)
    except Exception as exc:
        return DocumentContent(text=f"[MSG error: {exc}]")


def _load_html_for_extract(file_bytes: bytes) -> DocumentContent:
    """Strip HTML tags and return plain text (BeautifulSoup → regex fallback)."""
    import re
    html_str = file_bytes.decode("utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_str, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
    except Exception:
        text = re.sub(r"<(script|style)[^>]*>.*?</(script|style)>", " ", html_str, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
    return DocumentContent(text=text[:50000], page_count=1)


def _fitz_render_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    import base64
    try:
        import fitz
        doc        = fitz.open(stream=file_bytes)
        try:
            page_count = len(doc)
            if page_count > max_pages:
                logger.warning(
                    "Document has {} page(s) but max_pages={} — only first {} page(s) will be rendered",
                    page_count, max_pages, max_pages,
                )
            n = min(page_count, max_pages)
            images = []
            for i in range(n):
                pix = doc[i].get_pixmap(dpi=200)
                images.append(base64.b64encode(pix.tobytes("png")).decode())
                pix = None
        finally:
            doc.close()
        logger.info("Document rendered: {}/{} page(s) at 200 DPI", n, page_count)
        return DocumentContent(images=images, page_count=n)
    except Exception as exc:
        logger.opt(exception=True).error("Document render failed: {}", exc)
        return DocumentContent(text=f"[Render error: {exc}]")


# ── Quotation source resolver ──────────────────────────────────────────────────

def _resolve_quotation_sources(tgt_cs: str, pr_no: str) -> list:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    sources: list[dict] = []
    try:
        cur.execute("""
            SELECT ac.[attachment_classify_uuid_pk], ac.[file_path], ac.[attachment_id]
              FROM [ras_procurement].[attachment_classification] ac
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND ac.[doc_type] = 'Quotation'
        """, pr_no)
        for row in cur.fetchall():
            if row[1]:
                sources.append({
                    "blob_path": row[1],
                    "attachment_classify_fk": str(row[0]),
                    "embedded_classify_fk": None,
                    "attachment_id": str(row[2]),
                })

        cur.execute("""
            SELECT ec.[embedded_attachment_classification_id], ec.[file_path],
                   ec.[parent_attachment_id], ac.[attachment_classify_uuid_pk]
              FROM [ras_procurement].[embedded_attachment_classification] ec
              JOIN [ras_procurement].[attachment_classification] ac
                ON ec.[attachment_classification_id] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND ec.[doc_type] = 'Quotation'
        """, pr_no)
        for row in cur.fetchall():
            if row[1]:
                sources.append({
                    "blob_path": row[1],
                    "attachment_classify_fk": str(row[3]),
                    "embedded_classify_fk": str(row[0]),
                    "attachment_id": str(row[2]),
                })
    finally:
        conn.close()
    return sources


# ── Line items table builder ───────────────────────────────────────────────────

def _build_line_items_table(ctx: RASContext) -> str:
    def _na(v): return str(v) if v is not None else "N/A"
    header = ("| DTL_ID | Item No | Item Code | Description | Qty | UOM | Type "
              "| Unit Price | Req Value | Currency | Supplier | Delivery Date |")
    sep    = ("|--------|---------|-----------|-------------|-----|-----|------"
              "|------------|-----------|----------|----------|---------------|")
    rows   = [header, sep]
    for li in ctx.line_items:
        rows.append(
            f"| {_na(li.purchase_dtl_id)} | {_na(li.item_no)} | {_na(li.item_code)} "
            f"| {_na(li.item_description)} | {_na(li.quantity)} | {_na(li.uom)} "
            f"| {_na(li.item_type)} | {_na(li.unit_price)} | {_na(li.req_value)} "
            f"| {_na(li.currency)} | {_na(li.supplier_name)} | {_na(li.delivery_date)} |"
        )
    return "\n".join(rows)


def _build_raw_context(ctx: RASContext) -> str:
    parts: list[str] = []
    if ctx.raw_mst:
        lines = [f"  {k}: {v}" for k, v in ctx.raw_mst.items()]
        parts.append("#### purchase_req_mst (header)\n" + "\n".join(lines))
    if ctx.raw_dtl_rows:
        cols   = list(ctx.raw_dtl_rows[0].keys())
        header = "| " + " | ".join(cols) + " |"
        sep    = "|" + "|".join("---" for _ in cols) + "|"
        rows   = ["| " + " | ".join(str(r.get(c, "")) for c in cols) + " |"
                  for r in ctx.raw_dtl_rows]
        parts.append("#### purchase_req_detail\n" + "\n".join([header, sep] + rows))
    return "\n\n".join(parts)


# ── Extraction LLM call ────────────────────────────────────────────────────────

def _build_extraction_user_prompt(ctx: RASContext, doc: DocumentContent, prompts: dict | None = None) -> str:
    def _f(v): return str(v) if v is not None else "N/A"
    if doc.ocr_source and doc.text:
        doc_content_str = f"[OCR markdown from Azure Document Intelligence]\n\n{doc.text}"
    elif doc.images and doc.text:
        doc_content_str = f"[Extracted text — page images attached below]\n\n{doc.text}"
    elif doc.images:
        doc_content_str = "[Scanned document — page image(s) attached]"
    else:
        doc_content_str = doc.text or "[No content extracted]"

    ext_max = (prompts or {}).get("ext_max_chars")
    if ext_max and doc_content_str and len(doc_content_str) > ext_max:
        doc_content_str = doc_content_str[:ext_max]

    user_tmpl = (prompts or {}).get("ext_user", EXTRACTION_USER_TEMPLATE)
    return user_tmpl.format(
        purchase_req_no=_f(ctx.purchase_req_no),
        purchase_req_id=_f(ctx.purchase_req_id),
        justification=_f(ctx.justification),
        supplier_name=_f(ctx.supplier_name),
        currency=_f(ctx.currency),
        enquiry_no=_f(ctx.enquiry_no),
        classification=_f(ctx.classification),
        department=_f(ctx.department),
        negotiated_by=_f(ctx.negotiated_by),
        address=_f(ctx.address),
        contract_no=_f(ctx.contract_no),
        order_no=_f(ctx.order_no),
        purchase_value=_f(ctx.purchase_value),
        category=_f(ctx.category),
        sub_category=_f(ctx.sub_category),
        l3=_f(ctx.l3),
        l4=_f(ctx.l4),
        purchase_category=_f(ctx.purchase_category),
        ras_title=_f(ctx.ras_title),
        site_region=_f(ctx.site_region),
        site_country=_f(ctx.site_country),
        site=_f(ctx.site),
        division=_f(ctx.division),
        requisition_type=_f(ctx.requisition_type),
        parent_supplier=_f(ctx.parent_supplier),
        supplier_type=_f(ctx.supplier_type),
        supplier_country=_f(ctx.supplier_country),
        payment_days=_f(ctx.payment_days),
        po_date=_f(ctx.po_date),
        category_buyer=_f(ctx.category_buyer),
        line_items_table=_build_line_items_table(ctx),
        item_taxonomy=(prompts or {}).get("ext_taxonomy", ITEM_TAXONOMY),
        document_content=doc_content_str,
        raw_ras_context=_build_raw_context(ctx),
    )


def _call_extraction_llm(llm, ctx: RASContext, doc: DocumentContent, prompts: dict | None = None) -> str:
    from langchain_core.messages import HumanMessage, SystemMessage
    user_prompt = _build_extraction_user_prompt(ctx, doc, prompts)
    sys_prompt  = (prompts or {}).get("ext_system", EXTRACTION_SYSTEM_PROMPT)
    # Hard cap on images attached per LLM request — Azure OpenAI rejects
    # requests with more than ~50 images. Mirrors doc-intel _MAX_IMAGES.
    max_images = max(1, int((prompts or {}).get("ext_max_images", 50)))
    messages: list = [SystemMessage(content=sys_prompt)]
    img_count  = 0
    img_detail = "n/a"
    if doc.is_image_based and doc.images:
        total_imgs = len(doc.images)
        if total_imgs > max_images:
            logger.warning(
                "[PR={}] Quotation has {} image(s) but max_images={} — truncating to first {}",
                ctx.purchase_req_no, total_imgs, max_images, max_images,
            )
        images = doc.images[:max_images]
        img_count  = len(images)
        # Mirror doc-intel: low detail when text is also present, high otherwise.
        img_detail = "low" if doc.text else "high"
        content_parts: list = [{"type": "text", "text": user_prompt}]
        for b64 in images:
            content_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}", "detail": img_detail},
            })
        messages.append(HumanMessage(content=content_parts))
    else:
        messages.append(HumanMessage(content=user_prompt))
    logger.info(
        "[PR={}] Calling extraction LLM — {} image(s) detail={}, ~{:.0f} kB prompt",
        ctx.purchase_req_no, img_count, img_detail, len(user_prompt) / 1024,
    )
    # force_json=True so the extraction LLM emits a strict JSON object and
    # our _parse_extraction_response doesn't have to clean up prose/fences.
    response = _call_llm_with_retry(llm, messages, prompts, force_json=True)
    return (getattr(response, "content", None) or str(response)).strip()


# ── Supplier matching ──────────────────────────────────────────────────────────

_SELECTED_THRESHOLD     = Decimal("0.70")  # min overall score to mark as selected quote
_PRICE_MAX_BOOST        = Decimal("0.10")  # max confidence boost from price alignment
_CANONICALIZE_THRESHOLD = 0.82             # SequenceMatcher ratio to cluster two supplier names

def _compute_supplier_match(supplier: Optional[str], ctx: RASContext, thresholds: dict | None = None) -> tuple:
    if not supplier:
        return False, Decimal("0")
    known: set[str] = set()
    if ctx.supplier_name:   known.add(ctx.supplier_name.strip())
    if ctx.parent_supplier: known.add(ctx.parent_supplier.strip())
    for li in ctx.line_items:
        if li.supplier_name: known.add(li.supplier_name.strip())
    if not known:
        return False, Decimal("0")
    ext  = supplier.strip().lower()
    best = 0.0
    for n in known:
        nl = n.lower()
        if ext in nl or nl in ext:
            best = max(best, 0.90)
        else:
            best = max(best, SequenceMatcher(None, ext, nl).ratio())
    conf = Decimal(str(round(best, 4)))
    threshold = (thresholds or {}).get("selected_threshold", _SELECTED_THRESHOLD)
    return conf >= threshold, conf


# ── Extraction response parsing ────────────────────────────────────────────────

def _parse_extraction_response(raw: str, source: dict, ctx: RASContext) -> list:
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        logger.error(f"LLM returned invalid JSON: {exc}")
        return []

    header: dict = data if isinstance(data, dict) else {}
    items_raw: list = header.get("items", [])
    if not items_raw:
        return []

    h_supplier = header.get("supplier_name") or None
    _, match_conf = _compute_supplier_match(h_supplier, ctx)

    header_fields = {
        "supplier_name":      h_supplier,
        "supplier_address":   header.get("supplier_address") or None,
        # Store the LLM-extracted country verbatim (matches doc-intel
        # _parse_llm_response). Previously this was force-normalised to ISO
        # alpha-2 via pycountry which lost the readable country name in the
        # DB (e.g. "India" → "IN"). Downstream consumers (_compute_cpi_pct,
        # _estimate_inflation_via_llm) already handle both full names and
        # alpha-2/3 codes, so no functional dependency was lost.
        "supplier_country":   (header.get("supplier_country") or "").strip() or None,
        "quotation_ref_no":   header.get("quotation_ref_no") or None,
        "quotation_date":     header.get("quotation_date"),
        "currency":           _normalize_currency_code(header.get("currency") or None),
        "validity_date":      header.get("validity_date"),
        "validity_days":      header.get("validity_days"),
        "payment_terms":      header.get("payment_terms") or None,
        "supplier_match_conf": float(match_conf),
        "attachment_classify_fk": source["attachment_classify_fk"],
        "embedded_classify_fk":   source["embedded_classify_fk"],
        "is_selected_quote":      0,
        "quote_rank":             None,
    }

    results: list[dict] = []
    for raw_item in items_raw:
        item = dict(header_fields)
        item.update({k: (None if v == "" else v) for k, v in raw_item.items()})
        results.append(item)
    return results


# ── Item alignment to RAS line items ──────────────────────────────────────────

_IDENT_RE = re.compile(r"[A-Za-z0-9]*\d+[A-Za-z0-9]*")

def _ident_tokens(text: Optional[str]) -> set:
    if not text: return set()
    return {m.group(0).lower() for m in _IDENT_RE.finditer(text) if len(m.group(0)) >= 2}

def _data_score(item: dict) -> tuple:
    return (int(item.get("unit_price") is not None), int(item.get("total_price") is not None),
            int(item.get("item_name") is not None), int(item.get("item_description") is not None))

def _align_to_ras_line_items(items: list, ctx: RASContext, source: dict) -> list:
    valid_ids = {li.purchase_dtl_id for li in ctx.line_items}
    matched_by_dtl: dict = {}
    orphans: list[dict] = []

    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None:
            try: dtl_id = int(dtl_id)
            except Exception: dtl_id = None
        if dtl_id in valid_ids:
            if dtl_id not in matched_by_dtl:
                matched_by_dtl[dtl_id] = item
            elif _data_score(item) > _data_score(matched_by_dtl[dtl_id]):
                matched_by_dtl[dtl_id] = item
        elif any([item.get("item_name"), item.get("unit_price"), item.get("total_price")]):
            orphans.append(item)

    matched = list(matched_by_dtl.values())
    covered = {i.get("purchase_dtl_id") for i in matched}
    uncovered = [li for li in ctx.line_items if li.purchase_dtl_id not in covered]

    # Fuzzy assign orphans
    if orphans and uncovered:
        scored: list = []
        for oi, orp in enumerate(orphans):
            orp_text = f"{orp.get('item_name','')} {orp.get('item_description','')}".lower()
            for li_idx, li in enumerate(uncovered):
                ras_text = f"{li.item_description or ''} {li.item_code or ''}".lower()
                shared = _ident_tokens(orp_text) & _ident_tokens(ras_text)
                if shared:
                    score = 0.90
                elif li.item_code and li.item_code.lower() in orp_text:
                    score = 0.85
                else:
                    score = SequenceMatcher(None, orp_text, ras_text).ratio()
                try:
                    if orp.get("quantity") is not None and li.quantity is not None:
                        if Decimal(str(orp["quantity"])) == li.quantity:
                            score = min(1.0, score + 0.10)
                except Exception:
                    pass
                scored.append((score, oi, li_idx))
        scored.sort(key=lambda x: x[0], reverse=True)
        used_o: set = set(); used_l: set = set()
        for score, oi, li_idx in scored:
            if oi in used_o or li_idx in used_l: continue
            dtl_id = uncovered[li_idx].purchase_dtl_id
            orphans[oi]["purchase_dtl_id"] = dtl_id
            matched.append(orphans[oi])
            covered.add(dtl_id)
            used_o.add(oi); used_l.add(li_idx)

    # Stubs for uncovered
    still_uncovered = {li.purchase_dtl_id for li in uncovered if li.purchase_dtl_id not in covered}
    donor = matched[0] if matched else {}
    for li in ctx.line_items:
        if li.purchase_dtl_id not in still_uncovered:
            continue
        matched.append({
            "attachment_classify_fk": source["attachment_classify_fk"],
            "embedded_classify_fk":   source["embedded_classify_fk"],
            "purchase_dtl_id":        li.purchase_dtl_id,
            "is_selected_quote":      0,
            "supplier_match_conf":    0.0,
            "quote_rank":             None,
            "supplier_name":          donor.get("supplier_name"),
            "supplier_address":       donor.get("supplier_address"),
            "supplier_country":       donor.get("supplier_country"),
            "quotation_ref_no":       donor.get("quotation_ref_no"),
            "quotation_date":         donor.get("quotation_date"),
            "currency":               donor.get("currency"),
            "validity_date":          donor.get("validity_date"),
            "validity_days":          donor.get("validity_days"),
            "payment_terms":          donor.get("payment_terms"),
        })
    return matched


# ── Quote ranking + selection ──────────────────────────────────────────────────

def _compute_quote_ranks(all_items: list) -> None:
    from collections import defaultdict
    for item in all_items:
        item["quote_rank"] = None
    by_group: dict = defaultdict(list)
    for item in all_items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is None: continue
        key = ((item.get("supplier_name") or "").strip().lower() or "_unknown_", dtl_id)
        by_group[key].append(item)
    for group in by_group.values():
        group.sort(key=lambda i: (i.get("total_price") is None, float(i.get("total_price") or 0)))
        for rank, item in enumerate(group, 1):
            item["quote_rank"] = rank


def _select_best_quotes(all_items: list, ctx: RASContext) -> None:
    from collections import defaultdict
    for item in all_items:
        item["is_selected_quote"] = 0
    ras_by_dtl = {li.purchase_dtl_id: li for li in ctx.line_items}
    by_dtl: dict = defaultdict(list)
    for item in all_items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None:
            by_dtl[dtl_id].append(item)

    def _price_prox(item_val, ras_val):
        try:
            if item_val is None or ras_val is None or float(ras_val) <= 0: return 0.0
            diff = abs(float(item_val) - float(ras_val)) / float(ras_val)
            if diff <= 0.05: return 1.0
            if diff >= 0.25: return 0.0
            return 1.0 - (diff - 0.05) / 0.20
        except Exception:
            return 0.0

    for dtl_id, candidates in by_dtl.items():
        ras_line = ras_by_dtl.get(dtl_id)
        def _score(it):
            conf      = float(it.get("supplier_match_conf") or 0)
            price_fit = max(
                _price_prox(it.get("unit_price"),  ras_line.unit_price if ras_line else None),
                _price_prox(it.get("total_price"), ras_line.req_value  if ras_line else None),
            )
            has_price = int(it.get("unit_price") is not None)
            return (conf, price_fit, has_price)
        max(candidates, key=_score)["is_selected_quote"] = 1


# ── Currency conversion helper ─────────────────────────────────────────────────

_EUR_CUR_ID      = 3
_RATE_IS_MULTIPLY = True


def _convert_to_eur(tgt_cs: str, amount, currency_code: str | None, ref_date) -> "Decimal | None":
    """Return amount converted to EUR using EXCHANGE_RATE table.

    Returns None if amount is None or conversion data is unavailable.
    """
    if amount is None:
        return None
    try:
        from decimal import Decimal as _Dec
        amount_dec = _Dec(str(amount))
    except Exception:
        return None
    if not currency_code:
        return None
    try:
        import re as _re
        from datetime import date as _date_cls, datetime as _dt_cls
        if isinstance(ref_date, _date_cls):
            date_val = ref_date
        elif ref_date:
            m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(ref_date))
            if m:
                date_val = _date_cls(int(m[1]), int(m[2]), int(m[3]))
            else:
                date_val = _dt_cls.utcnow().date()
        else:
            from datetime import datetime as _dt_cls2
            date_val = _dt_cls2.utcnow().date()
        conn = _connect(tgt_cs)
        try:
            cur = conn.cursor()
            # look up source currency id
            cur.execute(
                "SELECT [CUR_ID] FROM [ras_procurement].[currency_mst] "
                "WHERE UPPER([CURRENCY]) = UPPER(?)",
                currency_code,
            )
            row = cur.fetchone()
            if row is None:
                logger.warning(f"Currency conversion: no currency_mst row for code={currency_code!r}")
                return None
            src_cur_id = row[0]
            if src_cur_id == _EUR_CUR_ID:
                return amount_dec
            # look up exchange rate valid for ref_date
            cur.execute(
                "SELECT TOP 1 [CONVERSION_RATE] "
                "FROM [ras_procurement].[EXCHANGE_RATE] "
                "WHERE [CUR_ID] = ? AND [BASE_CUR_ID] = ? "
                "  AND [STATUS_ID] = 10 "
                "  AND [FROM_DATE] <= ? AND [TO_DATE] >= ? "
                "ORDER BY [FROM_DATE] DESC",
                src_cur_id, _EUR_CUR_ID, date_val, date_val,
            )
            rate_row = cur.fetchone()
            if rate_row is None:
                logger.warning(
                    f"Currency conversion: no EXCHANGE_RATE row for "
                    f"FROM_CUR_ID={src_cur_id} TO_CUR_ID={_EUR_CUR_ID} date={date_val} STATUS_ID=10"
                )
                return None
            from decimal import Decimal as _Dec2
            rate = _Dec2(str(rate_row[0]))
            return amount_dec * rate if _RATE_IS_MULTIPLY else amount_dec / rate
        finally:
            conn.close()
    except Exception as exc:
        logger.warning(f"Currency conversion failed currency={currency_code!r} date={ref_date}: {exc}")
        return None


# ── Supplier name canonicalization (mirrors doc intel Union-Find logic) ────────

def _normalize_currency_code(code_or_name: str | None) -> str | None:
    """Normalize any currency string to ISO-4217 alpha-3 code using pycountry.

    Handles codes (USD), full names (US Dollar), and partial names via rapidfuzz.
    No hardcoded map — works for any currency in the ISO-4217 standard.
    Gracefully degrades to original value when pycountry is not installed.
    """
    if not code_or_name:
        return code_or_name
    val = code_or_name.strip()
    if not val:
        return None
    upper = val.upper()
    try:
        import pycountry as _pc
        # 1. Exact alpha-3 code match (most common — LLM usually returns ISO code)
        c = _pc.currencies.get(alpha_3=upper)
        if c:
            return c.alpha_3
        # 2. Exact name match (e.g. "US Dollar", "Euro")
        val_lower = val.lower()
        for c in _pc.currencies:
            if c.name.lower() == val_lower:
                return c.alpha_3
        # 3. Fuzzy name match via rapidfuzz (catches partial names, typos)
        try:
            from rapidfuzz import process as _fuzz
            all_names = [c.name for c in _pc.currencies]
            result = _fuzz.extractOne(val_lower, [n.lower() for n in all_names], score_cutoff=85)
            if result:
                idx = [n.lower() for n in all_names].index(result[0])
                return list(_pc.currencies)[idx].alpha_3
        except ImportError:
            pass
        # 4. Looks like a 3-char code but not in pycountry — return uppercased
        if len(upper) == 3 and upper.isalpha():
            logger.warning(f"Currency code {upper!r} not found in ISO-4217 — passing through as-is")
            return upper
    except ImportError:
        logger.warning("pycountry not installed — currency code not normalized (pip install pycountry)")
        if len(upper) == 3 and upper.isalpha():
            return upper
    return val


def _normalize_supplier_country(country_str: str | None) -> str | None:
    """Normalize free-text country name to ISO 3166-1 alpha-2 code using pycountry.

    Tries exact name/code match first, then rapidfuzz fuzzy match (threshold 75).
    Falls back to original string when pycountry is not installed or no match found.
    """
    if not country_str:
        return country_str
    text_lower = country_str.strip().lower()
    try:
        import pycountry as _pc
        for c in _pc.countries:
            if text_lower in (
                c.name.lower(),
                getattr(c, "official_name", "").lower(),
                c.alpha_2.lower(),
                c.alpha_3.lower(),
            ):
                return c.alpha_2
        try:
            from rapidfuzz import process as _fuzz
            all_names = [c.name for c in _pc.countries]
            result = _fuzz.extractOne(text_lower, [n.lower() for n in all_names])
            if result and result[1] > 75:
                matched = _pc.countries.get(name=all_names[[n.lower() for n in all_names].index(result[0])])
                if matched:
                    return matched.alpha_2
        except ImportError:
            pass
    except ImportError:
        logger.warning("pycountry not installed — supplier_country not normalized (pip install pycountry rapidfuzz)")
    return country_str


def _strip_contact_suffix(name: str) -> str:
    import re as _re
    return _re.sub(
        r'\s*[-–]\s*(contact|email|ph|phone|tel|mob)[:\s].*$',
        '', name, flags=_re.IGNORECASE,
    ).strip()


def _is_acronym_of(short: str, long_name: str) -> bool:
    """True if short is all-caps (≤ 6 chars) and matches the word initials of long_name."""
    if not short.isupper() or len(short) > 6:
        return False
    import re as _re
    initials = "".join(m[0].upper() for m in _re.findall(r'\b[A-Za-z]', long_name))
    return short == initials[:len(short)]


def _name_geo_tokens(cleaned: str) -> frozenset:
    import re as _re
    _GEO = frozenset({
        "india", "china", "japan", "usa", "us", "uk", "germany", "france",
        "italy", "korea", "taiwan", "singapore", "malaysia", "thailand",
        "vietnam", "indonesia", "australia", "canada", "brazil", "mexico",
        "uae", "dubai", "europe", "asia", "americas", "shanghai", "beijing",
        "mumbai", "delhi",
    })
    return frozenset(w for w in _re.findall(r'\b\w+\b', cleaned.lower()) if w in _GEO)


def _apply_price_alignment_boost(items: list[dict], ctx: RASContext, thresholds: dict | None = None) -> None:
    """Boost supplier_match_conf by up to price_max_boost when extracted prices align with RAS prices.

    If ≥50% of matched line items have a unit_price within _PRICE_TOLERANCE (5%) of the
    RAS line price, all items in this PR get a proportional boost. Matches doc-intel branch
    _apply_price_alignment_boost logic, adapted for dict items.
    """
    dtl_price: dict = {
        li.purchase_dtl_id: li.unit_price
        for li in ctx.line_items
        if li.unit_price is not None and li.unit_price > 0
    }
    if not dtl_price:
        return
    matches = comparable = 0
    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id in dtl_price:
            unit_price = item.get("unit_price")
            if unit_price is not None:
                comparable += 1
                try:
                    ras_p = dtl_price[dtl_id]
                    if abs(Decimal(str(unit_price)) - ras_p) / ras_p <= _PRICE_TOLERANCE:
                        matches += 1
                except Exception:
                    pass
    if comparable == 0 or matches / comparable < 0.5:
        return
    max_boost = (thresholds or {}).get("price_max_boost", _PRICE_MAX_BOOST)
    boost = Decimal(str(round((matches / comparable) * float(max_boost), 4)))
    for item in items:
        conf = item.get("supplier_match_conf")
        if conf is not None:
            try:
                item["supplier_match_conf"] = min(
                    Decimal("1.0"), Decimal(str(conf)) + boost
                )
            except Exception:
                pass
    logger.debug(f"Price alignment boost +{boost} applied ({matches}/{comparable} items within {int(_PRICE_TOLERANCE*100)}%)")


def _canonicalize_supplier_names(items: list[dict], ctx, thresholds: dict | None = None) -> None:
    """PR-wide Union-Find supplier name canonicalization with country guard.

    Clusters supplier_name variants across ALL DTLs in the current PR (a
    supplier that quotes on multiple line items can have slightly different
    name spellings on each quote — e.g. "Motherson Automotive Technologies
    & Engineering (MATE - ROBIS)" vs "Motherson Automotive Technologies &
    Engineering (A Division of Motherson Sumi Systems Ltd.)"). Five merge
    rules — applied in this order so cheaper checks short-circuit:
      1. Exact (after contact-suffix strip + lowercase).
      2. Substring.
      3. Paren-stripped exact ("X (Branch A)" + "X (Branch B)" → both
         strip to "X" and merge — handles the Motherson case where
         qualifying parentheticals differ but the legal name is identical).
      4. SequenceMatcher ratio >= _CANONICALIZE_THRESHOLD (0.82 default).
      5. Acronym ("MATE" matches "Motherson Automotive Tech & Eng").

    Country guard — two names do NOT merge when either:
      - both have non-empty supplier_country and they differ, OR
      - both have a geo-token (city/country) in the name itself and
        those tokens differ.
    So "ENGEL Machinery (Shanghai) Co., Ltd." and "ENGEL Machinery Korea
    LTD" stay separate even though every other rule would match them.

    Canonical preference within a cluster: RAS-known supplier_name first,
    else the shortest member. After clustering, supplier_match_conf is
    recomputed for ALL items against the RAS supplier set so ranking uses
    canonical names (this also wipes the earlier _apply_price_alignment_boost,
    same as doc-intel).
    """
    from difflib import SequenceMatcher

    # Build RAS-known name set for canonical preference (matches doc-intel)
    ras_known: set[str] = set()
    if getattr(ctx, "supplier_name", None):
        ras_known.add(ctx.supplier_name.strip().lower())
    if getattr(ctx, "parent_supplier", None):
        ras_known.add(ctx.parent_supplier.strip().lower())
    for li in getattr(ctx, "line_items", []):
        if getattr(li, "supplier_name", None):
            ras_known.add(li.supplier_name.strip().lower())

    threshold = (thresholds or {}).get("canonicalize_threshold", _CANONICALIZE_THRESHOLD)

    # Collect unique supplier names across the WHOLE PR (cross-DTL) along
    # with the first non-empty country seen for each name.
    raw_names: list[str] = []
    name_country: dict[str, str] = {}
    for it in items:
        n = (it.get("supplier_name") or "").strip()
        if not n:
            continue
        if n not in name_country:
            raw_names.append(n)
            name_country[n] = (it.get("supplier_country") or "").strip().lower()

    if len(raw_names) < 2:
        # Still recompute supplier_match_conf for consistency with prior
        # behaviour (fresh scores against RAS suppliers).
        for item in items:
            _, new_conf = _compute_supplier_match(item.get("supplier_name"), ctx, thresholds)
            item["supplier_match_conf"] = new_conf
        return

    uf: dict[str, str] = {n: n for n in raw_names}

    def _find(x: str) -> str:
        while uf[x] != x:
            uf[x] = uf[uf[x]]
            x = uf[x]
        return x

    def _union(a: str, b: str) -> None:
        uf[_find(a)] = _find(b)

    def _strip_parens(s: str) -> str:
        # Remove parenthetical qualifiers ("(MATE - ROBIS)", "(A Division
        # of Motherson Sumi Systems Ltd.)") so the legal core name can be
        # compared directly. Geo info that appears in parens is still
        # protected by the geo-token guard applied above.
        return re.sub(r"\([^)]*\)", "", s).strip()

    for i, a in enumerate(raw_names):
        a_clean   = _strip_contact_suffix(a).lower()
        a_country = name_country.get(a, "")
        a_geo     = _name_geo_tokens(a_clean)
        a_paren   = _strip_parens(a_clean)
        for b in raw_names[i + 1:]:
            b_clean   = _strip_contact_suffix(b).lower()
            b_country = name_country.get(b, "")
            b_geo     = _name_geo_tokens(b_clean)
            b_paren   = _strip_parens(b_clean)

            # Country-branch guard — keeps "ACME Corp India" and
            # "ACME Corp China" separate even when the LLM forgot to
            # extract supplier_country for one of them.
            if a_country and b_country and a_country != b_country:
                continue
            if a_geo and b_geo and a_geo != b_geo:
                continue

            if a_clean == b_clean:
                _union(a, b)
            elif a_clean in b_clean or b_clean in a_clean:
                _union(a, b)
            elif a_paren and b_paren and a_paren == b_paren:
                _union(a, b)
            elif SequenceMatcher(None, a_clean, b_clean).ratio() >= threshold:
                _union(a, b)
            elif _is_acronym_of(a.strip(), b) or _is_acronym_of(b.strip(), a):
                _union(a, b)

    from collections import defaultdict
    clusters: dict[str, list[str]] = defaultdict(list)
    for n in raw_names:
        clusters[_find(n)].append(n)

    canonical_names: dict[str, str] = {}
    for members in clusters.values():
        ras_match = next((m for m in members if m.lower() in ras_known), None)
        canonical = ras_match if ras_match else min(members, key=len)
        for m in members:
            canonical_names[m] = canonical

    # Apply canonical names + recompute supplier_match_conf for ALL items.
    changed = 0
    for item in items:
        orig = (item.get("supplier_name") or "").strip()
        if orig:
            canon = canonical_names.get(orig, orig)
            if canon != orig:
                item["supplier_name"] = canon
                changed += 1
        _, new_conf = _compute_supplier_match(item.get("supplier_name"), ctx, thresholds)
        item["supplier_match_conf"] = new_conf

    if changed:
        logger.info(f"Supplier canonicalization (PR-wide): {changed} name(s) updated")


# ── DB writer: extracted items ─────────────────────────────────────────────────

def _save_extracted_items(tgt_cs: str, items: list, fallback_date=None) -> int:
    if not items:
        return 0
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    saved = 0
    try:
        for item in items:
            _COL_LIMITS = {
                "item_description": 4000, "item_summary": 4000,
                "taxation_details": 2000, "payment_terms": 2000,
                "supplier_address": 500,  "item_name": 500,
                "supplier_name": 255,     "supplier_country": 100,
                "quotation_ref_no": 100,  "commodity_tag": 255,
                "unit": 50,
            }
            def _v(k, cast=None):
                v = item.get(k)
                if v is None: return None
                try:
                    v = cast(v) if cast else v
                except Exception:
                    return None
                if isinstance(v, str) and k in _COL_LIMITS:
                    limit = _COL_LIMITS[k]
                    if len(v) > limit:
                        logger.warning("Truncating '{}': {} → {} chars", k, len(v), limit)
                        v = v[:limit]
                return v
            def _d(k):
                v = item.get(k)
                if v is None: return None
                try: return Decimal(str(v))
                except Exception: return None
            def _date(k):
                v = item.get(k)
                if not v: return None
                try:
                    from datetime import date as date_cls
                    if isinstance(v, date_cls): return v
                    import re as _re
                    m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(v))
                    if m: return date_cls(int(m[1]), int(m[2]), int(m[3]))
                except Exception: pass
                return None
            unit_price_eur  = _convert_to_eur(tgt_cs, item.get("unit_price"),  item.get("currency"), fallback_date)
            total_price_eur = _convert_to_eur(tgt_cs, item.get("total_price"), item.get("currency"), fallback_date)
            cur.execute("""
                INSERT INTO [ras_procurement].[quotation_extracted_items] (
                    [attachment_classify_fk],[embedded_classify_fk],[purchase_dtl_id],
                    [is_selected_quote],[supplier_match_conf],[quote_rank],
                    [supplier_name],[supplier_address],[supplier_country],
                    [quotation_ref_no],[quotation_date],[currency],
                    [validity_date],[validity_days],[payment_terms],
                    [item_name],[item_description],[quantity],[unit],
                    [unit_price],[total_price],[discount],[taxation_details],
                    [delivery_date],[delivery_time_days],
                    [item_level_1],[item_level_2],[item_level_3],[item_level_4],
                    [item_level_5],[item_level_6],[item_level_7],[item_level_8],
                    [commodity_tag],[item_summary],
                    [unit_price_eur],[total_price_eur]
                ) VALUES (
                    ?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?
                )
            """,
                _v("attachment_classify_fk"), _v("embedded_classify_fk"),
                _v("purchase_dtl_id", int), int(item.get("is_selected_quote") or 0),
                _d("supplier_match_conf"), _v("quote_rank", int),
                _v("supplier_name"), _v("supplier_address"), _v("supplier_country"),
                _v("quotation_ref_no"), _date("quotation_date"), _v("currency"),
                _date("validity_date"), _v("validity_days", int), _v("payment_terms"),
                _v("item_name"), _v("item_description"), _d("quantity"), _v("unit"),
                _d("unit_price"), _d("total_price"), _d("discount"), _v("taxation_details"),
                _date("delivery_date"), _v("delivery_time_days", int),
                _v("item_level_1"), _v("item_level_2"), _v("item_level_3"), _v("item_level_4"),
                _v("item_level_5"), _v("item_level_6"), _v("item_level_7"), _v("item_level_8"),
                _v("commodity_tag"), _v("item_summary"),
                unit_price_eur, total_price_eur,
            )
            saved += 1
        conn.commit()
    except Exception as exc:
        logger.opt(exception=True).error(
            "DB INSERT into quotation_extracted_items failed after {} row(s): {}",
            saved, exc,
        )
        try:
            conn.rollback()
        except Exception as rb_exc:
            logger.warning("Rollback failed: {}", rb_exc)
        raise
    finally:
        conn.close()
    return saved


# ── Stage 5: run extraction for a PR ──────────────────────────────────────────

def _run_extraction(llm, tgt_cs: str, blob_cfg: dict, pr_no: str, prompts: dict | None = None) -> int:
    import os
    # Each abort path raises an ExtractionAbortError subclass so the per-PR
    # catch in _run_stages_48 records ras_pipeline_exceptions (stage_id=5),
    # sets ras_tracker.current_stage_fk = 99, and skips stages 6/7/8.
    # Returning 0 silently here would let the pipeline advance to stage 8
    # (COMPLETE) on a PR that had nothing to extract — which is wrong.
    ctx = _build_ras_context(tgt_cs, pr_no)
    if ctx is None:
        raise NoRASContextError(
            f"No purchase_req_mst row for PR={pr_no!r} — cannot extract."
        )
    if not ctx.line_items:
        raise NoLineItemsError(
            f"PR={pr_no!r} has no rows in purchase_req_detail — nothing to extract against."
        )

    sources = _resolve_quotation_sources(tgt_cs, pr_no)
    if not sources:
        raise NoQuotationFoundError(
            f"No quotation documents found for PR={pr_no!r}. "
            f"None of the attachments are classified as 'Quotation'."
        )

    # Each source is: download blob → load doc → LLM extract → parse → align.
    # This is I/O + LLM bound, so threads give real speedup when a PR has
    # multiple quotation files. Keep n_workers low (default 1) to avoid
    # multiplying LLM rate-limit pressure on top of parallel_workers.
    n_workers = max(1, int((prompts or {}).get("ext_parallel_sources", 1)))
    max_pages = max(1, int((prompts or {}).get("ext_max_pages", 20)))

    def _extract_one(src: dict) -> list[dict]:
        blob_path = src["blob_path"]
        filename  = os.path.basename(blob_path)
        try:
            file_bytes = _download_blob(blob_path, blob_cfg)
            doc        = _load_document(file_bytes, filename, max_pages=max_pages)
            raw        = _call_extraction_llm(llm, ctx, doc, prompts)
            items      = _parse_extraction_response(raw, src, ctx)
            items      = _align_to_ras_line_items(items, ctx, src)
            logger.info(
                "[{}] Extracted {} item(s) from {!r} ({} page(s) loaded, {} image(s))",
                pr_no, len(items), filename, doc.page_count,
                len(doc.images) if doc.images else 0,
            )
            return items
        except Exception as exc:
            logger.opt(exception=True).warning(
                "[{}] Extraction failed for {!r}: {}", pr_no, filename, exc
            )
            return []

    if n_workers == 1:
        all_items: list[dict] = []
        for src in sources:
            all_items.extend(_extract_one(src))
    else:
        import concurrent.futures
        logger.info(f"[{pr_no}] Extracting {len(sources)} source(s) with {n_workers} parallel worker(s)")
        with concurrent.futures.ThreadPoolExecutor(max_workers=n_workers) as pool:
            results = list(pool.map(_extract_one, sources))
        all_items = [item for batch in results for item in batch]

    if not all_items:
        # Every quotation source raised inside _extract_one (LLM error, parse
        # failure, blob 404, etc.) and was caught individually. With nothing to
        # save, abort so this PR is moved to exception state instead of
        # advancing through stages 6-8 with empty data.
        raise AllExtractionsFailedError(
            f"All {len(sources)} quotation source(s) for PR={pr_no!r} "
            f"failed extraction; no items produced."
        )

    _apply_price_alignment_boost(all_items, ctx, prompts)
    _canonicalize_supplier_names(all_items, ctx, prompts)
    _compute_quote_ranks(all_items)
    _select_best_quotes(all_items, ctx)
    fallback_date = ctx.req_start_date or ctx.c_datetime
    saved = _save_extracted_items(tgt_cs, all_items, fallback_date=fallback_date)
    logger.info(f"[{pr_no}] {saved} item(s) written to quotation_extracted_items")
    return saved


# ── Tracker helpers ────────────────────────────────────────────────────────────

def _advance_tracker(tgt_cs: str, pr_no: str, stage_id: int) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[ras_tracker]
               SET current_stage_fk=?, updated_at=SYSUTCDATETIME()
             WHERE purchase_req_no=?
        """, stage_id, pr_no)
        conn.commit()
    except Exception as exc:
        logger.warning(f"Tracker advance failed PR={pr_no!r} stage={stage_id}: {exc}")
    finally:
        conn.close()


def _set_last_processed_at(tgt_cs: str, pr_no: str) -> None:
    """Stamp last_processed_at = SYSUTCDATETIME() so SourceChangeDetector
    can detect future on-prem source changes for this PR."""
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[ras_tracker]
               SET last_processed_at = SYSUTCDATETIME(),
                   updated_at        = SYSUTCDATETIME()
             WHERE purchase_req_no   = ?
        """, pr_no)
        conn.commit()
    except Exception as exc:
        logger.warning(f"set_last_processed_at failed PR={pr_no!r}: {exc}")
    finally:
        conn.close()


def _record_exception(tgt_cs: str, pr_no: str, stage_id: int, error_msg: str) -> None:
    try:
        conn = _connect(tgt_cs)
        cur  = conn.cursor()
        try:
            # Also stamp last_processed_at so audits / reporting can see when
            # this PR last reached a terminal state (success OR exception).
            cur.execute("""
                UPDATE [ras_procurement].[ras_tracker]
                   SET current_stage_fk   = 99,
                       last_processed_at  = SYSUTCDATETIME(),
                       updated_at         = SYSUTCDATETIME()
                 WHERE purchase_req_no    = ?
            """, pr_no)
            cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no=?", pr_no)
            row = cur.fetchone()
            if row:
                cur.execute("""
                    INSERT INTO [ras_procurement].[ras_pipeline_exceptions]
                        (ras_tracker_id, stage_id, exception_message)
                    VALUES (?, ?, ?)
                """, row[0], stage_id, error_msg[:4000])
            conn.commit()
        finally:
            conn.close()
    except Exception as exc:
        logger.warning(f"[{pr_no}] Could not write exception record: {exc}")


# ── Embeddings + benchmark (Stage 6-7) ────────────────────────────────────────

def _build_embed_text(row_dict: dict) -> str:
    """Build the 12-field embedding text (identical order used by doc intel branch)."""
    _FIELDS = [
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag",
    ]
    return " | ".join(str(row_dict[f]) for f in _FIELDS if row_dict.get(f))


def _run_embeddings(tgt_cs: str, pr_no: str, embed_model, pinecone_index: str, pinecone_ns: str) -> None:
    # Structural errors (import, DB connect, index creation) propagate — caller records exception.
    from agentcore.services.pinecone_service_client import ensure_index_via_service, ingest_via_service
    _pinecone_ensure_with_retry(
        ensure_index_via_service,
        index_name=pinecone_index,
        embedding_dimension=3072,
    )
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8], qi.[commodity_tag],
                   prd.[C_DATETIME] AS [item_created_date]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        all_rows = cur.fetchall()
    finally:
        conn.close()

    # deduplicate by purchase_dtl_id — first row wins (same as doc intel branch)
    seen_dtl: dict = {}
    for row in all_rows:
        dtl_id = row[1]
        if dtl_id not in seen_dtl:
            seen_dtl[dtl_id] = row
    rows = list(seen_dtl.values())
    logger.info(f"[{pr_no}] Embedding {len(rows)} selected item(s) (deduped from {len(all_rows)} is_selected_quote=1 rows)")

    cols = [
        "extracted_item_uuid_pk", "purchase_dtl_id",
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag", "item_created_date",
    ]
    for row in rows:
        rd = dict(zip(cols, row))
        dtl_id    = rd["purchase_dtl_id"]
        item_uuid = rd["extracted_item_uuid_pk"]
        created   = rd.get("item_created_date")
        created_iso = (
            created.isoformat() if created and hasattr(created, "isoformat") else str(created or "")
        )
        content = _build_embed_text(rd)
        if not content:
            continue
        try:
            embedding = _embed_with_timeout(embed_model, content)
            _pinecone_ingest_with_timeout(
                ingest_via_service,
                index_name=pinecone_index,
                namespace=pinecone_ns,
                text_key="page_content",
                documents=[{
                    "page_content": content,
                    "metadata": {
                        "purchase_req_no":        pr_no,
                        "purchase_dtl_id":        int(dtl_id),
                        "extracted_item_uuid_pk": str(item_uuid or ""),
                        "commodity_tag":          str(rd.get("commodity_tag") or ""),
                        "item_created_date":      created_iso,
                    },
                }],
                embedding_vectors=[embedding],
                vector_ids=[f"dtl_{dtl_id}"],
                embedding_dimension=3072,
            )
            logger.info(f"[{pr_no}] Upserted vector dtl_{dtl_id} (item_created={created_iso})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Embedding failed for dtl_id {dtl_id}: {exc}")


def _fetch_historical_for_dtl_ids(tgt_cs: str, dtl_ids: list) -> list[dict]:
    """Fetch all extraction columns from quotation_extracted_items for the given dtl_ids
    (is_selected_quote = 1 only). Returns the full row so the pricing LLM has complete
    item context — name, description, category hierarchy, specs, pricing, and terms."""
    if not dtl_ids:
        return []
    placeholders = ", ".join(["?"] * len(dtl_ids))
    sql = f"""
        SELECT qi.[purchase_dtl_id], qi.[extracted_item_uuid_pk],
               qi.[item_name], qi.[item_description], qi.[item_summary],
               qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
               qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
               qi.[item_level_7], qi.[item_level_8],
               qi.[commodity_tag], qi.[purchase_category_llm], qi.[critical_attributes],
               qi.[supplier_name], qi.[supplier_address], qi.[supplier_country],
               qi.[quotation_ref_no], qi.[quotation_date], qi.[currency],
               qi.[validity_date], qi.[validity_days], qi.[payment_terms],
               qi.[quantity], qi.[unit],
               qi.[unit_price], qi.[total_price],
               qi.[unit_price_eur], qi.[total_price_eur],
               qi.[discount], qi.[taxation_details],
               qi.[delivery_date], qi.[delivery_time_days],
               qi.[supplier_match_conf], qi.[quote_rank],
               rt.[purchase_req_no],
               prd.[C_DATETIME] AS pr_created_date
          FROM [ras_procurement].[quotation_extracted_items] qi
          JOIN [ras_procurement].[attachment_classification] ac
            ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
          JOIN [ras_procurement].[ras_tracker] rt
            ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
          LEFT JOIN [ras_procurement].[purchase_req_detail] prd
            ON prd.[PURCHASE_DTL_ID] = qi.[purchase_dtl_id]
         WHERE qi.[purchase_dtl_id] IN ({placeholders})
           AND qi.[is_selected_quote] = 1
    """
    cols = [
        "purchase_dtl_id", "extracted_item_uuid_pk",
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag", "purchase_category_llm", "critical_attributes",
        "supplier_name", "supplier_address", "supplier_country",
        "quotation_ref_no", "quotation_date", "currency",
        "validity_date", "validity_days", "payment_terms",
        "quantity", "unit",
        "unit_price", "total_price",
        "unit_price_eur", "total_price_eur",
        "discount", "taxation_details",
        "delivery_date", "delivery_time_days",
        "supplier_match_conf", "quote_rank",
        "purchase_req_no", "pr_created_date",
    ]
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute(sql, *dtl_ids)
        return [dict(zip(cols, r)) for r in cur.fetchall()]
    finally:
        conn.close()


def _fetch_benchmark_for_dtl_ids(tgt_cs: str, dtl_ids: list) -> list[dict]:
    """Fetch existing benchmark_result rows for similar dtl_ids to use as LLM context."""
    if not dtl_ids:
        return []
    placeholders = ", ".join(["?"] * len(dtl_ids))
    sql = f"""
        SELECT br.[purchase_dtl_id],
               br.[bp_unit_price], br.[bp_total_price],
               br.[inflation_pct], br.[cpi_inflation_pct],
               br.[summary],
               prm.[PURCHASE_REQ_NO]
          FROM [ras_procurement].[benchmark_result] br
          JOIN [ras_procurement].[purchase_req_detail] prd
            ON prd.[PURCHASE_DTL_ID] = br.[purchase_dtl_id]
          JOIN [ras_procurement].[purchase_req_mst] prm
            ON prm.[PURCHASE_REQ_ID] = prd.[PURCHASE_REQ_ID]
         WHERE br.[purchase_dtl_id] IN ({placeholders})
           AND br.[bp_unit_price] IS NOT NULL
    """
    cols = [
        "purchase_dtl_id", "bp_unit_price", "bp_total_price",
        "inflation_pct", "cpi_inflation_pct", "summary", "purchase_req_no",
    ]
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute(sql, *dtl_ids)
        return [dict(zip(cols, r)) for r in cur.fetchall()]
    finally:
        conn.close()


def _compute_low_last(items: list[dict]) -> tuple:
    """Return (low_item, last_item) — cheapest EUR price and most recently created PR dtl.

    Guarantees a non-None result whenever `items` is non-empty: if no priced rows exist
    the cheapest item falls back to the first UUID-valid row (or first row); same for dated.
    This ensures low_hist_item_fk / last_hist_item_fk are always populated when the
    shortlist has data.
    """
    def _eur(it: dict):
        return it.get("unit_price_eur") or it.get("unit_price")

    # Prefer items that have a valid UUID so the FK write doesn't end up NULL.
    uuid_items = [it for it in items if it.get("extracted_item_uuid_pk") is not None]
    pool = uuid_items if uuid_items else items

    priced = [it for it in pool if _eur(it) is not None]
    dated  = [it for it in pool if it.get("pr_created_date") is not None]

    low_item  = min(priced, key=_eur) if priced else (pool[0] if pool else None)
    last_item = max(dated,  key=lambda it: it["pr_created_date"]) if dated else (pool[0] if pool else None)
    return low_item, last_item


def _hist_eur(it: dict) -> "Decimal | None":
    """Best-effort EUR unit price for a historical row (fallback to unit_price)."""
    v = it.get("unit_price_eur")
    if v is None:
        v = it.get("unit_price")
    try:
        return Decimal(str(v)) if v is not None else None
    except Exception:
        return None


def _filter_historical_uom(historical: list[dict], current_uom: str | None, strict: bool) -> list[dict]:
    """Drop historical items whose UOM does not match the current item.

    Why: a laptop priced 'per piece' cannot be benchmarked against one priced
    'per box of 10'; a steel coil priced 'MT' cannot be benchmarked against
    one priced 'KG'. Mixing UOMs silently corrupts the recommendation.
    Comparison is case-insensitive and trims whitespace; if either side has
    no UOM the row is kept (cannot prove they differ).
    """
    if not strict or not current_uom:
        return historical
    norm = (current_uom or "").strip().lower()
    if not norm:
        return historical
    kept: list[dict] = []
    dropped = 0
    for it in historical:
        h = (it.get("unit") or "").strip().lower()
        if not h or h == norm:
            kept.append(it)
        else:
            dropped += 1
    if dropped:
        logger.info("Benchmark UOM filter: dropped %d historical row(s) with UOM != %r", dropped, current_uom)
    return kept


def _filter_historical_outliers(historical: list[dict], factor: float) -> list[dict]:
    """Drop rows whose unit_price_eur is more than `factor`× away from the median.

    Why: a single typo'd quote (e.g. 1000 EUR instead of 100 EUR) skews the
    LLM's recommendation. We compute the median across the kept set and drop
    any row whose price is > median × factor or < median / factor. Set
    factor <= 0 to disable outlier removal entirely.
    """
    if factor is None or factor <= 0 or len(historical) < 3:
        return historical
    eurs = [_hist_eur(it) for it in historical]
    eurs = [float(e) for e in eurs if e is not None and e > 0]
    if len(eurs) < 3:
        return historical
    eurs_sorted = sorted(eurs)
    n = len(eurs_sorted)
    median = (eurs_sorted[n // 2] if n % 2 else (eurs_sorted[n // 2 - 1] + eurs_sorted[n // 2]) / 2)
    if median <= 0:
        return historical
    high = median * factor
    low  = median / factor
    kept: list[dict] = []
    dropped = 0
    for it in historical:
        e = _hist_eur(it)
        if e is None:
            kept.append(it)        # no price → can't be an outlier; keep
            continue
        ev = float(e)
        if low <= ev <= high:
            kept.append(it)
        else:
            dropped += 1
    if dropped:
        logger.info("Benchmark outlier filter: dropped %d row(s) outside [%.2f, %.2f] (median %.2f, factor %s)",
                    dropped, low, high, median, factor)
    return kept


def _filter_historical_age(historical: list[dict], current_iso: str, max_age_months: int) -> list[dict]:
    """Drop rows older than max_age_months relative to the current item's
    item_created_date. Set max_age_months <= 0 to disable."""
    if max_age_months is None or max_age_months <= 0:
        return historical
    if not current_iso:
        return historical
    import re as _re
    from datetime import date as _date
    m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(current_iso))
    if not m:
        return historical
    cur_date = _date(int(m[1]), int(m[2]), int(m[3]))
    cur_months = cur_date.year * 12 + cur_date.month
    kept: list[dict] = []
    dropped = 0
    for it in historical:
        qd = it.get("quotation_date")
        if qd is None:
            kept.append(it)
            continue
        try:
            qy = qd.year if hasattr(qd, "year") else int(str(qd)[:4])
            qm = qd.month if hasattr(qd, "month") else int(str(qd)[5:7])
        except Exception:
            kept.append(it)
            continue
        age_months = cur_months - (qy * 12 + qm)
        if age_months <= max_age_months:
            kept.append(it)
        else:
            dropped += 1
    if dropped:
        logger.info("Benchmark age filter: dropped %d row(s) older than %d month(s)", dropped, max_age_months)
    return kept


def _benchmark_aggregates(historical: list[dict], current_iso: str | None) -> dict:
    """Pre-compute aggregate stats for the LLM benchmark prompt.

    Returns a dict with:
      median, p25, p75, mean, min, max, count       (across unit_price_eur)
      time_buckets : { "≤6m": [...], "6-24m": [...], ">24m": [...] }
                     each value is the list of historical dicts in that bucket
      by_supplier  : { supplier_name: { count, min, max, median } }
      qty_min, qty_max, qty_median (current vs historical scale context)

    All numeric stats are floats, rounded to 4 dp.
    """
    if not historical:
        return {}
    import re as _re
    from datetime import date as _date

    eurs = [(_hist_eur(it), it) for it in historical]
    priced = [(float(e), it) for e, it in eurs if e is not None and e > 0]

    stats: dict = {"count": len(historical), "priced_count": len(priced)}
    if priced:
        vals = sorted(p for p, _ in priced)
        n = len(vals)
        def _quantile(arr: list[float], q: float) -> float:
            if not arr:
                return 0.0
            k = (len(arr) - 1) * q
            f, c = int(k), min(int(k) + 1, len(arr) - 1)
            if f == c:
                return arr[f]
            return arr[f] + (arr[c] - arr[f]) * (k - f)
        stats["min"]    = round(vals[0], 4)
        stats["max"]    = round(vals[-1], 4)
        stats["median"] = round(_quantile(vals, 0.50), 4)
        stats["p25"]    = round(_quantile(vals, 0.25), 4)
        stats["p75"]    = round(_quantile(vals, 0.75), 4)
        stats["mean"]   = round(sum(vals) / n, 4)

    # Time buckets relative to current_iso (or today if missing)
    cur_date: _date | None = None
    if current_iso:
        m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(current_iso))
        if m:
            cur_date = _date(int(m[1]), int(m[2]), int(m[3]))
    if cur_date is None:
        cur_date = _date.today()
    cur_months = cur_date.year * 12 + cur_date.month
    buckets: dict[str, list] = {"recent": [], "mid": [], "old": []}
    for it in historical:
        qd = it.get("quotation_date")
        if qd is None:
            buckets["old"].append(it)
            continue
        try:
            qy = qd.year if hasattr(qd, "year") else int(str(qd)[:4])
            qm = qd.month if hasattr(qd, "month") else int(str(qd)[5:7])
        except Exception:
            buckets["old"].append(it)
            continue
        age_months = cur_months - (qy * 12 + qm)
        if age_months <= 6:
            buckets["recent"].append(it)
        elif age_months <= 24:
            buckets["mid"].append(it)
        else:
            buckets["old"].append(it)
    stats["time_buckets"] = buckets

    # Per-supplier summary (top suppliers by count)
    by_sup: dict[str, list[float]] = {}
    for p, it in priced:
        s = (it.get("supplier_name") or "Unknown").strip() or "Unknown"
        by_sup.setdefault(s, []).append(p)
    sup_summary: dict[str, dict] = {}
    for s, prices in by_sup.items():
        prices.sort()
        n = len(prices)
        med = prices[n // 2] if n % 2 else (prices[n // 2 - 1] + prices[n // 2]) / 2
        sup_summary[s] = {
            "count":  n,
            "min":    round(prices[0], 4),
            "max":    round(prices[-1], 4),
            "median": round(med, 4),
        }
    stats["by_supplier"] = sup_summary

    # Quantity context
    qtys = [float(it.get("quantity")) for it in historical if it.get("quantity") is not None]
    if qtys:
        qtys.sort()
        n = len(qtys)
        med_q = qtys[n // 2] if n % 2 else (qtys[n // 2 - 1] + qtys[n // 2]) / 2
        stats["qty_min"]    = qtys[0]
        stats["qty_max"]    = qtys[-1]
        stats["qty_median"] = med_q

    return stats


def _format_bench_prompt(
    current: dict,
    historical: list[dict],
    stats: dict | None = None,
    prior_benchmarks: list[dict] | None = None,
) -> str:
    """Build a richer benchmark LLM prompt.

    Includes pre-computed aggregates (median, P25, P75, mean), per-supplier
    mini-summary, and three time-bucketed views of the historical data
    (recent ≤6m / mid 6-24m / old >24m). Asks for a price RANGE
    (bp_low / bp_unit_price / bp_high) plus a confidence score so the
    downstream system can communicate uncertainty.

    Backward-compatible JSON: callers only need bp_unit_price + summary,
    extra fields are folded into summary if not stored separately.
    """
    eur_unit  = current.get("unit_price_eur")
    raw_price = current.get("unit_price")
    currency  = current.get("currency") or ""
    if eur_unit is not None and currency.upper() != "EUR":
        price_str = f"{eur_unit} EUR  (original: {raw_price} {currency})"
    elif eur_unit is not None:
        price_str = f"{eur_unit} EUR"
    else:
        price_str = f"{raw_price} {currency}"

    category = " > ".join(
        str(current[f]) for f in [
            "item_level_1", "item_level_2", "item_level_3", "item_level_4",
            "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        ] if current.get(f)
    )
    current_block = "\n".join([
        f"- Item name        : {current.get('item_name') or 'N/A'}",
        f"- Description      : {current.get('item_description') or 'N/A'}",
        f"- Category         : {category or 'N/A'}",
        f"- Commodity tag    : {current.get('commodity_tag') or 'N/A'}",
        f"- Quantity         : {current.get('quantity')} {current.get('unit') or ''}",
        f"- Quoted unit price: {price_str}",
        f"- Supplier         : {current.get('supplier_name') or 'N/A'}",
        f"- Quotation date   : {current.get('quotation_date') or 'N/A'}",
    ])

    if not historical:
        return (
            "Recommend a benchmark unit price for this item. "
            "No historical purchase data is available for comparison — use your market knowledge.\n\n"
            f"CURRENT ITEM:\n{current_block}\n\n"
            "HISTORICAL DATA: none available — no comparable older purchases found in the system.\n\n"
            "Since no historical data exists, set confidence low (≤0.30) and state clearly in the "
            "summary that the estimate is based on market knowledge only.\n\n"
            "Return ONLY JSON (no markdown):\n"
            "{\n"
            '  "bp_unit_price": <number EUR or null>,\n'
            '  "bp_low":        <number EUR or null>,\n'
            '  "bp_high":       <number EUR or null>,\n'
            '  "confidence":    <0.0-0.30 — low, no historical data>,\n'
            '  "summary":       "<3-5 sentences: (1) note that no similar historical purchases were found, '
            "(2) basis for the market estimate (item specs, known market rates), "
            "(3) key assumptions made, "
            '(4) negotiation recommendation>"\n'
            "}"
        )

    stats = stats or {}

    # ── Aggregate stats block ────────────────────────────────────────────
    agg_lines: list[str] = []
    if stats.get("count") is not None:
        agg_lines.append(
            f"- Sample size : {stats['count']} historical row(s) "
            f"({stats.get('priced_count', 0)} priced)"
        )
    if "median" in stats:
        agg_lines.append(
            f"- Price stats (EUR) : "
            f"min={stats['min']}, P25={stats['p25']}, "
            f"median={stats['median']}, P75={stats['p75']}, "
            f"max={stats['max']}, mean={stats['mean']}"
        )
    if "qty_median" in stats:
        cur_qty = current.get("quantity") or "?"
        agg_lines.append(
            f"- Quantity context : current={cur_qty} {current.get('unit') or ''}, "
            f"historical qty (min / median / max) = "
            f"{stats.get('qty_min')} / {stats['qty_median']} / {stats.get('qty_max')}"
        )
    agg_block = "\n".join(agg_lines) if agg_lines else "- No aggregate stats available"

    # ── Per-supplier summary block ───────────────────────────────────────
    by_sup = stats.get("by_supplier") or {}
    if by_sup:
        # Sort by count desc, take top 8 to keep prompt size sane
        top_sups = sorted(by_sup.items(), key=lambda kv: -kv[1]["count"])[:8]
        sup_rows = ["| Supplier | N | Min EUR | Median EUR | Max EUR |",
                    "|----------|---|---------|------------|---------|"]
        for s, d in top_sups:
            sup_rows.append(
                f"| {s} | {d['count']} | {d['min']} | {d['median']} | {d['max']} |"
            )
        sup_block = "\n".join(sup_rows)
        if len(by_sup) > 8:
            sup_block += f"\n  (… and {len(by_sup) - 8} more supplier(s) with smaller samples)"
    else:
        sup_block = "No per-supplier breakdown available."

    # ── Time-bucketed historical rows ─────────────────────────────────────
    buckets = stats.get("time_buckets") or {}

    def _hist_block(rows: list[dict]) -> str:
        if not rows:
            return "  (none in this window)"
        out = []
        for i, it in enumerate(rows, 1):
            eur_p = it.get("unit_price_eur")
            raw_p = it.get("unit_price")
            price_str = (
                f"{eur_p} EUR (orig: {raw_p} {it.get('currency') or ''})"
                if eur_p is not None and str(it.get("currency") or "").upper() != "EUR"
                else (f"{eur_p} EUR" if eur_p is not None else f"{raw_p} {it.get('currency') or ''}")
            )
            cat_parts = [
                it.get(f"item_level_{n}") for n in range(1, 9)
                if it.get(f"item_level_{n}")
            ]
            cat = " > ".join(str(p) for p in cat_parts) if cat_parts else "N/A"
            crit = it.get("critical_attributes")
            if isinstance(crit, str):
                try: crit = json.loads(crit)
                except Exception: crit = []
            block = [
                f"[{i}] purchase_dtl_id={it.get('purchase_dtl_id')} | PR={it.get('purchase_req_no') or 'N/A'} | date={it.get('pr_created_date') or 'N/A'}",
                f"  Item name        : {it.get('item_name') or 'N/A'}",
                f"  Description      : {it.get('item_description') or 'N/A'}",
                f"  Summary          : {it.get('item_summary') or 'N/A'}",
                f"  Category         : {cat}",
                f"  Commodity tag    : {it.get('commodity_tag') or 'N/A'}",
                f"  LLM category     : {it.get('purchase_category_llm') or 'N/A'}",
                f"  Critical attrs   : {json.dumps(crit, ensure_ascii=False) if crit else 'N/A'}",
                f"  Supplier         : {it.get('supplier_name') or 'N/A'} | country={it.get('supplier_country') or 'N/A'}",
                f"  Quotation date   : {it.get('quotation_date') or 'N/A'} | ref={it.get('quotation_ref_no') or 'N/A'}",
                f"  Quantity / Unit  : {it.get('quantity') or 'N/A'} {it.get('unit') or ''}",
                f"  Unit price       : {price_str}",
                f"  Total price      : {it.get('total_price_eur') or it.get('total_price') or 'N/A'} EUR",
                f"  Discount         : {it.get('discount') or 'N/A'}",
                f"  Payment terms    : {it.get('payment_terms') or 'N/A'}",
                f"  Delivery         : {it.get('delivery_date') or 'N/A'} ({it.get('delivery_time_days') or 'N/A'} days)",
                f"  Validity         : {it.get('validity_date') or 'N/A'} ({it.get('validity_days') or 'N/A'} days)",
                f"  Taxation details : {it.get('taxation_details') or 'N/A'}",
            ]
            out.append("\n".join(block))
        return "\n\n".join(out)

    hist_block = (
        "RECENT (last 6 months) — weight these highest:\n"
        + _hist_block(buckets.get("recent") or []) + "\n\n"
        "MID (6-24 months) — weight moderately:\n"
        + _hist_block(buckets.get("mid") or []) + "\n\n"
        "OLDER (>24 months) — apply inflation adjustment if relevant:\n"
        + _hist_block(buckets.get("old") or [])
    )

    # ── Prior benchmark results block ────────────────────────────────────
    if prior_benchmarks:
        pb_rows = [
            "| PR No | Recommended (EUR) | Inflation % | CPI % | Summary |",
            "|-------|-------------------|-------------|-------|---------|",
        ]
        for pb in prior_benchmarks:
            pb_rows.append(
                f"| {pb.get('purchase_req_no') or 'N/A'} "
                f"| {pb.get('bp_unit_price') or 'N/A'} "
                f"| {pb.get('inflation_pct') or 'N/A'} "
                f"| {pb.get('cpi_inflation_pct') or 'N/A'} "
                f"| {(pb.get('summary') or '')[:120]} |"
            )
        prior_block = "\n".join(pb_rows)
        prior_section = (
            f"PRIOR BENCHMARK RESULTS FOR SIMILAR ITEMS (reference only):\n{prior_block}\n"
            "Use these as calibration context. Do NOT copy these values directly —\n"
            "they may reflect different specs, quantities, or market conditions.\n\n"
        )
    else:
        prior_section = ""

    return (
        "You are recommending a benchmark unit price for the CURRENT ITEM "
        "below, drawing on historical purchases of similar items. The "
        "historical set has already been pre-filtered for similarity, "
        "UOM compatibility, age, and price outliers — every row shown is "
        "a valid comparator.\n\n"
        "PRICING RULES — follow in order:\n"
        "  1. HISTORICAL ANCHOR: When the historical rows are genuinely comparable "
        "(same spec, same UOM, similar quantity scale), your bp_unit_price MUST NOT "
        "exceed the median historical price. The benchmark is a target/fair price — "
        "not a market estimate. If the quoted price is already below the median, "
        "anchor the benchmark at or below the quoted price.\n"
        "  2. INFLATION ADJUSTMENT: You may adjust upward from the oldest comparable "
        "price only to account for documented inflation — and only if recent data is "
        "absent. Never inflate above the most recent comparable price.\n"
        "  3. INCOMPARABLE DATA: If the historical rows are NOT comparable (different "
        "spec class, wrong UOM, order-of-magnitude size difference), state this "
        "clearly in the summary and provide a market-knowledge estimate with low "
        "confidence. Do NOT silently use incomparable prices as anchors.\n"
        "  4. PRICE CEILING: bp_high must never exceed the quoted price of the "
        "current item unless the historical data clearly shows the market is higher.\n\n"
        "How to weight the data:\n"
        "  • Recent (≤6m) prices reflect today's market — anchor on these.\n"
        "  • Mid (6-24m) prices are useful trend signals.\n"
        "  • Older (>24m) prices may need inflation adjustment.\n"
        "  • Per-supplier ranges show negotiation room with each vendor.\n"
        "  • Quantity context flags volume-vs-pilot pricing differences.\n"
        "  • The median and P25-P75 band are usually the safest anchor.\n\n"
        f"CURRENT ITEM:\n{current_block}\n\n"
        f"AGGREGATE STATS:\n{agg_block}\n\n"
        f"PER-SUPPLIER SUMMARY:\n{sup_block}\n\n"
        f"HISTORICAL DATA (time-bucketed):\n{hist_block}\n\n"
        f"{prior_section}"
        "Return ONLY JSON (no markdown):\n"
        "{\n"
        '  "bp_unit_price": <number EUR — your point estimate, must not exceed median when data is comparable>,\n'
        '  "bp_low":        <number EUR — floor of plausible fair range>,\n'
        '  "bp_high":       <number EUR — ceiling, must not exceed quoted price unless market evidence says otherwise>,\n'
        '  "confidence":    <0.0-1.0 — how reliable this benchmark is given the sample>,\n'
        '  "summary":       "<3-5 sentences: (1) which DTL IDs and suppliers were used as benchmark, '
        "(2) key spec differences vs current item if any, "
        "(3) how the price was derived from historical data, "
        '(4) negotiation recommendation>"\n'
        "}"
    )


def _compute_cpi_pct(country_str: str | None, start_year: int, end_year: int):
    """Cumulative CPI inflation % via World Bank API (FP.CPI.TOTL.ZG).

    Mirrors utils/cpi_inflation.py from the doc intel branch.
    Returns float or None on any failure.
    """
    if not country_str or start_year >= end_year:
        return None
    try:
        import httpx as _httpx
        import pycountry as _pc
        from rapidfuzz import process as _fuzz
        # Resolve country name → ISO alpha-2
        text = country_str.strip().lower()
        alpha2 = None
        for c in _pc.countries:
            if text in (c.name.lower(), getattr(c, "official_name", "").lower(),
                        c.alpha_2.lower(), c.alpha_3.lower()):
                alpha2 = c.alpha_2
                break
        if alpha2 is None:
            names = [c.name for c in _pc.countries]
            res = _fuzz.extractOne(text, names)
            if res and res[1] > 75:
                matched = _pc.countries.get(name=res[0])
                if matched:
                    alpha2 = matched.alpha_2
        if alpha2 is None:
            logger.info(f"CPI: could not resolve country {country_str!r}")
            return None
        # Fetch World Bank annual rates
        url = f"https://api.worldbank.org/v2/country/{alpha2}/indicator/FP.CPI.TOTL.ZG"
        resp = _httpx.get(url, params={"date": f"{start_year}:{end_year}", "format": "json", "per_page": 100}, timeout=30.0)
        resp.raise_for_status()
        data = resp.json()
        if len(data) < 2 or not data[1]:
            return None
        rates = {int(e["date"]): float(e["value"]) for e in data[1] if e.get("value") is not None}
        if not rates:
            return None
        factor = 1.0
        for yr in range(start_year + 1, end_year + 1):
            if yr in rates:
                factor *= (1 + rates[yr] / 100)
        return round((factor - 1) * 100, 4)
    except Exception as exc:
        logger.info(f"CPI: failed for {country_str!r} {start_year}-{end_year}: {exc}")
        return None


def _estimate_inflation_via_llm(llm, item_name: str | None, category: str | None,
                                 supplier_country: str | None, ref_year: int | None,
                                 current_year: int | None) -> float | None:
    """Ask LLM for estimated cumulative inflation % for an item category in a country.

    Returns float or None. Does NOT derive inflation from historical price delta.
    """
    if not supplier_country or not ref_year or not current_year or current_year <= ref_year:
        return None
    try:
        from langchain_core.messages import HumanMessage as _HM
        import json as _json, re as _re
        years = current_year - ref_year
        prompt = (
            f"I am buying the following item from {supplier_country}:\n"
            f"  Item     : {item_name or 'N/A'}\n"
            f"  Category : {category or 'N/A'}\n\n"
            f"What is the estimated cumulative inflation rate (%) for this type of item "
            f"in {supplier_country} from {ref_year} to {current_year} ({years} year(s))?\n\n"
            f"Base your estimate on macroeconomic inflation trends for this category in that country. "
            f"Do NOT derive the rate from any historical price data.\n"
            f'Respond ONLY with JSON: {{ "inflation_pct": <number or null> }}'
        )
        resp = _call_llm_with_retry(llm, [_HM(content=prompt)], prompts=None, force_json=True)
        raw  = (getattr(resp, "content", None) or str(resp)).strip()
        raw  = _re.sub(r"^```(?:json)?\s*", "", raw)
        raw  = _re.sub(r"\s*```$",          "", raw).strip()
        val  = _json.loads(raw).get("inflation_pct")
        return float(val) if val is not None else None
    except Exception as exc:
        logger.info(f"LLM inflation estimate failed ({supplier_country} {ref_year}-{current_year}): {exc}")
        return None


def _get_pr_master_date_for_dtl_id(tgt_cs: str, dtl_id: int):
    """Get purchase_req_mst.C_DATETIME for a purchase_dtl_id."""
    if not dtl_id:
        return None
    try:
        conn = _connect(tgt_cs)
        cur = conn.cursor()
        cur.execute("""
            SELECT TOP 1 prm.[C_DATETIME]
            FROM [ras_procurement].[purchase_req_detail] prd
            JOIN [ras_procurement].[purchase_req_mst] prm
              ON prd.[PURCHASE_REQ_ID] = prm.[PURCHASE_REQ_ID]
            WHERE prd.[PURCHASE_DTL_ID] = ?
        """, dtl_id)
        row = cur.fetchone()
        conn.close()
        return row[0] if row else None
    except Exception as exc:
        logger.warning(f"Could not fetch PR master date for dtl_id={dtl_id}: {exc}")
        return None


def _is_pinecone_server_error(exc: Exception) -> bool:
    """True for hard server/connectivity errors that mean Pinecone is unavailable."""
    msg = str(exc).lower()
    return any(kw in msg for kw in (
        "connection", "timeout", "503", "502", "500", "unavailable",
        "network", "refused", "reset", "ssl",
    ))


def _write_benchmark_stub(cur, dtl_id: int, item_uuid: str, summary: str) -> None:
    """Write a placeholder benchmark row so every selected item appears in vw_benchmark_summary."""
    try:
        cur.execute("""
            MERGE [ras_procurement].[benchmark_result] WITH (HOLDLOCK) AS target
            USING (SELECT ? AS purchase_dtl_id) AS src
               ON target.purchase_dtl_id = src.purchase_dtl_id
            WHEN MATCHED THEN
                UPDATE SET extracted_item_uuid_fk=?, bp_unit_price=NULL, bp_total_price=NULL,
                           low_hist_item_fk=NULL, last_hist_item_fk=NULL,
                           inflation_pct=NULL, cpi_inflation_pct=NULL,
                           similar_dtl_ids=NULL, summary=?, updated_at=SYSUTCDATETIME()
            WHEN NOT MATCHED THEN
                INSERT (purchase_dtl_id, extracted_item_uuid_fk, bp_unit_price, bp_total_price,
                        low_hist_item_fk, last_hist_item_fk, inflation_pct, cpi_inflation_pct,
                        similar_dtl_ids, summary)
                VALUES (?, ?, NULL, NULL, NULL, NULL, NULL, NULL, NULL, ?);
        """, dtl_id, item_uuid, summary, dtl_id, item_uuid, summary)
    except Exception as write_exc:
        logger.warning(f"Benchmark stub write failed dtl_id={dtl_id}: {write_exc}")


def _run_benchmark(
    llm,
    tgt_cs: str,
    pr_no: str,
    embed_model,
    pinecone_index: str,
    pinecone_ns: str,
    top_k: int,
    prompts: dict | None = None,
) -> None:
    # Structural errors (import, DB connect, query) propagate — caller records exception.
    from agentcore.services.pinecone_service_client import search_via_service
    from langchain_core.messages import HumanMessage

    # Read benchmark filtering knobs from prompts dict (canvas-tunable).
    p = prompts or {}
    min_score       = float(p.get("bench_min_similarity", 0.80))
    outlier_factor  = float(p.get("bench_outlier_factor", 3.0))
    max_age_months  = int(p.get("bench_max_age_months", 0))
    uom_strict      = bool(p.get("bench_uom_strict", False))
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8], qi.[commodity_tag],
                   qi.[unit_price], qi.[total_price], qi.[quantity], qi.[unit],
                   qi.[currency], qi.[quotation_date], qi.[supplier_name],
                   qi.[unit_price_eur], qi.[total_price_eur],
                   prd.[C_DATETIME] AS [item_created_date]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        all_rows = cur.fetchall()
    finally:
        conn.close()

    bench_cols = [
        "extracted_item_uuid_pk", "purchase_dtl_id",
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag",
        "unit_price", "total_price", "quantity", "unit",
        "currency", "quotation_date", "supplier_name",
        "unit_price_eur", "total_price_eur", "item_created_date",
    ]
    # deduplicate by purchase_dtl_id
    seen: dict = {}
    for row in all_rows:
        rd = dict(zip(bench_cols, row))
        dtl_id = rd["purchase_dtl_id"]
        if dtl_id not in seen:
            seen[dtl_id] = rd
    items = list(seen.values())

    conn2 = _connect(tgt_cs)
    cur2  = conn2.cursor()
    try:
        for rd in items:
            dtl_id    = rd["purchase_dtl_id"]
            item_uuid = str(rd["extracted_item_uuid_pk"] or "")
            qty       = rd.get("quantity") or Decimal("1")

            # Build the same 12-field text used at embedding time
            bench_text = _build_embed_text(rd)
            if not bench_text:
                logger.warning(f"[{pr_no}] dtl_id={dtl_id}: no embedding text — writing stub benchmark row")
                _write_benchmark_stub(cur2, dtl_id, item_uuid, "No embedding text — item descriptor fields are empty")
                continue

            # Item created_date as ISO string for date filtering
            created = rd.get("item_created_date")
            created_iso = (
                created.isoformat() if created and hasattr(created, "isoformat") else str(created or "")
            )

            try:
                embedding = _embed_with_timeout(embed_model, bench_text)
                raw_similar = _pinecone_search_with_timeout(
                    search_via_service,
                    index_name=pinecone_index,
                    namespace=pinecone_ns,
                    text_key="page_content",
                    query=bench_text,
                    query_embedding=embedding,
                    number_of_results=top_k * 3,  # fetch extra to absorb score + date filtered-out
                )
            except Exception as exc:
                if _is_pinecone_server_error(exc):
                    raise RuntimeError(f"Pinecone server error during benchmark for PR={pr_no}: {exc}") from exc
                logger.warning(f"[{pr_no}] Benchmark similarity search failed dtl_id={dtl_id}: {exc}")
                _write_benchmark_stub(cur2, dtl_id, item_uuid, f"Similarity search failed: {type(exc).__name__}")
                continue

            # Normalise response to a list of match dicts
            if isinstance(raw_similar, list):
                matches = raw_similar
            elif isinstance(raw_similar, dict):
                matches = raw_similar.get("results", raw_similar.get("matches", []))
            else:
                matches = []

            # ── Stage 1 filtering: at-vector-DB level ────────────────────
            # Configurable similarity threshold (was hardcoded 0.70).
            # Same-PR + future-date guards retained.
            filtered = []
            for m in matches:
                score = float(m.get("score", 0.0))
                if score < min_score:
                    continue
                meta = m.get("metadata") or {}
                if meta.get("purchase_req_no") == pr_no:
                    continue
                hist_date = meta.get("item_created_date", "")
                if created_iso and hist_date and hist_date >= created_iso:
                    continue
                filtered.append(m)

            similar_dtl_ids_raw = [
                int(m["metadata"]["purchase_dtl_id"])
                for m in filtered[:top_k]
                if m.get("metadata", {}).get("purchase_dtl_id") is not None
            ]
            logger.info(
                f"[{pr_no}] dtl_id={dtl_id}: {len(filtered)} Pinecone match(es) after similarity "
                f"filter (min_score={min_score}) → dtl_ids={similar_dtl_ids_raw}"
            )

            # Fetch full historical pricing from DB
            historical_raw = _fetch_historical_for_dtl_ids(tgt_cs, similar_dtl_ids_raw) if similar_dtl_ids_raw else []

            # Fetch prior benchmark results for similar items as LLM context
            prior_benchmarks = _fetch_benchmark_for_dtl_ids(tgt_cs, similar_dtl_ids_raw) if similar_dtl_ids_raw else []

            # ── Stage 2 filtering: post-fetch (uses DB-only fields) ─────
            # UOM / outlier / age filters take effect here because Pinecone
            # metadata doesn't carry the unit, eur price, or full date —
            # we only have those after the join with quotation_extracted_items.
            historical = historical_raw
            historical = _filter_historical_uom(historical, rd.get("unit"), uom_strict)
            historical = _filter_historical_age(historical, created_iso, max_age_months)
            historical = _filter_historical_outliers(historical, outlier_factor)
            if len(historical) != len(historical_raw):
                logger.info(
                    f"[{pr_no}] dtl_id={dtl_id}: filtered {len(historical_raw) - len(historical)}/{len(historical_raw)} "
                    f"historical row(s) (uom_strict={uom_strict}, max_age_months={max_age_months}, "
                    f"outlier_factor={outlier_factor})"
                )
            if similar_dtl_ids_raw and not historical:
                logger.warning(
                    f"[{pr_no}] dtl_id={dtl_id}: {len(similar_dtl_ids_raw)} Pinecone match(es) found "
                    f"but ALL eliminated by post-fetch filters (uom_strict={uom_strict}, "
                    f"max_age_months={max_age_months}, outlier_factor={outlier_factor}) — "
                    f"low_hist/last_hist will be NULL"
                )

            # Pre-compute aggregates for the LLM prompt (median, P25/P75,
            # time buckets, per-supplier mini-summary, qty context).
            agg = _benchmark_aggregates(historical, created_iso)

            low_item, last_item = _compute_low_last(historical)
            low_uuid  = str(low_item["extracted_item_uuid_pk"])  if low_item  else None
            last_uuid = str(last_item["extracted_item_uuid_pk"]) if last_item else None
            similar_dtl_ids_json = json.dumps(similar_dtl_ids_raw) if similar_dtl_ids_raw else None

            # Inflation estimates (non-fatal — both return None on failure)
            supplier_country = (low_item.get("supplier_country") if low_item else None) or rd.get("supplier_country")
            # Use purchase_req_mst.C_DATETIME instead of quotation_date for consistency
            ref_dt       = _get_pr_master_date_for_dtl_id(tgt_cs, low_item.get("purchase_dtl_id")) if low_item else None
            ref_year     = ref_dt.year if ref_dt and hasattr(ref_dt, "year") else None
            current_year = created.year if created and hasattr(created, "year") else None
            item_category = " > ".join(
                str(rd[f]) for f in ["item_level_1", "item_level_2", "item_level_3"] if rd.get(f)
            )
            infl_dec = cpi_dec = Decimal("0")
            infl_dec_last = cpi_dec_last = Decimal("0")
            if ref_year and current_year and ref_year < current_year:
                infl_raw = _estimate_inflation_via_llm(
                    llm, rd.get("item_name"), item_category or None,
                    supplier_country, ref_year, current_year,
                )
                if infl_raw is not None:
                    try: infl_dec = Decimal(str(infl_raw))
                    except Exception: pass
                cpi_raw = _compute_cpi_pct(supplier_country, ref_year, current_year)
                if cpi_raw is not None:
                    try: cpi_dec = Decimal(str(cpi_raw))
                    except Exception: pass
                logger.info(
                    f"[{pr_no}] dtl_id={dtl_id}: inflation_pct={infl_dec} "
                    f"cpi_pct={cpi_dec} (country={supplier_country!r} {ref_year}-{current_year})"
                )

            # Also calculate inflation for last_hist_item (most recent)
            if last_item and current_year:
                supplier_country_last = last_item.get("supplier_country")
                if not supplier_country_last:
                    logger.info(f"[{pr_no}] dtl_id={dtl_id}: last_hist_item missing supplier_country, skipping inflation_pct_last")
                else:
                    ref_dt_last = _get_pr_master_date_for_dtl_id(tgt_cs, last_item.get("purchase_dtl_id"))
                    ref_year_last = ref_dt_last.year if ref_dt_last and hasattr(ref_dt_last, "year") else None
                    if ref_year_last and current_year and ref_year_last < current_year:
                        infl_raw_last = _estimate_inflation_via_llm(
                            llm, rd.get("item_name"), item_category or None,
                            supplier_country_last, ref_year_last, current_year,
                        )
                        if infl_raw_last is not None:
                            try: infl_dec_last = Decimal(str(infl_raw_last))
                            except Exception: pass
                        cpi_raw_last = _compute_cpi_pct(supplier_country_last, ref_year_last, current_year)
                        if cpi_raw_last is not None:
                            try: cpi_dec_last = Decimal(str(cpi_raw_last))
                            except Exception: pass
                        logger.info(
                            f"[{pr_no}] dtl_id={dtl_id}: inflation_pct_last={infl_dec_last} "
                            f"cpi_pct_last={cpi_dec_last} (country={supplier_country_last!r} {ref_year_last}-{current_year})"
                        )

            # LLM benchmark analysis — uses pre-computed aggregates so the
            # LLM can anchor on the median + P25-P75 band + time buckets
            # without re-deriving them mentally from raw rows.
            bench_prompt = _format_bench_prompt(rd, historical[:top_k], stats=agg, prior_benchmarks=prior_benchmarks or None)
            bout = {}
            llm_succeeded = False
            try:
                resp = _call_llm_with_retry(llm, [HumanMessage(content=bench_prompt)], prompts=prompts, force_json=True)
                raw  = (getattr(resp, "content", None) or str(resp)).strip()
                raw  = re.sub(r"^```(?:json)?\s*", "", raw)
                raw  = re.sub(r"\s*```$", "", raw)
                bout = json.loads(raw)
                llm_succeeded = True
            except Exception as exc:
                logger.error(
                    f"[{pr_no}] LLM benchmark parse failed dtl_id={dtl_id}: {exc}",
                    exc_info=True,
                )

            bp_unit       = bout.get("bp_unit_price")
            bp_low        = bout.get("bp_low")
            bp_high       = bout.get("bp_high")
            bp_confidence = bout.get("confidence")
            llm_summary   = bout.get("summary", "")
            llm_succeeded = llm_succeeded and bp_unit is not None
            try:
                bp_dec   = Decimal(str(bp_unit)) if bp_unit is not None else Decimal("0")
                bp_total = round(float(bp_dec) * float(qty or 1), 2)
            except Exception:
                bp_dec = Decimal("0")
                bp_total = 0

            # Fold the new range + confidence + sample-size info into the
            # human-readable summary so we don't need a schema change to the
            # benchmark_result table. bp_unit_price stays the canonical
            # point estimate; the prefix gives consumers the band.
            extras = []
            if bp_low is not None or bp_high is not None:
                extras.append(f"range={bp_low}–{bp_high} EUR")
            if bp_confidence is not None:
                try:
                    extras.append(f"confidence={float(bp_confidence):.2f}")
                except Exception:
                    extras.append(f"confidence={bp_confidence}")
            if agg.get("count") is not None:
                extras.append(f"n={agg['count']} (priced={agg.get('priced_count', 0)})")
            if "median" in agg:
                extras.append(f"median={agg['median']} EUR")
            summary = (
                f"[{', '.join(extras)}] {llm_summary}".strip()
                if extras else llm_summary
            )

            # When LLM failed, pass None for price/summary in UPDATE so COALESCE
            # preserves any existing good values rather than overwriting with 0.
            bp_dec_u   = bp_dec   if llm_succeeded else None
            bp_total_u = bp_total if llm_succeeded else None
            summary_u  = summary  if (llm_succeeded and summary) else None

            try:
                cur2.execute("""
                    MERGE [ras_procurement].[benchmark_result] WITH (HOLDLOCK) AS target
                    USING (SELECT ? AS purchase_dtl_id) AS src
                       ON target.purchase_dtl_id = src.purchase_dtl_id
                    WHEN MATCHED THEN
                        UPDATE SET
                            extracted_item_uuid_fk = ?,
                            bp_unit_price          = COALESCE(?, target.bp_unit_price),
                            bp_total_price         = COALESCE(?, target.bp_total_price),
                            low_hist_item_fk       = ?,
                            last_hist_item_fk      = ?,
                            inflation_pct          = ?,
                            cpi_inflation_pct      = ?,
                            inflation_pct_last     = ?,
                            cpi_inflation_pct_last = ?,
                            similar_dtl_ids        = ?,
                            summary                = COALESCE(?, target.summary),
                            updated_at             = SYSUTCDATETIME()
                    WHEN NOT MATCHED THEN
                        INSERT (
                            purchase_dtl_id, extracted_item_uuid_fk,
                            bp_unit_price, bp_total_price,
                            low_hist_item_fk, last_hist_item_fk,
                            inflation_pct, cpi_inflation_pct,
                            inflation_pct_last, cpi_inflation_pct_last,
                            similar_dtl_ids, summary
                        )
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?);
                """,
                    dtl_id,
                    item_uuid, bp_dec_u, bp_total_u, low_uuid, last_uuid, infl_dec, cpi_dec, infl_dec_last, cpi_dec_last, similar_dtl_ids_json, summary_u,
                    dtl_id, item_uuid, bp_dec, bp_total, low_uuid, last_uuid, infl_dec, cpi_dec, infl_dec_last, cpi_dec_last, similar_dtl_ids_json, summary,
                )
            except Exception as exc:
                logger.warning(f"[{pr_no}] Benchmark write failed dtl_id={dtl_id}: {exc}")

        conn2.commit()
    except Exception:
        try: conn2.rollback()
        except Exception: pass
        raise
    finally:
        conn2.close()


# ── Fetch pending PRs ──────────────────────────────────────────────────────────

def _get_pr_current_stage(tgt_cs: str, pr_no: str) -> int:
    """Return the current_stage_fk for a PR, or 0 if not found."""
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute(
            "SELECT current_stage_fk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?",
            pr_no,
        )
        row = cur.fetchone()
        return int(row[0]) if row and row[0] is not None else 0
    finally:
        conn.close()


def _fetch_pending_prs(tgt_cs: str, pr_filter: str, batch_limit: int) -> list:
    if pr_filter:
        return [pr_filter]
    # Pick up PRs at any stage that needs stages 4-8 work:
    #   stage 3 = just finished Stage 1-3, ready for classification
    #   stages 4-7 = partially processed (e.g. failed mid-run), need to resume
    _WHERE = """
        FROM [ras_procurement].[purchase_req_mst] prm
        JOIN [ras_procurement].[ras_tracker] rt
          ON prm.PURCHASE_REQ_NO = rt.purchase_req_no
       WHERE rt.current_stage_fk IN (3, 4, 5, 6, 7)
         AND UPPER(prm.PURCHASEFINALAPPROVALSTATUS)
                 IN ('APPROVED BY ALL', 'APPROVED BY ALL EXCEPTION')
       ORDER BY prm.C_DATETIME ASC
    """
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        if batch_limit > 0:
            cur.execute(f"SELECT TOP (?) prm.PURCHASE_REQ_NO {_WHERE}", batch_limit)
        else:
            cur.execute(f"SELECT prm.PURCHASE_REQ_NO {_WHERE}")
        return [row[0] for row in cur.fetchall()]
    finally:
        conn.close()


# ── Excel-driven PR list (mirrors doc-intel run_pipeline_from_excel.py) ──

_AUTO_DETECT_PR_COLUMNS = {
    "purchase_req_no", "purchase req no",
    "pr_no", "pr no", "pr number",
    "ras_no", "ras no", "ras number",
}


def _parse_pr_list_from_excel_bytes(
    excel_bytes: bytes,
    column: str | None = None,
    sheet: str | None = None,
) -> list[str]:
    """Parse purchase requisition numbers out of an .xlsx blob.

    Mirrors doc-intel run_pipeline_from_excel.py._read_pr_nos_from_excel,
    adapted to read from in-memory bytes (Azure Blob / wired Data input)
    instead of a local file path.

    Column selection:
      1. explicit `column` arg → exact case-insensitive header match
      2. auto-detect from PURCHASE_REQ_NO / PR_NO / RAS_NO / "PR Number" etc.
      3. fall back to the first column

    Returns deduplicated, non-empty PR numbers in the order they appear.
    """
    import io as _io
    from openpyxl import load_workbook

    wb = load_workbook(_io.BytesIO(excel_bytes), read_only=True, data_only=True)
    try:
        sheet_name = (sheet or "").strip()
        if sheet_name:
            if sheet_name not in wb.sheetnames:
                raise ValueError(
                    f"Sheet {sheet_name!r} not found. Available: {wb.sheetnames}"
                )
            ws = wb[sheet_name]
        else:
            ws = wb.active

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return []

        headers = [
            str(h).strip() if h is not None else ""
            for h in rows[0]
        ]

        col_idx: int | None = None
        col_arg = (column or "").strip()
        if col_arg:
            for i, h in enumerate(headers):
                if h.lower() == col_arg.lower():
                    col_idx = i
                    break
            if col_idx is None:
                raise ValueError(
                    f"Column {col_arg!r} not found. Headers: {headers}"
                )
        else:
            for i, h in enumerate(headers):
                if h.lower() in _AUTO_DETECT_PR_COLUMNS:
                    col_idx = i
                    logger.info("Excel PR list — auto-detected column {!r}", h)
                    break
            if col_idx is None:
                col_idx = 0
                logger.info(
                    "Excel PR list — no recognised PR column header; using first column {!r}",
                    headers[0] if headers else "<empty>",
                )

        seen: set[str] = set()
        pr_nos: list[str] = []
        for row in rows[1:]:
            if col_idx >= len(row):
                continue
            val = row[col_idx]
            if val is None:
                continue
            pr_no = str(val).strip()
            if pr_no and pr_no not in seen:
                seen.add(pr_no)
                pr_nos.append(pr_no)
        return pr_nos
    finally:
        wb.close()


# ─────────────────────────────────────────────────────────────────────────────
# V2 helpers — category resolution, brand-leak scrubbing, validators
# ─────────────────────────────────────────────────────────────────────────────


# Legacy → canonical category aliases. Kept inline as a default argument
# in each function that uses it (NOT as a module-level constant) — MiCore
# Custom-Component exec contexts don't reliably expose module-level
# assignments to function bodies via LOAD_GLOBAL, but default args are
# captured at def-time which IS reliable in that context. The dict
# is therefore duplicated across three functions; keep them in sync
# when adding new aliases.


def _resolve_category_list(
    tgt_cs: str,
    _SEED: list = (
        # ── Standard procurement taxonomy ──────────────────────────────────
        "Moulding Machine",
        "Moulding Machine Auxiliaries",
        "Assembly / Secondary Equipment",
        "Robot",
        "IT and Communications",
        "Laboratory and Test Equipment",
        "Packaging",
        "Fire Protection",
        "Mould Change Equipment",
        "Paint Shop",
        "Plant equipment",
        "Tooling CapEx",
        "Wire Processing Equipment",
        "Facilities Management",
        "Fleet Management",
        "Machine Lease",
        "Maintenance Repair and Operation",
        "Material Handling Equipment (MHE) - Lease",
        "Material Handling Equipment (MHE) - Purchase",
        "Office Supplies Equipment",
        "Property/Land lease",
        "Tele Communication",
        # ── Other categories preserved from the original seed ─────────────
        # (kept distinct because they're NOT strict duplicates of any item
        # in the standard taxonomy — only the IMM-family is aliased below)
        "Compressor", "TBE Wiring Harness", "Tooling Rubber",
        "Central Portal Purchase", "IT Hardware", "Oil and Lubricants",
        "Measuring Instruments", "HVAC Air Conditioner", "Vehicle Car",
        "Group Company", "TBE Fee and Charges", "Tooling Plastic", "VMC",
        "Storage Equipment Trolley", "Robot and Gripper", "Wiring Harness",
        "Pump and Motor", "TBE Negotiated by Unit and Overseas Unit",
        "Gensets", "Environmental Chamber", "Software", "Transformer",
        "Tooling Die Casting", "Testing", "Electrical work",
        "Vehicle Motorcycle", "Fabrication", "CCTV", "Chillers",
        "Service and Repair", "Tool Room Machines", "Storage Equipment",
        "Consumable", "certification", "UPS", "Spares",
        "Storage Equipment Racks", "Post Facto",
        "TBE Customer and JV Partner directed", "Piping", "Electrical Panel",
        "Jigs and Fixtures", "MHE",
        "Wire Extruder and Auxiliary Equip", "SPM", "AMC",
        "Storage Equipment Bin", "Calibration",
        "Fire Safety and Alarm System", "TBE Negotiated by MCO Tooling",
        "Treatment Plant", "Moulding Machine Spares",
    ),
    _ALIASES: dict = {
        # Inline duplicate of category alias map (see comment above the
        # function — MiCore exec contexts require default-arg capture).
        "IMM":                          "Moulding Machine",
        "I.M.M.":                       "Moulding Machine",
        "Injection Moulding Machine":   "Moulding Machine",
        "Injection Molding Machine":    "Moulding Machine",
        "Molding Machine":              "Moulding Machine",
        "IMM Auxiliary":                "Moulding Machine Auxiliaries",
        "IMM Auxiliaries":              "Moulding Machine Auxiliaries",
        "Molding Machine Auxiliary":    "Moulding Machine Auxiliaries",
        "Molding Machine Auxiliaries":  "Moulding Machine Auxiliaries",
        "Moulding Machine Auxiliary":   "Moulding Machine Auxiliaries",
        "IMM Spare":                    "Moulding Machine Spares",
        "IMM Spares":                   "Moulding Machine Spares",
        "Moulding Machine Spare":       "Moulding Machine Spares",
        "Molding Machine Spare":        "Moulding Machine Spares",
        "Molding Machine Spares":       "Moulding Machine Spares",
    },
    _TTL: int = 300,
) -> list[str]:
    """Returns the seed list UNION every distinct purchase_category_llm
    previously persisted by V2, with legacy names normalised via _ALIASES
    so the LLM never sees duplicates (e.g. both "IMM" and "Moulding Machine").
    Cached in-process for 5 min so the LLM call doesn't hammer the DB on
    every PR. On DB error, falls back to seed list.

    Hardened against MiCore Custom-Component exec contexts where module-level
    `=` assignments may not be visible to function bodies via LOAD_GLOBAL:
      • The cache lives on the function object (`__dict__["_cache"]`) — set
        once on first call, never goes through LOAD_GLOBAL.
      • The seed list, alias map, and TTL are captured as default arguments
        at def-time, so the function body never has to resolve them from
        module globals.
    """
    cache = _resolve_category_list.__dict__.setdefault("_cache", {})
    now = time.time()
    cached = cache.get(tgt_cs)
    if cached and now - cached[0] < _TTL:
        return cached[1]
    db_categories: list[str] = []
    try:
        conn = _connect(tgt_cs)
        try:
            cur = conn.cursor()
            cur.execute("""
                SELECT DISTINCT [purchase_category_llm]
                  FROM [ras_procurement].[quotation_extracted_items]
                 WHERE [purchase_category_llm] IS NOT NULL
                   AND LTRIM(RTRIM([purchase_category_llm])) <> ''
            """)
            db_categories = [r[0] for r in cur.fetchall() if r and r[0]]
        finally:
            conn.close()
    except Exception as exc:
        logger.warning(f"V2 _resolve_category_list — DB lookup failed: {exc}")
    # Filter blanks and sentinel strings ("None", "null", "n/a", "Others")
    # so they never end up in the canonical category list passed to the LLM.
    _DROP = {"", "none", "null", "n/a", "na", "others", "other"}
    def _ok(v) -> bool:
        if v is None:
            return False
        s = str(v).strip()
        return bool(s) and s.lower() not in _DROP
    # Case-insensitive + whitespace-normalised alias lookup so "imm",
    # "  Injection   Molding Machine  ", "Molding machine" all hit the
    # same canonical bucket as "Moulding Machine".
    _alias_norm = {" ".join(k.lower().split()): v for k, v in _ALIASES.items()}
    seed_clean = [str(c).strip() for c in _SEED if _ok(c)]
    seed_norm  = {" ".join(c.lower().split()): c for c in seed_clean}
    def _canon(v: Any, *, use_fuzzy: bool = False) -> str:
        s = str(v).strip()
        if not s:
            return s
        n = " ".join(s.lower().split())
        # 1. Explicit alias map (IMM-family, etc.)
        if n in _alias_norm:
            return _alias_norm[n]
        # 2. Already canonical (case/whitespace-normalised)
        if n in seed_norm:
            return seed_norm[n]
        # 3. Fuzzy match against seed — only applied to DB values so the
        #    LLM never sees near-duplicates like "Compressors" vs "Compressor"
        if use_fuzzy and seed_norm:
            import difflib
            m = difflib.get_close_matches(n, list(seed_norm.keys()), n=1, cutoff=0.88)
            if m:
                snapped = seed_norm[m[0]]
                logger.info(
                    f"V2 _resolve_category_list: fuzzy-snap DB value "
                    f"'{s}' → '{snapped}'"
                )
                return snapped
        return s
    seed  = {_canon(c)                   for c in _SEED          if _ok(c)}
    dbset = {_canon(c, use_fuzzy=True)   for c in db_categories  if _ok(c)}
    combined = sorted(seed | dbset, key=str.lower)
    cache[tgt_cs] = (now, combined)
    return combined


def _category_with_aliases(
    category: str,
    _ALIASES: dict = {
        # Inline duplicate of category alias map — see _resolve_category_list.
        "IMM":                          "Moulding Machine",
        "I.M.M.":                       "Moulding Machine",
        "Injection Moulding Machine":   "Moulding Machine",
        "Injection Molding Machine":    "Moulding Machine",
        "Molding Machine":              "Moulding Machine",
        "IMM Auxiliary":                "Moulding Machine Auxiliaries",
        "IMM Auxiliaries":              "Moulding Machine Auxiliaries",
        "Molding Machine Auxiliary":    "Moulding Machine Auxiliaries",
        "Molding Machine Auxiliaries":  "Moulding Machine Auxiliaries",
        "Moulding Machine Auxiliary":   "Moulding Machine Auxiliaries",
        "IMM Spare":                    "Moulding Machine Spares",
        "IMM Spares":                   "Moulding Machine Spares",
        "Moulding Machine Spare":       "Moulding Machine Spares",
        "Molding Machine Spare":        "Moulding Machine Spares",
        "Molding Machine Spares":       "Moulding Machine Spares",
    },
) -> list[str]:
    """Return the canonical category plus any legacy aliases that map to it.
    Used by SQL Stage A so a search for "Moulding Machine" also matches
    historical rows still classified as "IMM" / "Injection Molding Machine"
    / etc. Lookup is case-insensitive and whitespace-normalised."""
    if not category:
        return []
    alias_norm = {" ".join(k.lower().split()): v for k, v in _ALIASES.items()}
    raw = category.strip()
    canon = alias_norm.get(" ".join(raw.lower().split()), raw)
    legacy = [k for k, v in _ALIASES.items() if v == canon]
    seen: set = set()
    out: list[str] = []
    for name in [canon] + legacy:
        key = name.lower()
        if key not in seen:
            seen.add(key)
            out.append(name)
    return out


def _normalise_category(
    value: Any,
    allowed: list[str],
    aliases: dict = {
        # Inline duplicate of category alias map — see _resolve_category_list.
        "IMM":                          "Moulding Machine",
        "I.M.M.":                       "Moulding Machine",
        "Injection Moulding Machine":   "Moulding Machine",
        "Injection Molding Machine":    "Moulding Machine",
        "Molding Machine":              "Moulding Machine",
        "IMM Auxiliary":                "Moulding Machine Auxiliaries",
        "IMM Auxiliaries":              "Moulding Machine Auxiliaries",
        "Molding Machine Auxiliary":    "Moulding Machine Auxiliaries",
        "Molding Machine Auxiliaries":  "Moulding Machine Auxiliaries",
        "Moulding Machine Auxiliary":   "Moulding Machine Auxiliaries",
        "IMM Spare":                    "Moulding Machine Spares",
        "IMM Spares":                   "Moulding Machine Spares",
        "Moulding Machine Spare":       "Moulding Machine Spares",
        "Molding Machine Spare":        "Moulding Machine Spares",
        "Molding Machine Spares":       "Moulding Machine Spares",
    },
    _FUZZY_CUTOFF: float = 0.85,
) -> str:
    """Snap an LLM-emitted category to the canonical name from `allowed`.

    Generic — works for any product category, not IMM-specific. The DB
    never accumulates near-duplicates regardless of how the LLM spells
    or punctuates the category. Tiered matching (most specific first):

      0a. Alias map exact   — case-insensitive, whitespace-collapsed key
                              ("IMM" → "Moulding Machine"; "IMM Auxiliary"
                              → "Moulding Machine Auxiliaries" — these two
                              stay DISTINCT because the alias map has
                              separate entries for each).
      0b. Alias map punct   — punctuation-stripped key
                              ("I.M.M." → "Moulding Machine"; "I.M.M.
                              Auxiliary" → "Moulding Machine Auxiliaries").
      1.  Canonical exact   — case + whitespace normalised.
      2.  Canonical punct   — punctuation stripped.
      3.  Canonical fuzzy   — difflib SequenceMatcher ratio >= 0.85
                              ("Compressors" → "Compressor"; "Moudling
                              Machine" → "Moulding Machine"; "Robot Arm"
                              vs "Robot" stays separate at ~0.62).

    Importantly, IMM-family aliases are checked BEFORE fuzzy matching so
    "IMM Auxiliary" cannot be mismerged into "IMM" or "Moulding Machine"
    by accidental string similarity — each alias key maps to its specific
    canonical target.

    Logs every snap so corrections are auditable.
    """
    if not value:
        return ""
    raw = str(value).strip()
    if not raw:
        return ""
    if not allowed:
        return raw

    def _norm(s: str) -> str:
        return " ".join(s.lower().split())

    def _strip_punct(s: str) -> str:
        return re.sub(r"[^a-z0-9]+", "", _norm(s))

    raw_norm  = _norm(raw)
    raw_punct = _strip_punct(raw)

    # 0. Alias map — check FIRST so "IMM" → "Moulding Machine" works even
    #    when fuzzy match would fail (IMM too lexically different from
    #    Moulding Machine). Distinct alias keys keep IMM vs IMM Auxiliary
    #    routed to their separate canonical targets.
    if aliases:
        alias_norm  = {_norm(k): v for k, v in aliases.items()}
        if raw_norm in alias_norm:
            return alias_norm[raw_norm]
        alias_punct = {_strip_punct(k): v for k, v in aliases.items()}
        if raw_punct and raw_punct in alias_punct:
            snapped = alias_punct[raw_punct]
            logger.info(f"_normalise_category: alias-punct-snap '{raw}' → '{snapped}'")
            return snapped

    # 1. Exact match against canonical list
    lookup_norm = {_norm(c): c for c in allowed}
    if raw_norm in lookup_norm:
        return lookup_norm[raw_norm]

    # 2. Punctuation-stripped match against canonical list
    lookup_punct = {_strip_punct(c): c for c in allowed}
    if raw_punct and raw_punct in lookup_punct:
        snapped = lookup_punct[raw_punct]
        logger.info(f"_normalise_category: punct-snap '{raw}' → '{snapped}'")
        return snapped

    # 3. Fuzzy match against canonical (high cutoff — keeps distinct
    #    concepts like "Robot Arm" vs "Robot" separate)
    import difflib
    matches = difflib.get_close_matches(
        raw_norm, list(lookup_norm.keys()), n=1, cutoff=_FUZZY_CUTOFF
    )
    if matches:
        snapped = lookup_norm[matches[0]]
        logger.info(
            f"_normalise_category: fuzzy-snap '{raw}' → '{snapped}' "
            f"(ratio>={_FUZZY_CUTOFF})"
        )
        return snapped

    # 4. Genuinely new category — caller treats as a proposal
    return raw


def _strip_brand_from_embed_content(text: str, brand_tokens: list[str]) -> tuple[str, list[str]]:
    """Best-effort regex strip of brand tokens from embed_content. Returns
    (cleaned_text, removed_tokens). Case-insensitive, whole-word match."""
    if not text:
        return text, []
    removed: list[str] = []
    out = text
    for tok in brand_tokens:
        if not tok or len(tok.strip()) < 2:
            continue
        pat = re.compile(rf"\b{re.escape(tok.strip())}\b", re.IGNORECASE)
        if pat.search(out):
            out = pat.sub("", out)
            removed.append(tok.strip())
    if removed:
        out = re.sub(r"\s+", " ", out).strip(" ,.")
    return out, removed


def _coerce_number_v2(v: Any) -> Optional[float]:
    if v is None:
        return None
    if isinstance(v, (int, float)):
        return float(v)
    if isinstance(v, str):
        m = re.search(r"-?\d+(?:\.\d+)?", v)
        if m:
            try:
                return float(m.group(0))
            except Exception:
                return None
    return None


def _validate_critical_attributes(raw: Any) -> list[dict]:
    """Coerce LLM-emitted critical_attributes into the canonical schema.
    Drops malformed entries; never raises."""
    if not isinstance(raw, list):
        return []
    out: list[dict] = []
    for entry in raw:
        if not isinstance(entry, dict):
            continue
        name = str(entry.get("name") or "").strip().lower()
        if not name:
            continue
        val = _coerce_number_v2(entry.get("value"))
        if val is None:
            sv = entry.get("value")
            if sv is None:
                continue
            value: Any = str(sv).strip()
        else:
            value = val
        unit = str(entry.get("unit") or "").strip()
        importance = str(entry.get("importance") or "important").strip().lower()
        if importance not in ("critical", "important", "informational"):
            importance = "important"
        out.append({"name": name, "value": value, "unit": unit, "importance": importance})
    return out


# ─────────────────────────────────────────────────────────────────────────────
# V2 Stage 5 — extraction with category + embed_content + critical_attributes
# ─────────────────────────────────────────────────────────────────────────────


def _build_extraction_user_prompt_v2(
    ctx: RASContext,
    doc: DocumentContent,
    prompts: dict | None = None,
    category_list: list[str] | None = None,
) -> str:
    """V2 user prompt — formats EXTRACTION_USER_TEMPLATE_V2 with the seed/DB
    category list interpolated into {category_list}."""
    def _f(v): return str(v) if v is not None else "N/A"

    if doc.ocr_source and doc.text:
        doc_content_str = f"[OCR markdown from Azure Document Intelligence]\n\n{doc.text}"
    elif doc.images and doc.text:
        doc_content_str = f"[Extracted text — page images attached below]\n\n{doc.text}"
    elif doc.images:
        doc_content_str = "[Scanned document — page image(s) attached]"
    else:
        doc_content_str = doc.text or "[No content extracted]"

    ext_max = (prompts or {}).get("ext_max_chars")
    if ext_max and doc_content_str and len(doc_content_str) > ext_max:
        doc_content_str = doc_content_str[:ext_max]

    cat_list = category_list or PURCHASE_CATEGORIES_SEED
    cat_bullet = "\n".join(f"  - {c}" for c in cat_list)

    user_tmpl = (prompts or {}).get("ext_user_v2", EXTRACTION_USER_TEMPLATE_V2)
    return user_tmpl.format(
        purchase_req_no=_f(ctx.purchase_req_no),
        purchase_req_id=_f(ctx.purchase_req_id),
        justification=_f(ctx.justification),
        supplier_name=_f(ctx.supplier_name),
        currency=_f(ctx.currency),
        enquiry_no=_f(ctx.enquiry_no),
        classification=_f(ctx.classification),
        department=_f(ctx.department),
        negotiated_by=_f(ctx.negotiated_by),
        address=_f(ctx.address),
        contract_no=_f(ctx.contract_no),
        order_no=_f(ctx.order_no),
        purchase_value=_f(ctx.purchase_value),
        category=_f(ctx.category),
        sub_category=_f(ctx.sub_category),
        l3=_f(ctx.l3),
        l4=_f(ctx.l4),
        purchase_category=_f(ctx.purchase_category),
        ras_title=_f(ctx.ras_title),
        site_region=_f(ctx.site_region),
        site_country=_f(ctx.site_country),
        site=_f(ctx.site),
        division=_f(ctx.division),
        requisition_type=_f(ctx.requisition_type),
        parent_supplier=_f(ctx.parent_supplier),
        supplier_type=_f(ctx.supplier_type),
        supplier_country=_f(ctx.supplier_country),
        payment_days=_f(ctx.payment_days),
        po_date=_f(ctx.po_date),
        category_buyer=_f(ctx.category_buyer),
        line_items_table=_build_line_items_table(ctx),
        item_taxonomy=(prompts or {}).get("ext_taxonomy", ITEM_TAXONOMY),
        document_content=doc_content_str,
        raw_ras_context=_build_raw_context(ctx),
        category_list=cat_bullet,
    )


def _call_extraction_llm_v2(
    llm,
    ctx: RASContext,
    doc: DocumentContent,
    prompts: dict | None,
    category_list: list[str],
) -> str:
    """V2 extraction LLM call — uses V2 system + user prompts."""
    from langchain_core.messages import HumanMessage, SystemMessage
    user_prompt = _build_extraction_user_prompt_v2(ctx, doc, prompts, category_list)
    sys_prompt  = (prompts or {}).get("ext_system_v2", EXTRACTION_SYSTEM_PROMPT_V2)
    max_images  = max(1, int((prompts or {}).get("ext_max_images", 50)))

    messages: list = [SystemMessage(content=sys_prompt)]
    img_count, img_detail = 0, "n/a"
    if doc.is_image_based and doc.images:
        if len(doc.images) > max_images:
            logger.warning(
                "[V2 PR={}] Quotation has {} image(s) but max_images={} — truncating to first {}",
                ctx.purchase_req_no, len(doc.images), max_images, max_images,
            )
        images = doc.images[:max_images]
        img_count  = len(images)
        img_detail = "low" if doc.text else "high"
        content_parts: list = [{"type": "text", "text": user_prompt}]
        for b64 in images:
            content_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}", "detail": img_detail},
            })
        messages.append(HumanMessage(content=content_parts))
    else:
        messages.append(HumanMessage(content=user_prompt))

    logger.info(
        "[V2 PR={}] Calling extraction LLM — {} image(s) detail={}, ~{:.0f} kB prompt",
        ctx.purchase_req_no, img_count, img_detail, len(user_prompt) / 1024,
    )
    response = _call_llm_with_retry(llm, messages, prompts, force_json=True)
    return (getattr(response, "content", None) or str(response)).strip()


def _parse_extraction_response_v2(
    raw: str,
    source: dict,
    ctx: RASContext,
    allowed_categories: list[str],
) -> list:
    """V2 parser — delegates to V1 _parse_extraction_response for the base
    schema then enriches each item with the 3 V2 fields, including a
    brand-leak scrub on embed_content."""
    base_items = _parse_extraction_response(raw, source, ctx)
    if not base_items:
        return base_items

    raw_clean = raw.strip()
    raw_clean = re.sub(r"^```(?:json)?\s*", "", raw_clean)
    raw_clean = re.sub(r"\s*```$", "", raw_clean)
    try:
        envelope = json.loads(raw_clean)
        raw_items = envelope.get("items", []) if isinstance(envelope, dict) else []
    except Exception as exc:
        logger.warning(f"[V2] Could not re-parse raw JSON for V2 fields: {exc}")
        raw_items = []

    raw_by_dtl: dict = {}
    for r in raw_items:
        if isinstance(r, dict):
            dtl = r.get("purchase_dtl_id")
            if dtl is not None:
                try:
                    raw_by_dtl[int(dtl)] = r
                except Exception:
                    pass

    for item in base_items:
        try:
            dtl = int(item.get("purchase_dtl_id") or 0)
        except Exception:
            dtl = 0
        v2_src = raw_by_dtl.get(dtl, {}) if dtl else {}

        # 1. purchase_category_llm — snap to canonical when case-insensitive match
        cat_canon = _normalise_category(v2_src.get("purchase_category_llm"), allowed_categories)
        if not cat_canon:
            # Defensive fallback chain — LLM didn't emit, so fall back to the
            # RAS purchase_category from on-prem; if that's also blank leave
            # it as Unspecified (matches the taxonomy fallback convention).
            # We no longer fall back to "Others" because that category was
            # removed from the canonical list.
            cat_canon = (str(ctx.purchase_category).strip()
                         if getattr(ctx, "purchase_category", None) else "")
            cat_canon = cat_canon or "Unspecified"
            logger.warning(
                f"[V2 PR={ctx.purchase_req_no}] dtl_id={dtl} missing purchase_category_llm — "
                f"defaulting to {cat_canon!r}"
            )
        item["purchase_category_llm"] = cat_canon

        # 2. embed_content — brand-leak scrub
        embed_text = (v2_src.get("embed_content") or "").strip()
        brand_tokens = []
        for f in ("item_level_4", "item_level_5", "supplier_name"):
            t = item.get(f) or v2_src.get(f)
            if t and isinstance(t, str):
                brand_tokens.extend(re.findall(r"[A-Za-z][A-Za-z0-9\-]+", t))
        if embed_text:
            cleaned, removed = _strip_brand_from_embed_content(embed_text, brand_tokens)
            if removed:
                logger.warning(
                    f"[V2 PR={ctx.purchase_req_no}] dtl_id={dtl} embed_content brand-leak scrubbed: "
                    f"removed {removed}"
                )
            embed_text = cleaned
        if not embed_text:
            embed_text = _build_embed_text({**item, **v2_src})
            logger.warning(
                f"[V2 PR={ctx.purchase_req_no}] dtl_id={dtl} no embed_content emitted — "
                f"falling back to 12-field concat"
            )
        item["embed_content"] = embed_text

        # 3. critical_attributes — coerce/validate
        item["critical_attributes"] = _validate_critical_attributes(v2_src.get("critical_attributes"))

    return base_items


def _save_extracted_items_v2(tgt_cs: str, items: list, fallback_date=None) -> int:
    """V2 writer — same shape as _save_extracted_items but adds 3 V2 columns
    plus the Stage-5b commercials columns (21 quote-level + 14 line-level).
    Commercial columns default to NULL when the second LLM call is disabled
    or fails."""
    if not items:
        return 0
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    saved = 0
    _COL_LIMITS = {
        "item_description": 4000, "item_summary": 4000,
        "taxation_details": 2000, "payment_terms": 2000,
        "supplier_address": 500,  "item_name": 500,
        "supplier_name": 255,     "supplier_country": 100,
        "quotation_ref_no": 100,  "commodity_tag": 255,
        "unit": 50,               "purchase_category_llm": 200,
        # ── Commercials (Stage 5b) string column limits — match migration
        "quote_incoterms": 50,             "quote_incoterms_named_place": 255,
        "quote_tax_type": 50,              "quote_grand_total_currency": 10,
        "quote_country_of_origin": 100,    "quote_port_of_loading": 255,
        "quote_port_of_discharge": 255,    "quote_mode_of_transport": 50,
        "quote_delivery_location": 500,
        "line_incoterms": 50,              "line_country_of_origin": 100,
        "line_hsn_sac_code": 50,           "line_charges_source": 30,
        "line_allocation_method": 40,
    }
    try:
        for item in items:
            def _v(k, cast=None):
                v = item.get(k)
                if v is None: return None
                try:
                    v = cast(v) if cast else v
                except Exception:
                    return None
                if isinstance(v, str) and k in _COL_LIMITS:
                    limit = _COL_LIMITS[k]
                    if len(v) > limit:
                        logger.warning("V2 truncating '{}': {} → {} chars", k, len(v), limit)
                        v = v[:limit]
                return v
            def _d(k):
                v = item.get(k)
                if v is None: return None
                if isinstance(v, Decimal): return v
                try: return Decimal(str(v))
                except Exception: return None
            def _date(k):
                v = item.get(k)
                if not v: return None
                try:
                    from datetime import date as date_cls
                    if isinstance(v, date_cls): return v
                    m = re.match(r"(\d{4})-(\d{2})-(\d{2})", str(v))
                    if m: return date_cls(int(m[1]), int(m[2]), int(m[3]))
                except Exception: pass
                return None

            unit_price_eur  = _convert_to_eur(tgt_cs, item.get("unit_price"),  item.get("currency"), fallback_date)
            total_price_eur = _convert_to_eur(tgt_cs, item.get("total_price"), item.get("currency"), fallback_date)

            crit_attrs = item.get("critical_attributes") or []
            critical_attributes_json = json.dumps(crit_attrs, ensure_ascii=False) if crit_attrs else None

            cur.execute("""
                INSERT INTO [ras_procurement].[quotation_extracted_items] (
                    [attachment_classify_fk],[embedded_classify_fk],[purchase_dtl_id],
                    [is_selected_quote],[supplier_match_conf],[quote_rank],
                    [supplier_name],[supplier_address],[supplier_country],
                    [quotation_ref_no],[quotation_date],[currency],
                    [validity_date],[validity_days],[payment_terms],
                    [item_name],[item_description],[quantity],[unit],
                    [unit_price],[total_price],[discount],[taxation_details],
                    [delivery_date],[delivery_time_days],
                    [item_level_1],[item_level_2],[item_level_3],[item_level_4],
                    [item_level_5],[item_level_6],[item_level_7],[item_level_8],
                    [commodity_tag],[item_summary],
                    [unit_price_eur],[total_price_eur],
                    [purchase_category_llm],[embed_content],[critical_attributes],
                    -- ── Commercials (Stage 5b) — quote-level (21)
                    [quote_incoterms],[quote_incoterms_named_place],
                    [quote_subtotal],[quote_discount_total],
                    [quote_packing_forwarding],[quote_freight],[quote_insurance],
                    [quote_customs_duties],[quote_installation],[quote_other_charges],
                    [quote_tax_type],[quote_tax_rate_pct],[quote_tax_amount_total],
                    [quote_grand_total],[quote_grand_total_currency],
                    [quote_country_of_origin],[quote_port_of_loading],
                    [quote_port_of_discharge],[quote_mode_of_transport],
                    [quote_delivery_location],[quote_charges_breakdown_json],
                    -- ── Commercials (Stage 5b) — line-level (14)
                    [line_incoterms],[line_tax_rate_pct],[line_tax_amount],
                    [line_freight],[line_insurance],[line_packing_forwarding],
                    [line_customs_duties],[line_installation],[line_other_charges],
                    [line_total_inclusive],[line_country_of_origin],
                    [line_hsn_sac_code],[line_charges_source],[line_allocation_method]
                ) VALUES (
                    ?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,
                    ?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,
                    ?,?,?,?,?,?,?,?,?,?,?,?,?,?
                )
            """,
                _v("attachment_classify_fk"), _v("embedded_classify_fk"),
                _v("purchase_dtl_id", int), int(item.get("is_selected_quote") or 0),
                _d("supplier_match_conf"), _v("quote_rank", int),
                _v("supplier_name"), _v("supplier_address"), _v("supplier_country"),
                _v("quotation_ref_no"), _date("quotation_date"), _v("currency"),
                _date("validity_date"), _v("validity_days", int), _v("payment_terms"),
                _v("item_name"), _v("item_description"), _d("quantity"), _v("unit"),
                _d("unit_price"), _d("total_price"), _d("discount"), _v("taxation_details"),
                _date("delivery_date"), _v("delivery_time_days", int),
                _v("item_level_1"), _v("item_level_2"), _v("item_level_3"), _v("item_level_4"),
                _v("item_level_5"), _v("item_level_6"), _v("item_level_7"), _v("item_level_8"),
                _v("commodity_tag"), _v("item_summary"),
                unit_price_eur, total_price_eur,
                _v("purchase_category_llm"), _v("embed_content"),
                critical_attributes_json,
                # ── Commercials — quote-level (21) ────────────────────────
                _v("quote_incoterms"), _v("quote_incoterms_named_place"),
                _d("quote_subtotal"), _d("quote_discount_total"),
                _d("quote_packing_forwarding"), _d("quote_freight"), _d("quote_insurance"),
                _d("quote_customs_duties"), _d("quote_installation"), _d("quote_other_charges"),
                _v("quote_tax_type"), _d("quote_tax_rate_pct"), _d("quote_tax_amount_total"),
                _d("quote_grand_total"), _v("quote_grand_total_currency"),
                _v("quote_country_of_origin"), _v("quote_port_of_loading"),
                _v("quote_port_of_discharge"), _v("quote_mode_of_transport"),
                _v("quote_delivery_location"), _v("quote_charges_breakdown_json"),
                # ── Commercials — line-level (14) ─────────────────────────
                _v("line_incoterms"), _d("line_tax_rate_pct"), _d("line_tax_amount"),
                _d("line_freight"), _d("line_insurance"), _d("line_packing_forwarding"),
                _d("line_customs_duties"), _d("line_installation"), _d("line_other_charges"),
                _d("line_total_inclusive"), _v("line_country_of_origin"),
                _v("line_hsn_sac_code"), _v("line_charges_source"), _v("line_allocation_method"),
            )
            saved += 1
        conn.commit()
    except Exception as exc:
        logger.opt(exception=True).error(
            "V2 DB INSERT into quotation_extracted_items failed after {} row(s): {}",
            saved, exc,
        )
        try: conn.rollback()
        except Exception as rb_exc: logger.warning("V2 rollback failed: {}", rb_exc)
        raise
    finally:
        conn.close()
    return saved


# ─────────────────────────────────────────────────────────────────────────────
# V2 Stage 5b — commercials extraction (second LLM call, fail-soft)
#
# Runs AFTER the line-item extraction succeeds. Sees the parsed items, calls
# the LLM once per quotation source to extract commercial/financial fields
# (Incoterms, freight, insurance, customs, taxes, grand total, etc.), maps
# quote-level totals down to each line by proportional allocation, and
# attaches the resulting fields to each item dict in-place.
#
# Failure here NEVER blocks the existing extraction — on any exception the
# function logs a warning and returns items unchanged with the new commercial
# columns left as None.
# ─────────────────────────────────────────────────────────────────────────────


# Commercial column names that this stage owns. Used by the saver to
# pull values from item dicts; if a key is absent, the column is set to
# NULL. Keep this list as the single source of truth so adding a new
# field only requires (a) the migration, (b) the prompt, (c) this list.
_COMMERCIALS_QUOTE_FIELDS: tuple = (
    "quote_incoterms", "quote_incoterms_named_place",
    "quote_subtotal", "quote_discount_total",
    "quote_packing_forwarding", "quote_freight", "quote_insurance",
    "quote_customs_duties", "quote_installation", "quote_other_charges",
    "quote_tax_type", "quote_tax_rate_pct", "quote_tax_amount_total",
    "quote_grand_total", "quote_grand_total_currency",
    "quote_country_of_origin", "quote_port_of_loading",
    "quote_port_of_discharge", "quote_mode_of_transport",
    "quote_delivery_location", "quote_charges_breakdown_json",
)
_COMMERCIALS_LINE_FIELDS: tuple = (
    "line_incoterms", "line_tax_rate_pct", "line_tax_amount",
    "line_freight", "line_insurance", "line_packing_forwarding",
    "line_customs_duties", "line_installation", "line_other_charges",
    "line_total_inclusive", "line_country_of_origin",
    "line_hsn_sac_code", "line_charges_source", "line_allocation_method",
)
_COMMERCIALS_QUOTE_DECIMAL_FIELDS: frozenset = frozenset({
    "quote_subtotal", "quote_discount_total",
    "quote_packing_forwarding", "quote_freight", "quote_insurance",
    "quote_customs_duties", "quote_installation", "quote_other_charges",
    "quote_tax_rate_pct", "quote_tax_amount_total", "quote_grand_total",
})
_COMMERCIALS_LINE_DECIMAL_FIELDS: frozenset = frozenset({
    "line_tax_rate_pct", "line_tax_amount", "line_freight",
    "line_insurance", "line_packing_forwarding", "line_customs_duties",
    "line_installation", "line_other_charges", "line_total_inclusive",
})
_COMMERCIALS_LINE_ALLOWED_SOURCES = {"line_explicit", "quote_allocated", "mixed", "none"}
_COMMERCIALS_LINE_ALLOWED_METHODS = {"proportional_by_total", "proportional_by_qty", "equal_split"}


def _build_commercials_user_prompt(
    ctx: RASContext,
    doc: DocumentContent,
    items: list,
    prompts: dict | None = None,
) -> str:
    """Build the user prompt for the commercials LLM call.

    Feeds the document content plus a compact JSON view of the already-extracted
    line items (purchase_dtl_id + price/quantity/identity columns only) so the
    LLM can map quote-level totals to specific dtl_ids.
    """
    if doc.ocr_source and doc.text:
        doc_content_str = f"[OCR markdown from Azure Document Intelligence]\n\n{doc.text}"
    elif doc.images and doc.text:
        doc_content_str = f"[Extracted text — page images attached below]\n\n{doc.text}"
    elif doc.images:
        doc_content_str = "[Scanned document — page image(s) attached]"
    else:
        doc_content_str = doc.text or "[No content extracted]"

    ext_max = (prompts or {}).get("ext_max_chars")
    if ext_max and doc_content_str and len(doc_content_str) > ext_max:
        doc_content_str = doc_content_str[:ext_max]

    items_min: list = []
    for it in items:
        try:
            dtl = it.get("purchase_dtl_id")
            dtl_int = int(dtl) if dtl is not None else None
        except Exception:
            dtl_int = None
        if dtl_int is None:
            continue
        items_min.append({
            "purchase_dtl_id": dtl_int,
            "item_name":       it.get("item_name"),
            "quantity":        _commercials_num_for_json(it.get("quantity")),
            "unit":            it.get("unit"),
            "unit_price":      _commercials_num_for_json(it.get("unit_price")),
            "total_price":     _commercials_num_for_json(it.get("total_price")),
            "currency":        it.get("currency"),
        })
    items_json = json.dumps(items_min, ensure_ascii=False, indent=2)

    user_tmpl = (prompts or {}).get("commercials_user_v2", COMMERCIALS_USER_TEMPLATE_V2)
    return user_tmpl.format(
        document_content=doc_content_str,
        items_json=items_json,
        purchase_req_no=str(ctx.purchase_req_no or "N/A"),
        currency=str(ctx.currency or "N/A"),
        site_country=str(ctx.site_country or "N/A"),
    )


def _commercials_num_for_json(v: Any) -> Any:
    """Convert Decimal/None/numeric to a JSON-safe number for the prompt."""
    if v is None:
        return None
    if isinstance(v, (int, float)):
        return v
    if isinstance(v, Decimal):
        try:
            return float(v)
        except Exception:
            return None
    try:
        return float(str(v))
    except Exception:
        return None


def _call_commercials_llm(
    llm,
    ctx: RASContext,
    doc: DocumentContent,
    items: list,
    prompts: dict | None,
) -> str:
    """Commercials LLM call — text-only (no images) for speed.
    Commercials extraction is text-based (incoterms, freight, taxes);
    images cause 79+ second timeouts due to large base64 encoding."""
    from langchain_core.messages import HumanMessage, SystemMessage
    user_prompt = _build_commercials_user_prompt(ctx, doc, items, prompts)
    sys_prompt  = (prompts or {}).get("commercials_system_v2", COMMERCIALS_SYSTEM_PROMPT_V2)

    messages: list = [SystemMessage(content=sys_prompt)]
    # Commercials extraction: text-only, skip images to avoid Azure slowness
    img_count = 0
    img_detail = "n/a"
    if False:  # Disabled: images cause 79+ second timeouts
        max_images  = max(1, int((prompts or {}).get("ext_max_images", 50)))
    if doc.is_image_based and doc.images and False:  # Never true now
        if len(doc.images) > max_images:
            logger.warning(
                "[V2-Commercials PR={}] Quotation has {} image(s) but max_images={} — truncating to first {}",
                ctx.purchase_req_no, len(doc.images), max_images, max_images,
            )
        images = doc.images[:max_images]
        img_count  = len(images)
        img_detail = "low" if doc.text else "high"
        content_parts: list = [{"type": "text", "text": user_prompt}]
        total_img_bytes = 0
        for b64 in images:
            total_img_bytes += len(b64)
            content_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}", "detail": img_detail},
            })
        messages.append(HumanMessage(content=content_parts))
        logger.info(
            "[V2-Commercials PR={}] Images: {} total base64 bytes ({} KB per image avg)",
            ctx.purchase_req_no, total_img_bytes, total_img_bytes // max(1, img_count) // 1024,
        )
    else:
        messages.append(HumanMessage(content=user_prompt))

    # Debug: log actual message sizes
    user_prompt_bytes = len(user_prompt.encode('utf-8'))
    sys_prompt = (prompts or {}).get("commercials_system_v2", COMMERCIALS_SYSTEM_PROMPT_V2)
    sys_prompt_bytes = len(sys_prompt.encode('utf-8'))
    total_bytes = user_prompt_bytes + sys_prompt_bytes

    logger.info(
        "[V2-Commercials PR={}] Calling commercials LLM — {} image(s) detail={}, user={:.1f}KB, sys={:.1f}KB, total={:.1f}KB",
        ctx.purchase_req_no, img_count, img_detail,
        user_prompt_bytes/1024, sys_prompt_bytes/1024, total_bytes/1024,
    )
    response = _call_llm_with_retry(llm, messages, prompts, force_json=True)
    return (getattr(response, "content", None) or str(response)).strip()


def _commercials_to_decimal(v: Any) -> Optional[Decimal]:
    """Coerce a JSON number / string to Decimal; return None on failure."""
    if v is None:
        return None
    if isinstance(v, Decimal):
        return v
    try:
        s = str(v).strip()
        if not s:
            return None
        # strip currency symbols / thousands separators the LLM may slip in
        s = re.sub(r"[^\d.\-]", "", s)
        if s in ("", "-", ".", "-."):
            return None
        return Decimal(s)
    except Exception:
        return None


def _parse_commercials_response(
    raw: str, pr_no: str,
    _quote_fields: tuple = (
        "quote_incoterms", "quote_incoterms_named_place",
        "quote_subtotal", "quote_discount_total",
        "quote_packing_forwarding", "quote_freight", "quote_insurance",
        "quote_customs_duties", "quote_installation", "quote_other_charges",
        "quote_tax_type", "quote_tax_rate_pct", "quote_tax_amount_total",
        "quote_grand_total", "quote_grand_total_currency",
        "quote_country_of_origin", "quote_port_of_loading",
        "quote_port_of_discharge", "quote_mode_of_transport",
        "quote_delivery_location", "quote_charges_breakdown_json",
    ),
    _line_fields: tuple = (
        "line_incoterms", "line_tax_rate_pct", "line_tax_amount",
        "line_freight", "line_insurance", "line_packing_forwarding",
        "line_customs_duties", "line_installation", "line_other_charges",
        "line_total_inclusive", "line_country_of_origin",
        "line_hsn_sac_code", "line_charges_source", "line_allocation_method",
    ),
    _quote_decimal_fields: frozenset = frozenset({
        "quote_subtotal", "quote_discount_total",
        "quote_packing_forwarding", "quote_freight", "quote_insurance",
        "quote_customs_duties", "quote_installation", "quote_other_charges",
        "quote_tax_amount_total", "quote_grand_total",
    }),
    _line_decimal_fields: frozenset = frozenset({
        "line_tax_rate_pct", "line_tax_amount", "line_freight",
        "line_insurance", "line_packing_forwarding", "line_customs_duties",
        "line_installation", "line_other_charges", "line_total_inclusive",
    }),
    _line_allowed_sources: frozenset = frozenset({"line_explicit", "quote_allocated", "mixed", "none"}),
    _line_allowed_methods: frozenset = frozenset({"proportional_by_total", "proportional_by_qty", "equal_split"}),
) -> dict:
    """Parse the commercials LLM JSON response into a normalized dict:

        {
          "quote": { <quote-level fields>, "quote_charges_breakdown_json": str|None },
          "lines_by_dtl": { dtl_id (int): { <line-level fields> } }
        }

    Never raises — returns {"quote": {}, "lines_by_dtl": {}} on any parse error.
    """
    out: dict = {"quote": {}, "lines_by_dtl": {}}
    if not raw:
        return out
    cleaned = raw.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    try:
        envelope = json.loads(cleaned)
    except Exception as exc:
        logger.warning(f"[V2-Commercials PR={pr_no}] invalid JSON: {exc}")
        return out
    if not isinstance(envelope, dict):
        return out

    quote: dict = {}
    for f in _quote_fields:
        if f == "quote_charges_breakdown_json":
            # The prompt asks for this as an object; we serialise it for DB.
            obj = envelope.get("quote_charges_breakdown")
            if isinstance(obj, dict) and obj:
                try:
                    quote[f] = json.dumps(obj, ensure_ascii=False)
                except Exception:
                    quote[f] = None
            else:
                quote[f] = None
            continue
        v = envelope.get(f)
        if v is None or (isinstance(v, str) and not v.strip()):
            quote[f] = None
            continue
        if f in _quote_decimal_fields:
            quote[f] = _commercials_to_decimal(v)
        else:
            quote[f] = str(v).strip() if isinstance(v, str) else v
    out["quote"] = quote

    lines_raw = envelope.get("lines")
    if isinstance(lines_raw, list):
        for entry in lines_raw:
            if not isinstance(entry, dict):
                continue
            try:
                dtl = int(entry.get("purchase_dtl_id"))
            except Exception:
                continue
            row: dict = {}
            for f in _line_fields:
                v = entry.get(f)
                if v is None or (isinstance(v, str) and not v.strip()):
                    row[f] = None
                    continue
                if f in _line_decimal_fields:
                    row[f] = _commercials_to_decimal(v)
                elif f == "line_charges_source":
                    sv = str(v).strip().lower()
                    row[f] = sv if sv in _line_allowed_sources else None
                elif f == "line_allocation_method":
                    sv = str(v).strip().lower()
                    row[f] = sv if sv in _line_allowed_methods else None
                else:
                    row[f] = str(v).strip()
            out["lines_by_dtl"][dtl] = row
    return out


def _allocate_quote_charges_to_lines(parsed: dict, items: list, pr_no: str) -> None:
    """Belt-and-suspenders post-processor.

    The prompt asks the LLM to allocate quote-level charges down to lines
    itself. This function defends against the case where the LLM emits the
    quote-level number but leaves the corresponding per-line values null —
    we then allocate them in proportion to each line's total_price (fallback
    to quantity, fallback to equal split) and stamp the audit trail.

    Mutates `parsed["lines_by_dtl"]` in place. Does NOT overwrite explicit
    per-line values the LLM already produced.
    """
    quote = parsed.get("quote") or {}
    lines = parsed.get("lines_by_dtl") or {}
    if not items:
        return

    valid_dtls: list[int] = []
    totals: dict = {}     # dtl_id -> Decimal(total_price) or None
    qtys: dict = {}       # dtl_id -> Decimal(quantity)    or None
    for it in items:
        dtl_raw = it.get("purchase_dtl_id")
        try:
            dtl = int(dtl_raw) if dtl_raw is not None else None
        except Exception:
            dtl = None
        if dtl is None:
            continue
        valid_dtls.append(dtl)
        totals[dtl] = _commercials_to_decimal(it.get("total_price"))
        qtys[dtl]   = _commercials_to_decimal(it.get("quantity"))

    if not valid_dtls:
        return

    # Decide allocation basis once per charge type
    charge_pairs = (
        ("quote_packing_forwarding", "line_packing_forwarding"),
        ("quote_freight",            "line_freight"),
        ("quote_insurance",          "line_insurance"),
        ("quote_customs_duties",     "line_customs_duties"),
        ("quote_installation",       "line_installation"),
        ("quote_other_charges",      "line_other_charges"),
        ("quote_tax_amount_total",   "line_tax_amount"),
    )

    sum_totals = sum((v for v in totals.values() if v is not None), Decimal(0))
    sum_qtys   = sum((v for v in qtys.values()   if v is not None), Decimal(0))
    if sum_totals > 0:
        basis: dict = {d: (totals[d] or Decimal(0)) / sum_totals for d in valid_dtls}
        method = "proportional_by_total"
    elif sum_qtys > 0:
        basis = {d: (qtys[d] or Decimal(0)) / sum_qtys for d in valid_dtls}
        method = "proportional_by_qty"
    else:
        share = Decimal(1) / Decimal(len(valid_dtls))
        basis = {d: share for d in valid_dtls}
        method = "equal_split"

    allocated_any = False
    for q_key, l_key in charge_pairs:
        q_val = quote.get(q_key)
        if q_val is None:
            continue
        # If ANY line already has an explicit non-null value for this charge,
        # trust the LLM's distribution and skip auto-allocation (mixed mode).
        existing_explicit = any(
            (lines.get(d) or {}).get(l_key) is not None for d in valid_dtls
        )
        if existing_explicit:
            continue
        for d in valid_dtls:
            row = lines.setdefault(d, {})
            row[l_key] = (q_val * basis[d]).quantize(Decimal("0.0001"))
            # mark this line as receiving allocated charges
            cur_src = row.get("line_charges_source")
            if cur_src in (None, "none"):
                row["line_charges_source"] = "quote_allocated"
            elif cur_src == "line_explicit":
                row["line_charges_source"] = "mixed"
            row.setdefault("line_allocation_method", method)
        allocated_any = True

    if allocated_any:
        logger.info(
            "[V2-Commercials PR={}] Allocated quote-level charges to {} line(s) via {}",
            pr_no, len(valid_dtls), method,
        )

    # If line_total_inclusive is missing for any line, compute it from the
    # available per-line components. Skip silently when total_price is None.
    for d in valid_dtls:
        row = lines.setdefault(d, {})
        if row.get("line_total_inclusive") is not None:
            continue
        line_total = totals.get(d)
        if line_total is None:
            continue
        components = [
            row.get("line_tax_amount"), row.get("line_freight"),
            row.get("line_insurance"), row.get("line_packing_forwarding"),
            row.get("line_customs_duties"), row.get("line_installation"),
            row.get("line_other_charges"),
        ]
        acc = line_total
        for c in components:
            if c is not None:
                acc = acc + c
        try:
            row["line_total_inclusive"] = Decimal(acc).quantize(Decimal("0.0001"))
        except Exception:
            pass


def _attach_commercials_to_items(
    parsed: dict, items: list,
    _quote_fields: tuple = (
        "quote_incoterms", "quote_incoterms_named_place",
        "quote_subtotal", "quote_discount_total",
        "quote_packing_forwarding", "quote_freight", "quote_insurance",
        "quote_customs_duties", "quote_installation", "quote_other_charges",
        "quote_tax_type", "quote_tax_rate_pct", "quote_tax_amount_total",
        "quote_grand_total", "quote_grand_total_currency",
        "quote_country_of_origin", "quote_port_of_loading",
        "quote_port_of_discharge", "quote_mode_of_transport",
        "quote_delivery_location", "quote_charges_breakdown_json",
    ),
    _line_fields: tuple = (
        "line_incoterms", "line_tax_rate_pct", "line_tax_amount",
        "line_freight", "line_insurance", "line_packing_forwarding",
        "line_customs_duties", "line_installation", "line_other_charges",
        "line_total_inclusive", "line_country_of_origin",
        "line_hsn_sac_code", "line_charges_source", "line_allocation_method",
    ),
) -> None:
    """Stamp the parsed commercial fields onto each item dict in-place.

    Quote-level fields are repeated on every item from the same quotation
    (matches how supplier_name / currency are already stored). Line-level
    fields are looked up by purchase_dtl_id; lines the LLM did not return
    for get None values.
    """
    quote = parsed.get("quote") or {}
    lines = parsed.get("lines_by_dtl") or {}
    for it in items:
        for f in _quote_fields:
            it[f] = quote.get(f)
        dtl_raw = it.get("purchase_dtl_id")
        try:
            dtl = int(dtl_raw) if dtl_raw is not None else None
        except Exception:
            dtl = None
        row = lines.get(dtl) if dtl is not None else None
        for f in _line_fields:
            it[f] = (row or {}).get(f)


def _enrich_with_commercials(
    llm, ctx: RASContext, doc: DocumentContent,
    items: list, prompts: dict | None,
) -> None:
    """Top-level entry point for the second LLM call. Mutates `items` in place
    with the new commercial fields. Fail-soft — never raises."""
    if not items:
        return
    try:
        raw = _call_commercials_llm(llm, ctx, doc, items, prompts)
        parsed = _parse_commercials_response(raw, str(ctx.purchase_req_no))
        _allocate_quote_charges_to_lines(parsed, items, str(ctx.purchase_req_no))
        _attach_commercials_to_items(parsed, items)
        logger.info(
            "[V2-Commercials PR={}] enriched {} item(s) — grand_total={}, incoterms={}",
            ctx.purchase_req_no, len(items),
            (parsed.get("quote") or {}).get("quote_grand_total"),
            (parsed.get("quote") or {}).get("quote_incoterms"),
        )
    except Exception as exc:
        logger.opt(exception=True).warning(
            "[V2-Commercials PR={}] commercials extraction failed (non-fatal): {}",
            ctx.purchase_req_no, exc,
        )


def _run_extraction_v2(llm, tgt_cs: str, blob_cfg: dict, pr_no: str, prompts: dict | None = None) -> int:
    """V2 extraction orchestrator — same shape as _run_extraction but uses
    V2 prompts, parser, and saver. Raises the same ExtractionAbortError
    subclasses so the orchestrator's exception handling is unchanged."""
    import os as _os
    ctx = _build_ras_context(tgt_cs, pr_no)
    if ctx is None:
        raise NoRASContextError(
            f"No purchase_req_mst row for PR={pr_no!r} — cannot extract."
        )
    if not ctx.line_items:
        raise NoLineItemsError(
            f"PR={pr_no!r} has no rows in purchase_req_detail — nothing to extract against."
        )
    sources = _resolve_quotation_sources(tgt_cs, pr_no)
    if not sources:
        raise NoQuotationFoundError(
            f"No quotation documents found for PR={pr_no!r}."
        )

    category_list = _resolve_category_list(tgt_cs)
    n_workers = max(1, int((prompts or {}).get("ext_parallel_sources", 1)))
    max_pages = max(1, int((prompts or {}).get("ext_max_pages", 20)))
    enable_commercials = bool((prompts or {}).get("enable_commercials_extraction", True))

    def _extract_one(src: dict) -> list[dict]:
        blob_path = src["blob_path"]
        filename  = _os.path.basename(blob_path)
        try:
            file_bytes = _download_blob(blob_path, blob_cfg)
            doc        = _load_document(file_bytes, filename, max_pages=max_pages)
            raw        = _call_extraction_llm_v2(llm, ctx, doc, prompts, category_list)
            items      = _parse_extraction_response_v2(raw, src, ctx, category_list)
            items      = _align_to_ras_line_items(items, ctx, src)
            logger.info(
                "[V2 {}] Extracted {} item(s) from {!r}", pr_no, len(items), filename
            )
            # Stage 5b — commercials extraction (fail-soft, second LLM call).
            # Runs in the same per-source worker thread, so it inherits the
            # existing ext_parallel_sources concurrency profile. Any exception
            # is swallowed inside _enrich_with_commercials; the line-item
            # rows are still saved even if the new columns end up NULL.
            if enable_commercials and items:
                _enrich_with_commercials(llm, ctx, doc, items, prompts)
            return items
        except Exception as exc:
            logger.opt(exception=True).warning(
                "[V2 {}] Extraction failed for {!r}: {}", pr_no, filename, exc
            )
            return []

    if n_workers == 1:
        all_items: list[dict] = []
        for src in sources:
            all_items.extend(_extract_one(src))
    else:
        import concurrent.futures
        logger.info(f"[V2 {pr_no}] Extracting {len(sources)} source(s) with {n_workers} parallel worker(s)")
        with concurrent.futures.ThreadPoolExecutor(max_workers=n_workers) as pool:
            results = list(pool.map(_extract_one, sources))
        all_items = [item for batch in results for item in batch]

    if not all_items:
        raise AllExtractionsFailedError(
            f"All {len(sources)} quotation source(s) for PR={pr_no!r} failed extraction."
        )

    _apply_price_alignment_boost(all_items, ctx, prompts)
    _canonicalize_supplier_names(all_items, ctx, prompts)
    _compute_quote_ranks(all_items)
    _select_best_quotes(all_items, ctx)
    fallback_date = ctx.req_start_date or ctx.c_datetime
    saved = _save_extracted_items_v2(tgt_cs, all_items, fallback_date=fallback_date)
    logger.info(f"[V2 {pr_no}] {saved} item(s) written to quotation_extracted_items")
    return saved


# ─────────────────────────────────────────────────────────────────────────────
# V2 Stage 6 — embedding with embed_content + richer Pinecone metadata
# ─────────────────────────────────────────────────────────────────────────────


def _run_embeddings_v2(tgt_cs: str, pr_no: str, embed_model, pinecone_index: str, pinecone_ns: str) -> None:
    """Embeds using embed_content (LLM-crafted, brand-free) instead of the
    12-field concat. Pinecone metadata gains purchase_category_llm +
    item_level_1/2 so Stage A's $eq filter is possible."""
    from agentcore.services.pinecone_service_client import ensure_index_via_service, ingest_via_service
    _pinecone_ensure_with_retry(
        ensure_index_via_service,
        index_name=pinecone_index,
        embedding_dimension=3072,
    )

    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8], qi.[commodity_tag],
                   qi.[purchase_category_llm], qi.[embed_content],
                   prd.[C_DATETIME] AS [item_created_date]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        all_rows = cur.fetchall()
    finally:
        conn.close()

    cols = [
        "extracted_item_uuid_pk", "purchase_dtl_id",
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag", "purchase_category_llm", "embed_content",
        "item_created_date",
    ]
    seen: dict = {}
    for row in all_rows:
        rd = dict(zip(cols, row))
        dtl_id = rd["purchase_dtl_id"]
        if dtl_id not in seen:
            seen[dtl_id] = rd
    rows = list(seen.values())
    logger.info(
        f"[V2 {pr_no}] Embedding {len(rows)} selected item(s) (deduped from {len(all_rows)})"
    )

    for rd in rows:
        dtl_id    = rd["purchase_dtl_id"]
        item_uuid = rd["extracted_item_uuid_pk"]
        created   = rd.get("item_created_date")
        created_iso = (
            created.isoformat() if created and hasattr(created, "isoformat")
            else str(created or "")
        )
        content = (rd.get("embed_content") or "").strip()
        if not content:
            content = _build_embed_text(rd)
            logger.warning(
                f"[V2 {pr_no}] dtl_id={dtl_id}: no embed_content — fallback to 12-field concat"
            )
        if not content:
            logger.warning(f"[V2 {pr_no}] dtl_id={dtl_id}: no text to embed, skipping")
            continue
        try:
            embedding = _embed_with_timeout(embed_model, content)
            _pinecone_ingest_with_timeout(
                ingest_via_service,
                index_name=pinecone_index,
                namespace=pinecone_ns,
                text_key="page_content",
                documents=[{
                    "page_content": content,
                    "metadata": {
                        "purchase_req_no":        pr_no,
                        "purchase_dtl_id":        int(dtl_id),
                        "extracted_item_uuid_pk": str(item_uuid or ""),
                        "commodity_tag":          str(rd.get("commodity_tag") or ""),
                        "item_created_date":      created_iso,
                        "purchase_category_llm":  str(rd.get("purchase_category_llm") or ""),
                        "item_level_1":           str(rd.get("item_level_1") or ""),
                        "item_level_2":           str(rd.get("item_level_2") or ""),
                    },
                }],
                embedding_vectors=[embedding],
                vector_ids=[f"dtl_{dtl_id}"],
                embedding_dimension=3072,
            )
            logger.info(f"[V2 {pr_no}] Upserted vector dtl_{dtl_id} (item_created={created_iso})")
        except Exception as exc:
            logger.warning(f"[V2 {pr_no}] Embedding failed for dtl_id {dtl_id}: {exc}")


# ─────────────────────────────────────────────────────────────────────────────
# V2 Stage 7 — 3-stage retrieval (SQL → Pinecone within pool → LLM rank)
# ─────────────────────────────────────────────────────────────────────────────


def _sql_pool_by_category(
    tgt_cs: str,
    src_category: str,
    src_created_date: Any,
    pool_size: int,
    exclude_dtl_ids: "list[int] | None" = None,
) -> list[int]:
    """Stage A — ALL historical dtl_ids in the same V2 category, older than
    source, ordered by C_DATETIME DESC (newest first). Empty if no peers.

    The previous TOP-N cap (default 100 most recent) was removed because it
    silently dropped older but still-relevant comparables from large
    categories — e.g. for an 1800T moulding machine source, a 2300T historical
    purchase from 18 months back could be excluded just because 100 newer
    items existed. Now every category-matching, older, priced peer is
    returned and the LLM re-ranker decides relevance. The `pool_size`
    parameter is kept for signature compatibility but is no longer applied.
    exclude_dtl_ids: dtl_ids from the current PR — always skipped.
    """
    if not src_category:
        return []
    # Expand to include legacy DB names so a search for "Moulding Machine"
    # also matches historical rows still classified as "IMM".
    cat_names = _category_with_aliases(src_category)
    if not cat_names:
        return []
    cat_placeholders = ",".join("?" * len(cat_names))
    excl = [int(x) for x in (exclude_dtl_ids or []) if x is not None]
    excl_clause = f"AND qi.[purchase_dtl_id] NOT IN ({','.join('?' * len(excl))})" if excl else ""
    conn = _connect(tgt_cs)
    try:
        cur = conn.cursor()
        cur.execute(f"""
            SELECT qi.[purchase_dtl_id]
              FROM [ras_procurement].[quotation_extracted_items] qi
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE qi.[is_selected_quote] = 1
               AND qi.[purchase_category_llm] IN ({cat_placeholders})
               AND prd.[C_DATETIME] < ?
               AND qi.[unit_price_eur] IS NOT NULL
               {excl_clause}
             ORDER BY prd.[C_DATETIME] DESC
        """, *cat_names, src_created_date, *excl)
        return [int(r[0]) for r in cur.fetchall()]
    finally:
        conn.close()


def _widen_pool_by_levels(
    tgt_cs: str,
    l1: str,
    l2: str,
    src_created_date: Any,
    pool_size: int,
    exclude_dtl_ids: "list[int] | None" = None,
) -> list[int]:
    """Sparse-pool fallback — accepts rows matching item_level_1 + item_level_2
    instead of LLM category. The TOP-N cap was removed for the same reason
    as `_sql_pool_by_category`; the `pool_size` parameter is retained for
    signature compatibility but no longer applied.
    exclude_dtl_ids: dtl_ids from the current PR — always skipped."""
    if not (l1 and l2):
        return []
    excl = [int(x) for x in (exclude_dtl_ids or []) if x is not None]
    excl_clause = f"AND qi.[purchase_dtl_id] NOT IN ({','.join('?' * len(excl))})" if excl else ""
    conn = _connect(tgt_cs)
    try:
        cur = conn.cursor()
        cur.execute(f"""
            SELECT qi.[purchase_dtl_id]
              FROM [ras_procurement].[quotation_extracted_items] qi
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE qi.[is_selected_quote] = 1
               AND qi.[item_level_1] = ?
               AND qi.[item_level_2] = ?
               AND prd.[C_DATETIME] < ?
               AND qi.[unit_price_eur] IS NOT NULL
               {excl_clause}
             ORDER BY prd.[C_DATETIME] DESC
        """, l1, l2, src_created_date, *excl)
        return [int(r[0]) for r in cur.fetchall()]
    finally:
        conn.close()


def _pinecone_narrow_within_pool(
    embedding: list[float],
    candidate_dtl_ids: list[int],
    top_k: int,
    pinecone_index: str,
    pinecone_ns: str,
    min_score: float = 0.70,
    query_text: str = "",
) -> list[dict]:
    """Stage B — Pinecone similarity search with optional pool restriction.
    Fetches top_k * 5 from Pinecone (filter kwarg not supported by
    search_via_service), applies min_score threshold, and optionally restricts
    to candidate_dtl_ids pool. Pass an empty list to search globally."""
    pool_set = set(candidate_dtl_ids)  # empty = no pool restriction
    from agentcore.services.pinecone_service_client import search_via_service
    # search_via_service requires query to be at least 1 char
    query = query_text.strip() or "item"
    try:
        raw = _pinecone_search_with_timeout(
            search_via_service,
            index_name=pinecone_index,
            namespace=pinecone_ns,
            text_key="page_content",
            query=query,
            query_embedding=embedding,
            number_of_results=top_k * 5,
        )
    except Exception as exc:
        if _is_pinecone_server_error(exc):
            raise
        logger.warning(f"V2 Pinecone within-pool search failed (non-server): {exc}")
        return []
    matches = raw if isinstance(raw, list) else (
        raw.get("matches", []) or raw.get("results", []) if isinstance(raw, dict) else []
    )
    filtered = []
    for m in matches:
        if float(m.get("score", 0.0)) < min_score:
            continue
        meta = (m.get("metadata") if isinstance(m, dict) else {}) or {}
        did = meta.get("purchase_dtl_id")
        try:
            did_int = int(did) if did is not None else None
            if did_int is None:
                continue
            # If pool_set is non-empty, restrict to pool; otherwise accept any match
            if pool_set and did_int not in pool_set:
                continue
            filtered.append(m)
        except Exception:
            continue
        if len(filtered) >= top_k:
            break
    return filtered


def _fetch_candidate_snapshots(tgt_cs: str, dtl_ids: list[int]) -> list[dict]:
    """Pull full extraction data for each candidate from quotation_extracted_items
    so the rank LLM can compare specs, category, and price against the source item."""
    if not dtl_ids:
        return []
    placeholders = ",".join("?" * len(dtl_ids))
    conn = _connect(tgt_cs)
    try:
        cur = conn.cursor()
        cur.execute(f"""
            SELECT qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8],
                   qi.[purchase_category_llm], qi.[commodity_tag],
                   qi.[critical_attributes],
                   qi.[supplier_name], qi.[supplier_country],
                   qi.[quotation_date], qi.[currency],
                   qi.[quantity], qi.[unit],
                   qi.[unit_price], qi.[unit_price_eur],
                   qi.[total_price], qi.[total_price_eur],
                   qi.[discount], qi.[payment_terms], qi.[delivery_time_days],
                   prd.[C_DATETIME] AS [item_created_date]
              FROM [ras_procurement].[quotation_extracted_items] qi
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE qi.[purchase_dtl_id] IN ({placeholders})
               AND qi.[is_selected_quote] = 1
        """, *dtl_ids)
        cols = [
            "purchase_dtl_id",
            "item_name", "item_description", "item_summary",
            "item_level_1", "item_level_2", "item_level_3", "item_level_4",
            "item_level_5", "item_level_6", "item_level_7", "item_level_8",
            "purchase_category_llm", "commodity_tag", "critical_attributes",
            "supplier_name", "supplier_country",
            "quotation_date", "currency",
            "quantity", "unit",
            "unit_price", "unit_price_eur",
            "total_price", "total_price_eur",
            "discount", "payment_terms", "delivery_time_days",
            "item_created_date",
        ]
        rows = []
        for r in cur.fetchall():
            rd = dict(zip(cols, r))
            # Parse critical_attributes JSON
            ca = rd.get("critical_attributes")
            if isinstance(ca, str) and ca.strip():
                try: rd["critical_attributes"] = json.loads(ca)
                except Exception: rd["critical_attributes"] = []
            else:
                rd["critical_attributes"] = []
            # Build category string from levels
            cat_parts = [rd.get(f"item_level_{n}") for n in range(1, 9) if rd.get(f"item_level_{n}")]
            rd["category"] = " > ".join(str(p) for p in cat_parts) if cat_parts else None
            # Coerce numerics
            for k in ("unit_price_eur", "unit_price", "total_price_eur", "total_price", "quantity", "discount"):
                v = rd.get(k)
                if v is not None:
                    try: rd[k] = float(v)
                    except Exception: rd[k] = None
            # Coerce dates
            for dk in ("item_created_date", "quotation_date"):
                dv = rd.get(dk)
                if dv and hasattr(dv, "isoformat"):
                    rd[dk] = dv.isoformat()
            rows.append(rd)
        return rows
    finally:
        conn.close()


def _pre_filter_candidates(source: dict, candidates: list[dict], prompts: dict | None) -> list[dict]:
    """Drop structurally incompatible candidates before LLM ranking.
    Catches obvious outliers (extreme price ratios) that would waste LLM tokens.
    Does NOT filter by category or item_level — those vary intentionally within
    purchase_category_llm and should be decided by the LLM.

    Default ratio bumped from 10× to 50× — the LLM has its own spec-to-price
    sanity rule (Part A rule 3) that catches category-drift outliers, so this
    pre-filter is now a coarse last-resort guard for genuinely broken data
    (mis-extracted prices, currency mishaps, accidental zero-padded values)
    rather than a tight first gate.
    """
    if not candidates:
        return []
    p = prompts or {}
    max_ratio = float(p.get("bench_max_price_ratio", 50.0))
    src_price = source.get("unit_price_eur")
    if not src_price:
        return candidates
    try:
        src_price = float(src_price)
    except (ValueError, TypeError):
        return candidates
    if src_price <= 0:
        return candidates
    filtered = []
    for c in candidates:
        cand_price = c.get("unit_price_eur")
        if not cand_price:
            filtered.append(c)
            continue
        try:
            cand_price = float(cand_price)
        except (ValueError, TypeError):
            filtered.append(c)
            continue
        if cand_price <= 0:
            filtered.append(c)
            continue
        ratio = max(src_price, cand_price) / min(src_price, cand_price)
        if ratio <= max_ratio:
            filtered.append(c)
    return filtered


def _build_short_reject_reasoning(verdict_row: dict, crit_pct: float = 10.0) -> str:
    """Template a 1-sentence buyer-readable rejection from the compact
    verdict-row fields. Used when the LLM emits a verdict=FAIL row in the
    new two-array schema; the post-processor synthesizes the prose so the
    LLM does not have to write 5–8 sentences for every reject.

    The helper is defensive about LLM mis-labeling: if `rule` says
    CRITICAL_BARRIER but the delta_pct is actually INSIDE the strict
    band AND a note is present, we treat the note as the real reason
    instead of claiming "exceeds strict band" (which would contradict
    the data). This handles the common LLM error of labeling a
    family-boundary rejection as CRITICAL_BARRIER on tonnage.
    """
    rule  = verdict_row.get("rule")  or "OTHER"
    attr  = verdict_row.get("attr")
    src   = verdict_row.get("src")
    cand  = verdict_row.get("cand")
    unit  = verdict_row.get("unit")
    delta = verdict_row.get("delta_pct")
    note  = (verdict_row.get("note") or "").strip()
    unit_str = unit or ""

    # Treat the numeric delta as "within band" only when we have a number
    # AND it's at or under the configured critical threshold. Categorical
    # rules (where delta is None) skip this branch.
    delta_is_within_band = (
        isinstance(delta, (int, float)) and abs(float(delta)) <= float(crit_pct)
    )

    if rule == "CRITICAL_BARRIER":
        if delta is not None and src is not None and cand is not None:
            if delta_is_within_band and note:
                # LLM labeled CRITICAL_BARRIER but the numeric delta is
                # within band — the note carries the actual reason.
                # Word the sentence so it doesn't claim the band was
                # exceeded.
                return (
                    f"REJECT — {attr} {cand}{unit_str} vs source {src}{unit_str} "
                    f"({delta}% — within ±{crit_pct}% band) but rejected: {note}."
                )
            if delta_is_within_band and not note:
                # Defensive fallback: LLM said CRITICAL_BARRIER on an
                # attribute that's actually inside band and gave no
                # justification. Surface that as a labeling concern
                # rather than echoing the contradictory claim.
                return (
                    f"REJECT — {attr} {cand}{unit_str} vs source {src}{unit_str} "
                    f"({delta}% — within ±{crit_pct}% band; LLM marked CRITICAL_BARRIER "
                    f"without a stated family/secondary reason)."
                )
            # Normal case: delta truly exceeds band.
            base = (
                f"REJECT — {attr} {cand}{unit_str} vs source {src}{unit_str} "
                f"({delta}% delta exceeds ±{crit_pct}% strict band)."
            )
            return f"{base} BUT {note}." if note else base
        # Categorical CRITICAL_BARRIER (no numeric delta — e.g. drive_type).
        if attr is not None and src is not None and cand is not None:
            base = (
                f"REJECT — {attr}: source={src!r} vs candidate={cand!r} "
                f"(family/categorical mismatch on a critical attribute)."
            )
            return f"{base} BUT {note}." if note else base
        return "REJECT — critical attribute barrier."

    if rule == "VARIANT_MISMATCH":
        # Variant-boundary rejection per Part B playbook (drive type,
        # injection units, fuel type, etc.). Lead with the variant fact.
        if attr is not None and src is not None and cand is not None:
            base = (
                f"REJECT — variant mismatch on {attr} "
                f"(source={src!r}, candidate={cand!r}); different machine "
                f"variant per category playbook."
            )
        elif note:
            base = f"REJECT — variant mismatch: {note}"
            return f"{base}."
        else:
            base = "REJECT — variant mismatch (rule: VARIANT_MISMATCH)."
        return f"{base} {note}." if note else base

    if rule == "PRODUCT_TYPE":
        base = (
            f"REJECT — product type mismatch ({attr or 'item_level'}: "
            f"{src!r} vs {cand!r})."
        )
        return f"{base} BUT {note}." if note else base

    if rule == "RATIO":
        base = (
            f"REJECT — spec-to-price ratio outlier on {attr or 'price'} "
            f"(value {cand} vs pool reference {src}, beyond sanity band)."
        )
        return f"{base} BUT {note}." if note else base

    if rule == "ABSENCE":
        base = (
            f"REJECT — candidate carries 3+ domain-specific attributes the "
            f"source lacks (absence principle on {attr or 'attribute set'})."
        )
        return f"{base} BUT {note}." if note else base

    # OTHER or unknown rule.
    if attr is not None and src is not None and cand is not None:
        base = f"REJECT — {attr}: source={src!r}, candidate={cand!r}."
        return f"{base} BUT {note}." if note else base
    if note:
        return f"REJECT — {note}."
    return "REJECT — outside acceptable bounds (rule: OTHER)."


def _extract_token_usage(response, llm) -> dict:
    """Pull input/output token counts and finish-reason from an LLM response
    across providers (Azure OpenAI, OpenAI, Anthropic). Returns a flat dict
    with values defaulting to None if a provider does not expose them.

    Used to diagnose silent-truncation failures: when finish_reason is
    'length' (OpenAI) or 'max_tokens' (Anthropic), the model hit its output
    cap and stopped mid-JSON. That is the smoking gun for the LLM_DROPPED
    candidate symptom.
    """
    md = getattr(response, "response_metadata", {}) or {}
    um = getattr(response, "usage_metadata",   {}) or {}

    # usage_metadata = newer langchain unified schema (preferred).
    input_tokens  = um.get("input_tokens")
    output_tokens = um.get("output_tokens")
    total_tokens  = um.get("total_tokens")

    # Fall back to provider-specific keys when usage_metadata is empty.
    if input_tokens is None:
        tu = md.get("token_usage") or md.get("usage") or {}
        input_tokens  = tu.get("prompt_tokens")     or tu.get("input_tokens")
        output_tokens = tu.get("completion_tokens") or tu.get("output_tokens")
        total_tokens  = tu.get("total_tokens")

    # finish_reason ('length') = OpenAI/Azure truncation;
    # stop_reason ('max_tokens') = Anthropic truncation.
    finish_reason = md.get("finish_reason") or md.get("stop_reason")
    truncated     = finish_reason in ("length", "max_tokens", "MAX_TOKENS")

    # Best-effort read of the model's configured max output tokens.
    max_output_tokens = None
    for attr in ("max_tokens", "max_completion_tokens", "max_new_tokens"):
        v = getattr(llm, attr, None)
        if v:
            try:
                max_output_tokens = int(v); break
            except (TypeError, ValueError):
                continue
    if max_output_tokens is None:
        mk = getattr(llm, "model_kwargs", None) or {}
        if isinstance(mk, dict):
            v = mk.get("max_tokens") or mk.get("max_completion_tokens")
            if v:
                try: max_output_tokens = int(v)
                except (TypeError, ValueError): pass

    tokens_remaining = None
    if max_output_tokens is not None and output_tokens is not None:
        try:
            tokens_remaining = int(max_output_tokens) - int(output_tokens)
        except (TypeError, ValueError):
            tokens_remaining = None

    return {
        "input_tokens":       input_tokens,
        "output_tokens":      output_tokens,
        "total_tokens":       total_tokens,
        "max_output_tokens":  max_output_tokens,
        "tokens_remaining":   tokens_remaining,
        "finish_reason":      finish_reason,
        "truncated":          truncated,
    }


def _llm_rank_candidates(
    llm,
    source_item: dict,
    candidates: list[dict],
    top_k: int,
    prompts: dict | None,
) -> tuple[list[dict], list[dict], dict]:
    """Stage C — Ask the LLM to rank candidates by relevance. Returns
    (selected, rejected, agent_reasoning). On failure returns ([], [], {})
    so the caller falls back to Pinecone order.

    `agent_reasoning` is a structured dict the caller will serialize into the
    benchmark_result.agent_reasoning column for interpretability:
      {
        "category_playbook": "<playbook applied at PR level>",
        "selected": [ <rich per-DTL row exactly as the LLM emitted> ],
        "rejected": [ <one entry per rejected candidate with rule+reason> ]
      }
    """
    if not candidates:
        return [], [], {}
    from langchain_core.messages import HumanMessage, SystemMessage
    p = prompts or {}
    sys_tmpl  = p.get("bench_rank_prompt_system_v2", RANK_PROMPT_SYSTEM_V2)
    user_tmpl = p.get("bench_rank_prompt_user_v2",   RANK_PROMPT_USER_V2)
    crit_pct  = int(p.get("bench_critical_threshold_pct",  10))
    imp_pct   = int(p.get("bench_important_threshold_pct", 20))
    ratio_pct = int(p.get("bench_ratio_band_pct",          50))

    # Determinism: present candidates to the LLM in a stable, content-
    # independent order so two runs over the same SOURCE + CANDIDATES set
    # produce the same prompt bytes. Candidates without purchase_dtl_id
    # (shouldn't happen in practice) sort to the end.
    def _cand_key(c: dict) -> tuple:
        did = c.get("purchase_dtl_id")
        try:
            return (0, int(did))
        except (TypeError, ValueError):
            return (1, str(did))
    candidates = sorted(candidates, key=_cand_key)

    sys_prompt  = sys_tmpl.format(
        critical_pct=crit_pct, important_pct=imp_pct, ratio_pct=ratio_pct
    )
    # Note: {top_k} is no longer in the user prompt template — the cap is
    # hard-coded as "20 to 25" in the prompt text directly. The earlier
    # bench_llm_shortlist_size knob was leaking into the prompt as the
    # caller's preferred cap (e.g. 10), which the LLM read as "emit only
    # 10 verdicts" and ignored the later top-25 instruction. Hard-coding
    # the range in the prompt removes that contradiction.
    user_prompt = user_tmpl.format(
        n_candidates=len(candidates),
        source_json=json.dumps(source_item, ensure_ascii=False, default=str),
        candidates_json=json.dumps(candidates, ensure_ascii=False, default=str),
        critical_pct=crit_pct, important_pct=imp_pct, ratio_pct=ratio_pct,
    )

    # ── Optional debug dump of the EXACT bytes sent to / received from the
    # rank LLM. Off by default. Enable by setting:
    #     prompts["bench_rank_debug_dump_dir"] = "<writable dir path>"
    # When set, two files are written per source DTL per call:
    #     rank_dtl_<id>_<ts>_in.txt   — full system + user prompts (the
    #                                    candidates_json is already embedded)
    #     rank_dtl_<id>_<ts>_out.txt  — raw response.content before any
    #                                    JSON parsing, fence-stripping, or
    #                                    schema adaptation
    # Failures are non-fatal — pipeline keeps running, only the dump is lost.
    dump_dir   = (p.get("bench_rank_debug_dump_dir") or "").strip()
    dump_in_path:  "str | None" = None
    dump_out_path: "str | None" = None
    if dump_dir:
        try:
            import os as _os
            from datetime import datetime as _dt
            _os.makedirs(dump_dir, exist_ok=True)
            src_dtl = source_item.get("purchase_dtl_id", "unknown")
            stamp   = _dt.now().strftime("%Y%m%d_%H%M%S_%f")
            dump_in_path  = _os.path.join(dump_dir, f"rank_dtl_{src_dtl}_{stamp}_in.txt")
            dump_out_path = _os.path.join(dump_dir, f"rank_dtl_{src_dtl}_{stamp}_out.txt")
            # Write the EXACT two messages that go to the LLM, nothing else.
            # The only addition is the two role markers needed to distinguish
            # the SystemMessage from the HumanMessage payload — without them
            # a reader cannot tell where one ends and the other begins.
            with open(dump_in_path, "w", encoding="utf-8") as _f:
                _f.write("[SYSTEM]\n")
                _f.write(sys_prompt)
                _f.write("\n[USER]\n")
                _f.write(user_prompt)
            logger.info(
                f"[V2 rank LLM] DEBUG: input prompts dumped to {dump_in_path}"
            )
        except Exception as exc:
            logger.warning(
                f"[V2 rank LLM] DEBUG: failed to write input dump to "
                f"{dump_dir!r}: {exc}"
            )
            dump_in_path = dump_out_path = None

    # Determinism: bind temperature=0 + a fixed seed for this ranking call,
    # without mutating the caller's llm object. Most Azure / OpenAI chat
    # models honor both; older models or non-OpenAI backends ignore unknown
    # bind kwargs. If .bind() itself fails (very old wrappers), fall back
    # to the unbound llm so the rank call still runs.
    rank_llm = llm
    try:
        rank_llm = llm.bind(
            temperature=int(p.get("bench_rank_temperature", 0)),
            seed=int(p.get("bench_rank_seed", 42)),
        )
    except Exception:
        rank_llm = llm

    try:
        response = _call_llm_with_retry(
            rank_llm, [SystemMessage(content=sys_prompt), HumanMessage(content=user_prompt)],
            prompts, force_json=True,
        )

        # ── DEBUG: dump raw LLM response before any parsing ────────────────
        # Writes ONLY response.content — the exact text the model returned.
        # No fence-stripping, no JSON parse, no metadata wrapping, nothing.
        # Only writes when bench_rank_debug_dump_dir was set AND the input
        # dump succeeded (dump_out_path is non-None).
        if dump_out_path:
            try:
                raw_content_for_dump = (
                    getattr(response, "content", None) or str(response)
                )
                with open(dump_out_path, "w", encoding="utf-8") as _f:
                    _f.write(raw_content_for_dump)
                logger.info(
                    f"[V2 rank LLM] DEBUG: raw output dumped to {dump_out_path}"
                )
            except Exception as _exc:
                logger.warning(
                    f"[V2 rank LLM] DEBUG: failed to write output dump to "
                    f"{dump_out_path!r}: {_exc}"
                )

        # ── Token usage diagnostics ─────────────────────────────────────
        # Pull token counts + finish_reason so we can see exactly how much
        # of the model's output budget was consumed, and detect the silent-
        # drop root cause (truncation when finish_reason == 'length' or
        # 'max_tokens'). Read max_tokens from the ORIGINAL llm, not the
        # bound copy, in case .bind() hid the attribute.
        usage = _extract_token_usage(response, llm)
        in_tok   = usage["input_tokens"]
        out_tok  = usage["output_tokens"]
        max_out  = usage["max_output_tokens"]
        remain   = usage["tokens_remaining"]
        finish   = usage["finish_reason"]
        usage_pct = None
        if max_out and out_tok is not None:
            try:
                usage_pct = round(100.0 * float(out_tok) / float(max_out), 1)
            except Exception:
                usage_pct = None
        logger.info(
            "[V2 rank LLM token usage] n_candidates={ncand} "
            "input_tokens={inp} output_tokens={out} "
            "max_output_tokens={mx} tokens_remaining={rem} "
            "output_used_pct={pct} finish_reason={fin!r}".format(
                ncand=len(candidates),
                inp=in_tok if in_tok is not None else "?",
                out=out_tok if out_tok is not None else "?",
                mx=max_out if max_out is not None else "?",
                rem=remain if remain is not None else "?",
                pct=usage_pct if usage_pct is not None else "?",
                fin=finish,
            )
        )
        if usage["truncated"]:
            logger.warning(
                "[V2 rank LLM] !!! TRUNCATION DETECTED !!! finish_reason={fin!r}, "
                "output_tokens={out}/{mx} ({pct}%), tokens_remaining={rem}. "
                "The model hit its output cap before writing verdicts for every "
                "candidate. This is the direct cause of LLM_DROPPED rejects in "
                "this run. Mitigations: (1) raise max_output_tokens on the LLM "
                "client, (2) reduce the candidate count sent to the LLM, "
                "(3) shorten per-rejection reasoning in the prompt.".format(
                    fin=finish,
                    out=out_tok if out_tok is not None else "?",
                    mx=max_out if max_out is not None else "?",
                    pct=usage_pct if usage_pct is not None else "?",
                    rem=remain if remain is not None else "?",
                )
            )
        else:
            logger.info(
                "[V2 rank LLM] no truncation (finish_reason={fin!r}); "
                "output budget used {pct}% ({out}/{mx} tokens).".format(
                    fin=finish,
                    pct=usage_pct if usage_pct is not None else "?",
                    out=out_tok if out_tok is not None else "?",
                    mx=max_out if max_out is not None else "?",
                )
            )

        raw = (getattr(response, "content", None) or str(response)).strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        try:
            parsed = json.loads(raw)
        except json.JSONDecodeError as je:
            # Truncation broke the JSON. Salvage: walk through the partial
            # text and extract every COMPLETE balanced `{ ... }` object that
            # looks like a verdict row, so PASS verdict dtl_ids are still
            # captured and persisted in `similar_dtl_ids` instead of being
            # silently lost.
            logger.warning(
                f"V2 rank LLM JSON parse failed (likely due to truncation): {je}. "
                f"Attempting partial-JSON salvage to recover verdict rows."
            )
            salvaged_verdicts:        list = []
            salvaged_selected_detail: list = []
            salvaged_playbook:        "str | None" = None

            # Best-effort playbook extraction from the partial text.
            pb_match = re.search(
                r'"category_playbook"\s*:\s*"([^"]+)"', raw
            )
            if pb_match:
                salvaged_playbook = pb_match.group(1)

            # Walk the text and pull every complete balanced {...} object,
            # at ANY nesting depth. We have to scan from every '{' because
            # the outermost '{' in a truncated payload never closes — so a
            # depth=0 anchor would never extract anything. By trying each
            # opening brace independently we still recover the well-formed
            # inner verdict rows even though the wrapping object is broken.
            def _extract_balanced_objects(text: str) -> list[str]:
                results: list[str] = []
                i = 0
                n = len(text)
                while i < n:
                    if text[i] != "{":
                        i += 1
                        continue
                    depth     = 0
                    in_string = False
                    escape    = False
                    end       = -1
                    for j in range(i, n):
                        ch = text[j]
                        if in_string:
                            if escape:
                                escape = False
                            elif ch == "\\":
                                escape = True
                            elif ch == '"':
                                in_string = False
                            continue
                        if ch == '"':
                            in_string = True
                        elif ch == "{":
                            depth += 1
                        elif ch == "}":
                            depth -= 1
                            if depth == 0:
                                end = j
                                break
                    if end > 0:
                        results.append(text[i:end + 1])
                        i = end + 1
                    else:
                        # This '{' never closes (truncated wrapper). Move on
                        # one character and keep trying smaller inner spans.
                        i += 1
                return results

            # The top-level `{...}` object will be missing in a truncated
            # payload, but inner verdict / detail objects will still be
            # complete. Scan and parse each one individually.
            for blob in _extract_balanced_objects(raw):
                try:
                    obj = json.loads(blob)
                except json.JSONDecodeError:
                    continue
                if not isinstance(obj, dict):
                    continue
                # A verdict row has these signature keys.
                if "verdict" in obj and "purchase_dtl_id" in obj:
                    salvaged_verdicts.append(obj)
                # A selected_detail row has reasoning + matched_attrs and
                # also a purchase_dtl_id.
                elif (
                    "purchase_dtl_id" in obj
                    and ("reasoning" in obj or "matched_attrs" in obj)
                ):
                    salvaged_selected_detail.append(obj)

            if salvaged_verdicts:
                logger.info(
                    f"V2 rank LLM partial-JSON salvage recovered "
                    f"{len(salvaged_verdicts)} verdict row(s) and "
                    f"{len(salvaged_selected_detail)} detail row(s) from "
                    f"the truncated payload."
                )
                parsed = {
                    "category_playbook": salvaged_playbook,
                    "verdicts":          salvaged_verdicts,
                    "selected_detail":   salvaged_selected_detail,
                    "_partial":          True,
                }
            else:
                # Nothing usable — preserve old behavior.
                return [], [], {"token_usage": usage, "selected": [], "rejected": []}
        if not isinstance(parsed, dict):
            return [], [], {"token_usage": usage, "selected": [], "rejected": []}

        # ── New two-array schema: verdicts[] + selected_detail[] ──────────
        # The LLM emits one COMPACT row per candidate in `verdicts` (forces
        # completeness) and rich reasoning only for PASS rows in
        # `selected_detail`. We adapt the pair back into the existing
        # selected[]/rejected[] shape so downstream code (DB MERGE, funnel,
        # logs) is untouched.
        #
        # Backwards compatibility: if the LLM emits the old shape with
        # `selected`/`rejected` arrays, fall back to that. Lets old prompts
        # set via the prompts dict keep working.
        verdicts_arr = parsed.get("verdicts")
        if isinstance(verdicts_arr, list):
            detail_arr = parsed.get("selected_detail", []) or []
            detail_by_id: dict[int, dict] = {}
            for d in detail_arr:
                if isinstance(d, dict) and d.get("purchase_dtl_id") is not None:
                    try:
                        detail_by_id[int(d["purchase_dtl_id"])] = d
                    except (TypeError, ValueError):
                        continue

            selected: list = []
            rejected: list = []
            for v in verdicts_arr:
                if not isinstance(v, dict):
                    continue
                did = v.get("purchase_dtl_id")
                if did is None:
                    continue
                try:
                    did_int = int(did)
                except (TypeError, ValueError):
                    continue
                verdict_value = (v.get("verdict") or "").strip().upper()
                if verdict_value == "PASS":
                    detail = detail_by_id.get(did_int, {})
                    merged = {**v, **detail, "purchase_dtl_id": did_int}
                    # Remove fields that belong only to the compact verdict
                    # row to keep the downstream shape clean for `selected`.
                    for k in ("verdict",):
                        # Keep verdict — it now holds STRONG/MODERATE/WEAK_MATCH
                        # from the detail row, not the PASS/FAIL flag.
                        pass
                    # If the detail row defined its own verdict (STRONG_MATCH
                    # etc.) prefer that; otherwise default to MODERATE_MATCH.
                    if detail.get("verdict"):
                        merged["verdict"] = detail["verdict"]
                    else:
                        merged["verdict"] = "MODERATE_MATCH"
                    selected.append(merged)
                else:
                    rejected.append({
                        "purchase_dtl_id":    did_int,
                        "rule":               v.get("rule")  or "OTHER",
                        "disqualifier_attr":  v.get("attr"),
                        "source_value":       v.get("src"),
                        "candidate_value":    v.get("cand"),
                        "unit":               v.get("unit"),
                        "delta_pct":          v.get("delta_pct"),
                        "note":               v.get("note") or "",
                        "reasoning":          _build_short_reject_reasoning(v, crit_pct=crit_pct),
                    })
        else:
            # Backwards-compatible path for old prompt shape.
            selected = parsed.get("selected", []) or []
            rejected = parsed.get("rejected", []) or []

        for i, s in enumerate(selected):
            if isinstance(s, dict):
                s["rank"] = int(s.get("rank", i + 1))
        agent_reasoning = {
            "category_playbook": parsed.get("category_playbook"),
            "thresholds": {
                "critical_pct":  crit_pct,
                "important_pct": imp_pct,
                "ratio_pct":     ratio_pct,
            },
            "token_usage": usage,
            "selected":    selected,
            "rejected":    rejected,
        }
        return selected, rejected, agent_reasoning
    except Exception as exc:
        logger.warning(f"V2 rank LLM failed: {exc}")
        return [], [], {}


def _build_source_snapshot_for_rank(rd: dict) -> dict:
    """Full snapshot the rank LLM receives as the SOURCE item — all fields from
    quotation_extracted_items that are relevant for relevance judgement."""
    crit = rd.get("critical_attributes")
    if isinstance(crit, str):
        try: crit = json.loads(crit)
        except Exception: crit = []
    cat_parts = [rd.get(f"item_level_{n}") for n in range(1, 9) if rd.get(f"item_level_{n}")]
    return {
        "purchase_dtl_id":       int(rd.get("purchase_dtl_id") or 0),
        "item_name":             rd.get("item_name"),
        "item_description":      rd.get("item_description"),
        "item_summary":          rd.get("item_summary"),
        "category":              " > ".join(str(p) for p in cat_parts) if cat_parts else None,
        "purchase_category_llm": rd.get("purchase_category_llm"),
        "commodity_tag":         rd.get("commodity_tag"),
        "critical_attributes":   crit or [],
        "supplier_name":         rd.get("supplier_name"),
        "quotation_date":        str(rd["quotation_date"]) if rd.get("quotation_date") else None,
        "quantity":              float(rd["quantity"]) if rd.get("quantity") is not None else None,
        "unit":                  rd.get("unit"),
        "unit_price":            float(rd["unit_price"]) if rd.get("unit_price") is not None else None,
        "unit_price_eur":        float(rd["unit_price_eur"]) if rd.get("unit_price_eur") is not None else None,
        "total_price_eur":       float(rd["total_price_eur"]) if rd.get("total_price_eur") is not None else None,
        "currency":              rd.get("currency"),
        "discount":              float(rd["discount"]) if rd.get("discount") is not None else None,
        "payment_terms":         rd.get("payment_terms"),
        "delivery_time_days":    rd.get("delivery_time_days"),
    }


def _run_benchmark_v2(
    llm,
    tgt_cs: str,
    pr_no: str,
    embed_model,
    pinecone_index: str,
    pinecone_ns: str,
    top_k_final: int,
    prompts: dict | None = None,
) -> None:
    """Stage 7 V2 — 3-stage retrieval (SQL → Pinecone → LLM rank) + V1-identical
    pricing (low/last/bp_unit_price/inflation) applied to the ranked shortlist.
    Current PR's own DTL IDs are excluded from the candidate pool at every stage.

    NOTE: benchmark_result.agent_reasoning (NVARCHAR(MAX), nullable) is assumed
    to exist. Run this once manually on the target DB before the first V2 run:
        ALTER TABLE [ras_procurement].[benchmark_result]
            ADD [agent_reasoning] NVARCHAR(MAX) NULL;
    """
    from langchain_core.messages import HumanMessage
    p = prompts or {}
    sql_pool_size     = int(p.get("bench_sql_pool_size",           100))
    pinecone_top_k    = int(p.get("bench_pinecone_top_k",          10))
    min_score         = float(p.get("bench_min_similarity",        0.80))
    llm_shortlist     = int(p.get("bench_llm_shortlist_size",      max(10, top_k_final)))
    widen_when_sparse = bool(p.get("bench_widen_l1l2_when_sparse", True))
    outlier_factor    = float(p.get("bench_outlier_factor",        3.0))
    max_age_months    = int(p.get("bench_max_age_months",          0))
    uom_strict        = bool(p.get("bench_uom_strict",             False))

    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[purchase_category_llm], qi.[embed_content],
                   qi.[critical_attributes], qi.[commodity_tag],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8],
                   qi.[unit_price], qi.[total_price], qi.[quantity], qi.[unit],
                   qi.[currency], qi.[quotation_date], qi.[supplier_name],
                   qi.[unit_price_eur], qi.[total_price_eur],
                   prd.[C_DATETIME] AS [item_created_date]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        rows = cur.fetchall()
    finally:
        conn.close()

    src_cols = [
        "extracted_item_uuid_pk","purchase_dtl_id","item_name","item_description","item_summary",
        "purchase_category_llm","embed_content","critical_attributes","commodity_tag",
        "item_level_1","item_level_2","item_level_3","item_level_4",
        "item_level_5","item_level_6","item_level_7","item_level_8",
        "unit_price","total_price","quantity","unit","currency","quotation_date","supplier_name",
        "unit_price_eur","total_price_eur","item_created_date",
    ]

    # Collect current PR's own DTL IDs so they are never returned as "similar"
    current_pr_dtl_ids = {int(r[1]) for r in rows if r[1] is not None}

    conn2 = _connect(tgt_cs)
    cur2  = conn2.cursor()
    try:
        for raw_row in rows:
            rd = dict(zip(src_cols, raw_row))
            dtl_id    = int(rd["purchase_dtl_id"])
            item_uuid = str(rd["extracted_item_uuid_pk"] or "")
            src_cat   = (rd.get("purchase_category_llm") or "").strip()
            src_date  = rd.get("item_created_date")
            qty       = rd.get("quantity") or Decimal("1")

            created = src_date
            created_iso = (
                created.isoformat() if created and hasattr(created, "isoformat") else str(created or "")
            )

            embed_text = (rd.get("embed_content") or "").strip() or _build_embed_text(rd)
            if not embed_text:
                _write_benchmark_stub(cur2, dtl_id, item_uuid,
                                      "V2 — no embed_content / no descriptor fields")
                continue

            # ── Funnel diagnostic — capture DTL IDs at every filter stage ──
            # Goes into benchmark_result.agent_reasoning under
            # `pipeline_funnel` so the buyer can see EXACTLY which candidates
            # entered the pool and where each one was dropped. Mirror of the
            # log lines below — same data, persisted instead of ephemeral.
            funnel: dict = {
                "source_dtl_id":          dtl_id,
                "source_category":        src_cat,
                "source_created_iso":     created_iso,
                "knobs": {
                    "sql_pool_size":          sql_pool_size,
                    "pinecone_top_k":         pinecone_top_k,
                    "min_similarity":         min_score,
                    "llm_shortlist_size":    llm_shortlist,
                    "widen_l1l2_when_sparse": widen_when_sparse,
                    "max_price_ratio":        float(p.get("bench_max_price_ratio", 50.0)),
                },
                "stage_a_sql_pool":         [],
                "stage_a_widened_l1l2":     [],
                "stage_b_pinecone":         [],
                "combined_unique":          [],
                "dropped_newer_than_src":   [],
                "dropped_pre_filter":       [],
                "sent_to_llm":              [],
                "llm_selected":             [],
                "llm_rejected":             [],
                "llm_dropped_silently":     [],
            }

            # ── Stage A — SQL category pool ───────────────────────────────
            own_excl = list(current_pr_dtl_ids)
            filter_chain = ""
            try:
                pool: list[int] = _sql_pool_by_category(tgt_cs, src_cat, src_date, sql_pool_size, own_excl)
                funnel["stage_a_sql_pool"] = list(pool)
                filter_chain = "sql_category"
                if len(pool) < 20 and widen_when_sparse:
                    widened = _widen_pool_by_levels(
                        tgt_cs, rd.get("item_level_1") or "", rd.get("item_level_2") or "",
                        src_date, sql_pool_size, own_excl,
                    )
                    seen_ids = set(pool)
                    new_widened: list[int] = []
                    for did in widened:
                        if did not in seen_ids:
                            pool.append(did); seen_ids.add(did); new_widened.append(did)
                    funnel["stage_a_widened_l1l2"] = new_widened
                    if widened:
                        filter_chain += "+widen_l1l2"
                        logger.info(
                            f"[V2 {pr_no}] dtl_id={dtl_id} sparse pool — "
                            f"widened to L1+L2, added {len(new_widened)} dtl_ids: "
                            f"{new_widened[:30]}{'...' if len(new_widened) > 30 else ''}"
                        )
                logger.info(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage A: {len(pool)} SQL candidate(s); "
                    f"ids={pool[:30]}{'...' if len(pool) > 30 else ''}"
                )
            except Exception as exc:
                logger.error(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage A (SQL pool) failed: {exc}",
                    exc_info=True,
                )
                pool = []
                filter_chain = "sql_failed"

            # ── Stage B — Pinecone vector search (independent of SQL pool) ──
            pinecone_dtl_ids: list[int] = []
            try:
                embedding = _embed_with_timeout(embed_model, embed_text)
                narrowed = _pinecone_narrow_within_pool(
                    embedding, [], pinecone_top_k, pinecone_index, pinecone_ns,
                    min_score=min_score, query_text=embed_text,
                )
                for m in narrowed:
                    meta = (m.get("metadata") if isinstance(m, dict) else {}) or {}
                    did = meta.get("purchase_dtl_id")
                    if did is not None:
                        try:
                            did_int = int(did)
                            if did_int not in current_pr_dtl_ids:
                                pinecone_dtl_ids.append(did_int)
                        except Exception:
                            continue
                funnel["stage_b_pinecone"] = list(pinecone_dtl_ids)
                filter_chain += "+pinecone"
                logger.info(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage B: {len(pinecone_dtl_ids)} "
                    f"Pinecone match(es) (min_score>={min_score}); ids={pinecone_dtl_ids}"
                )
            except Exception as exc:
                if _is_pinecone_server_error(exc):
                    raise RuntimeError(
                        f"Pinecone server error during V2 benchmark for PR={pr_no}: {exc}"
                    ) from exc
                logger.error(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage B (Pinecone) failed: {exc}",
                    exc_info=True,
                )
                filter_chain += "+pinecone_failed"

            # ── Combine Stage A + Stage B, deduplicate ────────────────────
            combined_dtl_ids: list[int] = list(pool)
            seen_combined: set[int] = set(pool)
            pinecone_only: list[int] = []
            for did in pinecone_dtl_ids:
                if did not in seen_combined:
                    combined_dtl_ids.append(did)
                    seen_combined.add(did)
                    pinecone_only.append(did)
            funnel["combined_unique"] = list(combined_dtl_ids)
            overlap_count = len(pinecone_dtl_ids) - len(pinecone_only)
            logger.info(
                f"[V2 {pr_no}] dtl_id={dtl_id} combined A+B: {len(combined_dtl_ids)} unique candidate(s) "
                f"(sql={len(pool)}, pinecone={len(pinecone_dtl_ids)}, overlap={overlap_count}, "
                f"pinecone_only_new={len(pinecone_only)})"
            )

            if not combined_dtl_ids:
                logger.warning(
                    f"[V2 {pr_no}] dtl_id={dtl_id} EMPTY POOL — Stage A and Stage B "
                    f"returned no candidates. category={src_cat!r}, created<{created_iso!r}. "
                    f"Likely causes: (1) category alias mismatch, (2) source PR predates "
                    f"all historical peers, (3) no peers have unit_price_eur."
                )
                _write_benchmark_stub(
                    cur2, dtl_id, item_uuid,
                    f"V2 chain={filter_chain} — Stage A and Stage B both returned no candidates. "
                    f"category={src_cat!r}, created<{created_iso!r}",
                )
                continue

            # ── Stage C — LLM relevance ranking ───────────────────────────
            shortlist_dtl_ids: list[int] = []
            rejected: list = []
            agent_reasoning_payload: dict = {}
            try:
                candidates = _fetch_candidate_snapshots(tgt_cs, combined_dtl_ids)
                # Hard date guard: only use items older than the current PR.
                # Stage A already enforces this via SQL, but Stage B (Pinecone) does
                # not, so newer items can slip through the combined pool.
                if created_iso:
                    before = []
                    skipped_newer = []
                    for c in candidates:
                        c_date = c.get("item_created_date")
                        if c_date is None:
                            before.append(c)
                            continue
                        c_date_str = (
                            c_date.isoformat() if hasattr(c_date, "isoformat") else str(c_date)
                        )
                        if c_date_str >= created_iso:
                            skipped_newer.append(c.get("purchase_dtl_id"))
                        else:
                            before.append(c)
                    funnel["dropped_newer_than_src"] = [
                        int(x) for x in skipped_newer if x is not None
                    ]
                    if skipped_newer:
                        logger.info(
                            f"[V2 {pr_no}] dtl_id={dtl_id} Stage C date guard: dropped "
                            f"{len(skipped_newer)} candidate(s) newer than source "
                            f"({created_iso}): {skipped_newer}"
                        )
                    candidates = before
                # Pre-filter: drop extreme price outliers before LLM ranking
                pre_filter_input_ids = {
                    int(c.get("purchase_dtl_id"))
                    for c in candidates if c.get("purchase_dtl_id") is not None
                }
                candidates = _pre_filter_candidates(rd, candidates, prompts)
                pre_filter_output_ids = {
                    int(c.get("purchase_dtl_id"))
                    for c in candidates if c.get("purchase_dtl_id") is not None
                }
                dropped_by_prefilter = sorted(pre_filter_input_ids - pre_filter_output_ids)
                funnel["dropped_pre_filter"] = dropped_by_prefilter
                if dropped_by_prefilter:
                    logger.info(
                        f"[V2 {pr_no}] dtl_id={dtl_id} pre-filter dropped "
                        f"{len(dropped_by_prefilter)} candidate(s) for extreme price ratio "
                        f"(>{funnel['knobs']['max_price_ratio']}×): {dropped_by_prefilter}"
                    )
                sent_to_llm_ids = [
                    int(c.get("purchase_dtl_id"))
                    for c in candidates if c.get("purchase_dtl_id") is not None
                ]
                funnel["sent_to_llm"] = sent_to_llm_ids
                logger.info(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage C INPUT to LLM: "
                    f"{len(sent_to_llm_ids)} candidate(s); ids={sent_to_llm_ids}"
                )
                source_snapshot = _build_source_snapshot_for_rank(rd)
                selected, rejected, agent_reasoning_payload = _llm_rank_candidates(
                    llm, source_snapshot, candidates, llm_shortlist, prompts,
                )
                for s in selected:
                    if not isinstance(s, dict):
                        continue
                    sid = s.get("purchase_dtl_id")
                    if sid is not None:
                        try:
                            sid_int = int(sid)
                            if sid_int not in current_pr_dtl_ids:
                                shortlist_dtl_ids.append(sid_int)
                        except Exception:
                            continue
                # Drop any rows the LLM emitted for the current PR's own DTL
                # IDs (defensive — candidates pool already excludes them).
                if agent_reasoning_payload and current_pr_dtl_ids:
                    agent_reasoning_payload["selected"] = [
                        row for row in (agent_reasoning_payload.get("selected") or [])
                        if isinstance(row, dict)
                        and row.get("purchase_dtl_id") not in current_pr_dtl_ids
                    ]
                # Completeness diagnostic ONLY — we no longer backfill the
                # rejected[] array with LLM_DROPPED placeholders, because the
                # prompt now hard-caps verdicts to 20–25 from a much larger
                # pool. Candidates outside the top 20–25 are intentionally
                # not given a verdict and MUST NOT appear in agent_reasoning
                # (this is customer-facing output). We still log which IDs
                # were silently dropped so an operator can spot prompt drift,
                # but the customer-facing JSON only contains the LLM's real
                # verdicts.
                emitted_sel_ids = {
                    int(r.get("purchase_dtl_id"))
                    for r in (agent_reasoning_payload.get("selected") or [])
                    if isinstance(r, dict) and r.get("purchase_dtl_id") is not None
                }
                emitted_rej_ids = {
                    int(r.get("purchase_dtl_id"))
                    for r in (agent_reasoning_payload.get("rejected") or [])
                    if isinstance(r, dict) and r.get("purchase_dtl_id") is not None
                }
                sent_set = set(sent_to_llm_ids)
                dropped_silently = sorted(sent_set - emitted_sel_ids - emitted_rej_ids)
                if dropped_silently:
                    logger.info(
                        f"[V2 {pr_no}] dtl_id={dtl_id} LLM did not emit verdicts "
                        f"for {len(dropped_silently)} candidate(s) outside its top "
                        f"20–25 selection (expected behavior — these are not "
                        f"backfilled into the customer-facing reasoning)"
                    )
                funnel["llm_selected"]         = sorted(emitted_sel_ids)
                funnel["llm_rejected"]         = sorted(emitted_rej_ids)
                funnel["llm_dropped_silently"] = dropped_silently
                # Read token_usage from the agent reasoning ONLY for log
                # decoration + filter_chain marker — do NOT persist a copy
                # inside `pipeline_funnel` (top-level `agent_reasoning
                # .token_usage` already carries it; persisting twice bloats
                # the JSON).
                top_level_token_usage = (
                    agent_reasoning_payload.get("token_usage")
                    if isinstance(agent_reasoning_payload, dict) else None
                )
                truncated_flag = bool(
                    top_level_token_usage and top_level_token_usage.get("truncated")
                )
                filter_chain += "+llm_rank"
                if truncated_flag:
                    filter_chain += "+truncated"
                logger.info(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage C OUTPUT: "
                    f"selected={sorted(emitted_sel_ids)} ({len(emitted_sel_ids)}), "
                    f"rejected={sorted(emitted_rej_ids)} ({len(emitted_rej_ids)})"
                    + (f", llm_dropped={dropped_silently}" if dropped_silently else "")
                    + (
                        f"  [!!! TRUNCATED: output {top_level_token_usage.get('output_tokens')}"
                        f"/{top_level_token_usage.get('max_output_tokens')} tokens, "
                        f"finish_reason={top_level_token_usage.get('finish_reason')!r} !!!]"
                        if truncated_flag and top_level_token_usage else ""
                    )
                )

                # Per-candidate verdict trace — one log line per candidate ID
                # that was sent to the LLM, naming the verdict the LLM gave it.
                # Use this to audit "did the LLM actually consider DTL X?" and
                # see WHY each one landed as SELECTED / REJECTED / DROPPED.
                sel_by_id = {
                    int(r["purchase_dtl_id"]): r
                    for r in (agent_reasoning_payload.get("selected") or [])
                    if isinstance(r, dict) and r.get("purchase_dtl_id") is not None
                }
                rej_by_id = {
                    int(r["purchase_dtl_id"]): r
                    for r in (agent_reasoning_payload.get("rejected") or [])
                    if isinstance(r, dict) and r.get("purchase_dtl_id") is not None
                }
                for cand_id in sent_to_llm_ids:
                    if cand_id in sel_by_id:
                        row = sel_by_id[cand_id]
                        logger.info(
                            f"[V2 {pr_no}] dtl_id={dtl_id} candidate {cand_id} "
                            f"→ SELECTED "
                            f"(rank={row.get('rank')}, "
                            f"verdict={row.get('verdict')}, "
                            f"confidence={row.get('confidence')})"
                        )
                    elif cand_id in rej_by_id:
                        row = rej_by_id[cand_id]
                        logger.info(
                            f"[V2 {pr_no}] dtl_id={dtl_id} candidate {cand_id} "
                            f"→ REJECTED "
                            f"(rule={row.get('rule')}, "
                            f"attr={row.get('disqualifier_attr')}, "
                            f"src={row.get('source_value')}, "
                            f"cand={row.get('candidate_value')}, "
                            f"unit={row.get('unit')}, "
                            f"delta_pct={row.get('delta_pct')})"
                        )
                    else:
                        # Should not happen because we backfill LLM_DROPPED
                        # above; defensive log in case of races.
                        logger.warning(
                            f"[V2 {pr_no}] dtl_id={dtl_id} candidate {cand_id} "
                            f"→ NO VERDICT (LLM silently dropped — see warning above)"
                        )
            except Exception as exc:
                logger.error(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage C (LLM rank) failed: {exc}",
                    exc_info=True,
                )
                filter_chain += "+llm_rank_failed"

            # Attach the funnel to the reasoning payload so it lands in
            # benchmark_result.agent_reasoning alongside selected/rejected.
            if isinstance(agent_reasoning_payload, dict):
                agent_reasoning_payload["pipeline_funnel"] = funnel

            # No similar IDs → still run pricing LLM with empty historical so it
            # can give a market-knowledge estimate with low confidence.
            if not shortlist_dtl_ids:
                logger.warning(
                    f"[V2 {pr_no}] dtl_id={dtl_id} Stage C returned no shortlist "
                    f"(chain={filter_chain}) — proceeding to pricing LLM with no historical data"
                )

            # ── Pricing ──────────────────────────────────────────────────
            historical_raw   = _fetch_historical_for_dtl_ids(tgt_cs, shortlist_dtl_ids) if shortlist_dtl_ids else []
            prior_benchmarks = _fetch_benchmark_for_dtl_ids(tgt_cs, shortlist_dtl_ids) if shortlist_dtl_ids else []

            historical = _filter_historical_uom(historical_raw, rd.get("unit"), uom_strict)
            historical = _filter_historical_age(historical, created_iso, max_age_months)
            historical = _filter_historical_outliers(historical, outlier_factor)
            if len(historical) != len(historical_raw):
                logger.info(
                    f"[V2 {pr_no}] dtl_id={dtl_id}: filtered "
                    f"{len(historical_raw) - len(historical)}/{len(historical_raw)} historical row(s)"
                )
            if shortlist_dtl_ids and not historical and historical_raw:
                logger.warning(
                    f"[V2 {pr_no}] dtl_id={dtl_id}: {len(shortlist_dtl_ids)} shortlist item(s) found "
                    f"but ALL {len(historical_raw)} row(s) eliminated by post-fetch filters — "
                    f"falling back to unfiltered data for low_hist/last_hist FK selection"
                )
                historical = historical_raw

            agg = _benchmark_aggregates(historical, created_iso)
            low_item, last_item = _compute_low_last(historical)
            low_uuid  = str(low_item["extracted_item_uuid_pk"])  if low_item  and low_item.get("extracted_item_uuid_pk")  else None
            last_uuid = str(last_item["extracted_item_uuid_pk"]) if last_item and last_item.get("extracted_item_uuid_pk") else None
            if shortlist_dtl_ids and low_uuid is None:
                if low_item is None:
                    logger.warning(
                        f"[V2 {pr_no}] dtl_id={dtl_id}: low_hist_item_fk is NULL — "
                        f"historical={len(historical)} row(s) but all have NULL unit_price_eur AND unit_price"
                    )
                else:
                    logger.warning(
                        f"[V2 {pr_no}] dtl_id={dtl_id}: low_hist_item_fk is NULL — "
                        f"low_item found (dtl_id={low_item.get('purchase_dtl_id')}) but extracted_item_uuid_pk is NULL in DB"
                    )
            if shortlist_dtl_ids and last_uuid is None:
                if last_item is None:
                    logger.warning(
                        f"[V2 {pr_no}] dtl_id={dtl_id}: last_hist_item_fk is NULL — "
                        f"historical={len(historical)} row(s) but all have NULL pr_created_date"
                    )
                else:
                    logger.warning(
                        f"[V2 {pr_no}] dtl_id={dtl_id}: last_hist_item_fk is NULL — "
                        f"last_item found (dtl_id={last_item.get('purchase_dtl_id')}) but extracted_item_uuid_pk is NULL in DB"
                    )
            similar_dtl_ids_json = json.dumps(shortlist_dtl_ids) if shortlist_dtl_ids else None
            agent_reasoning_json = (
                json.dumps(agent_reasoning_payload, ensure_ascii=False, default=str)
                if agent_reasoning_payload else None
            )

            # Inflation
            supplier_country = (low_item.get("supplier_country") if low_item else None)
            # Use purchase_req_mst.C_DATETIME instead of quotation_date for consistency
            ref_dt       = _get_pr_master_date_for_dtl_id(tgt_cs, low_item.get("purchase_dtl_id")) if low_item else None
            ref_year     = ref_dt.year if ref_dt and hasattr(ref_dt, "year") else None
            current_year = created.year if created and hasattr(created, "year") else None
            item_category = " > ".join(
                str(rd[f]) for f in [
                    "item_level_1", "item_level_2", "item_level_3",
                ] if rd.get(f)
            )
            infl_dec = cpi_dec = Decimal("0")
            infl_dec_last = cpi_dec_last = Decimal("0")
            if ref_year and current_year and ref_year < current_year:
                infl_raw = _estimate_inflation_via_llm(
                    llm, rd.get("item_name"), item_category or None,
                    supplier_country, ref_year, current_year,
                )
                if infl_raw is not None:
                    try: infl_dec = Decimal(str(infl_raw))
                    except Exception: pass
                cpi_raw = _compute_cpi_pct(supplier_country, ref_year, current_year)
                if cpi_raw is not None:
                    try: cpi_dec = Decimal(str(cpi_raw))
                    except Exception: pass
                logger.info(
                    f"[V2 {pr_no}] dtl_id={dtl_id}: inflation_pct={infl_dec} "
                    f"cpi_pct={cpi_dec} (country={supplier_country!r} {ref_year}-{current_year})"
                )

            # Also calculate inflation for last_hist_item (most recent)
            if last_item and current_year:
                supplier_country_last = last_item.get("supplier_country")
                if not supplier_country_last:
                    logger.info(f"[{pr_no}] dtl_id={dtl_id}: last_hist_item missing supplier_country, skipping inflation_pct_last")
                else:
                    ref_dt_last = _get_pr_master_date_for_dtl_id(tgt_cs, last_item.get("purchase_dtl_id"))
                    ref_year_last = ref_dt_last.year if ref_dt_last and hasattr(ref_dt_last, "year") else None
                    if ref_year_last and current_year and ref_year_last < current_year:
                        infl_raw_last = _estimate_inflation_via_llm(
                            llm, rd.get("item_name"), item_category or None,
                            supplier_country_last, ref_year_last, current_year,
                        )
                        if infl_raw_last is not None:
                            try: infl_dec_last = Decimal(str(infl_raw_last))
                            except Exception: pass
                        cpi_raw_last = _compute_cpi_pct(supplier_country_last, ref_year_last, current_year)
                        if cpi_raw_last is not None:
                            try: cpi_dec_last = Decimal(str(cpi_raw_last))
                            except Exception: pass
                        logger.info(
                            f"[V2 {pr_no}] dtl_id={dtl_id}: inflation_pct_last={infl_dec_last} "
                            f"cpi_pct_last={cpi_dec_last} (country={supplier_country_last!r} {ref_year_last}-{current_year})"
                        )

            # LLM benchmark price
            bench_prompt = _format_bench_prompt(
                rd, historical[:top_k_final], stats=agg, prior_benchmarks=prior_benchmarks or None,
            )
            bout = {}
            llm_succeeded = False
            try:
                resp = _call_llm_with_retry(llm, [HumanMessage(content=bench_prompt)], prompts=prompts, force_json=True)
                raw_txt = (getattr(resp, "content", None) or str(resp)).strip()
                raw_txt = re.sub(r"^```(?:json)?\s*", "", raw_txt)
                raw_txt = re.sub(r"\s*```$", "", raw_txt)
                bout = json.loads(raw_txt)
                llm_succeeded = True
            except Exception as exc:
                logger.error(
                    f"[V2 {pr_no}] LLM benchmark parse failed dtl_id={dtl_id}: {exc}",
                    exc_info=True,
                )

            bp_unit       = bout.get("bp_unit_price")
            bp_low        = bout.get("bp_low")
            bp_high       = bout.get("bp_high")
            bp_confidence = bout.get("confidence")
            llm_summary   = bout.get("summary", "")
            llm_succeeded = llm_succeeded and bp_unit is not None
            try:
                bp_dec   = Decimal(str(bp_unit)) if bp_unit is not None else Decimal("0")
                bp_total = round(float(bp_dec) * float(qty or 1), 2)
            except Exception:
                bp_dec = Decimal("0")
                bp_total = 0

            extras = []
            if bp_low is not None or bp_high is not None:
                extras.append(f"range={bp_low}-{bp_high} EUR")
            if bp_confidence is not None:
                try: extras.append(f"confidence={float(bp_confidence):.2f}")
                except Exception: extras.append(f"confidence={bp_confidence}")
            if agg.get("count") is not None:
                extras.append(f"n={agg['count']} (priced={agg.get('priced_count', 0)})")
            if "median" in agg:
                extras.append(f"median={agg['median']} EUR")
            summary = (
                f"[{', '.join(extras)}] {llm_summary}".strip()
                if extras else llm_summary
            )

            bp_dec_u   = bp_dec   if llm_succeeded else None
            bp_total_u = bp_total if llm_succeeded else None
            summary_u  = summary  if (llm_succeeded and summary) else None

            try:
                cur2.execute("""
                    MERGE [ras_procurement].[benchmark_result] WITH (HOLDLOCK) AS target
                    USING (SELECT ? AS purchase_dtl_id) AS src
                       ON target.purchase_dtl_id = src.purchase_dtl_id
                    WHEN MATCHED THEN
                        UPDATE SET
                            extracted_item_uuid_fk = ?,
                            bp_unit_price          = COALESCE(?, target.bp_unit_price),
                            bp_total_price         = COALESCE(?, target.bp_total_price),
                            low_hist_item_fk       = ?,
                            last_hist_item_fk      = ?,
                            inflation_pct          = ?,
                            cpi_inflation_pct      = ?,
                            inflation_pct_last     = ?,
                            cpi_inflation_pct_last = ?,
                            similar_dtl_ids        = ?,
                            agent_reasoning        = ?,
                            summary                = COALESCE(?, target.summary),
                            updated_at             = SYSUTCDATETIME()
                    WHEN NOT MATCHED THEN
                        INSERT (
                            purchase_dtl_id, extracted_item_uuid_fk,
                            bp_unit_price, bp_total_price,
                            low_hist_item_fk, last_hist_item_fk,
                            inflation_pct, cpi_inflation_pct,
                            inflation_pct_last, cpi_inflation_pct_last,
                            similar_dtl_ids, agent_reasoning, summary
                        )
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?);
                """,
                    dtl_id,
                    item_uuid, bp_dec_u, bp_total_u, low_uuid, last_uuid, infl_dec, cpi_dec, infl_dec_last, cpi_dec_last, similar_dtl_ids_json, agent_reasoning_json, summary_u,
                    dtl_id, item_uuid, bp_dec, bp_total, low_uuid, last_uuid, infl_dec, cpi_dec, infl_dec_last, cpi_dec_last, similar_dtl_ids_json, agent_reasoning_json, summary,
                )
                conn2.commit()
            except Exception as exc:
                logger.warning(f"[V2 {pr_no}] benchmark_result MERGE failed dtl_id={dtl_id}: {exc}")
    finally:
        conn2.close()


# ── Main Component ─────────────────────────────────────────────────────────────


# ── Main component ────────────────────────────────────────────────────────

class PipelineStage123NodeV2(Node):
    display_name = "Full Pipeline V2 (LLM Category + Smart Benchmark)"
    description  = (
        "V2 of the full procurement pipeline. Stages 1-4 and 8 are identical to V1; "
        "Stage 5 (Extraction) additionally emits purchase_category_llm, "
        "embed_content (brand/supplier-free), and critical_attributes. "
        "Stage 6 (Embeddings) feeds Pinecone with embed_content. "
        "Stage 7 (Benchmark) uses 3-stage retrieval: SQL category filter → "
        "Pinecone within pool → LLM relevance ranking. Requires the V2 DDL "
        "migration on ras_procurement.quotation_extracted_items."
    )
    icon = "Database"
    name = "PipelineStage123NodeV2"

    inputs = [
        HandleInput(
            name="source_connection",
            display_name="Source DB — On-Prem SQL Server",
            input_types=["Data"],
            info="Connection Config from the on-prem Database Connector node.",
        ),
        HandleInput(
            name="target_connection",
            display_name="Target DB — Azure SQL",
            input_types=["Data"],
            info="Connection Config from the Azure SQL Database Connector node.",
        ),
        MessageTextInput(
            name="blob_connector_name",
            display_name="Azure Blob Connector Name",
            value="",
            info="Exact name of your Azure Blob connector in Settings → Connectors.",
        ),
        MessageTextInput(
            name="pr_no_filter",
            display_name="PR Number Filter (optional)",
            value="",
            advanced=True,
            info=(
                "Leave blank to process all pending approved PRs. "
                "Enter one PURCHASE_REQ_NO to force a full reprocess from scratch."
            ),
        ),
        # ── Excel-driven PR list (mirrors doc-intel run_pipeline_from_excel.py) ──
        # Use ANY ONE of:
        #   - PR Number Filter            → single PR, full reprocess
        #   - Excel Blob Path             → download Excel from the configured
        #                                   Blob Connector container, parse PR list
        #   - PR List Data (HandleInput)  → wire any component (Blob Reader,
        #                                   Knowledge Base, etc.) that outputs
        #                                   either raw bytes or a list under
        #                                   data["pr_numbers"]
        # If none are set, the component falls back to fetching pending PRs
        # from ras_tracker (the standard batch flow).
        MessageTextInput(
            name="pr_excel_blob_path",
            display_name="PR List Excel — Blob Path (optional)",
            value="",
            advanced=True,
            info=(
                "Blob path inside the configured Blob Connector container "
                "pointing to an .xlsx with a PURCHASE_REQ_NO column "
                "(e.g. procurement/pr_lists/2026-04.xlsx). Each PR is "
                "force-reprocessed from scratch like the single-PR filter."
            ),
        ),
        HandleInput(
            name="pr_list_data",
            display_name="PR List Data (optional)",
            input_types=["Data"],
            required=False,
            info=(
                "Wire a Blob Reader / Knowledge Base / file-loader component "
                "that outputs Data here. Accepted shapes: "
                "(a) data['pr_numbers'] = [list of PURCHASE_REQ_NO strings] — "
                "consumed directly; or "
                "(b) data['content'] / data['bytes'] / data['file_bytes'] = "
                "raw .xlsx bytes — parsed inline."
            ),
        ),
        MessageTextInput(
            name="pr_excel_column",
            display_name="PR List Excel — Column Name (optional)",
            value="",
            advanced=True,
            info=(
                "Header of the column holding PURCHASE_REQ_NO values. "
                "Auto-detected from common names "
                "(PURCHASE_REQ_NO / PR_NO / RAS_NO / PR Number) "
                "if blank, else falls back to the first column."
            ),
        ),
        MessageTextInput(
            name="pr_excel_sheet",
            display_name="PR List Excel — Sheet Name (optional)",
            value="",
            advanced=True,
            info="Worksheet to read. Defaults to the first sheet.",
        ),
        # ── Stages 4-8 inputs ────────────────────────────────────────────────
        HandleInput(
            name="llm",
            display_name="LLM (GPT-4o or GPT-4o-mini)",
            input_types=["LanguageModel"],
            info="LLM used for Stage 4 (classification) and Stage 5 (extraction).",
        ),
        HandleInput(
            name="embed_model",
            display_name="Embeddings Model",
            input_types=["Embeddings"],
            info="Embeddings model used for Stage 6 (vector embeddings) and Stage 7 (benchmark search).",
        ),
        # ── Prompt template overrides (optional) ─────────────────────────────
        # Leave disconnected to use the default prompts baked into this component.
        # Wire a Prompt Template node to any of these to override the built-in prompts.
        HandleInput(
            name="cls_system_prompt",
            display_name="[Classification] System Prompt",
            input_types=["Message"],
            required=False,
            info="Override the classification system prompt. Connect a Prompt Template node.",
        ),
        HandleInput(
            name="cls_user_text_prompt",
            display_name="[Classification] User Prompt — Text/Tabular Files",
            input_types=["Message"],
            required=False,
            info="Override the user prompt for text/Excel/Word/PDF files. Must keep {filename}, {file_type}, {extra_metadata}, {extracted_content} placeholders.",
        ),
        HandleInput(
            name="cls_user_image_prompt",
            display_name="[Classification] User Prompt — Image/Scanned Files",
            input_types=["Message"],
            required=False,
            info="Override the user prompt sent with base64 image content. Must keep {filename}, {file_type}, {extra_metadata} placeholders.",
        ),
        HandleInput(
            name="ext_system_prompt",
            display_name="[Extraction] System Prompt",
            input_types=["Message"],
            required=False,
            info="Override the extraction system prompt. Connect a Prompt Template node.",
        ),
        HandleInput(
            name="ext_user_template",
            display_name="[Extraction] User Prompt Template",
            input_types=["Message"],
            required=False,
            info="Override the extraction user template. Must keep all {field} placeholders from the RAS context and {document_content}.",
        ),
        HandleInput(
            name="ext_item_taxonomy",
            display_name="[Extraction] Item Taxonomy",
            input_types=["Message"],
            required=False,
            info="Override the item taxonomy guidelines injected into {item_taxonomy} of the extraction user template.",
        ),
        # ── Advanced settings ─────────────────────────────────────────────────
        MessageTextInput(
            name="pinecone_index",
            display_name="Pinecone Index Name",
            value="ras-quotations",
            advanced=True,
            info="Pinecone index used for embeddings (Stage 6) and retry cleanup.",
        ),
        MessageTextInput(
            name="pinecone_namespace",
            display_name="Pinecone Namespace",
            value="procurement",
            advanced=True,
            info="Pinecone namespace used for embeddings (Stage 6) and retry cleanup.",
        ),
        IntInput(name="pinecone_top_k",    display_name="Benchmark Top-K (Stage 7)", value=10,    advanced=True),
        IntInput(name="batch_limit",       display_name="Max PRs per Run (0 = all)", value=0,     advanced=True,
                 info="0 = process all pending PRs. Set to N (e.g. 50) to cap at N PRs per run. Ignored when PR Number Filter is set."),
        IntInput(name="parallel_workers",  display_name="Parallel Workers",          value=4,     advanced=True),
        IntInput(name="max_content_chars", display_name="Max Chars per File (Stage 2)", value=80000, advanced=True),
        # ── Stage 4-7 tuning knobs ────────────────────────────────────────────
        IntInput(
            name="cls_max_chars",
            display_name="Max Chars — Classification Input (Stage 4)",
            value=24000,
            advanced=True,
            info=(
                "Truncate extracted text before sending to the classification LLM. "
                "Default 24000 ≈ 6000 tokens, matches doc-intel "
                "truncate_to_token_limit(max_content_tokens=6000). 0 = no limit. "
                "Lowering this may cause Quotation signals (payment terms, totals, "
                "signatures) at the end of long documents to be missed."
            ),
        ),
        IntInput(
            name="cls_max_pages_vision",
            display_name="Max Pages — Vision Fallback (Stage 4)",
            value=5,
            advanced=True,
            info=(
                "Number of pages to capture as high-res images for the vision LLM "
                "when text-based classification yields 'Other' or low confidence. "
                "Default is 5 (first 5 pages)."
            ),
        ),
        IntInput(
            name="ext_max_chars",
            display_name="Max Chars — Extraction Document (Stage 5)",
            value=50000,
            advanced=True,
            info="Truncate document content before sending to the extraction LLM. 0 = no limit.",
        ),
        IntInput(
            name="ext_max_pages",
            display_name="Max Pages — PDF/Doc Render (Stage 5)",
            value=20,
            advanced=True,
            info=(
                "Max pages rendered per quotation document during extraction. "
                "PDFs longer than this only have their first N pages processed; "
                "a warning is logged when truncation happens. Matches doc-intel "
                "MAX_PAGES default."
            ),
        ),
        IntInput(
            name="ext_max_images",
            display_name="Max Images per LLM Call (Stage 5)",
            value=50,
            advanced=True,
            info=(
                "Hard cap on images attached to one extraction LLM request. "
                "Azure OpenAI rejects requests with more than ~50 images per call, "
                "so the value should not be raised above 50."
            ),
        ),
        MessageTextInput(
            name="selected_threshold",
            display_name="Selected Quote Min Score (Stage 7)",
            value="0.70",
            advanced=True,
            info="Minimum supplier_match_conf (0.0–1.0) for a quote to be eligible for selection.",
        ),
        MessageTextInput(
            name="price_max_boost",
            display_name="Price Alignment Max Boost (Stage 7)",
            value="0.10",
            advanced=True,
            info="Maximum conf boost (0.0–1.0) applied when extracted prices align with RAS prices.",
        ),
        MessageTextInput(
            name="canonicalize_threshold",
            display_name="Supplier Canonicalize Threshold (Stage 7)",
            value="0.82",
            advanced=True,
            info="SequenceMatcher ratio (0.0–1.0) above which two supplier names are merged into one canonical name.",
        ),
        # ── Benchmark filtering / quality knobs (Stage 7) ────────────────────
        MessageTextInput(
            name="bench_min_similarity",
            display_name="Benchmark Min Similarity (Stage 7)",
            value="0.80",
            advanced=True,
            info=(
                "Minimum cosine similarity (0.0–1.0) for a Pinecone match to count "
                "as a benchmark candidate. Raise toward 0.85 for catalogs with very "
                "specific items; lower toward 0.60 for generic categories. "
                "Was hardcoded to 0.70 prior to exposing this knob."
            ),
        ),
        MessageTextInput(
            name="bench_outlier_factor",
            display_name="Benchmark Outlier Factor (Stage 7)",
            value="3.0",
            advanced=True,
            info=(
                "Drop any historical match whose unit_price_eur is more than this "
                "factor away from the median (so > median × factor or < median / factor). "
                "Typical: 3.0 = drop 3× outliers. Set 0 to disable outlier removal."
            ),
        ),
        IntInput(
            name="bench_max_age_months",
            display_name="Benchmark Max Age (months, Stage 7)",
            value=0,
            advanced=True,
            info=(
                "Drop any historical match whose item_created_date is older than "
                "N months. Helps filter out stale prices when tech / specs move on. "
                "Set 0 to disable age filtering (use all older items)."
            ),
        ),
        IntInput(
            name="bench_uom_strict",
            display_name="Benchmark Strict UOM Match (Stage 7)",
            value=1,
            advanced=True,
            info=(
                "1 = only benchmark against historical items with the SAME UOM "
                "(prevents per-piece vs per-box-of-10 mixing). 0 = ignore UOM, "
                "match by similarity score only."
            ),
        ),
        # ── LLM retry / rate-limit settings ──────────────────────────────────
        IntInput(
            name="llm_max_retries",
            display_name="LLM Max Retries (Rate Limit)",
            value=3,
            advanced=True,
            info="How many times to retry an LLM call after a 429 / token-quota error before failing the stage.",
        ),
        IntInput(
            name="llm_retry_cooldown",
            display_name="LLM Retry Cooldown (seconds)",
            value=60,
            advanced=True,
            info=(
                "Seconds to wait before the first retry. Each subsequent retry waits "
                "cooldown × attempt (60s → 120s → 180s with default 3 retries)."
            ),
        ),
        # ── DB deadlock retry settings ──────────────────────────────────────
        # Multi-worker runs can hit SQL Server deadlock-victim errors
        # (1205 / SQLSTATE 40001) on the cleanup transaction because parallel
        # cleanups all DELETE through shared join chains. _process_pr wraps
        # _cleanup_for_pr in a retry loop with exponential backoff + ±30%
        # jitter so the victim transaction is rerun automatically.
        IntInput(
            name="db_deadlock_max_retries",
            display_name="DB Deadlock Max Retries",
            value=3,
            advanced=True,
            info=(
                "How many times to retry a SQL Server deadlock-victim error "
                "(1205 / 40001) before aborting the PR. Default 3 → up to 4 total "
                "attempts (1 initial + 3 retries). Set 0 to disable retries."
            ),
        ),
        IntInput(
            name="db_deadlock_base_delay",
            display_name="DB Deadlock Base Delay (seconds)",
            value=2,
            advanced=True,
            info=(
                "Base seconds for exponential backoff between deadlock retries. "
                "Per-attempt sleep = base × 2^attempt × (1 + 0..0.3 jitter). "
                "Default 2 → 4s, 8s, 16s base before jitter."
            ),
        ),
        # ── Classification parallelism ────────────────────────────────────────
        IntInput(
            name="cls_parallel_sources",
            display_name="Classification Parallel Sources (Stage 4)",
            value=8,
            advanced=True,
            info=(
                "Max files classified in parallel for a single PR (parent + embedded). "
                "Default 8 matches the previous hardcoded cap. Reduce if hitting LLM "
                "rate limits; increase for PRs with many attachments."
            ),
        ),
        # ── Extraction parallelism ────────────────────────────────────────────
        IntInput(
            name="ext_parallel_sources",
            display_name="Extraction Parallel Sources (Stage 5)",
            value=1,
            advanced=True,
            info=(
                "Number of quotation files extracted in parallel for a single PR. "
                "1 = sequential (safe). Increase to 2-3 to speed up PRs with many "
                "quotation files, but keep low to avoid multiplying LLM rate-limit "
                "pressure on top of the parallel_workers setting."
            ),
        ),
        # ── Commercials (Stage 5b) ────────────────────────────────────────────
        # Toggles the second LLM call that extracts commercial / financial
        # fields (Incoterms, freight, insurance, customs, taxes, grand total,
        # etc.) AFTER the line-item extraction succeeds. Fail-soft: if the
        # call errors, line items are still saved with commercial columns NULL.
        BoolInput(
            name="enable_commercials_extraction",
            display_name="Enable Commercials Extraction (Stage 5b)",
            value=True,
            advanced=True,
            info=(
                "When ON (default), runs a second LLM call per quotation source "
                "AFTER line-item extraction to extract commercial/financial fields "
                "(Incoterms, freight, insurance, customs/duties, packing & "
                "forwarding, installation, taxes, grand total, etc.) and map "
                "quote-level totals to each line via proportional allocation. "
                "Requires migration "
                "quotation_extracted_items_add_commercials_cols.sql. Failures are "
                "logged but never block the line-item save."
            ),
        ),
        # ── DB connection pool ────────────────────────────────────────────
        IntInput(
            name="max_db_connections",
            display_name="Max DB Connections (Pool Size)",
            value=20,
            advanced=True,
            info=(
                "Hard cap on concurrent SQL Server connections opened by this Node. "
                "Connections are pooled and reused — idle ones go back to the pool "
                "on .close() instead of being torn down. Sized below Azure SQL's "
                "per-DB session cap (typically 75-100 on Standard tiers). Increase "
                "only if you see worker threads blocking on connection acquire; "
                "decrease if Azure reports 'maximum number of connections' errors. "
                "Default 20 is safe for 4 parallel workers × 8 inner classification "
                "threads."
            ),
        ),
        # ── V2 benchmark tuning ────────────────────────────────────────────
        IntInput(
            name="bench_sql_pool_size",
            display_name="V2 — SQL Pool Size (Stage A)",
            value=100,
            advanced=True,
            info=(
                "TOP N historical dtl_ids pulled by category match in Stage A. "
                "Default 100. Increase for highly populated categories; lower "
                "for sparse ones."
            ),
        ),
        IntInput(
            name="bench_pinecone_top_k",
            display_name="V2 — Pinecone Top-K (Stage B)",
            value=10,
            advanced=True,
            info=(
                "Top K returned from Pinecone within-pool search; passed to "
                "Stage C as the candidate set for the relevance LLM."
            ),
        ),
        IntInput(
            name="bench_llm_shortlist_size",
            display_name="V2 — LLM Shortlist Size (Stage C)",
            value=10,
            advanced=True,
            info=(
                "How many candidates the rank LLM keeps after applying critical-"
                "attribute barriers and spec-to-price sanity checks. Cap, not "
                "target — the LLM never pads with weak matches."
            ),
        ),
        BoolInput(
            name="bench_widen_l1l2_when_sparse",
            display_name="V2 — Widen to L1+L2 When Sparse",
            value=True,
            advanced=True,
            info=(
                "True: when the category pool has <20 rows, ALSO accept rows "
                "matching item_level_1 AND item_level_2 (unioned)."
            ),
        ),
        MessageTextInput(
            name="bench_rank_debug_dump_dir",
            display_name="V2 — Rank LLM Debug Dump Directory",
            value="",
            advanced=True,
            info=(
                "DEBUG ONLY. Absolute path to a writable directory. When set, "
                "each Stage C call writes two files: the full input prompts "
                "(system + user + embedded source/candidates JSON) and the "
                "raw LLM response (before any parsing). Filenames include "
                "purchase_dtl_id + microsecond timestamp so nothing is "
                "overwritten. Leave EMPTY in production — the directory will "
                "grow indefinitely if left on."
            ),
        ),
        IntInput(
            name="bench_critical_threshold_pct",
            display_name="V2 — Critical Attribute Threshold (%)",
            value=10,
            advanced=True,
            info=(
                "Stage C rejects a candidate if any \"critical\" attribute differs "
                "from the source by more than this %. Lower = stricter."
            ),
        ),
        IntInput(
            name="bench_important_threshold_pct",
            display_name="V2 — Important Attribute Threshold (%)",
            value=20,
            advanced=True,
            info=(
                "Penalty band for \"important\" attribute mismatches in Stage C "
                "(does not auto-reject)."
            ),
        ),
        IntInput(
            name="bench_ratio_band_pct",
            display_name="V2 — Spec-to-Price Ratio Band (%)",
            value=50,
            advanced=True,
            info=(
                "Stage C downranks candidates whose unit_price / critical_attribute "
                "ratio is outside ±this% of the surviving-pool median."
            ),
        ),
        MultilineInput(
            name="bench_rank_prompt_system_v2",
            display_name="V2 — Stage C Rank System Prompt",
            value="",
            advanced=True,
            info=(
                "Override the system prompt used by the Stage C relevance LLM. "
                "Leave blank for the built-in default."
            ),
        ),
        MultilineInput(
            name="bench_rank_prompt_user_v2",
            display_name="V2 — Stage C Rank User Prompt",
            value="",
            advanced=True,
            info=(
                "Override the user prompt template (placeholders: {source_json}, "
                "{candidates_json}, {top_k}, {n_candidates}, {critical_pct}, "
                "{important_pct}, {ratio_pct})."
            ),
        ),
    ]

    outputs = [
        Output(display_name="File Batch", name="file_batch", method="build_file_batch", types=["Message"]),
        Output(display_name="Processed PRs", name="processed_prs", method="get_processed_prs", types=["Data"]),
    ]

    # ── Connection helpers ────────────────────────────────────────────────

    def _safe_log(self, message: str) -> None:
        """Log routed via AgentCore in the main component thread, falling
        back to loguru in ThreadPoolExecutor workers.

        AgentCore's Node.log resolves the active component via a contextvar
        that's only populated on the original asyncio.to_thread call. Worker
        threads we spawn (e.g. inside `_process_pr` or `_run_classification`)
        do not inherit that contextvar, so calling Node.log there emits
        "called add_log but no component context found" warnings and the
        message is dropped.

        This helper detects worker threads by name and routes them through
        loguru so logs still reach stderr — and the main-thread path keeps
        surfacing in the AgentCore canvas as before.
        """
        import threading
        # build_file_batch stamps _main_thread_id with the AgentCore-invoked
        # thread before any bg work fires. Any other thread (ThreadPoolExecutor
        # workers, the fire-and-forget pipeline daemon, etc.) routes to loguru
        # to avoid the "no component context" warning.
        main_tid = getattr(self, "_main_thread_id", None)
        if main_tid is None or threading.get_ident() != main_tid:
            logger.info(message)
            return
        # getattr keeps this call invisible to text-based search-and-replace,
        # so it stays as the real AgentCore Node.log invocation instead of
        # recursing back into _safe_log.
        _ac_log = getattr(self, "log")
        try:
            _ac_log(message)
        except Exception:
            logger.info(message)

    def _conn_str(self, conn_data: Data) -> str:
        d      = conn_data.data or {}
        driver = d.get("driver", "ODBC Driver 18 for SQL Server")
        server = d.get("host", d.get("server", ""))
        port   = d.get("port", 1433)
        db     = d.get("database_name", d.get("database", ""))
        user   = d.get("username", "")
        pwd    = d.get("password", "")
        return (
            f"DRIVER={{{driver}}};SERVER={server},{port};"
            f"DATABASE={db};UID={user};PWD={pwd};TrustServerCertificate=yes;"
        )

    @staticmethod
    def _is_transient(exc: Exception) -> bool:
        kw = ["connection reset", "timeout", "throttl", "resource limit", "broken pipe", "transport-level", "login failed"]
        return any(k in str(exc).lower() for k in kw)

    def _connect(self, conn_str: str):
        """Delegate to the module-level pool-backed _connect so the Node's
        own _connect calls share the same connection budget as every other
        helper in this file."""
        return _connect(conn_str)

    def _blob_cfg(self) -> dict:
        if not hasattr(self, "_blob_config_cache"):
            self._blob_config_cache = _get_blob_config_by_name(self.blob_connector_name)
        return self._blob_config_cache

    # ── Fetch pending PRs ─────────────────────────────────────────────────

    def _resolve_excel_pr_list(self) -> list[str] | None:
        """Try the three Excel-driven input sources, in priority order.

        Returns
        -------
        None
            No Excel-driven input wired/configured — caller falls back to
            _fetch_pending_prs.
        list[str]
            Deduplicated PURCHASE_REQ_NO values to force-reprocess. Empty
            list means an Excel WAS provided but contained zero rows; the
            caller treats that as "nothing to do".
        """
        # Source 1: HandleInput Data wired from a Blob Reader / KB / loader.
        # Accepted shapes — checked in order:
        #   data["pr_numbers"]  → already-parsed list, no Excel parsing needed
        #   data["content"] / ["bytes"] / ["file_bytes"]  → raw .xlsx bytes
        data_in = getattr(self, "pr_list_data", None)
        if data_in is not None:
            payload = getattr(data_in, "data", None) or {}
            if isinstance(payload, dict):
                pr_list = payload.get("pr_numbers") or payload.get("pr_nos")
                if isinstance(pr_list, list) and pr_list:
                    cleaned = [str(p).strip() for p in pr_list if str(p).strip()]
                    self._safe_log(
                        f"Excel PR list — using {len(cleaned)} PR(s) from "
                        f"wired Data.pr_numbers (no parsing needed)"
                    )
                    return cleaned
                for key in ("content", "bytes", "file_bytes"):
                    raw = payload.get(key)
                    if raw:
                        if isinstance(raw, str):
                            try:
                                import base64 as _b64
                                raw = _b64.b64decode(raw)
                            except Exception:
                                raw = raw.encode("latin-1", errors="ignore")
                        try:
                            pr_nos = _parse_pr_list_from_excel_bytes(
                                bytes(raw),
                                column=(self.pr_excel_column or None),
                                sheet=(self.pr_excel_sheet or None),
                            )
                            self._safe_log(
                                f"Excel PR list — parsed {len(pr_nos)} PR(s) "
                                f"from wired Data.{key}"
                            )
                            return pr_nos
                        except Exception as exc:
                            logger.opt(exception=True).error(
                                "Failed to parse Excel from wired Data.{}: {}", key, exc,
                            )
                            raise

        # Source 2: blob path under the configured Blob Connector container.
        blob_path = (getattr(self, "pr_excel_blob_path", None) or "").strip()
        if blob_path:
            self._safe_log(f"Excel PR list — downloading from blob: {blob_path!r}")
            try:
                blob_cfg    = self._blob_cfg()
                excel_bytes = _download_blob(blob_path, blob_cfg)
                pr_nos = _parse_pr_list_from_excel_bytes(
                    excel_bytes,
                    column=(self.pr_excel_column or None),
                    sheet=(self.pr_excel_sheet or None),
                )
                self._safe_log(
                    f"Excel PR list — parsed {len(pr_nos)} PR(s) from {blob_path!r}"
                )
                return pr_nos
            except Exception as exc:
                logger.opt(exception=True).error(
                    "Failed to download/parse Excel from blob {!r}: {}", blob_path, exc,
                )
                raise

        return None

    def _batch_reset_for_reprocess(self, tgt_cs: str, pr_nos: list[str]) -> int:
        """Reset many tracker rows at once for a forced reprocess.

        Equivalent to looping _reset_for_reprocess(pr) for each PR but in
        2 statements per chunk instead of 2 round-trips per PR. Speeds up
        Excel-driven runs by O(N) DB calls.

        Steps per chunk (single transaction):
          1. DELETE ras_pipeline_exceptions for any tracker row in this set
          2. UPDATE ras_tracker SET current_stage_fk=NULL, retry_count+=1

        Returns number of PRs whose tracker row was actually updated. PRs
        without a tracker row are silently skipped (UPDATE rowcount=0)
        and processed as brand-new in _process_pr.
        """
        if not pr_nos:
            return 0
        BATCH = 500
        total_updated = 0
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            for i in range(0, len(pr_nos), BATCH):
                chunk = pr_nos[i : i + BATCH]
                placeholders = ",".join(["?"] * len(chunk))
                cur.execute(
                    f"""
                    DELETE FROM [ras_procurement].[ras_pipeline_exceptions]
                     WHERE ras_tracker_id IN (
                         SELECT ras_uuid_pk
                           FROM [ras_procurement].[ras_tracker]
                          WHERE purchase_req_no IN ({placeholders})
                     )
                    """,
                    *chunk,
                )
                cur.execute(
                    f"""
                    UPDATE [ras_procurement].[ras_tracker]
                       SET current_stage_fk = NULL,
                           retry_count      = COALESCE(retry_count, 0) + 1,
                           updated_at       = SYSUTCDATETIME()
                     WHERE purchase_req_no IN ({placeholders})
                    """,
                    *chunk,
                )
                total_updated += cur.rowcount
            conn.commit()
        except Exception:
            try:
                conn.rollback()
            except Exception:
                pass
            raise
        finally:
            conn.close()
        return total_updated

    def _sort_excel_pr_list_by_date(self, tgt_cs: str, pr_nos: list[str]) -> list[str]:
        """Return pr_nos reordered by purchase_req_mst.C_DATETIME ASC.

        Used in Excel mode to reduce the Stage 7 benchmark race where a
        newer PR's _run_benchmark can fire before an older PR finishes
        Stage 6 (embeddings) — leaving the newer PR's benchmark with
        thin historical context. Submitting the worker pool oldest-first
        means oldest PRs reach Pinecone first; by the time a newer PR's
        Stage 7 runs, the older PRs' items have already been upserted.

        Partial fix only: with heavily uneven PR sizes a small new PR
        can still overtake a large old one. Acceptable for typical
        workloads; for full guarantees the two-phase approach is needed.

        PRs missing from purchase_req_mst (will fail with
        NoRASContextError during processing anyway) are appended at the
        end in their original Excel order so they don't disturb the
        sort of valid PRs.
        """
        if len(pr_nos) <= 1:
            return list(pr_nos)

        BATCH = 500
        date_by_pr: dict[str, object] = {}
        conn = self._connect(tgt_cs)
        try:
            cur = conn.cursor()
            for i in range(0, len(pr_nos), BATCH):
                chunk = pr_nos[i : i + BATCH]
                placeholders = ",".join(["?"] * len(chunk))
                cur.execute(
                    f"SELECT [PURCHASE_REQ_NO], [C_DATETIME] "
                    f"FROM [ras_procurement].[purchase_req_mst] "
                    f"WHERE [PURCHASE_REQ_NO] IN ({placeholders})",
                    *chunk,
                )
                for row in cur.fetchall():
                    date_by_pr[row[0]] = row[1]
        finally:
            conn.close()

        found    = [pr for pr in pr_nos if date_by_pr.get(pr) is not None]
        missing  = [pr for pr in pr_nos if date_by_pr.get(pr) is None]
        found.sort(key=lambda p: date_by_pr[p])
        return found + missing

    def _fetch_current_stages(self, tgt_cs: str, pr_nos: list[str]) -> dict[str, int | None]:
        """Return {pr_no → current_stage_fk} for the given PRs.

        Missing tracker rows are mapped to None. Used by Excel-driven runs to
        skip PRs that already reached COMPLETE (stage 8) without re-fetching
        them per worker.
        """
        if not pr_nos:
            return {}
        # SQL Server caps IN(...) at ~2100 placeholders; chunk in batches of
        # 500 to be well under that limit.
        result: dict[str, int | None] = {pr: None for pr in pr_nos}
        BATCH = 500
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            for i in range(0, len(pr_nos), BATCH):
                chunk = pr_nos[i : i + BATCH]
                placeholders = ",".join(["?"] * len(chunk))
                cur.execute(
                    f"SELECT purchase_req_no, current_stage_fk "
                    f"FROM [ras_procurement].[ras_tracker] "
                    f"WHERE purchase_req_no IN ({placeholders})",
                    *chunk,
                )
                for row in cur.fetchall():
                    pr, stage = row[0], row[1]
                    result[pr] = int(stage) if stage is not None else None
        finally:
            conn.close()
        return result

    def _fetch_pending_prs(self, tgt_cs: str) -> list[str]:
        pr_filter = (self.pr_no_filter or "").strip()
        if pr_filter:
            return [pr_filter]
        limit = int(self.batch_limit)
        conn  = self._connect(tgt_cs)
        cur   = conn.cursor()
        # Stage 1-3 only owns stages 1 and 2.  PRs already at stage 3+ have
        # completed Stage 1-3 and belong to Stage 4-8 — do NOT re-pick them.
        # Only pick: new PRs (no tracker row), reset PRs (NULL stage), or PRs
        # stuck mid-Stage-1-3 (stages 1 or 2).
        _WHERE = """
            FROM [ras_procurement].[purchase_req_mst] prm
            LEFT JOIN [ras_procurement].[ras_tracker] rt
              ON prm.PURCHASE_REQ_NO = rt.purchase_req_no
           WHERE (rt.purchase_req_no IS NULL
                  OR rt.current_stage_fk IS NULL
                  OR rt.current_stage_fk IN (1, 2))
             AND UPPER(prm.PURCHASEFINALAPPROVALSTATUS)
                     IN ('APPROVED BY ALL', 'APPROVED BY ALL EXCEPTION')
           ORDER BY prm.C_DATETIME ASC
        """
        try:
            if limit > 0:
                cur.execute(f"SELECT TOP (?) prm.PURCHASE_REQ_NO {_WHERE}", limit)
            else:
                cur.execute(f"SELECT prm.PURCHASE_REQ_NO {_WHERE}")
            return [row[0] for row in cur.fetchall()]
        finally:
            conn.close()

    # ── Pre-run cleanup ───────────────────────────────────────────────────

    def _cleanup_for_pr(self, tgt_cs: str, pr_no: str) -> None:
        """Delete all prior pipeline output for a PR — runs only if the PR already
        exists in ras_tracker (i.e. it has been processed before). New PRs are
        skipped entirely so their data is never touched."""
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?", pr_no)
        existing = cur.fetchone()
        conn.close()
        if existing is None:
            # Brand-new PR — nothing to clean, proceed straight to Stage 1.
            self._safe_log(f"[{pr_no}] New PR — skipping cleanup")
            return
        self._safe_log(f"[{pr_no}] Existing PR detected — cleaning prior data before reprocessing")

        # Single connection for all DB cleanup work (collect IDs + delete tables)
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        pinecone_ids: list[str] = []
        try:
            # Collect Pinecone vector IDs FIRST, before deleting quotation_extracted_items.
            # We build BOTH formats so stale vectors from any run are cleaned:
            #   • old format: extracted_item_uuid_pk (used by pre-dtl_ code)
            #   • new format: dtl_{purchase_dtl_id}   (current code, stable per line item)
            cur.execute("""
                SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id], qi.[is_selected_quote]
                  FROM [ras_procurement].[quotation_extracted_items] qi
                  JOIN [ras_procurement].[attachment_classification] ac
                    ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
                  JOIN [ras_procurement].[ras_tracker] rt
                    ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
                 WHERE rt.[purchase_req_no] = ?
            """, pr_no)
            rows = cur.fetchall()
            selected_count = sum(1 for r in rows if r[2])
            # New format: only is_selected_quote=1 rows (current code only embeds these)
            new_ids = list({f"dtl_{r[1]}" for r in rows if r[1] and r[2]})
            # Old format: only is_selected_quote=1 UUIDs (legacy vectors only existed for selected rows too)
            old_ids = [str(r[0]) for r in rows if r[0] and r[2]]
            pinecone_ids = list(set(old_ids + new_ids))
            self._safe_log(
                f"[{pr_no}] Pinecone cleanup: {len(rows)} extracted rows, "
                f"{selected_count} is_selected_quote=1 — "
                f"deleting {len(pinecone_ids)} vector(s): "
                f"{len(new_ids)} dtl_ ID(s) + {len(old_ids)} legacy UUID(s)"
            )

            # 1. Null FK back-references in benchmark_result
            cur.execute("""
                UPDATE br SET br.low_hist_item_fk = NULL
                  FROM [ras_procurement].[benchmark_result] br
                  JOIN [ras_procurement].[quotation_extracted_items] qi ON br.low_hist_item_fk = qi.extracted_item_uuid_pk
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            cur.execute("""
                UPDATE br SET br.last_hist_item_fk = NULL
                  FROM [ras_procurement].[benchmark_result] br
                  JOIN [ras_procurement].[quotation_extracted_items] qi ON br.last_hist_item_fk = qi.extracted_item_uuid_pk
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 2. benchmark_result
            cur.execute("""
                DELETE br FROM [ras_procurement].[benchmark_result] br
                  JOIN [ras_procurement].[quotation_extracted_items] qi ON br.extracted_item_uuid_fk = qi.extracted_item_uuid_pk
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 3. quotation_extracted_items
            cur.execute("""
                DELETE qi FROM [ras_procurement].[quotation_extracted_items] qi
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 4. embedded_attachment_classification
            cur.execute("""
                DELETE ec FROM [ras_procurement].[embedded_attachment_classification] ec
                  JOIN [ras_procurement].[attachment_classification] ac ON ec.attachment_classification_id = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 5. attachment_classification
            cur.execute("""
                DELETE ac FROM [ras_procurement].[attachment_classification] ac
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 6. BI dashboard rows
            cur.execute("""
                DELETE FROM [ras_procurement].[vw_get_ras_data_for_bidashboard]
                 WHERE [PURCHASE_REQ_NO] = ?""", pr_no)
            # 7. Exception records
            cur.execute("""
                DELETE FROM [ras_procurement].[ras_pipeline_exceptions]
                 WHERE ras_tracker_id = (
                     SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker]
                      WHERE purchase_req_no = ?)""", pr_no)
            # 8. Any remaining data via SP
            cur.execute("EXEC [ras_procurement].[usp_cleanup_pr_data] ?", pr_no)
            conn.commit()
        except Exception as exc:
            conn.rollback()
            raise RuntimeError(f"Pre-run cleanup failed for {pr_no}: {exc}") from exc
        finally:
            conn.close()

        # Delete stale Pinecone vectors — CRITICAL (must succeed)
        if pinecone_ids:
            pinecone_index = (getattr(self, "pinecone_index", None) or "").strip()
            pinecone_ns    = (getattr(self, "pinecone_namespace", None) or "").strip()
            if pinecone_index and pinecone_ns:
                try:
                    from agentcore.services.pinecone_service_client import delete_vectors_via_service
                    preview = pinecone_ids[:5]
                    more    = len(pinecone_ids) - 5
                    self._safe_log(
                        f"[{pr_no}] Sending {len(pinecone_ids)} ID(s) to Pinecone delete "
                        f"(index={pinecone_index!r}, namespace={pinecone_ns!r}) — "
                        f"first 5: {preview}" + (f" … +{more} more" if more > 0 else "")
                    )
                    delete_vectors_via_service(
                        index_name=pinecone_index, namespace=pinecone_ns,
                        vector_ids=pinecone_ids,
                    )
                    self._safe_log(f"[{pr_no}] Pinecone delete completed — {len(pinecone_ids)} ID(s) sent for removal")
                except Exception as exc:
                    error_msg = (
                        f"Pinecone vector cleanup FAILED. Cannot delete old embeddings. "
                        f"Reason: {str(exc)}. "
                        f"Check: (1) Pinecone API key is valid, "
                        f"(2) Index name '{pinecone_index}' and namespace '{pinecone_ns}' exist, "
                        f"(3) Network connectivity to Pinecone. "
                        f"PR cannot proceed until cleanup succeeds."
                    )
                    self._safe_log(f"[{pr_no}] {error_msg}")
                    logger.opt(exception=True).error(f"[{pr_no}] {error_msg}")

                    # Record exception in DB
                    try:
                        self._record_exception(tgt_cs, pr_no, 99, error_msg)
                    except Exception as db_exc:
                        logger.warning(f"[{pr_no}] Could not record Pinecone exception in DB: {db_exc}")

                    # Raise to break pipeline for this PR
                    raise RuntimeError(f"[{pr_no}] Pinecone cleanup failed — PR cannot proceed") from exc
            else:
                self._safe_log(
                    f"[{pr_no}] Warning — pinecone_index or pinecone_namespace not set; "
                    f"{len(pinecone_ids)} vector(s) were NOT deleted from Pinecone. "
                    f"Check the 'Pinecone Index' and 'Pinecone Namespace' inputs on the component."
                )
        else:
            self._safe_log(f"[{pr_no}] No Pinecone vectors found in DB for this PR — skipping Pinecone delete")

        # Delete Azure Blob folder — CRITICAL (must succeed)
        try:
            self._delete_blob_folder(pr_no)
        except Exception as exc:
            error_msg = (
                f"Azure Blob Storage cleanup FAILED. Cannot delete old files. "
                f"Reason: {str(exc)}. "
                f"Check: (1) Azure CLI is installed and logged in (run: az login), "
                f"(2) Storage account access permissions, "
                f"(3) Network connectivity to Azure. "
                f"PR cannot proceed until cleanup succeeds."
            )
            self._safe_log(f"[{pr_no}] {error_msg}")
            logger.opt(exception=True).error(f"[{pr_no}] {error_msg}")

            # Record exception in DB
            try:
                self._record_exception(tgt_cs, pr_no, 99, error_msg)
            except Exception as db_exc:
                logger.warning(f"[{pr_no}] Could not record blob exception in DB: {db_exc}")

            # Raise to break pipeline for this PR
            raise RuntimeError(f"[{pr_no}] Blob cleanup failed — PR cannot proceed") from exc

        self._safe_log(f"[{pr_no}] Cleanup complete — all pipeline tables cleared")

    def _reset_for_reprocess(self, tgt_cs: str, pr_no: str) -> None:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("""
                DELETE FROM [ras_procurement].[ras_pipeline_exceptions]
                 WHERE ras_tracker_id = (SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?)
            """, pr_no)
            cur.execute("""
                UPDATE [ras_procurement].[ras_tracker]
                   SET current_stage_fk = NULL, retry_count = COALESCE(retry_count,0)+1, updated_at = SYSUTCDATETIME()
                 WHERE purchase_req_no = ?
            """, pr_no)
            conn.commit()
            self._safe_log(f"[{pr_no}] Tracker reset for reprocess")
        finally:
            conn.close()

    # ── Fetch attachments (Azure SQL for IDs+names, on-prem for binary) ───

    def _fetch_attachments(self, src_cs: str, tgt_cs: str, pr_no: str) -> list[dict]:
        tgt = self._connect(tgt_cs)
        cur = tgt.cursor()
        try:
            cur.execute("""
                SELECT pa.[ATTACHMENT_ID], pa.[FILES_NAME]
                  FROM [ras_procurement].[purchase_req_mst] prm
                  JOIN [ras_procurement].[purchase_attachments] pa ON prm.[PURCHASE_REQ_ID] = pa.[PURCHASE_ID]
                 WHERE prm.[PURCHASE_REQ_NO] = ?
                   AND pa.[ATTACHMENT_ID] IS NOT NULL AND pa.[FILES_NAME] IS NOT NULL
            """, pr_no)
            rows = cur.fetchall()
        finally:
            tgt.close()
        if not rows:
            return []
        id_to_filename = {str(r[0]): r[1] for r in rows}
        att_ids        = list(id_to_filename.keys())
        placeholders   = ",".join(["?"] * len(att_ids))
        src = self._connect(src_cs)
        cur = src.cursor()
        try:
            cur.execute(
                f"SELECT [attachment_id], [doc] FROM [dbo].[ras_attachments]"
                f" WHERE [attachment_id] IN ({placeholders}) AND [doc] IS NOT NULL",
                att_ids,
            )
            return [
                {
                    "attachment_id": str(r[0]),
                    "filename":      id_to_filename.get(str(r[0]), f"attachment_{r[0]}"),
                    "content":       bytes(r[1]),
                    "pr_no":         pr_no,
                }
                for r in cur.fetchall()
            ]
        finally:
            src.close()

    # ── Tracker helpers ───────────────────────────────────────────────────

    def _get_tracker_uuid(self, tgt_cs: str, pr_no: str) -> str:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?", pr_no)
            row = cur.fetchone()
            return str(row[0]) if row else ""
        finally:
            conn.close()

    def _advance_tracker(self, tgt_cs: str, pr_no: str, stage_id: int) -> None:
        # Wrapped in deadlock retry — when many parallel workers all advance
        # their own tracker rows the shared HOLDLOCK / UPDATE locks can race
        # and SQL Server picks a victim. Retry with backoff per the canvas
        # knobs (db_deadlock_max_retries / db_deadlock_base_delay).
        def _do() -> None:
            conn = self._connect(tgt_cs)
            cur  = conn.cursor()
            try:
                if stage_id == _STAGE_INGESTION:
                    cur.execute("""
                        MERGE [ras_procurement].[ras_tracker] WITH (HOLDLOCK) AS target
                        USING (
                            SELECT PURCHASE_REQ_NO, PURCHASEFINALAPPROVALSTATUS
                              FROM [ras_procurement].[purchase_req_mst] WHERE PURCHASE_REQ_NO = ?
                        ) AS src ON target.purchase_req_no = src.PURCHASE_REQ_NO
                        WHEN MATCHED THEN UPDATE SET current_stage_fk=?, updated_at=SYSUTCDATETIME()
                        WHEN NOT MATCHED THEN INSERT (purchase_req_no, ras_status, current_stage_fk)
                            VALUES (src.PURCHASE_REQ_NO, src.PURCHASEFINALAPPROVALSTATUS, ?);
                    """, pr_no, stage_id, stage_id)
                else:
                    cur.execute("""
                        UPDATE [ras_procurement].[ras_tracker]
                           SET current_stage_fk=?, updated_at=SYSUTCDATETIME()
                         WHERE purchase_req_no=?
                    """, stage_id, pr_no)
                conn.commit()
            except Exception:
                try:
                    conn.rollback()
                except Exception:
                    pass
                raise
            finally:
                conn.close()

        self._run_with_deadlock_retry(
            _do,
            op_name=f"advance_tracker(stage={stage_id})",
            pr_no=pr_no,
        )

    def _run_with_deadlock_retry(
        self,
        operation,
        *,
        op_name: str,
        pr_no: str,
        max_retries: int | None = None,
        base_delay: float | None = None,
    ):
        """Run a callable that touches the DB with deadlock-aware retries.

        SQL Server picks one transaction as the deadlock victim (error 1205
        / SQLSTATE 40001) when concurrent workers contend on shared rows.
        The server's documented guidance is to rerun the transaction, so
        this helper does that with exponential backoff + ±30% jitter so
        multiple victims don't synchronise their retries.

        Logs are deliberately verbose so the user can see, for any given
        PR, exactly how many retries were attempted, what happened, and
        whether it ultimately succeeded or gave up:

          • on victim   : warning showing attempt N/M, error code, sleep
          • on success after retry : success line "recovered after N retry"
          • on exhausted retries   : error line "giving up after M attempt(s)"
          • on non-deadlock error  : raises immediately (caller decides)
        """
        # Default knobs from canvas inputs when caller doesn't override.
        if max_retries is None:
            max_retries = max(0, int(getattr(self, "db_deadlock_max_retries", 3) or 0))
        if base_delay is None:
            base_delay  = max(0.5, float(getattr(self, "db_deadlock_base_delay", 2) or 2))

        attempt = 0
        while True:
            try:
                result = operation()
                if attempt > 0:
                    self._safe_log(
                        f"[{pr_no}] {op_name} recovered after {attempt} retry/retries"
                    )
                return result
            except Exception as exc:
                if not _is_deadlock_error(exc):
                    raise
                if attempt >= max_retries:
                    self._safe_log(
                        f"[{pr_no}] {op_name} — DEADLOCK retries exhausted "
                        f"(gave up after {max_retries + 1} attempt(s)): {exc}"
                    )
                    logger.opt(exception=True).error(
                        "[{}] {} — DEADLOCK retries exhausted: {}",
                        pr_no, op_name, exc,
                    )
                    raise
                attempt += 1
                delay = base_delay * (2 ** attempt) * (1 + random.random() * 0.3)
                self._safe_log(
                    f"[{pr_no}] {op_name} — DEADLOCK victim "
                    f"(attempt {attempt}/{max_retries}) — "
                    f"sleeping {delay:.1f}s then retrying"
                )
                logger.warning(
                    "[{}] {} hit SQL Server deadlock (1205); "
                    "retry {}/{} in {:.1f}s. underlying: {}",
                    pr_no, op_name, attempt, max_retries, delay, exc,
                )
                time.sleep(delay)

    def _record_exception(self, tgt_cs: str, pr_no: str, stage_id: int, error_msg: str) -> None:
        # Wrapped in deadlock retry — the MERGE on ras_tracker + INSERT on
        # ras_pipeline_exceptions touches the same shared chain that
        # _cleanup_for_pr does, and concurrent failure recording on
        # different PRs can pick this transaction as the deadlock victim.
        # Retries follow the canvas knobs (db_deadlock_max_retries /
        # db_deadlock_base_delay).
        def _do() -> None:
            conn = self._connect(tgt_cs)
            cur  = conn.cursor()
            try:
                # MERGE upsert. Updates current_stage_fk to 99 (EXCEPTION) and
                # also stamps last_processed_at so a PR's terminal time is
                # recorded for both success and failure paths (mirrors what
                # _set_last_processed_at does on Stage 8 success).
                cur.execute("""
                    MERGE [ras_procurement].[ras_tracker] WITH (HOLDLOCK) AS target
                    USING (SELECT PURCHASE_REQ_NO, PURCHASEFINALAPPROVALSTATUS
                             FROM [ras_procurement].[purchase_req_mst] WHERE PURCHASE_REQ_NO = ?
                    ) AS src ON target.purchase_req_no = src.PURCHASE_REQ_NO
                    WHEN MATCHED THEN UPDATE SET
                        current_stage_fk  = 99,
                        last_processed_at = SYSUTCDATETIME(),
                        updated_at        = SYSUTCDATETIME()
                    WHEN NOT MATCHED THEN INSERT
                        (purchase_req_no, ras_status, current_stage_fk, last_processed_at)
                        VALUES (
                            src.PURCHASE_REQ_NO,
                            src.PURCHASEFINALAPPROVALSTATUS,
                            99,
                            SYSUTCDATETIME()
                        );
                """, pr_no)
                cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no=?", pr_no)
                row = cur.fetchone()
                if row:
                    cur.execute("""
                        INSERT INTO [ras_procurement].[ras_pipeline_exceptions] (ras_tracker_id, stage_id, exception_message)
                        VALUES (?, ?, ?)
                    """, row[0], stage_id, error_msg[:4000])
                conn.commit()
            except Exception:
                try:
                    conn.rollback()
                except Exception:
                    pass
                raise
            finally:
                conn.close()

        try:
            self._run_with_deadlock_retry(
                _do,
                op_name=f"record_exception(stage={stage_id})",
                pr_no=pr_no,
            )
        except Exception as exc:
            # Even after retries the exception write failed. Don't propagate
            # — the per-PR processing path is already in an error state and
            # we don't want a secondary failure to mask the original cause.
            logger.warning(f"[{pr_no}] Could not write exception record: {exc}")

    # ── attachment_classification upserts ─────────────────────────────────

    def _upsert_parent_classification(self, tgt_cs, pr_no, ras_uuid_pk, att_id, file_path, embedded_flag, embedded_count) -> str:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("""
                MERGE [ras_procurement].[attachment_classification] WITH (HOLDLOCK) AS target
                USING (SELECT ? AS ras_uuid_pk, ? AS attachment_id) AS src
                  ON target.[attachment_id] = src.[attachment_id]
                WHEN MATCHED THEN
                    UPDATE SET [file_path]=?, [embedded_file_flag]=?, [embedded_file_count]=?, [updated_at]=SYSUTCDATETIME()
                WHEN NOT MATCHED THEN
                    INSERT ([ras_uuid_pk],[attachment_id],[file_path],[embedded_file_flag],[embedded_file_count])
                    VALUES (?,?,?,?,?);
            """, ras_uuid_pk, att_id,
                file_path, embedded_flag, embedded_count,
                ras_uuid_pk, att_id, file_path, embedded_flag, embedded_count)
            cur.execute(
                "SELECT [attachment_classify_uuid_pk] FROM [ras_procurement].[attachment_classification] WHERE [attachment_id]=?",
                att_id,
            )
            row = cur.fetchone()
            conn.commit()
            return str(row[0]) if row else ""
        finally:
            conn.close()

    def _upsert_embedded_classification(self, tgt_cs, parent_pk, parent_att_id, file_path) -> None:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("""
                MERGE [ras_procurement].[embedded_attachment_classification] WITH (HOLDLOCK) AS target
                USING (SELECT ? AS attachment_classification_id, ? AS file_path) AS src
                  ON target.[attachment_classification_id] = src.[attachment_classification_id]
                 AND target.[file_path] = src.[file_path]
                WHEN MATCHED THEN UPDATE SET [parent_attachment_id]=?, [updated_at]=SYSUTCDATETIME()
                WHEN NOT MATCHED THEN
                    INSERT ([attachment_classification_id],[parent_attachment_id],[file_path])
                    VALUES (?,?,?);
            """, parent_pk, file_path,
                parent_att_id,
                parent_pk, parent_att_id, file_path)
            conn.commit()
        finally:
            conn.close()

    # ── Text extraction ───────────────────────────────────────────────────

    def _extract_text(self, filename: str, raw: bytes) -> str:
        import os
        ext = os.path.splitext(filename.lower())[1]
        try:
            if ext in _TEXT_EXT:  return raw.decode("utf-8", errors="replace")
            if ext in _PDF_EXT:   return _extract_pdf(raw)
            if ext in _EXCEL_EXT: return _extract_excel(raw)
            if ext in _WORD_EXT:  return _extract_word(raw)
            if ext in _PPTX_EXT:  return _extract_pptx(raw)
            if ext in _MSG_EXT:   return _extract_msg_text(raw)
            return f"[Binary — {ext} not text-extractable]"
        except Exception as exc:
            return f"[Extraction error: {exc}]"

    # ── BI dashboard sync ─────────────────────────────────────────────────

    def _sync_bi_dashboard(self, src_cs: str, tgt_cs: str, pr_no: str) -> None:
        """Refresh the BI dashboard row for this PR.

        Reads from on-prem view [dbo].[vw_get_ras_data_for_bidashboard],
        deletes the old Azure row, inserts fresh data.
        Non-fatal — caller wraps in try/except.
        """
        import pyodbc
        # ── 1. read from on-prem ──────────────────────────────────────
        src_conn = self._connect(src_cs)
        try:
            cur = src_conn.cursor()
            cur.execute(
                "SELECT * FROM [dbo].[vw_get_ras_data_for_bidashboard] "
                "WHERE [PURCHASE_REQ_NO] = ?",
                pr_no,
            )
            columns = [d[0] for d in cur.description] if cur.description else []
            rows    = cur.fetchall()
        finally:
            src_conn.close()

        # ── 2. delete old + insert fresh in Azure ─────────────────────
        tgt_conn = self._connect(tgt_cs)
        try:
            tc = tgt_conn.cursor()
            tc.execute(
                "DELETE FROM [ras_procurement].[vw_get_ras_data_for_bidashboard] "
                "WHERE [PURCHASE_REQ_NO] = ?",
                pr_no,
            )
            deleted = tc.rowcount
            if rows and columns:
                col_list     = ", ".join(f"[{c}]" for c in columns)
                placeholders = ", ".join(["?"] * len(columns))
                ins_sql = (
                    f"INSERT INTO [ras_procurement].[vw_get_ras_data_for_bidashboard] "
                    f"({col_list}) VALUES ({placeholders})"
                )
                tc.fast_executemany = True
                tc.executemany(ins_sql, [list(r) for r in rows])
            tgt_conn.commit()
            self._safe_log(
                f"[{pr_no}] BI dashboard synced — "
                f"deleted {deleted}, inserted {len(rows)} row(s)"
            )
        except Exception:
            tgt_conn.rollback()
            raise
        finally:
            tgt_conn.close()

    # ── Blob helpers ──────────────────────────────────────────────────────

    def _container_client(self):
        """Return a cached ContainerClient — credential is initialised once per component instance."""
        if not hasattr(self, "_container_client_cache"):
            from azure.storage.blob import BlobServiceClient
            cfg        = self._blob_cfg()
            credential = _get_az_credential()
            self._container_client_cache = BlobServiceClient(
                account_url=cfg["account_url"],
                credential=credential,
            ).get_container_client(cfg["container_name"])
        return self._container_client_cache

    def _invalidate_container_client(self) -> None:
        """Drop the cached ContainerClient so the next call rebuilds it with
        a freshly-acquired credential. Called by _with_az_retry after a
        transient token-refresh failure — otherwise the cached client would
        keep using the stale credential it was constructed with."""
        cached = getattr(self, "_container_client_cache", None)
        if cached is None:
            return
        try:
            if hasattr(cached, "close"):
                cached.close()
        except Exception:
            pass
        try:
            del self._container_client_cache
        except Exception:
            pass

    def _delete_blob_folder(self, pr_no: str) -> None:
        """Delete every blob under procurement/{pr_no}/ in Azure Blob Storage."""
        safe_pr = pr_no.replace("/", "_")
        prefix  = f"procurement/{safe_pr}/"

        def _do() -> int:
            cc    = self._container_client()
            blobs = [b.name for b in cc.list_blobs(name_starts_with=prefix)]
            if not blobs:
                return 0
            _BATCH = 256
            for i in range(0, len(blobs), _BATCH):
                cc.delete_blobs(*blobs[i : i + _BATCH])
            return len(blobs)

        deleted = _with_az_retry(_do, on_rebuild=self._invalidate_container_client)
        if deleted:
            self._safe_log(f"[{pr_no}] Deleted {deleted} blob(s) from {prefix}")

    def _upload_blob(self, raw: bytes, blob_path: str) -> None:
        def _do() -> None:
            self._container_client().get_blob_client(blob_path).upload_blob(raw, overwrite=True)

        _with_az_retry(_do, on_rebuild=self._invalidate_container_client)

    # ── Source-change detection ───────────────────────────────────────────

    def _detect_and_requeue_changed_prs(self, tgt_cs: str) -> int:
        """Re-queue completed PRs whose on-prem source data changed since the
        pipeline last processed them.

        Compares purchase_req_mst.U_DATETIME (bumped by on-prem DB triggers
        whenever purchase_req_detail or purchase_attachments rows change) against
        ras_tracker.last_processed_at (stamped at Stage 8 completion).

        Re-queued PRs have current_stage_fk reset to 1 so _fetch_pending_prs
        picks them up on this run. Non-fatal — a failure is logged and skipped.
        """
        try:
            conn = self._connect(tgt_cs)
            cur  = conn.cursor()
            try:
                cur.execute("""
                    SELECT prm.PURCHASE_REQ_NO
                      FROM [ras_procurement].[purchase_req_mst] prm
                      JOIN [ras_procurement].[ras_tracker] rt
                        ON rt.purchase_req_no = prm.PURCHASE_REQ_NO
                     WHERE rt.current_stage_fk  = 8
                       AND rt.last_processed_at IS NOT NULL
                       AND prm.U_DATETIME        > rt.last_processed_at
                       AND UPPER(prm.PURCHASEFINALAPPROVALSTATUS)
                               IN ('APPROVED BY ALL', 'APPROVED BY ALL EXCEPTION')
                """)
                changed = [row[0] for row in cur.fetchall()]
                if not changed:
                    return 0
                for pr_no in changed:
                    cur.execute("""
                        UPDATE [ras_procurement].[ras_tracker]
                           SET current_stage_fk = 1,
                               retry_count      = COALESCE(retry_count, 0) + 1,
                               updated_at       = SYSUTCDATETIME()
                         WHERE purchase_req_no  = ?
                    """, pr_no)
                conn.commit()
                self._safe_log(f"Source-change detection: re-queued {len(changed)} PR(s) — {changed}")
                return len(changed)
            finally:
                conn.close()
        except Exception as exc:
            self._safe_log(f"Warning — source-change detection failed (non-fatal): {exc}")
            return 0

    # ── Full-pipeline continuation (stages 4-8) ──────────────────────────

    def _build_prompts(self) -> dict:
        """Build prompt overrides and tuning-knob dict from component inputs.
        Prompt keys: cls_system, cls_user_text, cls_user_image, ext_system, ext_user, ext_taxonomy.
        Threshold keys: cls_max_chars, ext_max_chars, selected_threshold, price_max_boost, canonicalize_threshold."""
        prompts: dict = {}
        _p = getattr(self, "cls_system_prompt",    None)
        if _p: prompts["cls_system"]    = _get_prompt_text(_p, CLASSIFICATION_SYSTEM_PROMPT)
        _p = getattr(self, "cls_user_text_prompt", None)
        if _p: prompts["cls_user_text"] = _get_prompt_text(_p, _CLASSIFY_USER_TEXT)
        _p = getattr(self, "cls_user_image_prompt", None)
        if _p: prompts["cls_user_image"] = _get_prompt_text(_p, _CLASSIFY_USER_IMAGE)
        _p = getattr(self, "ext_system_prompt",    None)
        if _p: prompts["ext_system"]    = _get_prompt_text(_p, EXTRACTION_SYSTEM_PROMPT)
        _p = getattr(self, "ext_user_template",    None)
        if _p: prompts["ext_user"]      = _get_prompt_text(_p, EXTRACTION_USER_TEMPLATE)
        _p = getattr(self, "ext_item_taxonomy",    None)
        if _p: prompts["ext_taxonomy"]  = _get_prompt_text(_p, ITEM_TAXONOMY)
        # ── Tuning knobs ──────────────────────────────────────────────────
        cls_max = getattr(self, "cls_max_chars", None)
        if cls_max: prompts["cls_max_chars"] = int(cls_max)
        cls_max_vision = getattr(self, "cls_max_pages_vision", None)
        if cls_max_vision is not None: prompts["cls_max_pages_vision"] = int(cls_max_vision)
        ext_max = getattr(self, "ext_max_chars", None)
        if ext_max: prompts["ext_max_chars"] = int(ext_max)
        ext_max_pages = getattr(self, "ext_max_pages", None)
        if ext_max_pages is not None: prompts["ext_max_pages"] = int(ext_max_pages)
        ext_max_images = getattr(self, "ext_max_images", None)
        if ext_max_images is not None: prompts["ext_max_images"] = int(ext_max_images)
        try:
            prompts["selected_threshold"] = Decimal(str(float(getattr(self, "selected_threshold", None) or "0.70")))
        except Exception:
            prompts["selected_threshold"] = _SELECTED_THRESHOLD
        try:
            prompts["price_max_boost"] = Decimal(str(float(getattr(self, "price_max_boost", None) or "0.10")))
        except Exception:
            prompts["price_max_boost"] = _PRICE_MAX_BOOST
        try:
            prompts["canonicalize_threshold"] = float(getattr(self, "canonicalize_threshold", None) or "0.82")
        except Exception:
            prompts["canonicalize_threshold"] = _CANONICALIZE_THRESHOLD
        # ── Benchmark filtering / quality knobs ──────────────────────────────
        try:
            prompts["bench_min_similarity"] = float(getattr(self, "bench_min_similarity", None) or "0.80")
        except Exception:
            prompts["bench_min_similarity"] = 0.80
        try:
            prompts["bench_outlier_factor"] = float(getattr(self, "bench_outlier_factor", None) or "3.0")
        except Exception:
            prompts["bench_outlier_factor"] = 3.0
        try:
            prompts["bench_max_age_months"] = int(getattr(self, "bench_max_age_months", None) or 0)
        except Exception:
            prompts["bench_max_age_months"] = 0
        try:
            prompts["bench_uom_strict"] = bool(int(getattr(self, "bench_uom_strict", None) or 0))
        except Exception:
            prompts["bench_uom_strict"] = True
        # ── LLM retry / rate-limit settings ──────────────────────────────────
        llm_retries = getattr(self, "llm_max_retries", None)
        if llm_retries is not None: prompts["llm_max_retries"]  = int(llm_retries)
        llm_cool    = getattr(self, "llm_retry_cooldown", None)
        if llm_cool is not None:    prompts["llm_retry_cooldown"] = int(llm_cool)
        cls_par     = getattr(self, "cls_parallel_sources", None)
        if cls_par is not None:     prompts["cls_parallel_sources"] = int(cls_par)
        ext_par     = getattr(self, "ext_parallel_sources", None)
        if ext_par is not None:     prompts["ext_parallel_sources"] = int(ext_par)
        # ── Commercials (Stage 5b) — default ON when attribute missing ──
        enable_comm = getattr(self, "enable_commercials_extraction", None)
        if enable_comm is None:
            prompts["enable_commercials_extraction"] = True
        else:
            try:
                prompts["enable_commercials_extraction"] = bool(enable_comm)
            except Exception:
                prompts["enable_commercials_extraction"] = True
        # ── V2 benchmark knobs ───────────────────────────────────────────────
        for src_attr, key in (
            ("bench_sql_pool_size",           "bench_sql_pool_size"),
            ("bench_pinecone_top_k",          "bench_pinecone_top_k"),
            ("bench_llm_shortlist_size",      "bench_llm_shortlist_size"),
            ("bench_widen_l1l2_when_sparse",  "bench_widen_l1l2_when_sparse"),
            ("bench_critical_threshold_pct",  "bench_critical_threshold_pct"),
            ("bench_important_threshold_pct", "bench_important_threshold_pct"),
            ("bench_ratio_band_pct",          "bench_ratio_band_pct"),
            ("bench_rank_debug_dump_dir",     "bench_rank_debug_dump_dir"),
        ):
            v = getattr(self, src_attr, None)
            if v is not None and v != "":
                prompts[key] = v
        sys_override  = (getattr(self, "bench_rank_prompt_system_v2", "") or "").strip()
        user_override = (getattr(self, "bench_rank_prompt_user_v2",   "") or "").strip()
        if sys_override:
            prompts["bench_rank_prompt_system_v2"] = sys_override
        if user_override:
            prompts["bench_rank_prompt_user_v2"] = user_override
        return prompts

    def _run_stages_48(self, pr_no: str, tgt_cs: str, result: dict, prompts: dict) -> None:
        """Run stages 4-8 for a PR that has already completed stages 1-3.
        All Stage 4-8 functions are inlined in this file — no cross-file imports needed.
        Exceptions are caught here and reflected in result["status"]/result["error"]
        so stages 1-3 results are always preserved even if stages 4-8 fail.
        """
        blob_cfg      = self._blob_cfg()
        current_stage = _STAGE_CLASSIFICATION
        top_k         = int(getattr(self, "pinecone_top_k", 10))

        try:
            # Stage 4 — Classification
            self._safe_log(f"[{pr_no}] Stage 4 — classifying attachments…")
            _run_classification(self.llm, tgt_cs, blob_cfg, pr_no, prompts)
            self._advance_tracker(tgt_cs, pr_no, _STAGE_CLASSIFICATION)
            self._safe_log(f"[{pr_no}] Stage 4 — classification complete")

            # Stage 5 — Extraction
            # Raises an ExtractionAbortError subclass (NoQuotationFoundError,
            # NoLineItemsError, NoRASContextError, AllExtractionsFailedError)
            # whenever the PR cannot produce extracted items. Caught below as
            # an expected outcome — records ras_pipeline_exceptions, sets
            # tracker.current_stage_fk = 99, and skips stages 6/7/8.
            current_stage = _STAGE_EXTRACTION
            self._safe_log(f"[V2 {pr_no}] Stage 5 — extracting quotation items (V2)…")
            n_items = _run_extraction_v2(self.llm, tgt_cs, blob_cfg, pr_no, prompts)
            self._advance_tracker(tgt_cs, pr_no, _STAGE_EXTRACTION)
            self._safe_log(f"[V2 {pr_no}] Stage 5 — {n_items} item(s) extracted (V2)")

            # Stage 6 — Embeddings (V2 — uses embed_content)
            current_stage = _STAGE_EMBEDDINGS
            _run_embeddings_v2(tgt_cs, pr_no, self.embed_model,
                               self.pinecone_index, self.pinecone_namespace)
            self._advance_tracker(tgt_cs, pr_no, _STAGE_EMBEDDINGS)
            self._safe_log(f"[V2 {pr_no}] Stage 6 — embeddings done (V2)")

            # Stage 7 — Benchmark (V2 3-stage retrieval)
            current_stage = _STAGE_PRICE_BENCHMARK
            _run_benchmark_v2(self.llm, tgt_cs, pr_no, self.embed_model,
                              self.pinecone_index, self.pinecone_namespace, top_k,
                              prompts=prompts)
            self._advance_tracker(tgt_cs, pr_no, _STAGE_PRICE_BENCHMARK)
            self._safe_log(f"[V2 {pr_no}] Stage 7 — benchmark done (V2)")

            # Stage 8 — Complete
            self._advance_tracker(tgt_cs, pr_no, _STAGE_COMPLETE)
            # _set_last_processed_at is the module-level helper; wrap its
            # call in deadlock retry so it shares the same retry budget as
            # the rest of the pipeline's tracker writes.
            self._run_with_deadlock_retry(
                lambda: _set_last_processed_at(tgt_cs, pr_no),
                op_name="set_last_processed_at",
                pr_no=pr_no,
            )
            self._safe_log(f"[{pr_no}] Stage 8 — pipeline complete")
            result["status"] = "complete"

        except ExtractionAbortError as exc:
            # Expected outcome: Stage 5 cannot proceed for a known reason
            # (NoQuotationFoundError, NoLineItemsError, NoRASContextError,
            # AllExtractionsFailedError). Log as warning (no traceback),
            # record in ras_pipeline_exceptions with stage_id=5, set
            # tracker=99 (EXCEPTION), and stop — stages 6/7/8 do not run.
            reason = type(exc).__name__
            self._safe_log(f"[{pr_no}] Stage 5 (Extraction) — {reason}: {exc}")
            logger.warning("[{}] Stage 5 (Extraction) — {}: {}", pr_no, reason, exc)
            result["status"] = f"failed_{reason}"
            result["error"]  = f"Stage 5 (Extraction) — {reason}: {exc}"
            self._record_exception(tgt_cs, pr_no, _STAGE_EXTRACTION, f"{reason}: {exc}")

        except Exception as exc:
            stage_name = {
                _STAGE_CLASSIFICATION:  "Stage 4 (Classification)",
                _STAGE_EXTRACTION:      "Stage 5 (Extraction)",
                _STAGE_EMBEDDINGS:      "Stage 6 (Embeddings)",
                _STAGE_PRICE_BENCHMARK: "Stage 7 (Benchmark)",
                _STAGE_COMPLETE:        "Stage 8 (Complete)",
            }.get(current_stage, f"Stage {current_stage}")
            logger.opt(exception=True).error(
                "[{}] {} failed: {}", pr_no, stage_name, exc
            )
            result["status"] = "failed_stages_48"
            result["error"]  = f"{stage_name}: {exc}"
            self._record_exception(tgt_cs, pr_no, current_stage, f"{stage_name}: {exc}")

    # ── Process single PR ─────────────────────────────────────────────────

    def _process_pr(
        self,
        pr_no: str,
        src_cs: str,
        tgt_cs: str,
        *,
        skip_if_complete: bool = False,
    ) -> dict:
        import os, shutil, tempfile
        result        = {"pr_no": pr_no, "files": [], "status": "failed", "error": ""}
        current_stage = _STAGE_INGESTION
        work_dir      = tempfile.mkdtemp()
        self._safe_log(f"[{pr_no}] Worker started — work_dir={work_dir}")

        # Excel mode: each worker first checks the tracker and short-circuits
        # if the PR is already at stage 8 (COMPLETE) or stage 99 (EXCEPTION)
        # so we do not redo work. Stage-99 PRs need an explicit single-PR
        # force-reprocess to be retried — bulk Excel runs should not keep
        # burning compute on PRs that already failed.
        # This DB read happens on the worker's own connection, so N PRs run
        # the check in *parallel* — no serial pre-pass needed before the pool.
        if skip_if_complete:
            try:
                stage_now = _get_pr_current_stage(tgt_cs, pr_no)
            except Exception as exc:
                # Failure to check stage is non-fatal — we proceed with full
                # processing rather than skipping a PR we shouldn't have.
                self._safe_log(
                    f"[{pr_no}] Stage check failed (will process anyway): {exc}"
                )
                stage_now = -1
            if stage_now in (_STAGE_COMPLETE, _STAGE_EXCEPTION):
                stage_name = "COMPLETE" if stage_now == _STAGE_COMPLETE else "EXCEPTION"
                status_key = "already_complete" if stage_now == _STAGE_COMPLETE else "already_exception"
                self._safe_log(
                    f"[{pr_no}] Already at stage {stage_now} ({stage_name}) "
                    f"— skipping (no cleanup, no reprocess)"
                )
                result["status"] = status_key
                shutil.rmtree(work_dir, ignore_errors=True)
                return result

        try:
            # ── Cleanup at start — only for PRs that already have prior data ──
            # Checked before stage 1: new PR → skipped; existing PR → blobs +
            # Pinecone vectors + quotation_extracted_items + benchmark_result deleted.
            # Explicit handling so cleanup failures are recorded with a clear message
            # before any stage tracker row exists.
            #
            # Cleanup retries on SQL Server deadlock victims (error 1205,
            # SQLSTATE 40001). With multiple parallel workers each running
            # their own multi-DELETE cleanup transaction, lock contention
            # on shared join tables (ras_tracker / attachment_classification)
            # can cause one transaction to be picked as the deadlock victim;
            # SQL Server explicitly asks us to rerun, so we do — with
            # exponential backoff and jitter to spread retries.
            self._safe_log(f"[{pr_no}] Pre-run cleanup check…")
            try:
                self._run_with_deadlock_retry(
                    lambda: self._cleanup_for_pr(tgt_cs, pr_no),
                    op_name="Pre-run cleanup",
                    pr_no=pr_no,
                )
            except Exception as cleanup_exc:
                # Either a non-deadlock error (raised immediately by the
                # helper) or deadlock retries fully exhausted. Both end the
                # PR at stage 99 with a descriptive exception record.
                logger.opt(exception=True).error(
                    "[{}] Pre-run cleanup failed — aborting PR: {}", pr_no, cleanup_exc
                )
                self._record_exception(
                    tgt_cs, pr_no, _STAGE_INGESTION,
                    f"PRE_RUN_CLEANUP failed: {cleanup_exc}",
                )
                result["error"] = str(cleanup_exc)
                return result

            # Stage 1 — INGESTION
            current_stage = _STAGE_INGESTION
            self._safe_log(f"[{pr_no}] Stage 1 — INGESTION starting…")
            self._advance_tracker(tgt_cs, pr_no, _STAGE_INGESTION)

            ras_uuid    = self._get_tracker_uuid(tgt_cs, pr_no)
            self._safe_log(f"[{pr_no}] Stage 1 — tracker uuid={ras_uuid}, fetching on-prem attachments…")
            attachments = self._fetch_attachments(src_cs, tgt_cs, pr_no)
            total_bytes = sum(len(a.get("content", b"")) for a in attachments)
            self._safe_log(
                f"[{pr_no}] Stage 1 — fetched {len(attachments)} attachment(s) "
                f"({total_bytes / 1024:.1f} KB total)"
            )

            # BI dashboard sync: read on-prem view → refresh Azure row.
            # Non-fatal — log and continue if it fails.
            self._safe_log(f"[{pr_no}] Stage 1 — syncing BI dashboard row…")
            try:
                self._sync_bi_dashboard(src_cs, tgt_cs, pr_no)
            except Exception as _bi_exc:
                self._safe_log(f"[{pr_no}] Warning — BI dashboard sync failed (non-fatal): {_bi_exc}")

            if not attachments:
                error_msg = "No attachments found in on-prem RAS. Cannot proceed without quotation documents. Check if PR has associated files in the source system."
                self._safe_log(f"[{pr_no}] {error_msg}")
                logger.opt(exception=True).error(f"[{pr_no}] {error_msg}")

                # Record exception in DB so user can investigate
                try:
                    self._record_exception(tgt_cs, pr_no, 99, error_msg)
                except Exception as db_exc:
                    logger.warning(f"[{pr_no}] Could not record no-attachments exception in DB: {db_exc}")

                result["status"] = "exception"
                result["error"] = error_msg
                return result

            self._safe_log(f"[{pr_no}] Stage 1 — INGESTION complete")

            safe_pr  = pr_no.replace("/", "_")
            extractor = FileExtractor()
            all_files = []   # dicts: filename, content, blob_path, att_id, is_embedded

            self._safe_log(
                f"[{pr_no}] Stage 2 — EMBED_DOC_EXTRACTION starting "
                f"({len(attachments)} parent attachment(s) to scan for embeds)…"
            )

            failed_atts: list[str] = []
            for att_idx, att in enumerate(attachments, 1):
                att_id   = att["attachment_id"]
                att_dir  = os.path.join(work_dir, att_id)
                emb_dir  = os.path.join(att_dir, "extracted")
                try:
                    os.makedirs(att_dir, exist_ok=True)
                    os.makedirs(emb_dir,  exist_ok=True)

                    # Save parent file to temp disk so FileExtractor can open it
                    parent_path = os.path.join(att_dir, att["filename"])
                    with open(parent_path, "wb") as fh:
                        fh.write(att["content"])

                    parent_size_kb = len(att["content"]) / 1024
                    parent_ext     = os.path.splitext(att["filename"])[1].lower() or "<no-ext>"
                    self._safe_log(
                        f"[{pr_no}] [{att_idx}/{len(attachments)}] [{att_id}] "
                        f"{att['filename']} ({parent_size_kb:.1f} KB, {parent_ext}) "
                        f"— scanning for embedded docs"
                    )

                    # If the parent attachment is itself an archive
                    # (.zip / .rar / .7z / .tar / .tar.gz / .tgz / .tar.bz2),
                    # expand its contents directly into emb_dir so each inner
                    # file becomes a separately-classifiable embedded record.
                    # Without this, a zipped quotation is invisible to Stage 4
                    # (the .zip itself is not in _SUPPORTED_CLASSIFY_EXTS).
                    # Mirrors doc-intel pipeline.stages.embed_doc_extraction
                    # pass-1 archive expansion, adapted to the AgentCore flow:
                    # the original archive bytes still serve as the parent
                    # blob, but its contents flow through as embedded files.
                    extractor.parent_prefix = os.path.splitext(att["filename"])[0]
                    if att["filename"].lower().endswith(ARCHIVE_EXTENSIONS):
                        self._safe_log(
                            f"[{pr_no}] [{att_id}] expanding archive "
                            f"{att['filename']!r} into embedded set"
                        )
                        extractor.extract_archive(parent_path, emb_dir)
                        # extract_archive deletes the archive on disk; rewrite
                        # it so the parent blob upload still has the bytes.
                        with open(parent_path, "wb") as fh:
                            fh.write(att["content"])
                        # Then process every newly extracted file that itself
                        # is a supported parent (xlsx/docx/pdf/etc.) for any
                        # nested embeds — matches doc-intel pass-2.
                        for inner_name in sorted(os.listdir(emb_dir)):
                            inner_path = os.path.join(emb_dir, inner_name)
                            if not os.path.isfile(inner_path):
                                continue
                            if inner_name.lower().endswith(SUPPORTED_PARENTS):
                                extractor.parent_prefix = os.path.splitext(inner_name)[0]
                                extractor.process_file(inner_path, emb_dir)
                    else:
                        # Standard parent — extract embedded media/docs.
                        extractor.process_file(parent_path, emb_dir)

                    # Collect embedded file contents
                    embedded = []
                    if os.path.isdir(emb_dir):
                        for emb_name in sorted(os.listdir(emb_dir)):
                            emb_path = os.path.join(emb_dir, emb_name)
                            if os.path.isfile(emb_path):
                                with open(emb_path, "rb") as fh:
                                    emb_content = fh.read()
                                embedded.append({
                                    "filename":    emb_name,
                                    "content":     emb_content,
                                    "blob_path":   f"procurement/{safe_pr}/{att_id}/extracted/{emb_name}",
                                    "att_id":      att_id,
                                    "is_embedded": True,
                                })

                    # Record parent in attachment_classification
                    parent_blob_path = f"procurement/{safe_pr}/{att_id}/{att['filename']}"
                    parent_pk = self._upsert_parent_classification(
                        tgt_cs, pr_no, ras_uuid, att_id,
                        parent_blob_path,
                        embedded_flag  = len(embedded) > 0,
                        embedded_count = len(embedded),
                    )

                    # Record embedded files
                    for emb in embedded:
                        self._upsert_embedded_classification(tgt_cs, parent_pk, att_id, emb["blob_path"])

                    all_files.append({
                        "filename":    att["filename"],
                        "content":     att["content"],
                        "blob_path":   f"procurement/{safe_pr}/{att_id}/{att['filename']}",
                        "att_id":      att_id,
                        "is_embedded": False,
                    })
                    all_files.extend(embedded)
                    self._safe_log(f"[{pr_no}] [{att_id}] {att['filename']} — {len(embedded)} embedded file(s) extracted")

                except Exception as att_exc:
                    # Log and continue — one bad attachment must not block the whole PR.
                    # The attachment is excluded from further processing but the PR continues.
                    logger.opt(exception=True).error(
                        "[{}] Attachment {} ({}) failed during ingestion — skipping: {}",
                        pr_no, att_id, att.get("filename", "?"), att_exc,
                    )
                    failed_atts.append(f"{att_id}({att.get('filename','?')}): {att_exc}")

            if failed_atts:
                self._safe_log(f"[{pr_no}] {len(failed_atts)} attachment(s) skipped due to errors: {'; '.join(failed_atts)}")

            # Stage 2 — EMBED_DOC_EXTRACTION
            current_stage = _STAGE_EMBED_DOC
            parent_count   = sum(1 for f in all_files if not f["is_embedded"])
            embedded_count = sum(1 for f in all_files if f["is_embedded"])
            self._advance_tracker(tgt_cs, pr_no, _STAGE_EMBED_DOC)
            self._safe_log(
                f"[{pr_no}] Stage 2 — EMBED_DOC_EXTRACTION complete "
                f"({parent_count} parent + {embedded_count} embedded = {len(all_files)} total file(s))"
            )

            # Stage 3 — BLOB_UPLOAD
            current_stage = _STAGE_BLOB_UPLOAD
            self._safe_log(f"[{pr_no}] Stage 3 — BLOB_UPLOAD starting ({len(all_files)} file(s) → Azure Blob)…")

            # Upload all files to blob + build text batch
            file_data = []
            uploaded_bytes = 0
            for upload_idx, f_info in enumerate(all_files, 1):
                f_size = len(f_info["content"])
                try:
                    self._upload_blob(f_info["content"], f_info["blob_path"])
                    uploaded_bytes += f_size
                except Exception as upload_exc:
                    logger.opt(exception=True).error(
                        "[{}] Blob upload failed for {!r}: {}",
                        pr_no, f_info["blob_path"], upload_exc,
                    )
                    raise
                # Per-file log only every 10 files (or for big files) to avoid noise
                # on PRs with many small embeds. Always log first/last.
                if (upload_idx % 10 == 0
                        or upload_idx == 1
                        or upload_idx == len(all_files)
                        or f_size > 1_000_000):
                    self._safe_log(
                        f"[{pr_no}] Stage 3 — uploaded {upload_idx}/{len(all_files)}: "
                        f"{f_info['filename']} ({f_size / 1024:.1f} KB)"
                    )
                text      = self._extract_text(f_info["filename"], f_info["content"])
                file_type = _detect_file_type(f_info["filename"])
                extra     = _build_extra_metadata(f_info["filename"], f_info["content"], text)
                if self.max_content_chars and len(text) > int(self.max_content_chars):
                    text = text[: int(self.max_content_chars)]
                file_data.append({
                    "filename":    f_info["filename"],
                    "pr_no":       pr_no,
                    "att_id":      f_info["att_id"],
                    "is_embedded": f_info["is_embedded"],
                    "file_type":   file_type,
                    "extra":       extra,
                    "text":        text,
                    "blob":        f_info["blob_path"],
                })
                # Free raw bytes — uploaded to blob + saved to work_dir already.
                # Stages 4-8 re-download from blob when they need the bytes again.
                f_info["content"] = None

            self._advance_tracker(tgt_cs, pr_no, _STAGE_BLOB_UPLOAD)
            self._safe_log(
                f"[{pr_no}] Stage 3 — BLOB_UPLOAD complete "
                f"({len(file_data)} file(s), {uploaded_bytes / 1024:.1f} KB total)"
            )

            result["files"]  = file_data
            result["status"] = "success"

            # ── Free Stage 1-3 working set before Stage 4-8 ───────────────────
            # Stages 4-8 re-query DB and re-download blobs — they do not use
            # attachments / all_files / file_data. Dropping these locals here
            # reclaims hundreds of MB per PR before the heaviest stage runs.
            # result["files"] still references file_data; it is cleared in the
            # finally block below once the PR is fully done.
            attachments = None
            all_files   = None
            file_data   = None

            # ── Continue to stages 4-8 in this same worker ──────────────────
            # Each parallel worker processes one PR end-to-end (stages 1→8)
            # without waiting for other PRs to finish Stage 1-3 first.
            self._run_stages_48(pr_no, tgt_cs, result, self._build_prompts())

        except Exception as exc:
            stage_name = {
                _STAGE_INGESTION:   "Stage 1 (Ingestion)",
                _STAGE_EMBED_DOC:   "Stage 2 (Embed Doc Extraction)",
                _STAGE_BLOB_UPLOAD: "Stage 3 (Blob Upload)",
            }.get(current_stage, f"Stage {current_stage}")
            logger.opt(exception=True).error(
                "[{}] {} failed: {}", pr_no, stage_name, exc
            )
            result["error"] = f"{stage_name}: {exc}"
            self._record_exception(tgt_cs, pr_no, current_stage, f"{stage_name}: {exc}")
        finally:
            # Drop the per-file extracted text — the outer batch only reads
            # result["pr_no"] and result["status"]. Keeping ~1 MB of text per
            # PR across 600 PRs accumulates to >500 MB in the run-level
            # `results` list otherwise.
            try:
                result["files"] = []
            except Exception:
                pass
            shutil.rmtree(work_dir, ignore_errors=True)
            # Return freed allocator pages to the OS. Cheap (~100ms) compared
            # to a PR's LLM latency, and the cumulative effect across 600 PRs
            # is what keeps RSS stable on long Excel runs.
            try:
                import gc
                gc.collect()
            except Exception:
                pass
        return result

    # ── Entry point ───────────────────────────────────────────────────────

    def build_file_batch(self) -> Message:
        """Fire-and-forget entry point.

        Returns immediately with a "started" Message so the AgentCore SSE
        stream closes well before the AKS ingress 300s idle timeout. The
        actual pipeline (stages 1-8 per PR) runs in a daemon thread and logs
        progress via loguru / ras_tracker. Frontend can poll
        `ras_procurement.ras_tracker` to see per-PR stage advancement.
        """
        import threading, uuid

        if hasattr(self, "_cached_result"):
            return self._cached_result

        # Stamp the main thread so _safe_log can distinguish it from the bg
        # daemon thread (which has no AgentCore component context).
        self._main_thread_id = threading.get_ident()

        # Resize the connection pool from the canvas knob BEFORE any
        # _connect() call fires. Default 20; user can raise/lower per-Node.
        try:
            pool_max = int(getattr(self, "max_db_connections", None) or _DEFAULT_POOL_SIZE)
        except Exception:
            pool_max = _DEFAULT_POOL_SIZE
        if pool_max != _CONN_POOL.max_size:
            _CONN_POOL.resize(pool_max)
        logger.info(
            "V2 connection pool ready | max={} (idle={}, busy={})",
            _CONN_POOL.max_size, _CONN_POOL.stats()["idle"], _CONN_POOL.stats()["busy"],
        )

        # Resolve connection strings up-front in the main thread so any
        # config error fails fast with a readable message rather than dying
        # silently in the background.
        try:
            src_cs = self._conn_str(self.source_connection)
            tgt_cs = self._conn_str(self.target_connection)
        except Exception as exc:
            msg = Message(text=f"Connection config error: {exc}")
            self._cached_result    = msg
            self._cached_pr_numbers = []
            return msg

        # Smoke-test both DB connections IN PARALLEL — source is on-prem
        # (high-latency over VPN) and target is Azure SQL (low-latency).
        # Running them sequentially via self._connect (which retries on
        # transient errors up to 4× with 30s login timeout each) made startup
        # wait the source DB out before even reaching the target. Probing
        # them concurrently with a single short-timeout login means total
        # smoke-test time is max(source, target) ≈ 5-15s instead of summing.
        try:
            import concurrent.futures, pyodbc
        except Exception as exc:
            # An ImportError here is structural (pyodbc missing) — surface it
            # clearly instead of letting the daemon thread fail silently.
            logger.opt(exception=True).error(f"V2 smoke-test bootstrap failed: {exc}")
            msg = Message(text=f"V2 smoke-test bootstrap failed: {exc}")
            self._cached_result    = msg
            self._cached_pr_numbers = []
            return msg

        def _probe_host(cs: str) -> str:
            """Pull SERVER=… out of the connection string for log lines —
            never logs UID/PWD."""
            m = re.search(r"SERVER=([^;]+)", cs or "", re.IGNORECASE)
            return m.group(1) if m else "?"

        def _probe(cs: str) -> float:
            t0 = time.time()
            c = pyodbc.connect(cs, timeout=15)  # one-shot, no retry loop
            c.close()
            return time.time() - t0

        smoke_t0 = time.time()
        self._safe_log("V2 — smoke-testing source + target DB connections in parallel…")
        logger.info(
            "V2 smoke-test starting | source={} target={} (timeout=15s, parallel=2)",
            _probe_host(src_cs), _probe_host(tgt_cs),
        )
        results: dict[str, tuple[Optional[float], Optional[Exception]]] = {
            "source": (None, None), "target": (None, None),
        }
        try:
            with concurrent.futures.ThreadPoolExecutor(
                max_workers=2, thread_name_prefix="db-probe",
            ) as pool:
                futures = {
                    pool.submit(_probe, src_cs): "source",
                    pool.submit(_probe, tgt_cs): "target",
                }
                for f in concurrent.futures.as_completed(futures):
                    label = futures[f]
                    try:
                        dt = f.result()
                        results[label] = (dt, None)
                        logger.info(f"V2 smoke-test {label} OK in {dt:.2f}s")
                    except Exception as exc:
                        results[label] = (None, exc)
                        logger.warning(
                            f"V2 smoke-test {label} FAILED after "
                            f"{time.time() - smoke_t0:.2f}s: {type(exc).__name__}: {exc}"
                        )
        except Exception as exc:
            # Pool-level / scheduling failure (e.g. OS thread limit). Treat as
            # a hard failure so the chat shows it and the daemon never starts.
            logger.opt(exception=True).error(
                f"V2 smoke-test pool failed before probes ran: {exc}"
            )
            msg = Message(text=f"V2 smoke-test pool failed: {exc}")
            self._cached_result    = msg
            self._cached_pr_numbers = []
            return msg

        errors = [(lbl, exc) for lbl, (_, exc) in results.items() if exc is not None]
        smoke_elapsed = time.time() - smoke_t0
        if errors:
            err_text = "\n".join(
                f"Cannot connect to {label} DB: {type(exc).__name__}: {exc}"
                for label, exc in errors
            )
            err_text += f"\n(smoke-test took {smoke_elapsed:.2f}s)"
            logger.error(
                f"V2 smoke-test aborting startup — {len(errors)} probe(s) failed: "
                f"{[lbl for lbl, _ in errors]}"
            )
            msg = Message(text=err_text)
            self._cached_result    = msg
            self._cached_pr_numbers = []
            return msg
        src_dt = results["source"][0] or 0.0
        tgt_dt = results["target"][0] or 0.0
        self._safe_log(
            f"V2 smoke-test passed — source {src_dt:.2f}s, target {tgt_dt:.2f}s "
            f"(wall {smoke_elapsed:.2f}s)"
        )

        job_id = str(uuid.uuid4())[:8]

        def _bg_run() -> None:
            try:
                self._run_full_pipeline(src_cs, tgt_cs, job_id)
            except Exception as exc:
                logger.opt(exception=True).error(
                    f"[pipeline job={job_id}] failed: {exc}"
                )

        t = threading.Thread(
            target=_bg_run, daemon=True, name=f"pipeline-{job_id}"
        )
        t.start()

        self._safe_log(
            f"Pipeline started in background | job={job_id} | "
            f"workers={self.parallel_workers}"
        )

        msg = Message(text=(
            f"✅ Pipeline started in background.\n"
            f"Job ID: {job_id}\n"
            f"Workers: {self.parallel_workers}\n\n"
            f"Stages 1-8 are running for each pending PR. Progress can be "
            f"monitored in `ras_procurement.ras_tracker` (current_stage_fk "
            f"advances 1 → 8; failures land at 99). Detailed per-PR logs "
            f"are visible in the AgentCore server logs (filter on "
            f"`[pipeline job={job_id}]`)."
        ))
        self._cached_result    = msg
        self._cached_pr_numbers = []
        return msg

    def _run_full_pipeline(self, src_cs: str, tgt_cs: str, job_id: str) -> None:
        """Heavy lifting — runs in the daemon thread spawned by
        build_file_batch. All logs in this method (and everything it calls)
        route through loguru since `_main_thread_id` no longer matches.

        Logging convention (so a single run can be traced with `grep`):
            [pipeline job=XXXXXXXX] <message>
        Filter the AgentCore server log with the literal job id printed in
        the "started" chat message and you will see the full timeline:
        start → mode + PR count → per-PR completions → final summary.
        """
        import time
        t_start = time.time()
        tag = f"[pipeline job={job_id}]"
        logger.info(f"{tag} starting")
        pr_filter = (self.pr_no_filter or "").strip()

        # Source resolution order:
        #   1. pr_no_filter           → single PR, force reprocess
        #   2. pr_list_data /         → Excel-driven (wired Data input or
        #      pr_excel_blob_path        blob path), each PR force-reprocessed
        #   3. _fetch_pending_prs     → standard batch from ras_tracker
        excel_pr_list: list[str] | None = None
        if not pr_filter:
            try:
                excel_pr_list = self._resolve_excel_pr_list()
            except Exception as exc:
                logger.error(
                    f"[pipeline job={job_id}] Excel PR list could not be read: {exc}"
                )
                return

        # Detect completed PRs whose on-prem source data changed (U_DATETIME >
        # last_processed_at) and reset them to stage 1 so they are re-processed
        # in this run. Non-fatal — skipped silently on failure.
        # Skipped when an explicit list (single-PR or Excel) is being run —
        # that flow already specifies exactly which PRs to process.
        if not pr_filter and excel_pr_list is None:
            self._detect_and_requeue_changed_prs(tgt_cs)

        if excel_pr_list is not None:
            pr_list = excel_pr_list
            run_mode = "excel"
        else:
            pr_list = self._fetch_pending_prs(tgt_cs)
            run_mode = "single" if pr_filter else "pending"

        if not pr_list:
            reason = (
                "Excel PR list was empty"
                if run_mode == "excel"
                else "no pending PRs to process"
            )
            elapsed = time.time() - t_start
            logger.info(f"{tag} done — {reason} (elapsed {elapsed:.1f}s)")
            return

        logger.info(
            f"{tag} resolved {len(pr_list)} PR(s) | mode={run_mode} | "
            f"first={pr_list[0]!r} last={pr_list[-1]!r}"
        )

        if pr_filter:
            # Single-PR mode = explicit force reprocess regardless of state.
            # Reset tracker from exception/any stage back to NULL so _process_pr
            # picks it up cleanly. Actual data cleanup (blobs, Pinecone, DB tables)
            # happens inside _process_pr at the start via _cleanup_for_pr.
            self._reset_for_reprocess(tgt_cs, pr_filter)
            self._safe_log(f"Single-PR reprocess: {pr_filter!r} — tracker reset, cleanup will run at processing start")
        elif run_mode == "excel":
            # Excel mode = workers self-filter inside _process_pr via
            # skip_if_complete=True. No serial pre-pass: stage-8 check
            # runs on each worker's own DB connection in parallel.
            #
            # Submit oldest-first to mitigate the Stage 7 benchmark race —
            # if PR_NEW is processed before PR_OLD reaches Stage 6,
            # PR_NEW's benchmark sees thin Pinecone history. Sorting by
            # purchase_req_mst.C_DATETIME ASC means older PRs start first
            # and their embeddings land in Pinecone before newer PRs reach
            # Stage 7. Single batched query (chunked at 500 PRs).
            try:
                sorted_excel = self._sort_excel_pr_list_by_date(tgt_cs, pr_list)
                if sorted_excel != pr_list:
                    self._safe_log(
                        f"Excel — reordered {len(sorted_excel)} PR(s) by "
                        f"C_DATETIME ASC (oldest first) so historical "
                        f"context is in Pinecone before newer PRs hit Stage 7"
                    )
                    pr_list = sorted_excel
                else:
                    self._safe_log("Excel — PR list already in C_DATETIME ASC order")
            except Exception as exc:
                # Non-fatal: if sort fails (e.g. DB hiccup), fall back to
                # original Excel order rather than aborting the run.
                self._safe_log(f"Excel — date sort skipped (non-fatal): {exc}")

            self._safe_log(
                f"Excel-driven run: {len(pr_list)} PR(s) queued — "
                f"each worker will skip if already at stage {_STAGE_COMPLETE} "
                f"(COMPLETE) or {_STAGE_EXCEPTION} (EXCEPTION)"
            )

        workers = max(1, int(self.parallel_workers))
        self._safe_log(
            f"Processing {len(pr_list)} PR(s) [mode={run_mode}] "
            f"with {workers} parallel worker(s)…"
        )
        logger.info(
            f"{tag} fan-out begins — {len(pr_list)} PR(s), "
            f"{workers} parallel worker(s)"
        )

        # Pre-warm shared caches before threads start so all workers reuse the
        # same already-initialised ContainerClient / blob config rather than
        # each racing to create their own copy.
        self._blob_cfg()
        self._container_client()

        # Worker-side skip is only enabled for Excel mode. Single-PR mode is
        # an explicit force reprocess (must run regardless of stage), and the
        # pending-batch flow already filters stage 8 out at fetch time.
        skip_if_complete = (run_mode == "excel")

        import concurrent.futures
        results: list[dict] = []
        total = len(pr_list)
        with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {
                pool.submit(
                    self._process_pr, pr, src_cs, tgt_cs,
                    skip_if_complete=skip_if_complete,
                ): pr
                for pr in pr_list
            }
            completed = 0
            for f in concurrent.futures.as_completed(futures):
                pr_no = futures[f]
                try:
                    res = f.result()
                except Exception as exc:
                    res = {"pr_no": pr_no, "files": [], "status": "failed",
                           "error": f"worker raised: {exc}"}
                    logger.opt(exception=True).error(
                        f"{tag} [{pr_no}] worker raised unhandled exception"
                    )
                results.append(res)
                completed += 1
                status = res.get("status", "failed")
                err = res.get("error") or ""
                err_suffix = f" — {err}" if err else ""
                logger.info(
                    f"{tag} progress {completed}/{total} | "
                    f"[{pr_no}] {status}{err_suffix}"
                )

        # Summary log — one line per status bucket. The "already_complete"
        # bucket is populated by workers that found their PR at stage 8 and
        # short-circuited (Excel mode). All buckets come from results since
        # there is no longer a serial pre-pass.
        status_counts: dict[str, int] = {}
        for r in results:
            s = r.get("status", "failed")
            status_counts[s] = status_counts.get(s, 0) + 1
        summary_parts = [f"{v} {k}" for k, v in sorted(status_counts.items())]
        elapsed = time.time() - t_start
        rate = (len(results) / elapsed) if elapsed > 0 else 0.0
        logger.info(
            f"{tag} RUN COMPLETE — {len(results)} PR(s) in {elapsed:.1f}s "
            f"({rate:.2f} PR/s) | {', '.join(summary_parts)}"
        )
        for r in results:
            if r.get("error"):
                logger.warning(
                    f"{tag} [{r['pr_no']}] {r['status']}: {r['error']}"
                )

        # Late-update _cached_pr_numbers so any poll of get_processed_prs()
        # after the run finishes returns the actual processed list. The
        # frontend-facing _cached_result Message was already set in
        # build_file_batch when the daemon was spawned and is not overwritten.
        self._cached_pr_numbers = [
            r["pr_no"] for r in results
            if r.get("status") in ("complete", "success", "skipped")
        ]
        return

    def get_processed_prs(self) -> Data:
        """Returns the PR numbers that reached stage 3 in this run.
        Wire this output to Stage 4-8's 'Stage 1-3 Result' input so Stage 4-8
        processes exactly those PRs synchronously after Stage 1-3 finishes."""
        if not hasattr(self, "_cached_pr_numbers"):
            self.build_file_batch()
        return Data(data={"pr_numbers": self._cached_pr_numbers})


# ── Text extraction helpers ───────────────────────────────────────────────

def _extract_pdf(raw: bytes) -> str:
    import io, fitz
    parts: list[str] = []
    with fitz.open(stream=raw, filetype="pdf") as doc:
        for page in doc:
            parts.append(page.get_text())
    return "\n".join(parts)


def _extract_excel(raw: bytes) -> str:
    import io, openpyxl
    wb    = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append("\t".join(cells))
        return "\n".join(parts)
    finally:
        try: wb.close()
        except Exception: pass


def _extract_word(raw: bytes) -> str:
    import io, docx
    doc = docx.Document(io.BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(raw: bytes) -> str:
    import io
    from pptx import Presentation
    prs   = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_msg_text(raw: bytes) -> str:
    import tempfile, os
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject:   parts.append(f"Subject: {msg.subject}")
            if msg.sender:    parts.append(f"From: {msg.sender}")
            if msg.body:      parts.append(msg.body)
            msg.close()
            return "\n".join(parts)
        finally:
            os.unlink(tmp_path)
    except Exception as exc:
        return f"[MSG extraction error: {exc}]"


def _detect_file_type(filename: str) -> str:
    import os
    ext = os.path.splitext(filename.lower())[1]
    return {
        ".pdf": "PDF",
        ".xlsx": "Excel", ".xls": "Excel", ".xlsm": "Excel", ".xlsb": "Excel",
        ".docx": "Word",  ".doc": "Word",
        ".pptx": "PowerPoint", ".ppt": "PowerPoint",
        ".txt": "Text", ".csv": "CSV", ".xml": "XML", ".json": "JSON",
        ".html": "HTML", ".htm": "HTML", ".msg": "Email", ".eml": "Email",
        ".png": "Image", ".jpg": "Image", ".jpeg": "Image",
    }.get(ext, f"Unknown ({ext})")


def _build_extra_metadata(filename: str, raw: bytes, text: str) -> str:
    import os, io
    ext   = os.path.splitext(filename.lower())[1]
    lines: list[str] = []
    if ext == ".pdf":
        try:
            import fitz
            with fitz.open(stream=raw, filetype="pdf") as doc:
                lines.append(f"- page_count: {doc.page_count}")
        except Exception:
            pass
    if ext in {".xlsx", ".xls", ".xlsm", ".xlsb"}:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            lines.append(f"- sheet_count: {len(wb.sheetnames)}")
            lines.append(f"- sheet_names: {', '.join(wb.sheetnames[:10])}")
        except Exception:
            pass
    if ext in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            lines.append(f"- slide_count: {len(Presentation(io.BytesIO(raw)).slides)}")
        except Exception:
            pass
    lines.append(f"- char_count: {len(text)}")
    lines.append(f"- size_bytes: {len(raw)}")
    return ("\n".join(lines) + "\n") if lines else ""

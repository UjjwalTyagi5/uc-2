"""Full Pipeline (Stages 1-8) — AgentCore custom-code component.

Logging
-------
This file does NOT call logger.add(...) or open a log file. loguru uses its
default stderr sink, and Node-level events are also surfaced via self.log()
so they appear in the AgentCore canvas log panel in real time. No log
directory is created on disk by this component.

Verbosity:
  - logger.info  : per-stage progress, page/image counts, LLM call summary
  - logger.warning : recoverable issues (page truncation, transient retries,
                     extraction failures for one file, no quotation found)
  - logger.error : unrecoverable per-PR failures (DB write, render error)
                   captured with stack trace via logger.opt(exception=True)
  - logger.debug : low-level diagnostics (suppressed at default level)
"""
from __future__ import annotations

import json
import random
import re
import time
from dataclasses import dataclass, field
from decimal import Decimal
from difflib import SequenceMatcher
from typing import Optional

from loguru import logger

from agentcore.custom import Node
from agentcore.io import HandleInput, IntInput, MessageTextInput, Output
from agentcore.schema.data import Data
from agentcore.schema.message import Message

_MAX_RETRIES   = 3
_BASE_DELAY    = 2.0
_STAGE_INGESTION  = 1
_STAGE_EMBED_DOC  = 2
_STAGE_BLOB_UPLOAD = 3
_STAGE_EXCEPTION  = 99


class ExtractionAbortError(RuntimeError):
    """Raised when Stage 5 (Extraction) cannot produce any items.

    The Stage 4-8 catch block treats every subclass as an expected pipeline
    outcome (not a crash): logs it as a warning without traceback, records a
    row in ras_pipeline_exceptions (stage_id=5), sets
    ras_tracker.current_stage_fk = 99 (EXCEPTION), and skips stages 6/7/8.

    Subclasses identify the specific reason so the exception_message in DB
    is self-explanatory.
    """
    _suppress_traceback = True


class NoQuotationFoundError(ExtractionAbortError):
    """No attachments classified as 'Quotation'."""


class NoLineItemsError(ExtractionAbortError):
    """RAS context has zero rows in purchase_req_detail."""


class NoRASContextError(ExtractionAbortError):
    """No purchase_req_mst row found for this PR."""


class AllExtractionsFailedError(ExtractionAbortError):
    """All quotation files failed to extract any items."""

_TEXT_EXT  = {".txt", ".csv", ".xml", ".json", ".html", ".htm", ".eml"}
_PDF_EXT   = {".pdf"}
_EXCEL_EXT = {".xlsx", ".xls", ".xlsm", ".xlsb"}
_WORD_EXT  = {".doc", ".docx"}
_PPTX_EXT  = {".ppt", ".pptx"}
_MSG_EXT   = {".msg"}

SUPPORTED_PARENTS   = (".xlsx", ".xls", ".docx", ".doc", ".pptx", ".ppt", ".pdf", ".msg")
ARCHIVE_EXTENSIONS  = (".zip", ".7z", ".rar", ".tar", ".tar.gz", ".tgz", ".tar.bz2")
EXTRACT_EXTENSIONS  = (
    ".pdf", ".xls", ".xlsx", ".doc", ".docx", ".ppt", ".pptx",
    ".jpg", ".jpeg", ".png", ".tif", ".tiff",
    ".zip", ".7z", ".rar", ".tar", ".tar.gz", ".tgz", ".tar.bz2",
    ".txt", ".msg",
)
SKIP_EXTENSIONS = (".emf", ".bin")
FILE_SIGNATURES = [
    (b"%PDF",               ".pdf"),
    (b"PK\x03\x04",         ".zip"),
    (b"7z\xBC\xAF\x27\x1C", ".7z"),
    (b"Rar!\x1A\x07",       ".rar"),
    (b"\xFF\xD8\xFF",       ".jpg"),
    (b"\x89PNG\r\n\x1A\n",  ".png"),
    (b"II*\x00",            ".tif"),  # TIFF little-endian
    (b"MM\x00*",            ".tif"),  # TIFF big-endian
]


# ── Embedded-doc extractor (inlined from embed_doc_extraction.extractor) ──

class FileExtractor:
    def __init__(self) -> None:
        self.extracted_count: int = 0
        self.parent_prefix:   str = ""

    def should_extract(self, filename: str) -> bool:
        fl = filename.lower()
        if any(fl.endswith(s) for s in SKIP_EXTENSIONS):
            return False
        return any(fl.endswith(e) for e in EXTRACT_EXTENSIONS)

    def sanitize(self, name: str) -> str:
        import os
        name = str(name).replace("\x00", "").replace("\x01", "").strip()
        name = os.path.basename(name)
        for ch in '<>:"/\\|?*':
            name = name.replace(ch, "_")
        name = "".join(ch for ch in name if ch.isprintable() and ord(ch) > 31)
        return name.strip(" .") or "file"

    def with_prefix(self, filename: str) -> str:
        filename = self.sanitize(filename)
        prefix   = self.sanitize(self.parent_prefix).strip()
        return f"{prefix}__{filename}" if prefix else filename

    def unique_path(self, directory: str, filename: str) -> str:
        import os
        fp = os.path.join(directory, filename)
        if not os.path.exists(fp):
            return fp
        base, ext = os.path.splitext(filename)
        c = 1
        while os.path.exists(os.path.join(directory, f"{base}_{c}{ext}")):
            c += 1
        return os.path.join(directory, f"{base}_{c}{ext}")

    def _find_payload(self, data: bytes):
        best = None
        for sig, ext in FILE_SIGNATURES:
            pos = data.find(sig)
            if pos != -1 and (best is None or pos < best[0]):
                best = (pos, ext)
        return best if best else (None, None)

    def _classify_zip(self, path: str) -> str:
        """If a saved .zip is actually an Office Open XML file, return the
        correct extension (.xlsx / .docx / .pptx). Mirrors doc-intel's
        embed_doc_extraction.extractor._classify_zip — without it, an XLSX
        payload extracted from an OLE container ends up on disk as `.zip`,
        which the classifier (limited to _SUPPORTED_CLASSIFY_EXTS) treats
        as 'Others' and any embedded Quotation is lost.
        """
        import zipfile
        try:
            with zipfile.ZipFile(path, "r") as z:
                names = set(z.namelist())
                if "[Content_Types].xml" in names:
                    if any(n.startswith("xl/")   for n in names): return ".xlsx"
                    if any(n.startswith("word/") for n in names): return ".docx"
                    if any(n.startswith("ppt/")  for n in names): return ".pptx"
        except Exception:
            pass
        return ".zip"

    def _maybe_fix_zip_extension(self, target: str) -> str:
        """If `target` ends in .zip but is actually OOXML, rename and return new path."""
        import os
        if not target.lower().endswith(".zip"):
            return target
        correct = self._classify_zip(target)
        if correct == ".zip":
            return target
        new_path = self.unique_path(
            os.path.dirname(target),
            os.path.splitext(os.path.basename(target))[0] + correct,
        )
        try:
            os.replace(target, new_path)
            logger.debug("Reclassified .zip → {}: {}", correct, os.path.basename(new_path))
            return new_path
        except Exception as exc:
            logger.warning("Could not rename {} → {}: {}", target, new_path, exc)
            return target

    def _save_file(self, data: bytes, filename: str, out: str) -> str:
        import os, shutil
        target = self.unique_path(out, self.with_prefix(filename))
        with open(target, "wb") as f:
            f.write(data)
        self.extracted_count += 1
        # Fix OOXML-misclassified-as-.zip BEFORE archive expansion so the
        # file isn't extracted as a plain ZIP archive when it's actually an
        # XLSX/DOCX/PPTX worth classifying directly.
        target = self._maybe_fix_zip_extension(target)
        if target.lower().endswith(ARCHIVE_EXTENSIONS):
            self.extract_archive(target, out)
        return target

    def extract_archive(self, path: str, out: str) -> None:
        import os, shutil, tarfile, tempfile, zipfile
        if not os.path.exists(path):
            return
        fl = path.lower()
        try:
            if fl.endswith(".zip"):
                with zipfile.ZipFile(path, "r") as z:
                    for m in z.infolist():
                        if m.is_dir():
                            continue
                        fn = os.path.basename(m.filename.replace("\\", "/").rstrip("/"))
                        if fn:
                            t = self.unique_path(out, self.with_prefix(fn))
                            with z.open(m) as s, open(t, "wb") as d:
                                shutil.copyfileobj(s, d)
                            self.extracted_count += 1
            elif fl.endswith(".7z"):
                import py7zr
                with tempfile.TemporaryDirectory() as tmp:
                    with py7zr.SevenZipFile(path, "r") as ar:
                        ar.extractall(path=tmp)
                    for root, _, files in os.walk(tmp):
                        for f in files:
                            t = self.unique_path(out, self.with_prefix(f))
                            shutil.copy2(os.path.join(root, f), t)
                            self.extracted_count += 1
            elif fl.endswith(".rar"):
                import rarfile
                with tempfile.TemporaryDirectory() as tmp:
                    with rarfile.RarFile(path, "r") as rar:
                        rar.extractall(tmp)
                    for root, _, files in os.walk(tmp):
                        for f in files:
                            t = self.unique_path(out, self.with_prefix(f))
                            shutil.copy2(os.path.join(root, f), t)
                            self.extracted_count += 1
            elif fl.endswith((".tar", ".tar.gz", ".tgz", ".tar.bz2")):
                with tarfile.open(path, "r:*") as tar:
                    for m in tar.getmembers():
                        if not m.isfile():
                            continue
                        fn = os.path.basename(m.name)
                        if fn:
                            t  = self.unique_path(out, self.with_prefix(fn))
                            src = tar.extractfile(m)
                            if src:
                                with src as s, open(t, "wb") as d:
                                    shutil.copyfileobj(s, d)
                                self.extracted_count += 1
        except Exception as exc:
            logger.warning(f"Archive extraction error ({os.path.basename(path)}): {exc}")
        try:
            os.remove(path)
        except Exception:
            pass

    def extract_from_ooxml(self, file_path: str, out: str) -> bool:
        import os, shutil, tempfile, zipfile
        try:
            with tempfile.TemporaryDirectory() as tmp:
                with zipfile.ZipFile(file_path, "r") as zr:
                    zr.extractall(tmp)
                for prefix in ("xl", "word", "ppt"):
                    for sub in ("media", "embeddings"):
                        folder = os.path.join(tmp, prefix, sub)
                        if not os.path.exists(folder):
                            continue
                        for fn in os.listdir(folder):
                            src = os.path.join(folder, fn)
                            if not os.path.isfile(src):
                                continue
                            if fn.lower().endswith(".bin"):
                                self._extract_from_bin(src, out)
                            elif self.should_extract(fn):
                                dst = self.unique_path(out, self.with_prefix(fn))
                                shutil.copy2(src, dst)
                                self.extracted_count += 1
                                # OOXML embeddings often come as .zip-stamped
                                # bytes that are really nested xlsx/docx/pptx.
                                dst = self._maybe_fix_zip_extension(dst)
                                if dst.lower().endswith(ARCHIVE_EXTENSIONS):
                                    self.extract_archive(dst, out)
            return True
        except Exception as exc:
            logger.error(f"OOXML error ({os.path.basename(file_path)}): {exc}")
            return False

    def _extract_from_bin(self, bin_path: str, out: str) -> None:
        try:
            import olefile
            if not olefile.isOleFile(bin_path):
                return
            ole = olefile.OleFileIO(bin_path)
            for entry in ole.listdir():
                try:
                    sname = entry[-1] if entry else ""
                    data  = ole.openstream(entry).read()
                    if not data or len(data) < 16:
                        continue
                    offset, ext = self._find_payload(data)
                    if offset is None:
                        continue
                    payload   = data[offset:]
                    suggested = None
                    if sname.lower() in ("\x01ole10native", "ole10native"):
                        try:
                            end = data.find(b"\x00", 4)
                            if end != -1 and end - 4 < 260:
                                suggested = data[4:end].decode(errors="ignore")
                        except Exception:
                            pass
                    suggested = suggested or f"embedded_file{ext}"
                    target = self.unique_path(out, self.with_prefix(self.sanitize(suggested)))
                    with open(target, "wb") as f:
                        f.write(payload)
                    self.extracted_count += 1
                    # Reclassify .zip → .xlsx/.docx/.pptx if the payload was
                    # actually an OOXML doc; then expand if it remains a zip.
                    target = self._maybe_fix_zip_extension(target)
                    if target.lower().endswith(ARCHIVE_EXTENSIONS):
                        self.extract_archive(target, out)
                except Exception:
                    pass
            ole.close()
        except Exception as exc:
            logger.warning(f"OLE bin error: {exc}")

    def extract_from_ole(self, file_path: str, out: str) -> bool:
        try:
            import olefile
            if not olefile.isOleFile(file_path):
                return False
            ole = olefile.OleFileIO(file_path)
            for entry in ole.listdir():
                try:
                    sname = entry[-1] if entry else ""
                    data  = ole.openstream(entry).read()
                    if not data or len(data) < 16:
                        continue
                    offset, ext = self._find_payload(data)
                    if offset is None:
                        continue
                    payload   = data[offset:]
                    suggested = None
                    if sname.lower() in ("\x01ole10native", "ole10native"):
                        try:
                            end = data.find(b"\x00", 4)
                            if end != -1 and end - 4 < 260:
                                suggested = data[4:end].decode(errors="ignore")
                        except Exception:
                            pass
                    suggested = suggested or f"embedded_file{ext}"
                    target = self.unique_path(out, self.with_prefix(self.sanitize(suggested)))
                    with open(target, "wb") as f:
                        f.write(payload)
                    self.extracted_count += 1
                    # Same OOXML-misclassified-as-.zip fix as in _extract_from_bin.
                    target = self._maybe_fix_zip_extension(target)
                    if target.lower().endswith(ARCHIVE_EXTENSIONS):
                        self.extract_archive(target, out)
                except Exception:
                    pass
            ole.close()
            return True
        except Exception as exc:
            logger.error(f"OLE error ({os.path.basename(file_path)}): {exc}")
            return False

    def extract_from_pdf(self, file_path: str, out: str) -> bool:
        try:
            import fitz
            doc = fitz.open(file_path)
            if doc.embfile_count() > 0:
                for i in range(doc.embfile_count()):
                    info = doc.embfile_info(i)
                    name = self.sanitize(info.get("filename", f"pdf_attachment_{i}"))
                    data = doc.embfile_get(i)
                    if data:
                        self._save_file(data, name, out)
            for pn, page in enumerate(doc, 1):
                annots = page.annots()
                if not annots:
                    continue
                for annot in annots:
                    if annot.type[0] == 17:
                        try:
                            fd      = annot.get_file()
                            name    = self.sanitize(fd.get("filename", f"page{pn}_attachment"))
                            content = fd.get("content", b"")
                            if content:
                                self._save_file(content, name, out)
                        except Exception:
                            pass
            doc.close()
            return True
        except Exception as exc:
            logger.error(f"PDF error ({os.path.basename(file_path)}): {exc}")
            return False

    def extract_from_msg(self, file_path: str, out: str) -> bool:
        try:
            import tempfile, os
            import extract_msg
            msg = extract_msg.Message(file_path)
            if not msg.attachments:
                msg.close()
                return True
            for att in msg.attachments:
                name = "attachment"
                try:
                    name = (
                        getattr(att, "longFilename", None)
                        or getattr(att, "shortFilename", None)
                        or getattr(att, "filename", None)
                        or "attachment"
                    )
                    name = self.sanitize(name)
                    if getattr(att, "type", None) == "msg" and not name.lower().endswith(".msg"):
                        name += ".msg"
                    data = getattr(att, "data", None) or b""
                    if not data:
                        with tempfile.TemporaryDirectory() as tmp:
                            att.save(customPath=tmp)
                            for fn in os.listdir(tmp):
                                src = os.path.join(tmp, fn)
                                if os.path.isfile(src):
                                    with open(src, "rb") as f:
                                        data = f.read()
                                    if name == "attachment":
                                        name = fn
                                    break
                    if data:
                        self._save_file(data, name, out)
                except Exception as exc:
                    logger.warning(f"MSG attachment '{name}' failed: {exc}")
            msg.close()
            return True
        except Exception as exc:
            logger.error(f"MSG error ({os.path.basename(file_path)}): {exc}")
            return False

    def detect_type(self, path: str):
        lower = path.lower()
        if lower.endswith((".xlsx", ".docx", ".pptx")): return "ooxml"
        if lower.endswith((".xls",  ".doc",  ".ppt")):  return "ole"
        if lower.endswith(".pdf"):                       return "pdf"
        if lower.endswith(".msg"):                       return "msg"
        try:
            with open(path, "rb") as f:
                h = f.read(8)
            if h.startswith(b"PK\x03\x04"):                        return "ooxml"
            if h.startswith(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"): return "ole"
            if h.startswith(b"%PDF"):                               return "pdf"
        except Exception:
            pass
        return None

    def process_file(self, file_path: str, out: str) -> bool:
        import os
        ft = self.detect_type(file_path)
        if not ft:
            return False
        os.makedirs(out, exist_ok=True)
        if ft == "ooxml": return self.extract_from_ooxml(file_path, out)
        if ft == "ole":   return self.extract_from_ole(file_path, out)
        if ft == "pdf":   return self.extract_from_pdf(file_path, out)
        if ft == "msg":   return self.extract_from_msg(file_path, out)
        return False


# ── Blob connector lookup ─────────────────────────────────────────────────

def _get_blob_config_by_name(connector_name: str) -> dict:
    name = (connector_name or "").strip()
    if not name:
        raise ValueError("blob_connector_name is empty. Enter the connector name from Settings → Connectors.")
    try:
        import asyncio
        import concurrent.futures as _cf
        from sqlalchemy import select

        async def _fetch():
            from agentcore.services.deps import get_db_service
            from agentcore.services.database.models.connector_catalogue.model import ConnectorCatalogue
            db_service = get_db_service()
            async with db_service.with_session() as session:
                stmt = (
                    select(ConnectorCatalogue)
                    .where(ConnectorCatalogue.name == name)
                    .where(ConnectorCatalogue.provider == "azure_blob")
                )
                result = await session.execute(stmt)
                row = result.scalars().first()
                if row is None:
                    raise ValueError(f"No azure_blob connector named {name!r} found. Check Settings → Connectors.")
                cfg            = row.provider_config or {}
                account_url    = (cfg.get("account_url") or row.host or "").strip()
                container_name = (cfg.get("container_name") or row.database_name or "").strip()
                from urllib.parse import urlparse
                if not urlparse(account_url).netloc:
                    raise ValueError(
                        f"Azure Blob connector {name!r} has an invalid Storage Account URL: {account_url!r}. "
                        f"Set the full URL in Settings → Connectors, e.g. https://youraccount.blob.core.windows.net"
                    )
                if not container_name:
                    raise ValueError(f"Azure Blob connector {name!r} has no container_name.")
                return {"account_url": account_url, "container_name": container_name}

        try:
            asyncio.get_running_loop()
            with _cf.ThreadPoolExecutor() as pool:
                return pool.submit(asyncio.run, _fetch()).result(timeout=10)
        except RuntimeError:
            return asyncio.run(_fetch())
    except Exception as exc:
        raise RuntimeError(f"Blob connector lookup failed: {exc}") from exc



# ── Stage 4-8 helpers (inlined) ─────────────────────────────────────────────

# ── Stage IDs ─────────────────────────────────────────────────────────────────
_STAGE_CLASSIFICATION  = 4
_STAGE_EXTRACTION      = 5
_STAGE_EMBEDDINGS      = 6
_STAGE_PRICE_BENCHMARK = 7
_STAGE_COMPLETE        = 8


# ── Classification constants ───────────────────────────────────────────────────
_CLASSIFICATION_MAP = {"E-Auction": "E-Auction Results", "Other": "Others"}
_SUPPORTED_CLASSIFY_EXTS = {
    ".xlsx", ".xls", ".csv", ".pdf", ".docx", ".doc",
    ".pptx", ".ppt", ".txt", ".html", ".htm",
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp",
}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"}
_MAX_CLASSIFY_CHARS = 100000  # safety cap inside each extractor only — the
                              # *active* classification truncation happens in
                              # _classify_file via the configurable cls_max_chars
                              # (default 24000 ≈ 6000 tokens, matches doc-intel
                              # truncate_to_token_limit max_content_tokens=6000).

_EXT_TO_TYPE = {
    ".xlsx": "excel",  ".xls": "excel",  ".xlsm": "excel",  ".xlsb": "excel",  ".csv": "csv",
    ".pdf":  "pdf",    ".docx": "word",  ".doc": "legacy_doc",
    ".pptx": "pptx",  ".ppt":  "legacy_doc",
    ".txt":  "text",   ".html": "html",  ".htm": "html",
    ".xml":  "text",   ".json": "text",  ".eml": "text",
    ".png":  "image",  ".jpg":  "image", ".jpeg": "image",
    ".tiff": "image",  ".tif":  "image", ".bmp":  "image",
    ".msg":  "msg",
}

VALID_CLASSIFICATIONS = {"RFQ", "Quotation", "MPBC", "BER", "E-Auction", "Other"}

# ── Classification system prompt (full original recovered from file_classifier) ──
CLASSIFICATION_SYSTEM_PROMPT = """You are an expert document classifier for Motherson's procurement team. You will be given the content of a file (extracted text, tables, sheet structure, or an image) and you must classify it into EXACTLY ONE of six categories.

============================================================
CRITICAL: MULTILINGUAL SUPPORT
============================================================
Documents may be in ANY language: English, German, Czech, Hindi, French, Spanish, Chinese, Japanese, Hungarian, and others. You MUST:
- Match fields by MEANING, not by exact English labels
- Recognize translated equivalents. Common examples:
  • German: "Preisspiegel" / "Angebotsvergleich" = Bid Comparison (MPBC), "Angebot" = Quotation/Offer, "Anfrage" = Inquiry/RFQ, "Lieferant" = Supplier, "Preis" = Price, "Lieferzeit" = Delivery Time, "Zahlungsbedingungen" = Payment Terms, "Genehmigung" = Approval, "Begründung" = Justification
  • Czech: "Nabídka" = Offer/Quotation, "Poptávka" = RFQ/Inquiry, "Dodavatel" = Supplier, "Cena" = Price, "Technický požadavek" / "Lastenheft" = Technical Specification
  • Hindi: "कोटेशन" = Quotation, "मूल्य" = Price, "आपूर्तिकर्ता" = Supplier
- Apply the same field-matching logic regardless of language
- If a document is in a language you recognize, translate field names mentally and match against the English category definitions

============================================================
CLASSIFICATION METHOD
============================================================
Classification is FIELD-BASED VALIDATION. Each category has a list of mandatory fields. You must verify the presence of those fields in the document content (allowing for naming/labeling variations, synonyms, AND translations in any language — the underlying meaning matters, not the exact label). The category whose mandatory fields are most completely satisfied wins.

============================================================
CATEGORIES & MANDATORY FIELDS
============================================================

1. **MPBC** — Motherson Purchase BID Comparison
   Official Motherson template that consolidates quotations from multiple suppliers (typically 3+) into a side-by-side comparison for procurement evaluation. Sheet usually titled "MPBC" or "Motherson Purchase BID Comparison". May be a multi-sheet workbook with sheets like "1. MPBC", "2. mandatory cells", "BER", "Supp X risk", "exchange_rates", etc.

   Mandatory fields (yellow-marked in the official MPBC template — must find most of these):
   PROJECT / HEADER:
     • Project Detail
     • RAS Number
     • Sheet No.
     • Indenter, Originator or RFQ Responsible
     • Contact Number
   PER-SUPPLIER (repeated for each of typically 3+ suppliers):
     • Name of the Supplier
     • Preference (Preferred source / Second Source / etc.)
     • Quote meet spec (Yes/No)
     • Supplier's Contact Person
     • Offer No.
     • Offer Date
     • Offer Validity
     • Supplier Country / Classification
     • Country of Origin / Classification
   PRICING:
     • Currency for Comparison
     • Total Amount in supplier currency
     • Total Amount in EUR
     • Saving %age
   SOURCING DECISION:
     • Insourcing
     • Hybrid
     • Single sourcing
     • Supplier justification (BER) signed
     • In case of Cust. funded: Revenue in EUR
     • Target against Budget
   COMMERCIAL TERMS (per supplier):
     • Incoterms
     • Packing and Forwarding
     • Freight
     • Insurance
     • Taxes
     • Customs / duties
     • Installation
     • Delivery Period (Weeks)
     • Advance payment (any payment before delivery to MOTHERSON)
     • SCF (Supply Chain Finance Program)
     • Payment Terms
   APPROVAL / FINAL:
     • GSP Purchase Saving %age
     • GSP Purchase Saving amount in EUR
     • Approved Supplier / Classification
     • Approved Cost (total amount)
     • Landed Cost (for Reference)
     • Approvals (Department / Name / Date / Approver Signature for Purchasing, Technics, Controlling, Sales, Plant Manager)

   Strongest distinguishing signals:
     • Title / sheet name contains "MPBC" or "Motherson Purchase BID Comparison"
     • Supplier 1 / Supplier 2 / Supplier 3 columns appearing side-by-side
     • Presence of RAS Number, GSP (Global Strategic Procurement) terminology
     • Multi-sheet Excel with sheets like "1. MPBC", "2. mandatory cells", "BER", "Supp X risk"
     • MULTILINGUAL: German "Preisspiegel" (price mirror) / "Angebotsvergleich" (offer comparison) / "Bid comparison (Purchased parts)" = MPBC. If 3+ suppliers are compared side-by-side in ANY language with pricing, it is MPBC even without the exact Motherson MPBC template fields like RAS Number.
   Allow naming variations / synonyms / translations; the MEANING of fields matters, not exact labels.

2. **Quotation** — A SINGLE vendor's price offer / response to an RFQ (vendor → buyer such as Motherson)
   Synonyms commonly used as the document title: "Quotation", "Quote", "Estimate", "Offer", "Proposal", "Price Bid". Treat all of these as Quotation candidates.

   Mandatory fields (naming may vary widely; identify by MEANING — Indian, German, and other vendor styles all appear):
     • Vendor Name (vendor's company name appears in letterhead at top AND in signature/footer like "For [Company Name]")
     • Vendor Address (postal address of the vendor)
     • Date of Quotation (e.g., "Date:", "Quotation Dt.", "Offer Date", "Est. Date")
     • Total Amount / Value (grand total / net amount / total in supplier currency)
     • Payment Terms (e.g., "30 days", "100% advance", "Against delivery", "100% after PO confirmation")
     • Delivery Time (e.g., "15-20 days", "within 2-3 weeks", "Delivery Period", "Lead Time", "Dispatch Time")
     • Validity (e.g., "Quotation valid until ___", "Offer Validity: 60 days", "Valid for one month")
     • Specifications (technical / product / scope specifications — may be a separate "TECH SPEC" sheet in Excel)
     • Item Description (line items with descriptions — columns like "Description of Goods", "Item Description", "Sr. No. + Description + Qty + Rate + Amount", "BOQ items")

   Additional supporting signals frequently observed (not all needed, but strengthen confidence):
     • Vendor letterhead at top: company logo / name + address + phone + email + GST No. / GSTIN / PAN No. / VAT No.
     • A reference / quotation number (e.g., "Quotation No.: PPAQ004560", "Ref: Q/875", "Offer No.", "Est. No.: SCL/Haryana/20-21/21", "Ref: MIPL/MATE/...")
     • Addressed to: "To, M/s [Customer Name]", "Client:", "Kind Attn: Mr./Ms. ___"
     • Formal letter language: "We are pleased to offer / quote", "With reference to your enquiry", "Thanking you", "Yours faithfully"
     • Closing: "For [Vendor Company]" + "Authorised Signatory" / "Proprietor" / "Sales Manager"
     • Tax columns: HSN/SAC Code, CGST %, SGST %, IGST %, Excise Duty, GST @ 18 %
     • Bank details / E. & O. E. / "Please mention our quotation number on your purchase order"

   CRITICAL DISAMBIGUATION:
     • A Quotation is from ONE vendor. If THREE OR MORE vendor names appear as parallel columns/sections being compared → it is MPBC, not Quotation.
     • A Quotation MAY be a multi-sheet Excel (e.g., "COMMERCIAL", "TECH. SPEC.", "Summary", "NPV", "Vendors" sheets) but with only ONE vendor. Do not misclassify these as MPBC — check vendor count, not sheet count.
     • If price columns are EMPTY (template for vendor to fill) → it is RFQ, not Quotation.
     • Filename is unreliable — e.g., "(875) CK Motherson Auto Hitech.docx" was issued BY Hi-Tech to CK Motherson; the customer name in the filename does not make it MPBC or RFQ.
     • PDF extraction may produce duplicated characters from font issues (e.g., "TTOO" instead of "TO", "QQUUOOTTAATTIIOONN" instead of "QUOTATION") — interpret semantically.
     • **A Quotation can use the buyer's RFQ TEMPLATE FORMAT.** When a supplier (e.g., HAITIAN, Engel, Sumitomo, KraussMaffei, Arburg) fills in an RFQ specification template with their technical responses ("STANDARD", "OPTIONAL", "OPTIONAL - AVAILABLE", option codes) AND populates pricing fields (Basic Price, Option Price, Net Price, Total Machine Cost, Delivery Price, etc. with actual monetary values in USD/EUR/INR), the document has become a Quotation — the supplier's price offer. Key differentiator: if ACTUAL PRICES are filled in (not blank "to be specified" placeholders) and a supplier contact person / company name is prominently featured, treat it as Quotation even if the original RFQ structure ("Pls Specify", "Required") is retained.
     . There can be many vendors/supplier not just HAITIAN, Engel, Sumitomo, KraussMaffei, Arburg etc. these are just examples.

3. **RFQ** — Request For Quotation (Motherson → vendors)
   A specification / scope document issued BY Motherson (the buyer) TO vendors, asking them to submit a quotation. The document defines WHAT Motherson needs and asks the vendor to respond with specs and pricing. May be an Excel template with columns for "Supplier Spec/Confirmation" that are blank (or have been partially filled by a responding vendor).

   Mandatory fields:
     • Project Name (e.g., "INJECTION MOULDING MACHINE SPECIFICATIONS", "30KLD STP", or a specific project title)
     • RFQ Number (reference / inquiry number — e.g., "MATEB/IMG 1300/2020/SEPT/1-Rev 01", "MATE SADDLES/SP/01-24/Rev03")
     • Date (issue date)
     • Specifications — the BULK of the document; detailed technical / scope specifications structured as tables with:
       - Sl. No / Item number
       - Description / technical parameter (e.g., "Screw diameter [mm]", "Clamping force [kN]")
       - "Required" / "Not Required" / "Std" flags (Motherson's requirement)
       - "To be specified by the supplier" / "Supplier Spec/Confirmation" / "Pls specify" columns (blank or for vendor to fill)
       - Multiple specification sections: e.g., Injection unit, Clamping unit, Electrical/Hydraulic, Controller, General features, Spares, Special features
     • Commercials — a section (typically at the bottom) asking for pricing breakdown:
       - Option Price, Basic Price, Gross price, Special discount, Net price
       - Spares Package, Sea worthy packing, CIF Cost + Insurance
       - Startup cost, Delivery Price, Price in INR
       - Payment terms, Taxes, Delivery period, Warranty, Guarantee

   Strong supporting signals:
     • Sheet name contains "RFQ" (e.g., " RFQ", "RFQ - STP")
     • "MATE-B Req" or "MATE Spec" column (Motherson Automotive Technologies & Engineering — a Motherson entity)
     • "Supplier Spec/Confirmation" or "Supplier Remarks" columns (empty = template, filled = vendor response captured in same template)
     • "Techno Commercial Comparison" as a title (still an RFQ when it's the request template, even if some vendor responses have been captured)
     • Accompanying sheet " Parts & Mold data " with part specifications (part name, size, weight, wall thickness)
     • Columns with "Fill this column by '0' if STD or Enter the cost if Optional" — template instruction language
     • "Company Name:", "Contact Person:", "Telephone:", "Email:" fields (for vendor to fill)

   CRITICAL DISAMBIGUATION:
     • An RFQ is the TEMPLATE / REQUEST document. Even if one or two vendors have filled in their spec responses, the document may still be an RFQ IF:
         (a) pricing fields remain EMPTY or say "to be specified", AND
         (b) no supplier contact/company is featured prominently.
       HOWEVER, if a supplier has populated pricing (Basic Price, Net Price, Total cost, Option prices with actual USD/EUR/INR values) AND the supplier's contact details appear prominently (name, phone, email), the document has become a QUOTATION — the supplier's completed price offer using the RFQ template format. Classify as Quotation, not RFQ.
     • If the document compares 3+ vendors side-by-side in a CONSOLIDATED EVALUATION format → it is MPBC, not RFQ.
     • If the document is a standalone vendor letter with prices (no Motherson spec template structure) → it is Quotation, not RFQ.
     • RFQ documents are typically MUCH longer / more detailed than Quotations (100+ rows of specifications with Required/Not Required flags). BUT length alone does not determine the category — a long spec sheet with filled-in prices is still a Quotation.

4. **BER** — Bid Exception Report
   SPECIFICALLY the Motherson "BID EXCEPTION REPORT" template form. This is NOT a catch-all for any waiver, justification, or single-source document. The document must use the Motherson BER template structure.

   Mandatory fields (ALL of these must be present or nearly all — this is a strict template match):
     • "BID EXCEPTION REPORT" appearing explicitly as a header / title (this exact phrase or very close equivalent is REQUIRED — a generic "Waiver of Competition" or "Single Source Justification" title is NOT sufficient for BER classification)
     • Reasoning for not obtaining at least three bids/quotes (NOTE: language may differ in 20–40% of cases)
     • Order Value field (e.g., "Order Value: 671,57 €", "Budget Line Ref. + amount")
     • Description of the product or service to be ordered (a labeled section)
     • Justification for waiver of competitive bidding (a labeled section)
     • Justification of the product or service to be ordered (a separate labeled section)

   Strong supporting signals (Motherson BER template-specific — presence of 2-3 is a near-certain indicator):
     • Header "Capital Equipment & Indirect Purchasing"
     • Budget Line Ref. (e.g., "Budget Line Ref.: ID 413128")
     • Reference to "LCC Suppliers" / "Low Cost Country" / "Motherson internal company"
     • Checkbox-style options A through E for reasoning (A: less than three potential bidders; B: sole-source / proprietary item; C: national/local supply contract; D: similar item purchased in past 6 months; E: Other reasons)
     • Three approval rows: "Prepared by:" + "Purchasing approval:" + "Managing Director / COO / EVP approval:"
     • "Approver comments:" section
     • Footer with template revision history (e.g., "Compiled by Mac Cheema", "Revised by Lousie Osgood")
     • Sheet named "BER" in an Excel workbook

   CRITICAL: Do NOT classify as BER if the document is:
     • A generic "Waiver of Competition" form (different template, different structure)
     • A "Single Source Justification" from a non-Motherson template
     • A "Double Source Waiver" or supplier approval form
     • Any competitive bidding exception document that does NOT use the specific Motherson "BID EXCEPTION REPORT" template with the A-E checkbox structure
     These should be classified as "Other" instead.

5. **E-Auction** — E-Auction results / reports / trackers
   Documents generated from or summarizing an online reverse-auction event where vendors bid in real time. May be raw auction output from an e-procurement platform OR a summary/tracker/presentation consolidating auction results.

   Mandatory fields (from client spec — match by meaning, naming may vary):
     • Event ID (e.g., "Doc3071019131", "Auction ID")
     • Event Name (e.g., "Live Japanese Reverse eAuction - R_268526-2026-EPP boxes...")
     • Publish Date
     • Open Date
     • Close Date
     • BID Id (e.g., "ID3240307353", "ID3244896110")
     • BID Status (e.g., "Accepted", "Default")
     • Participant (vendor / bidder name — e.g., "MORAplast, s.r.o(Sarka Bris...)")
     • Basic Price (per unit price)
     • Extended Price (total price for volume)

   Strong supporting signals (presence of 3+ is near-certain):
     • "eAuction", "e-Auction", "Reverse Auction", "Japanese Auction", "Japanese Reverse eAuction" in title/headers
     • Sheet names like "Overview Sheet", "Full Bid Data Sheet"
     • Event Type field (e.g., "Japanese Auction")
     • Owner field with eAuction email (e.g., "eAuction.GSP@motherson.com")
     • Rank column (bidder ranking: 1, 2, 3...)
     • Savings column / Total Cost column
     • Capacity Planning Volume
     • Report Generated Date
     • Currency field (e.g., "European Union Euro")
     • Submission Date with timestamps (time-stamped bids)
     • Pricing tiers: "Price 1" (initial quotation), "Price 2" (revised/negotiated), "Price 3" (after eAuction)
     • Saving definitions: "Pre Auction Saving", "eAuction Saving", "Total Saving"
     • Tracker / summary workbooks: sheets like "Pivot", "Summary", "Project Details", "Project Negotiation" with columns for Auction type, GSP Buyer, Group company, L1 price, Final price, Savings
     • Presentation slides with "eAuction Overview", "eAuction Status", monthly savings summaries
     • Project Names referencing RFQ numbers (e.g., "P352----RFQ16688---...")

   IMPORTANT — Two tiers of E-Auction documents (BOTH classify as E-Auction):
     TIER 1 (raw auction output): Contains most mandatory fields above (Event ID, BID Id, Participant, prices). High confidence.
     TIER 2 (summaries/trackers/presentations): May NOT contain individual event-level fields like Event ID or BID Id, but IS clearly ABOUT e-auction results — contains "eAuction" keyword prominently + pricing data (Price 1/2/3, savings, L1 price, Final price) + auction metadata (Auction type, Auction Month, GSP Buyer). Classify as E-Auction with moderate confidence (0.75-0.85) even if individual bid-level fields are missing.

   CRITICAL DISAMBIGUATION:
     • E-Auction documents focus on the AUCTION EVENT and BIDDING PROCESS — they have event metadata (IDs, dates, event type) and bid-level data (bid IDs, statuses, timestamps, ranks). This is fundamentally different from MPBC which is a static comparison table.
     • An MPBC compares vendor quotations side-by-side; an E-Auction shows a time-sequenced bidding process with event infrastructure fields.
     • E-Auction tracker/summary documents (multi-sheet workbooks tracking many auction events across a fiscal year) are still E-Auction, not "Other".
     • Presentation files (.pptx) summarizing eAuction results/savings are still E-Auction, not "Other".
     • If the document is ABOUT eAuction (mentions "eAuction" + savings/prices), classify as E-Auction even if not all 10 mandatory fields are present.

6. **Other** — Anything that does NOT satisfy the mandatory field set of any category above
   Examples: invoices, purchase orders, delivery notes, contracts, drawings, internal memos, generic emails, MSAs, NDAs in isolation, quality reports, etc. Use this only when the document genuinely fails the field checks for all five named categories.

============================================================
DECISION PROCESS (follow strictly)
============================================================
Step 1 — Scan the document and identify which mandatory fields from each category are PRESENT (allowing synonyms / paraphrases / equivalent column names).
Step 2 — For each candidate category, compute a coverage ratio: (fields present) / (total mandatory fields).
Step 3 — Pick the category with the highest coverage. To classify as that category, coverage must be ≥ 60% AND the strongest distinguishing signal must be present (e.g., for MPBC: ≥ 3 vendors compared; for BER: explicit BER title + waiver justification; for E-Auction: event + bid columns).
Step 4 — If no category reaches 60% coverage, classify as "Other".
Step 5 — Filename can be a hint but content always wins. Ignore filename if content contradicts it.
Step 6 — Confidence calibration:
   • ≥ 0.90 → all or nearly all mandatory fields present and unambiguous
   • 0.75 – 0.89 → most mandatory fields present; minor ambiguity
   • 0.60 – 0.74 → enough fields present to classify but with notable gaps
   • < 0.60 → genuine uncertainty — likely "Other"

============================================================
OUTPUT FORMAT (strict JSON, no markdown, no commentary)
============================================================
{
  "classification": "RFQ" | "Quotation" | "MPBC" | "BER" | "E-Auction" | "Other",
  "confidence": <float 0.0 - 1.0>,
  "reason": "<2-3 sentences explaining which mandatory fields you matched and which (if any) were missing>",
  "key_signals": ["<mandatory field matched 1>", "<mandatory field matched 2>", "<mandatory field matched 3>"],
  "fields_matched": ["<exact field/column/phrase observed in the document>", ...],
  "fields_missing": ["<mandatory fields that were NOT found>", ...]
}

Rules:
- Pick exactly ONE category.
- key_signals = the top 3-5 mandatory fields that drove the decision.
- fields_matched = up to 10 specific evidence items (column names, headers, phrases) you actually observed.
- fields_missing = mandatory fields for the chosen category that you could NOT find.
- Be specific in the reason — cite real evidence from the document, do not be vague."""

# ── Classification user prompt templates ──────────────────────────────────────
_CLASSIFY_USER_TEXT = """Classify the following file by checking its content against the mandatory fields for each category.

============================================================
FILE METADATA
============================================================
Filename: {filename}
File Type: {file_type}
{extra_metadata}
============================================================
EXTRACTED CONTENT
============================================================
{extracted_content}
============================================================

Apply the field-based decision process from the system prompt and return only the JSON object."""

_CLASSIFY_USER_IMAGE = """Classify the following file based on the image provided.

============================================================
FILE METADATA
============================================================
Filename: {filename}
File Type: {file_type}
{extra_metadata}
============================================================

Examine the image carefully — read every visible field, column, and label. Check the mandatory fields for each category against what you see. Apply the decision process from the system prompt and return only the JSON object."""

# ── Extraction prompt constants ────────────────────────────────────────────────
EXTRACTION_SYSTEM_PROMPT = """You are a senior procurement analyst with deep experience evaluating supplier quotations across every spend category. Your job here is to read each quotation thoroughly and produce a complete, accurate, decision-grade extraction in a strict JSON schema. The downstream system uses your output to benchmark prices, pick winning suppliers, and approve purchase requisitions, so completeness and correctness directly affect business decisions — missed line items, wrong DTL_ID mappings, and inflated confidence scores all propagate into bad recommendations.

The procurement system processes all types of purchase requisitions — industrial equipment, IT hardware and software, raw materials, consumables, engineering services, facility management, vehicle hiring, consulting, pharmaceuticals, construction, and any other category of goods or services. Your extraction must work equally well across all of these — do not assume any specific industry.

Key responsibilities:
- Read the entire quotation before answering. Do not stop at the first item table — quotations often have continuation pages, addenda, optional items, and supporting charges.
- Extract every item line, pricing, supplier information, and commercial terms from the document.
- Match each extracted item to the right purchase requisition DTL_ID using the strongest distinguishing signal available (model number, code, capacity, location, quantity, position).
- Translate all extracted text to English regardless of the source language of the document.
- Return data in the exact JSON schema specified — no markdown fences, no commentary, no analysis text, just the JSON object.
- Handle diverse document formats: formal quotations, proforma invoices, price lists, rate cards, cost estimates, email quotations, and scanned documents.

Extraction guidelines:
- Be precise with numbers: prices, quantities, dates.  Never hallucinate or infer figures not present in the document.
- Distinguish between unit price and total price.  total_price = unit_price × quantity unless the document states otherwise.
- Identify currency from the document (symbols like ₹, $, €, £, AED, ZAR, or explicit ISO codes).  Return the ISO-4217 three-letter code (INR, USD, EUR, GBP, ZAR, AED, …).
- For services, "quantity" may be hours, days, trips, or similar — capture the unit exactly as written.
- Extract payment terms verbatim then normalise (e.g. "100% advance" → "100% Advance", "net 30 days" → "Net 30").
- For the hierarchical item taxonomy (item_level_1 … item_level_8), classify from the broadest category down to the most specific attribute available.  See the taxonomy guidelines in the user message.
- Set supplier_match_conf honestly based on how well the extracted item matches the requisition line.  0.0 = completely unrelated, 1.0 = identical match certain.  A lower honest score is always better than an inflated wrong one.
- If a field genuinely cannot be determined from the document, set it to null.  Never guess or fill with placeholder text.
- Null data from the requisition context (shown as "N/A") means that field was not recorded — do not treat it as a matching signal."""

ITEM_TAXONOMY = """Guidelines for item_level_1 through item_level_8 hierarchical taxonomy:

The eight levels represent a progressively narrower classification of each item.  Fill as many levels as the quotation document supports; leave deeper levels as null when information is not available.

IMPORTANT: item_level_1, item_level_2, and item_level_3 MUST always be filled — you can always infer the broad category, sub-category, and product/service type from the item name and description.  Only item_level_4 through item_level_8 may be null when the information is genuinely absent.

Level 1 — Broad industry / domain category
  Examples: "Electronics", "Mechanical", "Civil", "IT Hardware", "Services",
            "Consumables", "Furniture", "Vehicles", "Raw Materials", "Safety Equipment"

Level 2 — Sub-category within the domain
  Examples: "Industrial Equipment", "Office Equipment", "Construction Materials",
            "Software Licensing", "Fleet Management", "Chemical Supplies"

Level 3 — Product type or service type
  Examples: "Temperature Controller", "Laptop", "Concrete Mix", "Annual Maintenance Contract",
            "Vehicle Hiring", "Mold Machine Component"

Level 4 — Brand / Manufacturer
  Examples: "Dell", "Siemens", "Matsui", "Bosch", "Toyota", "Tata Motors"

Level 5 — Model / Series / Part number
  Examples: "Latitude 5540", "MCLX-350A-0", "PowerEdge R740", "Hilux 2.8 GD-6"

Level 6 — Configuration / Variant
  Examples: "16 GB RAM / 512 GB SSD", "3-phase 440 V 50 Hz", "4×4 Double Cab",
            "Stainless Steel Grade 304"

Level 7 — Additional specification
  Examples: "With touchscreen", "IP65 rated", "CE certified", "Left-hand drive"

Level 8 — Any remaining distinguishing detail
  Examples: "Custom colour RAL 7035", "Extended warranty 5 yr", "Ex-works delivery"

Rules:
- Start at Level 1 and fill downward; never skip a level then fill a deeper one.
- Never repeat the same information across two levels.
- Use proper nouns for brands and models (Levels 4–5).
- Use descriptive noun phrases for categories (Levels 1–3).
- When the quotation is for a service rather than a physical product, adapt the hierarchy:
    L1=Services, L2=domain, L3=service type, L4=provider, L5=scope, …"""

EXTRACTION_USER_TEMPLATE = """## Purchase Requisition Context

The following data was recorded at the time the purchase requisition was raised.
Fields marked "N/A" were not captured in the source system — treat them as unknown.

### Header Information

| Field | Value |
|---|---|
| RAS Number | {purchase_req_no} |
| Requisition ID | {purchase_req_id} |
| RAS Title | {ras_title} |
| Requisition Type | {requisition_type} |
| Classification | {classification} |
| Justification | {justification} |
| Primary Supplier | {supplier_name} |
| Supplier Address | {address} |
| Parent Supplier | {parent_supplier} |
| Supplier Type | {supplier_type} |
| Supplier Country | {supplier_country} |
| Currency (RAS) | {currency} |
| Purchase Value | {purchase_value} |
| Enquiry No | {enquiry_no} |
| Contract No | {contract_no} |
| Order No | {order_no} |
| Department | {department} |
| Negotiated By | {negotiated_by} |
| Category Buyer | {category_buyer} |
| Purchase Category | {purchase_category} |
| Category L1 | {category} |
| Category L2 | {sub_category} |
| Category L3 | {l3} |
| Category L4 | {l4} |
| Site | {site} |
| Region / Country | {site_region} / {site_country} |
| Division | {division} |
| Payment Days | {payment_days} |
| PO Date | {po_date} |

### Line Items from the Requisition

The table below lists the items the buyer expects to find in the quotation.
Column legend:
- **DTL_ID** — unique line-item identifier; use this value for `purchase_dtl_id` when you match a quotation item
- **Item Code** — internal material/part code (may be N/A)
- **Unit Price / Original Value / Initial Offer / Negotiated** — price history; use for cross-validation only
- **Req Value** — total requisition value for the line (Qty × Unit Price)
- **Prepayment / Payment Terms** — buyer-side payment conditions

> **Matching tip:** Suppliers often use their own product codes (e.g. LP100, PT050SPEC, QMC122) that differ from the RAS description.
> Use secondary signals to match when names differ:
> - Machine tonnage / capacity (e.g. "650T", "350 tons") appearing in both the PDF and RAS description
> - Quantity, UOM, or price proximity to a RAS line
> - Machine model or site location if mentioned in both

{line_items_table}

---

### Additional RAS Reference Data

> **Important — treat as reference only.** The fields below come directly from
> the procurement system database. Users sometimes enter incomplete or incorrect
> data, so these values are hints, not ground truth.
> - **Always extract actual prices, quantities, dates, and descriptions from the
>   quotation document itself.**
> - Use this data only to help identify which `DTL_ID` corresponds to which item
>   in the quotation — look for matching machine specs, item codes, quantities,
>   tonnage, site, or any other common signal between the DB row and the PDF item.
> - If a DB field contradicts what the quotation clearly states, trust the quotation.

{raw_ras_context}

---

## Item Taxonomy Guidelines

{item_taxonomy}

---

## Quotation Document

{document_content}

---

## Required Output

Analyse the quotation document above against the requisition context.
Extract **exactly one item per DTL_ID** from the RAS line items table — the quotation row that best represents the supplier's quoted price and specs for that line.
For each item match it to the closest line item in the requisition table
(use Description, Item Code, Quantity, and UOM as primary matching signals).

> **Important:** A purchase requisition can cover any type of procurement —
> goods (equipment, components, materials, consumables), services (maintenance,
> consulting, hiring, IT), or mixed orders.  Adapt your extraction accordingly.

Return a **single JSON object** — no markdown fences, no extra text:

{{
  "supplier_name": "string or null",
  "supplier_address": "string or null",
  "supplier_country": "string or null",
  "quotation_ref_no": "string or null",
  "quotation_date": "YYYY-MM-DD or null",
  "currency": "ISO-4217 three-letter code or null",
  "validity_date": "YYYY-MM-DD or null",
  "validity_days": "integer or null",
  "payment_terms": "string or null",
  "items": [
    {{
      "purchase_dtl_id": "integer from DTL_ID column if matched, else null",
      "supplier_match_conf": "0.0 to 1.0 — your confidence this item matches the requisition line",
      "item_name": "canonical item name in English",
      "item_description": "full description with all specs in English",
      "quantity": "number or null",
      "unit": "unit of measurement or null",
      "unit_price": "number or null",
      "total_price": "number or null",
      "discount": "discount amount or percentage as number, or null",
      "taxation_details": "tax/GST/VAT info as string or null",
      "delivery_date": "YYYY-MM-DD or null",
      "delivery_time_days": "integer number of days or null",
      "item_level_1": "broadest category (e.g. Industrial Equipment, IT Hardware, Services)",
      "item_level_2": "sub-category (e.g. Pumps & Valves, Laptops, Facility Management)",
      "item_level_3": "product or service type",
      "item_level_4": "brand or manufacturer if known, else null",
      "item_level_5": "model or series if known, else null",
      "item_level_6": "configuration or variant if known, else null",
      "item_level_7": "key technical specification if known, else null",
      "item_level_8": "any additional distinguishing detail, else null",
      "commodity_tag": "lowercase-slug-tag (e.g. industrial-pump, it-laptop, vehicle-hiring)",
      "item_summary": "plain-English summary of the item and its intended use (max 20 words)"
    }}
  ]
}}

### Rules

**Required fields — never null:**
- `item_name` — always present; use the supplier's exact product/service name from the document.
- `item_description` — always present; include all specs, model numbers, and technical details visible in the document.
- `item_level_1`, `item_level_2`, `item_level_3` — always fill; you can always infer broad category, sub-category, and product type from the item name alone.
- `commodity_tag` — always present; derive from the item type.

**Pricing rules:**
- Extract `unit_price` and `total_price` separately whenever both are shown.
- If only a total price is shown and quantity > 1: `unit_price = total_price / quantity`.
- If only one price is shown and quantity is 1 or absent: treat it as both unit_price and total_price.
- Never leave both unit_price and total_price null if any price figure appears in the document for that item.

**General rules:**
- The quotation document is the primary source of truth. Extract all values (prices, quantities, dates, descriptions) from it — not from the RAS reference data.
- All text fields MUST be in English — translate from any source language.
- Dates must be ISO 8601 (YYYY-MM-DD).  Convert formats like "15-Apr-2026" or "15/04/26" accordingly.
- Prices as plain numbers — strip currency symbols, commas, and thousand separators.
- If a field genuinely cannot be determined from the document, return null (not an empty string, not "N/A").
- item_level_4 through item_level_8 may be null when information is genuinely absent — do not invent brand or model names.
- supplier_match_conf is confidence that the extracted item matches its mapped requisition line (0.0 = no match, 1.0 = certain).

**DTL_ID matching — CRITICAL:**

You MUST attempt to map every quotation item row to a `purchase_dtl_id` from the RAS line items table. Do not return `null` for `purchase_dtl_id` unless the item is clearly a non-line charge (shipping, packing, GST, freight, insurance, surcharge, taxes).

Use these signals to pick the right DTL_ID for each item — pick whichever signals are strongest in the documents you actually see. Do not invent signals that aren't there:

1. **Distinguishing identifiers shared between the item and the RAS line** — model numbers, part codes, sizes, capacities, dimensions, ratings, or any alphanumeric identifier that uniquely names the product/service variant. Match these exactly. If the RAS line says "650T machine" and the quotation item says "650T machine", that's a match. If the supplier uses their own internal code (e.g. their model "X1000-A") that maps to the same physical/logical thing as the RAS line, use that mapping.
2. **Site / location** — when both the RAS line and the quotation item name a site, region, or location, that's a strong signal.
3. **Quantity and UOM** — when only one RAS line matches the quantity / unit, use that DTL_ID.
4. **Item description overlap** — significant word overlap in the descriptive text.
5. **Position in the table** — if the quotation lists items in the same order as the RAS table and other signals are absent, use position as a last-resort signal.

For supporting items in the quotation (installation, packing, commissioning, training, freight when itemised, etc.), map each to the DTL_ID of the **product/service they support**, unless the RAS has a separate DTL_ID specifically for that supporting item — in which case prefer that one.

If the same quotation item legitimately covers multiple DTL_IDs (e.g. one quote line for "Installation for both Machine A and Machine B"), emit one row per DTL_ID with the same item details.

**One row per DTL_ID — STRICT:**
Return **at most one item per `purchase_dtl_id`**. If the supplier quoted multiple options, configurations, or variants for the same line (e.g. standard vs premium, different lead times, revised prices), pick the single row that best represents the primary quoted offer — the one with the clearest pricing and the closest match to the RAS description. Do not emit duplicate rows for the same DTL_ID.

**Set `supplier_match_conf` honestly:**
- `1.0` — exact match on a distinguishing identifier (model number, code, size).
- `0.7-0.9` — strong but not perfect match (one signal matches clearly).
- `0.3-0.6` — weak match (only position or partial description overlap).
- `< 0.3` — guessing — prefer to leave `purchase_dtl_id = null`.

**Coverage check before responding:**
After listing all items, verify:
1. For every DTL_ID in the RAS table that the supplier has clearly quoted on, you have emitted **exactly one** item with that `purchase_dtl_id`. If you forgot a DTL_ID, add it now. If you emitted more than one row for the same DTL_ID, remove the duplicates and keep only the best one.
2. No `purchase_dtl_id` appears more than once in your `items` array."""


# ── Dataclasses ────────────────────────────────────────────────────────────────

@dataclass
class LineItemContext:
    purchase_dtl_id:  int
    purchase_req_id:  int
    item_no:          int
    quantity:         Optional[Decimal]
    item_type:        Optional[str]
    item_description: Optional[str]
    unit_price:       Optional[Decimal] = None
    uom:              Optional[str]     = None
    supplier_name:    Optional[str]     = None
    discount:         Optional[Decimal] = None
    req_value:        Optional[Decimal] = None
    currency:         Optional[str]     = None
    delivery_date:    Optional[str]     = None
    payment_details:  Optional[str]     = None
    original_value:   Optional[Decimal] = None
    initial_offer:    Optional[Decimal] = None
    negotiation:      Optional[Decimal] = None
    comments:         Optional[str]     = None
    prepayment:       Optional[str]     = None
    item_code:        Optional[str]     = None


@dataclass
class RASContext:
    purchase_req_no:   str
    purchase_req_id:   int
    supplier_name:     Optional[str]
    justification:     Optional[str]
    currency:          Optional[str]
    enquiry_no:        Optional[str]     = None
    classification:    Optional[str]     = None
    department:        Optional[str]     = None
    negotiated_by:     Optional[str]     = None
    address:           Optional[str]     = None
    contract_no:       Optional[str]     = None
    order_no:          Optional[str]     = None
    purchase_value:    Optional[Decimal] = None
    category:          Optional[str]     = None
    sub_category:      Optional[str]     = None
    site_country:      Optional[str]     = None
    site_region:       Optional[str]     = None
    site:              Optional[str]     = None
    division:          Optional[str]     = None
    requisition_type:  Optional[str]     = None
    parent_supplier:   Optional[str]     = None
    supplier_type:     Optional[str]     = None
    supplier_country:  Optional[str]     = None
    payment_days:      Optional[str]     = None
    po_date:           Optional[str]     = None
    category_buyer:    Optional[str]     = None
    l3:                Optional[str]     = None
    l4:                Optional[str]     = None
    purchase_category: Optional[str]     = None
    ras_title:         Optional[str]     = None
    line_items:        list              = field(default_factory=list)
    raw_mst:           dict             = field(default_factory=dict)
    raw_dtl_rows:      list             = field(default_factory=list)
    raw_vw_rows:       list             = field(default_factory=list)


@dataclass
class DocumentContent:
    text:        Optional[str]      = None
    images:      Optional[list]     = None
    source_path: str                = ""
    page_count:  int                = 0
    ocr_source:  bool               = False

    @property
    def is_image_based(self) -> bool:
        return bool(self.images)


# ── Connection helpers ─────────────────────────────────────────────────────────

def _conn_str(conn_data: Data) -> str:
    d      = conn_data.data or {}
    driver = d.get("driver", "ODBC Driver 18 for SQL Server")
    server = d.get("host", d.get("server", ""))
    port   = d.get("port", 1433)
    db     = d.get("database_name", d.get("database", ""))
    user   = d.get("username", "")
    pwd    = d.get("password", "")
    return (
        f"DRIVER={{{driver}}};SERVER={server},{port};"
        f"DATABASE={db};UID={user};PWD={pwd};TrustServerCertificate=yes;"
    )


def _is_transient(exc: Exception) -> bool:
    kw = ["connection reset", "timeout", "throttl", "resource limit",
          "broken pipe", "transport-level", "login failed"]
    return any(k in str(exc).lower() for k in kw)


def _is_deadlock_error(exc: Exception) -> bool:
    """True for SQL Server deadlock-victim errors.

    Triggered when concurrent workers run _cleanup_for_pr against
    different PRs but the DELETEs join through shared rows in
    ras_tracker / attachment_classification and SQL Server picks the
    transaction as a deadlock victim.

    Error 1205 / SQLSTATE 40001 — the SQL Server message asks us
    to rerun the transaction, which is exactly what the retry wrapper
    in _process_pr does on the cleanup path.
    """
    msg = str(exc)
    if "40001" in msg or "1205" in msg:
        return True
    return "deadlock" in msg.lower()


def _connect(cs: str):
    import pyodbc
    for attempt in range(_MAX_RETRIES + 1):
        try:
            return pyodbc.connect(cs, timeout=30)
        except Exception as exc:
            if not _is_transient(exc) or attempt == _MAX_RETRIES:
                raise
            time.sleep(_BASE_DELAY * (2 ** attempt) * (1 + random.random() * 0.2))
    raise RuntimeError("unreachable")


# ── Blob helpers ───────────────────────────────────────────────────────────────

def _get_blob_config_by_name(connector_name: str) -> dict:
    name = (connector_name or "").strip()
    if not name:
        raise ValueError("blob_connector_name is empty.")
    try:
        import asyncio
        import concurrent.futures as _cf
        from sqlalchemy import select

        async def _fetch():
            from agentcore.services.deps import get_db_service
            from agentcore.services.database.models.connector_catalogue.model import ConnectorCatalogue
            db_service = get_db_service()
            async with db_service.with_session() as session:
                stmt = (
                    select(ConnectorCatalogue)
                    .where(ConnectorCatalogue.name == name)
                    .where(ConnectorCatalogue.provider == "azure_blob")
                )
                result = await session.execute(stmt)
                row = result.scalars().first()
                if row is None:
                    raise ValueError(f"No azure_blob connector named {name!r}.")
                cfg         = row.provider_config or {}
                account_url = (cfg.get("account_url") or row.host or "").strip()
                container   = (cfg.get("container_name") or row.database_name or "").strip()
                from urllib.parse import urlparse
                if not urlparse(account_url).netloc:
                    raise ValueError(f"Invalid Storage Account URL: {account_url!r}")
                if not container:
                    raise ValueError(f"No container_name in connector {name!r}.")
                return {"account_url": account_url, "container_name": container}

        try:
            asyncio.get_running_loop()
            with _cf.ThreadPoolExecutor() as pool:
                return pool.submit(asyncio.run, _fetch()).result(timeout=10)
        except RuntimeError:
            return asyncio.run(_fetch())
    except Exception as exc:
        raise RuntimeError(f"Blob connector lookup failed: {exc}") from exc


def _download_blob(blob_path: str, blob_cfg: dict) -> bytes:
    from azure.identity import DefaultAzureCredential
    from azure.storage.blob import BlobServiceClient
    credential = DefaultAzureCredential(
        exclude_environment_credential=True,
        exclude_interactive_browser_credential=True,
    )
    client = BlobServiceClient(
        account_url=blob_cfg["account_url"], credential=credential
    )
    blob = client.get_blob_client(container=blob_cfg["container_name"], blob=blob_path)
    return blob.download_blob().readall()


# ── File type detection ────────────────────────────────────────────────────────

def _detect_file_type(filename: str) -> str:
    import os
    ext = os.path.splitext(filename.lower())[1]
    return _EXT_TO_TYPE.get(ext, "unknown")


# ── Content extraction for classification ─────────────────────────────────────

def _extract_for_classification(file_bytes: bytes, filename: str) -> tuple:
    """Returns (text_content: str, image_b64: str|None, metadata_str: str)."""
    import os, io
    ext      = os.path.splitext(filename.lower())[1]
    ftype    = _EXT_TO_TYPE.get(ext, "unknown")
    meta_str = ""

    if ftype == "excel":
        return _extract_excel_classify(file_bytes, filename, meta_str)
    if ftype == "csv":
        return _extract_csv_classify(file_bytes, meta_str)
    if ftype == "pdf":
        return _extract_pdf_classify(file_bytes, filename, meta_str)
    if ftype == "word":
        return _extract_word_classify(file_bytes, filename, meta_str)
    if ftype == "pptx":
        return _extract_pptx_classify(file_bytes, filename, meta_str)
    if ftype == "legacy_doc":
        return _extract_fitz_classify(file_bytes, filename, meta_str)
    if ftype == "image":
        return _extract_image_classify(file_bytes, filename, meta_str)
    if ftype == "text":
        return _extract_text_classify(file_bytes, meta_str)
    if ftype == "html":
        return _extract_html_classify(file_bytes, meta_str)
    if ftype == "msg":
        return _extract_msg_classify(file_bytes, filename, meta_str)
    return "[Unsupported file type]", None, meta_str


def _extract_excel_classify(file_bytes, filename, meta_str):
    import io, pandas as pd
    _MAX_SHEETS = 8
    buf  = io.BytesIO(file_bytes)
    parts: list[str] = []
    try:
        for engine in ("openpyxl", "xlrd"):
            try:
                buf.seek(0)
                xls = pd.ExcelFile(buf, engine=engine)
                break
            except Exception:
                continue
        else:
            return "[Excel could not be opened]", None, meta_str

        sheets = xls.sheet_names
        meta_str = f"- sheet_count: {len(sheets)}\n- sheet_names: {sheets}\n- multi_sheet: {len(sheets) > 1}\n"
        parts.append(f"## Workbook Structure\nTotal Sheets: {len(sheets)}\nSheet Names: {sheets}\n")
        for idx, sheet in enumerate(sheets):
            try:
                df = pd.read_excel(xls, sheet_name=sheet, nrows=40, header=None)
                df = df.iloc[:, :30].dropna(how="all").dropna(axis=1, how="all")
                non_empty = len(df)
                if idx < _MAX_SHEETS and non_empty > 0:
                    parts.append(f"### Sheet {idx+1}: '{sheet}'")
                    parts.append(f"Non-empty rows in sample: {non_empty}")
                    try:
                        hdrs = [f"col_{c}" for c in range(df.shape[1])]
                        parts.append(df.fillna("").astype(str).to_markdown(index=False, headers=hdrs))
                    except Exception:
                        parts.append(df.fillna("").astype(str).to_string(index=False))
                    parts.append("")
                else:
                    parts.append(f"### Sheet {idx+1}: '{sheet}' (summary only — {non_empty} non-empty rows)\n")
            except Exception:
                pass
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Excel extraction error: {exc}]", None, meta_str


def _extract_csv_classify(file_bytes, meta_str):
    import io, pandas as pd
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), nrows=40, header=None).iloc[:, :30]
        try:
            text = df.fillna("").astype(str).to_markdown(index=False)
        except Exception:
            text = df.fillna("").astype(str).to_string(index=False)
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[CSV extraction error: {exc}]", None, meta_str


def _extract_pdf_classify(file_bytes, filename, meta_str):
    import io
    _MAX_PAGES = 10
    try:
        import pdfplumber
        buf  = io.BytesIO(file_bytes)
        with pdfplumber.open(buf) as pdf:
            total = len(pdf.pages)
            meta_str = f"- total_pages: {total}\n- pages_processed: {min(total, _MAX_PAGES)}\n"
            if total == 0:
                return "[Empty or corrupt PDF — no pages found]", None, meta_str

            # smart selection: first N-2 pages + last 2 pages (captures pricing summaries at end)
            if total <= _MAX_PAGES:
                page_idxs = list(range(total))
            else:
                page_idxs = sorted(set(list(range(_MAX_PAGES - 2)) + list(range(total - 2, total))))

            text_parts: list[str] = []
            table_parts: list[str] = []
            for i in page_idxs:
                page = pdf.pages[i]
                t = (page.extract_text() or "").strip()
                if t:
                    text_parts.append(f"--- Page {i+1} ---\n{t}")
                # extract tables from this page
                for t_idx, table in enumerate(page.extract_tables() or []):
                    if not table:
                        continue
                    header = table[0]
                    rows   = table[1:]
                    tstr   = " | ".join(str(c) for c in header) + "\n"
                    tstr  += " | ".join("---" for _ in header) + "\n"
                    for row in rows[:20]:
                        tstr += " | ".join(str(c) for c in row) + "\n"
                    table_parts.append(f"Table {t_idx+1} (Page {i+1}):\n{tstr}")

            combined = "\n\n".join(text_parts)
            if table_parts:
                combined += "\n\n### Extracted Tables\n" + "\n\n".join(table_parts)

            if len(combined.strip()) < 50:
                return _extract_pdf_as_image_classify(file_bytes, filename, meta_str, total)
            return combined[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception:
        return _extract_pdf_fitz_classify(file_bytes, filename, meta_str)


def _extract_pdf_fitz_classify(file_bytes, filename, meta_str):
    import io, base64
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total = len(doc)
        meta_str = f"- total_pages: {total}\n"
        pages = list(range(min(total, 8)))
        texts = [doc[i].get_text("text").strip() for i in pages]
        combined = "\n\n".join(f"--- Page {i+1} ---\n{t}" for i, t in zip(pages, texts) if t)
        if len(combined.strip()) < 50:
            pix = doc[0].get_pixmap(dpi=150)
            b64 = base64.b64encode(pix.tobytes("png")).decode()
            doc.close()
            return "[Scanned PDF - content sent as image]", b64, meta_str
        doc.close()
        return combined[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[PDF extraction error: {exc}]", None, meta_str


def _extract_pdf_as_image_classify(file_bytes, filename, meta_str, total):
    import io, base64
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            img = pdf.pages[0].to_image(resolution=200)
            buf = io.BytesIO()
            img.original.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode()
        return "[Scanned PDF - content sent as image]", b64, meta_str
    except Exception:
        return _extract_pdf_fitz_classify(file_bytes, filename, meta_str)


def _extract_word_classify(file_bytes, filename, meta_str):
    import io, zipfile, base64
    try:
        from docx import Document
        from PIL import Image
        doc   = Document(io.BytesIO(file_bytes))
        parts: list[str] = []
        para_count = 0
        for para in doc.paragraphs:
            if para_count >= 200:
                break
            text = para.text.strip()
            if not text:
                continue
            if para.style and para.style.name.startswith("Heading"):
                lvl = para.style.name.replace("Heading ", "").strip()
                try:
                    lvl = int(lvl)
                except ValueError:
                    lvl = 1
                parts.append(f"{'#' * lvl} {text}")
            else:
                parts.append(text)
            para_count += 1
        for i, tbl in enumerate(doc.tables):
            rows = []
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(cells)
            if rows:
                hdr = rows[0]
                tstr  = " | ".join(hdr) + "\n"
                tstr += " | ".join("---" for _ in hdr) + "\n"
                for row in rows[1:20]:
                    tstr += " | ".join(row) + "\n"
                parts.append(f"\n### Table {i+1}\n{tstr}")
        text = "\n\n".join(parts)
        meta_str = f"- paragraphs: {len(doc.paragraphs)}\n- tables: {len(doc.tables)}\n"
        if len(text.strip()) < 50:
            # try to extract largest embedded image from word/media/
            try:
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                    imgs = sorted([n for n in z.namelist() if n.startswith("word/media/")])
                    if imgs:
                        largest = max(imgs, key=lambda n: z.getinfo(n).file_size)
                        img_bytes = z.read(largest)
                        img = Image.open(io.BytesIO(img_bytes))
                        if img.mode not in ("RGB", "L"):
                            img = img.convert("RGB")
                        if max(img.size) > 2048:
                            img.thumbnail((2048, 2048), Image.LANCZOS)
                        buf2 = io.BytesIO()
                        img.save(buf2, format="PNG")
                        b64 = base64.b64encode(buf2.getvalue()).decode()
                        return "[Word document is image-based - content sent as image for visual analysis]", b64, meta_str
            except Exception:
                pass
            return _extract_fitz_classify(file_bytes, filename, meta_str)
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Word extraction error: {exc}]", None, meta_str


def _extract_pptx_classify(file_bytes, filename, meta_str):
    import io
    _MAX_SLIDES = 15
    try:
        from pptx import Presentation
        buf = io.BytesIO(file_bytes)
        prs = Presentation(buf)
        slides = list(prs.slides)
        total  = len(slides)
        parts  = [f"## Presentation: {total} slides\n"]
        for i, slide in enumerate(slides[:_MAX_SLIDES]):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_parts.append(t)
                if shape.has_table:
                    tbl = shape.table
                    rows = [[c.text.strip() for c in row.cells] for row in tbl.rows]
                    if rows:
                        hdr   = rows[0]
                        tstr  = " | ".join(hdr) + "\n"
                        tstr += " | ".join("---" for _ in hdr) + "\n"
                        for row in rows[1:20]:
                            tstr += " | ".join(row) + "\n"
                        slide_parts.append(tstr)
            if slide_parts:
                parts.append(f"### Slide {i+1}")
                parts.append("\n".join(slide_parts))
                parts.append("")
        meta_str = f"- total_slides: {total}\n- slides_extracted: {min(total, _MAX_SLIDES)}\n"
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return _extract_fitz_classify(file_bytes, filename, meta_str)


def _extract_fitz_classify(file_bytes, filename, meta_str):
    import base64
    try:
        import fitz
        doc = fitz.open(stream=file_bytes)
        total = len(doc)
        pages = list(range(min(total, 4)))
        images: list[str] = []
        for i in pages:
            pix = doc[i].get_pixmap(dpi=150)
            images.append(base64.b64encode(pix.tobytes("png")).decode())
        doc.close()
        if not images:
            return "[Document could not be rendered]", None, meta_str
        return "[Legacy document - content sent as images]", images[0], meta_str
    except Exception as exc:
        return f"[Document extraction error: {exc}]", None, meta_str


def _extract_image_classify(file_bytes, filename, meta_str):
    import io, base64
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        if max(img.size) > 2048:
            img.thumbnail((2048, 2048))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode()
        meta_str = f"- size: {img.size[0]}x{img.size[1]}\n"
        return "[Image file - content sent as image]", b64, meta_str
    except Exception as exc:
        return f"[Image extraction error: {exc}]", None, meta_str


def _extract_text_classify(file_bytes, meta_str):
    try:
        import chardet
        enc = (chardet.detect(file_bytes).get("encoding") or "utf-8")
        try:
            text = file_bytes.decode(enc)
        except Exception:
            text = file_bytes.decode("utf-8", errors="replace")
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Text extraction error: {exc}]", None, meta_str


def _extract_html_classify(file_bytes, meta_str):
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(file_bytes, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
    except Exception:
        raw  = file_bytes.decode("utf-8", errors="replace")
        text = re.sub(r"<[^>]+>", " ", raw)
        text = re.sub(r"\s+", " ", text).strip()
    return text[:_MAX_CLASSIFY_CHARS], None, meta_str


def _extract_msg_classify(file_bytes, filename, meta_str):
    import os, tempfile
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject: parts.append(f"Subject: {msg.subject}")
            if msg.sender:  parts.append(f"From: {msg.sender}")
            if msg.body:    parts.append(msg.body)
            msg.close()
        finally:
            os.unlink(tmp_path)
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[MSG extraction error: {exc}]", None, meta_str


# ── Trivial image heuristic ────────────────────────────────────────────────────

def _is_trivial_image(file_bytes: bytes, filename: str) -> tuple:
    import os
    if os.path.splitext(filename.lower())[1] not in _IMAGE_EXTS:
        return False, ""
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        img.load()
        w, h = img.size
        long_e  = max(w, h)
        short_e = max(min(w, h), 1)
        if long_e < 200:           return True, f"tiny {w}x{h}"
        if short_e < 150:          return True, f"short edge {w}x{h}"
        if long_e / short_e > 5:   return True, f"extreme aspect {w}x{h}"
        if long_e / short_e > 3 and short_e < 250:
            return True, f"banner {w}x{h}"
        if long_e < 400 and short_e < 200:
            return True, f"sub-thumbnail {w}x{h}"
        gray = img.convert("L")
        hist = gray.histogram()
        total = sum(hist)
        if total and sum(hist[230:]) / total > 0.97:
            return True, "mostly white"
        try:
            colors = img.convert("RGB").getcolors(maxcolors=256)
            if colors is not None and len(colors) < 16:
                return True, f"flat palette ({len(colors)} colors)"
        except Exception:
            pass
    except Exception as exc:
        return True, f"unreadable ({exc.__class__.__name__})"
    return False, ""


# ── Classification ─────────────────────────────────────────────────────────────

def _get_prompt_text(msg_input, default: str) -> str:
    """Extract text from an optional wired Message input; fall back to default."""
    if msg_input is None:
        return default
    text = getattr(msg_input, "text", None)
    if text and str(text).strip():
        return str(text).strip()
    return default


# ── LLM rate-limit retry helpers ──────────────────────────────────────────────

def _is_rate_limit_error(exc: Exception) -> bool:
    """Return True for 429 / token-quota / rate-limit errors from any LLM provider."""
    msg = str(exc).lower()
    return any(k in msg for k in (
        "429", "rate limit", "rate_limit", "ratelimit",
        "too many requests", "resource exhausted",
        "quota exceeded", "token quota", "tokens exceeded",
        "insufficient_quota", "token limit exceeded",
        "requests per minute", "tokens per minute",
    ))


def _is_json_format_error(exc: Exception) -> bool:
    """Detect when the LLM rejects the response_format=json_object kwarg.

    Older models / some Azure deployments don't support strict JSON mode and
    raise a 400-class BadRequestError mentioning response_format / json_object.
    """
    msg = str(exc).lower()
    return ("response_format" in msg) or ("json_object" in msg) or (
        "json mode" in msg
    )


def _call_llm_with_retry(
    llm,
    messages: list,
    prompts: dict | None = None,
    *,
    force_json: bool = False,
):
    """Invoke the LLM with automatic retry + cooldown on rate-limit / token-quota errors.

    Retry count and cooldown are read from the prompts dict so they can be set
    from the controller (llm_max_retries, llm_retry_cooldown keys).
    Cooldown is progressive: cooldown × (attempt+1) so each retry waits longer.

    When force_json=True, binds response_format={"type": "json_object"} on the
    LLM (mirrors doc-intel file_classifier/classifier/llm_client.py). If the
    model rejects the kwarg, falls back transparently to plain invoke so older
    deployments keep working.
    """
    import time
    max_retries = int((prompts or {}).get("llm_max_retries",  3))
    cooldown_s  = int((prompts or {}).get("llm_retry_cooldown", 60))

    # Try to bind strict JSON output. Most modern Azure / OpenAI chat models
    # support it; older ones reject the kwarg and we fall back below on first
    # invoke. .bind() itself is lazy so it won't fail here.
    active_llm = llm
    json_mode_enabled = False
    if force_json:
        try:
            active_llm = llm.bind(response_format={"type": "json_object"})
            json_mode_enabled = True
        except Exception:
            active_llm = llm  # bind unsupported → use plain LLM

    for attempt in range(max_retries + 1):
        try:
            return active_llm.invoke(messages)
        except Exception as exc:
            # If JSON mode tripped a model-not-supported error, drop it and
            # retry once with the unbound LLM. This is a recovery path, not a
            # rate-limit retry, so it doesn't consume a retry slot.
            if json_mode_enabled and _is_json_format_error(exc):
                logger.warning(
                    "LLM rejected response_format=json_object — "
                    "retrying without JSON mode: {}",
                    exc,
                )
                active_llm = llm
                json_mode_enabled = False
                continue
            if _is_rate_limit_error(exc) and attempt < max_retries:
                wait = cooldown_s * (attempt + 1)  # 60s → 120s → 180s
                logger.warning(
                    "LLM rate-limit / token-quota (attempt {}/{}) — "
                    "cooling down {}s before retry: {}",
                    attempt + 1, max_retries, wait, exc,
                )
                time.sleep(wait)
            else:
                raise


def _classify_file(llm, file_bytes: bytes, filename: str, prompts: dict | None = None) -> tuple:
    """Returns (doc_type: str, confidence: float)."""
    import os
    from langchain_core.messages import HumanMessage, SystemMessage

    p = prompts or {}
    sys_prompt      = p.get("cls_system", CLASSIFICATION_SYSTEM_PROMPT)
    user_text_tmpl  = p.get("cls_user_text", _CLASSIFY_USER_TEXT)
    user_image_tmpl = p.get("cls_user_image", _CLASSIFY_USER_IMAGE)

    ext = os.path.splitext(filename.lower())[1]
    if ext not in _SUPPORTED_CLASSIFY_EXTS:
        return "Others", 0.0

    trivial, reason = _is_trivial_image(file_bytes, filename)
    if trivial:
        logger.info(f"Trivial image {filename!r} ({reason}) → Others")
        return "Others", 0.0

    file_type = _EXT_TO_TYPE.get(ext, "unknown")
    text_content, image_b64, meta_str = _extract_for_classification(file_bytes, filename)

    cls_max = (prompts or {}).get("cls_max_chars")
    if cls_max and text_content and len(text_content) > cls_max:
        text_content = text_content[:cls_max]

    try:
        if image_b64:
            user_prompt = user_image_tmpl.format(
                filename=filename, file_type=file_type, extra_metadata=meta_str
            )
            content = [
                {"type": "text", "text": user_prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}", "detail": "high"}},
            ]
            messages = [SystemMessage(content=sys_prompt), HumanMessage(content=content)]
        else:
            user_prompt = user_text_tmpl.format(
                filename=filename, file_type=file_type,
                extra_metadata=meta_str, extracted_content=text_content,
            )
            messages = [SystemMessage(content=sys_prompt), HumanMessage(content=user_prompt)]

        # force_json=True mirrors doc-intel's response_format={"type":"json_object"}.
        # Without this the LLM may emit prose / markdown and json.loads() below
        # fails, dropping the file to "Others" — a major source of misclassification.
        response = _call_llm_with_retry(llm, messages, prompts, force_json=True)
        raw = (getattr(response, "content", None) or str(response)).strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        result = json.loads(raw)

        raw_cls = result.get("classification", "Other")
        if raw_cls not in VALID_CLASSIFICATIONS:
            raw_cls = "Other"
        confidence = float(result.get("confidence", 0.0))
        doc_type   = _CLASSIFICATION_MAP.get(raw_cls, raw_cls)
        return doc_type, confidence

    except json.JSONDecodeError as exc:
        # Surface JSON parse failures explicitly so the user can see in the
        # canvas log that the LLM output was unparseable (not just a generic
        # classification miss).
        logger.opt(exception=True).warning(
            "Classification JSON parse failed for {!r}: {} — raw response was: {!r}",
            filename, exc, (raw[:500] if 'raw' in locals() else '<no response>'),
        )
        return "Others", 0.0
    except Exception as exc:
        logger.opt(exception=True).warning(
            "Classification failed for {!r}: {}", filename, exc
        )
        return "Others", 0.0


# ── Stage 4: run classification for a PR ──────────────────────────────────────

def _run_classification(llm, tgt_cs: str, blob_cfg: dict, pr_no: str, prompts: dict | None = None) -> None:
    import os
    from concurrent.futures import ThreadPoolExecutor, as_completed

    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT ac.[attachment_classify_uuid_pk], ac.[file_path], ac.[attachment_id]
              FROM [ras_procurement].[attachment_classification] ac
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
        """, pr_no)
        parent_rows = cur.fetchall()

        cur.execute("""
            SELECT ec.[embedded_attachment_classification_id], ec.[file_path],
                   ac.[attachment_classify_uuid_pk]
              FROM [ras_procurement].[embedded_attachment_classification] ec
              JOIN [ras_procurement].[attachment_classification] ac
                ON ec.[attachment_classification_id] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
        """, pr_no)
        embedded_rows = cur.fetchall()
    finally:
        conn.close()

    # Each file: download blob + LLM call + single-row DB update — fully independent.
    # All three operations use separate connections / different rows, so concurrent
    # execution is safe. LangChain LLMs use httpx connection pooling (thread-safe).
    # ProcessPoolExecutor is NOT used because llm/embed_model are not picklable.

    def _classify_parent(row):
        att_pk, blob_path, att_id = str(row[0]), row[1], row[2]
        if not blob_path:
            return
        try:
            file_bytes     = _download_blob(blob_path, blob_cfg)
            filename       = os.path.basename(blob_path)
            doc_type, conf = _classify_file(llm, file_bytes, filename, prompts)
            _update_parent_classification(tgt_cs, att_id, doc_type, conf)
            logger.info(f"[{pr_no}] Parent {filename!r}: {doc_type} (conf={conf:.2f})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Parent classification failed ({blob_path!r}): {exc}")

    def _classify_embedded(row):
        emb_pk, blob_path, parent_pk = str(row[0]), row[1], str(row[2])
        if not blob_path:
            return
        try:
            file_bytes     = _download_blob(blob_path, blob_cfg)
            filename       = os.path.basename(blob_path)
            doc_type, conf = _classify_file(llm, file_bytes, filename, prompts)
            _update_embedded_classification(tgt_cs, parent_pk, blob_path, doc_type, conf)
            logger.info(f"[{pr_no}] Embedded {filename!r}: {doc_type} (conf={conf:.2f})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Embedded classification failed ({blob_path!r}): {exc}")

    tasks = (
        [(_classify_parent,   r) for r in parent_rows] +
        [(_classify_embedded, r) for r in embedded_rows]
    )
    if not tasks:
        return

    # Cap inner workers at 8 — each task is an LLM HTTP call (I/O-bound, not CPU).
    # With N outer PR workers, peak threads = N × 8; all blocked on network I/O.
    with ThreadPoolExecutor(max_workers=min(len(tasks), 8)) as inner_pool:
        futures = [inner_pool.submit(fn, row) for fn, row in tasks]
        for f in as_completed(futures):
            f.result()  # exceptions already caught inside fn; this just awaits all


def _update_parent_classification(tgt_cs: str, att_id: str, doc_type: str, conf: float) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[attachment_classification]
               SET [doc_type]            = ?,
                   [classification_conf] = ?,
                   [updated_at]          = SYSUTCDATETIME()
             WHERE [attachment_id] = ?
        """, doc_type, conf, att_id)
        conn.commit()
    finally:
        conn.close()


def _update_embedded_classification(tgt_cs: str, parent_pk: str, blob_path: str, doc_type: str, conf: float) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[embedded_attachment_classification]
               SET [doc_type]            = ?,
                   [classification_conf] = ?,
                   [updated_at]          = SYSUTCDATETIME()
             WHERE [attachment_classification_id] = ?
               AND [file_path] = ?
        """, doc_type, conf, parent_pk, blob_path)
        conn.commit()
    finally:
        conn.close()


# ── RAS context builder ────────────────────────────────────────────────────────

def _build_ras_context(tgt_cs: str, pr_no: str) -> Optional[RASContext]:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT TOP 1
                   prm.[PURCHASE_REQ_ID], prm.[SUPPLIER_NAME], prm.[JUSTIFICATION],
                   prm.[CURRENCY], prm.[ENQUIRY_NO], prm.[CLASSIFICATION],
                   prm.[Department], prm.[NEGOTIATED_BY], prm.[ADDRESS],
                   prm.[CONTRACT_NO], prm.[ORDER_NO], prm.[PURCHASE_VALUE]
              FROM [ras_procurement].[purchase_req_mst] prm
             WHERE prm.[PURCHASE_REQ_NO] = ?
        """, pr_no)
        mst = cur.fetchone()
        if not mst:
            return None

        req_id = mst[0]
        def _d(v): return Decimal(str(v)) if v is not None else None

        cur.execute("""
            SELECT prd.[PURCHASE_DTL_ID], prd.[PURCHASE_REQ_ID], prd.[ITEM_NO],
                   prd.[QUANTITY], prd.[ITEM_TYPE], prd.[ITEMDESCRIPTION],
                   prd.[PRICE], prd.[UOM], prd.[DISCOUNT], prd.[REQ_VALUE],
                   prd.[CURRENCY], prd.[DELIVERY_DATE], prd.[SUPPLIER_NAME],
                   prd.[PAYMENT_DETAILS], prd.[ORIGINAL_VALUE], prd.[Initial_Offer],
                   prd.[Negotiation], prd.[CommentsforItem], prd.[PREPAYMENT], prd.[ITEM_CODE]
              FROM [ras_procurement].[purchase_req_detail] prd
             WHERE prd.[PURCHASE_REQ_ID] = ?
             ORDER BY prd.[ITEM_NO]
        """, req_id)
        dtl_rows = cur.fetchall()

        line_items: list[LineItemContext] = []
        for r in dtl_rows:
            raw_del = r[11]
            del_str = (
                raw_del.date().isoformat() if hasattr(raw_del, "date")
                else str(raw_del) if raw_del else None
            )
            line_items.append(LineItemContext(
                purchase_dtl_id=r[0], purchase_req_id=r[1], item_no=r[2] or 0,
                quantity=_d(r[3]), item_type=r[4], item_description=r[5],
                unit_price=_d(r[6]), uom=r[7], discount=_d(r[8]),
                req_value=_d(r[9]), currency=r[10], delivery_date=del_str,
                supplier_name=r[12] or mst[1], payment_details=r[13],
                original_value=_d(r[14]), initial_offer=_d(r[15]),
                negotiation=_d(r[16]), comments=r[17], prepayment=r[18], item_code=r[19],
            ))

        # BI dashboard enrichment (best-effort)
        cat=sub_cat=site_c=site_r=div=site=req_t=l3=l4=pur_cat=ras_t=par_s=sup_t=sup_c=pay_d=po_d=cat_b=None
        try:
            cur.execute("""
                SELECT TOP 1 vw.[L1], vw.[Sub_Category_Type], vw.[Site_Country],
                             vw.[Site_Region], vw.[Division], vw.[L3], vw.[L4],
                             vw.[Purchase_Category], vw.[Title], vw.[Site],
                             vw.[Requisition_Type], vw.[Parent_Supplier], vw.[Supplier_Type],
                             vw.[Suplier_country], vw.[Payment_Days], vw.[PO_Date],
                             vw.[Category_Buyer], vw.[L2]
                  FROM vw_get_ras_data_for_bidashboard vw
                 WHERE vw.[PURCHASE_REQ_ID] = ?
            """, req_id)
            vw = cur.fetchone()
            if vw:
                cat,sub_cat,site_c,site_r,div,l3,l4 = vw[0],vw[1]or vw[17],vw[2],vw[3],vw[4],vw[5],vw[6]
                pur_cat,ras_t,site,req_t = vw[7],vw[8],vw[9],vw[10]
                par_s,sup_t,sup_c = vw[11],vw[12],vw[13]
                pay_d,po_d,cat_b  = vw[14],vw[15],vw[16]
        except Exception:
            pass

        return RASContext(
            purchase_req_no=pr_no, purchase_req_id=req_id,
            supplier_name=mst[1], justification=mst[2], currency=mst[3],
            enquiry_no=mst[4], classification=mst[5], department=mst[6],
            negotiated_by=mst[7], address=mst[8], contract_no=mst[9],
            order_no=mst[10], purchase_value=_d(mst[11]),
            category=cat, sub_category=sub_cat, site_country=site_c,
            site_region=site_r, site=site, division=div, requisition_type=req_t,
            parent_supplier=par_s, supplier_type=sup_t, supplier_country=sup_c,
            payment_days=pay_d, po_date=po_d, category_buyer=cat_b,
            l3=l3, l4=l4, purchase_category=pur_cat, ras_title=ras_t,
            line_items=line_items,
        )
    finally:
        conn.close()


# ── Document loader for extraction ─────────────────────────────────────────────

def _load_document(file_bytes: bytes, filename: str, max_pages: int = 20) -> DocumentContent:
    import os
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pdf":                                  return _load_pdf_for_extract(file_bytes, max_pages)
    if ext in (".xlsx", ".xls", ".xlsm", ".xlsb"):    return _load_spreadsheet_for_extract(file_bytes, ext)
    if ext == ".docx":                                 return _load_docx_for_extract(file_bytes)
    if ext == ".doc":                                  return _load_doc_legacy_for_extract(file_bytes, max_pages)
    if ext in (".pptx", ".ppt"):                       return _load_pptx_for_extract(file_bytes, max_pages)
    if ext in _IMAGE_EXTS:                             return _load_image_for_extract(file_bytes, ext)
    if ext in (".txt", ".csv", ".xml", ".json", ".eml"):
        return DocumentContent(text=file_bytes.decode("utf-8", errors="replace")[:50000], page_count=1)
    if ext in (".html", ".htm"):                       return _load_html_for_extract(file_bytes)
    if ext == ".msg":                                  return _load_msg_for_extract(file_bytes)
    # Fallback: fitz render (handles most office formats via PyMuPDF)
    return _fitz_render_for_extract(file_bytes, max_pages)


def _load_pdf_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    """Hybrid PDF loader (matches doc-intel branch non-OCR strategy).

    Digital page (≥50 extracted chars): captured as text + 72 DPI image for layout context.
    Scanned page (<50 chars): captured as 150 DPI image only.
    Returns both text and images so the extraction LLM has full context without Azure Doc Intel.
    """
    import base64
    _SCANNED_THRESHOLD = 50
    try:
        import fitz
        doc        = fitz.open(stream=file_bytes, filetype="pdf")
        page_count = len(doc)
        if page_count > max_pages:
            # Match doc-intel: warn explicitly so the user can see truncation in the canvas log.
            logger.warning(
                "PDF has {} page(s) but max_pages={} — only first {} page(s) will be processed",
                page_count, max_pages, max_pages,
            )
        total      = min(page_count, max_pages)
        text_pages: list[str] = []
        images:     list[str] = []
        text_count = scanned_count = 0

        for i in range(total):
            page      = doc[i]
            page_text = page.get_text("text").strip()
            if len(page_text) >= _SCANNED_THRESHOLD:
                # Digital page — extract text; low-res image preserves layout context
                text_pages.append(f"[Page {i + 1}]\n{page_text}")
                text_count   += 1
                pix = page.get_pixmap(dpi=72)
            else:
                # Scanned/image page — high-res render is the only usable content
                scanned_count += 1
                pix = page.get_pixmap(dpi=150)
            images.append(base64.b64encode(pix.tobytes("png")).decode())

        doc.close()
        logger.info(
            "PDF loaded: {}/{} page(s) — {} digital (72 DPI text+image), {} scanned (150 DPI image)",
            total, page_count, text_count, scanned_count,
        )
        text = "\n\n".join(text_pages) if text_pages else None
        return DocumentContent(text=text, images=images, page_count=total)
    except Exception as exc:
        logger.opt(exception=True).error("PDF load failed: {}", exc)
        return DocumentContent(text=f"[PDF error: {exc}]", page_count=0)


def _load_spreadsheet_for_extract(file_bytes: bytes, ext: str) -> DocumentContent:
    """Load Excel workbook as a markdown table string.

    .xlsx / .xlsm — openpyxl (handles macro-enabled workbooks too)
    .xls           — xlrd (legacy binary format)
    .xlsb          — openpyxl attempt; falls back to fitz image render
    """
    import io
    parts: list[str] = []

    def _openpyxl_parts(raw: bytes) -> list[str]:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        result: list[str] = []
        for ws in wb.worksheets:
            lines = [f"### Sheet: {ws.title}"]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    lines.append("| " + " | ".join(cells) + " |")
            if len(lines) > 1:
                result.append("\n".join(lines))
        wb.close()
        return result

    try:
        if ext in (".xlsx", ".xlsm"):
            parts = _openpyxl_parts(file_bytes)
        elif ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(file_contents=file_bytes)
            for sheet in wb.sheets():
                lines = [f"### Sheet: {sheet.name}"]
                for rx in range(sheet.nrows):
                    cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
                    if any(c.strip() for c in cells):
                        lines.append("| " + " | ".join(cells) + " |")
                if len(lines) > 1:
                    parts.append("\n".join(lines))
        else:  # .xlsb — binary Excel; openpyxl may handle it, else render as images
            try:
                parts = _openpyxl_parts(file_bytes)
            except Exception:
                return _fitz_render_for_extract(file_bytes, 20)

        text = "\n\n".join(parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, 20)
    except Exception as exc:
        return DocumentContent(text=f"[Spreadsheet error: {exc}]", page_count=0)


def _load_docx_for_extract(file_bytes: bytes) -> DocumentContent:
    import io
    try:
        from docx import Document
        doc   = Document(io.BytesIO(file_bytes))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        table_parts: list[str] = []
        for tbl in doc.tables:
            rows = ["| " + " | ".join(c.text.strip() for c in row.cells) + " |"
                    for row in tbl.rows]
            if rows:
                table_parts.append("\n".join(rows))
        text = "\n".join(paras)
        if table_parts:
            text += "\n\n" + "\n\n".join(table_parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, 20)
    except Exception as exc:
        return DocumentContent(text=f"[DOCX error: {exc}]", page_count=0)


def _load_doc_legacy_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    try:
        return _fitz_render_for_extract(file_bytes, max_pages)
    except Exception:
        pass
    try:
        import io, olefile
        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        if ole.exists("WordDocument"):
            import re as _re
            stream = ole.openstream("WordDocument").read()
            text = stream.decode("utf-8", errors="ignore")
            text = _re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
            if text.strip():
                return DocumentContent(text=text, page_count=1)
        ole.close()
    except Exception:
        pass
    return DocumentContent(text="[Document could not be read]")


def _load_pptx_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    import io
    try:
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(file_bytes))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_t: list[str] = [f"--- Slide {idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_t.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        slide_t.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
            if len(slide_t) > 1:
                parts.append("\n".join(slide_t))
        text = "\n\n".join(parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, max_pages)
    except Exception:
        return _fitz_render_for_extract(file_bytes, max_pages)


def _load_image_for_extract(file_bytes: bytes, ext: str) -> DocumentContent:
    import io, base64
    raw = file_bytes
    if ext in (".tif", ".tiff"):
        try:
            from PIL import Image
            buf = io.BytesIO()
            Image.open(io.BytesIO(raw)).save(buf, format="PNG")
            raw = buf.getvalue()
        except Exception:
            pass
    return DocumentContent(images=[base64.b64encode(raw).decode()], page_count=1)


def _load_msg_for_extract(file_bytes: bytes) -> DocumentContent:
    import os, tempfile
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject: parts.append(f"Subject: {msg.subject}")
            if msg.body:    parts.append(msg.body)
            msg.close()
        finally:
            os.unlink(tmp_path)
        return DocumentContent(text="\n\n".join(parts), page_count=1)
    except Exception as exc:
        return DocumentContent(text=f"[MSG error: {exc}]")


def _load_html_for_extract(file_bytes: bytes) -> DocumentContent:
    """Strip HTML tags and return plain text (BeautifulSoup → regex fallback)."""
    import re
    html_str = file_bytes.decode("utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_str, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
    except Exception:
        text = re.sub(r"<(script|style)[^>]*>.*?</(script|style)>", " ", html_str, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
    return DocumentContent(text=text[:50000], page_count=1)


def _fitz_render_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    import base64
    try:
        import fitz
        doc        = fitz.open(stream=file_bytes)
        page_count = len(doc)
        if page_count > max_pages:
            logger.warning(
                "Document has {} page(s) but max_pages={} — only first {} page(s) will be rendered",
                page_count, max_pages, max_pages,
            )
        n = min(page_count, max_pages)
        images = []
        for i in range(n):
            pix = doc[i].get_pixmap(dpi=200)
            images.append(base64.b64encode(pix.tobytes("png")).decode())
        doc.close()
        logger.info("Document rendered: {}/{} page(s) at 200 DPI", n, page_count)
        return DocumentContent(images=images, page_count=n)
    except Exception as exc:
        logger.opt(exception=True).error("Document render failed: {}", exc)
        return DocumentContent(text=f"[Render error: {exc}]")


# ── Quotation source resolver ──────────────────────────────────────────────────

def _resolve_quotation_sources(tgt_cs: str, pr_no: str) -> list:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    sources: list[dict] = []
    try:
        cur.execute("""
            SELECT ac.[attachment_classify_uuid_pk], ac.[file_path], ac.[attachment_id]
              FROM [ras_procurement].[attachment_classification] ac
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND ac.[doc_type] = 'Quotation'
        """, pr_no)
        for row in cur.fetchall():
            if row[1]:
                sources.append({
                    "blob_path": row[1],
                    "attachment_classify_fk": str(row[0]),
                    "embedded_classify_fk": None,
                    "attachment_id": str(row[2]),
                })

        cur.execute("""
            SELECT ec.[embedded_attachment_classification_id], ec.[file_path],
                   ec.[parent_attachment_id], ac.[attachment_classify_uuid_pk]
              FROM [ras_procurement].[embedded_attachment_classification] ec
              JOIN [ras_procurement].[attachment_classification] ac
                ON ec.[attachment_classification_id] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND ec.[doc_type] = 'Quotation'
        """, pr_no)
        for row in cur.fetchall():
            if row[1]:
                sources.append({
                    "blob_path": row[1],
                    "attachment_classify_fk": str(row[3]),
                    "embedded_classify_fk": str(row[0]),
                    "attachment_id": str(row[2]),
                })
    finally:
        conn.close()
    return sources


# ── Line items table builder ───────────────────────────────────────────────────

def _build_line_items_table(ctx: RASContext) -> str:
    def _na(v): return str(v) if v is not None else "N/A"
    header = ("| DTL_ID | Item No | Item Code | Description | Qty | UOM | Type "
              "| Unit Price | Req Value | Currency | Supplier | Delivery Date |")
    sep    = ("|--------|---------|-----------|-------------|-----|-----|------"
              "|------------|-----------|----------|----------|---------------|")
    rows   = [header, sep]
    for li in ctx.line_items:
        rows.append(
            f"| {_na(li.purchase_dtl_id)} | {_na(li.item_no)} | {_na(li.item_code)} "
            f"| {_na(li.item_description)} | {_na(li.quantity)} | {_na(li.uom)} "
            f"| {_na(li.item_type)} | {_na(li.unit_price)} | {_na(li.req_value)} "
            f"| {_na(li.currency)} | {_na(li.supplier_name)} | {_na(li.delivery_date)} |"
        )
    return "\n".join(rows)


def _build_raw_context(ctx: RASContext) -> str:
    parts: list[str] = []
    if ctx.raw_mst:
        lines = [f"  {k}: {v}" for k, v in ctx.raw_mst.items()]
        parts.append("#### purchase_req_mst (header)\n" + "\n".join(lines))
    if ctx.raw_dtl_rows:
        cols   = list(ctx.raw_dtl_rows[0].keys())
        header = "| " + " | ".join(cols) + " |"
        sep    = "|" + "|".join("---" for _ in cols) + "|"
        rows   = ["| " + " | ".join(str(r.get(c, "")) for c in cols) + " |"
                  for r in ctx.raw_dtl_rows]
        parts.append("#### purchase_req_detail\n" + "\n".join([header, sep] + rows))
    return "\n\n".join(parts)


# ── Extraction LLM call ────────────────────────────────────────────────────────

def _build_extraction_user_prompt(ctx: RASContext, doc: DocumentContent, prompts: dict | None = None) -> str:
    def _f(v): return str(v) if v is not None else "N/A"
    if doc.ocr_source and doc.text:
        doc_content_str = f"[OCR markdown from Azure Document Intelligence]\n\n{doc.text}"
    elif doc.images and doc.text:
        doc_content_str = f"[Extracted text — page images attached below]\n\n{doc.text}"
    elif doc.images:
        doc_content_str = "[Scanned document — page image(s) attached]"
    else:
        doc_content_str = doc.text or "[No content extracted]"

    ext_max = (prompts or {}).get("ext_max_chars")
    if ext_max and doc_content_str and len(doc_content_str) > ext_max:
        doc_content_str = doc_content_str[:ext_max]

    user_tmpl = (prompts or {}).get("ext_user", EXTRACTION_USER_TEMPLATE)
    return user_tmpl.format(
        purchase_req_no=_f(ctx.purchase_req_no),
        purchase_req_id=_f(ctx.purchase_req_id),
        justification=_f(ctx.justification),
        supplier_name=_f(ctx.supplier_name),
        currency=_f(ctx.currency),
        enquiry_no=_f(ctx.enquiry_no),
        classification=_f(ctx.classification),
        department=_f(ctx.department),
        negotiated_by=_f(ctx.negotiated_by),
        address=_f(ctx.address),
        contract_no=_f(ctx.contract_no),
        order_no=_f(ctx.order_no),
        purchase_value=_f(ctx.purchase_value),
        category=_f(ctx.category),
        sub_category=_f(ctx.sub_category),
        l3=_f(ctx.l3),
        l4=_f(ctx.l4),
        purchase_category=_f(ctx.purchase_category),
        ras_title=_f(ctx.ras_title),
        site_region=_f(ctx.site_region),
        site_country=_f(ctx.site_country),
        site=_f(ctx.site),
        division=_f(ctx.division),
        requisition_type=_f(ctx.requisition_type),
        parent_supplier=_f(ctx.parent_supplier),
        supplier_type=_f(ctx.supplier_type),
        supplier_country=_f(ctx.supplier_country),
        payment_days=_f(ctx.payment_days),
        po_date=_f(ctx.po_date),
        category_buyer=_f(ctx.category_buyer),
        line_items_table=_build_line_items_table(ctx),
        item_taxonomy=(prompts or {}).get("ext_taxonomy", ITEM_TAXONOMY),
        document_content=doc_content_str,
        raw_ras_context=_build_raw_context(ctx),
    )


def _call_extraction_llm(llm, ctx: RASContext, doc: DocumentContent, prompts: dict | None = None) -> str:
    from langchain_core.messages import HumanMessage, SystemMessage
    user_prompt = _build_extraction_user_prompt(ctx, doc, prompts)
    sys_prompt  = (prompts or {}).get("ext_system", EXTRACTION_SYSTEM_PROMPT)
    # Hard cap on images attached per LLM request — Azure OpenAI rejects
    # requests with more than ~50 images. Mirrors doc-intel _MAX_IMAGES.
    max_images = max(1, int((prompts or {}).get("ext_max_images", 50)))
    messages: list = [SystemMessage(content=sys_prompt)]
    img_count  = 0
    img_detail = "n/a"
    if doc.is_image_based and doc.images:
        total_imgs = len(doc.images)
        if total_imgs > max_images:
            logger.warning(
                "[PR={}] Quotation has {} image(s) but max_images={} — truncating to first {}",
                ctx.purchase_req_no, total_imgs, max_images, max_images,
            )
        images = doc.images[:max_images]
        img_count  = len(images)
        # Mirror doc-intel: low detail when text is also present, high otherwise.
        img_detail = "low" if doc.text else "high"
        content_parts: list = [{"type": "text", "text": user_prompt}]
        for b64 in images:
            content_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}", "detail": img_detail},
            })
        messages.append(HumanMessage(content=content_parts))
    else:
        messages.append(HumanMessage(content=user_prompt))
    logger.info(
        "[PR={}] Calling extraction LLM — {} image(s) detail={}, ~{:.0f} kB prompt",
        ctx.purchase_req_no, img_count, img_detail, len(user_prompt) / 1024,
    )
    # force_json=True so the extraction LLM emits a strict JSON object and
    # our _parse_extraction_response doesn't have to clean up prose/fences.
    response = _call_llm_with_retry(llm, messages, prompts, force_json=True)
    return (getattr(response, "content", None) or str(response)).strip()


# ── Supplier matching ──────────────────────────────────────────────────────────

_SELECTED_THRESHOLD     = Decimal("0.70")  # min overall score to mark as selected quote
_PRICE_MAX_BOOST        = Decimal("0.10")  # max confidence boost from price alignment
_CANONICALIZE_THRESHOLD = 0.82             # SequenceMatcher ratio to cluster two supplier names

def _compute_supplier_match(supplier: Optional[str], ctx: RASContext, thresholds: dict | None = None) -> tuple:
    if not supplier:
        return False, Decimal("0")
    known: set[str] = set()
    if ctx.supplier_name:   known.add(ctx.supplier_name.strip())
    if ctx.parent_supplier: known.add(ctx.parent_supplier.strip())
    for li in ctx.line_items:
        if li.supplier_name: known.add(li.supplier_name.strip())
    if not known:
        return False, Decimal("0")
    ext  = supplier.strip().lower()
    best = 0.0
    for n in known:
        nl = n.lower()
        if ext in nl or nl in ext:
            best = max(best, 0.90)
        else:
            best = max(best, SequenceMatcher(None, ext, nl).ratio())
    conf = Decimal(str(round(best, 4)))
    threshold = (thresholds or {}).get("selected_threshold", _SELECTED_THRESHOLD)
    return conf >= threshold, conf


# ── Extraction response parsing ────────────────────────────────────────────────

def _parse_extraction_response(raw: str, source: dict, ctx: RASContext) -> list:
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        logger.error(f"LLM returned invalid JSON: {exc}")
        return []

    header: dict = data if isinstance(data, dict) else {}
    items_raw: list = header.get("items", [])
    if not items_raw:
        return []

    h_supplier = header.get("supplier_name") or None
    _, match_conf = _compute_supplier_match(h_supplier, ctx)

    header_fields = {
        "supplier_name":      h_supplier,
        "supplier_address":   header.get("supplier_address") or None,
        # Store the LLM-extracted country verbatim (matches doc-intel
        # _parse_llm_response). Previously this was force-normalised to ISO
        # alpha-2 via pycountry which lost the readable country name in the
        # DB (e.g. "India" → "IN"). Downstream consumers (_compute_cpi_pct,
        # _estimate_inflation_via_llm) already handle both full names and
        # alpha-2/3 codes, so no functional dependency was lost.
        "supplier_country":   (header.get("supplier_country") or "").strip() or None,
        "quotation_ref_no":   header.get("quotation_ref_no") or None,
        "quotation_date":     header.get("quotation_date"),
        "currency":           _normalize_currency_code(header.get("currency") or None),
        "validity_date":      header.get("validity_date"),
        "validity_days":      header.get("validity_days"),
        "payment_terms":      header.get("payment_terms") or None,
        "supplier_match_conf": float(match_conf),
        "attachment_classify_fk": source["attachment_classify_fk"],
        "embedded_classify_fk":   source["embedded_classify_fk"],
        "is_selected_quote":      0,
        "quote_rank":             None,
    }

    results: list[dict] = []
    for raw_item in items_raw:
        item = dict(header_fields)
        item.update({k: (None if v == "" else v) for k, v in raw_item.items()})
        results.append(item)
    return results


# ── Item alignment to RAS line items ──────────────────────────────────────────

_IDENT_RE = re.compile(r"[A-Za-z0-9]*\d+[A-Za-z0-9]*")

def _ident_tokens(text: Optional[str]) -> set:
    if not text: return set()
    return {m.group(0).lower() for m in _IDENT_RE.finditer(text) if len(m.group(0)) >= 2}

def _data_score(item: dict) -> tuple:
    return (int(item.get("unit_price") is not None), int(item.get("total_price") is not None),
            int(item.get("item_name") is not None), int(item.get("item_description") is not None))

def _align_to_ras_line_items(items: list, ctx: RASContext, source: dict) -> list:
    valid_ids = {li.purchase_dtl_id for li in ctx.line_items}
    matched_by_dtl: dict = {}
    orphans: list[dict] = []

    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None:
            try: dtl_id = int(dtl_id)
            except Exception: dtl_id = None
        if dtl_id in valid_ids:
            if dtl_id not in matched_by_dtl:
                matched_by_dtl[dtl_id] = item
            elif _data_score(item) > _data_score(matched_by_dtl[dtl_id]):
                matched_by_dtl[dtl_id] = item
        elif any([item.get("item_name"), item.get("unit_price"), item.get("total_price")]):
            orphans.append(item)

    matched = list(matched_by_dtl.values())
    covered = {i.get("purchase_dtl_id") for i in matched}
    uncovered = [li for li in ctx.line_items if li.purchase_dtl_id not in covered]

    # Fuzzy assign orphans
    if orphans and uncovered:
        scored: list = []
        for oi, orp in enumerate(orphans):
            orp_text = f"{orp.get('item_name','')} {orp.get('item_description','')}".lower()
            for li_idx, li in enumerate(uncovered):
                ras_text = f"{li.item_description or ''} {li.item_code or ''}".lower()
                shared = _ident_tokens(orp_text) & _ident_tokens(ras_text)
                if shared:
                    score = 0.90
                elif li.item_code and li.item_code.lower() in orp_text:
                    score = 0.85
                else:
                    score = SequenceMatcher(None, orp_text, ras_text).ratio()
                try:
                    if orp.get("quantity") is not None and li.quantity is not None:
                        if Decimal(str(orp["quantity"])) == li.quantity:
                            score = min(1.0, score + 0.10)
                except Exception:
                    pass
                scored.append((score, oi, li_idx))
        scored.sort(key=lambda x: x[0], reverse=True)
        used_o: set = set(); used_l: set = set()
        for score, oi, li_idx in scored:
            if oi in used_o or li_idx in used_l: continue
            dtl_id = uncovered[li_idx].purchase_dtl_id
            orphans[oi]["purchase_dtl_id"] = dtl_id
            matched.append(orphans[oi])
            covered.add(dtl_id)
            used_o.add(oi); used_l.add(li_idx)

    # Stubs for uncovered
    still_uncovered = {li.purchase_dtl_id for li in uncovered if li.purchase_dtl_id not in covered}
    donor = matched[0] if matched else {}
    for li in ctx.line_items:
        if li.purchase_dtl_id not in still_uncovered:
            continue
        matched.append({
            "attachment_classify_fk": source["attachment_classify_fk"],
            "embedded_classify_fk":   source["embedded_classify_fk"],
            "purchase_dtl_id":        li.purchase_dtl_id,
            "is_selected_quote":      0,
            "supplier_match_conf":    0.0,
            "quote_rank":             None,
            "supplier_name":          donor.get("supplier_name"),
            "supplier_address":       donor.get("supplier_address"),
            "supplier_country":       donor.get("supplier_country"),
            "quotation_ref_no":       donor.get("quotation_ref_no"),
            "quotation_date":         donor.get("quotation_date"),
            "currency":               donor.get("currency"),
            "validity_date":          donor.get("validity_date"),
            "validity_days":          donor.get("validity_days"),
            "payment_terms":          donor.get("payment_terms"),
        })
    return matched


# ── Quote ranking + selection ──────────────────────────────────────────────────

def _compute_quote_ranks(all_items: list) -> None:
    from collections import defaultdict
    for item in all_items:
        item["quote_rank"] = None
    by_group: dict = defaultdict(list)
    for item in all_items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is None: continue
        key = ((item.get("supplier_name") or "").strip().lower() or "_unknown_", dtl_id)
        by_group[key].append(item)
    for group in by_group.values():
        group.sort(key=lambda i: (i.get("total_price") is None, float(i.get("total_price") or 0)))
        for rank, item in enumerate(group, 1):
            item["quote_rank"] = rank


def _select_best_quotes(all_items: list, ctx: RASContext) -> None:
    from collections import defaultdict
    for item in all_items:
        item["is_selected_quote"] = 0
    ras_by_dtl = {li.purchase_dtl_id: li for li in ctx.line_items}
    by_dtl: dict = defaultdict(list)
    for item in all_items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None:
            by_dtl[dtl_id].append(item)

    def _price_prox(item_val, ras_val):
        try:
            if item_val is None or ras_val is None or float(ras_val) <= 0: return 0.0
            diff = abs(float(item_val) - float(ras_val)) / float(ras_val)
            if diff <= 0.05: return 1.0
            if diff >= 0.25: return 0.0
            return 1.0 - (diff - 0.05) / 0.20
        except Exception:
            return 0.0

    for dtl_id, candidates in by_dtl.items():
        ras_line = ras_by_dtl.get(dtl_id)
        def _score(it):
            conf      = float(it.get("supplier_match_conf") or 0)
            price_fit = max(
                _price_prox(it.get("unit_price"),  ras_line.unit_price if ras_line else None),
                _price_prox(it.get("total_price"), ras_line.req_value  if ras_line else None),
            )
            has_price = int(it.get("unit_price") is not None)
            return (conf, price_fit, has_price)
        max(candidates, key=_score)["is_selected_quote"] = 1


# ── Currency conversion helper ─────────────────────────────────────────────────

_EUR_CUR_ID      = 3
_RATE_IS_MULTIPLY = True


def _convert_to_eur(tgt_cs: str, amount, currency_code: str | None, ref_date) -> "Decimal | None":
    """Return amount converted to EUR using EXCHANGE_RATE table.

    Returns None if amount is None or conversion data is unavailable.
    """
    if amount is None:
        return None
    try:
        from decimal import Decimal as _Dec
        amount_dec = _Dec(str(amount))
    except Exception:
        return None
    if not currency_code:
        return None
    try:
        import re as _re
        from datetime import date as _date_cls, datetime as _dt_cls
        if isinstance(ref_date, _date_cls):
            date_val = ref_date
        elif ref_date:
            m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(ref_date))
            if m:
                date_val = _date_cls(int(m[1]), int(m[2]), int(m[3]))
            else:
                date_val = _dt_cls.utcnow().date()
        else:
            from datetime import datetime as _dt_cls2
            date_val = _dt_cls2.utcnow().date()
        conn = _connect(tgt_cs)
        try:
            cur = conn.cursor()
            # look up source currency id
            cur.execute(
                "SELECT [CUR_ID] FROM [ras_procurement].[currency_mst] "
                "WHERE UPPER([CURRENCY]) = UPPER(?)",
                currency_code,
            )
            row = cur.fetchone()
            if row is None:
                logger.warning(f"Currency conversion: no currency_mst row for code={currency_code!r}")
                return None
            src_cur_id = row[0]
            if src_cur_id == _EUR_CUR_ID:
                return amount_dec
            # look up exchange rate valid for ref_date
            cur.execute(
                "SELECT TOP 1 [CONVERSION_RATE] "
                "FROM [ras_procurement].[EXCHANGE_RATE] "
                "WHERE [CUR_ID] = ? AND [BASE_CUR_ID] = ? "
                "  AND [STATUS_ID] = 10 "
                "  AND [FROM_DATE] <= ? AND [TO_DATE] >= ? "
                "ORDER BY [FROM_DATE] DESC",
                src_cur_id, _EUR_CUR_ID, date_val, date_val,
            )
            rate_row = cur.fetchone()
            if rate_row is None:
                logger.warning(
                    f"Currency conversion: no EXCHANGE_RATE row for "
                    f"FROM_CUR_ID={src_cur_id} TO_CUR_ID={_EUR_CUR_ID} date={date_val} STATUS_ID=10"
                )
                return None
            from decimal import Decimal as _Dec2
            rate = _Dec2(str(rate_row[0]))
            return amount_dec * rate if _RATE_IS_MULTIPLY else amount_dec / rate
        finally:
            conn.close()
    except Exception as exc:
        logger.warning(f"Currency conversion failed currency={currency_code!r} date={ref_date}: {exc}")
        return None


# ── Supplier name canonicalization (mirrors doc intel Union-Find logic) ────────

def _normalize_currency_code(code_or_name: str | None) -> str | None:
    """Normalize any currency string to ISO-4217 alpha-3 code using pycountry.

    Handles codes (USD), full names (US Dollar), and partial names via rapidfuzz.
    No hardcoded map — works for any currency in the ISO-4217 standard.
    Gracefully degrades to original value when pycountry is not installed.
    """
    if not code_or_name:
        return code_or_name
    val = code_or_name.strip()
    if not val:
        return None
    upper = val.upper()
    try:
        import pycountry as _pc
        # 1. Exact alpha-3 code match (most common — LLM usually returns ISO code)
        c = _pc.currencies.get(alpha_3=upper)
        if c:
            return c.alpha_3
        # 2. Exact name match (e.g. "US Dollar", "Euro")
        val_lower = val.lower()
        for c in _pc.currencies:
            if c.name.lower() == val_lower:
                return c.alpha_3
        # 3. Fuzzy name match via rapidfuzz (catches partial names, typos)
        try:
            from rapidfuzz import process as _fuzz
            all_names = [c.name for c in _pc.currencies]
            result = _fuzz.extractOne(val_lower, [n.lower() for n in all_names], score_cutoff=85)
            if result:
                idx = [n.lower() for n in all_names].index(result[0])
                return list(_pc.currencies)[idx].alpha_3
        except ImportError:
            pass
        # 4. Looks like a 3-char code but not in pycountry — return uppercased
        if len(upper) == 3 and upper.isalpha():
            logger.warning(f"Currency code {upper!r} not found in ISO-4217 — passing through as-is")
            return upper
    except ImportError:
        logger.warning("pycountry not installed — currency code not normalized (pip install pycountry)")
        if len(upper) == 3 and upper.isalpha():
            return upper
    return val


def _normalize_supplier_country(country_str: str | None) -> str | None:
    """Normalize free-text country name to ISO 3166-1 alpha-2 code using pycountry.

    Tries exact name/code match first, then rapidfuzz fuzzy match (threshold 75).
    Falls back to original string when pycountry is not installed or no match found.
    """
    if not country_str:
        return country_str
    text_lower = country_str.strip().lower()
    try:
        import pycountry as _pc
        for c in _pc.countries:
            if text_lower in (
                c.name.lower(),
                getattr(c, "official_name", "").lower(),
                c.alpha_2.lower(),
                c.alpha_3.lower(),
            ):
                return c.alpha_2
        try:
            from rapidfuzz import process as _fuzz
            all_names = [c.name for c in _pc.countries]
            result = _fuzz.extractOne(text_lower, [n.lower() for n in all_names])
            if result and result[1] > 75:
                matched = _pc.countries.get(name=all_names[[n.lower() for n in all_names].index(result[0])])
                if matched:
                    return matched.alpha_2
        except ImportError:
            pass
    except ImportError:
        logger.warning("pycountry not installed — supplier_country not normalized (pip install pycountry rapidfuzz)")
    return country_str


def _strip_contact_suffix(name: str) -> str:
    import re as _re
    return _re.sub(
        r'\s*[-–]\s*(contact|email|ph|phone|tel|mob)[:\s].*$',
        '', name, flags=_re.IGNORECASE,
    ).strip()


def _is_acronym_of(short: str, long_name: str) -> bool:
    """True if short is all-caps (≤ 6 chars) and matches the word initials of long_name."""
    if not short.isupper() or len(short) > 6:
        return False
    import re as _re
    initials = "".join(m[0].upper() for m in _re.findall(r'\b[A-Za-z]', long_name))
    return short == initials[:len(short)]


def _name_geo_tokens(cleaned: str) -> frozenset:
    import re as _re
    _GEO = frozenset({
        "india", "china", "japan", "usa", "us", "uk", "germany", "france",
        "italy", "korea", "taiwan", "singapore", "malaysia", "thailand",
        "vietnam", "indonesia", "australia", "canada", "brazil", "mexico",
        "uae", "dubai", "europe", "asia", "americas", "shanghai", "beijing",
        "mumbai", "delhi",
    })
    return frozenset(w for w in _re.findall(r'\b\w+\b', cleaned.lower()) if w in _GEO)


def _apply_price_alignment_boost(items: list[dict], ctx: RASContext, thresholds: dict | None = None) -> None:
    """Boost supplier_match_conf by up to price_max_boost when extracted prices align with RAS prices.

    If ≥50% of matched line items have a unit_price within _PRICE_TOLERANCE (5%) of the
    RAS line price, all items in this PR get a proportional boost. Matches doc-intel branch
    _apply_price_alignment_boost logic, adapted for dict items.
    """
    dtl_price: dict = {
        li.purchase_dtl_id: li.unit_price
        for li in ctx.line_items
        if li.unit_price is not None and li.unit_price > 0
    }
    if not dtl_price:
        return
    matches = comparable = 0
    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id in dtl_price:
            unit_price = item.get("unit_price")
            if unit_price is not None:
                comparable += 1
                try:
                    ras_p = dtl_price[dtl_id]
                    if abs(Decimal(str(unit_price)) - ras_p) / ras_p <= _PRICE_TOLERANCE:
                        matches += 1
                except Exception:
                    pass
    if comparable == 0 or matches / comparable < 0.5:
        return
    max_boost = (thresholds or {}).get("price_max_boost", _PRICE_MAX_BOOST)
    boost = Decimal(str(round((matches / comparable) * float(max_boost), 4)))
    for item in items:
        conf = item.get("supplier_match_conf")
        if conf is not None:
            try:
                item["supplier_match_conf"] = min(
                    Decimal("1.0"), Decimal(str(conf)) + boost
                )
            except Exception:
                pass
    logger.debug(f"Price alignment boost +{boost} applied ({matches}/{comparable} items within {int(_PRICE_TOLERANCE*100)}%)")


def _canonicalize_supplier_names(items: list[dict], ctx, thresholds: dict | None = None) -> None:
    """Per-DTL Union-Find supplier name canonicalization with supplier_country guard.

    Mirrors doc-intel ExtractionWriter.canonicalize_supplier_names_in_db
    (the production variant — not the in-memory geo-token variant):
      1. Group items by purchase_dtl_id.
      2. Within each DTL_ID, cluster supplier_name variants using four merge
         rules: exact, substring, SequenceMatcher ratio ≥ 0.82, acronym.
      3. Country guard: if both items have a non-empty supplier_country and
         those countries differ, they are NOT merged — different country
         branches stay separate suppliers.
      4. Canonical preference: RAS-known supplier_name first, else shortest.

    After clustering, supplier_match_conf is recomputed for ALL items against
    the RAS supplier set (matches doc-intel run.py step 6 — fresh scores so
    ranking uses canonical names; this also wipes the earlier
    _apply_price_alignment_boost, same as doc-intel).
    """
    from difflib import SequenceMatcher
    from collections import defaultdict

    # Build RAS-known name set for canonical preference (matches doc-intel)
    ras_known: set[str] = set()
    if getattr(ctx, "supplier_name", None):
        ras_known.add(ctx.supplier_name.strip().lower())
    if getattr(ctx, "parent_supplier", None):
        ras_known.add(ctx.parent_supplier.strip().lower())
    for li in getattr(ctx, "line_items", []):
        if getattr(li, "supplier_name", None):
            ras_known.add(li.supplier_name.strip().lower())

    threshold = (thresholds or {}).get("canonicalize_threshold", _CANONICALIZE_THRESHOLD)

    # Group items by purchase_dtl_id (None → not clustered, keeps raw name)
    by_dtl: dict[int, list[dict]] = defaultdict(list)
    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None:
            by_dtl[dtl_id].append(item)

    # (dtl_id, raw_name) → canonical_name
    canonical_per_dtl: dict[tuple[int, str], str] = {}

    for dtl_id, dtl_items in by_dtl.items():
        # Collect unique names + first non-empty country seen per name
        raw_names: list[str] = []
        name_country: dict[str, str] = {}
        for it in dtl_items:
            n = (it.get("supplier_name") or "").strip()
            if not n:
                continue
            if n not in name_country:
                raw_names.append(n)
                name_country[n] = (it.get("supplier_country") or "").strip().lower()

        if len(raw_names) < 2:
            continue

        uf: dict[str, str] = {n: n for n in raw_names}

        def _find(x: str) -> str:
            while uf[x] != x:
                uf[x] = uf[uf[x]]
                x = uf[x]
            return x

        def _union(a: str, b: str) -> None:
            uf[_find(a)] = _find(b)

        for i, a in enumerate(raw_names):
            a_clean   = _strip_contact_suffix(a).lower()
            a_country = name_country.get(a, "")
            a_geo     = _name_geo_tokens(a_clean)
            for b in raw_names[i + 1:]:
                b_clean   = _strip_contact_suffix(b).lower()
                b_country = name_country.get(b, "")
                b_geo     = _name_geo_tokens(b_clean)

                # Country-branch guard — combines BOTH doc-intel variants:
                #   1. supplier_country column (writer.py canonicalize_in_db)
                #   2. geo-token in the supplier_name itself
                #      (extractor.py canonicalize_supplier_names)
                # Either signal alone is enough to keep two names separate,
                # so "ACME Corp India" / "ACME Corp China" don't merge even
                # when the LLM forgot to extract supplier_country.
                if a_country and b_country and a_country != b_country:
                    continue
                if a_geo and b_geo and a_geo != b_geo:
                    continue

                if a_clean == b_clean:
                    _union(a, b)
                elif a_clean in b_clean or b_clean in a_clean:
                    _union(a, b)
                elif SequenceMatcher(None, a_clean, b_clean).ratio() >= threshold:
                    _union(a, b)
                elif _is_acronym_of(a.strip(), b) or _is_acronym_of(b.strip(), a):
                    _union(a, b)

        clusters: dict[str, list[str]] = defaultdict(list)
        for n in raw_names:
            clusters[_find(n)].append(n)

        for members in clusters.values():
            ras_match = next((m for m in members if m.lower() in ras_known), None)
            canonical = ras_match if ras_match else min(members, key=len)
            for m in members:
                canonical_per_dtl[(dtl_id, m)] = canonical

    # Apply canonical names + recompute supplier_match_conf for ALL items
    # (matches doc-intel run.py step 6 — fresh scores against RAS suppliers).
    changed = 0
    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        orig   = (item.get("supplier_name") or "").strip()
        if dtl_id is not None and orig:
            canon = canonical_per_dtl.get((dtl_id, orig), orig)
            if canon != orig:
                item["supplier_name"] = canon
                changed += 1
        # Recompute conf for every item against the (possibly canonical) name
        _, new_conf = _compute_supplier_match(item.get("supplier_name"), ctx, thresholds)
        item["supplier_match_conf"] = new_conf

    if changed:
        logger.info(f"Supplier canonicalization (per-DTL): {changed} name(s) updated")


# ── DB writer: extracted items ─────────────────────────────────────────────────

def _save_extracted_items(tgt_cs: str, items: list) -> int:
    if not items:
        return 0
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    saved = 0
    try:
        for item in items:
            def _v(k, cast=None):
                v = item.get(k)
                if v is None: return None
                try: return cast(v) if cast else v
                except Exception: return None
            def _d(k):
                v = item.get(k)
                if v is None: return None
                try: return Decimal(str(v))
                except Exception: return None
            def _date(k):
                v = item.get(k)
                if not v: return None
                try:
                    from datetime import date as date_cls
                    if isinstance(v, date_cls): return v
                    import re as _re
                    m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(v))
                    if m: return date_cls(int(m[1]), int(m[2]), int(m[3]))
                except Exception: pass
                return None
            unit_price_eur  = _convert_to_eur(tgt_cs, item.get("unit_price"),  item.get("currency"), item.get("quotation_date"))
            total_price_eur = _convert_to_eur(tgt_cs, item.get("total_price"), item.get("currency"), item.get("quotation_date"))
            cur.execute("""
                INSERT INTO [ras_procurement].[quotation_extracted_items] (
                    [attachment_classify_fk],[embedded_classify_fk],[purchase_dtl_id],
                    [is_selected_quote],[supplier_match_conf],[quote_rank],
                    [supplier_name],[supplier_address],[supplier_country],
                    [quotation_ref_no],[quotation_date],[currency],
                    [validity_date],[validity_days],[payment_terms],
                    [item_name],[item_description],[quantity],[unit],
                    [unit_price],[total_price],[discount],[taxation_details],
                    [delivery_date],[delivery_time_days],
                    [item_level_1],[item_level_2],[item_level_3],[item_level_4],
                    [item_level_5],[item_level_6],[item_level_7],[item_level_8],
                    [commodity_tag],[item_summary],
                    [unit_price_eur],[total_price_eur]
                ) VALUES (
                    ?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?
                )
            """,
                _v("attachment_classify_fk"), _v("embedded_classify_fk"),
                _v("purchase_dtl_id", int), int(item.get("is_selected_quote") or 0),
                _d("supplier_match_conf"), _v("quote_rank", int),
                _v("supplier_name"), _v("supplier_address"), _v("supplier_country"),
                _v("quotation_ref_no"), _date("quotation_date"), _v("currency"),
                _date("validity_date"), _v("validity_days", int), _v("payment_terms"),
                _v("item_name"), _v("item_description"), _d("quantity"), _v("unit"),
                _d("unit_price"), _d("total_price"), _d("discount"), _v("taxation_details"),
                _date("delivery_date"), _v("delivery_time_days", int),
                _v("item_level_1"), _v("item_level_2"), _v("item_level_3"), _v("item_level_4"),
                _v("item_level_5"), _v("item_level_6"), _v("item_level_7"), _v("item_level_8"),
                _v("commodity_tag"), _v("item_summary"),
                unit_price_eur, total_price_eur,
            )
            saved += 1
        conn.commit()
    except Exception as exc:
        logger.opt(exception=True).error(
            "DB INSERT into quotation_extracted_items failed after {} row(s): {}",
            saved, exc,
        )
        try:
            conn.rollback()
        except Exception as rb_exc:
            logger.warning("Rollback failed: {}", rb_exc)
        raise
    finally:
        conn.close()
    return saved


# ── Stage 5: run extraction for a PR ──────────────────────────────────────────

def _run_extraction(llm, tgt_cs: str, blob_cfg: dict, pr_no: str, prompts: dict | None = None) -> int:
    import os
    # Each abort path raises an ExtractionAbortError subclass so the per-PR
    # catch in _run_stages_48 records ras_pipeline_exceptions (stage_id=5),
    # sets ras_tracker.current_stage_fk = 99, and skips stages 6/7/8.
    # Returning 0 silently here would let the pipeline advance to stage 8
    # (COMPLETE) on a PR that had nothing to extract — which is wrong.
    ctx = _build_ras_context(tgt_cs, pr_no)
    if ctx is None:
        raise NoRASContextError(
            f"No purchase_req_mst row for PR={pr_no!r} — cannot extract."
        )
    if not ctx.line_items:
        raise NoLineItemsError(
            f"PR={pr_no!r} has no rows in purchase_req_detail — nothing to extract against."
        )

    sources = _resolve_quotation_sources(tgt_cs, pr_no)
    if not sources:
        raise NoQuotationFoundError(
            f"No quotation documents found for PR={pr_no!r}. "
            f"None of the attachments are classified as 'Quotation'."
        )

    # Each source is: download blob → load doc → LLM extract → parse → align.
    # This is I/O + LLM bound, so threads give real speedup when a PR has
    # multiple quotation files. Keep n_workers low (default 1) to avoid
    # multiplying LLM rate-limit pressure on top of parallel_workers.
    n_workers = max(1, int((prompts or {}).get("ext_parallel_sources", 1)))
    max_pages = max(1, int((prompts or {}).get("ext_max_pages", 20)))

    def _extract_one(src: dict) -> list[dict]:
        blob_path = src["blob_path"]
        filename  = os.path.basename(blob_path)
        try:
            file_bytes = _download_blob(blob_path, blob_cfg)
            doc        = _load_document(file_bytes, filename, max_pages=max_pages)
            raw        = _call_extraction_llm(llm, ctx, doc, prompts)
            items      = _parse_extraction_response(raw, src, ctx)
            items      = _align_to_ras_line_items(items, ctx, src)
            logger.info(
                "[{}] Extracted {} item(s) from {!r} ({} page(s) loaded, {} image(s))",
                pr_no, len(items), filename, doc.page_count,
                len(doc.images) if doc.images else 0,
            )
            return items
        except Exception as exc:
            logger.opt(exception=True).warning(
                "[{}] Extraction failed for {!r}: {}", pr_no, filename, exc
            )
            return []

    if n_workers == 1:
        all_items: list[dict] = []
        for src in sources:
            all_items.extend(_extract_one(src))
    else:
        import concurrent.futures
        logger.info(f"[{pr_no}] Extracting {len(sources)} source(s) with {n_workers} parallel worker(s)")
        with concurrent.futures.ThreadPoolExecutor(max_workers=n_workers) as pool:
            results = list(pool.map(_extract_one, sources))
        all_items = [item for batch in results for item in batch]

    if not all_items:
        # Every quotation source raised inside _extract_one (LLM error, parse
        # failure, blob 404, etc.) and was caught individually. With nothing to
        # save, abort so this PR is moved to exception state instead of
        # advancing through stages 6-8 with empty data.
        raise AllExtractionsFailedError(
            f"All {len(sources)} quotation source(s) for PR={pr_no!r} "
            f"failed extraction; no items produced."
        )

    _apply_price_alignment_boost(all_items, ctx, prompts)
    _canonicalize_supplier_names(all_items, ctx, prompts)
    _compute_quote_ranks(all_items)
    _select_best_quotes(all_items, ctx)
    saved = _save_extracted_items(tgt_cs, all_items)
    logger.info(f"[{pr_no}] {saved} item(s) written to quotation_extracted_items")
    return saved


# ── Tracker helpers ────────────────────────────────────────────────────────────

def _advance_tracker(tgt_cs: str, pr_no: str, stage_id: int) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[ras_tracker]
               SET current_stage_fk=?, updated_at=SYSUTCDATETIME()
             WHERE purchase_req_no=?
        """, stage_id, pr_no)
        conn.commit()
    except Exception as exc:
        logger.warning(f"Tracker advance failed PR={pr_no!r} stage={stage_id}: {exc}")
    finally:
        conn.close()


def _set_last_processed_at(tgt_cs: str, pr_no: str) -> None:
    """Stamp last_processed_at = SYSUTCDATETIME() so SourceChangeDetector
    can detect future on-prem source changes for this PR."""
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[ras_tracker]
               SET last_processed_at = SYSUTCDATETIME(),
                   updated_at        = SYSUTCDATETIME()
             WHERE purchase_req_no   = ?
        """, pr_no)
        conn.commit()
    except Exception as exc:
        logger.warning(f"set_last_processed_at failed PR={pr_no!r}: {exc}")
    finally:
        conn.close()


def _record_exception(tgt_cs: str, pr_no: str, stage_id: int, error_msg: str) -> None:
    try:
        conn = _connect(tgt_cs)
        cur  = conn.cursor()
        try:
            # Also stamp last_processed_at so audits / reporting can see when
            # this PR last reached a terminal state (success OR exception).
            cur.execute("""
                UPDATE [ras_procurement].[ras_tracker]
                   SET current_stage_fk   = 99,
                       last_processed_at  = SYSUTCDATETIME(),
                       updated_at         = SYSUTCDATETIME()
                 WHERE purchase_req_no    = ?
            """, pr_no)
            cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no=?", pr_no)
            row = cur.fetchone()
            if row:
                cur.execute("""
                    INSERT INTO [ras_procurement].[ras_pipeline_exceptions]
                        (ras_tracker_id, stage_id, exception_message)
                    VALUES (?, ?, ?)
                """, row[0], stage_id, error_msg[:4000])
            conn.commit()
        finally:
            conn.close()
    except Exception as exc:
        logger.warning(f"[{pr_no}] Could not write exception record: {exc}")


# ── Embeddings + benchmark (Stage 6-7) ────────────────────────────────────────

def _build_embed_text(row_dict: dict) -> str:
    """Build the 12-field embedding text (identical order used by doc intel branch)."""
    _FIELDS = [
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag",
    ]
    return " | ".join(str(row_dict[f]) for f in _FIELDS if row_dict.get(f))


def _run_embeddings(tgt_cs: str, pr_no: str, embed_model, pinecone_index: str, pinecone_ns: str) -> None:
    # Structural errors (import, DB connect, index creation) propagate — caller records exception.
    from agentcore.services.pinecone_service_client import ensure_index_via_service, ingest_via_service
    ensure_index_via_service(index_name=pinecone_index, embedding_dimension=3072)
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8], qi.[commodity_tag],
                   prd.[C_DATETIME] AS [item_created_date]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        all_rows = cur.fetchall()
    finally:
        conn.close()

    # deduplicate by purchase_dtl_id — first row wins (same as doc intel branch)
    seen_dtl: dict = {}
    for row in all_rows:
        dtl_id = row[1]
        if dtl_id not in seen_dtl:
            seen_dtl[dtl_id] = row
    rows = list(seen_dtl.values())
    logger.info(f"[{pr_no}] Embedding {len(rows)} selected item(s) (deduped from {len(all_rows)} is_selected_quote=1 rows)")

    cols = [
        "extracted_item_uuid_pk", "purchase_dtl_id",
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag", "item_created_date",
    ]
    for row in rows:
        rd = dict(zip(cols, row))
        dtl_id    = rd["purchase_dtl_id"]
        item_uuid = rd["extracted_item_uuid_pk"]
        created   = rd.get("item_created_date")
        created_iso = (
            created.isoformat() if created and hasattr(created, "isoformat") else str(created or "")
        )
        content = _build_embed_text(rd)
        if not content:
            continue
        try:
            embedding = embed_model.embed_query(content)
            ingest_via_service(
                index_name=pinecone_index,
                namespace=pinecone_ns,
                text_key="page_content",
                documents=[{
                    "page_content": content,
                    "metadata": {
                        "purchase_req_no":        pr_no,
                        "purchase_dtl_id":        int(dtl_id),
                        "extracted_item_uuid_pk": str(item_uuid or ""),
                        "commodity_tag":          str(rd.get("commodity_tag") or ""),
                        "item_created_date":      created_iso,
                    },
                }],
                embedding_vectors=[embedding],
                vector_ids=[f"dtl_{dtl_id}"],
                embedding_dimension=3072,
            )
            logger.info(f"[{pr_no}] Upserted vector dtl_{dtl_id} (item_created={created_iso})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Embedding failed for dtl_id {dtl_id}: {exc}")


def _fetch_historical_for_dtl_ids(tgt_cs: str, dtl_ids: list) -> list[dict]:
    """Fetch pricing rows from quotation_extracted_items for a list of purchase_dtl_id values."""
    if not dtl_ids:
        return []
    placeholders = ", ".join(["?"] * len(dtl_ids))
    sql = f"""
        SELECT qi.[purchase_dtl_id], qi.[extracted_item_uuid_pk],
               qi.[unit_price], qi.[total_price], qi.[quantity], qi.[unit],
               qi.[currency], qi.[quotation_date], qi.[supplier_name], qi.[supplier_country],
               qi.[unit_price_eur], qi.[total_price_eur],
               rt.[purchase_req_no]
          FROM [ras_procurement].[quotation_extracted_items] qi
          JOIN [ras_procurement].[attachment_classification] ac
            ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
          JOIN [ras_procurement].[ras_tracker] rt
            ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
         WHERE qi.[purchase_dtl_id] IN ({placeholders})
           AND qi.[is_selected_quote] = 1
    """
    cols = [
        "purchase_dtl_id", "extracted_item_uuid_pk",
        "unit_price", "total_price", "quantity", "unit",
        "currency", "quotation_date", "supplier_name", "supplier_country",
        "unit_price_eur", "total_price_eur", "purchase_req_no",
    ]
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute(sql, *dtl_ids)
        return [dict(zip(cols, r)) for r in cur.fetchall()]
    finally:
        conn.close()


def _compute_low_last(items: list[dict]) -> tuple:
    """Return (low_item, last_item) — cheapest EUR price and most recent quotation_date."""
    def _eur(it: dict):
        return it.get("unit_price_eur") or it.get("unit_price")

    priced = [it for it in items if _eur(it) is not None]
    dated  = [it for it in items if it.get("quotation_date") is not None]
    low_item  = min(priced, key=_eur) if priced else None
    last_item = max(dated,  key=lambda it: it["quotation_date"]) if dated else None
    return low_item, last_item


def _format_bench_prompt(current: dict, historical: list[dict]) -> str:
    """Build structured LLM prompt matching doc intel branch style."""
    eur_unit = current.get("unit_price_eur")
    raw_price = current.get("unit_price")
    currency  = current.get("currency") or ""
    if eur_unit is not None and currency.upper() != "EUR":
        price_str = f"{eur_unit} EUR  (original: {raw_price} {currency})"
    elif eur_unit is not None:
        price_str = f"{eur_unit} EUR"
    else:
        price_str = f"{raw_price} {currency}"

    category = " > ".join(
        str(current[f]) for f in [
            "item_level_1", "item_level_2", "item_level_3", "item_level_4",
            "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        ] if current.get(f)
    )
    current_block = "\n".join([
        f"- Item name        : {current.get('item_name') or 'N/A'}",
        f"- Description      : {current.get('item_description') or 'N/A'}",
        f"- Category         : {category or 'N/A'}",
        f"- Commodity tag    : {current.get('commodity_tag') or 'N/A'}",
        f"- Quantity         : {current.get('quantity')} {current.get('unit') or ''}",
        f"- Quoted unit price: {price_str}",
        f"- Supplier         : {current.get('supplier_name') or 'N/A'}",
        f"- Quotation date   : {current.get('quotation_date') or 'N/A'}",
    ])

    if not historical:
        hist_block = "No historical data available."
    else:
        header = "| # | Supplier | Date | Qty | Unit | Unit Price (EUR) | Orig. Currency |"
        sep    = "|---|----------|------|-----|------|-----------------|----------------|"
        rows_out = []
        for i, it in enumerate(historical, 1):
            eur_p = it.get("unit_price_eur")
            raw_p = it.get("unit_price")
            eur_label = str(eur_p) if eur_p is not None else (f"{raw_p} *" if raw_p is not None else "N/A")
            rows_out.append(
                f"| {i} | {it.get('supplier_name') or 'N/A'} "
                f"| {it.get('quotation_date') or 'N/A'} "
                f"| {it.get('quantity') or 'N/A'} "
                f"| {it.get('unit') or 'N/A'} "
                f"| {eur_label} "
                f"| {it.get('currency') or 'N/A'} |"
            )
        hist_block = "\n".join([header, sep] + rows_out)

    return (
        f"Recommend a benchmark unit price for this item based on historical similar purchases.\n\n"
        f"CURRENT ITEM:\n{current_block}\n\n"
        f"HISTORICAL SIMILAR ITEMS:\n{hist_block}\n\n"
        f"Return ONLY JSON (no markdown): "
        f'{{ "bp_unit_price": <number in EUR or null>, "summary": "<2-3 sentences>" }}'
    )


def _compute_cpi_pct(country_str: str | None, start_year: int, end_year: int):
    """Cumulative CPI inflation % via World Bank API (FP.CPI.TOTL.ZG).

    Mirrors utils/cpi_inflation.py from the doc intel branch.
    Returns float or None on any failure.
    """
    if not country_str or start_year >= end_year:
        return None
    try:
        import httpx as _httpx
        import pycountry as _pc
        from rapidfuzz import process as _fuzz
        # Resolve country name → ISO alpha-2
        text = country_str.strip().lower()
        alpha2 = None
        for c in _pc.countries:
            if text in (c.name.lower(), getattr(c, "official_name", "").lower(),
                        c.alpha_2.lower(), c.alpha_3.lower()):
                alpha2 = c.alpha_2
                break
        if alpha2 is None:
            names = [c.name for c in _pc.countries]
            res = _fuzz.extractOne(text, names)
            if res and res[1] > 75:
                matched = _pc.countries.get(name=res[0])
                if matched:
                    alpha2 = matched.alpha_2
        if alpha2 is None:
            logger.info(f"CPI: could not resolve country {country_str!r}")
            return None
        # Fetch World Bank annual rates
        url = f"https://api.worldbank.org/v2/country/{alpha2}/indicator/FP.CPI.TOTL.ZG"
        resp = _httpx.get(url, params={"date": f"{start_year}:{end_year}", "format": "json", "per_page": 100}, timeout=30.0)
        resp.raise_for_status()
        data = resp.json()
        if len(data) < 2 or not data[1]:
            return None
        rates = {int(e["date"]): float(e["value"]) for e in data[1] if e.get("value") is not None}
        if not rates:
            return None
        factor = 1.0
        for yr in range(start_year + 1, end_year + 1):
            if yr in rates:
                factor *= (1 + rates[yr] / 100)
        return round((factor - 1) * 100, 4)
    except Exception as exc:
        logger.info(f"CPI: failed for {country_str!r} {start_year}-{end_year}: {exc}")
        return None


def _estimate_inflation_via_llm(llm, item_name: str | None, category: str | None,
                                 supplier_country: str | None, ref_year: int | None,
                                 current_year: int | None) -> float | None:
    """Ask LLM for estimated cumulative inflation % for an item category in a country.

    Returns float or None. Does NOT derive inflation from historical price delta.
    """
    if not supplier_country or not ref_year or not current_year or current_year <= ref_year:
        return None
    try:
        from langchain_core.messages import HumanMessage as _HM
        import json as _json, re as _re
        years = current_year - ref_year
        prompt = (
            f"I am buying the following item from {supplier_country}:\n"
            f"  Item     : {item_name or 'N/A'}\n"
            f"  Category : {category or 'N/A'}\n\n"
            f"What is the estimated cumulative inflation rate (%) for this type of item "
            f"in {supplier_country} from {ref_year} to {current_year} ({years} year(s))?\n\n"
            f"Base your estimate on macroeconomic inflation trends for this category in that country. "
            f"Do NOT derive the rate from any historical price data.\n"
            f'Respond ONLY with JSON: {{ "inflation_pct": <number or null> }}'
        )
        resp = llm.invoke([_HM(content=prompt)])
        raw  = (getattr(resp, "content", None) or str(resp)).strip()
        raw  = _re.sub(r"^```(?:json)?\s*", "", raw)
        raw  = _re.sub(r"\s*```$",          "", raw).strip()
        val  = _json.loads(raw).get("inflation_pct")
        return float(val) if val is not None else None
    except Exception as exc:
        logger.info(f"LLM inflation estimate failed ({supplier_country} {ref_year}-{current_year}): {exc}")
        return None


def _run_benchmark(llm, tgt_cs: str, pr_no: str, embed_model, pinecone_index: str, pinecone_ns: str, top_k: int) -> None:
    # Structural errors (import, DB connect, query) propagate — caller records exception.
    from agentcore.services.pinecone_service_client import search_via_service
    from langchain_core.messages import HumanMessage
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8], qi.[commodity_tag],
                   qi.[unit_price], qi.[total_price], qi.[quantity], qi.[unit],
                   qi.[currency], qi.[quotation_date], qi.[supplier_name],
                   qi.[unit_price_eur], qi.[total_price_eur],
                   prd.[C_DATETIME] AS [item_created_date]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
              LEFT JOIN [ras_procurement].[purchase_req_detail] prd
                ON qi.[purchase_dtl_id] = prd.[PURCHASE_DTL_ID]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        all_rows = cur.fetchall()
    finally:
        conn.close()

    bench_cols = [
        "extracted_item_uuid_pk", "purchase_dtl_id",
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag",
        "unit_price", "total_price", "quantity", "unit",
        "currency", "quotation_date", "supplier_name",
        "unit_price_eur", "total_price_eur", "item_created_date",
    ]
    # deduplicate by purchase_dtl_id
    seen: dict = {}
    for row in all_rows:
        rd = dict(zip(bench_cols, row))
        dtl_id = rd["purchase_dtl_id"]
        if dtl_id not in seen:
            seen[dtl_id] = rd
    items = list(seen.values())

    conn2 = _connect(tgt_cs)
    cur2  = conn2.cursor()
    try:
        for rd in items:
            dtl_id    = rd["purchase_dtl_id"]
            item_uuid = str(rd["extracted_item_uuid_pk"] or "")
            qty       = rd.get("quantity") or Decimal("1")

            # Build the same 12-field text used at embedding time
            bench_text = _build_embed_text(rd)
            if not bench_text:
                continue

            # Item created_date as ISO string for date filtering
            created = rd.get("item_created_date")
            created_iso = (
                created.isoformat() if created and hasattr(created, "isoformat") else str(created or "")
            )

            try:
                embedding = embed_model.embed_query(bench_text)
                raw_similar = search_via_service(
                    index_name=pinecone_index,
                    namespace=pinecone_ns,
                    text_key="page_content",
                    query=bench_text,
                    query_embedding=embedding,
                    number_of_results=top_k * 3,  # fetch extra to absorb score + date filtered-out
                )
            except Exception as exc:
                logger.warning(f"[{pr_no}] Benchmark similarity search failed dtl_id={dtl_id}: {exc}")
                continue

            # Normalise response to a list of match dicts
            if isinstance(raw_similar, list):
                matches = raw_similar
            elif isinstance(raw_similar, dict):
                matches = raw_similar.get("results", raw_similar.get("matches", []))
            else:
                matches = []

            # Filter: minimum similarity score, exclude same PR, only items from older PRs
            _MIN_SCORE = 0.70
            filtered = []
            for m in matches:
                score = float(m.get("score", 0.0))
                if score < _MIN_SCORE:
                    continue
                meta = m.get("metadata") or {}
                if meta.get("purchase_req_no") == pr_no:
                    continue
                hist_date = meta.get("item_created_date", "")
                if created_iso and hist_date and hist_date >= created_iso:
                    continue
                filtered.append(m)

            similar_dtl_ids_raw = [
                int(m["metadata"]["purchase_dtl_id"])
                for m in filtered[:top_k]
                if m.get("metadata", {}).get("purchase_dtl_id") is not None
            ]
            logger.info(
                f"[{pr_no}] dtl_id={dtl_id}: {len(filtered)} Pinecone match(es) after filter "
                f"→ dtl_ids={similar_dtl_ids_raw}"
            )

            # Fetch full historical pricing from DB
            historical = _fetch_historical_for_dtl_ids(tgt_cs, similar_dtl_ids_raw) if similar_dtl_ids_raw else []
            low_item, last_item = _compute_low_last(historical)
            low_uuid  = str(low_item["extracted_item_uuid_pk"])  if low_item  else None
            last_uuid = str(last_item["extracted_item_uuid_pk"]) if last_item else None
            similar_dtl_ids_json = json.dumps(similar_dtl_ids_raw) if similar_dtl_ids_raw else None

            # Inflation estimates (non-fatal — both return None on failure)
            supplier_country = (low_item.get("supplier_country") if low_item else None) or rd.get("supplier_country")
            ref_dt       = low_item.get("quotation_date") if low_item else None
            ref_year     = ref_dt.year if ref_dt and hasattr(ref_dt, "year") else None
            current_year = created.year if created and hasattr(created, "year") else None
            item_category = " > ".join(
                str(rd[f]) for f in ["item_level_1", "item_level_2", "item_level_3"] if rd.get(f)
            )
            infl_dec = cpi_dec = None
            if ref_year and current_year and ref_year < current_year:
                infl_raw = _estimate_inflation_via_llm(
                    llm, rd.get("item_name"), item_category or None,
                    supplier_country, ref_year, current_year,
                )
                if infl_raw is not None:
                    try: infl_dec = Decimal(str(infl_raw))
                    except Exception: pass
                cpi_raw = _compute_cpi_pct(supplier_country, ref_year, current_year)
                if cpi_raw is not None:
                    try: cpi_dec = Decimal(str(cpi_raw))
                    except Exception: pass
                logger.info(
                    f"[{pr_no}] dtl_id={dtl_id}: inflation_pct={infl_dec} "
                    f"cpi_pct={cpi_dec} (country={supplier_country!r} {ref_year}-{current_year})"
                )

            # LLM benchmark analysis
            bench_prompt = _format_bench_prompt(rd, historical[:top_k])
            bout = {}
            try:
                resp = llm.invoke([HumanMessage(content=bench_prompt)])
                raw  = (getattr(resp, "content", None) or str(resp)).strip()
                raw  = re.sub(r"^```(?:json)?\s*", "", raw)
                raw  = re.sub(r"\s*```$", "", raw)
                bout = json.loads(raw)
            except Exception as exc:
                logger.warning(f"[{pr_no}] LLM benchmark parse failed dtl_id={dtl_id}: {exc}")

            bp_unit = bout.get("bp_unit_price")
            summary = bout.get("summary", "")
            try:
                bp_dec   = Decimal(str(bp_unit)) if bp_unit is not None else None
                bp_total = round(float(bp_dec) * float(qty or 1), 2) if bp_dec is not None else None
            except Exception:
                bp_dec = bp_total = None

            try:
                cur2.execute("""
                    MERGE [ras_procurement].[benchmark_result] WITH (HOLDLOCK) AS target
                    USING (SELECT ? AS purchase_dtl_id) AS src
                       ON target.purchase_dtl_id = src.purchase_dtl_id
                    WHEN MATCHED THEN
                        UPDATE SET
                            extracted_item_uuid_fk = ?,
                            bp_unit_price          = ?,
                            bp_total_price         = ?,
                            low_hist_item_fk       = ?,
                            last_hist_item_fk      = ?,
                            inflation_pct          = ?,
                            cpi_inflation_pct      = ?,
                            similar_dtl_ids        = ?,
                            summary                = ?,
                            updated_at             = SYSUTCDATETIME()
                    WHEN NOT MATCHED THEN
                        INSERT (
                            purchase_dtl_id, extracted_item_uuid_fk,
                            bp_unit_price, bp_total_price,
                            low_hist_item_fk, last_hist_item_fk,
                            inflation_pct, cpi_inflation_pct,
                            similar_dtl_ids, summary
                        )
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?);
                """,
                    dtl_id,
                    item_uuid, bp_dec, bp_total, low_uuid, last_uuid, infl_dec, cpi_dec, similar_dtl_ids_json, summary,
                    dtl_id, item_uuid, bp_dec, bp_total, low_uuid, last_uuid, infl_dec, cpi_dec, similar_dtl_ids_json, summary,
                )
            except Exception as exc:
                logger.warning(f"[{pr_no}] Benchmark write failed dtl_id={dtl_id}: {exc}")

        conn2.commit()
    except Exception:
        try: conn2.rollback()
        except Exception: pass
        raise
    finally:
        conn2.close()


# ── Fetch pending PRs ──────────────────────────────────────────────────────────

def _get_pr_current_stage(tgt_cs: str, pr_no: str) -> int:
    """Return the current_stage_fk for a PR, or 0 if not found."""
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute(
            "SELECT current_stage_fk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?",
            pr_no,
        )
        row = cur.fetchone()
        return int(row[0]) if row and row[0] is not None else 0
    finally:
        conn.close()


def _fetch_pending_prs(tgt_cs: str, pr_filter: str, batch_limit: int) -> list:
    if pr_filter:
        return [pr_filter]
    # Pick up PRs at any stage that needs stages 4-8 work:
    #   stage 3 = just finished Stage 1-3, ready for classification
    #   stages 4-7 = partially processed (e.g. failed mid-run), need to resume
    _WHERE = """
        FROM [ras_procurement].[purchase_req_mst] prm
        JOIN [ras_procurement].[ras_tracker] rt
          ON prm.PURCHASE_REQ_NO = rt.purchase_req_no
       WHERE rt.current_stage_fk IN (3, 4, 5, 6, 7)
         AND UPPER(prm.PURCHASEFINALAPPROVALSTATUS)
                 IN ('APPROVED BY ALL', 'APPROVED BY ALL EXCEPTION')
       ORDER BY prm.C_DATETIME ASC
    """
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        if batch_limit > 0:
            cur.execute(f"SELECT TOP (?) prm.PURCHASE_REQ_NO {_WHERE}", batch_limit)
        else:
            cur.execute(f"SELECT prm.PURCHASE_REQ_NO {_WHERE}")
        return [row[0] for row in cur.fetchall()]
    finally:
        conn.close()


# ── Excel-driven PR list (mirrors doc-intel run_pipeline_from_excel.py) ──

_AUTO_DETECT_PR_COLUMNS = {
    "purchase_req_no", "purchase req no",
    "pr_no", "pr no", "pr number",
    "ras_no", "ras no", "ras number",
}


def _parse_pr_list_from_excel_bytes(
    excel_bytes: bytes,
    column: str | None = None,
    sheet: str | None = None,
) -> list[str]:
    """Parse purchase requisition numbers out of an .xlsx blob.

    Mirrors doc-intel run_pipeline_from_excel.py._read_pr_nos_from_excel,
    adapted to read from in-memory bytes (Azure Blob / wired Data input)
    instead of a local file path.

    Column selection:
      1. explicit `column` arg → exact case-insensitive header match
      2. auto-detect from PURCHASE_REQ_NO / PR_NO / RAS_NO / "PR Number" etc.
      3. fall back to the first column

    Returns deduplicated, non-empty PR numbers in the order they appear.
    """
    import io as _io
    from openpyxl import load_workbook

    wb = load_workbook(_io.BytesIO(excel_bytes), read_only=True, data_only=True)
    try:
        sheet_name = (sheet or "").strip()
        if sheet_name:
            if sheet_name not in wb.sheetnames:
                raise ValueError(
                    f"Sheet {sheet_name!r} not found. Available: {wb.sheetnames}"
                )
            ws = wb[sheet_name]
        else:
            ws = wb.active

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return []

        headers = [
            str(h).strip() if h is not None else ""
            for h in rows[0]
        ]

        col_idx: int | None = None
        col_arg = (column or "").strip()
        if col_arg:
            for i, h in enumerate(headers):
                if h.lower() == col_arg.lower():
                    col_idx = i
                    break
            if col_idx is None:
                raise ValueError(
                    f"Column {col_arg!r} not found. Headers: {headers}"
                )
        else:
            for i, h in enumerate(headers):
                if h.lower() in _AUTO_DETECT_PR_COLUMNS:
                    col_idx = i
                    logger.info("Excel PR list — auto-detected column {!r}", h)
                    break
            if col_idx is None:
                col_idx = 0
                logger.info(
                    "Excel PR list — no recognised PR column header; using first column {!r}",
                    headers[0] if headers else "<empty>",
                )

        seen: set[str] = set()
        pr_nos: list[str] = []
        for row in rows[1:]:
            if col_idx >= len(row):
                continue
            val = row[col_idx]
            if val is None:
                continue
            pr_no = str(val).strip()
            if pr_no and pr_no not in seen:
                seen.add(pr_no)
                pr_nos.append(pr_no)
        return pr_nos
    finally:
        wb.close()


# ── Main Component ─────────────────────────────────────────────────────────────


# ── Main component ────────────────────────────────────────────────────────

class PipelineStage123Node(Node):
    display_name = "Full Pipeline (Stages 1-8)"
    description  = (
        "Runs the complete procurement pipeline end-to-end per PR in parallel workers: "
        "Stage 1 (Ingestion) → Stage 2 (Embed Doc Extraction) → Stage 3 (Blob Upload) → "
        "Stage 4 (Classification) → Stage 5 (Extraction) → Stage 6 (Embeddings) → "
        "Stage 7 (Price Benchmark) → Stage 8 (Complete). "
        "Wire Prompt Template nodes to the optional prompt inputs to override defaults."
    )
    icon = "Database"
    name = "PipelineStage123Node"

    inputs = [
        HandleInput(
            name="source_connection",
            display_name="Source DB — On-Prem SQL Server",
            input_types=["Data"],
            info="Connection Config from the on-prem Database Connector node.",
        ),
        HandleInput(
            name="target_connection",
            display_name="Target DB — Azure SQL",
            input_types=["Data"],
            info="Connection Config from the Azure SQL Database Connector node.",
        ),
        MessageTextInput(
            name="blob_connector_name",
            display_name="Azure Blob Connector Name",
            value="",
            info="Exact name of your Azure Blob connector in Settings → Connectors.",
        ),
        MessageTextInput(
            name="pr_no_filter",
            display_name="PR Number Filter (optional)",
            value="",
            advanced=True,
            info=(
                "Leave blank to process all pending approved PRs. "
                "Enter one PURCHASE_REQ_NO to force a full reprocess from scratch."
            ),
        ),
        # ── Excel-driven PR list (mirrors doc-intel run_pipeline_from_excel.py) ──
        # Use ANY ONE of:
        #   - PR Number Filter            → single PR, full reprocess
        #   - Excel Blob Path             → download Excel from the configured
        #                                   Blob Connector container, parse PR list
        #   - PR List Data (HandleInput)  → wire any component (Blob Reader,
        #                                   Knowledge Base, etc.) that outputs
        #                                   either raw bytes or a list under
        #                                   data["pr_numbers"]
        # If none are set, the component falls back to fetching pending PRs
        # from ras_tracker (the standard batch flow).
        MessageTextInput(
            name="pr_excel_blob_path",
            display_name="PR List Excel — Blob Path (optional)",
            value="",
            advanced=True,
            info=(
                "Blob path inside the configured Blob Connector container "
                "pointing to an .xlsx with a PURCHASE_REQ_NO column "
                "(e.g. procurement/pr_lists/2026-04.xlsx). Each PR is "
                "force-reprocessed from scratch like the single-PR filter."
            ),
        ),
        HandleInput(
            name="pr_list_data",
            display_name="PR List Data (optional)",
            input_types=["Data"],
            required=False,
            info=(
                "Wire a Blob Reader / Knowledge Base / file-loader component "
                "that outputs Data here. Accepted shapes: "
                "(a) data['pr_numbers'] = [list of PURCHASE_REQ_NO strings] — "
                "consumed directly; or "
                "(b) data['content'] / data['bytes'] / data['file_bytes'] = "
                "raw .xlsx bytes — parsed inline."
            ),
        ),
        MessageTextInput(
            name="pr_excel_column",
            display_name="PR List Excel — Column Name (optional)",
            value="",
            advanced=True,
            info=(
                "Header of the column holding PURCHASE_REQ_NO values. "
                "Auto-detected from common names "
                "(PURCHASE_REQ_NO / PR_NO / RAS_NO / PR Number) "
                "if blank, else falls back to the first column."
            ),
        ),
        MessageTextInput(
            name="pr_excel_sheet",
            display_name="PR List Excel — Sheet Name (optional)",
            value="",
            advanced=True,
            info="Worksheet to read. Defaults to the first sheet.",
        ),
        # ── Stages 4-8 inputs ────────────────────────────────────────────────
        HandleInput(
            name="llm",
            display_name="LLM (GPT-4o or GPT-4o-mini)",
            input_types=["LanguageModel"],
            info="LLM used for Stage 4 (classification) and Stage 5 (extraction).",
        ),
        HandleInput(
            name="embed_model",
            display_name="Embeddings Model",
            input_types=["Embeddings"],
            info="Embeddings model used for Stage 6 (vector embeddings) and Stage 7 (benchmark search).",
        ),
        # ── Prompt template overrides (optional) ─────────────────────────────
        # Leave disconnected to use the default prompts baked into this component.
        # Wire a Prompt Template node to any of these to override the built-in prompts.
        HandleInput(
            name="cls_system_prompt",
            display_name="[Classification] System Prompt",
            input_types=["Message"],
            required=False,
            info="Override the classification system prompt. Connect a Prompt Template node.",
        ),
        HandleInput(
            name="cls_user_text_prompt",
            display_name="[Classification] User Prompt — Text/Tabular Files",
            input_types=["Message"],
            required=False,
            info="Override the user prompt for text/Excel/Word/PDF files. Must keep {filename}, {file_type}, {extra_metadata}, {extracted_content} placeholders.",
        ),
        HandleInput(
            name="cls_user_image_prompt",
            display_name="[Classification] User Prompt — Image/Scanned Files",
            input_types=["Message"],
            required=False,
            info="Override the user prompt sent with base64 image content. Must keep {filename}, {file_type}, {extra_metadata} placeholders.",
        ),
        HandleInput(
            name="ext_system_prompt",
            display_name="[Extraction] System Prompt",
            input_types=["Message"],
            required=False,
            info="Override the extraction system prompt. Connect a Prompt Template node.",
        ),
        HandleInput(
            name="ext_user_template",
            display_name="[Extraction] User Prompt Template",
            input_types=["Message"],
            required=False,
            info="Override the extraction user template. Must keep all {field} placeholders from the RAS context and {document_content}.",
        ),
        HandleInput(
            name="ext_item_taxonomy",
            display_name="[Extraction] Item Taxonomy",
            input_types=["Message"],
            required=False,
            info="Override the item taxonomy guidelines injected into {item_taxonomy} of the extraction user template.",
        ),
        # ── Advanced settings ─────────────────────────────────────────────────
        MessageTextInput(
            name="pinecone_index",
            display_name="Pinecone Index Name",
            value="ras-quotations",
            advanced=True,
            info="Pinecone index used for embeddings (Stage 6) and retry cleanup.",
        ),
        MessageTextInput(
            name="pinecone_namespace",
            display_name="Pinecone Namespace",
            value="procurement",
            advanced=True,
            info="Pinecone namespace used for embeddings (Stage 6) and retry cleanup.",
        ),
        IntInput(name="pinecone_top_k",    display_name="Benchmark Top-K (Stage 7)", value=5,     advanced=True),
        IntInput(name="batch_limit",       display_name="Max PRs per Run (0 = all)", value=0,     advanced=True,
                 info="0 = process all pending PRs. Set to N (e.g. 50) to cap at N PRs per run. Ignored when PR Number Filter is set."),
        IntInput(name="parallel_workers",  display_name="Parallel Workers",          value=4,     advanced=True),
        IntInput(name="max_content_chars", display_name="Max Chars per File (Stage 2)", value=80000, advanced=True),
        # ── Stage 4-7 tuning knobs ────────────────────────────────────────────
        IntInput(
            name="cls_max_chars",
            display_name="Max Chars — Classification Input (Stage 4)",
            value=24000,
            advanced=True,
            info=(
                "Truncate extracted text before sending to the classification LLM. "
                "Default 24000 ≈ 6000 tokens, matches doc-intel "
                "truncate_to_token_limit(max_content_tokens=6000). 0 = no limit. "
                "Lowering this may cause Quotation signals (payment terms, totals, "
                "signatures) at the end of long documents to be missed."
            ),
        ),
        IntInput(
            name="ext_max_chars",
            display_name="Max Chars — Extraction Document (Stage 5)",
            value=50000,
            advanced=True,
            info="Truncate document content before sending to the extraction LLM. 0 = no limit.",
        ),
        IntInput(
            name="ext_max_pages",
            display_name="Max Pages — PDF/Doc Render (Stage 5)",
            value=20,
            advanced=True,
            info=(
                "Max pages rendered per quotation document during extraction. "
                "PDFs longer than this only have their first N pages processed; "
                "a warning is logged when truncation happens. Matches doc-intel "
                "MAX_PAGES default."
            ),
        ),
        IntInput(
            name="ext_max_images",
            display_name="Max Images per LLM Call (Stage 5)",
            value=50,
            advanced=True,
            info=(
                "Hard cap on images attached to one extraction LLM request. "
                "Azure OpenAI rejects requests with more than ~50 images per call, "
                "so the value should not be raised above 50."
            ),
        ),
        MessageTextInput(
            name="selected_threshold",
            display_name="Selected Quote Min Score (Stage 7)",
            value="0.70",
            advanced=True,
            info="Minimum supplier_match_conf (0.0–1.0) for a quote to be eligible for selection.",
        ),
        MessageTextInput(
            name="price_max_boost",
            display_name="Price Alignment Max Boost (Stage 7)",
            value="0.10",
            advanced=True,
            info="Maximum conf boost (0.0–1.0) applied when extracted prices align with RAS prices.",
        ),
        MessageTextInput(
            name="canonicalize_threshold",
            display_name="Supplier Canonicalize Threshold (Stage 7)",
            value="0.82",
            advanced=True,
            info="SequenceMatcher ratio (0.0–1.0) above which two supplier names are merged into one canonical name.",
        ),
        # ── LLM retry / rate-limit settings ──────────────────────────────────
        IntInput(
            name="llm_max_retries",
            display_name="LLM Max Retries (Rate Limit)",
            value=3,
            advanced=True,
            info="How many times to retry an LLM call after a 429 / token-quota error before failing the stage.",
        ),
        IntInput(
            name="llm_retry_cooldown",
            display_name="LLM Retry Cooldown (seconds)",
            value=60,
            advanced=True,
            info=(
                "Seconds to wait before the first retry. Each subsequent retry waits "
                "cooldown × attempt (60s → 120s → 180s with default 3 retries)."
            ),
        ),
        # ── Extraction parallelism ────────────────────────────────────────────
        IntInput(
            name="ext_parallel_sources",
            display_name="Extraction Parallel Sources (Stage 5)",
            value=1,
            advanced=True,
            info=(
                "Number of quotation files extracted in parallel for a single PR. "
                "1 = sequential (safe). Increase to 2-3 to speed up PRs with many "
                "quotation files, but keep low to avoid multiplying LLM rate-limit "
                "pressure on top of the parallel_workers setting."
            ),
        ),
    ]

    outputs = [
        Output(display_name="File Batch", name="file_batch", method="build_file_batch", types=["Message"]),
        Output(display_name="Processed PRs", name="processed_prs", method="get_processed_prs", types=["Data"]),
    ]

    # ── Connection helpers ────────────────────────────────────────────────

    def _conn_str(self, conn_data: Data) -> str:
        d      = conn_data.data or {}
        driver = d.get("driver", "ODBC Driver 18 for SQL Server")
        server = d.get("host", d.get("server", ""))
        port   = d.get("port", 1433)
        db     = d.get("database_name", d.get("database", ""))
        user   = d.get("username", "")
        pwd    = d.get("password", "")
        return (
            f"DRIVER={{{driver}}};SERVER={server},{port};"
            f"DATABASE={db};UID={user};PWD={pwd};TrustServerCertificate=yes;"
        )

    @staticmethod
    def _is_transient(exc: Exception) -> bool:
        kw = ["connection reset", "timeout", "throttl", "resource limit", "broken pipe", "transport-level", "login failed"]
        return any(k in str(exc).lower() for k in kw)

    def _connect(self, conn_str: str):
        import pyodbc, random, time
        for attempt in range(_MAX_RETRIES + 1):
            try:
                return pyodbc.connect(conn_str, timeout=30)
            except Exception as exc:
                if not self._is_transient(exc) or attempt == _MAX_RETRIES:
                    raise
                time.sleep(_BASE_DELAY * (2 ** attempt) * (1 + random.random() * 0.2))
        raise RuntimeError("unreachable")

    def _blob_cfg(self) -> dict:
        if not hasattr(self, "_blob_config_cache"):
            self._blob_config_cache = _get_blob_config_by_name(self.blob_connector_name)
        return self._blob_config_cache

    # ── Fetch pending PRs ─────────────────────────────────────────────────

    def _resolve_excel_pr_list(self) -> list[str] | None:
        """Try the three Excel-driven input sources, in priority order.

        Returns
        -------
        None
            No Excel-driven input wired/configured — caller falls back to
            _fetch_pending_prs.
        list[str]
            Deduplicated PURCHASE_REQ_NO values to force-reprocess. Empty
            list means an Excel WAS provided but contained zero rows; the
            caller treats that as "nothing to do".
        """
        # Source 1: HandleInput Data wired from a Blob Reader / KB / loader.
        # Accepted shapes — checked in order:
        #   data["pr_numbers"]  → already-parsed list, no Excel parsing needed
        #   data["content"] / ["bytes"] / ["file_bytes"]  → raw .xlsx bytes
        data_in = getattr(self, "pr_list_data", None)
        if data_in is not None:
            payload = getattr(data_in, "data", None) or {}
            if isinstance(payload, dict):
                pr_list = payload.get("pr_numbers") or payload.get("pr_nos")
                if isinstance(pr_list, list) and pr_list:
                    cleaned = [str(p).strip() for p in pr_list if str(p).strip()]
                    self.log(
                        f"Excel PR list — using {len(cleaned)} PR(s) from "
                        f"wired Data.pr_numbers (no parsing needed)"
                    )
                    return cleaned
                for key in ("content", "bytes", "file_bytes"):
                    raw = payload.get(key)
                    if raw:
                        if isinstance(raw, str):
                            try:
                                import base64 as _b64
                                raw = _b64.b64decode(raw)
                            except Exception:
                                raw = raw.encode("latin-1", errors="ignore")
                        try:
                            pr_nos = _parse_pr_list_from_excel_bytes(
                                bytes(raw),
                                column=(self.pr_excel_column or None),
                                sheet=(self.pr_excel_sheet or None),
                            )
                            self.log(
                                f"Excel PR list — parsed {len(pr_nos)} PR(s) "
                                f"from wired Data.{key}"
                            )
                            return pr_nos
                        except Exception as exc:
                            logger.opt(exception=True).error(
                                "Failed to parse Excel from wired Data.{}: {}", key, exc,
                            )
                            raise

        # Source 2: blob path under the configured Blob Connector container.
        blob_path = (getattr(self, "pr_excel_blob_path", None) or "").strip()
        if blob_path:
            self.log(f"Excel PR list — downloading from blob: {blob_path!r}")
            try:
                blob_cfg    = self._blob_cfg()
                excel_bytes = _download_blob(blob_path, blob_cfg)
                pr_nos = _parse_pr_list_from_excel_bytes(
                    excel_bytes,
                    column=(self.pr_excel_column or None),
                    sheet=(self.pr_excel_sheet or None),
                )
                self.log(
                    f"Excel PR list — parsed {len(pr_nos)} PR(s) from {blob_path!r}"
                )
                return pr_nos
            except Exception as exc:
                logger.opt(exception=True).error(
                    "Failed to download/parse Excel from blob {!r}: {}", blob_path, exc,
                )
                raise

        return None

    def _batch_reset_for_reprocess(self, tgt_cs: str, pr_nos: list[str]) -> int:
        """Reset many tracker rows at once for a forced reprocess.

        Equivalent to looping _reset_for_reprocess(pr) for each PR but in
        2 statements per chunk instead of 2 round-trips per PR. Speeds up
        Excel-driven runs by O(N) DB calls.

        Steps per chunk (single transaction):
          1. DELETE ras_pipeline_exceptions for any tracker row in this set
          2. UPDATE ras_tracker SET current_stage_fk=NULL, retry_count+=1

        Returns number of PRs whose tracker row was actually updated. PRs
        without a tracker row are silently skipped (UPDATE rowcount=0)
        and processed as brand-new in _process_pr.
        """
        if not pr_nos:
            return 0
        BATCH = 500
        total_updated = 0
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            for i in range(0, len(pr_nos), BATCH):
                chunk = pr_nos[i : i + BATCH]
                placeholders = ",".join(["?"] * len(chunk))
                cur.execute(
                    f"""
                    DELETE FROM [ras_procurement].[ras_pipeline_exceptions]
                     WHERE ras_tracker_id IN (
                         SELECT ras_uuid_pk
                           FROM [ras_procurement].[ras_tracker]
                          WHERE purchase_req_no IN ({placeholders})
                     )
                    """,
                    *chunk,
                )
                cur.execute(
                    f"""
                    UPDATE [ras_procurement].[ras_tracker]
                       SET current_stage_fk = NULL,
                           retry_count      = COALESCE(retry_count, 0) + 1,
                           updated_at       = SYSUTCDATETIME()
                     WHERE purchase_req_no IN ({placeholders})
                    """,
                    *chunk,
                )
                total_updated += cur.rowcount
            conn.commit()
        except Exception:
            try:
                conn.rollback()
            except Exception:
                pass
            raise
        finally:
            conn.close()
        return total_updated

    def _fetch_current_stages(self, tgt_cs: str, pr_nos: list[str]) -> dict[str, int | None]:
        """Return {pr_no → current_stage_fk} for the given PRs.

        Missing tracker rows are mapped to None. Used by Excel-driven runs to
        skip PRs that already reached COMPLETE (stage 8) without re-fetching
        them per worker.
        """
        if not pr_nos:
            return {}
        # SQL Server caps IN(...) at ~2100 placeholders; chunk in batches of
        # 500 to be well under that limit.
        result: dict[str, int | None] = {pr: None for pr in pr_nos}
        BATCH = 500
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            for i in range(0, len(pr_nos), BATCH):
                chunk = pr_nos[i : i + BATCH]
                placeholders = ",".join(["?"] * len(chunk))
                cur.execute(
                    f"SELECT purchase_req_no, current_stage_fk "
                    f"FROM [ras_procurement].[ras_tracker] "
                    f"WHERE purchase_req_no IN ({placeholders})",
                    *chunk,
                )
                for row in cur.fetchall():
                    pr, stage = row[0], row[1]
                    result[pr] = int(stage) if stage is not None else None
        finally:
            conn.close()
        return result

    def _fetch_pending_prs(self, tgt_cs: str) -> list[str]:
        pr_filter = (self.pr_no_filter or "").strip()
        if pr_filter:
            return [pr_filter]
        limit = int(self.batch_limit)
        conn  = self._connect(tgt_cs)
        cur   = conn.cursor()
        # Stage 1-3 only owns stages 1 and 2.  PRs already at stage 3+ have
        # completed Stage 1-3 and belong to Stage 4-8 — do NOT re-pick them.
        # Only pick: new PRs (no tracker row), reset PRs (NULL stage), or PRs
        # stuck mid-Stage-1-3 (stages 1 or 2).
        _WHERE = """
            FROM [ras_procurement].[purchase_req_mst] prm
            LEFT JOIN [ras_procurement].[ras_tracker] rt
              ON prm.PURCHASE_REQ_NO = rt.purchase_req_no
           WHERE (rt.purchase_req_no IS NULL
                  OR rt.current_stage_fk IS NULL
                  OR rt.current_stage_fk IN (1, 2))
             AND UPPER(prm.PURCHASEFINALAPPROVALSTATUS)
                     IN ('APPROVED BY ALL', 'APPROVED BY ALL EXCEPTION')
           ORDER BY prm.C_DATETIME ASC
        """
        try:
            if limit > 0:
                cur.execute(f"SELECT TOP (?) prm.PURCHASE_REQ_NO {_WHERE}", limit)
            else:
                cur.execute(f"SELECT prm.PURCHASE_REQ_NO {_WHERE}")
            return [row[0] for row in cur.fetchall()]
        finally:
            conn.close()

    # ── Pre-run cleanup ───────────────────────────────────────────────────

    def _cleanup_for_pr(self, tgt_cs: str, pr_no: str) -> None:
        """Delete all prior pipeline output for a PR — runs only if the PR already
        exists in ras_tracker (i.e. it has been processed before). New PRs are
        skipped entirely so their data is never touched."""
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?", pr_no)
        existing = cur.fetchone()
        conn.close()
        if existing is None:
            # Brand-new PR — nothing to clean, proceed straight to Stage 1.
            self.log(f"[{pr_no}] New PR — skipping cleanup")
            return
        self.log(f"[{pr_no}] Existing PR detected — cleaning prior data before reprocessing")

        # Single connection for all DB cleanup work (collect IDs + delete tables)
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        pinecone_ids: list[str] = []
        try:
            # Collect Pinecone vector IDs FIRST, before deleting quotation_extracted_items.
            # We build BOTH formats so stale vectors from any run are cleaned:
            #   • old format: extracted_item_uuid_pk (used by pre-dtl_ code)
            #   • new format: dtl_{purchase_dtl_id}   (current code, stable per line item)
            cur.execute("""
                SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id], qi.[is_selected_quote]
                  FROM [ras_procurement].[quotation_extracted_items] qi
                  JOIN [ras_procurement].[attachment_classification] ac
                    ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
                  JOIN [ras_procurement].[ras_tracker] rt
                    ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
                 WHERE rt.[purchase_req_no] = ?
            """, pr_no)
            rows = cur.fetchall()
            selected_count = sum(1 for r in rows if r[2])
            # New format: only is_selected_quote=1 rows (current code only embeds these)
            new_ids = list({f"dtl_{r[1]}" for r in rows if r[1] and r[2]})
            # Old format: only is_selected_quote=1 UUIDs (legacy vectors only existed for selected rows too)
            old_ids = [str(r[0]) for r in rows if r[0] and r[2]]
            pinecone_ids = list(set(old_ids + new_ids))
            self.log(
                f"[{pr_no}] Pinecone cleanup: {len(rows)} extracted rows, "
                f"{selected_count} is_selected_quote=1 — "
                f"deleting {len(pinecone_ids)} vector(s): "
                f"{len(new_ids)} dtl_ ID(s) + {len(old_ids)} legacy UUID(s)"
            )

            # 1. Null FK back-references in benchmark_result
            cur.execute("""
                UPDATE br SET br.low_hist_item_fk = NULL
                  FROM [ras_procurement].[benchmark_result] br
                  JOIN [ras_procurement].[quotation_extracted_items] qi ON br.low_hist_item_fk = qi.extracted_item_uuid_pk
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            cur.execute("""
                UPDATE br SET br.last_hist_item_fk = NULL
                  FROM [ras_procurement].[benchmark_result] br
                  JOIN [ras_procurement].[quotation_extracted_items] qi ON br.last_hist_item_fk = qi.extracted_item_uuid_pk
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 2. benchmark_result
            cur.execute("""
                DELETE br FROM [ras_procurement].[benchmark_result] br
                  JOIN [ras_procurement].[quotation_extracted_items] qi ON br.extracted_item_uuid_fk = qi.extracted_item_uuid_pk
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 3. quotation_extracted_items
            cur.execute("""
                DELETE qi FROM [ras_procurement].[quotation_extracted_items] qi
                  JOIN [ras_procurement].[attachment_classification] ac ON qi.attachment_classify_fk = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 4. embedded_attachment_classification
            cur.execute("""
                DELETE ec FROM [ras_procurement].[embedded_attachment_classification] ec
                  JOIN [ras_procurement].[attachment_classification] ac ON ec.attachment_classification_id = ac.attachment_classify_uuid_pk
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 5. attachment_classification
            cur.execute("""
                DELETE ac FROM [ras_procurement].[attachment_classification] ac
                  JOIN [ras_procurement].[ras_tracker] rt ON ac.ras_uuid_pk = rt.ras_uuid_pk
                 WHERE rt.purchase_req_no = ?""", pr_no)
            # 6. BI dashboard rows
            cur.execute("""
                DELETE FROM [ras_procurement].[vw_get_ras_data_for_bidashboard]
                 WHERE [PURCHASE_REQ_NO] = ?""", pr_no)
            # 7. Exception records
            cur.execute("""
                DELETE FROM [ras_procurement].[ras_pipeline_exceptions]
                 WHERE ras_tracker_id = (
                     SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker]
                      WHERE purchase_req_no = ?)""", pr_no)
            # 8. Any remaining data via SP
            cur.execute("EXEC [ras_procurement].[usp_cleanup_pr_data] ?", pr_no)
            conn.commit()
        except Exception as exc:
            conn.rollback()
            raise RuntimeError(f"Pre-run cleanup failed for {pr_no}: {exc}") from exc
        finally:
            conn.close()

        # Delete stale Pinecone vectors (non-fatal)
        if pinecone_ids:
            pinecone_index = (getattr(self, "pinecone_index", None) or "").strip()
            pinecone_ns    = (getattr(self, "pinecone_namespace", None) or "").strip()
            if pinecone_index and pinecone_ns:
                try:
                    from agentcore.services.pinecone_service_client import delete_vectors_via_service
                    preview = pinecone_ids[:5]
                    more    = len(pinecone_ids) - 5
                    self.log(
                        f"[{pr_no}] Sending {len(pinecone_ids)} ID(s) to Pinecone delete "
                        f"(index={pinecone_index!r}, namespace={pinecone_ns!r}) — "
                        f"first 5: {preview}" + (f" … +{more} more" if more > 0 else "")
                    )
                    delete_vectors_via_service(
                        index_name=pinecone_index, namespace=pinecone_ns,
                        vector_ids=pinecone_ids,
                    )
                    self.log(f"[{pr_no}] Pinecone delete completed — {len(pinecone_ids)} ID(s) sent for removal")
                except Exception as exc:
                    self.log(f"[{pr_no}] Warning — Pinecone cleanup failed (non-fatal): {exc}")
            else:
                self.log(
                    f"[{pr_no}] Warning — pinecone_index or pinecone_namespace not set; "
                    f"{len(pinecone_ids)} vector(s) were NOT deleted from Pinecone. "
                    f"Check the 'Pinecone Index' and 'Pinecone Namespace' inputs on the component."
                )
        else:
            self.log(f"[{pr_no}] No Pinecone vectors found in DB for this PR — skipping Pinecone delete")

        # Delete Azure Blob folder (non-fatal)
        try:
            self._delete_blob_folder(pr_no)
        except Exception as exc:
            self.log(f"[{pr_no}] Warning — Blob folder cleanup failed (non-fatal): {exc}")

        self.log(f"[{pr_no}] Cleanup complete — all pipeline tables cleared")

    def _reset_for_reprocess(self, tgt_cs: str, pr_no: str) -> None:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("""
                DELETE FROM [ras_procurement].[ras_pipeline_exceptions]
                 WHERE ras_tracker_id = (SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?)
            """, pr_no)
            cur.execute("""
                UPDATE [ras_procurement].[ras_tracker]
                   SET current_stage_fk = NULL, retry_count = COALESCE(retry_count,0)+1, updated_at = SYSUTCDATETIME()
                 WHERE purchase_req_no = ?
            """, pr_no)
            conn.commit()
            self.log(f"[{pr_no}] Tracker reset for reprocess")
        finally:
            conn.close()

    # ── Fetch attachments (Azure SQL for IDs+names, on-prem for binary) ───

    def _fetch_attachments(self, src_cs: str, tgt_cs: str, pr_no: str) -> list[dict]:
        tgt = self._connect(tgt_cs)
        cur = tgt.cursor()
        try:
            cur.execute("""
                SELECT pa.[ATTACHMENT_ID], pa.[FILES_NAME]
                  FROM [ras_procurement].[purchase_req_mst] prm
                  JOIN [ras_procurement].[purchase_attachments] pa ON prm.[PURCHASE_REQ_ID] = pa.[PURCHASE_ID]
                 WHERE prm.[PURCHASE_REQ_NO] = ?
                   AND pa.[ATTACHMENT_ID] IS NOT NULL AND pa.[FILES_NAME] IS NOT NULL
            """, pr_no)
            rows = cur.fetchall()
        finally:
            tgt.close()
        if not rows:
            return []
        id_to_filename = {str(r[0]): r[1] for r in rows}
        att_ids        = list(id_to_filename.keys())
        placeholders   = ",".join(["?"] * len(att_ids))
        src = self._connect(src_cs)
        cur = src.cursor()
        try:
            cur.execute(
                f"SELECT [attachment_id], [doc] FROM [dbo].[ras_attachments]"
                f" WHERE [attachment_id] IN ({placeholders}) AND [doc] IS NOT NULL",
                att_ids,
            )
            return [
                {
                    "attachment_id": str(r[0]),
                    "filename":      id_to_filename.get(str(r[0]), f"attachment_{r[0]}"),
                    "content":       bytes(r[1]),
                    "pr_no":         pr_no,
                }
                for r in cur.fetchall()
            ]
        finally:
            src.close()

    # ── Tracker helpers ───────────────────────────────────────────────────

    def _get_tracker_uuid(self, tgt_cs: str, pr_no: str) -> str:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no = ?", pr_no)
            row = cur.fetchone()
            return str(row[0]) if row else ""
        finally:
            conn.close()

    def _advance_tracker(self, tgt_cs: str, pr_no: str, stage_id: int) -> None:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            if stage_id == _STAGE_INGESTION:
                cur.execute("""
                    MERGE [ras_procurement].[ras_tracker] WITH (HOLDLOCK) AS target
                    USING (
                        SELECT PURCHASE_REQ_NO, PURCHASEFINALAPPROVALSTATUS
                          FROM [ras_procurement].[purchase_req_mst] WHERE PURCHASE_REQ_NO = ?
                    ) AS src ON target.purchase_req_no = src.PURCHASE_REQ_NO
                    WHEN MATCHED THEN UPDATE SET current_stage_fk=?, updated_at=SYSUTCDATETIME()
                    WHEN NOT MATCHED THEN INSERT (purchase_req_no, ras_status, current_stage_fk)
                        VALUES (src.PURCHASE_REQ_NO, src.PURCHASEFINALAPPROVALSTATUS, ?);
                """, pr_no, stage_id, stage_id)
            else:
                cur.execute("""
                    UPDATE [ras_procurement].[ras_tracker]
                       SET current_stage_fk=?, updated_at=SYSUTCDATETIME()
                     WHERE purchase_req_no=?
                """, stage_id, pr_no)
            conn.commit()
        finally:
            conn.close()

    def _record_exception(self, tgt_cs: str, pr_no: str, stage_id: int, error_msg: str) -> None:
        try:
            conn = self._connect(tgt_cs)
            cur  = conn.cursor()
            try:
                # MERGE upsert. Updates current_stage_fk to 99 (EXCEPTION) and
                # also stamps last_processed_at so a PR's terminal time is
                # recorded for both success and failure paths (mirrors what
                # _set_last_processed_at does on Stage 8 success).
                cur.execute("""
                    MERGE [ras_procurement].[ras_tracker] WITH (HOLDLOCK) AS target
                    USING (SELECT PURCHASE_REQ_NO, PURCHASEFINALAPPROVALSTATUS
                             FROM [ras_procurement].[purchase_req_mst] WHERE PURCHASE_REQ_NO = ?
                    ) AS src ON target.purchase_req_no = src.PURCHASE_REQ_NO
                    WHEN MATCHED THEN UPDATE SET
                        current_stage_fk  = 99,
                        last_processed_at = SYSUTCDATETIME(),
                        updated_at        = SYSUTCDATETIME()
                    WHEN NOT MATCHED THEN INSERT
                        (purchase_req_no, ras_status, current_stage_fk, last_processed_at)
                        VALUES (
                            src.PURCHASE_REQ_NO,
                            src.PURCHASEFINALAPPROVALSTATUS,
                            99,
                            SYSUTCDATETIME()
                        );
                """, pr_no)
                cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no=?", pr_no)
                row = cur.fetchone()
                if row:
                    cur.execute("""
                        INSERT INTO [ras_procurement].[ras_pipeline_exceptions] (ras_tracker_id, stage_id, exception_message)
                        VALUES (?, ?, ?)
                    """, row[0], stage_id, error_msg[:4000])
                conn.commit()
            finally:
                conn.close()
        except Exception as exc:
            logger.warning(f"[{pr_no}] Could not write exception record: {exc}")

    # ── attachment_classification upserts ─────────────────────────────────

    def _upsert_parent_classification(self, tgt_cs, pr_no, ras_uuid_pk, att_id, file_path, embedded_flag, embedded_count) -> str:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("""
                MERGE [ras_procurement].[attachment_classification] WITH (HOLDLOCK) AS target
                USING (SELECT ? AS ras_uuid_pk, ? AS attachment_id) AS src
                  ON target.[attachment_id] = src.[attachment_id]
                WHEN MATCHED THEN
                    UPDATE SET [file_path]=?, [embedded_file_flag]=?, [embedded_file_count]=?, [updated_at]=SYSUTCDATETIME()
                WHEN NOT MATCHED THEN
                    INSERT ([ras_uuid_pk],[attachment_id],[file_path],[embedded_file_flag],[embedded_file_count])
                    VALUES (?,?,?,?,?);
            """, ras_uuid_pk, att_id,
                file_path, embedded_flag, embedded_count,
                ras_uuid_pk, att_id, file_path, embedded_flag, embedded_count)
            cur.execute(
                "SELECT [attachment_classify_uuid_pk] FROM [ras_procurement].[attachment_classification] WHERE [attachment_id]=?",
                att_id,
            )
            row = cur.fetchone()
            conn.commit()
            return str(row[0]) if row else ""
        finally:
            conn.close()

    def _upsert_embedded_classification(self, tgt_cs, parent_pk, parent_att_id, file_path) -> None:
        conn = self._connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("""
                MERGE [ras_procurement].[embedded_attachment_classification] WITH (HOLDLOCK) AS target
                USING (SELECT ? AS attachment_classification_id, ? AS file_path) AS src
                  ON target.[attachment_classification_id] = src.[attachment_classification_id]
                 AND target.[file_path] = src.[file_path]
                WHEN MATCHED THEN UPDATE SET [parent_attachment_id]=?, [updated_at]=SYSUTCDATETIME()
                WHEN NOT MATCHED THEN
                    INSERT ([attachment_classification_id],[parent_attachment_id],[file_path])
                    VALUES (?,?,?);
            """, parent_pk, file_path,
                parent_att_id,
                parent_pk, parent_att_id, file_path)
            conn.commit()
        finally:
            conn.close()

    # ── Text extraction ───────────────────────────────────────────────────

    def _extract_text(self, filename: str, raw: bytes) -> str:
        import os
        ext = os.path.splitext(filename.lower())[1]
        try:
            if ext in _TEXT_EXT:  return raw.decode("utf-8", errors="replace")
            if ext in _PDF_EXT:   return _extract_pdf(raw)
            if ext in _EXCEL_EXT: return _extract_excel(raw)
            if ext in _WORD_EXT:  return _extract_word(raw)
            if ext in _PPTX_EXT:  return _extract_pptx(raw)
            if ext in _MSG_EXT:   return _extract_msg_text(raw)
            return f"[Binary — {ext} not text-extractable]"
        except Exception as exc:
            return f"[Extraction error: {exc}]"

    # ── BI dashboard sync ─────────────────────────────────────────────────

    def _sync_bi_dashboard(self, src_cs: str, tgt_cs: str, pr_no: str) -> None:
        """Refresh the BI dashboard row for this PR.

        Reads from on-prem view [dbo].[vw_get_ras_data_for_bidashboard],
        deletes the old Azure row, inserts fresh data.
        Non-fatal — caller wraps in try/except.
        """
        import pyodbc
        # ── 1. read from on-prem ──────────────────────────────────────
        src_conn = self._connect(src_cs)
        try:
            cur = src_conn.cursor()
            cur.execute(
                "SELECT * FROM [dbo].[vw_get_ras_data_for_bidashboard] "
                "WHERE [PURCHASE_REQ_NO] = ?",
                pr_no,
            )
            columns = [d[0] for d in cur.description] if cur.description else []
            rows    = cur.fetchall()
        finally:
            src_conn.close()

        # ── 2. delete old + insert fresh in Azure ─────────────────────
        tgt_conn = self._connect(tgt_cs)
        try:
            tc = tgt_conn.cursor()
            tc.execute(
                "DELETE FROM [ras_procurement].[vw_get_ras_data_for_bidashboard] "
                "WHERE [PURCHASE_REQ_NO] = ?",
                pr_no,
            )
            deleted = tc.rowcount
            if rows and columns:
                col_list     = ", ".join(f"[{c}]" for c in columns)
                placeholders = ", ".join(["?"] * len(columns))
                ins_sql = (
                    f"INSERT INTO [ras_procurement].[vw_get_ras_data_for_bidashboard] "
                    f"({col_list}) VALUES ({placeholders})"
                )
                tc.fast_executemany = True
                tc.executemany(ins_sql, [list(r) for r in rows])
            tgt_conn.commit()
            self.log(
                f"[{pr_no}] BI dashboard synced — "
                f"deleted {deleted}, inserted {len(rows)} row(s)"
            )
        except Exception:
            tgt_conn.rollback()
            raise
        finally:
            tgt_conn.close()

    # ── Blob helpers ──────────────────────────────────────────────────────

    def _container_client(self):
        """Return a cached ContainerClient — credential is initialised once per component instance."""
        if not hasattr(self, "_container_client_cache"):
            from azure.identity import DefaultAzureCredential
            from azure.storage.blob import BlobServiceClient
            cfg        = self._blob_cfg()
            credential = DefaultAzureCredential(
                exclude_environment_credential=True,
                exclude_interactive_browser_credential=True,
            )
            self._container_client_cache = BlobServiceClient(
                account_url=cfg["account_url"],
                credential=credential,
            ).get_container_client(cfg["container_name"])
        return self._container_client_cache

    def _delete_blob_folder(self, pr_no: str) -> None:
        """Delete every blob under procurement/{pr_no}/ in Azure Blob Storage."""
        safe_pr = pr_no.replace("/", "_")
        prefix  = f"procurement/{safe_pr}/"
        cc      = self._container_client()
        blobs   = [b.name for b in cc.list_blobs(name_starts_with=prefix)]
        if not blobs:
            return
        _BATCH = 256
        for i in range(0, len(blobs), _BATCH):
            cc.delete_blobs(*blobs[i : i + _BATCH])
        self.log(f"[{pr_no}] Deleted {len(blobs)} blob(s) from {prefix}")

    def _upload_blob(self, raw: bytes, blob_path: str) -> None:
        self._container_client().get_blob_client(blob_path).upload_blob(raw, overwrite=True)

    # ── Source-change detection ───────────────────────────────────────────

    def _detect_and_requeue_changed_prs(self, tgt_cs: str) -> int:
        """Re-queue completed PRs whose on-prem source data changed since the
        pipeline last processed them.

        Compares purchase_req_mst.U_DATETIME (bumped by on-prem DB triggers
        whenever purchase_req_detail or purchase_attachments rows change) against
        ras_tracker.last_processed_at (stamped at Stage 8 completion).

        Re-queued PRs have current_stage_fk reset to 1 so _fetch_pending_prs
        picks them up on this run. Non-fatal — a failure is logged and skipped.
        """
        try:
            conn = self._connect(tgt_cs)
            cur  = conn.cursor()
            try:
                cur.execute("""
                    SELECT prm.PURCHASE_REQ_NO
                      FROM [ras_procurement].[purchase_req_mst] prm
                      JOIN [ras_procurement].[ras_tracker] rt
                        ON rt.purchase_req_no = prm.PURCHASE_REQ_NO
                     WHERE rt.current_stage_fk  = 8
                       AND rt.last_processed_at IS NOT NULL
                       AND prm.U_DATETIME        > rt.last_processed_at
                       AND UPPER(prm.PURCHASEFINALAPPROVALSTATUS)
                               IN ('APPROVED BY ALL', 'APPROVED BY ALL EXCEPTION')
                """)
                changed = [row[0] for row in cur.fetchall()]
                if not changed:
                    return 0
                for pr_no in changed:
                    cur.execute("""
                        UPDATE [ras_procurement].[ras_tracker]
                           SET current_stage_fk = 1,
                               retry_count      = COALESCE(retry_count, 0) + 1,
                               updated_at       = SYSUTCDATETIME()
                         WHERE purchase_req_no  = ?
                    """, pr_no)
                conn.commit()
                self.log(f"Source-change detection: re-queued {len(changed)} PR(s) — {changed}")
                return len(changed)
            finally:
                conn.close()
        except Exception as exc:
            self.log(f"Warning — source-change detection failed (non-fatal): {exc}")
            return 0

    # ── Full-pipeline continuation (stages 4-8) ──────────────────────────

    def _build_prompts(self) -> dict:
        """Build prompt overrides and tuning-knob dict from component inputs.
        Prompt keys: cls_system, cls_user_text, cls_user_image, ext_system, ext_user, ext_taxonomy.
        Threshold keys: cls_max_chars, ext_max_chars, selected_threshold, price_max_boost, canonicalize_threshold."""
        prompts: dict = {}
        _p = getattr(self, "cls_system_prompt",    None)
        if _p: prompts["cls_system"]    = _get_prompt_text(_p, CLASSIFICATION_SYSTEM_PROMPT)
        _p = getattr(self, "cls_user_text_prompt", None)
        if _p: prompts["cls_user_text"] = _get_prompt_text(_p, _CLASSIFY_USER_TEXT)
        _p = getattr(self, "cls_user_image_prompt", None)
        if _p: prompts["cls_user_image"] = _get_prompt_text(_p, _CLASSIFY_USER_IMAGE)
        _p = getattr(self, "ext_system_prompt",    None)
        if _p: prompts["ext_system"]    = _get_prompt_text(_p, EXTRACTION_SYSTEM_PROMPT)
        _p = getattr(self, "ext_user_template",    None)
        if _p: prompts["ext_user"]      = _get_prompt_text(_p, EXTRACTION_USER_TEMPLATE)
        _p = getattr(self, "ext_item_taxonomy",    None)
        if _p: prompts["ext_taxonomy"]  = _get_prompt_text(_p, ITEM_TAXONOMY)
        # ── Tuning knobs ──────────────────────────────────────────────────
        cls_max = getattr(self, "cls_max_chars", None)
        if cls_max: prompts["cls_max_chars"] = int(cls_max)
        ext_max = getattr(self, "ext_max_chars", None)
        if ext_max: prompts["ext_max_chars"] = int(ext_max)
        ext_max_pages = getattr(self, "ext_max_pages", None)
        if ext_max_pages is not None: prompts["ext_max_pages"] = int(ext_max_pages)
        ext_max_images = getattr(self, "ext_max_images", None)
        if ext_max_images is not None: prompts["ext_max_images"] = int(ext_max_images)
        try:
            prompts["selected_threshold"] = Decimal(str(float(getattr(self, "selected_threshold", None) or "0.70")))
        except Exception:
            prompts["selected_threshold"] = _SELECTED_THRESHOLD
        try:
            prompts["price_max_boost"] = Decimal(str(float(getattr(self, "price_max_boost", None) or "0.10")))
        except Exception:
            prompts["price_max_boost"] = _PRICE_MAX_BOOST
        try:
            prompts["canonicalize_threshold"] = float(getattr(self, "canonicalize_threshold", None) or "0.82")
        except Exception:
            prompts["canonicalize_threshold"] = _CANONICALIZE_THRESHOLD
        # ── LLM retry / rate-limit settings ──────────────────────────────────
        llm_retries = getattr(self, "llm_max_retries", None)
        if llm_retries is not None: prompts["llm_max_retries"]  = int(llm_retries)
        llm_cool    = getattr(self, "llm_retry_cooldown", None)
        if llm_cool is not None:    prompts["llm_retry_cooldown"] = int(llm_cool)
        ext_par     = getattr(self, "ext_parallel_sources", None)
        if ext_par is not None:     prompts["ext_parallel_sources"] = int(ext_par)
        return prompts

    def _run_stages_48(self, pr_no: str, tgt_cs: str, result: dict, prompts: dict) -> None:
        """Run stages 4-8 for a PR that has already completed stages 1-3.
        All Stage 4-8 functions are inlined in this file — no cross-file imports needed.
        Exceptions are caught here and reflected in result["status"]/result["error"]
        so stages 1-3 results are always preserved even if stages 4-8 fail.
        """
        blob_cfg      = self._blob_cfg()
        current_stage = _STAGE_CLASSIFICATION
        top_k         = int(getattr(self, "pinecone_top_k", 5))

        try:
            # Stage 4 — Classification
            self.log(f"[{pr_no}] Stage 4 — classifying attachments…")
            _run_classification(self.llm, tgt_cs, blob_cfg, pr_no, prompts)
            _advance_tracker(tgt_cs, pr_no, _STAGE_CLASSIFICATION)
            self.log(f"[{pr_no}] Stage 4 — classification complete")

            # Stage 5 — Extraction
            # Raises an ExtractionAbortError subclass (NoQuotationFoundError,
            # NoLineItemsError, NoRASContextError, AllExtractionsFailedError)
            # whenever the PR cannot produce extracted items. Caught below as
            # an expected outcome — records ras_pipeline_exceptions, sets
            # tracker.current_stage_fk = 99, and skips stages 6/7/8.
            current_stage = _STAGE_EXTRACTION
            self.log(f"[{pr_no}] Stage 5 — extracting quotation items…")
            n_items = _run_extraction(self.llm, tgt_cs, blob_cfg, pr_no, prompts)
            _advance_tracker(tgt_cs, pr_no, _STAGE_EXTRACTION)
            self.log(f"[{pr_no}] Stage 5 — {n_items} item(s) extracted")

            # Stage 6 — Embeddings
            current_stage = _STAGE_EMBEDDINGS
            _run_embeddings(tgt_cs, pr_no, self.embed_model,
                            self.pinecone_index, self.pinecone_namespace)
            _advance_tracker(tgt_cs, pr_no, _STAGE_EMBEDDINGS)
            self.log(f"[{pr_no}] Stage 6 — embeddings done")

            # Stage 7 — Benchmark
            current_stage = _STAGE_PRICE_BENCHMARK
            _run_benchmark(self.llm, tgt_cs, pr_no, self.embed_model,
                           self.pinecone_index, self.pinecone_namespace, top_k)
            _advance_tracker(tgt_cs, pr_no, _STAGE_PRICE_BENCHMARK)
            self.log(f"[{pr_no}] Stage 7 — benchmark done")

            # Stage 8 — Complete
            _advance_tracker(tgt_cs, pr_no, _STAGE_COMPLETE)
            _set_last_processed_at(tgt_cs, pr_no)
            self.log(f"[{pr_no}] Stage 8 — pipeline complete")
            result["status"] = "complete"

        except ExtractionAbortError as exc:
            # Expected outcome: Stage 5 cannot proceed for a known reason
            # (NoQuotationFoundError, NoLineItemsError, NoRASContextError,
            # AllExtractionsFailedError). Log as warning (no traceback),
            # record in ras_pipeline_exceptions with stage_id=5, set
            # tracker=99 (EXCEPTION), and stop — stages 6/7/8 do not run.
            reason = type(exc).__name__
            self.log(f"[{pr_no}] Stage 5 (Extraction) — {reason}: {exc}")
            logger.warning("[{}] Stage 5 (Extraction) — {}: {}", pr_no, reason, exc)
            result["status"] = f"failed_{reason}"
            result["error"]  = f"Stage 5 (Extraction) — {reason}: {exc}"
            self._record_exception(tgt_cs, pr_no, _STAGE_EXTRACTION, f"{reason}: {exc}")

        except Exception as exc:
            stage_name = {
                _STAGE_CLASSIFICATION:  "Stage 4 (Classification)",
                _STAGE_EXTRACTION:      "Stage 5 (Extraction)",
                _STAGE_EMBEDDINGS:      "Stage 6 (Embeddings)",
                _STAGE_PRICE_BENCHMARK: "Stage 7 (Benchmark)",
                _STAGE_COMPLETE:        "Stage 8 (Complete)",
            }.get(current_stage, f"Stage {current_stage}")
            logger.opt(exception=True).error(
                "[{}] {} failed: {}", pr_no, stage_name, exc
            )
            result["status"] = "failed_stages_48"
            result["error"]  = f"{stage_name}: {exc}"
            self._record_exception(tgt_cs, pr_no, current_stage, f"{stage_name}: {exc}")

    # ── Process single PR ─────────────────────────────────────────────────

    def _process_pr(
        self,
        pr_no: str,
        src_cs: str,
        tgt_cs: str,
        *,
        skip_if_complete: bool = False,
    ) -> dict:
        import os, shutil, tempfile
        result        = {"pr_no": pr_no, "files": [], "status": "failed", "error": ""}
        current_stage = _STAGE_INGESTION
        work_dir      = tempfile.mkdtemp()
        self.log(f"[{pr_no}] Worker started — work_dir={work_dir}")

        # Excel mode: each worker first checks the tracker and short-circuits
        # if the PR is already at stage 8 (COMPLETE) so we do not redo work.
        # This DB read happens on the worker's own connection, so N PRs run
        # the check in *parallel* — no serial pre-pass needed before the pool.
        if skip_if_complete:
            try:
                stage_now = _get_pr_current_stage(tgt_cs, pr_no)
            except Exception as exc:
                # Failure to check stage is non-fatal — we proceed with full
                # processing rather than skipping a PR we shouldn't have.
                self.log(
                    f"[{pr_no}] Stage check failed (will process anyway): {exc}"
                )
                stage_now = -1
            if stage_now == _STAGE_COMPLETE:
                self.log(
                    f"[{pr_no}] Already at stage {_STAGE_COMPLETE} (COMPLETE) "
                    f"— skipping (no cleanup, no reprocess)"
                )
                result["status"] = "already_complete"
                shutil.rmtree(work_dir, ignore_errors=True)
                return result

        try:
            # ── Cleanup at start — only for PRs that already have prior data ──
            # Checked before stage 1: new PR → skipped; existing PR → blobs +
            # Pinecone vectors + quotation_extracted_items + benchmark_result deleted.
            # Explicit handling so cleanup failures are recorded with a clear message
            # before any stage tracker row exists.
            #
            # Cleanup retries on SQL Server deadlock victims (error 1205,
            # SQLSTATE 40001). With multiple parallel workers each running
            # their own multi-DELETE cleanup transaction, lock contention
            # on shared join tables (ras_tracker / attachment_classification)
            # can cause one transaction to be picked as the deadlock victim;
            # SQL Server explicitly asks us to rerun, so we do — with
            # exponential backoff and jitter to spread retries.
            self.log(f"[{pr_no}] Pre-run cleanup check…")
            _CLEANUP_MAX_RETRIES = 3
            _CLEANUP_BASE_DELAY  = 2.0
            cleanup_attempt = 0
            cleanup_exc: Exception | None = None
            while True:
                try:
                    self._cleanup_for_pr(tgt_cs, pr_no)
                    cleanup_exc = None
                    break
                except Exception as exc:
                    cleanup_exc = exc
                    if _is_deadlock_error(exc) and cleanup_attempt < _CLEANUP_MAX_RETRIES:
                        cleanup_attempt += 1
                        delay = _CLEANUP_BASE_DELAY * (2 ** cleanup_attempt) * (1 + random.random() * 0.3)
                        self.log(
                            f"[{pr_no}] Cleanup hit deadlock victim "
                            f"(attempt {cleanup_attempt}/{_CLEANUP_MAX_RETRIES}) "
                            f"— retrying in {delay:.1f}s"
                        )
                        time.sleep(delay)
                        continue
                    break
            if cleanup_exc is not None:
                logger.opt(exception=True).error(
                    "[{}] Pre-run cleanup failed — aborting PR: {}", pr_no, cleanup_exc
                )
                self._record_exception(
                    tgt_cs, pr_no, _STAGE_INGESTION,
                    f"PRE_RUN_CLEANUP failed: {cleanup_exc}",
                )
                result["error"] = str(cleanup_exc)
                return result

            # Stage 1 — INGESTION
            current_stage = _STAGE_INGESTION
            self.log(f"[{pr_no}] Stage 1 — INGESTION starting…")
            self._advance_tracker(tgt_cs, pr_no, _STAGE_INGESTION)

            ras_uuid    = self._get_tracker_uuid(tgt_cs, pr_no)
            self.log(f"[{pr_no}] Stage 1 — tracker uuid={ras_uuid}, fetching on-prem attachments…")
            attachments = self._fetch_attachments(src_cs, tgt_cs, pr_no)
            total_bytes = sum(len(a.get("content", b"")) for a in attachments)
            self.log(
                f"[{pr_no}] Stage 1 — fetched {len(attachments)} attachment(s) "
                f"({total_bytes / 1024:.1f} KB total)"
            )

            # BI dashboard sync: read on-prem view → refresh Azure row.
            # Non-fatal — log and continue if it fails.
            self.log(f"[{pr_no}] Stage 1 — syncing BI dashboard row…")
            try:
                self._sync_bi_dashboard(src_cs, tgt_cs, pr_no)
            except Exception as _bi_exc:
                self.log(f"[{pr_no}] Warning — BI dashboard sync failed (non-fatal): {_bi_exc}")

            if not attachments:
                self.log(f"[{pr_no}] No attachments — skipping (PR has no on-prem files)")
                result["status"] = "skipped"
                return result

            self.log(f"[{pr_no}] Stage 1 — INGESTION complete")

            safe_pr  = pr_no.replace("/", "_")
            extractor = FileExtractor()
            all_files = []   # dicts: filename, content, blob_path, att_id, is_embedded

            self.log(
                f"[{pr_no}] Stage 2 — EMBED_DOC_EXTRACTION starting "
                f"({len(attachments)} parent attachment(s) to scan for embeds)…"
            )

            failed_atts: list[str] = []
            for att_idx, att in enumerate(attachments, 1):
                att_id   = att["attachment_id"]
                att_dir  = os.path.join(work_dir, att_id)
                emb_dir  = os.path.join(att_dir, "extracted")
                try:
                    os.makedirs(att_dir, exist_ok=True)
                    os.makedirs(emb_dir,  exist_ok=True)

                    # Save parent file to temp disk so FileExtractor can open it
                    parent_path = os.path.join(att_dir, att["filename"])
                    with open(parent_path, "wb") as fh:
                        fh.write(att["content"])

                    parent_size_kb = len(att["content"]) / 1024
                    parent_ext     = os.path.splitext(att["filename"])[1].lower() or "<no-ext>"
                    self.log(
                        f"[{pr_no}] [{att_idx}/{len(attachments)}] [{att_id}] "
                        f"{att['filename']} ({parent_size_kb:.1f} KB, {parent_ext}) "
                        f"— scanning for embedded docs"
                    )

                    # If the parent attachment is itself an archive
                    # (.zip / .rar / .7z / .tar / .tar.gz / .tgz / .tar.bz2),
                    # expand its contents directly into emb_dir so each inner
                    # file becomes a separately-classifiable embedded record.
                    # Without this, a zipped quotation is invisible to Stage 4
                    # (the .zip itself is not in _SUPPORTED_CLASSIFY_EXTS).
                    # Mirrors doc-intel pipeline.stages.embed_doc_extraction
                    # pass-1 archive expansion, adapted to the AgentCore flow:
                    # the original archive bytes still serve as the parent
                    # blob, but its contents flow through as embedded files.
                    extractor.parent_prefix = os.path.splitext(att["filename"])[0]
                    if att["filename"].lower().endswith(ARCHIVE_EXTENSIONS):
                        self.log(
                            f"[{pr_no}] [{att_id}] expanding archive "
                            f"{att['filename']!r} into embedded set"
                        )
                        extractor.extract_archive(parent_path, emb_dir)
                        # extract_archive deletes the archive on disk; rewrite
                        # it so the parent blob upload still has the bytes.
                        with open(parent_path, "wb") as fh:
                            fh.write(att["content"])
                        # Then process every newly extracted file that itself
                        # is a supported parent (xlsx/docx/pdf/etc.) for any
                        # nested embeds — matches doc-intel pass-2.
                        for inner_name in sorted(os.listdir(emb_dir)):
                            inner_path = os.path.join(emb_dir, inner_name)
                            if not os.path.isfile(inner_path):
                                continue
                            if inner_name.lower().endswith(SUPPORTED_PARENTS):
                                extractor.parent_prefix = os.path.splitext(inner_name)[0]
                                extractor.process_file(inner_path, emb_dir)
                    else:
                        # Standard parent — extract embedded media/docs.
                        extractor.process_file(parent_path, emb_dir)

                    # Collect embedded file contents
                    embedded = []
                    if os.path.isdir(emb_dir):
                        for emb_name in sorted(os.listdir(emb_dir)):
                            emb_path = os.path.join(emb_dir, emb_name)
                            if os.path.isfile(emb_path):
                                with open(emb_path, "rb") as fh:
                                    emb_content = fh.read()
                                embedded.append({
                                    "filename":    emb_name,
                                    "content":     emb_content,
                                    "blob_path":   f"procurement/{safe_pr}/{att_id}/extracted/{emb_name}",
                                    "att_id":      att_id,
                                    "is_embedded": True,
                                })

                    # Record parent in attachment_classification
                    parent_blob_path = f"procurement/{safe_pr}/{att_id}/{att['filename']}"
                    parent_pk = self._upsert_parent_classification(
                        tgt_cs, pr_no, ras_uuid, att_id,
                        parent_blob_path,
                        embedded_flag  = len(embedded) > 0,
                        embedded_count = len(embedded),
                    )

                    # Record embedded files
                    for emb in embedded:
                        self._upsert_embedded_classification(tgt_cs, parent_pk, att_id, emb["blob_path"])

                    all_files.append({
                        "filename":    att["filename"],
                        "content":     att["content"],
                        "blob_path":   f"procurement/{safe_pr}/{att_id}/{att['filename']}",
                        "att_id":      att_id,
                        "is_embedded": False,
                    })
                    all_files.extend(embedded)
                    self.log(f"[{pr_no}] [{att_id}] {att['filename']} — {len(embedded)} embedded file(s) extracted")

                except Exception as att_exc:
                    # Log and continue — one bad attachment must not block the whole PR.
                    # The attachment is excluded from further processing but the PR continues.
                    logger.opt(exception=True).error(
                        "[{}] Attachment {} ({}) failed during ingestion — skipping: {}",
                        pr_no, att_id, att.get("filename", "?"), att_exc,
                    )
                    failed_atts.append(f"{att_id}({att.get('filename','?')}): {att_exc}")

            if failed_atts:
                self.log(f"[{pr_no}] {len(failed_atts)} attachment(s) skipped due to errors: {'; '.join(failed_atts)}")

            # Stage 2 — EMBED_DOC_EXTRACTION
            current_stage = _STAGE_EMBED_DOC
            parent_count   = sum(1 for f in all_files if not f["is_embedded"])
            embedded_count = sum(1 for f in all_files if f["is_embedded"])
            self._advance_tracker(tgt_cs, pr_no, _STAGE_EMBED_DOC)
            self.log(
                f"[{pr_no}] Stage 2 — EMBED_DOC_EXTRACTION complete "
                f"({parent_count} parent + {embedded_count} embedded = {len(all_files)} total file(s))"
            )

            # Stage 3 — BLOB_UPLOAD
            current_stage = _STAGE_BLOB_UPLOAD
            self.log(f"[{pr_no}] Stage 3 — BLOB_UPLOAD starting ({len(all_files)} file(s) → Azure Blob)…")

            # Upload all files to blob + build text batch
            file_data = []
            uploaded_bytes = 0
            for upload_idx, f_info in enumerate(all_files, 1):
                f_size = len(f_info["content"])
                try:
                    self._upload_blob(f_info["content"], f_info["blob_path"])
                    uploaded_bytes += f_size
                except Exception as upload_exc:
                    logger.opt(exception=True).error(
                        "[{}] Blob upload failed for {!r}: {}",
                        pr_no, f_info["blob_path"], upload_exc,
                    )
                    raise
                # Per-file log only every 10 files (or for big files) to avoid noise
                # on PRs with many small embeds. Always log first/last.
                if (upload_idx % 10 == 0
                        or upload_idx == 1
                        or upload_idx == len(all_files)
                        or f_size > 1_000_000):
                    self.log(
                        f"[{pr_no}] Stage 3 — uploaded {upload_idx}/{len(all_files)}: "
                        f"{f_info['filename']} ({f_size / 1024:.1f} KB)"
                    )
                text      = self._extract_text(f_info["filename"], f_info["content"])
                file_type = _detect_file_type(f_info["filename"])
                extra     = _build_extra_metadata(f_info["filename"], f_info["content"], text)
                if self.max_content_chars and len(text) > int(self.max_content_chars):
                    text = text[: int(self.max_content_chars)]
                file_data.append({
                    "filename":    f_info["filename"],
                    "pr_no":       pr_no,
                    "att_id":      f_info["att_id"],
                    "is_embedded": f_info["is_embedded"],
                    "file_type":   file_type,
                    "extra":       extra,
                    "text":        text,
                    "blob":        f_info["blob_path"],
                })

            self._advance_tracker(tgt_cs, pr_no, _STAGE_BLOB_UPLOAD)
            self.log(
                f"[{pr_no}] Stage 3 — BLOB_UPLOAD complete "
                f"({len(file_data)} file(s), {uploaded_bytes / 1024:.1f} KB total)"
            )

            result["files"]  = file_data
            result["status"] = "success"

            # ── Continue to stages 4-8 in this same worker ──────────────────
            # Each parallel worker processes one PR end-to-end (stages 1→8)
            # without waiting for other PRs to finish Stage 1-3 first.
            self._run_stages_48(pr_no, tgt_cs, result, self._build_prompts())

        except Exception as exc:
            stage_name = {
                _STAGE_INGESTION:   "Stage 1 (Ingestion)",
                _STAGE_EMBED_DOC:   "Stage 2 (Embed Doc Extraction)",
                _STAGE_BLOB_UPLOAD: "Stage 3 (Blob Upload)",
            }.get(current_stage, f"Stage {current_stage}")
            logger.opt(exception=True).error(
                "[{}] {} failed: {}", pr_no, stage_name, exc
            )
            result["error"] = f"{stage_name}: {exc}"
            self._record_exception(tgt_cs, pr_no, current_stage, f"{stage_name}: {exc}")
        finally:
            shutil.rmtree(work_dir, ignore_errors=True)
        return result

    # ── Entry point ───────────────────────────────────────────────────────

    def build_file_batch(self) -> Message:
        if hasattr(self, "_cached_result"):
            return self._cached_result

        src_cs    = self._conn_str(self.source_connection)
        tgt_cs    = self._conn_str(self.target_connection)
        pr_filter = (self.pr_no_filter or "").strip()

        # Source resolution order:
        #   1. pr_no_filter           → single PR, force reprocess
        #   2. pr_list_data /         → Excel-driven (wired Data input or
        #      pr_excel_blob_path        blob path), each PR force-reprocessed
        #   3. _fetch_pending_prs     → standard batch from ras_tracker
        excel_pr_list: list[str] | None = None
        if not pr_filter:
            try:
                excel_pr_list = self._resolve_excel_pr_list()
            except Exception as exc:
                msg_text = f"[Excel PR list could not be read: {exc}]"
                self.log(msg_text)
                msg = Message(text=msg_text)
                self._cached_result = msg
                self._cached_pr_numbers = []
                return msg

        # Detect completed PRs whose on-prem source data changed (U_DATETIME >
        # last_processed_at) and reset them to stage 1 so they are re-processed
        # in this run. Non-fatal — skipped silently on failure.
        # Skipped when an explicit list (single-PR or Excel) is being run —
        # that flow already specifies exactly which PRs to process.
        if not pr_filter and excel_pr_list is None:
            self._detect_and_requeue_changed_prs(tgt_cs)

        if excel_pr_list is not None:
            pr_list = excel_pr_list
            run_mode = "excel"
        else:
            pr_list = self._fetch_pending_prs(tgt_cs)
            run_mode = "single" if pr_filter else "pending"

        if not pr_list:
            text = (
                "[Excel PR list was empty — nothing to process]"
                if run_mode == "excel"
                else "[No pending PRs to process]"
            )
            msg = Message(text=text)
            self._cached_result = msg
            self._cached_pr_numbers = []
            return msg

        if pr_filter:
            # Single-PR mode = explicit force reprocess regardless of state.
            # Reset tracker from exception/any stage back to NULL so _process_pr
            # picks it up cleanly. Actual data cleanup (blobs, Pinecone, DB tables)
            # happens inside _process_pr at the start via _cleanup_for_pr.
            self._reset_for_reprocess(tgt_cs, pr_filter)
            self.log(f"Single-PR reprocess: {pr_filter!r} — tracker reset, cleanup will run at processing start")
        elif run_mode == "excel":
            # Excel mode = workers self-filter inside _process_pr via
            # skip_if_complete=True. No serial pre-pass: stage-8 check
            # runs on each worker's own DB connection in parallel.
            self.log(
                f"Excel-driven run: {len(pr_list)} PR(s) queued — "
                f"each worker will skip if already at stage {_STAGE_COMPLETE}"
            )

        workers = max(1, int(self.parallel_workers))
        self.log(
            f"Processing {len(pr_list)} PR(s) [mode={run_mode}] "
            f"with {workers} parallel worker(s)…"
        )

        # Pre-warm shared caches before threads start so all workers reuse the
        # same already-initialised ContainerClient / blob config rather than
        # each racing to create their own copy.
        self._blob_cfg()
        self._container_client()

        # Worker-side skip is only enabled for Excel mode. Single-PR mode is
        # an explicit force reprocess (must run regardless of stage), and the
        # pending-batch flow already filters stage 8 out at fetch time.
        skip_if_complete = (run_mode == "excel")

        import concurrent.futures
        results: list[dict] = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=workers) as pool:
            futures = {
                pool.submit(
                    self._process_pr, pr, src_cs, tgt_cs,
                    skip_if_complete=skip_if_complete,
                ): pr
                for pr in pr_list
            }
            for f in concurrent.futures.as_completed(futures):
                results.append(f.result())

        parts: list[str] = []
        for r in results:
            for fd in r.get("files", []):
                label = f"{fd['filename']} (PR: {fd['pr_no']}, att_id: {fd['att_id']}"
                label += ", embedded" if fd["is_embedded"] else ", parent"
                label += ")"
                parts.append(
                    f"=== FILE: {label} ===\n"
                    f"============================================================\n"
                    f"FILE METADATA\n"
                    f"============================================================\n"
                    f"Filename: {fd['filename']}\n"
                    f"File Type: {fd['file_type']}\n"
                    f"{fd['extra']}"
                    f"============================================================\n"
                    f"EXTRACTED CONTENT\n"
                    f"============================================================\n"
                    f"{fd['text']}\n"
                    f"============================================================\n\n"
                )

        # Summary log — one line per status bucket. The "already_complete"
        # bucket is populated by workers that found their PR at stage 8 and
        # short-circuited (Excel mode). All buckets come from results since
        # there is no longer a serial pre-pass.
        status_counts: dict[str, int] = {}
        for r in results:
            s = r.get("status", "failed")
            status_counts[s] = status_counts.get(s, 0) + 1
        summary_parts = [f"{v} {k}" for k, v in sorted(status_counts.items())]
        self.log(f"Run complete — {len(results)} PR(s): {', '.join(summary_parts)}")
        for r in results:
            if r.get("error"):
                self.log(f"  [{r['pr_no']}] {r['status']}: {r['error']}")

        batch_text = "".join(parts) if parts else "[No files extracted]"
        msg = Message(text=batch_text)
        self._cached_result = msg
        self._cached_pr_numbers = [
            r["pr_no"] for r in results
            if r.get("status") in ("complete", "success", "skipped")
        ]
        return msg

    def get_processed_prs(self) -> Data:
        """Returns the PR numbers that reached stage 3 in this run.
        Wire this output to Stage 4-8's 'Stage 1-3 Result' input so Stage 4-8
        processes exactly those PRs synchronously after Stage 1-3 finishes."""
        if not hasattr(self, "_cached_pr_numbers"):
            self.build_file_batch()
        return Data(data={"pr_numbers": self._cached_pr_numbers})


# ── Text extraction helpers ───────────────────────────────────────────────

def _extract_pdf(raw: bytes) -> str:
    import io, fitz
    parts: list[str] = []
    with fitz.open(stream=raw, filetype="pdf") as doc:
        for page in doc:
            parts.append(page.get_text())
    return "\n".join(parts)


def _extract_excel(raw: bytes) -> str:
    import io, openpyxl
    wb    = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_word(raw: bytes) -> str:
    import io, docx
    doc = docx.Document(io.BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(raw: bytes) -> str:
    import io
    from pptx import Presentation
    prs   = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_msg_text(raw: bytes) -> str:
    import tempfile, os
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(raw)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject:   parts.append(f"Subject: {msg.subject}")
            if msg.sender:    parts.append(f"From: {msg.sender}")
            if msg.body:      parts.append(msg.body)
            msg.close()
            return "\n".join(parts)
        finally:
            os.unlink(tmp_path)
    except Exception as exc:
        return f"[MSG extraction error: {exc}]"


def _detect_file_type(filename: str) -> str:
    import os
    ext = os.path.splitext(filename.lower())[1]
    return {
        ".pdf": "PDF",
        ".xlsx": "Excel", ".xls": "Excel", ".xlsm": "Excel", ".xlsb": "Excel",
        ".docx": "Word",  ".doc": "Word",
        ".pptx": "PowerPoint", ".ppt": "PowerPoint",
        ".txt": "Text", ".csv": "CSV", ".xml": "XML", ".json": "JSON",
        ".html": "HTML", ".htm": "HTML", ".msg": "Email", ".eml": "Email",
        ".png": "Image", ".jpg": "Image", ".jpeg": "Image",
    }.get(ext, f"Unknown ({ext})")


def _build_extra_metadata(filename: str, raw: bytes, text: str) -> str:
    import os, io
    ext   = os.path.splitext(filename.lower())[1]
    lines: list[str] = []
    if ext == ".pdf":
        try:
            import fitz
            with fitz.open(stream=raw, filetype="pdf") as doc:
                lines.append(f"- page_count: {doc.page_count}")
        except Exception:
            pass
    if ext in {".xlsx", ".xls", ".xlsm", ".xlsb"}:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            lines.append(f"- sheet_count: {len(wb.sheetnames)}")
            lines.append(f"- sheet_names: {', '.join(wb.sheetnames[:10])}")
        except Exception:
            pass
    if ext in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            lines.append(f"- slide_count: {len(Presentation(io.BytesIO(raw)).slides)}")
        except Exception:
            pass
    lines.append(f"- char_count: {len(text)}")
    lines.append(f"- size_bytes: {len(raw)}")
    return ("\n".join(lines) + "\n") if lines else ""

"""Pipeline Stage 4-8: Classification → Extraction → Embeddings → Benchmark → Complete.

Stage 4 — Downloads every parent + embedded file from Azure Blob, classifies
           each using the full file_classifier logic (inlined), updates
           attachment_classification.doc_type + classification_conf.

Stage 5 — Downloads files where doc_type='Quotation', builds RAS context
           from DB, calls LLM for structured extraction, aligns items to
           RAS line items, writes quotation_extracted_items.

Stage 6 — Embeds extracted text into Pinecone via AgentCore service.
Stage 7 — Price benchmark (Pinecone similarity + LLM).
Stage 8 — Marks PR complete in ras_tracker.

Canvas wiring
─────────────
  Azure SQL Connector   ──► target_connection
  LLM node              ──► llm
  Embeddings Model node ──► embed_model
  (Blob connector name is entered in the text field below)
"""
from __future__ import annotations

import json
import random
import re
import time
from dataclasses import dataclass, field
from decimal import Decimal
from difflib import SequenceMatcher
from typing import Optional

from loguru import logger

from agentcore.custom import Node
from agentcore.io import HandleInput, IntInput, MessageTextInput, Output
from agentcore.schema.data import Data
from agentcore.schema.message import Message


# ── Stage IDs ─────────────────────────────────────────────────────────────────
_STAGE_CLASSIFICATION  = 4
_STAGE_EXTRACTION      = 5
_STAGE_EMBEDDINGS      = 6
_STAGE_PRICE_BENCHMARK = 7
_STAGE_COMPLETE        = 8

_MAX_RETRIES = 3
_BASE_DELAY  = 2.0

# ── Classification constants ───────────────────────────────────────────────────
_CLASSIFICATION_MAP = {"E-Auction": "E-Auction Results", "Other": "Others"}
_SUPPORTED_CLASSIFY_EXTS = {
    ".xlsx", ".xls", ".csv", ".pdf", ".docx", ".doc",
    ".pptx", ".ppt", ".txt", ".html", ".htm",
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp",
}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"}
_MAX_CLASSIFY_CHARS = 12000   # ~3500 tokens for GPT-4o

_EXT_TO_TYPE = {
    ".xlsx": "excel",  ".xls": "excel",  ".csv": "csv",
    ".pdf":  "pdf",    ".docx": "word",  ".doc": "legacy_doc",
    ".pptx": "pptx",  ".ppt":  "legacy_doc",
    ".txt":  "text",   ".html": "html",  ".htm": "html",
    ".png":  "image",  ".jpg":  "image", ".jpeg": "image",
    ".tiff": "image",  ".tif":  "image", ".bmp":  "image",
    ".msg":  "msg",
}

VALID_CLASSIFICATIONS = {"RFQ", "Quotation", "MPBC", "BER", "E-Auction", "Other"}

# ── Classification system prompt (full original recovered from file_classifier) ──
CLASSIFICATION_SYSTEM_PROMPT = """You are an expert document classifier for Motherson's procurement team. You will be given the content of a file (extracted text, tables, sheet structure, or an image) and you must classify it into EXACTLY ONE of six categories.

============================================================
CRITICAL: MULTILINGUAL SUPPORT
============================================================
Documents may be in ANY language: English, German, Czech, Hindi, French, Spanish, Chinese, Japanese, Hungarian, and others. You MUST:
- Match fields by MEANING, not by exact English labels
- Recognize translated equivalents. Common examples:
  • German: "Preisspiegel" / "Angebotsvergleich" = Bid Comparison (MPBC), "Angebot" = Quotation/Offer, "Anfrage" = Inquiry/RFQ, "Lieferant" = Supplier, "Preis" = Price, "Lieferzeit" = Delivery Time, "Zahlungsbedingungen" = Payment Terms, "Genehmigung" = Approval, "Begründung" = Justification
  • Czech: "Nabídka" = Offer/Quotation, "Poptávka" = RFQ/Inquiry, "Dodavatel" = Supplier, "Cena" = Price, "Technický požadavek" / "Lastenheft" = Technical Specification
  • Hindi: "कोटेशन" = Quotation, "मूल्य" = Price, "आपूर्तिकर्ता" = Supplier
- Apply the same field-matching logic regardless of language
- If a document is in a language you recognize, translate field names mentally and match against the English category definitions

============================================================
CLASSIFICATION METHOD
============================================================
Classification is FIELD-BASED VALIDATION. Each category has a list of mandatory fields. You must verify the presence of those fields in the document content (allowing for naming/labeling variations, synonyms, AND translations in any language — the underlying meaning matters, not the exact label). The category whose mandatory fields are most completely satisfied wins.

============================================================
CATEGORIES & MANDATORY FIELDS
============================================================

1. **MPBC** — Motherson Purchase BID Comparison
   Official Motherson template that consolidates quotations from multiple suppliers (typically 3+) into a side-by-side comparison for procurement evaluation. Sheet usually titled "MPBC" or "Motherson Purchase BID Comparison". May be a multi-sheet workbook with sheets like "1. MPBC", "2. mandatory cells", "BER", "Supp X risk", "exchange_rates", etc.

   Mandatory fields (yellow-marked in the official MPBC template — must find most of these):
   PROJECT / HEADER:
     • Project Detail
     • RAS Number
     • Sheet No.
     • Indenter, Originator or RFQ Responsible
     • Contact Number
   PER-SUPPLIER (repeated for each of typically 3+ suppliers):
     • Name of the Supplier
     • Preference (Preferred source / Second Source / etc.)
     • Quote meet spec (Yes/No)
     • Supplier's Contact Person
     • Offer No.
     • Offer Date
     • Offer Validity
     • Supplier Country / Classification
     • Country of Origin / Classification
   PRICING:
     • Currency for Comparison
     • Total Amount in supplier currency
     • Total Amount in EUR
     • Saving %age
   SOURCING DECISION:
     • Insourcing
     • Hybrid
     • Single sourcing
     • Supplier justification (BER) signed
     • In case of Cust. funded: Revenue in EUR
     • Target against Budget
   COMMERCIAL TERMS (per supplier):
     • Incoterms
     • Packing and Forwarding
     • Freight
     • Insurance
     • Taxes
     • Customs / duties
     • Installation
     • Delivery Period (Weeks)
     • Advance payment (any payment before delivery to MOTHERSON)
     • SCF (Supply Chain Finance Program)
     • Payment Terms
   APPROVAL / FINAL:
     • GSP Purchase Saving %age
     • GSP Purchase Saving amount in EUR
     • Approved Supplier / Classification
     • Approved Cost (total amount)
     • Landed Cost (for Reference)
     • Approvals (Department / Name / Date / Approver Signature for Purchasing, Technics, Controlling, Sales, Plant Manager)

   Strongest distinguishing signals:
     • Title / sheet name contains "MPBC" or "Motherson Purchase BID Comparison"
     • Supplier 1 / Supplier 2 / Supplier 3 columns appearing side-by-side
     • Presence of RAS Number, GSP (Global Strategic Procurement) terminology
     • Multi-sheet Excel with sheets like "1. MPBC", "2. mandatory cells", "BER", "Supp X risk"
     • MULTILINGUAL: German "Preisspiegel" (price mirror) / "Angebotsvergleich" (offer comparison) / "Bid comparison (Purchased parts)" = MPBC. If 3+ suppliers are compared side-by-side in ANY language with pricing, it is MPBC even without the exact Motherson MPBC template fields like RAS Number.
   Allow naming variations / synonyms / translations; the MEANING of fields matters, not exact labels.

2. **Quotation** — A SINGLE vendor's price offer / response to an RFQ (vendor → buyer such as Motherson)
   Synonyms commonly used as the document title: "Quotation", "Quote", "Estimate", "Offer", "Proposal", "Price Bid". Treat all of these as Quotation candidates.

   Mandatory fields (naming may vary widely; identify by MEANING — Indian, German, and other vendor styles all appear):
     • Vendor Name (vendor's company name appears in letterhead at top AND in signature/footer like "For [Company Name]")
     • Vendor Address (postal address of the vendor)
     • Date of Quotation (e.g., "Date:", "Quotation Dt.", "Offer Date", "Est. Date")
     • Total Amount / Value (grand total / net amount / total in supplier currency)
     • Payment Terms (e.g., "30 days", "100% advance", "Against delivery", "100% after PO confirmation")
     • Delivery Time (e.g., "15-20 days", "within 2-3 weeks", "Delivery Period", "Lead Time", "Dispatch Time")
     • Validity (e.g., "Quotation valid until ___", "Offer Validity: 60 days", "Valid for one month")
     • Specifications (technical / product / scope specifications — may be a separate "TECH SPEC" sheet in Excel)
     • Item Description (line items with descriptions — columns like "Description of Goods", "Item Description", "Sr. No. + Description + Qty + Rate + Amount", "BOQ items")

   Additional supporting signals frequently observed (not all needed, but strengthen confidence):
     • Vendor letterhead at top: company logo / name + address + phone + email + GST No. / GSTIN / PAN No. / VAT No.
     • A reference / quotation number (e.g., "Quotation No.: PPAQ004560", "Ref: Q/875", "Offer No.", "Est. No.: SCL/Haryana/20-21/21", "Ref: MIPL/MATE/...")
     • Addressed to: "To, M/s [Customer Name]", "Client:", "Kind Attn: Mr./Ms. ___"
     • Formal letter language: "We are pleased to offer / quote", "With reference to your enquiry", "Thanking you", "Yours faithfully"
     • Closing: "For [Vendor Company]" + "Authorised Signatory" / "Proprietor" / "Sales Manager"
     • Tax columns: HSN/SAC Code, CGST %, SGST %, IGST %, Excise Duty, GST @ 18 %
     • Bank details / E. & O. E. / "Please mention our quotation number on your purchase order"

   CRITICAL DISAMBIGUATION:
     • A Quotation is from ONE vendor. If THREE OR MORE vendor names appear as parallel columns/sections being compared → it is MPBC, not Quotation.
     • A Quotation MAY be a multi-sheet Excel (e.g., "COMMERCIAL", "TECH. SPEC.", "Summary", "NPV", "Vendors" sheets) but with only ONE vendor. Do not misclassify these as MPBC — check vendor count, not sheet count.
     • If price columns are EMPTY (template for vendor to fill) → it is RFQ, not Quotation.
     • Filename is unreliable — e.g., "(875) CK Motherson Auto Hitech.docx" was issued BY Hi-Tech to CK Motherson; the customer name in the filename does not make it MPBC or RFQ.
     • PDF extraction may produce duplicated characters from font issues (e.g., "TTOO" instead of "TO", "QQUUOOTTAATTIIOONN" instead of "QUOTATION") — interpret semantically.
     • **A Quotation can use the buyer's RFQ TEMPLATE FORMAT.** When a supplier fills in an RFQ specification template with their technical responses AND populates pricing fields with actual monetary values, the document has become a Quotation. Key differentiator: if ACTUAL PRICES are filled in (not blank placeholders) and a supplier contact person / company name is prominently featured, treat it as Quotation even if the original RFQ structure is retained.

3. **RFQ** — Request For Quotation (Motherson → vendors)
   A specification / scope document issued BY Motherson (the buyer) TO vendors, asking them to submit a quotation.

   Mandatory fields:
     • Project Name
     • RFQ Number (reference / inquiry number)
     • Date (issue date)
     • Specifications — detailed technical / scope specifications structured as tables with:
       - Sl. No / Item number
       - Description / technical parameter
       - "Required" / "Not Required" / "Std" flags
       - "To be specified by the supplier" / "Supplier Spec/Confirmation" / "Pls specify" columns (blank or for vendor to fill)
     • Commercials section asking for pricing breakdown (blank fields for vendors to fill)

   Strong supporting signals:
     • Sheet name contains "RFQ"
     • "MATE-B Req" or "MATE Spec" column
     • "Supplier Spec/Confirmation" or "Supplier Remarks" columns (empty = template)
     • "Techno Commercial Comparison" as title
     • Columns with "Fill this column by '0' if STD or Enter the cost if Optional"
     • "Company Name:", "Contact Person:", "Telephone:", "Email:" fields for vendor to fill

   CRITICAL DISAMBIGUATION:
     • Even if vendors have filled some spec responses, it remains RFQ IF pricing fields are EMPTY and no supplier is prominently featured.
     • If a supplier has populated actual pricing (USD/EUR/INR values) AND supplier contact details appear prominently → classify as Quotation, not RFQ.
     • If 3+ vendors compared side-by-side in consolidated evaluation → MPBC, not RFQ.

4. **BER** — Bid Exception Report
   SPECIFICALLY the Motherson "BID EXCEPTION REPORT" template form.

   Mandatory fields (ALL must be present or nearly all — strict template match):
     • "BID EXCEPTION REPORT" appearing explicitly as a header / title (REQUIRED — generic "Waiver of Competition" or "Single Source Justification" is NOT sufficient)
     • Reasoning for not obtaining at least three bids/quotes
     • Order Value field (e.g., "Order Value: 671,57 €")
     • Description of the product or service to be ordered
     • Justification for waiver of competitive bidding

   Strong supporting signals:
     • Header "Capital Equipment & Indirect Purchasing"
     • Budget Line Ref.
     • Reference to "LCC Suppliers" / "Low Cost Country"
     • Checkbox options A through E (A: less than three potential bidders; B: sole-source; C: national supply contract; D: similar item purchased in past 6 months; E: Other)
     • Three approval rows: "Prepared by:" + "Purchasing approval:" + "Managing Director / COO / EVP approval:"
     • Sheet named "BER" in an Excel workbook

   CRITICAL: Do NOT classify as BER if the document is a generic waiver form, Single Source Justification from a non-Motherson template, or any exception document that does NOT use the specific Motherson BER template with the A-E checkbox structure. Classify those as "Other".

5. **E-Auction** — E-Auction results / reports / trackers
   Documents generated from or summarizing an online reverse-auction event.

   Mandatory fields:
     • Event ID
     • Event Name
     • Publish Date / Open Date / Close Date
     • BID Id
     • BID Status (e.g., "Accepted", "Default")
     • Participant (vendor / bidder name)
     • Basic Price (per unit price)
     • Extended Price (total price for volume)

   Strong supporting signals:
     • "eAuction", "e-Auction", "Reverse Auction", "Japanese Auction" in title/headers
     • Sheet names "Overview Sheet", "Full Bid Data Sheet"
     • Rank column; Savings column; Capacity Planning Volume
     • Pricing tiers: "Price 1" / "Price 2" / "Price 3" with savings definitions
     • Tracker workbooks with sheets "Pivot", "Summary", "Project Details", "Project Negotiation"
     • Presentation slides with "eAuction Overview", "eAuction Status", monthly savings summaries

   IMPORTANT — Two tiers of E-Auction documents (BOTH classify as E-Auction):
     TIER 1 (raw auction output): Contains most mandatory fields (Event ID, BID Id, Participant, prices). High confidence.
     TIER 2 (summaries/trackers/presentations): May NOT contain individual event-level fields like Event ID or BID Id, but IS clearly ABOUT e-auction results — contains "eAuction" keyword prominently + pricing data (Price 1/2/3, savings, L1 price, Final price) + auction metadata (Auction type, Auction Month, GSP Buyer). Classify as E-Auction with moderate confidence (0.75-0.85) even if individual bid-level fields are missing.

   CRITICAL DISAMBIGUATION:
     • E-Auction focuses on AUCTION EVENT and BIDDING PROCESS with event metadata and time-sequenced bids. MPBC is a static comparison table — fundamentally different.
     • E-Auction tracker/summary workbooks and .pptx presentations summarizing eAuction results are still E-Auction, not "Other".
     • If the document is ABOUT eAuction (mentions "eAuction" + savings/prices), classify as E-Auction even if not all mandatory fields are present.

6. **Other** — Anything that does NOT satisfy the mandatory field set of any category above
   Examples: invoices, purchase orders, delivery notes, contracts, drawings, internal memos, generic emails, MSAs, NDAs, quality reports, etc.

============================================================
DECISION PROCESS (follow strictly)
============================================================
Step 1 — Scan the document and identify which mandatory fields from each category are PRESENT (allowing synonyms / paraphrases / equivalent column names).
Step 2 — For each candidate category, compute a coverage ratio: (fields present) / (total mandatory fields).
Step 3 — Pick the category with the highest coverage. To classify as that category, coverage must be ≥ 60% AND the strongest distinguishing signal must be present (e.g., for MPBC: ≥ 3 vendors compared; for BER: explicit BER title + waiver justification; for E-Auction: event + bid columns).
Step 4 — If no category reaches 60% coverage, classify as "Other".
Step 5 — Filename can be a hint but content always wins. Ignore filename if content contradicts it.
Step 6 — Confidence calibration:
   • ≥ 0.90 → all or nearly all mandatory fields present and unambiguous
   • 0.75 – 0.89 → most mandatory fields present; minor ambiguity
   • 0.60 – 0.74 → enough fields present to classify but with notable gaps
   • < 0.60 → genuine uncertainty — likely "Other"

============================================================
OUTPUT FORMAT (strict JSON, no markdown, no commentary)
============================================================
{
  "classification": "RFQ" | "Quotation" | "MPBC" | "BER" | "E-Auction" | "Other",
  "confidence": <float 0.0 - 1.0>,
  "reason": "<2-3 sentences explaining which mandatory fields you matched and which were missing>",
  "key_signals": ["<mandatory field matched 1>", "<mandatory field matched 2>", "<mandatory field matched 3>"],
  "fields_matched": ["<exact field/column/phrase observed in the document>", ...],
  "fields_missing": ["<mandatory fields for the chosen category that were NOT found>", ...]
}

Rules:
- Pick exactly ONE category.
- key_signals = the top 3-5 mandatory fields that drove the decision.
- fields_matched = up to 10 specific evidence items (column names, headers, phrases) you actually observed.
- fields_missing = mandatory fields for the chosen category that you could NOT find.
- Be specific in the reason — cite real evidence from the document, do not be vague."""

# ── Classification user prompt templates ──────────────────────────────────────
_CLASSIFY_USER_TEXT = """Classify the following file by checking its content against the mandatory fields for each category.

============================================================
FILE METADATA
============================================================
Filename: {filename}
File Type: {file_type}
{extra_metadata}
============================================================
EXTRACTED CONTENT
============================================================
{extracted_content}
============================================================

Apply the field-based decision process from the system prompt and return only the JSON object."""

_CLASSIFY_USER_IMAGE = """Classify the following file based on the image provided.

============================================================
FILE METADATA
============================================================
Filename: {filename}
File Type: {file_type}
{extra_metadata}
============================================================

Examine the image carefully — read every visible field, column, and label. Check the mandatory fields for each category against what you see. Apply the decision process from the system prompt and return only the JSON object."""

# ── Extraction prompt constants ────────────────────────────────────────────────
EXTRACTION_SYSTEM_PROMPT = """You are a senior procurement analyst with deep experience evaluating supplier quotations across every spend category. Your job here is to read each quotation thoroughly and produce a complete, accurate, decision-grade extraction in a strict JSON schema. The downstream system uses your output to benchmark prices, pick winning suppliers, and approve purchase requisitions, so completeness and correctness directly affect business decisions.

Key responsibilities:
- Read the entire quotation before answering. Do not stop at the first item table.
- Extract every item line, pricing, supplier information, and commercial terms.
- Match each extracted item to the right purchase requisition DTL_ID.
- Translate all extracted text to English regardless of the source language.
- Return data in the exact JSON schema specified — no markdown fences, no commentary.
- Handle diverse document formats: formal quotations, proforma invoices, price lists, rate cards, cost estimates, email quotations, and scanned documents.

Extraction guidelines:
- Be precise with numbers: prices, quantities, dates. Never hallucinate or infer figures not present.
- Distinguish between unit price and total price. total_price = unit_price × quantity unless stated otherwise.
- Identify currency from the document (₹, $, €, £, AED, ZAR or ISO codes). Return ISO-4217 code.
- Extract payment terms verbatim then normalise.
- Set supplier_match_conf honestly. 0.0 = unrelated, 1.0 = identical. A lower honest score is always better.
- If a field genuinely cannot be determined, set it to null. Never guess."""

ITEM_TAXONOMY = """Guidelines for item_level_1 through item_level_8 hierarchical taxonomy:

Level 1 — Broad industry/domain: "Electronics", "Mechanical", "IT Hardware", "Services", "Consumables", "Raw Materials"
Level 2 — Sub-category: "Industrial Equipment", "Office Equipment", "Software Licensing", "Fleet Management"
Level 3 — Product/service type: "Temperature Controller", "Laptop", "Annual Maintenance Contract", "Vehicle Hiring"
Level 4 — Brand/Manufacturer: "Dell", "Siemens", "Bosch", "Toyota"
Level 5 — Model/Series/Part: "Latitude 5540", "MCLX-350A-0", "PowerEdge R740"
Level 6 — Configuration/Variant: "16 GB RAM / 512 GB SSD", "3-phase 440 V 50 Hz"
Level 7 — Additional spec: "With touchscreen", "IP65 rated", "CE certified"
Level 8 — Remaining detail: "Custom colour RAL 7035", "Extended warranty 5 yr"

Rules: item_level_1/2/3 MUST always be filled. Levels 4-8 may be null when information is absent."""

EXTRACTION_USER_TEMPLATE = """## Purchase Requisition Context

### Header Information

| Field | Value |
|---|---|
| RAS Number | {purchase_req_no} |
| Requisition ID | {purchase_req_id} |
| RAS Title | {ras_title} |
| Requisition Type | {requisition_type} |
| Classification | {classification} |
| Justification | {justification} |
| Primary Supplier | {supplier_name} |
| Supplier Address | {address} |
| Parent Supplier | {parent_supplier} |
| Supplier Type | {supplier_type} |
| Supplier Country | {supplier_country} |
| Currency (RAS) | {currency} |
| Purchase Value | {purchase_value} |
| Enquiry No | {enquiry_no} |
| Contract No | {contract_no} |
| Order No | {order_no} |
| Department | {department} |
| Negotiated By | {negotiated_by} |
| Category Buyer | {category_buyer} |
| Purchase Category | {purchase_category} |
| Category L1 | {category} |
| Category L2 | {sub_category} |
| Category L3 | {l3} |
| Category L4 | {l4} |
| Site | {site} |
| Region / Country | {site_region} / {site_country} |
| Division | {division} |
| Payment Days | {payment_days} |
| PO Date | {po_date} |

### Line Items from the Requisition

{line_items_table}

---

### Additional RAS Reference Data

{raw_ras_context}

---

### Item Taxonomy Guidelines

{item_taxonomy}

---

## Quotation Document

{document_content}

---

## Required Output

Extract exactly one item per DTL_ID. Return a **single JSON object** — no markdown fences, no extra text:

{{
  "supplier_name": "string or null",
  "supplier_address": "string or null",
  "supplier_country": "string or null",
  "quotation_ref_no": "string or null",
  "quotation_date": "YYYY-MM-DD or null",
  "currency": "ISO-4217 three-letter code or null",
  "validity_date": "YYYY-MM-DD or null",
  "validity_days": "integer or null",
  "payment_terms": "string or null",
  "items": [
    {{
      "purchase_dtl_id": "integer from DTL_ID column if matched, else null",
      "supplier_match_conf": "0.0 to 1.0",
      "item_name": "canonical item name in English",
      "item_description": "full description with all specs in English",
      "quantity": "number or null",
      "unit": "unit of measurement or null",
      "unit_price": "number or null",
      "total_price": "number or null",
      "discount": "number or null",
      "taxation_details": "string or null",
      "delivery_date": "YYYY-MM-DD or null",
      "delivery_time_days": "integer or null",
      "item_level_1": "broadest category (REQUIRED)",
      "item_level_2": "sub-category (REQUIRED)",
      "item_level_3": "product/service type (REQUIRED)",
      "item_level_4": "brand if known, else null",
      "item_level_5": "model if known, else null",
      "item_level_6": "configuration if known, else null",
      "item_level_7": "key spec if known, else null",
      "item_level_8": "additional detail if known, else null",
      "commodity_tag": "lowercase-slug-tag",
      "item_summary": "plain-English summary max 20 words"
    }}
  ]
}}"""


# ── Dataclasses ────────────────────────────────────────────────────────────────

@dataclass
class LineItemContext:
    purchase_dtl_id:  int
    purchase_req_id:  int
    item_no:          int
    quantity:         Optional[Decimal]
    item_type:        Optional[str]
    item_description: Optional[str]
    unit_price:       Optional[Decimal] = None
    uom:              Optional[str]     = None
    supplier_name:    Optional[str]     = None
    discount:         Optional[Decimal] = None
    req_value:        Optional[Decimal] = None
    currency:         Optional[str]     = None
    delivery_date:    Optional[str]     = None
    payment_details:  Optional[str]     = None
    original_value:   Optional[Decimal] = None
    initial_offer:    Optional[Decimal] = None
    negotiation:      Optional[Decimal] = None
    comments:         Optional[str]     = None
    prepayment:       Optional[str]     = None
    item_code:        Optional[str]     = None


@dataclass
class RASContext:
    purchase_req_no:   str
    purchase_req_id:   int
    supplier_name:     Optional[str]
    justification:     Optional[str]
    currency:          Optional[str]
    enquiry_no:        Optional[str]     = None
    classification:    Optional[str]     = None
    department:        Optional[str]     = None
    negotiated_by:     Optional[str]     = None
    address:           Optional[str]     = None
    contract_no:       Optional[str]     = None
    order_no:          Optional[str]     = None
    purchase_value:    Optional[Decimal] = None
    category:          Optional[str]     = None
    sub_category:      Optional[str]     = None
    site_country:      Optional[str]     = None
    site_region:       Optional[str]     = None
    site:              Optional[str]     = None
    division:          Optional[str]     = None
    requisition_type:  Optional[str]     = None
    parent_supplier:   Optional[str]     = None
    supplier_type:     Optional[str]     = None
    supplier_country:  Optional[str]     = None
    payment_days:      Optional[str]     = None
    po_date:           Optional[str]     = None
    category_buyer:    Optional[str]     = None
    l3:                Optional[str]     = None
    l4:                Optional[str]     = None
    purchase_category: Optional[str]     = None
    ras_title:         Optional[str]     = None
    line_items:        list              = field(default_factory=list)
    raw_mst:           dict             = field(default_factory=dict)
    raw_dtl_rows:      list             = field(default_factory=list)
    raw_vw_rows:       list             = field(default_factory=list)


@dataclass
class DocumentContent:
    text:        Optional[str]      = None
    images:      Optional[list]     = None
    source_path: str                = ""
    page_count:  int                = 0
    ocr_source:  bool               = False

    @property
    def is_image_based(self) -> bool:
        return bool(self.images)


# ── Connection helpers ─────────────────────────────────────────────────────────

def _conn_str(conn_data: Data) -> str:
    d      = conn_data.data or {}
    driver = d.get("driver", "ODBC Driver 18 for SQL Server")
    server = d.get("host", d.get("server", ""))
    port   = d.get("port", 1433)
    db     = d.get("database_name", d.get("database", ""))
    user   = d.get("username", "")
    pwd    = d.get("password", "")
    return (
        f"DRIVER={{{driver}}};SERVER={server},{port};"
        f"DATABASE={db};UID={user};PWD={pwd};TrustServerCertificate=yes;"
    )


def _is_transient(exc: Exception) -> bool:
    kw = ["connection reset", "timeout", "throttl", "resource limit",
          "broken pipe", "transport-level", "login failed"]
    return any(k in str(exc).lower() for k in kw)


def _connect(cs: str):
    import pyodbc
    for attempt in range(_MAX_RETRIES + 1):
        try:
            return pyodbc.connect(cs, timeout=30)
        except Exception as exc:
            if not _is_transient(exc) or attempt == _MAX_RETRIES:
                raise
            time.sleep(_BASE_DELAY * (2 ** attempt) * (1 + random.random() * 0.2))
    raise RuntimeError("unreachable")


# ── Blob helpers ───────────────────────────────────────────────────────────────

def _get_blob_config_by_name(connector_name: str) -> dict:
    name = (connector_name or "").strip()
    if not name:
        raise ValueError("blob_connector_name is empty.")
    try:
        import asyncio
        import concurrent.futures as _cf
        from sqlalchemy import select

        async def _fetch():
            from agentcore.services.deps import get_db_service
            from agentcore.services.database.models.connector_catalogue.model import ConnectorCatalogue
            db_service = get_db_service()
            async with db_service.with_session() as session:
                stmt = (
                    select(ConnectorCatalogue)
                    .where(ConnectorCatalogue.name == name)
                    .where(ConnectorCatalogue.provider == "azure_blob")
                )
                result = await session.execute(stmt)
                row = result.scalars().first()
                if row is None:
                    raise ValueError(f"No azure_blob connector named {name!r}.")
                cfg         = row.provider_config or {}
                account_url = (cfg.get("account_url") or row.host or "").strip()
                container   = (cfg.get("container_name") or row.database_name or "").strip()
                from urllib.parse import urlparse
                if not urlparse(account_url).netloc:
                    raise ValueError(f"Invalid Storage Account URL: {account_url!r}")
                if not container:
                    raise ValueError(f"No container_name in connector {name!r}.")
                return {"account_url": account_url, "container_name": container}

        try:
            asyncio.get_running_loop()
            with _cf.ThreadPoolExecutor() as pool:
                return pool.submit(asyncio.run, _fetch()).result(timeout=10)
        except RuntimeError:
            return asyncio.run(_fetch())
    except Exception as exc:
        raise RuntimeError(f"Blob connector lookup failed: {exc}") from exc


def _download_blob(blob_path: str, blob_cfg: dict) -> bytes:
    from azure.identity import DefaultAzureCredential
    from azure.storage.blob import BlobServiceClient
    credential = DefaultAzureCredential(
        exclude_environment_credential=True,
        exclude_interactive_browser_credential=True,
    )
    client = BlobServiceClient(
        account_url=blob_cfg["account_url"], credential=credential
    )
    blob = client.get_blob_client(container=blob_cfg["container_name"], blob=blob_path)
    return blob.download_blob().readall()


# ── File type detection ────────────────────────────────────────────────────────

def _detect_file_type(filename: str) -> str:
    import os
    ext = os.path.splitext(filename.lower())[1]
    return _EXT_TO_TYPE.get(ext, "unknown")


# ── Content extraction for classification ─────────────────────────────────────

def _extract_for_classification(file_bytes: bytes, filename: str) -> tuple:
    """Returns (text_content: str, image_b64: str|None, metadata_str: str)."""
    import os, io
    ext      = os.path.splitext(filename.lower())[1]
    ftype    = _EXT_TO_TYPE.get(ext, "unknown")
    meta_str = ""

    if ftype == "excel":
        return _extract_excel_classify(file_bytes, filename, meta_str)
    if ftype == "csv":
        return _extract_csv_classify(file_bytes, meta_str)
    if ftype == "pdf":
        return _extract_pdf_classify(file_bytes, filename, meta_str)
    if ftype == "word":
        return _extract_word_classify(file_bytes, filename, meta_str)
    if ftype == "pptx":
        return _extract_pptx_classify(file_bytes, filename, meta_str)
    if ftype == "legacy_doc":
        return _extract_fitz_classify(file_bytes, filename, meta_str)
    if ftype == "image":
        return _extract_image_classify(file_bytes, filename, meta_str)
    if ftype == "text":
        return _extract_text_classify(file_bytes, meta_str)
    if ftype == "html":
        return _extract_html_classify(file_bytes, meta_str)
    if ftype == "msg":
        return _extract_msg_classify(file_bytes, filename, meta_str)
    return "[Unsupported file type]", None, meta_str


def _extract_excel_classify(file_bytes, filename, meta_str):
    import io, pandas as pd
    _MAX_SHEETS = 8
    buf  = io.BytesIO(file_bytes)
    parts: list[str] = []
    try:
        for engine in ("openpyxl", "xlrd"):
            try:
                buf.seek(0)
                xls = pd.ExcelFile(buf, engine=engine)
                break
            except Exception:
                continue
        else:
            return "[Excel could not be opened]", None, meta_str

        sheets = xls.sheet_names
        meta_str = f"- sheet_count: {len(sheets)}\n- sheet_names: {sheets}\n- multi_sheet: {len(sheets) > 1}\n"
        parts.append(f"## Workbook Structure\nTotal Sheets: {len(sheets)}\nSheet Names: {sheets}\n")
        for idx, sheet in enumerate(sheets):
            try:
                df = pd.read_excel(xls, sheet_name=sheet, nrows=40, header=None)
                df = df.iloc[:, :30].dropna(how="all").dropna(axis=1, how="all")
                non_empty = len(df)
                if idx < _MAX_SHEETS and non_empty > 0:
                    parts.append(f"### Sheet {idx+1}: '{sheet}'")
                    parts.append(f"Non-empty rows in sample: {non_empty}")
                    try:
                        hdrs = [f"col_{c}" for c in range(df.shape[1])]
                        parts.append(df.fillna("").astype(str).to_markdown(index=False, headers=hdrs))
                    except Exception:
                        parts.append(df.fillna("").astype(str).to_string(index=False))
                    parts.append("")
                else:
                    parts.append(f"### Sheet {idx+1}: '{sheet}' (summary only — {non_empty} non-empty rows)\n")
            except Exception:
                pass
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Excel extraction error: {exc}]", None, meta_str


def _extract_csv_classify(file_bytes, meta_str):
    import io, pandas as pd
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), nrows=40, header=None).iloc[:, :30]
        try:
            text = df.fillna("").astype(str).to_markdown(index=False)
        except Exception:
            text = df.fillna("").astype(str).to_string(index=False)
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[CSV extraction error: {exc}]", None, meta_str


def _extract_pdf_classify(file_bytes, filename, meta_str):
    import io
    _MAX_PAGES = 10
    try:
        import pdfplumber
        buf  = io.BytesIO(file_bytes)
        with pdfplumber.open(buf) as pdf:
            total = len(pdf.pages)
            meta_str = f"- total_pages: {total}\n- pages_processed: {min(total, _MAX_PAGES)}\n"
            if total == 0:
                return "[Empty or corrupt PDF — no pages found]", None, meta_str

            # smart selection: first N-2 pages + last 2 pages (captures pricing summaries at end)
            if total <= _MAX_PAGES:
                page_idxs = list(range(total))
            else:
                page_idxs = sorted(set(list(range(_MAX_PAGES - 2)) + list(range(total - 2, total))))

            text_parts: list[str] = []
            table_parts: list[str] = []
            for i in page_idxs:
                page = pdf.pages[i]
                t = (page.extract_text() or "").strip()
                if t:
                    text_parts.append(f"--- Page {i+1} ---\n{t}")
                # extract tables from this page
                for t_idx, table in enumerate(page.extract_tables() or []):
                    if not table:
                        continue
                    header = table[0]
                    rows   = table[1:]
                    tstr   = " | ".join(str(c) for c in header) + "\n"
                    tstr  += " | ".join("---" for _ in header) + "\n"
                    for row in rows[:20]:
                        tstr += " | ".join(str(c) for c in row) + "\n"
                    table_parts.append(f"Table {t_idx+1} (Page {i+1}):\n{tstr}")

            combined = "\n\n".join(text_parts)
            if table_parts:
                combined += "\n\n### Extracted Tables\n" + "\n\n".join(table_parts)

            if len(combined.strip()) < 50:
                return _extract_pdf_as_image_classify(file_bytes, filename, meta_str, total)
            return combined[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception:
        return _extract_pdf_fitz_classify(file_bytes, filename, meta_str)


def _extract_pdf_fitz_classify(file_bytes, filename, meta_str):
    import io, base64
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total = len(doc)
        meta_str = f"- total_pages: {total}\n"
        pages = list(range(min(total, 8)))
        texts = [doc[i].get_text("text").strip() for i in pages]
        combined = "\n\n".join(f"--- Page {i+1} ---\n{t}" for i, t in zip(pages, texts) if t)
        if len(combined.strip()) < 50:
            pix = doc[0].get_pixmap(dpi=150)
            b64 = base64.b64encode(pix.tobytes("png")).decode()
            doc.close()
            return "[Scanned PDF - content sent as image]", b64, meta_str
        doc.close()
        return combined[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[PDF extraction error: {exc}]", None, meta_str


def _extract_pdf_as_image_classify(file_bytes, filename, meta_str, total):
    import io, base64
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            img = pdf.pages[0].to_image(resolution=200)
            buf = io.BytesIO()
            img.original.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode()
        return "[Scanned PDF - content sent as image]", b64, meta_str
    except Exception:
        return _extract_pdf_fitz_classify(file_bytes, filename, meta_str)


def _extract_word_classify(file_bytes, filename, meta_str):
    import io, zipfile, base64
    try:
        from docx import Document
        from PIL import Image
        doc   = Document(io.BytesIO(file_bytes))
        parts: list[str] = []
        para_count = 0
        for para in doc.paragraphs:
            if para_count >= 200:
                break
            text = para.text.strip()
            if not text:
                continue
            if para.style and para.style.name.startswith("Heading"):
                lvl = para.style.name.replace("Heading ", "").strip()
                try:
                    lvl = int(lvl)
                except ValueError:
                    lvl = 1
                parts.append(f"{'#' * lvl} {text}")
            else:
                parts.append(text)
            para_count += 1
        for i, tbl in enumerate(doc.tables):
            rows = []
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(cells)
            if rows:
                hdr = rows[0]
                tstr  = " | ".join(hdr) + "\n"
                tstr += " | ".join("---" for _ in hdr) + "\n"
                for row in rows[1:20]:
                    tstr += " | ".join(row) + "\n"
                parts.append(f"\n### Table {i+1}\n{tstr}")
        text = "\n\n".join(parts)
        meta_str = f"- paragraphs: {len(doc.paragraphs)}\n- tables: {len(doc.tables)}\n"
        if len(text.strip()) < 50:
            # try to extract largest embedded image from word/media/
            try:
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                    imgs = sorted([n for n in z.namelist() if n.startswith("word/media/")])
                    if imgs:
                        largest = max(imgs, key=lambda n: z.getinfo(n).file_size)
                        img_bytes = z.read(largest)
                        img = Image.open(io.BytesIO(img_bytes))
                        if img.mode not in ("RGB", "L"):
                            img = img.convert("RGB")
                        if max(img.size) > 2048:
                            img.thumbnail((2048, 2048), Image.LANCZOS)
                        buf2 = io.BytesIO()
                        img.save(buf2, format="PNG")
                        b64 = base64.b64encode(buf2.getvalue()).decode()
                        return "[Word document is image-based - content sent as image for visual analysis]", b64, meta_str
            except Exception:
                pass
            return _extract_fitz_classify(file_bytes, filename, meta_str)
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Word extraction error: {exc}]", None, meta_str


def _extract_pptx_classify(file_bytes, filename, meta_str):
    import io
    _MAX_SLIDES = 15
    try:
        from pptx import Presentation
        buf = io.BytesIO(file_bytes)
        prs = Presentation(buf)
        slides = list(prs.slides)
        total  = len(slides)
        parts  = [f"## Presentation: {total} slides\n"]
        for i, slide in enumerate(slides[:_MAX_SLIDES]):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_parts.append(t)
                if shape.has_table:
                    tbl = shape.table
                    rows = [[c.text.strip() for c in row.cells] for row in tbl.rows]
                    if rows:
                        hdr   = rows[0]
                        tstr  = " | ".join(hdr) + "\n"
                        tstr += " | ".join("---" for _ in hdr) + "\n"
                        for row in rows[1:20]:
                            tstr += " | ".join(row) + "\n"
                        slide_parts.append(tstr)
            if slide_parts:
                parts.append(f"### Slide {i+1}")
                parts.append("\n".join(slide_parts))
                parts.append("")
        meta_str = f"- total_slides: {total}\n- slides_extracted: {min(total, _MAX_SLIDES)}\n"
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return _extract_fitz_classify(file_bytes, filename, meta_str)


def _extract_fitz_classify(file_bytes, filename, meta_str):
    import base64
    try:
        import fitz
        doc = fitz.open(stream=file_bytes)
        total = len(doc)
        pages = list(range(min(total, 4)))
        images: list[str] = []
        for i in pages:
            pix = doc[i].get_pixmap(dpi=150)
            images.append(base64.b64encode(pix.tobytes("png")).decode())
        doc.close()
        if not images:
            return "[Document could not be rendered]", None, meta_str
        return "[Legacy document - content sent as images]", images[0], meta_str
    except Exception as exc:
        return f"[Document extraction error: {exc}]", None, meta_str


def _extract_image_classify(file_bytes, filename, meta_str):
    import io, base64
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        if max(img.size) > 2048:
            img.thumbnail((2048, 2048))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode()
        meta_str = f"- size: {img.size[0]}x{img.size[1]}\n"
        return "[Image file - content sent as image]", b64, meta_str
    except Exception as exc:
        return f"[Image extraction error: {exc}]", None, meta_str


def _extract_text_classify(file_bytes, meta_str):
    try:
        import chardet
        enc = (chardet.detect(file_bytes).get("encoding") or "utf-8")
        try:
            text = file_bytes.decode(enc)
        except Exception:
            text = file_bytes.decode("utf-8", errors="replace")
        return text[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[Text extraction error: {exc}]", None, meta_str


def _extract_html_classify(file_bytes, meta_str):
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(file_bytes, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
    except Exception:
        raw  = file_bytes.decode("utf-8", errors="replace")
        text = re.sub(r"<[^>]+>", " ", raw)
        text = re.sub(r"\s+", " ", text).strip()
    return text[:_MAX_CLASSIFY_CHARS], None, meta_str


def _extract_msg_classify(file_bytes, filename, meta_str):
    import os, tempfile
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject: parts.append(f"Subject: {msg.subject}")
            if msg.sender:  parts.append(f"From: {msg.sender}")
            if msg.body:    parts.append(msg.body)
            msg.close()
        finally:
            os.unlink(tmp_path)
        return "\n".join(parts)[:_MAX_CLASSIFY_CHARS], None, meta_str
    except Exception as exc:
        return f"[MSG extraction error: {exc}]", None, meta_str


# ── Trivial image heuristic ────────────────────────────────────────────────────

def _is_trivial_image(file_bytes: bytes, filename: str) -> tuple:
    import os
    if os.path.splitext(filename.lower())[1] not in _IMAGE_EXTS:
        return False, ""
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        img.load()
        w, h = img.size
        long_e  = max(w, h)
        short_e = max(min(w, h), 1)
        if long_e < 200:           return True, f"tiny {w}x{h}"
        if short_e < 150:          return True, f"short edge {w}x{h}"
        if long_e / short_e > 5:   return True, f"extreme aspect {w}x{h}"
        if long_e / short_e > 3 and short_e < 250:
            return True, f"banner {w}x{h}"
        if long_e < 400 and short_e < 200:
            return True, f"sub-thumbnail {w}x{h}"
        gray = img.convert("L")
        hist = gray.histogram()
        total = sum(hist)
        if total and sum(hist[230:]) / total > 0.97:
            return True, "mostly white"
        try:
            colors = img.convert("RGB").getcolors(maxcolors=256)
            if colors is not None and len(colors) < 16:
                return True, f"flat palette ({len(colors)} colors)"
        except Exception:
            pass
    except Exception as exc:
        return True, f"unreadable ({exc.__class__.__name__})"
    return False, ""


# ── Classification ─────────────────────────────────────────────────────────────

def _get_prompt_text(msg_input, default: str) -> str:
    """Extract text from an optional wired Message input; fall back to default."""
    if msg_input is None:
        return default
    text = getattr(msg_input, "text", None)
    if text and str(text).strip():
        return str(text).strip()
    return default


def _classify_file(llm, file_bytes: bytes, filename: str, prompts: dict | None = None) -> tuple:
    """Returns (doc_type: str, confidence: float)."""
    import os
    from langchain_core.messages import HumanMessage, SystemMessage

    p = prompts or {}
    sys_prompt      = p.get("cls_system", CLASSIFICATION_SYSTEM_PROMPT)
    user_text_tmpl  = p.get("cls_user_text", _CLASSIFY_USER_TEXT)
    user_image_tmpl = p.get("cls_user_image", _CLASSIFY_USER_IMAGE)

    ext = os.path.splitext(filename.lower())[1]
    if ext not in _SUPPORTED_CLASSIFY_EXTS:
        return "Others", 0.0

    trivial, reason = _is_trivial_image(file_bytes, filename)
    if trivial:
        logger.info(f"Trivial image {filename!r} ({reason}) → Others")
        return "Others", 0.0

    file_type = _EXT_TO_TYPE.get(ext, "unknown")
    text_content, image_b64, meta_str = _extract_for_classification(file_bytes, filename)

    try:
        if image_b64:
            user_prompt = user_image_tmpl.format(
                filename=filename, file_type=file_type, extra_metadata=meta_str
            )
            content = [
                {"type": "text", "text": user_prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}", "detail": "high"}},
            ]
            messages = [SystemMessage(content=sys_prompt), HumanMessage(content=content)]
        else:
            user_prompt = user_text_tmpl.format(
                filename=filename, file_type=file_type,
                extra_metadata=meta_str, extracted_content=text_content,
            )
            messages = [SystemMessage(content=sys_prompt), HumanMessage(content=user_prompt)]

        response = llm.invoke(messages)
        raw = (getattr(response, "content", None) or str(response)).strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        result = json.loads(raw)

        raw_cls = result.get("classification", "Other")
        if raw_cls not in VALID_CLASSIFICATIONS:
            raw_cls = "Other"
        confidence = float(result.get("confidence", 0.0))
        doc_type   = _CLASSIFICATION_MAP.get(raw_cls, raw_cls)
        return doc_type, confidence

    except Exception as exc:
        logger.warning(f"Classification failed for {filename!r}: {exc}")
        return "Others", 0.0


# ── Stage 4: run classification for a PR ──────────────────────────────────────

def _run_classification(llm, tgt_cs: str, blob_cfg: dict, pr_no: str, prompts: dict | None = None) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT ac.[attachment_classify_uuid_pk], ac.[file_path], ac.[attachment_id]
              FROM [ras_procurement].[attachment_classification] ac
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
        """, pr_no)
        parent_rows = cur.fetchall()

        cur.execute("""
            SELECT ec.[embedded_attachment_classification_id], ec.[file_path],
                   ac.[attachment_classify_uuid_pk]
              FROM [ras_procurement].[embedded_attachment_classification] ec
              JOIN [ras_procurement].[attachment_classification] ac
                ON ec.[attachment_classification_id] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
        """, pr_no)
        embedded_rows = cur.fetchall()
    finally:
        conn.close()

    # Classify parent attachments
    for row in parent_rows:
        att_pk, blob_path, att_id = str(row[0]), row[1], row[2]
        if not blob_path:
            continue
        try:
            file_bytes = _download_blob(blob_path, blob_cfg)
            import os
            filename   = os.path.basename(blob_path)
            doc_type, conf = _classify_file(llm, file_bytes, filename, prompts)
            _update_parent_classification(tgt_cs, att_id, doc_type, conf)
            logger.info(f"[{pr_no}] Parent {filename!r}: {doc_type} (conf={conf:.2f})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Parent classification failed ({blob_path!r}): {exc}")

    # Classify embedded files
    for row in embedded_rows:
        emb_pk, blob_path, parent_pk = str(row[0]), row[1], str(row[2])
        if not blob_path:
            continue
        try:
            file_bytes = _download_blob(blob_path, blob_cfg)
            import os
            filename   = os.path.basename(blob_path)
            doc_type, conf = _classify_file(llm, file_bytes, filename, prompts)
            _update_embedded_classification(tgt_cs, parent_pk, blob_path, doc_type, conf)
            logger.info(f"[{pr_no}] Embedded {filename!r}: {doc_type} (conf={conf:.2f})")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Embedded classification failed ({blob_path!r}): {exc}")


def _update_parent_classification(tgt_cs: str, att_id: str, doc_type: str, conf: float) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[attachment_classification]
               SET [doc_type]            = ?,
                   [classification_conf] = ?,
                   [updated_at]          = SYSUTCDATETIME()
             WHERE [attachment_id] = ?
        """, doc_type, conf, att_id)
        conn.commit()
    finally:
        conn.close()


def _update_embedded_classification(tgt_cs: str, parent_pk: str, blob_path: str, doc_type: str, conf: float) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[embedded_attachment_classification]
               SET [doc_type]            = ?,
                   [classification_conf] = ?,
                   [updated_at]          = SYSUTCDATETIME()
             WHERE [attachment_classification_id] = ?
               AND [file_path] = ?
        """, doc_type, conf, parent_pk, blob_path)
        conn.commit()
    finally:
        conn.close()


# ── RAS context builder ────────────────────────────────────────────────────────

def _build_ras_context(tgt_cs: str, pr_no: str) -> Optional[RASContext]:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT TOP 1
                   prm.[PURCHASE_REQ_ID], prm.[SUPPLIER_NAME], prm.[JUSTIFICATION],
                   prm.[CURRENCY], prm.[ENQUIRY_NO], prm.[CLASSIFICATION],
                   prm.[Department], prm.[NEGOTIATED_BY], prm.[ADDRESS],
                   prm.[CONTRACT_NO], prm.[ORDER_NO], prm.[PURCHASE_VALUE]
              FROM [ras_procurement].[purchase_req_mst] prm
             WHERE prm.[PURCHASE_REQ_NO] = ?
        """, pr_no)
        mst = cur.fetchone()
        if not mst:
            return None

        req_id = mst[0]
        def _d(v): return Decimal(str(v)) if v is not None else None

        cur.execute("""
            SELECT prd.[PURCHASE_DTL_ID], prd.[PURCHASE_REQ_ID], prd.[ITEM_NO],
                   prd.[QUANTITY], prd.[ITEM_TYPE], prd.[ITEMDESCRIPTION],
                   prd.[PRICE], prd.[UOM], prd.[DISCOUNT], prd.[REQ_VALUE],
                   prd.[CURRENCY], prd.[DELIVERY_DATE], prd.[SUPPLIER_NAME],
                   prd.[PAYMENT_DETAILS], prd.[ORIGINAL_VALUE], prd.[Initial_Offer],
                   prd.[Negotiation], prd.[CommentsforItem], prd.[PREPAYMENT], prd.[ITEM_CODE]
              FROM [ras_procurement].[purchase_req_detail] prd
             WHERE prd.[PURCHASE_REQ_ID] = ?
             ORDER BY prd.[ITEM_NO]
        """, req_id)
        dtl_rows = cur.fetchall()

        line_items: list[LineItemContext] = []
        for r in dtl_rows:
            raw_del = r[11]
            del_str = (
                raw_del.date().isoformat() if hasattr(raw_del, "date")
                else str(raw_del) if raw_del else None
            )
            line_items.append(LineItemContext(
                purchase_dtl_id=r[0], purchase_req_id=r[1], item_no=r[2] or 0,
                quantity=_d(r[3]), item_type=r[4], item_description=r[5],
                unit_price=_d(r[6]), uom=r[7], discount=_d(r[8]),
                req_value=_d(r[9]), currency=r[10], delivery_date=del_str,
                supplier_name=r[12] or mst[1], payment_details=r[13],
                original_value=_d(r[14]), initial_offer=_d(r[15]),
                negotiation=_d(r[16]), comments=r[17], prepayment=r[18], item_code=r[19],
            ))

        # BI dashboard enrichment (best-effort)
        cat=sub_cat=site_c=site_r=div=site=req_t=l3=l4=pur_cat=ras_t=par_s=sup_t=sup_c=pay_d=po_d=cat_b=None
        try:
            cur.execute("""
                SELECT TOP 1 vw.[L1], vw.[Sub_Category_Type], vw.[Site_Country],
                             vw.[Site_Region], vw.[Division], vw.[L3], vw.[L4],
                             vw.[Purchase_Category], vw.[Title], vw.[Site],
                             vw.[Requisition_Type], vw.[Parent_Supplier], vw.[Supplier_Type],
                             vw.[Suplier_country], vw.[Payment_Days], vw.[PO_Date],
                             vw.[Category_Buyer], vw.[L2]
                  FROM vw_get_ras_data_for_bidashboard vw
                 WHERE vw.[PURCHASE_REQ_ID] = ?
            """, req_id)
            vw = cur.fetchone()
            if vw:
                cat,sub_cat,site_c,site_r,div,l3,l4 = vw[0],vw[1]or vw[17],vw[2],vw[3],vw[4],vw[5],vw[6]
                pur_cat,ras_t,site,req_t = vw[7],vw[8],vw[9],vw[10]
                par_s,sup_t,sup_c = vw[11],vw[12],vw[13]
                pay_d,po_d,cat_b  = vw[14],vw[15],vw[16]
        except Exception:
            pass

        return RASContext(
            purchase_req_no=pr_no, purchase_req_id=req_id,
            supplier_name=mst[1], justification=mst[2], currency=mst[3],
            enquiry_no=mst[4], classification=mst[5], department=mst[6],
            negotiated_by=mst[7], address=mst[8], contract_no=mst[9],
            order_no=mst[10], purchase_value=_d(mst[11]),
            category=cat, sub_category=sub_cat, site_country=site_c,
            site_region=site_r, site=site, division=div, requisition_type=req_t,
            parent_supplier=par_s, supplier_type=sup_t, supplier_country=sup_c,
            payment_days=pay_d, po_date=po_d, category_buyer=cat_b,
            l3=l3, l4=l4, purchase_category=pur_cat, ras_title=ras_t,
            line_items=line_items,
        )
    finally:
        conn.close()


# ── Document loader for extraction ─────────────────────────────────────────────

def _load_document(file_bytes: bytes, filename: str, max_pages: int = 20) -> DocumentContent:
    import io, os
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pdf":        return _load_pdf_for_extract(file_bytes, max_pages)
    if ext in (".xlsx", ".xls"): return _load_spreadsheet_for_extract(file_bytes, ext)
    if ext == ".docx":       return _load_docx_for_extract(file_bytes)
    if ext == ".doc":        return _load_doc_legacy_for_extract(file_bytes, max_pages)
    if ext in (".pptx", ".ppt"): return _load_pptx_for_extract(file_bytes, max_pages)
    if ext in _IMAGE_EXTS:   return _load_image_for_extract(file_bytes, ext)
    if ext in (".txt", ".csv"): return DocumentContent(
        text=file_bytes.decode("utf-8", errors="replace")[:50000], page_count=1
    )
    if ext == ".msg":        return _load_msg_for_extract(file_bytes)
    # Fallback: fitz render
    return _fitz_render_for_extract(file_bytes, max_pages)


def _load_pdf_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    """Render every page as a 200-dpi image — matches original pipeline behaviour.
    PDFs are always sent to the LLM as vision (images), never as extracted text,
    because scanned/mixed PDFs lose critical layout/table information when OCR'd."""
    import base64
    try:
        import fitz
        doc    = fitz.open(stream=file_bytes, filetype="pdf")
        total  = len(doc)
        n      = min(total, max_pages)
        images = []
        for i in range(n):
            pix = doc[i].get_pixmap(dpi=200)
            images.append(base64.b64encode(pix.tobytes("png")).decode())
        doc.close()
        return DocumentContent(images=images, page_count=n)
    except Exception as exc:
        return DocumentContent(text=f"[PDF error: {exc}]", page_count=0)


def _load_spreadsheet_for_extract(file_bytes: bytes, ext: str) -> DocumentContent:
    import io
    parts: list[str] = []
    try:
        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            for ws in wb.worksheets:
                lines = [f"### Sheet: {ws.title}"]
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        lines.append("| " + " | ".join(cells) + " |")
                if len(lines) > 1:
                    parts.append("\n".join(lines))
            wb.close()
        else:
            import xlrd
            wb = xlrd.open_workbook(file_contents=file_bytes)
            for sheet in wb.sheets():
                lines = [f"### Sheet: {sheet.name}"]
                for rx in range(sheet.nrows):
                    cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
                    if any(c.strip() for c in cells):
                        lines.append("| " + " | ".join(cells) + " |")
                if len(lines) > 1:
                    parts.append("\n".join(lines))
        text = "\n\n".join(parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, 20)
    except Exception as exc:
        return DocumentContent(text=f"[Spreadsheet error: {exc}]", page_count=0)


def _load_docx_for_extract(file_bytes: bytes) -> DocumentContent:
    import io
    try:
        from docx import Document
        doc   = Document(io.BytesIO(file_bytes))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        table_parts: list[str] = []
        for tbl in doc.tables:
            rows = ["| " + " | ".join(c.text.strip() for c in row.cells) + " |"
                    for row in tbl.rows]
            if rows:
                table_parts.append("\n".join(rows))
        text = "\n".join(paras)
        if table_parts:
            text += "\n\n" + "\n\n".join(table_parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, 20)
    except Exception as exc:
        return DocumentContent(text=f"[DOCX error: {exc}]", page_count=0)


def _load_doc_legacy_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    try:
        return _fitz_render_for_extract(file_bytes, max_pages)
    except Exception:
        pass
    try:
        import io, olefile
        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        if ole.exists("WordDocument"):
            import re as _re
            stream = ole.openstream("WordDocument").read()
            text = stream.decode("utf-8", errors="ignore")
            text = _re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
            if text.strip():
                return DocumentContent(text=text, page_count=1)
        ole.close()
    except Exception:
        pass
    return DocumentContent(text="[Document could not be read]")


def _load_pptx_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    import io
    try:
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(file_bytes))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_t: list[str] = [f"--- Slide {idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_t.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        slide_t.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
            if len(slide_t) > 1:
                parts.append("\n".join(slide_t))
        text = "\n\n".join(parts)
        if text.strip():
            return DocumentContent(text=text, page_count=1)
        return _fitz_render_for_extract(file_bytes, max_pages)
    except Exception:
        return _fitz_render_for_extract(file_bytes, max_pages)


def _load_image_for_extract(file_bytes: bytes, ext: str) -> DocumentContent:
    import io, base64
    raw = file_bytes
    if ext in (".tif", ".tiff"):
        try:
            from PIL import Image
            buf = io.BytesIO()
            Image.open(io.BytesIO(raw)).save(buf, format="PNG")
            raw = buf.getvalue()
        except Exception:
            pass
    return DocumentContent(images=[base64.b64encode(raw).decode()], page_count=1)


def _load_msg_for_extract(file_bytes: bytes) -> DocumentContent:
    import os, tempfile
    try:
        import extract_msg
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            msg   = extract_msg.Message(tmp_path)
            parts = []
            if msg.subject: parts.append(f"Subject: {msg.subject}")
            if msg.body:    parts.append(msg.body)
            msg.close()
        finally:
            os.unlink(tmp_path)
        return DocumentContent(text="\n\n".join(parts), page_count=1)
    except Exception as exc:
        return DocumentContent(text=f"[MSG error: {exc}]")


def _fitz_render_for_extract(file_bytes: bytes, max_pages: int) -> DocumentContent:
    import base64
    try:
        import fitz
        doc    = fitz.open(stream=file_bytes)
        n      = min(len(doc), max_pages)
        images = []
        for i in range(n):
            pix = doc[i].get_pixmap(dpi=200)
            images.append(base64.b64encode(pix.tobytes("png")).decode())
        doc.close()
        return DocumentContent(images=images, page_count=n)
    except Exception as exc:
        return DocumentContent(text=f"[Render error: {exc}]")


# ── Quotation source resolver ──────────────────────────────────────────────────

def _resolve_quotation_sources(tgt_cs: str, pr_no: str) -> list:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    sources: list[dict] = []
    try:
        cur.execute("""
            SELECT ac.[attachment_classify_uuid_pk], ac.[file_path], ac.[attachment_id]
              FROM [ras_procurement].[attachment_classification] ac
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND ac.[doc_type] = 'Quotation'
        """, pr_no)
        for row in cur.fetchall():
            if row[1]:
                sources.append({
                    "blob_path": row[1],
                    "attachment_classify_fk": str(row[0]),
                    "embedded_classify_fk": None,
                    "attachment_id": str(row[2]),
                })

        cur.execute("""
            SELECT ec.[embedded_attachment_classification_id], ec.[file_path],
                   ec.[parent_attachment_id], ac.[attachment_classify_uuid_pk]
              FROM [ras_procurement].[embedded_attachment_classification] ec
              JOIN [ras_procurement].[attachment_classification] ac
                ON ec.[attachment_classification_id] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND ec.[doc_type] = 'Quotation'
        """, pr_no)
        for row in cur.fetchall():
            if row[1]:
                sources.append({
                    "blob_path": row[1],
                    "attachment_classify_fk": str(row[3]),
                    "embedded_classify_fk": str(row[0]),
                    "attachment_id": str(row[2]),
                })
    finally:
        conn.close()
    return sources


# ── Line items table builder ───────────────────────────────────────────────────

def _build_line_items_table(ctx: RASContext) -> str:
    def _na(v): return str(v) if v is not None else "N/A"
    header = ("| DTL_ID | Item No | Item Code | Description | Qty | UOM | Type "
              "| Unit Price | Req Value | Currency | Supplier | Delivery Date |")
    sep    = ("|--------|---------|-----------|-------------|-----|-----|------"
              "|------------|-----------|----------|----------|---------------|")
    rows   = [header, sep]
    for li in ctx.line_items:
        rows.append(
            f"| {_na(li.purchase_dtl_id)} | {_na(li.item_no)} | {_na(li.item_code)} "
            f"| {_na(li.item_description)} | {_na(li.quantity)} | {_na(li.uom)} "
            f"| {_na(li.item_type)} | {_na(li.unit_price)} | {_na(li.req_value)} "
            f"| {_na(li.currency)} | {_na(li.supplier_name)} | {_na(li.delivery_date)} |"
        )
    return "\n".join(rows)


def _build_raw_context(ctx: RASContext) -> str:
    parts: list[str] = []
    if ctx.raw_mst:
        lines = [f"  {k}: {v}" for k, v in ctx.raw_mst.items()]
        parts.append("#### purchase_req_mst (header)\n" + "\n".join(lines))
    if ctx.raw_dtl_rows:
        cols   = list(ctx.raw_dtl_rows[0].keys())
        header = "| " + " | ".join(cols) + " |"
        sep    = "|" + "|".join("---" for _ in cols) + "|"
        rows   = ["| " + " | ".join(str(r.get(c, "")) for c in cols) + " |"
                  for r in ctx.raw_dtl_rows]
        parts.append("#### purchase_req_detail\n" + "\n".join([header, sep] + rows))
    return "\n\n".join(parts)


# ── Extraction LLM call ────────────────────────────────────────────────────────

def _build_extraction_user_prompt(ctx: RASContext, doc: DocumentContent, prompts: dict | None = None) -> str:
    def _f(v): return str(v) if v is not None else "N/A"
    if doc.ocr_source and doc.text:
        doc_content_str = f"[OCR markdown from Azure Document Intelligence]\n\n{doc.text}"
    elif doc.images and doc.text:
        doc_content_str = f"[Extracted text — page images attached below]\n\n{doc.text}"
    elif doc.images:
        doc_content_str = "[Scanned document — page image(s) attached]"
    else:
        doc_content_str = doc.text or "[No content extracted]"

    user_tmpl = (prompts or {}).get("ext_user", EXTRACTION_USER_TEMPLATE)
    return user_tmpl.format(
        purchase_req_no=_f(ctx.purchase_req_no),
        purchase_req_id=_f(ctx.purchase_req_id),
        justification=_f(ctx.justification),
        supplier_name=_f(ctx.supplier_name),
        currency=_f(ctx.currency),
        enquiry_no=_f(ctx.enquiry_no),
        classification=_f(ctx.classification),
        department=_f(ctx.department),
        negotiated_by=_f(ctx.negotiated_by),
        address=_f(ctx.address),
        contract_no=_f(ctx.contract_no),
        order_no=_f(ctx.order_no),
        purchase_value=_f(ctx.purchase_value),
        category=_f(ctx.category),
        sub_category=_f(ctx.sub_category),
        l3=_f(ctx.l3),
        l4=_f(ctx.l4),
        purchase_category=_f(ctx.purchase_category),
        ras_title=_f(ctx.ras_title),
        site_region=_f(ctx.site_region),
        site_country=_f(ctx.site_country),
        site=_f(ctx.site),
        division=_f(ctx.division),
        requisition_type=_f(ctx.requisition_type),
        parent_supplier=_f(ctx.parent_supplier),
        supplier_type=_f(ctx.supplier_type),
        supplier_country=_f(ctx.supplier_country),
        payment_days=_f(ctx.payment_days),
        po_date=_f(ctx.po_date),
        category_buyer=_f(ctx.category_buyer),
        line_items_table=_build_line_items_table(ctx),
        item_taxonomy=(prompts or {}).get("ext_taxonomy", ITEM_TAXONOMY),
        document_content=doc_content_str,
        raw_ras_context=_build_raw_context(ctx),
    )


def _call_extraction_llm(llm, ctx: RASContext, doc: DocumentContent, prompts: dict | None = None) -> str:
    from langchain_core.messages import HumanMessage, SystemMessage
    user_prompt = _build_extraction_user_prompt(ctx, doc, prompts)
    sys_prompt  = (prompts or {}).get("ext_system", EXTRACTION_SYSTEM_PROMPT)
    messages: list = [SystemMessage(content=sys_prompt)]
    if doc.is_image_based and doc.images:
        images = doc.images[:50]
        content_parts: list = [{"type": "text", "text": user_prompt}]
        for b64 in images:
            content_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"},
            })
        messages.append(HumanMessage(content=content_parts))
    else:
        messages.append(HumanMessage(content=user_prompt))
    response = llm.invoke(messages)
    return (getattr(response, "content", None) or str(response)).strip()


# ── Supplier matching ──────────────────────────────────────────────────────────

_SELECTED_THRESHOLD = Decimal("0.70")

def _compute_supplier_match(supplier: Optional[str], ctx: RASContext) -> tuple:
    if not supplier:
        return False, Decimal("0")
    known: set[str] = set()
    if ctx.supplier_name:   known.add(ctx.supplier_name.strip())
    if ctx.parent_supplier: known.add(ctx.parent_supplier.strip())
    for li in ctx.line_items:
        if li.supplier_name: known.add(li.supplier_name.strip())
    if not known:
        return False, Decimal("0")
    ext  = supplier.strip().lower()
    best = 0.0
    for n in known:
        nl = n.lower()
        if ext in nl or nl in ext:
            best = max(best, 0.90)
        else:
            best = max(best, SequenceMatcher(None, ext, nl).ratio())
    conf = Decimal(str(round(best, 4)))
    return conf >= _SELECTED_THRESHOLD, conf


# ── Extraction response parsing ────────────────────────────────────────────────

def _parse_extraction_response(raw: str, source: dict, ctx: RASContext) -> list:
    raw = raw.strip()
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        logger.error(f"LLM returned invalid JSON: {exc}")
        return []

    header: dict = data if isinstance(data, dict) else {}
    items_raw: list = header.get("items", [])
    if not items_raw:
        return []

    h_supplier = header.get("supplier_name") or None
    _, match_conf = _compute_supplier_match(h_supplier, ctx)

    header_fields = {
        "supplier_name":      h_supplier,
        "supplier_address":   header.get("supplier_address") or None,
        "supplier_country":   _normalize_supplier_country(header.get("supplier_country") or None),
        "quotation_ref_no":   header.get("quotation_ref_no") or None,
        "quotation_date":     header.get("quotation_date"),
        "currency":           _normalize_currency_code(header.get("currency") or None),
        "validity_date":      header.get("validity_date"),
        "validity_days":      header.get("validity_days"),
        "payment_terms":      header.get("payment_terms") or None,
        "supplier_match_conf": float(match_conf),
        "attachment_classify_fk": source["attachment_classify_fk"],
        "embedded_classify_fk":   source["embedded_classify_fk"],
        "is_selected_quote":      False,
        "quote_rank":             None,
    }

    results: list[dict] = []
    for raw_item in items_raw:
        item = dict(header_fields)
        item.update({k: (None if v == "" else v) for k, v in raw_item.items()})
        results.append(item)
    return results


# ── Item alignment to RAS line items ──────────────────────────────────────────

_IDENT_RE = re.compile(r"[A-Za-z0-9]*\d+[A-Za-z0-9]*")

def _ident_tokens(text: Optional[str]) -> set:
    if not text: return set()
    return {m.group(0).lower() for m in _IDENT_RE.finditer(text) if len(m.group(0)) >= 2}

def _data_score(item: dict) -> tuple:
    return (int(item.get("unit_price") is not None), int(item.get("total_price") is not None),
            int(item.get("item_name") is not None), int(item.get("item_description") is not None))

def _align_to_ras_line_items(items: list, ctx: RASContext, source: dict) -> list:
    valid_ids = {li.purchase_dtl_id for li in ctx.line_items}
    matched_by_dtl: dict = {}
    orphans: list[dict] = []

    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None:
            try: dtl_id = int(dtl_id)
            except Exception: dtl_id = None
        if dtl_id in valid_ids:
            if dtl_id not in matched_by_dtl:
                matched_by_dtl[dtl_id] = item
            elif _data_score(item) > _data_score(matched_by_dtl[dtl_id]):
                matched_by_dtl[dtl_id] = item
        elif any([item.get("item_name"), item.get("unit_price"), item.get("total_price")]):
            orphans.append(item)

    matched = list(matched_by_dtl.values())
    covered = {i.get("purchase_dtl_id") for i in matched}
    uncovered = [li for li in ctx.line_items if li.purchase_dtl_id not in covered]

    # Fuzzy assign orphans
    if orphans and uncovered:
        scored: list = []
        for oi, orp in enumerate(orphans):
            orp_text = f"{orp.get('item_name','')} {orp.get('item_description','')}".lower()
            for li_idx, li in enumerate(uncovered):
                ras_text = f"{li.item_description or ''} {li.item_code or ''}".lower()
                shared = _ident_tokens(orp_text) & _ident_tokens(ras_text)
                if shared:
                    score = 0.90
                elif li.item_code and li.item_code.lower() in orp_text:
                    score = 0.85
                else:
                    score = SequenceMatcher(None, orp_text, ras_text).ratio()
                try:
                    if orp.get("quantity") is not None and li.quantity is not None:
                        if Decimal(str(orp["quantity"])) == li.quantity:
                            score = min(1.0, score + 0.10)
                except Exception:
                    pass
                scored.append((score, oi, li_idx))
        scored.sort(key=lambda x: x[0], reverse=True)
        used_o: set = set(); used_l: set = set()
        for score, oi, li_idx in scored:
            if oi in used_o or li_idx in used_l: continue
            dtl_id = uncovered[li_idx].purchase_dtl_id
            orphans[oi]["purchase_dtl_id"] = dtl_id
            matched.append(orphans[oi])
            covered.add(dtl_id)
            used_o.add(oi); used_l.add(li_idx)

    # Stubs for uncovered
    still_uncovered = {li.purchase_dtl_id for li in uncovered if li.purchase_dtl_id not in covered}
    donor = matched[0] if matched else {}
    for li in ctx.line_items:
        if li.purchase_dtl_id not in still_uncovered:
            continue
        matched.append({
            "attachment_classify_fk": source["attachment_classify_fk"],
            "embedded_classify_fk":   source["embedded_classify_fk"],
            "purchase_dtl_id":        li.purchase_dtl_id,
            "is_selected_quote":      False,
            "supplier_match_conf":    0.0,
            "quote_rank":             None,
            "supplier_name":          donor.get("supplier_name"),
            "supplier_address":       donor.get("supplier_address"),
            "supplier_country":       donor.get("supplier_country"),
            "quotation_ref_no":       donor.get("quotation_ref_no"),
            "quotation_date":         donor.get("quotation_date"),
            "currency":               donor.get("currency"),
            "validity_date":          donor.get("validity_date"),
            "validity_days":          donor.get("validity_days"),
            "payment_terms":          donor.get("payment_terms"),
        })
    return matched


# ── Quote ranking + selection ──────────────────────────────────────────────────

def _compute_quote_ranks(all_items: list) -> None:
    from collections import defaultdict
    for item in all_items:
        item["quote_rank"] = None
    by_group: dict = defaultdict(list)
    for item in all_items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is None: continue
        key = ((item.get("supplier_name") or "").strip().lower() or "_unknown_", dtl_id)
        by_group[key].append(item)
    for group in by_group.values():
        group.sort(key=lambda i: (i.get("total_price") is None, float(i.get("total_price") or 0)))
        for rank, item in enumerate(group, 1):
            item["quote_rank"] = rank


def _select_best_quotes(all_items: list, ctx: RASContext) -> None:
    from collections import defaultdict
    for item in all_items:
        item["is_selected_quote"] = False
    ras_by_dtl = {li.purchase_dtl_id: li for li in ctx.line_items}
    by_dtl: dict = defaultdict(list)
    for item in all_items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None:
            by_dtl[dtl_id].append(item)

    def _price_prox(item_val, ras_val):
        try:
            if item_val is None or ras_val is None or float(ras_val) <= 0: return 0.0
            diff = abs(float(item_val) - float(ras_val)) / float(ras_val)
            if diff <= 0.05: return 1.0
            if diff >= 0.25: return 0.0
            return 1.0 - (diff - 0.05) / 0.20
        except Exception:
            return 0.0

    for dtl_id, candidates in by_dtl.items():
        ras_line = ras_by_dtl.get(dtl_id)
        def _score(it):
            conf      = float(it.get("supplier_match_conf") or 0)
            price_fit = max(
                _price_prox(it.get("unit_price"),  ras_line.unit_price if ras_line else None),
                _price_prox(it.get("total_price"), ras_line.req_value  if ras_line else None),
            )
            has_price = int(it.get("unit_price") is not None)
            return (conf, price_fit, has_price)
        max(candidates, key=_score)["is_selected_quote"] = True


# ── Currency conversion helper ─────────────────────────────────────────────────

_EUR_CUR_ID      = 3
_RATE_IS_MULTIPLY = True


def _convert_to_eur(tgt_cs: str, amount, currency_code: str | None, ref_date) -> "Decimal | None":
    """Return amount converted to EUR using EXCHANGE_RATE table.

    Returns None if amount is None or conversion data is unavailable.
    """
    if amount is None:
        return None
    try:
        from decimal import Decimal as _Dec
        amount_dec = _Dec(str(amount))
    except Exception:
        return None
    if not currency_code:
        return None
    try:
        import re as _re
        from datetime import date as _date_cls, datetime as _dt_cls
        if isinstance(ref_date, _date_cls):
            date_val = ref_date
        elif ref_date:
            m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(ref_date))
            if m:
                date_val = _date_cls(int(m[1]), int(m[2]), int(m[3]))
            else:
                date_val = _dt_cls.utcnow().date()
        else:
            from datetime import datetime as _dt_cls2
            date_val = _dt_cls2.utcnow().date()
        conn = _connect(tgt_cs)
        try:
            cur = conn.cursor()
            # look up source currency id
            cur.execute(
                "SELECT [CUR_ID] FROM [ras_procurement].[currency_mst] "
                "WHERE UPPER([CURRENCY]) = UPPER(?)",
                currency_code,
            )
            row = cur.fetchone()
            if row is None:
                logger.warning(f"Currency conversion: no currency_mst row for code={currency_code!r}")
                return None
            src_cur_id = row[0]
            if src_cur_id == _EUR_CUR_ID:
                return amount_dec
            # look up exchange rate valid for ref_date
            cur.execute(
                "SELECT TOP 1 [CONVERSION_RATE] "
                "FROM [ras_procurement].[EXCHANGE_RATE] "
                "WHERE [CUR_ID] = ? AND [BASE_CUR_ID] = ? "
                "  AND [STATUS_ID] = 10 "
                "  AND [FROM_DATE] <= ? AND [TO_DATE] >= ? "
                "ORDER BY [FROM_DATE] DESC",
                src_cur_id, _EUR_CUR_ID, date_val, date_val,
            )
            rate_row = cur.fetchone()
            if rate_row is None:
                logger.warning(
                    f"Currency conversion: no EXCHANGE_RATE row for "
                    f"FROM_CUR_ID={src_cur_id} TO_CUR_ID={_EUR_CUR_ID} date={date_val} STATUS_ID=10"
                )
                return None
            from decimal import Decimal as _Dec2
            rate = _Dec2(str(rate_row[0]))
            return amount_dec * rate if _RATE_IS_MULTIPLY else amount_dec / rate
        finally:
            conn.close()
    except Exception as exc:
        logger.warning(f"Currency conversion failed currency={currency_code!r} date={ref_date}: {exc}")
        return None


# ── Supplier name canonicalization (mirrors doc intel Union-Find logic) ────────

import re as _re_supplier

_CANONICALIZE_THRESHOLD = 0.82

# Mirrors doc intel branch — geo tokens used to prevent merging different country branches
_GEO_TOKENS: frozenset = frozenset({
    "india", "china", "japan", "usa", "us", "uk", "germany", "france",
    "italy", "korea", "taiwan", "singapore", "malaysia", "thailand",
    "vietnam", "indonesia", "australia", "canada", "brazil", "mexico",
    "uae", "dubai", "europe", "asia", "americas", "shanghai", "beijing",
    "mumbai", "delhi",
})

def _normalize_currency_code(code_or_name: str | None) -> str | None:
    """Normalize currency string to ISO-4217 alpha-3 code.

    Priority: symbol map → pycountry alpha_3 exact → pycountry name match → uppercase if 3-char.
    Gracefully degrades when pycountry is not installed.
    """
    # Defined inside the function so it's always in scope in agentcore's exec context
    _SYMBOL_MAP: dict = {
        "₹": "INR", "$": "USD", "€": "EUR", "£": "GBP", "¥": "JPY",
        "₩": "KRW", "₫": "VND", "฿": "THB", "Rp": "IDR", "RM": "MYR",
        "S$": "SGD", "A$": "AUD", "C$": "CAD", "R": "ZAR",
    }
    if not code_or_name:
        return code_or_name
    val = code_or_name.strip()
    if val in _SYMBOL_MAP:
        return _SYMBOL_MAP[val]
    upper = val.upper()
    try:
        import pycountry as _pc
        c = _pc.currencies.get(alpha_3=upper)
        if c:
            return c.alpha_3
        val_lower = val.lower()
        for c in _pc.currencies:
            if c.name.lower() == val_lower:
                return c.alpha_3
    except ImportError:
        pass
    return upper if len(upper) == 3 else val


def _normalize_supplier_country(country_str: str | None) -> str | None:
    """Normalize free-text country name to ISO 3166-1 alpha-2 code using pycountry.

    Tries exact name/code match first, then rapidfuzz fuzzy match (threshold 75).
    Falls back to original string when pycountry is not installed or no match found.
    """
    if not country_str:
        return country_str
    text_lower = country_str.strip().lower()
    try:
        import pycountry as _pc
        for c in _pc.countries:
            if text_lower in (
                c.name.lower(),
                getattr(c, "official_name", "").lower(),
                c.alpha_2.lower(),
                c.alpha_3.lower(),
            ):
                return c.alpha_2
        try:
            from rapidfuzz import process as _fuzz
            all_names = [c.name for c in _pc.countries]
            result = _fuzz.extractOne(text_lower, [n.lower() for n in all_names])
            if result and result[1] > 75:
                matched = _pc.countries.get(name=all_names[[n.lower() for n in all_names].index(result[0])])
                if matched:
                    return matched.alpha_2
        except ImportError:
            pass
    except ImportError:
        logger.warning("pycountry not installed — supplier_country not normalized (pip install pycountry rapidfuzz)")
    return country_str


def _strip_contact_suffix(name: str) -> str:
    return _re_supplier.sub(
        r'\s*[-–]\s*(contact|email|ph|phone|tel|mob)[:\s].*$',
        '', name, flags=_re_supplier.IGNORECASE,
    ).strip()


def _is_acronym_of(short: str, long_name: str) -> bool:
    """True if short is all-caps (≤ 6 chars) and matches the word initials of long_name."""
    if not short.isupper() or len(short) > 6:
        return False
    initials = "".join(m[0].upper() for m in _re_supplier.findall(r'\b[A-Za-z]', long_name))
    return short == initials[:len(short)]


def _name_geo_tokens(cleaned: str) -> frozenset:
    return frozenset(w for w in _re_supplier.findall(r'\b\w+\b', cleaned.lower()) if w in _GEO_TOKENS)


def _canonicalize_supplier_names(items: list[dict]) -> None:
    """In-memory Union-Find supplier name canonicalization — same rules as doc intel branch.

    Groups items by purchase_dtl_id, clusters name variants using four merge
    rules (exact match, substring, fuzzy ratio ≥ 0.82, acronym), with a
    geo-token guard that prevents merging different country branches.
    Mutates supplier_name in-place for non-canonical rows.
    """
    from difflib import SequenceMatcher
    from collections import defaultdict

    by_dtl: dict[int, list[dict]] = defaultdict(list)
    for item in items:
        dtl_id = item.get("purchase_dtl_id")
        if dtl_id is not None and item.get("supplier_name"):
            by_dtl[int(dtl_id)].append(item)

    for dtl_id, dtl_items in by_dtl.items():
        raw_names = list({i["supplier_name"] for i in dtl_items})
        if len(raw_names) < 2:
            continue

        uf: dict[str, str] = {n: n for n in raw_names}

        def _find(x: str) -> str:
            while uf[x] != x:
                uf[x] = uf[uf[x]]
                x = uf[x]
            return x

        def _union(a: str, b: str) -> None:
            uf[_find(a)] = _find(b)

        for i, a in enumerate(raw_names):
            a_clean = _strip_contact_suffix(a).lower()
            a_geo   = _name_geo_tokens(a_clean)
            for b in raw_names[i + 1:]:
                b_clean = _strip_contact_suffix(b).lower()
                b_geo   = _name_geo_tokens(b_clean)
                if a_geo and b_geo and a_geo != b_geo:
                    continue
                if a_clean == b_clean:
                    _union(a, b)
                elif a_clean in b_clean or b_clean in a_clean:
                    _union(a, b)
                elif SequenceMatcher(None, a_clean, b_clean).ratio() >= _CANONICALIZE_THRESHOLD:
                    _union(a, b)
                elif _is_acronym_of(a.strip(), b) or _is_acronym_of(b.strip(), a):
                    _union(a, b)

        clusters: dict[str, list[str]] = defaultdict(list)
        for n in raw_names:
            clusters[_find(n)].append(n)

        canonical_map: dict[str, str] = {}
        for members in clusters.values():
            canonical = min(members, key=len)
            for m in members:
                canonical_map[m] = canonical

        changed = 0
        for item in dtl_items:
            orig = item.get("supplier_name", "")
            canon = canonical_map.get(orig, orig)
            if canon != orig:
                item["supplier_name"] = canon
                changed += 1
        if changed:
            logger.info(f"Supplier canonicalization DTL {dtl_id}: {changed} name(s) updated")


# ── DB writer: extracted items ─────────────────────────────────────────────────

def _save_extracted_items(tgt_cs: str, items: list) -> int:
    if not items:
        return 0
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    saved = 0
    try:
        for item in items:
            def _v(k, cast=None):
                v = item.get(k)
                if v is None: return None
                try: return cast(v) if cast else v
                except Exception: return None
            def _d(k):
                v = item.get(k)
                if v is None: return None
                try: return Decimal(str(v))
                except Exception: return None
            def _date(k):
                v = item.get(k)
                if not v: return None
                try:
                    from datetime import date as date_cls
                    if isinstance(v, date_cls): return v
                    import re as _re
                    m = _re.match(r"(\d{4})-(\d{2})-(\d{2})", str(v))
                    if m: return date_cls(int(m[1]), int(m[2]), int(m[3]))
                except Exception: pass
                return None
            unit_price_eur  = _convert_to_eur(tgt_cs, item.get("unit_price"),  item.get("currency"), item.get("quotation_date"))
            total_price_eur = _convert_to_eur(tgt_cs, item.get("total_price"), item.get("currency"), item.get("quotation_date"))
            cur.execute("""
                INSERT INTO [ras_procurement].[quotation_extracted_items] (
                    [attachment_classify_fk],[embedded_classify_fk],[purchase_dtl_id],
                    [is_selected_quote],[supplier_match_conf],[quote_rank],
                    [supplier_name],[supplier_address],[supplier_country],
                    [quotation_ref_no],[quotation_date],[currency],
                    [validity_date],[validity_days],[payment_terms],
                    [item_name],[item_description],[quantity],[unit],
                    [unit_price],[total_price],[discount],[taxation_details],
                    [delivery_date],[delivery_time_days],
                    [item_level_1],[item_level_2],[item_level_3],[item_level_4],
                    [item_level_5],[item_level_6],[item_level_7],[item_level_8],
                    [commodity_tag],[item_summary],
                    [unit_price_eur],[total_price_eur]
                ) VALUES (
                    ?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?
                )
            """,
                _v("attachment_classify_fk"), _v("embedded_classify_fk"),
                _v("purchase_dtl_id", int), bool(item.get("is_selected_quote")),
                _d("supplier_match_conf"), _v("quote_rank", int),
                _v("supplier_name"), _v("supplier_address"), _v("supplier_country"),
                _v("quotation_ref_no"), _date("quotation_date"), _v("currency"),
                _date("validity_date"), _v("validity_days", int), _v("payment_terms"),
                _v("item_name"), _v("item_description"), _d("quantity"), _v("unit"),
                _d("unit_price"), _d("total_price"), _d("discount"), _v("taxation_details"),
                _date("delivery_date"), _v("delivery_time_days", int),
                _v("item_level_1"), _v("item_level_2"), _v("item_level_3"), _v("item_level_4"),
                _v("item_level_5"), _v("item_level_6"), _v("item_level_7"), _v("item_level_8"),
                _v("commodity_tag"), _v("item_summary"),
                unit_price_eur, total_price_eur,
            )
            saved += 1
        conn.commit()
    except Exception as exc:
        try: conn.rollback()
        except Exception: pass
        raise
    finally:
        conn.close()
    return saved


# ── Stage 5: run extraction for a PR ──────────────────────────────────────────

def _run_extraction(llm, tgt_cs: str, blob_cfg: dict, pr_no: str, prompts: dict | None = None) -> int:
    ctx = _build_ras_context(tgt_cs, pr_no)
    if ctx is None:
        logger.warning(f"[{pr_no}] No RAS context found — skipping extraction")
        return 0
    if not ctx.line_items:
        logger.warning(f"[{pr_no}] No line items — skipping extraction")
        return 0

    sources = _resolve_quotation_sources(tgt_cs, pr_no)
    if not sources:
        logger.warning(f"[{pr_no}] No Quotation attachments found — check classification")
        return 0

    all_items: list[dict] = []
    for src in sources:
        import os
        blob_path = src["blob_path"]
        filename  = os.path.basename(blob_path)
        try:
            file_bytes = _download_blob(blob_path, blob_cfg)
            doc        = _load_document(file_bytes, filename)
            raw        = _call_extraction_llm(llm, ctx, doc, prompts)
            items      = _parse_extraction_response(raw, src, ctx)
            items      = _align_to_ras_line_items(items, ctx, src)
            all_items.extend(items)
            logger.info(f"[{pr_no}] Extracted {len(items)} item(s) from {filename!r}")
        except Exception as exc:
            logger.warning(f"[{pr_no}] Extraction failed for {filename!r}: {exc}")
            continue

    if not all_items:
        return 0

    _compute_quote_ranks(all_items)
    _select_best_quotes(all_items, ctx)
    _canonicalize_supplier_names(all_items)
    saved = _save_extracted_items(tgt_cs, all_items)
    logger.info(f"[{pr_no}] {saved} item(s) written to quotation_extracted_items")
    return saved


# ── Tracker helpers ────────────────────────────────────────────────────────────

def _advance_tracker(tgt_cs: str, pr_no: str, stage_id: int) -> None:
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            UPDATE [ras_procurement].[ras_tracker]
               SET current_stage_fk=?, updated_at=SYSUTCDATETIME()
             WHERE purchase_req_no=?
        """, stage_id, pr_no)
        conn.commit()
    except Exception as exc:
        logger.warning(f"Tracker advance failed PR={pr_no!r} stage={stage_id}: {exc}")
    finally:
        conn.close()


def _record_exception(tgt_cs: str, pr_no: str, stage_id: int, error_msg: str) -> None:
    try:
        conn = _connect(tgt_cs)
        cur  = conn.cursor()
        try:
            cur.execute("""
                UPDATE [ras_procurement].[ras_tracker]
                   SET current_stage_fk=99, updated_at=SYSUTCDATETIME()
                 WHERE purchase_req_no=?
            """, pr_no)
            cur.execute("SELECT ras_uuid_pk FROM [ras_procurement].[ras_tracker] WHERE purchase_req_no=?", pr_no)
            row = cur.fetchone()
            if row:
                cur.execute("""
                    INSERT INTO [ras_procurement].[ras_pipeline_exceptions]
                        (ras_tracker_id, stage_id, exception_message)
                    VALUES (?, ?, ?)
                """, row[0], stage_id, error_msg[:4000])
            conn.commit()
        finally:
            conn.close()
    except Exception as exc:
        logger.warning(f"[{pr_no}] Could not write exception record: {exc}")


# ── Embeddings + benchmark (Stage 6-7) ────────────────────────────────────────

def _run_embeddings(tgt_cs: str, pr_no: str, embed_model, pinecone_index: str, pinecone_ns: str) -> None:
    # Structural errors (import, DB connect, index creation) propagate — caller records exception.
    from agentcore.services.pinecone_service_client import ensure_index_via_service, ingest_via_service
    ensure_index_via_service(index_name=pinecone_index, embedding_dimension=3072)
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description], qi.[item_summary],
                   qi.[item_level_1], qi.[item_level_2], qi.[item_level_3],
                   qi.[item_level_4], qi.[item_level_5], qi.[item_level_6],
                   qi.[item_level_7], qi.[item_level_8], qi.[commodity_tag]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        all_rows = cur.fetchall()
    finally:
        conn.close()

    # deduplicate by purchase_dtl_id — first row wins (same as doc intel branch)
    seen_dtl: dict = {}
    for row in all_rows:
        dtl_id = row[1]
        if dtl_id not in seen_dtl:
            seen_dtl[dtl_id] = row
    rows = list(seen_dtl.values())
    logger.info(f"[{pr_no}] Embedding {len(rows)} selected item(s) (deduped from {len(all_rows)} is_selected_quote=1 rows)")

    _EMBED_FIELDS_ORDER = [
        "item_name", "item_description", "item_summary",
        "item_level_1", "item_level_2", "item_level_3", "item_level_4",
        "item_level_5", "item_level_6", "item_level_7", "item_level_8",
        "commodity_tag",
    ]
    for row in rows:
        item_uuid, dtl_id, item_name, item_desc, item_summary, *rest = row
        levels = rest[:8]
        commodity_tag = rest[8] if len(rest) > 8 else None
        field_vals = {
            "item_name":        item_name,
            "item_description": item_desc,
            "item_summary":     item_summary,
            "item_level_1":     levels[0], "item_level_2": levels[1],
            "item_level_3":     levels[2], "item_level_4": levels[3],
            "item_level_5":     levels[4], "item_level_6": levels[5],
            "item_level_7":     levels[6], "item_level_8": levels[7],
            "commodity_tag":    commodity_tag,
        }
        content = " | ".join(str(field_vals[f]) for f in _EMBED_FIELDS_ORDER if field_vals.get(f))
        if not content:
            continue
        try:
            embedding = embed_model.embed_query(content)
            ingest_via_service(
                index_name=pinecone_index,
                namespace=pinecone_ns,
                text_key="page_content",
                documents=[{
                    "page_content": content,
                    "purchase_req_no": pr_no,
                    "purchase_dtl_id": str(dtl_id or ""),
                }],
                embedding_vectors=[embedding],
                vector_ids=[f"dtl_{dtl_id}"],
                embedding_dimension=3072,
            )
        except Exception as exc:
            logger.warning(f"[{pr_no}] Embedding failed for dtl_id {dtl_id}: {exc}")


def _run_benchmark(llm, tgt_cs: str, pr_no: str, embed_model, pinecone_index: str, pinecone_ns: str, top_k: int) -> None:
    # Structural errors (import, DB connect, query) propagate — caller records exception.
    from agentcore.services.pinecone_service_client import search_via_service
    from langchain_core.messages import HumanMessage
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT qi.[extracted_item_uuid_pk], qi.[purchase_dtl_id],
                   qi.[item_name], qi.[item_description],
                   qi.[unit_price], qi.[quantity]
              FROM [ras_procurement].[quotation_extracted_items] qi
              JOIN [ras_procurement].[attachment_classification] ac
                ON qi.[attachment_classify_fk] = ac.[attachment_classify_uuid_pk]
              JOIN [ras_procurement].[ras_tracker] rt
                ON ac.[ras_uuid_pk] = rt.[ras_uuid_pk]
             WHERE rt.[purchase_req_no] = ?
               AND qi.[is_selected_quote] = 1
               AND qi.[purchase_dtl_id] IS NOT NULL
        """, pr_no)
        items = cur.fetchall()
    finally:
        conn.close()

    conn2 = _connect(tgt_cs)
    cur2  = conn2.cursor()
    try:
        for row in items:
            item_uuid, dtl_id, name, desc, unit_price, qty = row
            item_text = f"{name or ''} {desc or ''}".strip()
            if not item_text: continue
            try:
                embedding = embed_model.embed_query(item_text[:500])
                raw_similar = search_via_service(
                    index_name=pinecone_index,
                    namespace=pinecone_ns,
                    text_key="content",
                    query=item_text[:500],
                    query_embedding=embedding,
                    number_of_results=top_k,
                )
                # normalise to list — service may return list or {"results": [...]}
                if isinstance(raw_similar, list):
                    similar = raw_similar
                elif isinstance(raw_similar, dict):
                    similar = raw_similar.get("results", raw_similar.get("matches", []))
                else:
                    similar = []
                # exclude vectors from the same PR (no server-side filter in sync API)
                similar = [s for s in similar
                           if (s.get("metadata") or {}).get("purchase_req_no") != pr_no]
            except Exception as exc:
                logger.warning(f"[{pr_no}] Benchmark similarity failed dtl_id={dtl_id}: {exc}")
                continue
            if not similar: continue
            bench_prompt = (
                f"Recommend a benchmark unit price for this item based on historical similar purchases.\n"
                f"Item: {item_text}\nCurrent unit price: {unit_price}\nQuantity: {qty}\n"
                f"Historical similar items (top {min(3, len(similar))}):\n"
                f"{json.dumps(similar[:3], indent=2)}\n\n"
                f'Return ONLY JSON: {{"bp_unit_price": <number or null>, "summary": "<2-3 sentences>"}}'
            )
            try:
                resp = llm.invoke([HumanMessage(content=bench_prompt)])
                raw  = (getattr(resp, "content", None) or str(resp)).strip()
                raw  = re.sub(r"^```(?:json)?\s*", "", raw)
                raw  = re.sub(r"\s*```$", "", raw)
                bout = json.loads(raw)
            except Exception:
                bout = {}
            bp_unit = bout.get("bp_unit_price")
            summary = bout.get("summary", "")
            try:
                bp_dec   = Decimal(str(bp_unit)) if bp_unit is not None else None
                bp_total = round(float(bp_dec) * float(qty or 1), 2) if bp_dec is not None else None
            except Exception:
                bp_dec = bp_total = None
            try:
                cur2.execute("""
                    MERGE [ras_procurement].[benchmark_result] WITH (HOLDLOCK) AS target
                    USING (SELECT ? AS purchase_dtl_id) AS src
                       ON target.purchase_dtl_id = src.purchase_dtl_id
                    WHEN MATCHED THEN
                        UPDATE SET extracted_item_uuid_fk=?, bp_unit_price=?, bp_total_price=?,
                                   summary=?, updated_at=SYSUTCDATETIME()
                    WHEN NOT MATCHED THEN
                        INSERT (purchase_dtl_id, extracted_item_uuid_fk, bp_unit_price, bp_total_price, summary)
                        VALUES (?, ?, ?, ?, ?);
                """, dtl_id,
                     item_uuid, bp_dec, bp_total, summary,
                     dtl_id, item_uuid, bp_dec, bp_total, summary)
            except Exception as exc:
                logger.warning(f"Benchmark write failed dtl_id={dtl_id}: {exc}")
        conn2.commit()
    except Exception:
        try: conn2.rollback()
        except Exception: pass
        raise
    finally:
        conn2.close()


# ── Fetch pending PRs ──────────────────────────────────────────────────────────

def _fetch_pending_prs(tgt_cs: str, pr_filter: str, batch_limit: int) -> list:
    if pr_filter:
        return [pr_filter]
    conn = _connect(tgt_cs)
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT TOP (?) prm.PURCHASE_REQ_NO
              FROM [ras_procurement].[purchase_req_mst] prm
              JOIN [ras_procurement].[ras_tracker] rt
                ON prm.PURCHASE_REQ_NO = rt.purchase_req_no
             WHERE rt.current_stage_fk = 3
               AND UPPER(prm.PURCHASEFINALAPPROVALSTATUS)
                       IN ('APPROVED BY ALL', 'APPROVED BY ALL EXCEPTION')
             ORDER BY prm.C_DATETIME ASC
        """, batch_limit)
        return [row[0] for row in cur.fetchall()]
    finally:
        conn.close()


# ── Main Component ─────────────────────────────────────────────────────────────

class PipelineStage4567Node(Node):
    display_name = "Pipeline Stage 4-8"
    description  = (
        "Stage 4: Classify all attachments (parent + embedded) using the full "
        "file_classifier logic — downloads files from blob, extracts content, "
        "calls LLM with field-based SYSTEM_PROMPT. "
        "Stage 5: LLM extraction for Quotation files — builds RAS context from DB, "
        "aligns items to line items, writes quotation_extracted_items. "
        "Stage 6: Embeddings. Stage 7: Benchmark. Stage 8: Complete."
    )
    icon = "Cpu"
    name = "PipelineStage4567Node"

    inputs = [
        HandleInput(
            name="stage123_result",
            display_name="Stage 1-3 Processed PRs",
            input_types=["Data"],
            required=False,
            info="Wire the 'Processed PRs' Data output of the Stage 1-3 component here. "
                 "Stage 4-8 will process exactly those PR numbers synchronously. "
                 "If left disconnected, Stage 4-8 queries the DB for all stage=3 PRs.",
        ),
        HandleInput(
            name="target_connection",
            display_name="Target DB — Azure SQL",
            input_types=["Data"],
            info="Connection Config from the Azure SQL Database Connector node.",
        ),
        HandleInput(
            name="llm",
            display_name="LLM (GPT-4o or GPT-4o-mini)",
            input_types=["LanguageModel"],
        ),
        HandleInput(
            name="embed_model",
            display_name="Embeddings Model",
            input_types=["Embeddings"],
        ),
        # ── Optional prompt overrides ──────────────────────────────────────────
        # Wire a Prompt Template node to any of these to override the built-in prompts.
        # Leave disconnected to use the default prompts baked into this component.
        HandleInput(
            name="cls_system_prompt",
            display_name="[Classification] System Prompt",
            input_types=["Message"],
            required=False,
            info="Override the classification system prompt. Connect a Prompt Template node.",
        ),
        HandleInput(
            name="cls_user_text_prompt",
            display_name="[Classification] User Prompt — Text/Tabular Files",
            input_types=["Message"],
            required=False,
            info="Override the user prompt for text/Excel/Word/PDF files. Must keep {filename}, {file_type}, {extra_metadata}, {extracted_content} placeholders.",
        ),
        HandleInput(
            name="cls_user_image_prompt",
            display_name="[Classification] User Prompt — Image/Scanned Files",
            input_types=["Message"],
            required=False,
            info="Override the user prompt sent with base64 image content. Must keep {filename}, {file_type}, {extra_metadata} placeholders.",
        ),
        HandleInput(
            name="ext_system_prompt",
            display_name="[Extraction] System Prompt",
            input_types=["Message"],
            required=False,
            info="Override the extraction system prompt. Connect a Prompt Template node.",
        ),
        HandleInput(
            name="ext_user_template",
            display_name="[Extraction] User Prompt Template",
            input_types=["Message"],
            required=False,
            info="Override the extraction user template. Must keep all {field} placeholders from the RAS context and {document_content}.",
        ),
        HandleInput(
            name="ext_item_taxonomy",
            display_name="[Extraction] Item Taxonomy",
            input_types=["Message"],
            required=False,
            info="Override the item taxonomy guidelines injected into {item_taxonomy} of the extraction user template.",
        ),
        MessageTextInput(
            name="blob_connector_name",
            display_name="Azure Blob Connector Name",
            value="",
            info="Exact name of your Azure Blob connector in Settings → Connectors.",
        ),
        MessageTextInput(
            name="pr_no_filter",
            display_name="PR Number Filter (optional)",
            value="",
            advanced=True,
            info="Leave blank to process all PRs at stage 3. Enter one PR to restrict.",
        ),
        IntInput(name="batch_limit",      display_name="Max PRs per Run",             value=50,  advanced=True),
        MessageTextInput(name="pinecone_index",    display_name="Pinecone Index Name",  value="ras-quotations", advanced=True),
        MessageTextInput(name="pinecone_namespace", display_name="Pinecone Namespace",  value="procurement",    advanced=True),
        IntInput(name="pinecone_top_k",   display_name="Benchmark Similarity Top-K",   value=5,  advanced=True),
    ]

    outputs = [
        Output(display_name="Sync Status", name="sync_status", method="run_pipeline", types=["Message"]),
    ]

    def run_pipeline(self) -> Message:
        tgt_cs    = _conn_str(self.target_connection)
        pr_filter = (self.pr_no_filter or "").strip()
        blob_cfg  = _get_blob_config_by_name(self.blob_connector_name)

        # If Stage 1-3 is wired, use its exact PR numbers (synchronous per-PR flow).
        # Otherwise fall back to querying the DB for all stage=3 PRs.
        stage123 = getattr(self, "stage123_result", None)
        if stage123 is not None and isinstance(stage123, Data):
            pr_numbers_from_123 = (stage123.data or {}).get("pr_numbers", [])
            if pr_filter:
                pr_numbers_from_123 = [p for p in pr_numbers_from_123 if p == pr_filter]
            pr_list = pr_numbers_from_123
        else:
            pr_list = _fetch_pending_prs(tgt_cs, pr_filter, int(self.batch_limit))

        # Build prompt overrides from wired Prompt Template nodes (None = use default)
        prompts: dict = {}
        _p = getattr(self, "cls_system_prompt",   None)
        if _p: prompts["cls_system"]    = _get_prompt_text(_p, CLASSIFICATION_SYSTEM_PROMPT)
        _p = getattr(self, "cls_user_text_prompt", None)
        if _p: prompts["cls_user_text"] = _get_prompt_text(_p, _CLASSIFY_USER_TEXT)
        _p = getattr(self, "cls_user_image_prompt", None)
        if _p: prompts["cls_user_image"] = _get_prompt_text(_p, _CLASSIFY_USER_IMAGE)
        _p = getattr(self, "ext_system_prompt",   None)
        if _p: prompts["ext_system"]    = _get_prompt_text(_p, EXTRACTION_SYSTEM_PROMPT)
        _p = getattr(self, "ext_user_template",   None)
        if _p: prompts["ext_user"]      = _get_prompt_text(_p, EXTRACTION_USER_TEMPLATE)
        _p = getattr(self, "ext_item_taxonomy",   None)
        if _p: prompts["ext_taxonomy"]  = _get_prompt_text(_p, ITEM_TAXONOMY)

        if not pr_list:
            return Message(text="No PRs at stage 3 to process.")

        self.log(f"Processing {len(pr_list)} PR(s)…")
        ok_lines: list[str] = []
        err_lines: list[str] = []

        for pr_no in pr_list:
            current_stage = _STAGE_CLASSIFICATION
            try:
                # Stage 4 — Classification
                self.log(f"[{pr_no}] Stage 4 — classifying attachments…")
                _run_classification(self.llm, tgt_cs, blob_cfg, pr_no, prompts)
                _advance_tracker(tgt_cs, pr_no, _STAGE_CLASSIFICATION)
                self.log(f"[{pr_no}] Stage 4 — classification complete")

                # Stage 5 — Extraction
                current_stage = _STAGE_EXTRACTION
                self.log(f"[{pr_no}] Stage 5 — extracting quotation items…")
                n_items = _run_extraction(self.llm, tgt_cs, blob_cfg, pr_no, prompts)
                _advance_tracker(tgt_cs, pr_no, _STAGE_EXTRACTION)
                self.log(f"[{pr_no}] Stage 5 — {n_items} item(s) extracted")

                # Stage 6 — Embeddings
                current_stage = _STAGE_EMBEDDINGS
                _run_embeddings(tgt_cs, pr_no, self.embed_model,
                                self.pinecone_index, self.pinecone_namespace)
                _advance_tracker(tgt_cs, pr_no, _STAGE_EMBEDDINGS)
                self.log(f"[{pr_no}] Stage 6 — embeddings done")

                # Stage 7 — Benchmark
                current_stage = _STAGE_PRICE_BENCHMARK
                _run_benchmark(self.llm, tgt_cs, pr_no, self.embed_model,
                               self.pinecone_index, self.pinecone_namespace, int(self.pinecone_top_k))
                _advance_tracker(tgt_cs, pr_no, _STAGE_PRICE_BENCHMARK)
                self.log(f"[{pr_no}] Stage 7 — benchmark done")

                # Stage 8 — Complete
                _advance_tracker(tgt_cs, pr_no, _STAGE_COMPLETE)
                self.log(f"[{pr_no}] Stage 8 — complete")

                ok_lines.append(f"  OK  {pr_no}: {n_items} item(s)")

            except Exception as exc:
                logger.opt(exception=True).error("[{}] Stage 4-8 failed at stage {}: {}", pr_no, current_stage, exc)
                _record_exception(tgt_cs, pr_no, current_stage, str(exc))
                err_lines.append(f"  ERR {pr_no}: {exc}")

        summary = f"Pipeline Stage 4-8 complete — {len(ok_lines)} OK, {len(err_lines)} errors\n"
        return Message(text=summary + "\n".join(ok_lines + err_lines))
